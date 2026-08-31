import os
import smtplib
from email.message import EmailMessage
import json
import secrets
import unicodedata
import re
import zipfile
import io
import base64
import uuid
import mimetypes
import hashlib
import gzip
from difflib import SequenceMatcher
from pathlib import Path
from datetime import datetime, timedelta, date
from zoneinfo import ZoneInfo
import time
import tempfile
import shutil
import threading
import html as html_lib
from functools import wraps

from flask import Flask, has_request_context, render_template, request, redirect, url_for, session, jsonify, flash, send_from_directory, Response, send_file, make_response, g, abort
from flask_sqlalchemy import SQLAlchemy
from sqlalchemy import UniqueConstraint, Index, func, case, text, and_, event
from sqlalchemy.exc import IntegrityError
from werkzeug.security import generate_password_hash, check_password_hash
from werkzeug.utils import secure_filename
from botocore.exceptions import ClientError, BotoCoreError
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Font, PatternFill, Alignment
from openpyxl.utils import get_column_letter

BASE_DIR = Path(__file__).resolve().parent
DATA_DIR = BASE_DIR / "data"
STATIC_DIR = BASE_DIR / "static"
BASE_DATA_VERSION = "1408-5"
APP_RELEASE = "V71.3"
DASHBOARD_RELEASE = APP_RELEASE
TEAMS_RELEASE = APP_RELEASE
FIELD_NEARBY_RADIUS_M = int(os.getenv("FIELD_NEARBY_RADIUS_M", "3000"))
FIELD_GPS_GOOD_ACCURACY_M = float(os.getenv("FIELD_GPS_GOOD_ACCURACY_M", "30"))
FIELD_GPS_MAX_ACCURACY_M = float(os.getenv("FIELD_GPS_MAX_ACCURACY_M", "80"))
_GPS_LAST_RETENTION_CLEANUP = 0.0
# Denominadores executivos oficiais informados para o parque contratado.
OFFICIAL_PARK = {
    "ATM": 590,
    "POS": 972,
    "VALIDADOR": 629,  # Recarga
    "BLOQUEIO": 1610,
}
OFFICIAL_PARK_TOTAL = sum(OFFICIAL_PARK.values())  # 3.801
TECHNICAL_TDI_TOTAL = int(os.getenv("TECHNICAL_TDI_TOTAL", "80"))
EXPECTED_CACHE_TTL_SECONDS = 600
_expected_cache = {"at": 0.0, "data": None}
# V56-D: cache curtíssimo da visão consolidada de localidades para reduzir recomputações concorrentes.
_LOCATIONS_API_CACHE = {"light": {"at": 0.0, "payload": None}, "observed": {"at": 0.0, "payload": None}}
_LOCATIONS_API_CACHE_TTL = int(os.getenv("LOCATIONS_API_CACHE_TTL", "900"))
_LOCATIONS_API_CACHE_LOCK = threading.Lock()

# V63 CORE 2.0 — parâmetros de desempenho não destrutivos.
V63_JSON_GZIP_MIN_BYTES = int(os.getenv("JSON_GZIP_MIN_BYTES", "16384"))
V63_EMV_CACHE_TTL = int(os.getenv("EMV_API_CACHE_TTL", "60"))
V63_GARAGE_CACHE_TTL = int(os.getenv("GARAGE_API_CACHE_TTL", "60"))
_V63_EMV_CACHE = {"full": {"at": 0.0, "payload": None}, "slim": {"at": 0.0, "payload": None}}
_V63_EMV_CACHE_LOCK = threading.Lock()
_V63_GARAGE_CACHE = {"at": 0.0, "payload": None}
_V63_GARAGE_CACHE_LOCK = threading.Lock()

def _v63_invalidate_emv_cache():
    with _V63_EMV_CACHE_LOCK:
        for slot in ("full", "slim"):
            _V63_EMV_CACHE[slot] = {"at": 0.0, "payload": None}

def _v63_invalidate_garage_cache():
    with _V63_GARAGE_CACHE_LOCK:
        _V63_GARAGE_CACHE["at"] = 0.0; _V63_GARAGE_CACHE["payload"] = None
# V59: Operação 2.0 — consolidação da escala, histórico operacional e leitura em lote.
_TEAM_PROFILE_SYNC_STATE = {"at": 0.0}
_TEAM_PROFILE_SYNC_TTL = int(os.getenv("TEAM_PROFILE_SYNC_TTL", "60"))
_TEAM_PROFILE_SYNC_LOCK = threading.Lock()
_FIN_TERMINALS_CACHE = {"at": 0.0, "payload": None}
_FIN_TERMINALS_CACHE_TTL = int(os.getenv("FIN_TERMINALS_CACHE_TTL", "900"))
FIN_COST_CENTERS = [
    {"key":"SUPORTE_CAMPO","id":"CVD0011","label":"SUPORTE E CAMPO"},
    {"key":"ASSISTENCIA_TECNICA","id":"CFD0024","label":"ASSISTENCIA TECNICA"},
    {"key":"IMPLANTACAO_HARDWARE","id":"CVD0016","label":"IMPLANTAÇÃO HW"},
    {"key":"MECANICA","id":"CFD0025","label":"MECANICA"},
    {"key":"ENGENHARIA_HW","id":"CVD0017","label":"ENGENHARIA HW"},
    {"key":"LINHA_17","id":"CVD0020","label":"LINHA 17"},
]
FIN_COST_CENTER_BY_KEY = {x["key"]:x for x in FIN_COST_CENTERS}
FIN_COST_CENTER_BY_ID = {x["id"]:x for x in FIN_COST_CENTERS}

UPLOAD_DIR = BASE_DIR / "uploads"
UPLOAD_DIR.mkdir(exist_ok=True)

# Render fornece DATABASE_URL. Para teste local, usa SQLite automaticamente.
database_url = os.environ.get("DATABASE_URL", "").strip()
if database_url.startswith("postgres://"):
    database_url = database_url.replace("postgres://", "postgresql://", 1)
if not database_url:
    database_url = f"sqlite:///{BASE_DIR / 'inventario_local.db'}"

app = Flask(__name__)
app.secret_key = os.environ.get("INVENTARIO_SECRET_KEY", "chave-local-apenas-para-desenvolvimento")

# V52.8: acompanhamento de importações TopDesk em background.
# O workspace Hobby normalmente roda uma única instância; o estado é apenas operacional
# e não substitui o registro persistente TopDeskImportBatch concluído.
TOPDESK_IMPORT_JOBS = {}
TOPDESK_IMPORT_LOCK = threading.Lock()

# V56-B REV2 — importação financeira desacoplada da requisição HTTP.
FIN_IMPORT_JOBS = {}
FIN_IMPORT_LOCK = threading.Lock()
FIN_IMPORT_DIR = UPLOAD_DIR / "finance_import_jobs"
FIN_IMPORT_DIR.mkdir(exist_ok=True)

# V62 REV6 — exportações pesadas fora da requisição principal.
PANORAMA_EXPORT_JOBS = {}
PANORAMA_EXPORT_LOCK = threading.Lock()
PANORAMA_EXPORT_DIR = Path(tempfile.gettempdir()) / "inventario_panorama_exports"
PANORAMA_EXPORT_DIR.mkdir(parents=True, exist_ok=True)
PANORAMA_EXPORT_MAX_AGE_SECONDS = int(os.getenv("PANORAMA_EXPORT_MAX_AGE_SECONDS", "7200"))
_STORAGE_CACHE = {"at": 0.0, "data": None}
_STORAGE_CACHE_TTL = int(os.getenv("STORAGE_TELEMETRY_TTL", "900"))
_STORAGE_CACHE_LOCK = threading.Lock()

def _fin_job_update(job_id, **changes):
    with FIN_IMPORT_LOCK:
        job=FIN_IMPORT_JOBS.get(job_id)
        if job:
            job.update(changes)
            job["updated_at"]=datetime.utcnow().isoformat()+"Z"

def _fin_job_snapshot(job_id):
    with FIN_IMPORT_LOCK:
        job=FIN_IMPORT_JOBS.get(job_id)
        return dict(job) if job else None

def _td_job_update(job_id, **changes):
    with TOPDESK_IMPORT_LOCK:
        job = TOPDESK_IMPORT_JOBS.setdefault(job_id, {})
        job.update(changes)
        job["updated_at"] = datetime.utcnow().isoformat() + "Z"

def _td_job_snapshot(job_id):
    with TOPDESK_IMPORT_LOCK:
        job = TOPDESK_IMPORT_JOBS.get(job_id)
        return dict(job) if job else None

# V55: cache curto de analytics TopDesk para evitar recalcular dezenas de milhares
# de chamados a cada repaint/filtro idêntico. É invalidado ao fim de cada importação.
TOPDESK_ANALYTICS_CACHE = {}
TOPDESK_ANALYTICS_CACHE_LOCK = threading.Lock()
TOPDESK_ANALYTICS_TTL = int(os.getenv("TOPDESK_ANALYTICS_TTL", "600"))

# V56-A.1: estado operacional do backfill TopDesk.
# Deve existir antes das rotas, migrações e do startup que o consultam.
_V56A_BACKFILL = {"running": False, "processed": 0, "error": None}
_V56A_BACKFILL_LOCK = threading.Lock()

def _td_cache_get(key):
    now=time.time()
    with TOPDESK_ANALYTICS_CACHE_LOCK:
        row=TOPDESK_ANALYTICS_CACHE.get(key)
        if not row or now-row[0] > TOPDESK_ANALYTICS_TTL:
            if row: TOPDESK_ANALYTICS_CACHE.pop(key,None)
            return None
        return row[1]

def _td_cache_put(key, payload):
    with TOPDESK_ANALYTICS_CACHE_LOCK:
        if len(TOPDESK_ANALYTICS_CACHE) >= 64:
            oldest=min(TOPDESK_ANALYTICS_CACHE.items(),key=lambda kv:kv[1][0])[0]
            TOPDESK_ANALYTICS_CACHE.pop(oldest,None)
        TOPDESK_ANALYTICS_CACHE[key]=(time.time(),payload)

def _td_cache_clear():
    with TOPDESK_ANALYTICS_CACHE_LOCK:
        TOPDESK_ANALYTICS_CACHE.clear()
app.config["SQLALCHEMY_DATABASE_URI"] = database_url
app.config["SQLALCHEMY_TRACK_MODIFICATIONS"] = False
app.config["MAX_CONTENT_LENGTH"] = 160 * 1024 * 1024
app.config["PERMANENT_SESSION_LIFETIME"] = timedelta(hours=12)
app.config["SESSION_COOKIE_HTTPONLY"] = True
app.config["SESSION_COOKIE_SAMESITE"] = "Lax"
app.config["SESSION_COOKIE_SECURE"] = os.getenv("RENDER", "").lower() in ("true","1","yes")

# Mantém conexões saudáveis em hospedagens gerenciadas.
app.config["SQLALCHEMY_ENGINE_OPTIONS"] = {
    "pool_pre_ping": True,
    "pool_recycle": 280,
}

db = SQLAlchemy(app)

# V56-B REV — mede custo SQL por requisição sem gravar payloads/dados sensíveis.
# IMPORTANTE: db.engine exige application context no Flask-SQLAlchemy 3.x.
# Os listeners são registrados uma única vez dentro do contexto para não impedir o boot do Gunicorn.
def _perf_sql_before(conn, cursor, statement, parameters, context, executemany):
    context._v56_sql_started = time.perf_counter()

def _perf_sql_after(conn, cursor, statement, parameters, context, executemany):
    try:
        started=getattr(context,"_v56_sql_started",None)
        if started is not None and has_request_context():
            g._perf_sql_ms=float(getattr(g,"_perf_sql_ms",0) or 0)+(time.perf_counter()-started)*1000.0
            g._perf_query_count=int(getattr(g,"_perf_query_count",0) or 0)+1
    except Exception:
        pass

with app.app_context():
    _perf_engine = db.engine
    if not event.contains(_perf_engine, "before_cursor_execute", _perf_sql_before):
        event.listen(_perf_engine, "before_cursor_execute", _perf_sql_before)
    if not event.contains(_perf_engine, "after_cursor_execute", _perf_sql_after):
        event.listen(_perf_engine, "after_cursor_execute", _perf_sql_after)

# V56-B — telemetria leve. Não registra arquivos estáticos nem a própria API de telemetria.
@app.before_request
def _v56b_perf_start():
    g._perf_started = time.perf_counter()
    g._perf_sql_ms = 0.0
    g._perf_query_count = 0

@app.after_request
def _v56b_perf_finish(response):
    try:
        path=request.path or ""
        if path.startswith(("/static/","/uploads/","/api/telemetria")) or path in ("/sw.js","/favicon.ico"):
            return response
        started=getattr(g,"_perf_started",None)
        if started is not None and 'PerformanceMetric' in globals():
            ms=(time.perf_counter()-started)*1000.0
            # amostragem integral de rotas lentas/erros e 1/4 das rápidas para reduzir overhead.
            keep = ms >= 750 or response.status_code >= 400 or (int(time.time()*1000) % 4 == 0)
            if keep:
                db.session.add(PerformanceMetric(route=(request.url_rule.rule if request.url_rule else path)[:220],method=request.method,status_code=response.status_code,duration_ms=round(ms,2),sql_ms=round(float(getattr(g,"_perf_sql_ms",0) or 0),2),query_count=int(getattr(g,"_perf_query_count",0) or 0),user_id=session.get("user_id")))
                db.session.commit()
    except Exception:
        try: db.session.rollback()
        except Exception: pass
    # V63: Server-Timing para diagnóstico e gzip apenas em JSON grande.
    try:
        total_ms=(time.perf_counter()-getattr(g,"_perf_started",time.perf_counter()))*1000.0
        response.headers["Server-Timing"] = f"app;dur={total_ms:.1f}, sql;dur={float(getattr(g,'_perf_sql_ms',0) or 0):.1f}"
        response.headers["X-Autopass-Release"] = APP_RELEASE
        ae=(request.headers.get("Accept-Encoding") or "").lower(); ct=(response.headers.get("Content-Type") or "").lower()
        if "gzip" in ae and response.status_code==200 and not response.direct_passthrough and ("application/json" in ct or "text/json" in ct):
            raw=response.get_data()
            if len(raw)>=V63_JSON_GZIP_MIN_BYTES and not response.headers.get("Content-Encoding"):
                packed=gzip.compress(raw,compresslevel=5)
                if len(packed)<len(raw)*0.92:
                    response.set_data(packed); response.headers["Content-Encoding"]="gzip"; response.headers["Content-Length"]=str(len(packed)); response.headers["Vary"]="Accept-Encoding"; response.headers["X-Autopass-Compressed"]="gzip"
    except Exception:
        pass
    return response


class User(db.Model):
    __tablename__ = "users"
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(180), nullable=False)
    username = db.Column(db.String(120), nullable=False, unique=True, index=True)
    password_hash = db.Column(db.String(255), nullable=False)
    role = db.Column(db.String(30), nullable=False)
    active = db.Column(db.Boolean, nullable=False, default=True)
    
    user_code = db.Column(db.String(30), unique=True, index=True)
    email = db.Column(db.String(180), unique=True, index=True)
    phone = db.Column(db.String(30), unique=True, index=True)
    photo_url = db.Column(db.String(500))
    archived_at = db.Column(db.DateTime)
    company = db.Column(db.String(180))
    work_schedule_type = db.Column(db.String(30))
    work_shift = db.Column(db.String(30))
    work_anchor_date = db.Column(db.Date)
    work_anchor_status = db.Column(db.String(20))
    job_title = db.Column(db.String(120))
    personnel_status = db.Column(db.String(30), nullable=False, default="ATIVO")
    personnel_status_note = db.Column(db.String(240))
    access_json = db.Column(db.Text)
    gps_required = db.Column(db.Boolean, nullable=False, default=False)
    customer_company_ids = db.Column(db.Text)  # JSON: empresas liberadas para perfil Cliente
    system_profile_id = db.Column(db.Integer, db.ForeignKey("system_profiles.id"), index=True)



class SystemProfile(db.Model):
    __tablename__ = "system_profiles"
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(120), nullable=False, unique=True, index=True)
    base_role = db.Column(db.String(40), nullable=False, default="technician")
    access_json = db.Column(db.Text, nullable=False, default="[]")
    active = db.Column(db.Boolean, nullable=False, default=True, index=True)
    created_by = db.Column(db.Integer, db.ForeignKey("users.id"))
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)
    updated_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow, onupdate=datetime.utcnow)

class BuiltinDashboardSetting(db.Model):
    __tablename__ = "builtin_dashboard_settings"
    id = db.Column(db.Integer, primary_key=True)
    dashboard_key = db.Column(db.String(80), nullable=False, unique=True, index=True)
    visible = db.Column(db.Boolean, nullable=False, default=True)
    order_index = db.Column(db.Integer, nullable=False, default=100)
    allowed_roles_json = db.Column(db.Text, nullable=False, default="[]")
    updated_by = db.Column(db.Integer, db.ForeignKey("users.id"))
    updated_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow, onupdate=datetime.utcnow)

class CustomerCompany(db.Model):
    __tablename__ = "customer_companies"
    id = db.Column(db.Integer, primary_key=True)
    legal_name = db.Column(db.String(180), nullable=False)
    trade_name = db.Column(db.String(180))
    cnpj = db.Column(db.String(30), unique=True, index=True)
    state_registration = db.Column(db.String(60))
    contact_name = db.Column(db.String(180))
    contact_role = db.Column(db.String(120))
    phone = db.Column(db.String(120))
    mobile = db.Column(db.String(120))
    email = db.Column(db.String(180))
    address = db.Column(db.String(300))
    city = db.Column(db.String(120))
    state = db.Column(db.String(10))
    zip_code = db.Column(db.String(20))
    notes = db.Column(db.Text)
    active = db.Column(db.Boolean, nullable=False, default=True)
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)



# V69 — Portal do Cliente / Agendamentos de equipamentos
class CustomerAppointment(db.Model):
    __tablename__ = "customer_appointments"
    id = db.Column(db.Integer, primary_key=True)
    code = db.Column(db.String(40), unique=True, index=True)
    customer_company = db.Column(db.String(180), nullable=False, index=True)
    responsible_name = db.Column(db.String(180), nullable=False)
    responsible_email = db.Column(db.String(180))
    responsible_phone = db.Column(db.String(40))
    scheduled_date = db.Column(db.Date, index=True)
    # V71 — solicitação e programação logística ficam separadas.
    request_date = db.Column(db.Date, index=True)
    expected_date = db.Column(db.Date, index=True)
    programmed_date = db.Column(db.Date, index=True)
    alternate_date_requested = db.Column(db.Boolean, nullable=False, default=False, index=True)
    alternate_reason = db.Column(db.String(300))
    notes = db.Column(db.Text)
    status = db.Column(db.String(30), nullable=False, default="RASCUNHO", index=True)
    accepted_name = db.Column(db.String(180))
    accepted_at = db.Column(db.DateTime)
    signature_file = db.Column(db.String(600))
    pdf_file = db.Column(db.String(600))
    invoice_number = db.Column(db.String(120))
    invoice_file = db.Column(db.String(600))
    invoice_original_name = db.Column(db.String(255))
    created_by = db.Column(db.Integer, db.ForeignKey("users.id"), nullable=False, index=True)
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow, index=True)
    updated_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow, onupdate=datetime.utcnow)
    received_by = db.Column(db.Integer, db.ForeignKey("users.id"), index=True)
    received_at = db.Column(db.DateTime)
    email_status = db.Column(db.String(30))
    email_detail = db.Column(db.String(500))
    # V71.2 — cancelamento administrativo com rastreabilidade.
    cancelled_at = db.Column(db.DateTime)
    cancelled_by = db.Column(db.Integer, db.ForeignKey("users.id"), index=True)
    cancellation_reason = db.Column(db.String(500))

class CustomerAppointmentEquipment(db.Model):
    __tablename__ = "customer_appointment_equipments"
    id = db.Column(db.Integer, primary_key=True)
    appointment_id = db.Column(db.Integer, db.ForeignKey("customer_appointments.id", ondelete="CASCADE"), nullable=False, index=True)
    item_no = db.Column(db.Integer, nullable=False, default=1)
    serial_number = db.Column(db.String(120), nullable=False, index=True)
    equipment = db.Column(db.String(120))
    version = db.Column(db.String(80))
    eod = db.Column(db.String(120))
    defect = db.Column(db.String(500), nullable=False)
    notes = db.Column(db.Text)
    photo_file = db.Column(db.String(600))
    received = db.Column(db.Boolean, nullable=False, default=False, index=True)
    received_at = db.Column(db.DateTime)
    received_by = db.Column(db.Integer, db.ForeignKey("users.id"))
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)
    __table_args__=(UniqueConstraint("appointment_id","item_no",name="uq_customer_appt_item"),)


# V71 — Matriz Logística Leva e Traz / calendário operacional
class LogisticsGarageRoute(db.Model):
    __tablename__ = "logistics_garage_routes"
    id = db.Column(db.Integer, primary_key=True)
    garage_name = db.Column(db.String(180), nullable=False, unique=True, index=True)
    customer_company_id = db.Column(db.Integer, db.ForeignKey("customer_companies.id"), index=True)
    weekday = db.Column(db.Integer, nullable=False, index=True)  # 0=segunda ... 4=sexta
    contact_name = db.Column(db.String(180))
    address = db.Column(db.String(350))
    region = db.Column(db.String(160))
    active = db.Column(db.Boolean, nullable=False, default=True, index=True)
    source_import = db.Column(db.String(180))
    created_by = db.Column(db.Integer, db.ForeignKey("users.id"))
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)
    updated_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow, onupdate=datetime.utcnow)

class LogisticsBlockedDate(db.Model):
    __tablename__ = "logistics_blocked_dates"
    id = db.Column(db.Integer, primary_key=True)
    blocked_date = db.Column(db.Date, nullable=False, unique=True, index=True)
    description = db.Column(db.String(220))
    active = db.Column(db.Boolean, nullable=False, default=True, index=True)
    created_by = db.Column(db.Integer, db.ForeignKey("users.id"))
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)



# V67 — Dossiê do Colaborador / Materiais e Ferramentas
class MaterialCatalogItem(db.Model):
    __tablename__ = "material_catalog_items"
    id = db.Column(db.Integer, primary_key=True)
    code = db.Column(db.String(50), unique=True, nullable=False, index=True)
    category = db.Column(db.String(40), nullable=False, default="FERRAMENTA", index=True)
    description = db.Column(db.String(220), nullable=False, index=True)
    brand = db.Column(db.String(120))
    model = db.Column(db.String(120))
    unit = db.Column(db.String(20), nullable=False, default="UN")
    control_type = db.Column(db.String(30), nullable=False, default="DEVOLVIVEL")
    quantity_mode = db.Column(db.String(20), nullable=False, default="INTEIRO")
    active = db.Column(db.Boolean, nullable=False, default=True, index=True)
    created_by = db.Column(db.Integer, db.ForeignKey("users.id"))
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)
    updated_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow, onupdate=datetime.utcnow)

class MaterialKit(db.Model):
    __tablename__ = "material_kits"
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(160), nullable=False, unique=True, index=True)
    description = db.Column(db.Text)
    active = db.Column(db.Boolean, nullable=False, default=True)
    created_by = db.Column(db.Integer, db.ForeignKey("users.id"))
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)

class MaterialKitItem(db.Model):
    __tablename__ = "material_kit_items"
    id = db.Column(db.Integer, primary_key=True)
    kit_id = db.Column(db.Integer, db.ForeignKey("material_kits.id", ondelete="CASCADE"), nullable=False, index=True)
    material_id = db.Column(db.Integer, db.ForeignKey("material_catalog_items.id"), nullable=False, index=True)
    quantity = db.Column(db.Float, nullable=False, default=1)
    __table_args__=(UniqueConstraint("kit_id","material_id",name="uq_material_kit_item"),)

class CollaboratorDocument(db.Model):
    __tablename__ = "collaborator_documents"
    id = db.Column(db.Integer, primary_key=True)
    document_code = db.Column(db.String(40), unique=True, index=True)
    document_type = db.Column(db.String(60), nullable=False, default="TERMO_RECEBIMENTO_FERRAMENTAS", index=True)
    user_id = db.Column(db.Integer, db.ForeignKey("users.id"), nullable=False, index=True)
    status = db.Column(db.String(40), nullable=False, default="RASCUNHO", index=True)
    delivery_date = db.Column(db.Date)
    title = db.Column(db.String(220), nullable=False, default="Termo de Recebimento de Materiais / Ferramentas")
    notes = db.Column(db.Text)
    correction_note = db.Column(db.Text)
    signature_file = db.Column(db.String(600))
    return_receiver_signature_file = db.Column(db.String(600))
    return_receiver_id = db.Column(db.Integer, db.ForeignKey("users.id"), index=True)
    return_received_at = db.Column(db.DateTime)
    pdf_file = db.Column(db.String(600))
    invoice_number = db.Column(db.String(120))
    invoice_file = db.Column(db.String(600))
    invoice_original_name = db.Column(db.String(255))
    created_by = db.Column(db.Integer, db.ForeignKey("users.id"), nullable=False, index=True)
    sent_at = db.Column(db.DateTime)
    signed_at = db.Column(db.DateTime)
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)
    updated_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow, onupdate=datetime.utcnow)

class CollaboratorDocumentItem(db.Model):
    __tablename__ = "collaborator_document_items"
    id = db.Column(db.Integer, primary_key=True)
    document_id = db.Column(db.Integer, db.ForeignKey("collaborator_documents.id", ondelete="CASCADE"), nullable=False, index=True)
    material_id = db.Column(db.Integer, db.ForeignKey("material_catalog_items.id"), index=True)
    description = db.Column(db.String(220), nullable=False)
    brand = db.Column(db.String(120))
    model = db.Column(db.String(120))
    quantity = db.Column(db.Float, nullable=False, default=1)
    unit = db.Column(db.String(20), default="UN")
    condition = db.Column(db.String(40), default="BOM")
    notes = db.Column(db.Text)

class MaterialMovement(db.Model):
    __tablename__ = "material_movements"
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey("users.id"), nullable=False, index=True)
    material_id = db.Column(db.Integer, db.ForeignKey("material_catalog_items.id"), nullable=False, index=True)
    document_id = db.Column(db.Integer, db.ForeignKey("collaborator_documents.id"), index=True)
    movement_type = db.Column(db.String(30), nullable=False, index=True)
    quantity = db.Column(db.Float, nullable=False)
    condition = db.Column(db.String(40))
    identifier = db.Column(db.String(160))
    notes = db.Column(db.Text)
    created_by = db.Column(db.Integer, db.ForeignKey("users.id"), nullable=False)
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow, index=True)

class MaterialRequest(db.Model):
    __tablename__ = "material_requests"
    id = db.Column(db.Integer, primary_key=True)
    request_code = db.Column(db.String(40), unique=True, index=True)
    user_id = db.Column(db.Integer, db.ForeignKey("users.id"), nullable=False, index=True)
    material_id = db.Column(db.Integer, db.ForeignKey("material_catalog_items.id"), nullable=False, index=True)
    quantity = db.Column(db.Float, nullable=False, default=1)
    reason = db.Column(db.Text)
    urgency = db.Column(db.String(20), nullable=False, default="NORMAL")
    notes = db.Column(db.Text)
    status = db.Column(db.String(30), nullable=False, default="SOLICITADO", index=True)
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)
    updated_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow, onupdate=datetime.utcnow)


class FinancialSupplier(db.Model):
    __tablename__ = "financial_suppliers"
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(180), nullable=False, unique=True, index=True)
    active = db.Column(db.Boolean, nullable=False, default=True)
    trade_name = db.Column(db.String(180))
    cnpj = db.Column(db.String(30), index=True)
    primary_cost_center = db.Column(db.String(60))
    cost_center_id = db.Column(db.String(20), index=True)
    contact_name = db.Column(db.String(180))
    phone = db.Column(db.String(40))
    email = db.Column(db.String(180))
    pending_profile = db.Column(db.Boolean, nullable=False, default=False)
    created_by = db.Column(db.Integer, db.ForeignKey("users.id"))
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)

class FinancialService(db.Model):
    __tablename__ = "financial_services"
    id = db.Column(db.Integer, primary_key=True)
    supplier_id = db.Column(db.Integer, db.ForeignKey("financial_suppliers.id"), nullable=False, index=True)
    name = db.Column(db.String(220), nullable=False)
    description = db.Column(db.Text)
    category = db.Column(db.String(80), default="OUTROS")
    active = db.Column(db.Boolean, nullable=False, default=True)
    created_by = db.Column(db.Integer, db.ForeignKey("users.id"))
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)

class FinancialMonthlyCost(db.Model):
    __tablename__ = "financial_monthly_costs"
    id = db.Column(db.Integer, primary_key=True)
    competence = db.Column(db.String(7), nullable=False, index=True)
    supplier_id = db.Column(db.Integer, db.ForeignKey("financial_suppliers.id"), nullable=False, index=True)
    service_id = db.Column(db.Integer, db.ForeignKey("financial_services.id"), nullable=False, index=True)
    amount = db.Column(db.Float, nullable=False, default=0)
    forecast_amount = db.Column(db.Float)
    cost_center = db.Column(db.String(60), nullable=False, default="SUPORTE_CAMPO", index=True)
    cost_center_id = db.Column(db.String(20), index=True)
    project = db.Column(db.String(220))
    service_text = db.Column(db.String(300))
    allocation_json = db.Column(db.Text, nullable=False, default="{}")
    invoice_number = db.Column(db.String(120))
    notes = db.Column(db.Text)
    created_by = db.Column(db.Integer, db.ForeignKey("users.id"))
    updated_by = db.Column(db.Integer, db.ForeignKey("users.id"))
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)
    updated_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)

class FinancialCashCollection(db.Model):
    __tablename__ = "financial_cash_collections"
    id = db.Column(db.Integer, primary_key=True)
    terminal = db.Column(db.String(40), nullable=False, index=True)
    point_name = db.Column(db.String(240))
    collection_date = db.Column(db.Date, nullable=False, index=True)
    start_at = db.Column(db.DateTime, index=True)
    end_at = db.Column(db.DateTime, nullable=False, index=True)
    collected_amount = db.Column(db.Float, nullable=False, default=0)
    gtv = db.Column(db.String(60), index=True)
    route = db.Column(db.String(40))
    municipality = db.Column(db.String(120))
    declared_amount = db.Column(db.Float)
    processed_amount = db.Column(db.Float)
    processed_at = db.Column(db.DateTime)
    processed_note_count = db.Column(db.Integer)
    processed_media_type = db.Column(db.String(40))
    processing_charge = db.Column(db.Float)
    source_file = db.Column(db.String(255))
    source_hash = db.Column(db.String(64), nullable=False, unique=True, index=True)
    imported_by = db.Column(db.Integer, db.ForeignKey("users.id"))
    imported_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)

class FinancialATMTransaction(db.Model):
    __tablename__ = "financial_atm_transactions"
    id = db.Column(db.Integer, primary_key=True)
    terminal = db.Column(db.String(40), nullable=False, index=True)
    transaction_at = db.Column(db.DateTime, nullable=False, index=True)
    status = db.Column(db.String(20), index=True)
    value = db.Column(db.Float, nullable=False, default=0)
    cpm_id = db.Column(db.String(60))
    source_file = db.Column(db.String(255))
    source_hash = db.Column(db.String(64), nullable=False, unique=True, index=True)
    imported_by = db.Column(db.Integer, db.ForeignKey("users.id"))
    imported_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)
    __table_args__ = (Index("ix_fin_atm_tx_terminal_datetime", "terminal", "transaction_at"),)

class PerformanceMetric(db.Model):
    __tablename__ = "performance_metrics"
    id = db.Column(db.Integer, primary_key=True)
    route = db.Column(db.String(220), nullable=False, index=True)
    method = db.Column(db.String(12), nullable=False)
    status_code = db.Column(db.Integer, nullable=False, index=True)
    duration_ms = db.Column(db.Float, nullable=False, index=True)
    sql_ms = db.Column(db.Float, nullable=False, default=0)
    query_count = db.Column(db.Integer, nullable=False, default=0)
    user_id = db.Column(db.Integer, db.ForeignKey("users.id"), index=True)
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow, index=True)
    __table_args__ = (Index("ix_perf_created_route", "created_at", "route"),)

class SchemaMigration(db.Model):
    __tablename__ = "schema_migrations"
    id = db.Column(db.Integer, primary_key=True)
    version = db.Column(db.String(80), nullable=False, unique=True, index=True)
    description = db.Column(db.String(240))
    applied_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow, index=True)

class OperationalAction(db.Model):
    __tablename__ = "operational_actions"
    id = db.Column(db.Integer, primary_key=True)
    action_key = db.Column(db.String(220), nullable=False, unique=True, index=True)
    category = db.Column(db.String(60), nullable=False, index=True)
    title = db.Column(db.String(240), nullable=False)
    detail = db.Column(db.Text)
    severity = db.Column(db.String(20), nullable=False, default="MEDIA", index=True)
    status = db.Column(db.String(30), nullable=False, default="NOVO", index=True)
    owner_user_id = db.Column(db.Integer, db.ForeignKey("users.id"), index=True)
    due_date = db.Column(db.Date)
    source_url = db.Column(db.String(500))
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)
    updated_at = db.Column(db.DateTime)


class AuditEvent(db.Model):
    __tablename__ = "audit_events"
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey("users.id"), index=True)
    event_type = db.Column(db.String(80), nullable=False, index=True)
    entity_type = db.Column(db.String(80), nullable=False, index=True)
    entity_id = db.Column(db.String(80), index=True)
    detail = db.Column(db.Text)
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow, index=True)


class TeamScheduleProfile(db.Model):
    __tablename__ = "team_schedule_profiles"
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey("users.id"), index=True)
    name = db.Column(db.String(180), nullable=False, unique=True, index=True)
    active = db.Column(db.Boolean, nullable=False, default=True, index=True)
    category = db.Column(db.String(30), nullable=False, default="TECNICO", index=True)
    schedule_type = db.Column(db.String(30), nullable=False, default="12x36", index=True)
    shift = db.Column(db.String(30), nullable=False)
    supervision = db.Column(db.String(180))
    entry = db.Column(db.String(180))
    lines_json = db.Column(db.Text)
    anchor_date = db.Column(db.Date, nullable=False, index=True)
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)
    updated_at = db.Column(db.DateTime)

    @property
    def lines(self):
        try:
            return json.loads(self.lines_json or "[]")
        except Exception:
            return []


class TechnicianPosition(db.Model):
    __tablename__ = "technician_positions"
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey("users.id"), nullable=False, index=True)
    latitude = db.Column(db.Float, nullable=False)
    longitude = db.Column(db.Float, nullable=False)
    accuracy = db.Column(db.Float)
    captured_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow, index=True)
    source = db.Column(db.String(40), nullable=False, default="browser")

class TechnicianCheckin(db.Model):
    __tablename__ = "technician_checkins"
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey("users.id"), nullable=False, index=True)
    location_id = db.Column(db.Integer, db.ForeignKey("locations.id"), nullable=False, index=True)
    latitude = db.Column(db.Float, nullable=False)
    longitude = db.Column(db.Float, nullable=False)
    accuracy = db.Column(db.Float)
    distance_m = db.Column(db.Float)
    status = db.Column(db.String(40), nullable=False, default="REGISTRADO", index=True)
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow, index=True)




class SessionEvent(db.Model):
    __tablename__ = "session_events"
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey("users.id"), nullable=False, index=True)
    event_type = db.Column(db.String(30), nullable=False, index=True)
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow, index=True)


class TopDeskTicket(db.Model):
    __tablename__ = "topdesk_tickets"
    id = db.Column(db.Integer, primary_key=True)
    ticket_number = db.Column(db.String(120), nullable=False, unique=True, index=True)
    object_id = db.Column(db.String(260), index=True)
    category = db.Column(db.String(180), index=True)
    subcategory = db.Column(db.String(220), index=True)
    incident_type = db.Column(db.String(120), index=True)
    status = db.Column(db.String(120), index=True)
    operator = db.Column(db.String(180))
    created_at_text = db.Column(db.String(120))
    # V56-A: dimensões normalizadas para filtros/analytics no PostgreSQL.
    # Mantemos os campos legados para auditoria e compatibilidade.
    created_at = db.Column(db.DateTime, index=True)
    line_code = db.Column(db.String(40), index=True)
    station_code = db.Column(db.String(180), index=True)
    model_code = db.Column(db.String(80), index=True)
    sla_target_text = db.Column(db.String(120))
    requester = db.Column(db.String(220))
    request_text = db.Column(db.Text)
    action_text = db.Column(db.Text)
    attachments_text = db.Column(db.Text)
    source_file = db.Column(db.String(300))
    source_kind = db.Column(db.String(40), nullable=False, default="TOPDESK_EXCEL")
    equipment_type = db.Column(db.String(80), index=True)
    location_id = db.Column(db.Integer, db.ForeignKey("locations.id"), index=True)
    assigned_technician_id = db.Column(db.Integer, db.ForeignKey("users.id"), index=True)
    work_status = db.Column(db.String(60), nullable=False, default="RECEBIDO", index=True)
    priority = db.Column(db.String(40), nullable=False, default="NORMAL", index=True)
    imported_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)
    last_import_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)
    updated_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)


class TopDeskActivity(db.Model):
    __tablename__ = "topdesk_activities"
    id = db.Column(db.Integer, primary_key=True)
    ticket_id = db.Column(db.Integer, db.ForeignKey("topdesk_tickets.id", ondelete="CASCADE"), nullable=False, index=True)
    technician_id = db.Column(db.Integer, db.ForeignKey("users.id"), index=True)
    actor_user_id = db.Column(db.Integer, db.ForeignKey("users.id"), index=True)
    event_type = db.Column(db.String(60), nullable=False, index=True)
    status = db.Column(db.String(60))
    note = db.Column(db.Text)
    latitude = db.Column(db.Float)
    longitude = db.Column(db.Float)
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow, index=True)


class TopDeskImportBatch(db.Model):
    __tablename__ = "topdesk_import_batches"
    id = db.Column(db.Integer, primary_key=True)
    filename = db.Column(db.String(300), nullable=False)
    imported_by = db.Column(db.Integer, db.ForeignKey("users.id"), nullable=False)
    row_count = db.Column(db.Integer, nullable=False, default=0)
    inserted_count = db.Column(db.Integer, nullable=False, default=0)
    updated_count = db.Column(db.Integer, nullable=False, default=0)
    error_count = db.Column(db.Integer, nullable=False, default=0)
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow, index=True)


class PreventiveRequest(db.Model):
    __tablename__ = "preventive_requests"
    id = db.Column(db.Integer, primary_key=True)
    protocol = db.Column(db.String(40), unique=True, nullable=False, index=True)
    requester_id = db.Column(db.Integer, db.ForeignKey("users.id"), nullable=False, index=True)
    location_id = db.Column(db.Integer, db.ForeignKey("locations.id"), nullable=False, index=True)
    base_asset_id = db.Column(db.Integer, db.ForeignKey("base_assets.id"), nullable=False, index=True)
    company = db.Column(db.String(180))
    locality = db.Column(db.String(220), nullable=False)
    asset_identifier = db.Column(db.String(220), nullable=False, index=True)
    request_type = db.Column(db.String(40), nullable=False, default="PREVENTIVA")
    service = db.Column(db.String(220), nullable=False, index=True)
    description = db.Column(db.Text, nullable=False)
    origin = db.Column(db.String(40), nullable=False, default="SISTEMA_CAMPO")
    status = db.Column(db.String(50), nullable=False, default="AGUARDANDO_INTEGRACAO", index=True)
    topdesk_number = db.Column(db.String(80), index=True)
    topdesk_sla = db.Column(db.String(120))
    latitude = db.Column(db.Float)
    longitude = db.Column(db.Float)
    gps_accuracy = db.Column(db.Float)
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow, index=True)
    updated_at = db.Column(db.DateTime)

class PreventiveAttachment(db.Model):
    __tablename__ = "preventive_attachments"
    id = db.Column(db.Integer, primary_key=True)
    preventive_id = db.Column(db.Integer, db.ForeignKey("preventive_requests.id", ondelete="CASCADE"), nullable=False, index=True)
    original_name = db.Column(db.String(255), nullable=False)
    stored_name = db.Column(db.String(500), nullable=False)
    mime_type = db.Column(db.String(120))
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)

class Location(db.Model):
    __tablename__ = "locations"
    id = db.Column(db.Integer, primary_key=True)
    company = db.Column(db.String(180), nullable=False)
    line = db.Column(db.String(180), nullable=False)
    location = db.Column(db.String(220), nullable=False)
    base_status = db.Column(db.String(80))
    expected_atm = db.Column(db.Integer, nullable=False, default=0)
    expected_validator = db.Column(db.Integer, nullable=False, default=0)
    expected_pos = db.Column(db.Integer, nullable=False, default=0)
    survey_status = db.Column(db.String(30), nullable=False, default="PENDENTE", index=True)
    panorama_status_override = db.Column(db.String(30))
    started_at = db.Column(db.DateTime)
    completed_at = db.Column(db.DateTime)
    completed_by = db.Column(db.Integer, db.ForeignKey("users.id"))
    reference_latitude = db.Column(db.Float)
    reference_longitude = db.Column(db.Float)
    reference_source = db.Column(db.String(120))
    reference_updated_at = db.Column(db.DateTime)

    __table_args__ = (
        UniqueConstraint("company", "line", "location", name="uq_location_company_line_name"),
    )


class BaseAsset(db.Model):
    __tablename__ = "base_assets"
    id = db.Column(db.Integer, primary_key=True)
    asset_key = db.Column(db.String(255), unique=True, index=True)
    description = db.Column(db.String(500))
    company = db.Column(db.String(180))
    station_code = db.Column(db.String(80))
    line = db.Column(db.String(180))
    locality = db.Column(db.String(220))
    serial = db.Column(db.String(180))
    qrcode_id = db.Column(db.String(180))
    top_id = db.Column(db.String(180))
    products = db.Column(db.String(255))
    model = db.Column(db.String(180))
    supplier = db.Column(db.String(180))
    transactions = db.Column(db.String(255))
    pix = db.Column(db.String(80))
    mount = db.Column(db.String(80))
    base_status = db.Column(db.String(80))
    equipment_type = db.Column(db.String(50), index=True)
    location_code = db.Column(db.String(80))
    terminal_number = db.Column(db.String(120))
    application = db.Column(db.String(180))
    bom_id = db.Column(db.String(120))
    bu_id = db.Column(db.String(120))
    software_version = db.Column(db.String(120))
    quantity = db.Column(db.Integer)
    base_notes = db.Column(db.Text)
    leasing_status = db.Column(db.String(120))
    contract_end = db.Column(db.String(80))
    installation_type = db.Column(db.String(180))
    installation_date = db.Column(db.String(80))
    # V42.4 — complemento técnico/cadastral ATM. Identificadores são texto por definição.
    teamviewer_enabled = db.Column(db.String(20))
    teamviewer_id = db.Column(db.String(120), index=True)
    address = db.Column(db.String(500))
    ip_address = db.Column(db.String(120))
    city_id = db.Column(db.String(120))
    praja_id = db.Column(db.String(120))
    cielo_code = db.Column(db.String(120))
    printer_model = db.Column(db.String(180))
    acceptor_model = db.Column(db.String(180))
    motherboard = db.Column(db.String(220))
    ownership_type = db.Column(db.String(80))
    contract_name = db.Column(db.String(180), index=True)


class Inventory(db.Model):
    __tablename__ = "inventory"
    id = db.Column(db.Integer, primary_key=True)
    location_id = db.Column(db.Integer, db.ForeignKey("locations.id"), nullable=False, index=True)
    equipment_type = db.Column(db.String(100), nullable=False)
    base_asset_id = db.Column(db.Integer, db.ForeignKey("base_assets.id"))
    asset_identifier = db.Column(db.String(220), nullable=False)
    serial = db.Column(db.String(220))
    supplier = db.Column(db.String(180))
    model = db.Column(db.String(180))
    application = db.Column(db.String(180))
    bom_id = db.Column(db.String(120))
    bu_id = db.Column(db.String(120))
    validator_top_id = db.Column(db.String(120))
    software_version = db.Column(db.String(120))
    exact_position = db.Column(db.Text)
    mount = db.Column(db.String(100))
    operational_status = db.Column(db.String(120), nullable=False)
    connectivity = db.Column(db.String(120))
    network_id = db.Column(db.String(220))
    teamviewer_id = db.Column(db.String(120))
    label_status = db.Column(db.String(80))
    in_base = db.Column(db.String(80))
    divergence = db.Column(db.String(180))
    notes = db.Column(db.Text)
    technician_id = db.Column(db.Integer, db.ForeignKey("users.id"), nullable=False)
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)
    updated_at = db.Column(db.DateTime)
    latitude = db.Column(db.Float)
    longitude = db.Column(db.Float)
    gps_accuracy = db.Column(db.Float)
    gps_captured_at = db.Column(db.DateTime)
    sync_uuid = db.Column(db.String(80), unique=True, index=True)

    __table_args__ = (
        UniqueConstraint(
            "location_id", "equipment_type", "asset_identifier",
            name="uq_inventory_location_type_identifier"
        ),
    )


class AssetLifecycleEvent(db.Model):
    __tablename__ = "asset_lifecycle_events"
    id = db.Column(db.Integer, primary_key=True)
    inventory_id = db.Column(db.Integer, db.ForeignKey("inventory.id", ondelete="CASCADE"), nullable=False, index=True)
    event_type = db.Column(db.String(50), nullable=False, index=True)
    from_location_id = db.Column(db.Integer, db.ForeignKey("locations.id"), index=True)
    to_location_id = db.Column(db.Integer, db.ForeignKey("locations.id"), index=True)
    status = db.Column(db.String(80))
    notes = db.Column(db.Text)
    user_id = db.Column(db.Integer, db.ForeignKey("users.id"), nullable=False, index=True)
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow, index=True)


class Attachment(db.Model):
    __tablename__ = "attachments"
    id = db.Column(db.Integer, primary_key=True)
    inventory_id = db.Column(db.Integer, db.ForeignKey("inventory.id", ondelete="CASCADE"), nullable=False)
    original_name = db.Column(db.String(300), nullable=False)
    stored_name = db.Column(db.String(400), nullable=False)
    mime_type = db.Column(db.String(180))




class ChipSwap(db.Model):
    __tablename__ = "chip_swaps"
    id = db.Column(db.Integer, primary_key=True)
    location_id = db.Column(db.Integer, db.ForeignKey("locations.id", ondelete="CASCADE"), nullable=False, index=True)
    base_asset_id = db.Column(db.Integer, db.ForeignKey("base_assets.id"), nullable=False, index=True)
    technician_id = db.Column(db.Integer, db.ForeignKey("users.id"), nullable=False, index=True)
    status = db.Column(db.String(40), nullable=False, default="EM ANDAMENTO", index=True)
    notes = db.Column(db.Text)
    latitude = db.Column(db.Float)
    longitude = db.Column(db.Float)
    gps_accuracy = db.Column(db.Float)
    test_result = db.Column(db.String(80), index=True)
    test_notes = db.Column(db.Text)
    started_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow, index=True)
    completed_at = db.Column(db.DateTime, index=True)
    completed_by_id = db.Column(db.Integer, db.ForeignKey("users.id"), index=True)
    updated_at = db.Column(db.DateTime)
    __table_args__ = (UniqueConstraint("location_id", "base_asset_id", name="uq_chip_swap_location_asset"),)

class ChipSwapPhoto(db.Model):
    __tablename__ = "chip_swap_photos"
    id = db.Column(db.Integer, primary_key=True)
    chip_swap_id = db.Column(db.Integer, db.ForeignKey("chip_swaps.id", ondelete="CASCADE"), nullable=False, index=True)
    original_name = db.Column(db.String(300), nullable=False)
    stored_name = db.Column(db.String(700), nullable=False)
    mime_type = db.Column(db.String(180))
    uploaded_by = db.Column(db.Integer, db.ForeignKey("users.id"), nullable=False, index=True)
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow, index=True)


class GarageChipBase(db.Model):
    __tablename__ = "garage_chip_base"
    id = db.Column(db.Integer, primary_key=True)
    company = db.Column(db.String(220), nullable=False, index=True)
    terminal = db.Column(db.String(120), nullable=False, unique=True, index=True)
    model = db.Column(db.String(80), index=True)
    ip = db.Column(db.String(80))
    sam_type = db.Column(db.String(80), default="NÃO MIGRADO", index=True)
    active = db.Column(db.Boolean, nullable=False, default=True, index=True)

class GarageChipSwap(db.Model):
    __tablename__ = "garage_chip_swaps"
    id = db.Column(db.Integer, primary_key=True)
    base_id = db.Column(db.Integer, db.ForeignKey("garage_chip_base.id", ondelete="CASCADE"), nullable=False, unique=True, index=True)
    technician_id = db.Column(db.Integer, db.ForeignKey("users.id"), nullable=False, index=True)
    status = db.Column(db.String(40), nullable=False, default="EM ANDAMENTO", index=True)
    test_result = db.Column(db.String(80), index=True)
    notes = db.Column(db.Text)
    started_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow, index=True)
    completed_at = db.Column(db.DateTime, index=True)
    updated_at = db.Column(db.DateTime)

class GarageChipPhoto(db.Model):
    __tablename__ = "garage_chip_photos"
    id = db.Column(db.Integer, primary_key=True)
    swap_id = db.Column(db.Integer, db.ForeignKey("garage_chip_swaps.id", ondelete="CASCADE"), nullable=False, index=True)
    original_name = db.Column(db.String(300), nullable=False)
    stored_name = db.Column(db.String(700), nullable=False)
    mime_type = db.Column(db.String(180))
    uploaded_by = db.Column(db.Integer, db.ForeignKey("users.id"), nullable=False)
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow, index=True)


class OperationalBaseItem(db.Model):
    """V63 REV2 — base operacional importável, sem apagar o histórico das atividades."""
    __tablename__ = "operational_base_items"
    id = db.Column(db.Integer, primary_key=True)
    module = db.Column(db.String(30), nullable=False, index=True)
    terminal = db.Column(db.String(120), nullable=False, index=True)
    company = db.Column(db.String(220), index=True)
    station = db.Column(db.String(220), index=True)
    line = db.Column(db.String(180), index=True)
    desired_status = db.Column(db.String(40), nullable=False, default="PENDENTE", index=True)
    active = db.Column(db.Boolean, nullable=False, default=True, index=True)
    updated_by = db.Column(db.Integer, db.ForeignKey("users.id"), index=True)
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)
    updated_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow, onupdate=datetime.utcnow)
    __table_args__ = (UniqueConstraint("module", "terminal", name="uq_operational_base_module_terminal"),)

class DashboardDefinition(db.Model):
    __tablename__ = "dashboard_definitions"
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(180), nullable=False)
    slug = db.Column(db.String(180), nullable=False, unique=True, index=True)
    data_source = db.Column(db.String(80), nullable=False, default="TOPDESK", index=True)
    config_json = db.Column(db.Text, nullable=False, default="{}")
    published = db.Column(db.Boolean, nullable=False, default=False, index=True)
    tv_enabled = db.Column(db.Boolean, nullable=False, default=False, index=True)
    tv_order = db.Column(db.Integer, nullable=False, default=0)
    tv_seconds = db.Column(db.Integer, nullable=False, default=30)
    allowed_roles_json = db.Column(db.Text, nullable=False, default="[]")
    created_by = db.Column(db.Integer, db.ForeignKey("users.id"), nullable=False, index=True)
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)
    updated_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow, onupdate=datetime.utcnow)

class EmvChipSwap(db.Model):
    __tablename__ = "emv_chip_swaps"
    id = db.Column(db.Integer, primary_key=True)
    terminal = db.Column(db.String(120), nullable=False, unique=True, index=True)
    technician_id = db.Column(db.Integer, db.ForeignKey("users.id"), nullable=False, index=True)
    status = db.Column(db.String(40), nullable=False, default="EM ANDAMENTO", index=True)
    test_result = db.Column(db.String(80), index=True)
    notes = db.Column(db.Text)
    latitude = db.Column(db.Float); longitude = db.Column(db.Float); gps_accuracy = db.Column(db.Float)
    started_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)
    completed_at = db.Column(db.DateTime); completed_by_id = db.Column(db.Integer, db.ForeignKey("users.id"), index=True); updated_at = db.Column(db.DateTime)
    # V56-A.1: bloqueio informado em campo quando não existe na base EMV.
    manual_entry = db.Column(db.Boolean, nullable=False, default=False, index=True)
    company = db.Column(db.String(120)); line = db.Column(db.String(120)); station = db.Column(db.String(180)); block_number = db.Column(db.String(80))

class EmvChipSwapPhoto(db.Model):
    __tablename__ = "emv_chip_swap_photos"
    id = db.Column(db.Integer, primary_key=True)
    swap_id = db.Column(db.Integer, db.ForeignKey("emv_chip_swaps.id", ondelete="CASCADE"), nullable=False, index=True)
    original_name = db.Column(db.String(300), nullable=False); stored_name = db.Column(db.String(700), nullable=False)
    mime_type = db.Column(db.String(180)); uploaded_by = db.Column(db.Integer, db.ForeignKey("users.id"), nullable=False)
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)

class HardwareFieldVisit(db.Model):
    __tablename__ = "hardware_field_visits"
    id = db.Column(db.Integer, primary_key=True)
    report_code = db.Column(db.String(40), unique=True, index=True)
    client = db.Column(db.String(180), nullable=False, index=True)
    project = db.Column(db.String(180), nullable=False, index=True)
    report_group = db.Column(db.String(40), nullable=False, default="AUTOPASS", index=True)
    requester = db.Column(db.String(180))
    has_topdesk = db.Column(db.Boolean, nullable=False, default=False)
    topdesk_ticket = db.Column(db.String(80), index=True)
    location_type = db.Column(db.String(80))
    location_name = db.Column(db.String(180), nullable=False, index=True)
    city = db.Column(db.String(120), index=True)
    state = db.Column(db.String(10))
    address = db.Column(db.String(300))
    latitude = db.Column(db.Float)
    longitude = db.Column(db.Float)
    gps_accuracy = db.Column(db.Float)
    visit_date = db.Column(db.Date, nullable=False, index=True)
    start_time = db.Column(db.String(10))
    end_time = db.Column(db.String(10))
    reason = db.Column(db.String(120), index=True)
    activities = db.Column(db.Text)
    activity_notes = db.Column(db.Text)
    technical_details = db.Column(db.Text)
    conclusion_status = db.Column(db.String(60), nullable=False, default="EM ANDAMENTO", index=True)
    conclusion = db.Column(db.Text)
    pending_items = db.Column(db.Text)
    client_contact = db.Column(db.String(180))
    client_company = db.Column(db.String(180))
    client_role = db.Column(db.String(120))
    client_email = db.Column(db.String(180))
    client_phone = db.Column(db.String(40))
    # V56-A.3: múltiplos acompanhantes/destinatários preservados no relatório.
    contacts_json = db.Column(db.Text)
    client_observations = db.Column(db.Text)
    client_accepted = db.Column(db.Boolean, nullable=False, default=False)
    signature_file = db.Column(db.String(500))
    signed_at = db.Column(db.DateTime)
    technician_id = db.Column(db.Integer, db.ForeignKey("users.id"), nullable=False, index=True)
    status = db.Column(db.String(40), nullable=False, default="RASCUNHO", index=True)
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)
    updated_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow, onupdate=datetime.utcnow)

class HardwareFieldVisitPhoto(db.Model):
    __tablename__ = "hardware_field_visit_photos"
    id = db.Column(db.Integer, primary_key=True)
    visit_id = db.Column(db.Integer, db.ForeignKey("hardware_field_visits.id"), nullable=False, index=True)
    stored_name = db.Column(db.String(500), nullable=False)
    original_name = db.Column(db.String(255))
    category = db.Column(db.String(80))
    description = db.Column(db.Text)
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)

class PanoramaPoint(db.Model):
    __tablename__ = "panorama_points"
    id = db.Column(db.Integer, primary_key=True)
    location_id = db.Column(db.Integer, db.ForeignKey("locations.id", ondelete="CASCADE"), nullable=False, index=True)
    point_name = db.Column(db.String(220), nullable=False, default="Visão geral")
    notes = db.Column(db.Text)
    created_by = db.Column(db.Integer, db.ForeignKey("users.id"), nullable=False, index=True)
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)
    updated_at = db.Column(db.DateTime)
    __table_args__ = (UniqueConstraint("location_id", "point_name", name="uq_panorama_location_point"),)

class PanoramaPhoto(db.Model):
    __tablename__ = "panorama_photos"
    id = db.Column(db.Integer, primary_key=True)
    point_id = db.Column(db.Integer, db.ForeignKey("panorama_points.id", ondelete="CASCADE"), nullable=False, index=True)
    original_name = db.Column(db.String(300), nullable=False)
    stored_name = db.Column(db.String(700), nullable=False)
    mime_type = db.Column(db.String(180))
    uploaded_by = db.Column(db.Integer, db.ForeignKey("users.id"), nullable=False, index=True)
    latitude = db.Column(db.Float)
    longitude = db.Column(db.Float)
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow, index=True)

Index("idx_inventory_location", Inventory.location_id)


class FieldEvidenceVisit(db.Model):
    __tablename__ = "field_evidence_visits"
    id = db.Column(db.Integer, primary_key=True)
    source_key = db.Column(db.String(80), nullable=False, unique=True, index=True)
    source_batch = db.Column(db.String(120), index=True)
    source_date = db.Column(db.String(20))
    source_time = db.Column(db.String(20))
    author = db.Column(db.String(180))
    station_raw = db.Column(db.String(220))
    line_raw = db.Column(db.String(120))
    location_id = db.Column(db.Integer, db.ForeignKey("locations.id"), index=True)
    match_confidence = db.Column(db.String(40), index=True)
    match_score = db.Column(db.Float)
    report_text = db.Column(db.Text)
    competition_text = db.Column(db.Text)
    storage_source = db.Column(db.String(40))
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)


class FieldEvidenceItem(db.Model):
    __tablename__ = "field_evidence_items"
    id = db.Column(db.Integer, primary_key=True)
    visit_id = db.Column(db.Integer, db.ForeignKey("field_evidence_visits.id", ondelete="CASCADE"), nullable=False, index=True)
    equipment_type = db.Column(db.String(80), nullable=False, index=True)
    identifier = db.Column(db.String(220), nullable=False)
    model = db.Column(db.String(180))
    serial = db.Column(db.String(220))
    patrimony = db.Column(db.String(120))
    operational_status = db.Column(db.String(120))
    source_line = db.Column(db.Text)
    base_asset_id = db.Column(db.Integer, db.ForeignKey("base_assets.id"))
    inventory_id = db.Column(db.Integer, db.ForeignKey("inventory.id"))
    audit_status = db.Column(db.String(80), index=True)
    audit_detail = db.Column(db.Text)

    __table_args__ = (
        UniqueConstraint("visit_id", "equipment_type", "identifier", name="uq_field_evidence_visit_item"),
    )


class FieldEvidenceMedia(db.Model):
    __tablename__ = "field_evidence_media"
    id = db.Column(db.Integer, primary_key=True)
    visit_id = db.Column(db.Integer, db.ForeignKey("field_evidence_visits.id", ondelete="CASCADE"), nullable=False, index=True)
    sha256 = db.Column(db.String(64), nullable=False, unique=True, index=True)
    original_name = db.Column(db.String(300), nullable=False)
    mime_type = db.Column(db.String(180))
    storage_kind = db.Column(db.String(30), nullable=False)
    storage_key = db.Column(db.String(700), nullable=False)
    created_at = db.Column(db.DateTime, nullable=False, default=datetime.utcnow)


Index("idx_field_evidence_location", FieldEvidenceVisit.location_id)
Index("idx_field_evidence_item_audit", FieldEvidenceItem.audit_status)


def normalize(value):
    value = unicodedata.normalize("NFD", value or "").encode("ascii", "ignore").decode()
    return " ".join(value.upper().strip().split())


def seed_data():
    # Usuários iniciais somente quando a tabela está vazia.
    if User.query.count() == 0:
        db.session.add_all([
            User(
                name="Administrador",
                username="admin",
                password_hash=generate_password_hash("Admin@123"),
                role="manager",
                active=True
            ),
            User(
                name="Técnico Field",
                username="tecnico",
                password_hash=generate_password_hash("Tecnico@123"),
                role="technician",
                active=True
            )
        ])
        db.session.commit()

    if Location.query.count() == 0:
        data = json.loads((DATA_DIR / "locations.json").read_text(encoding="utf-8"))
        for x in data:
            db.session.add(Location(
                company=x["company"],
                line=x["line"],
                location=x["location"],
                base_status=x.get("base_status", ""),
                expected_atm=x.get("expected_atm", 0),
                expected_validator=x.get("expected_validator", 0),
                expected_pos=x.get("expected_pos", 0),
                survey_status="PENDENTE"
            ))
        db.session.commit()

    if BaseAsset.query.count() == 0:
        data = json.loads((DATA_DIR / "atm_assets.json").read_text(encoding="utf-8"))
        for a in data:
            db.session.add(BaseAsset(
                asset_key=a.get("asset_key", ""),
                description=a.get("description", ""),
                company=a.get("company", ""),
                station_code=a.get("station_code", ""),
                line=a.get("line", ""),
                locality=a.get("locality", ""),
                serial=a.get("serial", ""),
                qrcode_id=a.get("qrcode_id", ""),
                top_id=a.get("top_id", ""),
                products=a.get("products", ""),
                model=a.get("model", ""),
                supplier=a.get("supplier", ""),
                transactions=a.get("transactions", ""),
                pix=a.get("pix", ""),
                mount=a.get("mount", ""),
                base_status=a.get("base_status", ""),
                equipment_type="ATM"
            ))
        db.session.commit()



def r2_client():
    import boto3
    endpoint = os.environ.get("R2_ENDPOINT_URL", "").strip()
    access_key = os.environ.get("R2_ACCESS_KEY_ID", "").strip()
    secret_key = os.environ.get("R2_SECRET_ACCESS_KEY", "").strip()

    if not endpoint or not access_key or not secret_key:
        raise RuntimeError("Variáveis do R2 não configuradas no Render.")

    return boto3.client(
        "s3",
        endpoint_url=endpoint,
        aws_access_key_id=access_key,
        aws_secret_access_key=secret_key,
        region_name="auto"
    )


def r2_test_connection():
    bucket = os.environ.get("R2_BUCKET_NAME", "").strip()
    if not bucket:
        return False, "R2_BUCKET_NAME não configurado."

    key = "diagnostico/render-r2-test.txt"
    body = b"Inventario Autopass - teste de conexao Render para Cloudflare R2"

    try:
        client = r2_client()
        client.put_object(
            Bucket=bucket,
            Key=key,
            Body=body,
            ContentType="text/plain; charset=utf-8"
        )
        obj = client.get_object(Bucket=bucket, Key=key)
        returned = obj["Body"].read()

        if returned != body:
            return False, "O arquivo foi gravado, mas a leitura retornou conteúdo diferente."

        client.delete_object(Bucket=bucket, Key=key)
        return True, "Conexão validada: gravação, leitura e exclusão concluídas."
    except (ClientError, BotoCoreError, RuntimeError) as exc:
        return False, str(exc)
    except Exception as exc:
        return False, f"{type(exc).__name__}: {exc}"

def login_required(fn):
    @wraps(fn)
    def inner(*args, **kwargs):
        if not session.get("user_id"):
            return redirect(url_for("login"))
        return fn(*args, **kwargs)
    return inner


ACCESS_GROUPS = {
    "dashboard": ("Dashboard Geral", ("dashboard.general",)),
    "field": ("Field", (
        "field.dashboard","field.inventory","field.calls","field.preventive","field.equipment","field.evidence","field.panorama","field.chip_recarga"
    )),
    "implantation": ("Implantação de Hardware", (
        "implantation.dashboard","implantation.visits","implantation.reports","implantation.emv","implantation.garage"
    )),
    "teams": ("RH / Equipes", (
        "teams.map","teams.today","teams.schedule","teams.manage","teams.export"
    )),
    "users": ("RH / Usuários", (
        "users.view","users.create","users.edit","users.activate","users.delete","users.password","users.export"
    )),
    "finance": ("Financeiro", (
        "finance.support","finance.collection","finance.apuracao","finance.assistance","finance.implantation","finance.entries","finance.suppliers","finance.import","finance.edit","finance.delete"
    )),
    "finance_dashboard": ("Dashboard Financeira", ("finance.dashboard",)),
    "management": ("Gestão", (
        "management.calls","management.360","management.notifications","management.diagnostics","management.health","management.settings","management.dashboard_config","management.profiles"
    )),
    "materials": ("Dossiê / Materiais", (
        "materials.my_documents","materials.request","materials.catalog.view","materials.catalog.manage","materials.kits.manage","materials.delivery.create","materials.delivery.manage","materials.dossier.view"
    )),
    "portal": ("Portal do Cliente", ("portal.appointments","portal.receive","portal.manage")),
    "about_versions": ("Sobre / Versões", ("about.versions",)),
}
ACCESS_MODULES = tuple(ACCESS_GROUPS.keys())
ACCESS_SUBMODULES = tuple(k for _g,(_label,children) in ACCESS_GROUPS.items() for k in children)
ACCESS_ALL = set(ACCESS_MODULES) | set(ACCESS_SUBMODULES)
ACCESS_LABELS = {
 "dashboard.general":"Dashboard Geral",
 "field.dashboard":"Dashboard Field","field.inventory":"Inventário / Lançamento","field.calls":"Chamados","field.preventive":"Solicitação Preventiva ATM","field.equipment":"Equipamentos","field.evidence":"Evidências","field.panorama":"Visão Panorâmica","field.chip_recarga":"Troca de Chips – Recarga",
 "implantation.dashboard":"Dashboard Implantação","implantation.visits":"Visita a Campo / Relatório de Visita","implantation.reports":"Relatórios / Visitas recentes","implantation.emv":"Troca de Chips EMV – Trilhos","implantation.garage":"Troca de Chips Garagem",
 "teams.map":"Mapa operacional","teams.today":"Operação de Hoje","teams.schedule":"Escala por dias","teams.manage":"Gestão de equipes / escala","teams.export":"Exportar dados",
 "users.view":"Visualizar usuários","users.create":"Criar usuário","users.edit":"Editar usuário","users.activate":"Ativar / Desativar","users.delete":"Excluir / Arquivar","users.password":"Redefinir senha","users.export":"Exportar Excel",
 "finance.dashboard":"Dashboard Financeira","finance.support":"Suporte a Campo","finance.collection":"Coleta de Valores","finance.apuracao":"Apuração de Numerário","finance.assistance":"Assistência Técnica","finance.implantation":"Implantação de Hardware","finance.entries":"Lançamentos","finance.suppliers":"Empresas / Fornecedores","finance.import":"Importar planilha","finance.edit":"Editar lançamentos","finance.delete":"Excluir lançamentos",
 "management.calls":"Chamados","management.360":"Central 360","management.notifications":"Notificações","management.diagnostics":"Diagnóstico","management.health":"Saúde da Plataforma","management.settings":"Configurações","management.dashboard_config":"Configuração de Dashboards","management.profiles":"Perfis & Permissões",
 "materials.my_documents":"Meus documentos / Minha carga","materials.request":"Solicitar material","materials.catalog.view":"Visualizar catálogo","materials.catalog.manage":"Cadastrar / editar / inativar materiais","materials.kits.manage":"Gerenciar kits","materials.delivery.create":"Criar e enviar entregas","materials.delivery.manage":"Gerenciar aceites / correções","materials.dossier.view":"Dossiê dos colaboradores",
 "portal.appointments":"Criar / consultar agendamentos","portal.receive":"Receber agendamentos","portal.manage":"Administrar Portal do Cliente",
 "about.versions":"Histórico / Versões",
}

def _expand_legacy_access(values):
    values=set(values or ())
    out=set(values)
    # Compatibilidade: uma permissão antiga de módulo significa acesso a todas as subatividades do módulo.
    for group in list(values):
        if group in ACCESS_GROUPS:
            out.update(ACCESS_GROUPS[group][1])
    return out

def _default_access_for_role(role):
    defaults={
      "manager":set(ACCESS_SUBMODULES),
      "manager_field":{"materials.my_documents","materials.request","materials.catalog.view","materials.catalog.manage","materials.kits.manage","materials.delivery.create","materials.delivery.manage","materials.dossier.view","dashboard.general","field.dashboard","field.inventory","field.calls","field.preventive","field.equipment","field.evidence","field.panorama","field.chip_recarga","implantation.dashboard","implantation.visits","implantation.reports","implantation.emv","implantation.garage","teams.map","teams.today","teams.schedule","teams.manage","teams.export","finance.dashboard","management.calls","management.360","management.notifications","management.diagnostics","portal.receive","portal.manage","about.versions"},
      "technician":{"materials.my_documents","materials.request","field.dashboard","field.inventory","field.calls","field.preventive","field.equipment","field.evidence","field.panorama","field.chip_recarga","about.versions"},
      "technician_implantation":{"materials.my_documents","materials.request","field.inventory","field.equipment","field.evidence","field.panorama","field.chip_recarga","implantation.dashboard","implantation.visits","implantation.reports","implantation.emv","implantation.garage","about.versions"},
      "consultation":{"dashboard.general","field.dashboard","field.inventory","field.equipment","field.evidence","field.panorama","field.chip_recarga","teams.map","teams.today","teams.schedule","about.versions"},
      "hr":{"materials.catalog.view","materials.dossier.view","teams.map","teams.today","teams.schedule","teams.manage","teams.export","users.view","users.create","users.edit","users.activate","users.password","users.export","about.versions"},
      "dispatcher":{"dashboard.general","field.calls","field.chip_recarga","teams.map","teams.today","teams.schedule","management.calls","about.versions"},
      "customer":{"portal.appointments"},
      "atm_financial_admin":{"finance.dashboard","finance.support","finance.collection","finance.apuracao","finance.assistance","finance.implantation","finance.entries","finance.suppliers","finance.import","finance.edit","finance.delete","about.versions"},
    }
    return defaults.get(role,set())

def _user_access_set(user=None):
    # V66 REV3: o cache de permissões pertence SOMENTE ao usuário autenticado.
    # Ao renderizar RH/Usuários também consultamos user_access(u) para cada usuário
    # listado. A versão anterior sobrescrevia o cache global da requisição com as
    # permissões do primeiro técnico e fazia os botões Editar/Ativar desaparecerem
    # das linhas seguintes para RH.
    current_lookup = user is None
    if current_lookup and has_request_context() and hasattr(g, "_autopass_access_set"):
        return g._autopass_access_set
    if current_lookup:
        uid=session.get("user_id")
        user=db.session.get(User,uid) if uid else None
    if not user:return set()
    try:
        if getattr(user,"system_profile_id",None):
            prof=db.session.get(SystemProfile,user.system_profile_id)
            if prof and prof.active:
                custom=json.loads(prof.access_json or "[]")
                access=_expand_legacy_access({x for x in custom if x in ACCESS_ALL})
                if current_lookup and has_request_context(): g._autopass_access_set=access
                return access
    except Exception:
        pass
    try:
        custom=json.loads(user.access_json or "null")
        if isinstance(custom,list):
            access=_expand_legacy_access({x for x in custom if x in ACCESS_ALL})
            # RH sempre mantém as visualizações operacionais de Equipes em modo leitura,
            # mesmo quando o access_json foi salvo antes da criação das subpermissões atuais.
            if user.role=="hr": access.update({"teams.map","teams.today","teams.schedule"})
            if current_lookup and has_request_context(): g._autopass_access_set=access
            return access
    except Exception: pass
    access=_default_access_for_role(user.role)
    if current_lookup and has_request_context(): g._autopass_access_set=access
    return access

def _has_access(permission):
    if not session.get("user_id"): return False
    if session.get("role")=="manager": return True
    access=_user_access_set()
    if permission in ACCESS_GROUPS:
        return permission in access or any(x in access for x in ACCESS_GROUPS[permission][1])
    if permission in ACCESS_SUBMODULES:
        group=permission.split('.',1)[0]
        # access_json legado com o pai continua liberando todas as subatividades.
        return permission in access or group in access
    return False

def _parse_access_form(role):
    raw=request.form.getlist("access_modules")
    allowed={x for x in raw if x in ACCESS_ALL}
    # Se o formulário novo mandar filhos, salvamos os filhos; pais antigos permanecem aceitos por compatibilidade.
    if role=="atm_financial_admin": allowed={x for x in allowed if x.startswith("finance.") or x in ("finance","finance_dashboard","about_versions","about.versions")}
    if role=="hr": allowed={x for x in allowed if x.startswith("teams.") or x.startswith("users.") or x.startswith("materials.") or x in ("teams","users","materials","about_versions","about.versions")}
    if role in ("technician","technician_implantation"):
        allowed={x for x in allowed if not (x.startswith("users.") or x.startswith("finance.") or x.startswith("management.")) and x not in ("users","finance","finance_dashboard","management")}
    return sorted(allowed if request.form.get("access_config_present")=="1" else _default_access_for_role(role))

@app.context_processor
def inject_access_helpers():
    return {"can_view": _has_access, "access_modules": ACCESS_MODULES, "access_groups": ACCESS_GROUPS, "access_labels": ACCESS_LABELS, "user_access": _user_access_set}

def manager_required(fn):
    @wraps(fn)
    def inner(*args, **kwargs):
        if not session.get("user_id"):
            return redirect(url_for("login"))
        if session.get("role") != "manager":
            if session.get("role") == "hr":
                return redirect(url_for("teams_page"))
            return redirect(url_for("manager" if session.get("role") == "consultation" else "technician"))
        return fn(*args, **kwargs)
    return inner


def dashboard_required(fn):
    """Painel gerencial completo: Gestor, Gestor Field, Consulta e Dispatcher."""
    @wraps(fn)
    def inner(*args, **kwargs):
        if not session.get("user_id"):
            return redirect(url_for("login"))
        if not _has_access("dashboard"):
            return redirect(url_for("dashboard_landing"))
        return fn(*args, **kwargs)
    return inner


def teams_view_required(fn):
    """Visualização de equipes: Gestor, Consulta, RH e Dispatcher."""
    @wraps(fn)
    def inner(*args, **kwargs):
        if not session.get("user_id"):
            return redirect(url_for("login"))
        if not _has_access("teams"):
            return redirect(url_for("dashboard_landing"))
        return fn(*args, **kwargs)
    return inner


def topdesk_required(fn):
    """TopDesk: Gestor, Gestor Field e Dispatcher."""
    @wraps(fn)
    def inner(*args, **kwargs):
        if not session.get("user_id"):
            return redirect(url_for("login"))
        if not _has_access("management"):
            if request.path.startswith("/api/"):
                return jsonify({"ok": False, "error": "Sem permissão para operação TopDesk."}), 403
            return redirect(url_for("manager" if session.get("role") == "consultation" else "technician_work"))
        return fn(*args, **kwargs)
    return inner


def user_admin_required(fn):
    """Administração de usuários: Gestor e RH."""
    @wraps(fn)
    def inner(*args, **kwargs):
        if not session.get("user_id"):
            return redirect(url_for("login"))
        if not _has_access("users"):
            return redirect(url_for("dashboard_landing"))
        return fn(*args, **kwargs)
    return inner


def _current_user_is_superadmin():
    """Somente o administrador principal pode atribuir perfis sensíveis.
    Por padrão, considera o login `admin`; pode ser alterado por SUPERADMIN_USERNAME.
    O usuário ID 1/manager é mantido como fallback de compatibilidade do projeto.
    """
    uid = session.get("user_id")
    if not uid:
        return False
    user = db.session.get(User, uid)
    if not user or user.role != "manager":
        return False
    expected = os.getenv("SUPERADMIN_USERNAME", "admin").strip().lower()
    return (user.username or "").strip().lower() == expected or user.id == 1


def _role_assignment_allowed(role):
    role = (role or "").strip()
    # RH administra somente perfis operacionais. Perfis sensíveis ficam restritos ao ADM.
    if session.get("role") == "hr":
        return role in ("technician", "technician_implantation")
    if role in ("manager", "manager_field", "consultation", "dispatcher", "atm_financial_admin", "customer"):
        return _current_user_is_superadmin()
    return role in ("technician", "technician_implantation", "hr")

def _hr_target_allowed(user):
    return bool(user) and (session.get("role") != "hr" or user.role in ("technician", "technician_implantation"))


def hardware_implantation_required(fn):
    """Implantação de Hardware: somente Gestor e Técnico Implantação."""
    @wraps(fn)
    def inner(*args, **kwargs):
        if not session.get("user_id"):
            return redirect(url_for("login"))
        if not _has_access("implantation"):
            if request.path.startswith("/api/"):
                return jsonify({"ok": False, "error": "Acesso restrito à Implantação de Hardware."}), 403
            return redirect(url_for("manager" if session.get("role") in ("manager","consultation","dispatcher") else "activities_page"))
        return fn(*args, **kwargs)
    return inner

def emv_field_required(fn):
    """Troca de Chips EMV: Gestor e Técnico Implantação."""
    @wraps(fn)
    def inner(*args, **kwargs):
        if not session.get("user_id"):
            return redirect(url_for("login"))
        if not _has_access("implantation"):
            if request.path.startswith("/api/"):
                return jsonify({"ok": False, "error": "Acesso restrito à equipe de Implantação."}), 403
            return redirect(url_for("activities_page"))
        return fn(*args, **kwargs)
    return inner


def field_required(fn):
    """Operação Field: somente Gestor e Técnico Field."""
    @wraps(fn)
    def inner(*args, **kwargs):
        if not session.get("user_id"):
            return redirect(url_for("login"))
        if not _has_access("field"):
            if request.path.startswith("/api/"):
                return jsonify({"ok": False, "error": "Acesso restrito ao Field."}), 403
            return redirect(url_for("activities_page"))
        return fn(*args, **kwargs)
    return inner


def field_dashboard_required(fn):
    """Dashboard Field: visão operacional para Técnico Field e visão ampla para gestores Field/ADM."""
    @wraps(fn)
    def inner(*args, **kwargs):
        if not session.get("user_id"):
            return redirect(url_for("login"))
        if not _has_access("field"):
            if request.path.startswith("/api/"):
                return jsonify({"ok": False, "error": "Sem permissão para a Dashboard Field."}), 403
            return redirect(url_for("dashboard_landing"))
        return fn(*args, **kwargs)
    return inner


@app.get("/dashboard")
@login_required
def dashboard_landing():
    # V69.3.2: existe uma única central de dashboards: /gerencial.
    # Elimina a tela intermediária (dashboard_hub) e preserva o menu lateral.
    return redirect(url_for("manager"))

@app.route("/")
def index():
    if not session.get("user_id"):
        return redirect(url_for("login"))
    role = session.get("role")
    if role == "hr":
        return redirect(url_for("teams_page"))
    if role == "atm_financial_admin":
        return redirect(url_for("financial_cost_management_page"))
    if role == "customer":
        return redirect(url_for("portal_cliente_page"))
    if role in ("manager", "manager_field", "consultation", "dispatcher", "technician", "technician_implantation"):
        return redirect(url_for("dashboard_landing"))
    return redirect(url_for("my_profile_page"))


@app.route("/login", methods=["GET", "POST"])
def login():
    if request.method == "POST":
        username = request.form.get("username", "").strip().lower()
        password = request.form.get("password", "")
        user = User.query.filter(func.lower(User.username) == username, User.active.is_(True)).first()
        if user and check_password_hash(user.password_hash, password):
            session.clear()
            session.permanent = True
            session.update(user_id=user.id, name=user.name, role=user.role, gps_session_token=uuid.uuid4().hex)
            try:
                db.session.add(SessionEvent(user_id=user.id, event_type="LOGIN"))
                db.session.commit()
            except Exception:
                db.session.rollback()
            if user.role == "hr":
                return redirect(url_for("teams_page"))
            if user.role == "atm_financial_admin":
                return redirect(url_for("financial_cost_management_page"))
            if user.role == "customer":
                return redirect(url_for("portal_cliente_page"))
            return redirect(url_for("manager" if user.role in ("manager", "manager_field", "consultation", "dispatcher") else "activities_page"))
        flash("Usuário ou senha inválidos.")
    return render_template("login.html")


@app.route("/logout")
def logout():
    uid = session.get("user_id")
    if uid:
        try:
            db.session.add(SessionEvent(user_id=uid, event_type="LOGOUT"))
            db.session.commit()
        except Exception:
            db.session.rollback()
    session.clear()
    return redirect(url_for("login"))


def _preventive_asset_identifier(a):
    return (a.terminal_number or a.top_id or a.asset_key or a.description or str(a.id)).strip()

@app.get("/preventivas")
@field_required
def preventive_page():
    if not (_has_access("field.preventive") or _has_access("field.calls")): abort(403)
    return render_template("preventive.html", app_release=APP_RELEASE)

@app.get("/api/preventivas/contexto")
@field_required
def preventive_context_api():
    if not (_has_access("field.preventive") or _has_access("field.calls")): abort(403)
    locs=Location.query.order_by(Location.company,Location.line,Location.location).all()
    # Só oferece localidades com ATM cadastrado na base ou inventariado.
    atm_localities={str(x.locality or '').strip().casefold() for x in BaseAsset.query.filter(func.upper(BaseAsset.equipment_type).like('%ATM%')).all() if x.locality}
    rows=[{"id":l.id,"company":l.company,"line":l.line,"location":l.location} for l in locs if l.location.strip().casefold() in atm_localities]
    services=["Bobina - Pouco Papel","Limpeza preventiva","Inspeção geral","Teste de comunicação","Verificação de periféricos","Outro"]
    return jsonify({"ok":True,"locations":rows,"services":services,"request_type":"PREVENTIVA","integration":"TOPDESK_PENDENTE"})

@app.get("/api/preventivas/ativos")
@field_required
def preventive_assets_api():
    if not (_has_access("field.preventive") or _has_access("field.calls")): abort(403)
    location_id=request.args.get("location_id",type=int); loc=db.session.get(Location,location_id) if location_id else None
    if not loc:return jsonify({"ok":False,"error":"Localidade inválida."}),400
    q=BaseAsset.query.filter(func.upper(BaseAsset.equipment_type).like('%ATM%'))
    assets=[a for a in q.all() if (a.locality or '').strip().casefold()==loc.location.strip().casefold()]
    rows=[]
    for a in sorted(assets,key=lambda x:_preventive_asset_identifier(x)):
        rows.append({"id":a.id,"identifier":_preventive_asset_identifier(a),"model":a.model or "","serial":a.serial or "","ip":a.ip_address or "","company":a.company or loc.company})
    return jsonify({"ok":True,"location":{"id":loc.id,"company":loc.company,"line":loc.line,"location":loc.location},"assets":rows})

@app.get("/api/preventivas")
@field_required
def preventive_list_api():
    if not (_has_access("field.preventive") or _has_access("field.calls")): abort(403)
    q=PreventiveRequest.query.order_by(PreventiveRequest.created_at.desc())
    if session.get('role')=='technician':q=q.filter_by(requester_id=session.get('user_id'))
    out=[]
    for x in q.limit(100).all():
        out.append({"id":x.id,"protocol":x.protocol,"locality":x.locality,"asset_identifier":x.asset_identifier,"service":x.service,"status":x.status,"topdesk_number":x.topdesk_number,"created_at":x.created_at.isoformat()+"Z"})
    return jsonify({"ok":True,"rows":out})

@app.post("/api/preventivas")
@field_required
def preventive_create_api():
    if not (_has_access("field.preventive") or _has_access("field.calls")): abort(403)
    location_id=request.form.get("location_id",type=int); asset_id=request.form.get("base_asset_id",type=int)
    service=(request.form.get("service") or "").strip(); description=(request.form.get("description") or "").strip()
    loc=db.session.get(Location,location_id) if location_id else None; asset=db.session.get(BaseAsset,asset_id) if asset_id else None
    if not loc or not asset or not service or not description:return jsonify({"ok":False,"error":"Preencha localidade, ATM, serviço e descrição."}),400
    if (asset.locality or '').strip().casefold()!=loc.location.strip().casefold():return jsonify({"ok":False,"error":"O ATM selecionado não pertence à localidade informada."}),400
    ident=_preventive_asset_identifier(asset)
    active=("AGUARDANDO_INTEGRACAO","ENVIANDO_TOPDESK","ABERTA_TOPDESK","EM_ATENDIMENTO")
    dup=PreventiveRequest.query.filter(PreventiveRequest.base_asset_id==asset.id,PreventiveRequest.service==service,PreventiveRequest.status.in_(active)).order_by(PreventiveRequest.created_at.desc()).first()
    if dup:return jsonify({"ok":False,"duplicate":True,"error":f"Já existe preventiva ativa para este ATM/serviço: {dup.topdesk_number or dup.protocol}.","existing":{"id":dup.id,"protocol":dup.protocol,"topdesk_number":dup.topdesk_number,"status":dup.status}}),409
    now=datetime.now(ZoneInfo("America/Sao_Paulo")); protocol=f"PV-{now:%y%m%d}-{uuid.uuid4().hex[:6].upper()}"
    lat=request.form.get("latitude",type=float); lon=request.form.get("longitude",type=float); acc=request.form.get("gps_accuracy",type=float)
    row=PreventiveRequest(protocol=protocol,requester_id=session.get('user_id'),location_id=loc.id,base_asset_id=asset.id,company=loc.company,locality=loc.location,asset_identifier=ident,service=service,description=description,latitude=lat,longitude=lon,gps_accuracy=acc)
    db.session.add(row); db.session.flush()
    allowed={"image/jpeg","image/png","image/webp","application/pdf"}; folder=UPLOAD_DIR/"preventive"; folder.mkdir(parents=True,exist_ok=True)
    for f in request.files.getlist("attachments")[:8]:
        if not f or not f.filename:continue
        mime=(f.mimetype or "").lower()
        if mime not in allowed:continue
        ext=Path(f.filename).suffix.lower()[:10]; stored=f"preventive/{row.id}_{uuid.uuid4().hex[:12]}{ext}"; target=UPLOAD_DIR/stored; target.parent.mkdir(parents=True,exist_ok=True); f.save(target)
        db.session.add(PreventiveAttachment(preventive_id=row.id,original_name=f.filename[:255],stored_name=stored,mime_type=mime))
    db.session.commit()
    return jsonify({"ok":True,"id":row.id,"protocol":row.protocol,"status":row.status,"message":"Preventiva registrada. Integração TOPdesk ainda não ativada."})

@app.route("/trabalho")
@field_required
def technician_work():
    return render_template("technician_work.html", app_release=APP_RELEASE)


@app.route("/tecnico")
@field_required
def technician():
    return render_template("technician.html")


@app.route("/gerencial")
@dashboard_required
def manager():
    if session.get("role")=="technician":
        return redirect(url_for("field_dashboard_page"))
    return render_template("manager.html", app_release=APP_RELEASE)

@app.route("/gerencial/tv")
@dashboard_required
def manager_tv():
    if session.get("role")=="technician":
        return redirect(url_for("field_dashboard_page"))
    return render_template("manager_tv.html", app_release=APP_RELEASE)


@app.get("/dashboard/atm")
@dashboard_required
def atm_dashboard_page():
    # V50.1: rota canônica para a Dashboard ATM 2.0 dentro do gerencial.
    # Mantém parâmetros de investigação (ex.: teamviewer_missing=1).
    qs=request.query_string.decode("utf-8")
    suffix=("&"+qs) if qs else ""
    return redirect(f"/gerencial?view=atm-inventory{suffix}")


@app.get("/dashboard/field")
@field_dashboard_required
def field_dashboard_page():
    return render_template("field_dashboard.html", app_release=APP_RELEASE)


@app.get("/api/dashboard/field")
@field_dashboard_required
def field_dashboard_api():
    uid=session.get("user_id"); role=session.get("role")
    company=(request.args.get("company") or "").strip(); line=(request.args.get("line") or "").strip()
    location=(request.args.get("location") or "").strip(); equipment=(request.args.get("equipment") or "").strip()
    status=(request.args.get("status") or "").strip(); technician=(request.args.get("technician") or "").strip()
    date_from=(request.args.get("date_from") or "").strip(); date_to=(request.args.get("date_to") or "").strip()
    def dt(v, end=False):
        if not v: return None
        try:
            d=datetime.strptime(v,"%Y-%m-%d")
            return d.replace(hour=23,minute=59,second=59) if end else d
        except Exception: return None
    d1,d2=dt(date_from),dt(date_to,True)
    # opções vêm do escopo permitido, sem depender do filtro corrente.
    lq=Location.query
    locations=lq.order_by(Location.company,Location.line,Location.location).all()
    inv_base=Inventory.query
    if role=="technician": inv_base=inv_base.filter(Inventory.technician_id==uid)
    tech_ids={x[0] for x in inv_base.with_entities(Inventory.technician_id).distinct().all() if x[0]}
    users=User.query.filter(User.id.in_(tech_ids)).order_by(User.name).all() if tech_ids else []
    inv_q=inv_base.join(Location,Inventory.location_id==Location.id)
    tq=TopDeskTicket.query.outerjoin(Location,TopDeskTicket.location_id==Location.id)
    cq=ChipSwap.query.join(Location,ChipSwap.location_id==Location.id)
    if role=="technician":
        tq=tq.filter(TopDeskTicket.assigned_technician_id==uid); cq=cq.filter(ChipSwap.technician_id==uid)
    if company:
        inv_q=inv_q.filter(Location.company==company); tq=tq.filter(Location.company==company); cq=cq.filter(Location.company==company)
    if line:
        inv_q=inv_q.filter(Location.line==line); tq=tq.filter(Location.line==line); cq=cq.filter(Location.line==line)
    if location:
        inv_q=inv_q.filter(Location.location==location); tq=tq.filter(Location.location==location); cq=cq.filter(Location.location==location)
    if technician and technician.isdigit():
        tid=int(technician); inv_q=inv_q.filter(Inventory.technician_id==tid); tq=tq.filter(TopDeskTicket.assigned_technician_id==tid); cq=cq.filter(ChipSwap.technician_id==tid)
    if equipment: inv_q=inv_q.filter(Inventory.equipment_type.ilike(f"%{equipment}%"))
    if status: inv_q=inv_q.filter(Inventory.operational_status==status)
    if d1:
        inv_q=inv_q.filter(Inventory.created_at>=d1); tq=tq.filter(TopDeskTicket.created_at>=d1); cq=cq.filter(ChipSwap.started_at>=d1)
    if d2:
        inv_q=inv_q.filter(Inventory.created_at<=d2); tq=tq.filter(TopDeskTicket.created_at<=d2); cq=cq.filter(ChipSwap.started_at<=d2)
    inv=inv_q.all(); tickets=tq.all(); chips=cq.all()
    by_type={}
    for x in inv:
        k=_canonical_equipment_type(x.equipment_type); by_type[k]=by_type.get(k,0)+1
    inop=[x for x in inv if normalize(x.operational_status) in ("INOPERANTE","DEFEITO","FORA DE OPERACAO")]
    div=[x for x in inv if (x.divergence or "").strip()]
    open_tickets=[x for x in tickets if normalize(x.work_status) not in ("CONCLUIDO","RESOLVIDO","FECHADO")]
    done=[x for x in chips if normalize(getattr(x,"status","") or "") in ("CONCLUIDO","CONCLUIDA")]
    locmap={x.id:x for x in locations}; usermap={u.id:u.name for u in User.query.filter(User.id.in_({*(x.technician_id for x in inv),*(x.assigned_technician_id for x in tickets if x.assigned_technician_id),*(x.technician_id for x in chips)})).all()} if (inv or tickets or chips) else {}
    def invrow(x):
        l=locmap.get(x.location_id); return {"id":x.id,"type":_canonical_equipment_type(x.equipment_type),"identifier":x.asset_identifier,"serial":x.serial or "—","status":x.operational_status,"divergence":x.divergence or "","company":getattr(l,"company","") or "","line":getattr(l,"line","") or "","location":getattr(l,"location","") or "","technician":usermap.get(x.technician_id,"")}
    def tickrow(x):
        l=locmap.get(x.location_id); return {"id":x.id,"ticket":x.ticket_number,"status":x.work_status,"company":getattr(l,"company","") or "","line":getattr(l,"line","") or "","location":getattr(l,"location","") or "","technician":usermap.get(x.assigned_technician_id,"")}
    def chiprow(x):
        l=locmap.get(x.location_id); return {"id":x.id,"status":x.status,"result":x.test_result or "","company":getattr(l,"company","") or "","line":getattr(l,"line","") or "","location":getattr(l,"location","") or "","technician":usermap.get(x.technician_id,"")}
    equipment_opts=sorted({_canonical_equipment_type(x[0]) for x in inv_base.with_entities(Inventory.equipment_type).distinct().all() if x[0]})
    status_opts=sorted({x[0] for x in inv_base.with_entities(Inventory.operational_status).distinct().all() if x[0]})
    return jsonify({"ok":True,"inventory":{"total":len(inv),"inoperative":len(inop),"divergences":len(div),"by_type":by_type},"tickets":{"total":len(tickets),"open":len(open_tickets)},"chips":{"total":len(chips),"done":len(done)},"filters":{"companies":sorted({x.company for x in locations if x.company}),"lines":sorted({x.line for x in locations if x.line and (not company or x.company==company)}),"locations":sorted({x.location for x in locations if x.location and (not company or x.company==company) and (not line or x.line==line)}),"equipment":equipment_opts,"statuses":status_opts,"technicians":[{"id":u.id,"name":u.name} for u in users]},"details":{"inventory":[invrow(x) for x in inv[:200]],"inoperative":[invrow(x) for x in inop[:200]],"divergences":[invrow(x) for x in div[:200]],"tickets":[tickrow(x) for x in tickets[:200]],"open_tickets":[tickrow(x) for x in open_tickets[:200]],"chips":[chiprow(x) for x in chips[:200]],"chip_done":[chiprow(x) for x in done[:200]]}})



def migrate_user_archive_column():
    try:
        inspector = db.inspect(db.engine)
        existing = {c["name"] for c in inspector.get_columns("users")}
        if "archived_at" not in existing:
            db.session.execute(db.text('ALTER TABLE users ADD COLUMN archived_at TIMESTAMP'))
            db.session.commit()
    except Exception:
        db.session.rollback()
        raise


def migrate_user_v23_columns():
    """V23: empresa e referência individual da escala do colaborador."""
    try:
        inspector = db.inspect(db.engine)
        existing = {c["name"] for c in inspector.get_columns("users")}
        commands = []
        if "company" not in existing:
            commands.append('ALTER TABLE users ADD COLUMN company VARCHAR(180)')
        if "work_schedule_type" not in existing:
            commands.append('ALTER TABLE users ADD COLUMN work_schedule_type VARCHAR(30)')
        if "work_shift" not in existing:
            commands.append('ALTER TABLE users ADD COLUMN work_shift VARCHAR(30)')
        if "work_anchor_date" not in existing:
            commands.append('ALTER TABLE users ADD COLUMN work_anchor_date DATE')
        if "work_anchor_status" not in existing:
            commands.append('ALTER TABLE users ADD COLUMN work_anchor_status VARCHAR(20)')
        if "job_title" not in existing:
            commands.append('ALTER TABLE users ADD COLUMN job_title VARCHAR(120)')
        if "personnel_status" not in existing:
            commands.append("ALTER TABLE users ADD COLUMN personnel_status VARCHAR(30) DEFAULT 'ATIVO'")
        if "personnel_status_note" not in existing:
            commands.append('ALTER TABLE users ADD COLUMN personnel_status_note VARCHAR(240)')
        for command in commands:
            db.session.execute(db.text(command))
        if commands:
            db.session.commit()
    except Exception:
        db.session.rollback()
        raise



def migrate_user_gps_required_column():
    """V66 REV2: GPS obrigatório passa a ser configurável individualmente."""
    try:
        inspector=db.inspect(db.engine)
        existing={c["name"] for c in inspector.get_columns("users")}
        if "gps_required" not in existing:
            db.session.execute(db.text("ALTER TABLE users ADD COLUMN gps_required BOOLEAN DEFAULT FALSE NOT NULL"))
            # Preserva o comportamento da REV1 para os perfis operacionais já existentes.
            db.session.execute(db.text("UPDATE users SET gps_required=TRUE WHERE role IN ('technician','technician_implantation','manager_field','dispatcher')"))
            db.session.commit()
    except Exception:
        db.session.rollback(); raise


def _sync_user_schedule_profile(user):
    """Mantém a grade operacional alinhada ao cadastro do usuário."""
    if not user or user.role != "technician":
        return
    profile = TeamScheduleProfile.query.filter_by(user_id=user.id).first()
    if not profile:
        profile = TeamScheduleProfile.query.filter(func.upper(TeamScheduleProfile.name) == normalize(user.name)).first()
    schedule_type = (user.work_schedule_type or "12x36").strip()
    shift = (user.work_shift or "05:00-17:00").strip()
    ref_date = user.work_anchor_date or datetime.now(ZoneInfo("America/Sao_Paulo")).date()
    status = normalize(user.work_anchor_status or "TRABALHA")
    anchor = ref_date if status != "FOLGA" else ref_date - timedelta(days=1)
    if not profile:
        profile = TeamScheduleProfile(
            user_id=user.id, name=user.name, active=bool(user.active), category="TECNICO",
            schedule_type=schedule_type, shift=shift, supervision="", entry="",
            lines_json="[]", anchor_date=anchor
        )
        db.session.add(profile)
    else:
        profile.user_id=user.id
        profile.name=user.name
        profile.active=bool(user.active)
        profile.schedule_type=schedule_type
        profile.shift=shift
        profile.anchor_date=anchor
    profile.updated_at=datetime.utcnow()


def migrate_inventory_sync_uuid():
    """V8: idempotency key for offline/PWA retries."""
    try:
        inspector = db.inspect(db.engine)
        existing = {c["name"] for c in inspector.get_columns("inventory")}
        if "sync_uuid" not in existing:
            db.session.execute(db.text('ALTER TABLE inventory ADD COLUMN sync_uuid VARCHAR(80)'))
            db.session.commit()
        try:
            db.session.execute(db.text('CREATE UNIQUE INDEX IF NOT EXISTS ix_inventory_sync_uuid ON inventory (sync_uuid)'))
            db.session.commit()
        except Exception:
            db.session.rollback()
    except Exception:
        db.session.rollback()
        raise


def migrate_team_schedule_columns():
    """Adds V6 schedule columns without losing the editable roster already in PostgreSQL."""
    try:
        inspector = db.inspect(db.engine)
        existing = {c["name"] for c in inspector.get_columns("team_schedule_profiles")}
        commands = []
        if "category" not in existing:
            commands.append('ALTER TABLE team_schedule_profiles ADD COLUMN category VARCHAR(30)')
        if "schedule_type" not in existing:
            commands.append('ALTER TABLE team_schedule_profiles ADD COLUMN schedule_type VARCHAR(30)')
        for command in commands:
            db.session.execute(db.text(command))
        if commands:
            db.session.execute(db.text(
                "UPDATE team_schedule_profiles SET category='TECNICO' WHERE category IS NULL OR category=''"
            ))
            db.session.execute(db.text(
                "UPDATE team_schedule_profiles SET schedule_type='12x36' WHERE schedule_type IS NULL OR schedule_type=''"
            ))
            db.session.commit()
    except Exception:
        db.session.rollback()
        raise



def _load_technician_schedule():
    path = DATA_DIR / "technician_schedule_v5.json"
    if not path.exists():
        return {"technicians": [], "support": []}
    return json.loads(path.read_text(encoding="utf-8"))


def _schedule_default_anchor(group, schedule):
    techs = schedule.get("technicians", [])
    anchors = sorted({
        (x.get("days") or [None])[0]
        for x in techs if (x.get("days") or [None])[0]
    })
    # The source plan has two alternate 12x36 groups.
    if not anchors:
        return datetime.now(ZoneInfo("America/Sao_Paulo")).date()
    if normalize(group) == "A":
        raw = anchors[-1]
    elif normalize(group) == "B":
        raw = anchors[0]
    else:
        raw = anchors[0]
    return datetime.strptime(raw, "%Y-%m-%d").date()


def _ensure_team_schedule_profiles(force=False):
    """V58: sincroniza perfis de escala em lote, com TTL e sem consultas N+1.

    A escala continua derivada do cadastro de usuários, porém as telas de leitura não
    fazem mais SELECT por colaborador. Uma sincronização completa custa poucas queries
    e só roda novamente após o TTL (ou quando force=True).
    """
    now=time.time()
    if not force and now-float(_TEAM_PROFILE_SYNC_STATE.get("at") or 0) < _TEAM_PROFILE_SYNC_TTL:
        return
    with _TEAM_PROFILE_SYNC_LOCK:
        now=time.time()
        if not force and now-float(_TEAM_PROFILE_SYNC_STATE.get("at") or 0) < _TEAM_PROFILE_SYNC_TTL:
            return
        users=(User.query.filter(User.active.is_(True),User.role.in_(("technician","technician_implantation","dispatcher","manager"))).order_by(User.name).all())
        profiles=TeamScheduleProfile.query.order_by(TeamScheduleProfile.id).all()
        by_uid={p.user_id:p for p in profiles if p.user_id}
        by_name={}
        for p in profiles:
            by_name.setdefault((p.name or "").strip().casefold(),p)
        valid_ids=set()
        try:
            today=datetime.now(ZoneInfo("America/Sao_Paulo")).date()
            for u in users:
                if normalize(u.personnel_status or "ATIVO") != "ATIVO":
                    continue
                valid_ids.add(u.id)
                row=by_uid.get(u.id) or by_name.get((u.name or "").strip().casefold())
                sched=(u.work_schedule_type or "12x36").strip()
                shift=(u.work_shift or ("08:00-18:00" if normalize(sched)=="5X2" else "05:00-17:00")).strip()
                ref_date=u.work_anchor_date or today
                ref_status=normalize(u.work_anchor_status or "TRABALHA")
                anchor=ref_date if ref_status != "FOLGA" else ref_date-timedelta(days=1)
                jt=normalize(u.job_title or "")
                category="SUPERVISOR" if ("SUPERV" in jt or u.role=="dispatcher") else ("APOIO" if "APOIO" in jt else "TECNICO")
                if row is None:
                    row=TeamScheduleProfile(user_id=u.id,name=u.name,active=True,category=category,schedule_type=sched,shift=shift,supervision="",entry="",lines_json="[]",anchor_date=anchor)
                    db.session.add(row)
                    profiles.append(row); by_uid[u.id]=row; by_name[(u.name or "").strip().casefold()]=row
                else:
                    row.user_id=u.id; row.name=u.name; row.active=True; row.category=category; row.schedule_type=sched; row.shift=shift; row.anchor_date=anchor; row.updated_at=datetime.utcnow()
            for row in profiles:
                if row.user_id and row.user_id not in valid_ids:
                    row.active=False
            db.session.commit()
            _TEAM_PROFILE_SYNC_STATE["at"]=time.time()
        except IntegrityError:
            db.session.rollback()
            # legado com nomes duplicados: deixa a próxima sincronização tentar após saneamento manual.
            _TEAM_PROFILE_SYNC_STATE["at"]=time.time()


def _team_profile_is_scheduled(profile, target_date):
    if not profile.active:
        return False
    schedule_type = normalize(profile.schedule_type or "12x36")
    if schedule_type == "5X2":
        # Operational support: Monday-Friday.
        return target_date.weekday() < 5
    if not profile.anchor_date:
        return False
    delta = (target_date - profile.anchor_date).days
    return delta >= 0 and delta % 2 == 0


def _profile_to_dict(profile, user=None):
    if user is None and profile.user_id:
        user = db.session.get(User, profile.user_id)
    # Perfil vinculado a usuário inativo deixa de participar da visão operacional.
    linked_user_active = bool(user and user.active)
    return {
        "profile_id": profile.id,
        "user_id": profile.user_id,
        "linked_user_name": user.name if linked_user_active else None,
        "linked": linked_user_active,
        "name": profile.name,
        "category": profile.category or "TECNICO",
        "schedule_type": profile.schedule_type or "12x36",
        "shift": profile.shift,
        "supervision": profile.supervision or "",
        "entry": profile.entry or "",
        "lines": profile.lines,
        "anchor_date": profile.anchor_date.isoformat() if profile.anchor_date else None,
        "active": bool(profile.active),
        "company": (user.company if user else "") or "",
        "job_title": (user.job_title if user else "") or "",
        "personnel_status": (user.personnel_status if user else "ATIVO") or "ATIVO",
        "personnel_status_note": (user.personnel_status_note if user else "") or "",
        "source": "CADASTRO_USUARIO" if user else "LEGADO_ESCALA",
    }


def _schedule_today_db(target_date=None):
    _ensure_team_schedule_profiles()
    target_date = target_date or datetime.now(ZoneInfo("America/Sao_Paulo")).date()
    profiles=TeamScheduleProfile.query.filter_by(active=True).order_by(TeamScheduleProfile.category,TeamScheduleProfile.name).all()
    user_ids={p.user_id for p in profiles if p.user_id}
    users={u.id:u for u in User.query.filter(User.id.in_(user_ids)).all()} if user_ids else {}
    out=[]
    for p in profiles:
        if not _team_profile_is_scheduled(p,target_date): continue
        u=users.get(p.user_id) if p.user_id else None
        if p.user_id and (not u or not u.active or normalize(u.personnel_status or "ATIVO")!="ATIVO"): continue
        d={"profile_id":p.id,"user_id":p.user_id,"linked_user_name":u.name if u else None,"linked":bool(u and u.active),"name":p.name,"category":p.category or "TECNICO","schedule_type":p.schedule_type or "12x36","shift":p.shift,"supervision":p.supervision or "","entry":p.entry or "","lines":p.lines,"anchor_date":p.anchor_date.isoformat() if p.anchor_date else None,"active":bool(p.active),"company":(u.company if u else "") or "","job_title":(u.job_title if u else "") or "","personnel_status":(u.personnel_status if u else "ATIVO") or "ATIVO","personnel_status_note":(u.personnel_status_note if u else "") or "","source":"CADASTRO_USUARIO" if u else "LEGADO_ESCALA"}
        out.append(d)
    return out


def _team_latest_position(user_id, only_today=True):
    if not user_id:
        return None
    q=TechnicianPosition.query.filter_by(user_id=user_id)
    if only_today:
        local_day=datetime.now(ZoneInfo("America/Sao_Paulo")).date()
        start_local=datetime.combine(local_day, datetime.min.time(), tzinfo=ZoneInfo("America/Sao_Paulo"))
        end_local=start_local+timedelta(days=1)
        start_utc=start_local.astimezone(ZoneInfo("UTC")).replace(tzinfo=None)
        end_utc=end_local.astimezone(ZoneInfo("UTC")).replace(tzinfo=None)
        q=q.filter(TechnicianPosition.captured_at>=start_utc,TechnicianPosition.captured_at<end_utc)
    return q.order_by(TechnicianPosition.captured_at.desc()).first()


@app.post("/api/tecnico/position")
@field_required
def technician_position_update():
    data = request.get_json(silent=True) or {}
    try:
        lat = float(data.get("latitude"))
        lon = float(data.get("longitude"))
        acc = float(data.get("accuracy")) if data.get("accuracy") is not None else None
    except (TypeError, ValueError):
        return jsonify({"ok": False, "error": "Coordenadas inválidas"}), 400
    if not (-90 <= lat <= 90 and -180 <= lon <= 180):
        return jsonify({"ok": False, "error": "Coordenadas fora do intervalo"}), 400

    row = TechnicianPosition(
        user_id=session["user_id"], latitude=lat, longitude=lon, accuracy=acc,
        captured_at=datetime.utcnow(), source=str(data.get("source") or "session_periodic")[:40]
    )
    db.session.add(row)
    # V66 GPS Operacional 2.0: retenção móvel de 7 dias. A limpeza é amortizada
    # (no máximo uma vez/hora por processo) para não penalizar cada captura.
    global _GPS_LAST_RETENTION_CLEANUP
    now_ts=time.time()
    if now_ts-_GPS_LAST_RETENTION_CLEANUP > 3600:
        cutoff=datetime.utcnow()-timedelta(days=max(1,int(os.getenv("TEAM_GPS_RETENTION_DAYS","7"))))
        TechnicianPosition.query.filter(TechnicianPosition.captured_at < cutoff).delete(synchronize_session=False)
        _GPS_LAST_RETENTION_CLEANUP=now_ts
    db.session.commit()
    return jsonify({"ok": True, "captured_at": row.captured_at.isoformat() + "Z", "retention_days":7})


@app.get("/api/campo/config")
@field_required
def field_config_api():
    return jsonify({
        "ok": True,
        "release": APP_RELEASE,
        "nearby_radius_m": FIELD_NEARBY_RADIUS_M,
        "gps_good_accuracy_m": FIELD_GPS_GOOD_ACCURACY_M,
        "gps_max_accuracy_m": FIELD_GPS_MAX_ACCURACY_M,
        "gps_warn_distance_m": float(os.getenv("FIELD_GPS_WARN_DISTANCE_M", "250")),
        "gps_max_distance_m": float(os.getenv("FIELD_GPS_MAX_DISTANCE_M", "600")),
        "gps_override_min_chars": 3,
    })


@app.post("/api/tecnico/checkin")
@field_required
def technician_checkin():
    data = request.get_json(silent=True) or {}
    location_id = data.get("location_id")
    try:
        location_id = int(location_id)
        lat = float(data.get("latitude"))
        lon = float(data.get("longitude"))
        acc = float(data.get("accuracy")) if data.get("accuracy") is not None else None
    except (TypeError, ValueError):
        return jsonify({"ok": False, "error": "Dados de check-in inválidos."}), 400

    loc = db.session.get(Location, location_id)
    if not loc:
        return jsonify({"ok": False, "error": "Localidade inválida."}), 404

    distance_m = None
    status = "SEM_REFERENCIA"
    if acc is not None and acc > FIELD_GPS_MAX_ACCURACY_M:
        status = "BAIXA_PRECISAO"
    elif loc.reference_latitude is not None and loc.reference_longitude is not None:
        distance_m = _haversine_m(lat, lon, loc.reference_latitude, loc.reference_longitude)
        warn_m = float(os.getenv("FIELD_GPS_WARN_DISTANCE_M", "250"))
        max_m = float(os.getenv("FIELD_GPS_MAX_DISTANCE_M", "600"))
        if distance_m <= warn_m:
            status = "CONFIRMADO"
        elif distance_m <= max_m:
            status = "PROXIMO"
        else:
            status = "FORA_DA_AREA"

    row = TechnicianCheckin(
        user_id=session["user_id"], location_id=loc.id, latitude=lat, longitude=lon,
        accuracy=acc, distance_m=distance_m, status=status, created_at=datetime.utcnow()
    )
    db.session.add(row)
    db.session.commit()
    return jsonify({
        "ok": True, "status": status, "location": loc.location, "company": loc.company, "line": loc.line,
        "distance_m": round(distance_m) if distance_m is not None else None,
        "accuracy": round(acc) if acc is not None else None,
        "created_at": row.created_at.isoformat()+"Z"
    })


@app.get("/api/equipes/status")
@teams_view_required
def teams_status_api():
    if not (_has_access("teams.today") or _has_access("teams.map")): abort(403)
    _ensure_team_schedule_profiles()
    local_now=datetime.now(ZoneInfo("America/Sao_Paulo"))
    date_raw=(request.args.get("date") or "").strip()
    try:
        target_date=datetime.strptime(date_raw,"%Y-%m-%d").date() if date_raw else local_now.date()
    except ValueError:
        return jsonify({"ok":False,"error":"Data operacional inválida."}),400
    scheduled=_schedule_today_db(target_date)
    # V59: para datas passadas, a análise considera o encerramento daquele dia;
    # para hoje, usa o relógio real de São Paulo. Datas futuras não marcam ausência/atraso.
    is_today=(target_date==local_now.date())
    is_future=(target_date>local_now.date())
    effective_local_now=(local_now if is_today else datetime.combine(target_date,datetime.max.time(),tzinfo=ZoneInfo("America/Sao_Paulo")))
    now_utc=effective_local_now.astimezone(ZoneInfo("UTC")).replace(tzinfo=None)
    start_local=datetime.combine(target_date,datetime.min.time(),tzinfo=ZoneInfo("America/Sao_Paulo")); start_utc=start_local.astimezone(ZoneInfo("UTC")).replace(tzinfo=None); end_utc=(start_local+timedelta(days=1)).astimezone(ZoneInfo("UTC")).replace(tzinfo=None)
    user_ids={int(m["user_id"]) for m in scheduled if m.get("user_id")}
    users={u.id:u for u in User.query.filter(User.id.in_(user_ids)).all()} if user_ids else {}
    # Última posição do dia por usuário em uma consulta + join.
    pos_map={}
    if user_ids:
        sub=(db.session.query(TechnicianPosition.user_id,func.max(TechnicianPosition.captured_at).label("mx")).filter(TechnicianPosition.user_id.in_(user_ids),TechnicianPosition.captured_at>=start_utc,TechnicianPosition.captured_at<end_utc).group_by(TechnicianPosition.user_id).subquery())
        for p in db.session.query(TechnicianPosition).join(sub,and_(TechnicianPosition.user_id==sub.c.user_id,TechnicianPosition.captured_at==sub.c.mx)).all(): pos_map[p.user_id]=p
    login_map={}; login_counts={}; gps_counts={}
    if user_ids:
        for uid,first_at,n in db.session.query(SessionEvent.user_id,func.min(SessionEvent.created_at),func.count(SessionEvent.id)).filter(SessionEvent.user_id.in_(user_ids),SessionEvent.event_type=="LOGIN",SessionEvent.created_at>=start_utc,SessionEvent.created_at<end_utc).group_by(SessionEvent.user_id).all(): login_map[uid]=first_at; login_counts[uid]=int(n)
        gps_counts={uid:int(n) for uid,n in db.session.query(TechnicianPosition.user_id,func.count(TechnicianPosition.id)).filter(TechnicianPosition.user_id.in_(user_ids),TechnicianPosition.captured_at>=start_utc,TechnicianPosition.captured_at<end_utc).group_by(TechnicianPosition.user_id).all()}
    stations=Location.query.filter(Location.reference_latitude.isnot(None),Location.reference_longitude.isnot(None)).all()
    def nearest_station(lat,lon):
        if lat is None or lon is None:return None
        best=None
        for loc in stations:
            try: dist=_haversine_m(float(lat),float(lon),float(loc.reference_latitude),float(loc.reference_longitude))
            except Exception: continue
            if best is None or dist<best[0]: best=(dist,loc)
        if not best:return None
        dist,loc=best; return {"id":loc.id,"name":loc.location,"company":loc.company or "","line":loc.line or "","distance_m":round(dist),"relation":"NA ESTAÇÃO" if dist<=500 else "MAIS PRÓXIMA"}
    rows=[]; summary={"in_operation":0,"late":0,"not_logged":0,"stale_gt10":0,"no_gps":0,"outside_locality":0,"not_started":0}
    for member in scheduled:
        uid=member.get("user_id"); user=users.get(uid); pos=pos_map.get(uid); minutes=max(0,int((now_utc-pos.captured_at).total_seconds()//60)) if pos else None
        shift=(member.get("shift") or member.get("entry") or "").strip(); m=re.search(r'(\d{1,2}):(\d{2})',shift); expected_local=None
        if m: expected_local=datetime.combine(target_date,datetime.min.time(),tzinfo=ZoneInfo("America/Sao_Paulo")).replace(hour=int(m.group(1)),minute=int(m.group(2)))
        first_at=login_map.get(uid); login_local=first_at.replace(tzinfo=ZoneInfo("UTC")).astimezone(ZoneInfo("America/Sao_Paulo")) if first_at else None; late_minutes=max(0,int((login_local-expected_local).total_seconds()//60)) if login_local and expected_local else 0; station=nearest_station(pos.latitude,pos.longitude) if pos else None
        if (is_future or (expected_local and effective_local_now<expected_local)) and not login_local: operation_status="AINDA NÃO INICIOU"; summary["not_started"]+=1
        elif not login_local: operation_status="NÃO LOGOU"; summary["not_logged"]+=1
        elif not pos: operation_status="SEM GPS"; summary["no_gps"]+=1
        elif minutes is not None and minutes>10: operation_status="SEM POSIÇÃO >10 MIN"; summary["stale_gt10"]+=1
        elif late_minutes>0: operation_status=f"ATRASADO {late_minutes} MIN"; summary["late"]+=1
        else: operation_status="EM OPERAÇÃO"; summary["in_operation"]+=1
        if station and station["relation"]=="MAIS PRÓXIMA": summary["outside_locality"]+=1
        freshness="SEM SINAL" if minutes is None else ("ATUAL" if minutes<=5 else ("ATENÇÃO" if minutes<=15 else "ATRASADO"))
        rows.append({**member,"gps_points_today":gps_counts.get(uid,0),"session_events_today":login_counts.get(uid,0),"photo_url":(f"/usuarios/{user.id}/foto" if user and user.photo_url else None),"photo_version":(str(user.photo_url) if user and user.photo_url else None),"latitude":pos.latitude if pos else None,"longitude":pos.longitude if pos else None,"accuracy":pos.accuracy if pos else None,"captured_at":(pos.captured_at.isoformat()+"Z") if pos else None,"minutes_since":minutes,"freshness":freshness,"first_login":login_local.strftime("%H:%M") if login_local else None,"late_minutes":late_minutes,"operation_status":operation_status,"nearest_station":station,"current_location":station["name"] if station else None})
    counts={}
    for row in rows: counts[row["category"]]=counts.get(row["category"],0)+1
    return jsonify({"ok":True,"date":target_date.isoformat(),"time":(local_now.strftime("%H:%M") if is_today else "23:59"),"is_today":is_today,"is_future":is_future,"scheduled":len(rows),"counts_by_category":counts,"summary":summary,"technicians":rows})


@app.get("/api/equipes/colaboradores")
@teams_view_required
def teams_collaborators_api():
    users=User.query.filter(User.active.is_(True)).order_by(User.name).all()
    profile_by_user = {p.user_id: p for p in TeamScheduleProfile.query.filter(TeamScheduleProfile.user_id.isnot(None), TeamScheduleProfile.active.is_(True)).all()}
    out=[]
    for u in users:
        if normalize(u.personnel_status or "ATIVO")!="ATIVO":
            continue
        p=profile_by_user.get(u.id)
        out.append({
            "id":u.id,"name":u.name,"role":u.role,"job_title":u.job_title or "",
            "company":u.company or "","schedule_type": (p.schedule_type if p else ""),
            "shift": (p.shift if p else ""),"entry": (p.entry if p else ""),
            "lines": (p.lines if p else []),"profile_id": (p.id if p else None)
        })
    return jsonify({"ok":True,"users":out})


@app.get("/api/equipes/calendario")
@teams_view_required
def teams_calendar_api():
    if not _has_access("teams.schedule"): abort(403)
    """V22.1 — grade multi-dia sempre responde JSON, inclusive em erro controlado."""
    try:
        _ensure_team_schedule_profiles()
        start_raw = request.args.get("start", "").strip()
        days = max(1, min(31, request.args.get("days", type=int) or 14))
        category = normalize(request.args.get("category", ""))
        cargos = {(x or "").strip().casefold() for x in (request.args.get("cargo", "") or "").split(",") if (x or "").strip()}

        try:
            start_date = (
                datetime.strptime(start_raw, "%Y-%m-%d").date()
                if start_raw else datetime.now(ZoneInfo("America/Sao_Paulo")).date()
            )
        except ValueError:
            return jsonify({"ok": False, "error": "Data inicial inválida."}), 400

        profiles = TeamScheduleProfile.query.filter_by(active=True).order_by(
            TeamScheduleProfile.category, TeamScheduleProfile.name
        ).all()

        # V56-D: carrega os usuários vinculados em lote. Evita 1-2 SELECTs por perfil
        # em cada atualização da grade de escala.
        profile_user_ids = {p.user_id for p in profiles if p.user_id}
        calendar_users = {u.id: u for u in User.query.filter(User.id.in_(profile_user_ids)).all()} if profile_user_ids else {}
        active_profiles = []
        for p in profiles:
            user = calendar_users.get(p.user_id) if p.user_id else None
            if p.user_id and (not user or not user.active):
                continue
            if category and normalize(p.category) != category:
                continue
            if cargos and (not user or (user.job_title or "").strip().casefold() not in cargos):
                continue
            active_profiles.append(p)

        date_list = [start_date + timedelta(days=i) for i in range(days)]
        rows = []
        for p in active_profiles:
            day_rows = []
            linked_user = calendar_users.get(p.user_id) if p.user_id else None
            personnel_status = normalize(linked_user.personnel_status or "ATIVO") if linked_user else "ATIVO"
            for d in date_list:
                scheduled = bool(_team_profile_is_scheduled(p, d)) if personnel_status == "ATIVO" else False
                day_rows.append({
                    "date": d.isoformat(),
                    "scheduled": scheduled,
                    "shift": (p.shift or "") if scheduled else "FOLGA",
                    "status_override": personnel_status if personnel_status != "ATIVO" else None,
                })
            row = _profile_to_dict(p, linked_user)
            row["days"] = day_rows
            rows.append(row)

        return jsonify({
            "ok": True,
            "release": TEAMS_RELEASE,
            "start": start_date.isoformat(),
            "requested_days": days,
            "dates": [d.isoformat() for d in date_list],
            "members": rows,
            "technicians": rows,
        })
    except Exception as exc:
        app.logger.exception("V22.1: falha no calendário de equipes: %s", exc)
        db.session.rollback()
        return jsonify({
            "ok": False,
            "release": TEAMS_RELEASE,
            "error": "Não foi possível carregar a escala. Tente novamente.",
        }), 500


@app.get("/api/equipes/perfis")
@teams_view_required
def teams_profiles_api():
    if not _has_access("teams.schedule"): abort(403)
    _ensure_team_schedule_profiles()
    profiles = TeamScheduleProfile.query.order_by(
        TeamScheduleProfile.active.desc(), TeamScheduleProfile.category, TeamScheduleProfile.name
    ).all()
    users = User.query.filter(User.active.is_(True)).order_by(User.name).all()
    users_by_id={u.id:u for u in users}
    profile_payload=[]
    for p in profiles:
        u=users_by_id.get(p.user_id) if p.user_id else None
        if session.get("role")=="dispatcher" and u and u.role in ("manager","hr"):
            continue
        if session.get("role")=="hr" and u and u.role in ("manager","consultation","dispatcher","atm_financial_admin","hr"):
            continue
        profile_payload.append({
            "profile_id":p.id,"user_id":p.user_id,"linked_user_name":u.name if u else None,
            "linked":bool(u and u.active),"name":p.name,"category":p.category or "TECNICO",
            "schedule_type":p.schedule_type or "12x36","shift":p.shift,"supervision":p.supervision or "",
            "entry":p.entry or "","lines":p.lines,"anchor_date":p.anchor_date.isoformat() if p.anchor_date else None,
            "active":bool(p.active),"company":(u.company if u else "") or "","job_title":(u.job_title if u else "") or "",
            "personnel_status":(u.personnel_status if u else "ATIVO") or "ATIVO",
            "personnel_status_note":(u.personnel_status_note if u else "") or "",
            "source":"CADASTRO_USUARIO" if u else "LEGADO_ESCALA"
        })
    return jsonify({
        "ok": True,
        "profiles": profile_payload,
        "users": [
            {
                "id": u.id,
                "name": u.name,
                "role": u.role,
                "user_code": u.user_code or "",
                "company": u.company or "",
                "job_title": u.job_title or "",
                "personnel_status": u.personnel_status or "ATIVO",
            }
            for u in users
        ],
    })


@app.post("/api/equipes/perfis")
@manager_required
def teams_profile_create_api():
    if not _has_access("teams.manage"): abort(403)
    _ensure_team_schedule_profiles()
    data = request.get_json(silent=True) or {}

    name = (data.get("name") or "").strip()
    user_id = data.get("user_id")
    if user_id:
        user = db.session.get(User, int(user_id))
        if not user:
            return jsonify({"ok": False, "error": "Usuário selecionado não existe."}), 400
        name = user.name

    if not name:
        return jsonify({"ok": False, "error": "Informe o integrante da equipe."}), 400

    existing = TeamScheduleProfile.query.filter(
        func.upper(TeamScheduleProfile.name) == normalize(name)
    ).first()
    if existing and existing.active:
        return jsonify({"ok": False, "error": "Este integrante já possui uma escala ativa."}), 409

    category = normalize(data.get("category") or "TECNICO")
    if category not in ("TECNICO", "SUPERVISOR", "APOIO"):
        category = "TECNICO"

    schedule_type = data.get("schedule_type") or ("5x2" if category == "APOIO" else "12x36")
    if schedule_type not in ("12x36", "5x2"):
        return jsonify({"ok": False, "error": "Tipo de escala inválido."}), 400

    allowed_shifts = ("05:00-17:00", "11:00-23:00", "08:00-18:00")
    shift = (data.get("shift") or ("08:00-18:00" if schedule_type == "5x2" else "05:00-17:00")).strip()
    if shift not in allowed_shifts:
        return jsonify({"ok": False, "error": "Turno inválido."}), 400

    try:
        anchor = datetime.strptime(
            data.get("anchor_date") or datetime.now(ZoneInfo("America/Sao_Paulo")).date().isoformat(),
            "%Y-%m-%d"
        ).date()
    except ValueError:
        return jsonify({"ok": False, "error": "Primeiro dia da escala inválido."}), 400

    lines = data.get("lines") or []
    if isinstance(lines, str):
        lines = [x.strip() for x in re.split(r"[,;/]+", lines) if x.strip()]

    if existing and not existing.active:
        row = existing
        row.active = True
        row.user_id = int(user_id) if user_id else None
        row.category = category
        row.schedule_type = schedule_type
        row.shift = shift
        row.supervision = (data.get("supervision") or "").strip()
        row.entry = (data.get("entry") or "").strip()
        row.lines_json = json.dumps(lines, ensure_ascii=False)
        row.anchor_date = anchor
        row.updated_at = datetime.utcnow()
    else:
        row = TeamScheduleProfile(
            user_id=int(user_id) if user_id else None,
            name=name,
            active=True,
            category=category,
            schedule_type=schedule_type,
            shift=shift,
            supervision=(data.get("supervision") or "").strip(),
            entry=(data.get("entry") or "").strip(),
            lines_json=json.dumps(lines, ensure_ascii=False),
            anchor_date=anchor,
        )
        db.session.add(row)

    db.session.commit()
    return jsonify({"ok": True, "profile": _profile_to_dict(row)})


@app.put("/api/equipes/perfis/<int:profile_id>")
@manager_required
def teams_profile_update_api(profile_id):
    if not _has_access("teams.manage"): abort(403)
    _ensure_team_schedule_profiles()
    row = db.session.get(TeamScheduleProfile, profile_id)
    if not row:
        return jsonify({"ok": False, "error": "Integrante não encontrado na escala."}), 404

    data = request.get_json(silent=True) or {}

    # V6: explicit link/unlink with Users.
    if "user_id" in data:
        user_id = data.get("user_id")
        if user_id:
            user = db.session.get(User, int(user_id))
            if not user:
                return jsonify({"ok": False, "error": "Usuário selecionado não existe."}), 400
            conflict = TeamScheduleProfile.query.filter(
                TeamScheduleProfile.id != row.id,
                TeamScheduleProfile.user_id == user.id,
                TeamScheduleProfile.active.is_(True),
            ).first()
            if conflict:
                return jsonify({
                    "ok": False,
                    "error": f"Este usuário já está vinculado à escala de {conflict.name}."
                }), 409
            row.user_id = user.id
            row.name = user.name
        else:
            row.user_id = None

    if "category" in data:
        category = normalize(data.get("category"))
        if category in ("TECNICO", "SUPERVISOR", "APOIO"):
            row.category = category

    if "schedule_type" in data:
        schedule_type = data.get("schedule_type")
        if schedule_type not in ("12x36", "5x2"):
            return jsonify({"ok": False, "error": "Tipo de escala inválido."}), 400
        row.schedule_type = schedule_type

    if "shift" in data:
        shift = (data.get("shift") or "").strip()
        if shift not in ("05:00-17:00", "11:00-23:00", "08:00-18:00"):
            return jsonify({"ok": False, "error": "Turno inválido."}), 400
        row.shift = shift

    if "anchor_date" in data:
        try:
            row.anchor_date = datetime.strptime(data.get("anchor_date") or "", "%Y-%m-%d").date()
        except ValueError:
            return jsonify({"ok": False, "error": "Primeiro dia da escala inválido."}), 400

    # V6: entry can be changed independently.
    if "entry" in data:
        row.entry = (data.get("entry") or "").strip()

    if "supervision" in data:
        row.supervision = (data.get("supervision") or "").strip()

    if "lines" in data:
        lines = data.get("lines") or []
        if isinstance(lines, str):
            lines = [x.strip() for x in re.split(r"[,;/]+", lines) if x.strip()]
        row.lines_json = json.dumps(lines, ensure_ascii=False)

    if "active" in data:
        row.active = bool(data.get("active"))

    row.updated_at = datetime.utcnow()
    db.session.commit()
    return jsonify({"ok": True, "profile": _profile_to_dict(row)})


@app.delete("/api/equipes/perfis/<int:profile_id>")
@manager_required
def teams_profile_remove_api(profile_id):
    if not _has_access("teams.manage"): abort(403)
    _ensure_team_schedule_profiles()
    row = db.session.get(TeamScheduleProfile, profile_id)
    if not row:
        return jsonify({"ok": False, "error": "Integrante não encontrado na escala."}), 404
    row.active = False
    row.updated_at = datetime.utcnow()
    db.session.commit()
    return jsonify({"ok": True})


@app.get("/api/equipes/export/excel")
@dashboard_required
def teams_export_excel_api():
    if not _has_access("teams.export"): abort(403)
    _ensure_team_schedule_profiles()
    start_raw = request.args.get("start", "").strip()
    days = max(1, min(31, request.args.get("days", type=int) or 14))
    category = normalize(request.args.get("category", ""))

    try:
        start_date = (
            datetime.strptime(start_raw, "%Y-%m-%d").date()
            if start_raw else datetime.now(ZoneInfo("America/Sao_Paulo")).date()
        )
    except ValueError:
        return jsonify({"ok": False, "error": "Data inicial inválida."}), 400

    profiles = TeamScheduleProfile.query.filter_by(active=True).order_by(
        TeamScheduleProfile.category, TeamScheduleProfile.name
    ).all()
    if category:
        profiles = [p for p in profiles if normalize(p.category) == category]

    dates = [start_date + timedelta(days=i) for i in range(days)]
    wb = Workbook()
    ws = wb.active
    ws.title = "Escala por Dia"

    headers = ["Categoria", "Nome", "Usuário vinculado", "Escala", "Turno", "Entrada", "Linhas", "Supervisão"]
    headers += [d.strftime("%d/%m/%Y") for d in dates]
    ws.append(headers)

    for cell in ws[1]:
        cell.font = Font(bold=True)
        cell.fill = PatternFill("solid", fgColor="DCEAF7")
        cell.alignment = Alignment(horizontal="center")

    for p in profiles:
        user = db.session.get(User, p.user_id) if p.user_id else None
        row = [
            p.category or "TECNICO",
            p.name,
            user.name if user else "NÃO VINCULADO",
            p.schedule_type or "12x36",
            p.shift,
            p.entry or "",
            " / ".join(p.lines),
            p.supervision or "",
        ]
        for d in dates:
            row.append(p.shift if _team_profile_is_scheduled(p, d) else "FOLGA")
        ws.append(row)

    # Color scheduled / off.
    for r in range(2, ws.max_row + 1):
        for c in range(9, ws.max_column + 1):
            value = ws.cell(r, c).value
            if value == "FOLGA":
                ws.cell(r, c).fill = PatternFill("solid", fgColor="F2F4F7")
            else:
                ws.cell(r, c).fill = PatternFill("solid", fgColor="E5F6EC")

    ws.freeze_panes = "I2"
    ws.auto_filter.ref = ws.dimensions
    for idx, width in {
        1: 15, 2: 32, 3: 30, 4: 12, 5: 16, 6: 20, 7: 18, 8: 24
    }.items():
        ws.column_dimensions[get_column_letter(idx)].width = width
    for c in range(9, ws.max_column + 1):
        ws.column_dimensions[get_column_letter(c)].width = 13

    # Summary.
    ws2 = wb.create_sheet("Resumo")
    ws2.append(["Categoria", "Ativos na escala"])
    for cell in ws2[1]:
        cell.font = Font(bold=True)
        cell.fill = PatternFill("solid", fgColor="DCEAF7")
    for cat in ("TECNICO", "SUPERVISOR", "APOIO"):
        ws2.append([cat.title(), sum(1 for p in profiles if normalize(p.category) == cat)])
    ws2.append([])
    ws2.append(["Período", f"{dates[0].strftime('%d/%m/%Y')} a {dates[-1].strftime('%d/%m/%Y')}"])
    ws2.append(["Gerado em", datetime.now(ZoneInfo("America/Sao_Paulo")).strftime("%d/%m/%Y %H:%M")])

    # Link audit.
    ws3 = wb.create_sheet("Vínculo Usuários")
    ws3.append(["Categoria", "Nome escala", "Usuário", "ID usuário", "Situação vínculo"])
    for cell in ws3[1]:
        cell.font = Font(bold=True)
        cell.fill = PatternFill("solid", fgColor="DCEAF7")
    for p in profiles:
        user = db.session.get(User, p.user_id) if p.user_id else None
        ws3.append([
            p.category or "TECNICO",
            p.name,
            user.name if user else "",
            user.id if user else "",
            "VINCULADO" if user else "NÃO VINCULADO",
        ])

    buffer = io.BytesIO()
    wb.save(buffer)
    buffer.seek(0)
    filename = f"autopass_escala_{start_date.strftime('%Y%m%d')}_{days}dias.xlsx"
    return send_file(
        buffer,
        mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        as_attachment=True,
        download_name=filename,
    )



# ============================================================
# V7.0 — PATRIMÔNIO 360
# ============================================================

def _location_360_payload(loc):
    inventories = Inventory.query.filter_by(location_id=loc.id).order_by(Inventory.created_at.desc()).all()
    visits = FieldEvidenceVisit.query.filter_by(location_id=loc.id).order_by(FieldEvidenceVisit.id.desc()).all()
    expected_bucket = _expected_assets_by_location().get(loc.id, {})
    expected = sum(int(expected_bucket.get(t, 0) or 0) for t in ("ATM","VALIDADOR","POS","TDI","BLOQUEIO"))
    inventoried = len(inventories)
    outside_base = sum(1 for x in inventories if normalize(x.in_base) in ("NAO","NÃO","NO","FORA DA BASE") or _canonical_equipment_type(x.equipment_type)=="OUTRO")
    divergences = sum(1 for x in inventories if (x.divergence or "").strip())
    inoperative = sum(1 for x in inventories if normalize(x.operational_status) not in ("", "OPERACIONAL", "OK"))
    media_count = 0
    if visits:
        visit_ids=[v.id for v in visits]
        media_count = FieldEvidenceMedia.query.filter(FieldEvidenceMedia.visit_id.in_(visit_ids)).count()
    return {
        "id":loc.id, "company":loc.company, "line":loc.line, "location":loc.location,
        "survey_status":loc.survey_status,
        "expected":expected, "inventoried":inventoried, "missing":max(0,expected-inventoried),
        "coverage_pct":round((inventoried/expected*100),1) if expected else 0,
        "divergences":divergences, "inoperative":inoperative,
        "visits":len(visits), "media":media_count, "outside_base":outside_base,
        "latitude":loc.reference_latitude, "longitude":loc.reference_longitude,
        "can_manage": session.get("role") == "manager",
        "equipment":[{
            "id":x.id, "type":x.equipment_type, "identifier":x.asset_identifier,
            "serial":x.serial, "model":x.model, "supplier":x.supplier,
            "status":x.operational_status, "divergence":x.divergence, "in_base":x.in_base,
            "notes":x.notes, "technician_id":x.technician_id,
            "creator": ((db.session.get(User, x.technician_id).name if db.session.get(User, x.technician_id) else None) or "—"),
            "creator_username": ((db.session.get(User, x.technician_id).username if db.session.get(User, x.technician_id) else None) or ""),
            "created_at":x.created_at.isoformat()+"Z" if x.created_at else None,
            "media":[{
                "id":a.id,"name":a.original_name,"mime":a.mime_type,
                "url":url_for("uploaded", name=a.stored_name)
            } for a in Attachment.query.filter_by(inventory_id=x.id).order_by(Attachment.id).all()]
        } for x in inventories[:250]],
        "evidence_visits":[{
            "id":v.id, "date":v.source_date, "time":v.source_time, "author":v.author,
            "confidence":v.match_confidence, "report":v.report_text,
            "media":[{
                "id":m.id,
                "name":m.original_name,
                "mime":m.mime_type,
                "url":url_for("field_evidence_media", media_id=m.id)
            } for m in FieldEvidenceMedia.query.filter_by(visit_id=v.id).order_by(FieldEvidenceMedia.id).all()]
        } for v in visits[:100]]
    }


@app.get("/api/v7/localidades")
@dashboard_required
def v7_locations_api():
    q=(request.args.get("q") or "").strip()
    query=Location.query
    if q:
        like=f"%{q}%"
        query=query.filter(db.or_(
            Location.location.ilike(like), Location.line.ilike(like), Location.company.ilike(like)
        ))
    rows=query.order_by(Location.company,Location.line,Location.location).limit(500).all()
    return jsonify({"ok":True,"locations":[_location_360_payload(x) for x in rows]})


@app.get("/api/v7/localidades/<int:location_id>")
@dashboard_required
def v7_location_detail_api(location_id):
    loc=db.session.get(Location,location_id)
    if not loc:
        return jsonify({"ok":False,"error":"Localidade não encontrada."}),404
    return jsonify({"ok":True,"location":_location_360_payload(loc)})


def _asset_360_payload(asset):
    inventories=Inventory.query.filter(
        db.or_(Inventory.base_asset_id==asset.id,
               Inventory.serial==asset.serial if asset.serial else db.false())
    ).order_by(Inventory.created_at.desc()).all()
    evidence=FieldEvidenceItem.query.filter(
        db.or_(FieldEvidenceItem.base_asset_id==asset.id,
               FieldEvidenceItem.serial==asset.serial if asset.serial else db.false())
    ).order_by(FieldEvidenceItem.id.desc()).all()
    latest=inventories[0] if inventories else None
    loc=db.session.get(Location,latest.location_id) if latest else None
    return {
        "id":asset.id,"asset_key":asset.asset_key,"type":asset.equipment_type,
        "serial":asset.serial,"qrcode_id":asset.qrcode_id,"top_id":asset.top_id,
        "model":asset.model,"supplier":asset.supplier,
        "company":asset.company,"line":asset.line,"locality":asset.locality,
        "base_status":asset.base_status,
        "latest_inventory":{
            "id":latest.id,"location_id":latest.location_id,
            "location":loc.location if loc else None,
            "line":loc.line if loc else None,
            "status":latest.operational_status,
            "divergence":latest.divergence,
            "technician_id":latest.technician_id,
            "created_at":latest.created_at.isoformat()+"Z" if latest.created_at else None
        } if latest else None,
        "history":[{
            "inventory_id":x.id,"location_id":x.location_id,"status":x.operational_status,
            "divergence":x.divergence,"created_at":x.created_at.isoformat()+"Z" if x.created_at else None
        } for x in inventories[:100]],
        "evidence":[{
            "id":e.id,"visit_id":e.visit_id,"status":e.audit_status,
            "detail":e.audit_detail,"identifier":e.identifier
        } for e in evidence[:100]]
    }


@app.get("/api/v7/equipamentos/<int:asset_id>")
@dashboard_required
def v7_asset_detail_api(asset_id):
    asset=db.session.get(BaseAsset,asset_id)
    if not asset:
        return jsonify({"ok":False,"error":"Equipamento não encontrado."}),404
    return jsonify({"ok":True,"asset":_asset_360_payload(asset)})


@app.get("/api/v7/busca")
@dashboard_required
def v7_global_search_api():
    q=(request.args.get("q") or "").strip()
    if len(q)<2:
        return jsonify({"ok":True,"query":q,"locations":[],"assets":[],"inventory":[]})
    like=f"%{q}%"
    locations=Location.query.filter(db.or_(
        Location.location.ilike(like),Location.line.ilike(like),Location.company.ilike(like)
    )).limit(15).all()
    assets=BaseAsset.query.filter(db.or_(
        BaseAsset.serial.ilike(like),BaseAsset.asset_key.ilike(like),
        BaseAsset.qrcode_id.ilike(like),BaseAsset.top_id.ilike(like),
        BaseAsset.locality.ilike(like),BaseAsset.description.ilike(like)
    )).limit(25).all()
    matching_location_ids=[x.id for x in locations]
    creator_rows=User.query.filter(db.or_(User.name.ilike(like),User.username.ilike(like),User.user_code.ilike(like))).limit(50).all()
    creator_ids=[u.id for u in creator_rows]
    inventory_filters=[
        Inventory.serial.ilike(like),Inventory.asset_identifier.ilike(like),
        Inventory.model.ilike(like),Inventory.supplier.ilike(like),
        Inventory.equipment_type.ilike(like), Inventory.notes.ilike(like)
    ]
    if creator_ids:
        inventory_filters.append(Inventory.technician_id.in_(creator_ids))
    if matching_location_ids:
        inventory_filters.append(Inventory.location_id.in_(matching_location_ids))
    inventory=Inventory.query.filter(db.or_(*inventory_filters)).order_by(Inventory.created_at.desc()).limit(50).all()
    creator_map={u.id:u for u in User.query.filter(User.id.in_([x.technician_id for x in inventory] or [-1])).all()}
    return jsonify({
        "ok":True,"query":q,"can_admin":_current_user_is_superadmin(),
        "locations":[{"id":x.id,"company":x.company,"line":x.line,"location":x.location} for x in locations],
        "assets":[{"id":x.id,"type":x.equipment_type,"serial":x.serial,"asset_key":x.asset_key,
                   "locality":x.locality,"line":x.line,"model":x.model} for x in assets],
        "inventory":[{"id":x.id,"type":x.equipment_type,"identifier":x.asset_identifier,
                      "serial":x.serial,"location_id":x.location_id,"status":x.operational_status,
                      "in_base":x.in_base,"model":x.model,
                      "creator":creator_map.get(x.technician_id).name if creator_map.get(x.technician_id) else "—",
                      "creator_username":creator_map.get(x.technician_id).username if creator_map.get(x.technician_id) else "",
                      "created_at":x.created_at.isoformat()+"Z" if x.created_at else None} for x in inventory]
    })


def _v12_location_intelligence(loc):
    inventories = Inventory.query.filter_by(location_id=loc.id).order_by(Inventory.created_at.desc()).all()
    expected = int(loc.expected_atm or 0) + int(loc.expected_validator or 0) + int(loc.expected_pos or 0)
    found = len(inventories)
    missing = max(0, expected - found)
    divergences = sum(1 for x in inventories if (x.divergence or "").strip())
    outside_base = sum(1 for x in inventories if normalize(x.in_base or "") in ("NAO", "NÃO", "FORA DA BASE"))
    bad_gps = sum(1 for x in inventories if x.gps_accuracy is not None and x.gps_accuracy > FIELD_GPS_MAX_ACCURACY_M)
    no_gps = sum(1 for x in inventories if x.latitude is None or x.longitude is None)
    media = Attachment.query.join(Inventory, Attachment.inventory_id == Inventory.id).filter(Inventory.location_id == loc.id).count()
    coverage = min(100.0, (found / expected * 100.0)) if expected else (100.0 if found else 0.0)
    # Score explicável: cobertura 60%, integridade 20%, GPS 10%, evidência 10%.
    integrity = max(0.0, 100.0 - (divergences + outside_base) * 12.5)
    gps_quality = 100.0 if found == 0 else max(0.0, 100.0 - ((bad_gps + no_gps) / max(found,1) * 100.0))
    evidence_quality = 100.0 if found and media >= found else (media / max(found,1) * 100.0 if found else 0.0)
    score = round(coverage * .60 + integrity * .20 + gps_quality * .10 + evidence_quality * .10, 1)
    if score >= 85: level = "CONFIAVEL"
    elif score >= 65: level = "ATENCAO"
    else: level = "CRITICO"
    return {"id":loc.id,"company":loc.company,"line":loc.line,"location":loc.location,
            "expected":expected,"found":found,"missing":missing,"coverage":round(coverage,1),
            "divergences":divergences,"outside_base":outside_base,"bad_gps":bad_gps,"media":media,
            "score":score,"level":level,"survey_status":loc.survey_status}

@app.get("/api/v12/resumo")
@dashboard_required
def v12_summary_api():
    rows=[_v12_location_intelligence(x) for x in Location.query.order_by(Location.company,Location.line,Location.location).all()]
    return jsonify({"ok":True,"release":APP_RELEASE,"locations":len(rows),
        "reliable":sum(1 for x in rows if x["level"]=="CONFIAVEL"),
        "attention":sum(1 for x in rows if x["level"]=="ATENCAO"),
        "critical":sum(1 for x in rows if x["level"]=="CRITICO"),
        "missing":sum(x["missing"] for x in rows),"divergences":sum(x["divergences"] for x in rows),
        "outside_base":sum(x["outside_base"] for x in rows),
        "average_score":round(sum(x["score"] for x in rows)/len(rows),1) if rows else 0})

@app.get("/api/v12/localidades")
@dashboard_required
def v12_locations_api():
    rows=[_v12_location_intelligence(x) for x in Location.query.order_by(Location.company,Location.line,Location.location).all()]
    rows.sort(key=lambda x:(x["score"],-x["missing"],x["location"]))
    return jsonify({"ok":True,"locations":rows})



def _parse_contract_date(value):
    raw=(value or "").strip()
    if not raw:
        return None
    for fmt in ("%Y-%m-%d","%d/%m/%Y","%d-%m-%Y"):
        try:
            return datetime.strptime(raw,fmt).date()
        except Exception:
            pass
    return None


def _v551_contract_reference():
    """Referência contratual validada na planilha ATM com vencimentos atualizados.
    Esperado: 412 ATMs em Dez/2026, 100 em Fev/2028 e demais sem contrato.
    """
    try:
        rows=json.loads((DATA_DIR / "atm_contract_reference_v551.json").read_text(encoding="utf-8"))
    except Exception:
        rows=[]
    by_asset={}; by_top={}
    for r in rows:
        if str(r.get("asset_key") or "").strip(): by_asset[str(r.get("asset_key")).strip()]=r
        if str(r.get("id_top") or "").strip(): by_top[str(r.get("id_top")).strip()]=r
    return rows,by_asset,by_top

def _v551_apply_contract_reference(rows):
    _,by_asset,by_top=_v551_contract_reference()
    out=[]
    for src in rows:
        a=dict(src)
        ref=by_asset.get(str(a.get("asset_key") or "").strip()) or by_top.get(str(a.get("id_top") or "").strip())
        if ref:
            a["contract"]=ref.get("contract_label") or "Sem contrato"
            a["contract_end"]=ref.get("contract_end") or ""
        else:
            a["contract"]="Sem contrato"
            a["contract_end"]=""
        out.append(a)
    return out

@app.get("/api/dashboard/inventory-atm")
@dashboard_required
def inventory_atm_dashboard_api():
    if session.get("role")=="technician":
        return jsonify({"ok":False,"error":"Dashboard ATM restrita à gestão."}),403
    # V42.4.1: a Dashboard ATM usa a planilha oficial 08/2026 como universo mestre.
    # Base oficial: 590 alocados + 12 estoque = 602 ATMs. Levantamentos de campo não alteram esse universo.
    official_path=DATA_DIR / "atm_official_082026.json"
    try:
        all_rows=json.loads(official_path.read_text(encoding="utf-8"))
    except Exception:
        all_rows=[]
    all_rows=_v551_apply_contract_reference(all_rows)
    # V55.4: modelo ATM é uma dimensão controlada. IDs/terminais numéricos não podem contaminar o filtro Modelo.
    valid_atm_models={"TCI","MK","MKNEO","TCINEO","MINIWALL","TCIPLUS","DCASH"}
    for _a in all_rows:
        _raw_model=str(_a.get("model") or "").strip().upper()
        _a["model_raw"]=_raw_model
        _a["model"]=_raw_model if _raw_model in valid_atm_models else "Modelo não identificado"
    filters={k:(request.args.get(k) or "").strip() for k in ("company","line","locality","model","contract","ownership","status")}
    teamviewer_missing=(request.args.get("teamviewer_missing") or "").strip() in ("1","true","TRUE","sim","SIM")
    field_map={"company":"company","line":"line","locality":"locality","model":"model","contract":"contract","ownership":"ownership","status":"status"}
    rows=[x for x in all_rows if all(not filters[k] or str(x.get(field_map[k],""))==filters[k] for k in filters)]
    if teamviewer_missing:
        rows=[x for x in rows if not str(x.get("teamviewer_id") or "").strip()]
    def agg(attr):
        out={}
        for a in rows:
            v=(str(a.get(attr) or "Não informado")).strip() or "Não informado"
            out[v]=out.get(v,0)+1
        return out
    # V44.0: filtros facetados/bidirecionais. Cada combo oferece apenas valores
    # compatíveis com os demais filtros ativos, como em uma ferramenta de BI.
    def facet_options(filter_key, attr):
        candidates=[]
        for a in all_rows:
            if teamviewer_missing and str(a.get("teamviewer_id") or "").strip():
                continue
            ok=True
            for k,v in filters.items():
                if k==filter_key or not v:
                    continue
                if str(a.get(field_map[k],"")).strip()!=v:
                    ok=False; break
            if ok:
                val=str(a.get(attr) or "").strip()
                if val: candidates.append(val)
        return sorted(set(candidates),key=lambda x:x.casefold())
    return jsonify({"ok":True,"release":APP_RELEASE,"source":"INVENTARIO AUTOPASS - EQUIPAMENTOS DE CAMPO - 082026.xlsm / aba ATM",
        "official_total":602,"official_allocated":590,"official_stock":12,"total":len(rows),
        "allocated":sum(1 for x in rows if not x.get("stock")),"stock":sum(1 for x in rows if x.get("stock")),
        "operators":agg("company"),"models":agg("model"),"contracts":agg("contract"),"ownership":agg("ownership"),"locations":agg("locality"),"lines":agg("line"),
        "cptm_stations":len({str(x.get("locality") or "").strip() for x in rows if str(x.get("company") or "").upper()=="CPTM" and str(x.get("locality") or "").strip()}),
        "metro_stations":len({str(x.get("locality") or "").strip() for x in rows if str(x.get("company") or "").upper() in ("METRÔ","METRO") and str(x.get("locality") or "").strip()}),
        "teamviewer_count":sum(1 for x in rows if str(x.get("teamviewer_id") or "").strip()),"assets":rows,
        "options":{"companies":facet_options("company","company"),"lines":facet_options("line","line"),"localities":facet_options("locality","locality"),"models":facet_options("model","model"),"contracts":facet_options("contract","contract"),"ownership":facet_options("ownership","ownership"),"statuses":facet_options("status","status")}})

@app.get("/api/dashboard/atm-financial")
@login_required
def atm_financial_dashboard_api():
    # V46: visão financeira restrita ao ADM/Gestor principal.
    if session.get("role") not in ("manager", "manager_field", "atm_financial_admin"):
        return jsonify({"ok":False,"error":"Dashboard financeira restrita aos perfis autorizados."}),403
    path=DATA_DIR / "atm_financial_082026.json"
    try:
        payload=json.loads(path.read_text(encoding="utf-8"))
    except Exception as exc:
        return jsonify({"ok":False,"error":f"Base financeira ATM indisponível: {exc}"}),500
    # V51: custos operacionais por fornecedor, preservando origem e descrição da atividade.
    supplier_path=DATA_DIR / "atm_supplier_costs_v51.json"
    try:
        supplier_payload=json.loads(supplier_path.read_text(encoding="utf-8"))
    except Exception:
        supplier_payload={"suppliers":[],"totals":{}}
    competences=[x.strip() for x in request.args.getlist("competence") if x.strip()]
    if not competences:
        raw=(request.args.get("competences") or request.args.get("competence") or "").strip()
        competences=[x.strip() for x in raw.split(",") if x.strip()]
    q=FinancialMonthlyCost.query
    if competences: q=q.filter(FinancialMonthlyCost.competence.in_(competences))
    dynamic=[]; allocated={"ATM":0.0,"POS":0.0,"RECARGA":0.0,"RACK":0.0,"OUTROS":0.0}; forecast_allocated={"ATM":0.0,"POS":0.0,"RECARGA":0.0,"RACK":0.0,"OUTROS":0.0}
    sups={x.id:x.name for x in FinancialSupplier.query.all()}; svcs={x.id:x.name for x in FinancialService.query.all()}
    for x in q.order_by(FinancialMonthlyCost.competence.desc()).all():
        try: alloc=json.loads(x.allocation_json or "{}")
        except: alloc={}
        for k,pct in alloc.items():
            key=str(k).upper(); allocated[key]=allocated.get(key,0.0)+float(x.amount or 0)*float(pct or 0)/100.0
            if getattr(x,"forecast_amount",None) is not None:
                forecast_allocated[key]=forecast_allocated.get(key,0.0)+float(x.forecast_amount or 0)*float(pct or 0)/100.0
        dynamic.append({"id":x.id,"competence":x.competence,"supplier":sups.get(x.supplier_id,""),"service":getattr(x,"service_text",None) or svcs.get(x.service_id,""),"amount":round(float(x.amount or 0),2),"forecast_amount":None if getattr(x,"forecast_amount",None) is None else round(float(x.forecast_amount or 0),2),"cost_center":getattr(x,"cost_center",None) or "SUPORTE_CAMPO","project":getattr(x,"project",None) or "","allocation":alloc})
    all_competences=sorted({x[0] for x in db.session.query(FinancialMonthlyCost.competence).filter(FinancialMonthlyCost.competence.isnot(None)).all() if x and x[0]})
    # V55.1: havendo lançamentos internos, a dashboard deixa de depender da planilha histórica de fornecedores.
    if dynamic:
        by_supplier={}
        for row in dynamic:
            atm_pct=float((row.get("allocation") or {}).get("ATM") or 0)
            atm_value=float(row.get("amount") or 0)*atm_pct/100.0
            key=(row.get("supplier") or "Sem fornecedor", row.get("service") or "Sem serviço")
            d=by_supplier.setdefault(key,{"supplier":key[0],"description":key[1],"period_value":0.0,"avg_jan_jun_2026":0.0,"category":"ATM","allocation_rule":"RATEIO","allocation_percentages":{"ATM":atm_pct}})
            d["period_value"]+=atm_value
        supplier_payload={"suppliers":[{**v,"period_value":round(v["period_value"],2)} for v in by_supplier.values()],"totals":{"period_value":round(sum(v["period_value"] for v in by_supplier.values()),2)},"source":"LANÇAMENTOS_FINANCEIROS"}
    payload.update({"ok":True,"release":APP_RELEASE,"supplier_costs":supplier_payload,"monthly_costs":dynamic,"monthly_allocated":allocated,"monthly_forecast_allocated":forecast_allocated,"monthly_total":round(sum(float(x.get("amount") or 0) for x in dynamic),2),"monthly_forecast_total":round(sum(float(x.get("forecast_amount") or 0) for x in dynamic if x.get("forecast_amount") is not None),2),"competences":all_competences})
    return jsonify(payload)

@app.get("/api/v30/atm-contracts")
@dashboard_required
def v30_atm_contracts():
    company=(request.args.get("company") or "").strip(); line=(request.args.get("line") or "").strip(); contract=(request.args.get("contract") or "").strip(); horizon=(request.args.get("horizon") or "").strip()
    # V52.6: a situação contratual usa somente a base oficial ATM importada; nenhuma data é inferida.
    try: rows=json.loads((DATA_DIR / "atm_official_082026.json").read_text(encoding="utf-8"))
    except Exception: rows=[]
    rows=_v551_apply_contract_reference(rows)
    today=datetime.utcnow().date(); out=[]
    for idx,a in enumerate(rows,1):
        if company and str(a.get("company") or "")!=company: continue
        if line and str(a.get("line") or "")!=line: continue
        c=(str(a.get("contract") or "Sem contrato")).strip() or "Sem contrato"
        if contract and c!=contract: continue
        raw_end=(a.get("contract_end") or a.get("vencimento_contrato") or a.get("venc_contrato") or "")
        end=_parse_contract_date(str(raw_end)) if raw_end else None; days=(end-today).days if end else None
        status="SEM CONTRATO" if not end else ("VENCIDO" if days<0 else ("ATÉ 30 DIAS" if days<=30 else ("31–60 DIAS" if days<=60 else ("61–90 DIAS" if days<=90 else "ACIMA DE 90 DIAS"))))
        if horizon=="expired" and status!="VENCIDO": continue
        if horizon=="none" and status!="SEM CONTRATO": continue
        out.append({"id":idx,"asset_key":a.get("asset_key") or a.get("id_top") or "","company":a.get("company") or "","line":a.get("line") or "","locality":a.get("locality") or "","serial":a.get("serial") or "","model":a.get("model") or "","supplier":a.get("supplier") or "","contract":c,"contract_end":str(raw_end or ""),"days_to_expire":days,"contract_status":status})
    contracts=sorted({x["contract"] for x in out if x["contract"]},key=lambda x:x.casefold())
    return jsonify({"ok":True,"release":APP_RELEASE,"count":len(out),"contracts":contracts,"assets":out,"source":"PLANILHA_ATM_VENCIMENTOS_ATUALIZADA","date_rule":"412 Dez/2026 · 100 Fev/2028 · demais sem contrato; sem inferência"})

@app.get("/api/v30/atm-contracts/export")
@dashboard_required
def v30_atm_contracts_export():
    company=(request.args.get("company") or "").strip(); line=(request.args.get("line") or "").strip()
    contract=(request.args.get("contract") or "").strip(); horizon=(request.args.get("horizon") or "").strip()
    try: rows=json.loads((DATA_DIR / "atm_official_082026.json").read_text(encoding="utf-8"))
    except Exception: rows=[]
    rows=_v551_apply_contract_reference(rows); today=datetime.utcnow().date(); selected=[]
    for a in rows:
        if company and str(a.get("company") or "")!=company: continue
        if line and str(a.get("line") or "")!=line: continue
        c=(str(a.get("contract") or "Sem contrato")).strip() or "Sem contrato"
        if contract and c!=contract: continue
        raw_end=str(a.get("contract_end") or ""); end=_parse_contract_date(raw_end) if raw_end else None; days=(end-today).days if end else None
        status="SEM CONTRATO" if not end else ("VENCIDO" if days<0 else ("ATÉ 30 DIAS" if days<=30 else ("31–60 DIAS" if days<=60 else ("61–90 DIAS" if days<=90 else "ACIMA DE 90 DIAS"))))
        if horizon=="expired" and status!="VENCIDO": continue
        if horizon=="none" and status!="SEM CONTRATO": continue
        selected.append((a,c,raw_end,days,status))
    wb=Workbook(); ws=wb.active; ws.title="Contratos ATM"
    headers=["Empresa","Linha","Localidade","ATM","Série","Modelo","Fornecedor","Referência contratual","Vencimento","Dias para vencer","Status contrato","Status base"]
    ws.append(headers); fill=PatternFill("solid",fgColor="17365D"); font=Font(color="FFFFFF",bold=True)
    for cell in ws[1]: cell.fill=fill; cell.font=font
    for a,c,raw_end,days,status in selected:
        ws.append([a.get("company"),a.get("line"),a.get("locality"),a.get("asset_key") or a.get("id_top"),a.get("serial"),a.get("model"),a.get("supplier"),c,raw_end,days,status,a.get("status")])
    for col in range(1,ws.max_column+1): ws.column_dimensions[get_column_letter(col)].width=20
    ws.freeze_panes="A2"; out=io.BytesIO(); wb.save(out); out.seek(0)
    return send_file(out,as_attachment=True,download_name=f"autopass_contratos_atm_{datetime.now().strftime('%Y%m%d_%H%M')}.xlsx",mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")


@app.get("/api/v35/equipments/export")
@dashboard_required
def v35_equipments_export():
    """Exporta a base de equipamentos respeitando o recorte da tela Equipamentos.

    Para ATM, os filtros de contrato/vencimento também são aplicados. Para os
    demais tipos esses filtros são ignorados, permitindo usar o mesmo botão para
    ATM, Recarga, POS, TDI, Bloqueio ou todos os tipos.
    """
    company=(request.args.get("company") or "").strip()
    line=(request.args.get("line") or "").strip()
    location_id=(request.args.get("location") or "").strip()
    equipment_type=(request.args.get("type") or "").strip().upper()
    contract=(request.args.get("contract") or "").strip()
    horizon=(request.args.get("horizon") or "").strip()

    q=BaseAsset.query
    if company:
        q=q.filter(BaseAsset.company==company)
    if line:
        q=q.filter(BaseAsset.line==line)
    if location_id:
        try:
            loc=db.session.get(Location,int(location_id))
        except Exception:
            loc=None
        if loc:
            q=q.filter(BaseAsset.locality==loc.location)

    type_variants={
        "ATM":{"ATM"},
        "VALIDADOR":{"VALIDADOR","RECARGA","VALIDADOR DE RECARGA"},
        "RECARGA":{"VALIDADOR","RECARGA","VALIDADOR DE RECARGA"},
        "POS":{"POS"},
        "TDI":{"TDI","TDI TECNICO","TDI TÉCNICO"},
        "BLOQUEIO":{"BLOQUEIO","BLOQUEIO DE ACESSO"},
        "OUTRO":{"OUTRO"},
    }
    if equipment_type:
        variants=type_variants.get(equipment_type,{equipment_type})
        q=q.filter(func.upper(func.coalesce(BaseAsset.equipment_type,"" )).in_(variants))

    today=datetime.utcnow().date()
    selected=[]
    for a in q.order_by(BaseAsset.company,BaseAsset.line,BaseAsset.locality,BaseAsset.equipment_type,BaseAsset.asset_key).all():
        # Contrato/vencimento são atributos específicos de ATM.
        if equipment_type=="ATM":
            if contract and (a.leasing_status or "")!=contract:
                continue
            end=_parse_contract_date(a.contract_end)
            days=(end-today).days if end else None
            if horizon=="expired" and not(end and days<0): continue
            if horizon=="30" and not(end and 0<=days<=30): continue
            if horizon=="60" and not(end and 0<=days<=60): continue
            if horizon=="90" and not(end and 0<=days<=90): continue
        else:
            end=_parse_contract_date(a.contract_end)
            days=(end-today).days if end else None
        selected.append((a,days))

    wb=Workbook(); ws=wb.active; ws.title="Equipamentos"
    headers=["Empresa","Linha","Localidade","Tipo","Ativo / patrimônio","Série","QR Code","TOP ID","Modelo","Fornecedor","Aplicação","Terminal","Versão software","Quantidade","Status base","Contrato ATM","Vencimento","Dias para vencer","Observações da base"]
    ws.append(headers)
    fill=PatternFill("solid",fgColor="17365D"); font=Font(color="FFFFFF",bold=True)
    for c in ws[1]: c.fill=fill; c.font=font
    for a,days in selected:
        ws.append([a.company,a.line,a.locality,a.equipment_type,a.asset_key,a.serial,a.qrcode_id,a.top_id,a.model,a.supplier,a.application,a.terminal_number,a.software_version,a.quantity,a.base_status,a.leasing_status,a.contract_end,days,a.base_notes])
    for col in range(1,ws.max_column+1):
        ws.column_dimensions[get_column_letter(col)].width=22
    ws.freeze_panes="A2"; ws.auto_filter.ref=ws.dimensions
    out=io.BytesIO(); wb.save(out); out.seek(0)
    label=(equipment_type or "todos").lower().replace(" ","-")
    return send_file(out,as_attachment=True,download_name=f"autopass_equipamentos_{label}_{datetime.now().strftime('%Y%m%d_%H%M')}.xlsx",mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

@app.get("/api/v12/inventario/<int:inventory_id>/ciclo")
@dashboard_required
def v12_lifecycle_api(inventory_id):
    inv=db.session.get(Inventory,inventory_id)
    if not inv: return jsonify({"ok":False,"error":"Equipamento não encontrado."}),404
    events=AssetLifecycleEvent.query.filter_by(inventory_id=inventory_id).order_by(AssetLifecycleEvent.created_at.desc()).all()
    return jsonify({"ok":True,"inventory_id":inventory_id,"events":[{"id":e.id,"type":e.event_type,"status":e.status,
        "from_location_id":e.from_location_id,"to_location_id":e.to_location_id,"notes":e.notes,
        "user_id":e.user_id,"created_at":e.created_at.isoformat()+"Z" if e.created_at else None} for e in events]})

@app.post("/api/v12/inventario/<int:inventory_id>/ciclo")
@manager_required
def v12_lifecycle_create_api(inventory_id):
    inv=db.session.get(Inventory,inventory_id)
    if not inv: return jsonify({"ok":False,"error":"Equipamento não encontrado."}),404
    data=request.get_json(silent=True) or {}
    event_type=normalize(data.get("event_type") or "")
    allowed={"ENCONTRADO","INSTALADO","MOVIMENTADO","SUBSTITUIDO","RETIRADO","MANUTENCAO","AUDITORIA"}
    if event_type not in allowed: return jsonify({"ok":False,"error":"Tipo de evento inválido."}),400
    row=AssetLifecycleEvent(inventory_id=inventory_id,event_type=event_type,
        from_location_id=data.get("from_location_id"),to_location_id=data.get("to_location_id") or inv.location_id,
        status=(data.get("status") or "REGISTRADO")[:80],notes=(data.get("notes") or "")[:4000],user_id=session["user_id"])
    db.session.add(row); db.session.commit()
    return jsonify({"ok":True,"id":row.id})

@app.route("/meu-perfil", methods=["GET","POST"])
@login_required
def my_profile_page():
    user=db.session.get(User,session["user_id"])
    if request.method=="POST":
        action=(request.form.get("action") or "photo").strip()
        if action=="password":
            current=request.form.get("current_password") or ""
            new_password=request.form.get("new_password") or ""
            confirm=request.form.get("confirm_password") or ""
            if not check_password_hash(user.password_hash,current):
                flash("Senha atual incorreta."); return redirect(url_for("my_profile_page"))
            if len(new_password)<8:
                flash("A nova senha deve ter pelo menos 8 caracteres."); return redirect(url_for("my_profile_page"))
            if new_password!=confirm:
                flash("A confirmação da nova senha não confere."); return redirect(url_for("my_profile_page"))
            user.password_hash=generate_password_hash(new_password); db.session.commit()
            flash("Senha alterada com sucesso."); return redirect(url_for("my_profile_page"))
        photo=request.files.get("photo")
        if not photo or not photo.filename:
            flash("Selecione uma imagem."); return redirect(url_for("my_profile_page"))
        if not (photo.mimetype or "").startswith("image/"):
            flash("A foto deve ser uma imagem."); return redirect(url_for("my_profile_page"))
        data=photo.read()
        if len(data)>3*1024*1024:
            flash("A foto deve ter no máximo 3 MB."); return redirect(url_for("my_profile_page"))
        old=user.photo_url
        safe_name=secure_filename(photo.filename) or "foto.jpg"
        key=f"usuarios/{user.user_code or user.id}/{uuid.uuid4().hex}-{safe_name}"
        _r2_put_bytes(key,data,photo.mimetype or "image/jpeg")
        user.photo_url=key; db.session.commit()
        if old and old!=key:
            try: r2_client().delete_object(Bucket=os.environ["R2_BUCKET_NAME"],Key=old)
            except Exception: pass
        flash("Foto atualizada com sucesso."); return redirect(url_for("my_profile_page"))
    return render_template("profile.html",user=user,customer_companies=_customer_companies_for_user(user) if user.role=="customer" else [])

@app.get("/patrimonio")
@dashboard_required
def patrimonio_page():
    if not _has_access("field.equipment"): abort(403)
    return render_template("patrimonio.html")


@app.get("/api/v8/operacao")
@dashboard_required
def v8_operation_api():
    local_now = datetime.now(ZoneInfo("America/Sao_Paulo"))
    today_start_local = local_now.replace(hour=0, minute=0, second=0, microsecond=0)
    today_start_utc = today_start_local.astimezone(ZoneInfo("UTC")).replace(tzinfo=None)

    today_inventory = Inventory.query.filter(Inventory.created_at >= today_start_utc).all()
    visited_location_ids = {x.location_id for x in today_inventory}
    divergences_today = sum(1 for x in today_inventory if (x.divergence or "").strip())
    inoperative_today = sum(
        1 for x in today_inventory
        if normalize(x.operational_status) not in ("", "OPERACIONAL", "OK")
    )

    scheduled = _schedule_today_db(local_now.date())
    now_utc = datetime.utcnow()
    with_signal = 0
    stale = 0
    for member in scheduled:
        pos = _team_latest_position(member.get("user_id"))
        if not pos:
            continue
        minutes = max(0, int((now_utc - pos.captured_at).total_seconds() // 60))
        if minutes <= 15:
            with_signal += 1
        else:
            stale += 1

    return jsonify({
        "ok": True,
        "date": local_now.date().isoformat(),
        "time": local_now.strftime("%H:%M"),
        "today": {
            "inventory": len(today_inventory),
            "locations": len(visited_location_ids),
            "divergences": divergences_today,
            "inoperative": inoperative_today,
        },
        "teams": {
            "scheduled": len(scheduled),
            "recent_signal": with_signal,
            "stale_signal": stale,
            "no_signal": max(0, len(scheduled)-with_signal-stale),
        },
        "offline": {
            "mode": "PWA",
            "idempotency": True,
            "server_queue": 0
        }
    })


@app.get("/api/equipes/rail-network")
@teams_view_required
def teams_rail_network_api():
    if not _has_access("teams.map"): abort(403)
    """V39.7.3 — payload leve e dedicado ao mapa de Equipes.

    Evita usar /api/locations, que também calcula inventário, divergências e referências
    observadas. Aqui enviamos apenas os campos necessários para trilhos/estações.
    """
    rows = (
        db.session.query(
            Location.id, Location.company, Location.line, Location.location,
            Location.reference_latitude, Location.reference_longitude
        )
        .filter(Location.reference_latitude.isnot(None), Location.reference_longitude.isnot(None))
        .order_by(Location.line, Location.location)
        .all()
    )
    return jsonify([{
        "id": r.id, "company": r.company or "", "line": r.line or "",
        "location": r.location or "",
        "reference_latitude": float(r.reference_latitude),
        "reference_longitude": float(r.reference_longitude),
    } for r in rows])


@app.get("/equipes")
@teams_view_required
def teams_page():
    if not _has_access("teams"): abort(403)
    return render_template("teams.html")


@app.get("/central-operacional")
@teams_view_required
def central_operacional_page():
    # V60: painel executivo operacional. Usa a mesma fonte de verdade de Equipes,
    # evitando uma segunda regra de escala/GPS e números divergentes.
    if not _has_access("teams.today"): abort(403)
    return render_template("central_operacional.html", app_release=APP_RELEASE)



@app.get("/implantacao-hardware")
@hardware_implantation_required
def hardware_implantation_page():
    if not (_has_access("implantation.visits") or _has_access("implantation.reports")): abort(403)
    return render_template("hardware_implantation.html", app_release=APP_RELEASE)

@app.get("/implantacao-hardware/visita-campo")
@hardware_implantation_required
def hardware_field_visit_page():
    if not _has_access("implantation.visits"): abort(403)
    return render_template("hardware_field_visit.html", app_release=APP_RELEASE)

@app.get("/dashboard/implantacao")
@login_required
def hardware_implantation_dashboard_canonical():
    if session.get("role") not in ("manager","manager_field","technician_implantation"):
        abort(403)
    # V55.2: gestores visualizam a dashboard dentro do shell gerencial; não como atividade.
    if session.get("role") in ("manager","manager_field"):
        return redirect("/gerencial?view=implantation-dashboard")
    return render_template("hardware_implantation_dashboard.html", app_release=APP_RELEASE)

@app.get("/dashboard/implantacao/embed")
@login_required
def hardware_implantation_dashboard_embed():
    if session.get("role") not in ("manager","manager_field"):
        abort(403)
    return render_template("hardware_implantation_dashboard.html", app_release=APP_RELEASE, embedded=True)

@app.get("/implantacao-hardware/dashboard")
@login_required
def hardware_implantation_dashboard_page():
    # V55.1: rota legada redireciona uma única vez para a rota canônica, evitando conflito com a atividade.
    return redirect(url_for("hardware_implantation_dashboard_canonical"), code=302)

@app.get("/api/implantacao-hardware/visitas")
@hardware_implantation_required
def hardware_field_visits_api():
    q=HardwareFieldVisit.query
    for field,col in (("client",HardwareFieldVisit.client),("project",HardwareFieldVisit.project),("status",HardwareFieldVisit.status),("conclusion_status",HardwareFieldVisit.conclusion_status)):
        value=(request.args.get(field) or "").strip()
        if value: q=q.filter(col==value)
    location=(request.args.get("location") or "").strip()
    if location: q=q.filter(HardwareFieldVisit.location_name.ilike(f"%{location}%"))
    date_from=(request.args.get("date_from") or "").strip(); date_to=(request.args.get("date_to") or "").strip(); technician=(request.args.get("technician") or "").strip()
    try:
        if date_from: q=q.filter(HardwareFieldVisit.visit_date >= datetime.strptime(date_from,"%Y-%m-%d").date())
        if date_to: q=q.filter(HardwareFieldVisit.visit_date <= datetime.strptime(date_to,"%Y-%m-%d").date())
    except ValueError: return jsonify({"ok":False,"error":"Período inválido."}),400
    if technician:
        try: q=q.filter(HardwareFieldVisit.technician_id == int(technician))
        except ValueError: return jsonify({"ok":False,"error":"Técnico inválido."}),400
    rows=q.order_by(HardwareFieldVisit.created_at.desc()).limit(500).all()
    # V56-A.1: evita duplicidade visual de visitas idênticas sem apagar histórico/auditoria.
    deduped=[]; seen=set()
    for x in rows:
        sig=((x.client or '').strip().upper(),(x.project or '').strip().upper(),(x.location_name or '').strip().upper(),x.visit_date,x.technician_id,(x.status or '').strip().upper())
        if sig in seen: continue
        seen.add(sig); deduped.append(x)
    rows=deduped
    ids={x.technician_id for x in rows if x.technician_id}; tech_map={u.id:u.name for u in User.query.filter(User.id.in_(ids)).all()} if ids else {}
    visit_ids=[x.id for x in rows]; photo_counts={}
    if visit_ids:
        for vid,count in db.session.query(HardwareFieldVisitPhoto.visit_id,func.count(HardwareFieldVisitPhoto.id)).filter(HardwareFieldVisitPhoto.visit_id.in_(visit_ids)).group_by(HardwareFieldVisitPhoto.visit_id).all(): photo_counts[vid]=count
    return jsonify({"ok":True,"visits":[{"id":x.id,"report_code":f"RV-{x.id:06d}","client":x.client,"project":x.project,"location_name":x.location_name,"city":x.city,"visit_date":x.visit_date.isoformat(),"reason":x.reason,"status":x.status,"conclusion_status":x.conclusion_status,"client_accepted":x.client_accepted,"has_signature":bool(x.signature_file),"photo_count":photo_counts.get(x.id,0),"topdesk_ticket":x.topdesk_ticket,"report_group":x.report_group or "AUTOPASS","technician_id":x.technician_id,"technician_name":tech_map.get(x.technician_id,"—")} for x in rows]})

@app.get("/api/implantacao-hardware/tecnicos")
@hardware_implantation_required
def hardware_implantation_technicians_api():
    ids=[r[0] for r in db.session.query(HardwareFieldVisit.technician_id).distinct().all() if r[0]]
    rows=User.query.filter(User.id.in_(ids)).order_by(User.name).all() if ids else []
    return jsonify({"ok":True,"technicians":[{"id":u.id,"name":u.name} for u in rows]})

@app.get("/api/implantacao-hardware/resumo")
@hardware_implantation_required
def hardware_implantation_summary_api():
    rows=HardwareFieldVisit.query.all(); total=len(rows)
    by_status={}
    by_client={}
    for x in rows:
        k=x.conclusion_status or x.status or "PENDENTE"; by_status[k]=by_status.get(k,0)+1
        by_client[x.client]=by_client.get(x.client,0)+1
    return jsonify({"ok":True,"total":total,"finalized":sum(1 for x in rows if x.status=="FINALIZADO"),"in_progress":sum(1 for x in rows if x.status in ("RASCUNHO","PAUSADO")),"accepted":sum(1 for x in rows if x.client_accepted),"with_pending":sum(1 for x in rows if x.conclusion_status=="CONCLUÍDA COM PENDÊNCIAS"),"by_status":by_status,"by_client":by_client})

@app.post("/api/implantacao-hardware/visitas")
@hardware_implantation_required
def hardware_field_visit_save_api():
    if _activity_request_too_large(): return jsonify({"ok":False,"error":f"Envio excede {_ACTIVITY_REQUEST_MAX_MB} MB. Envie menos fotos por vez."}),413
    f=request.form
    required=("client","project","location_name","visit_date","reason")
    missing=[k for k in required if not (f.get(k) or "").strip()]
    if missing: return jsonify({"ok":False,"error":"Preencha os campos obrigatórios: "+", ".join(missing)}),400
    try: visit_date=datetime.strptime(f.get("visit_date"),"%Y-%m-%d").date()
    except Exception: return jsonify({"ok":False,"error":"Data inválida."}),400
    visit=HardwareFieldVisit(
      report_code="VIS-"+datetime.utcnow().strftime("%Y%m%d-%H%M%S")+"-"+secrets.token_hex(2).upper(),
      client=f.get("client").strip(), project=f.get("project").strip(), report_group=(f.get("report_group") or "AUTOPASS").strip().upper(),
      requester=f.get("requester"),
      has_topdesk=f.get("has_topdesk")=="1", topdesk_ticket=f.get("topdesk_ticket"), location_type=f.get("location_type"),
      location_name=f.get("location_name").strip(), city=f.get("city"), state=f.get("state"), address=f.get("address"),
      latitude=float(f.get("latitude")) if f.get("latitude") else None, longitude=float(f.get("longitude")) if f.get("longitude") else None,
      gps_accuracy=float(f.get("gps_accuracy")) if f.get("gps_accuracy") else None, visit_date=visit_date,start_time=f.get("start_time"),end_time=f.get("end_time"),
      reason=f.get("reason"), activities=f.get("activities"), activity_notes=f.get("activity_notes"), technical_details=f.get("technical_details"),
      conclusion_status=f.get("conclusion_status") or "EM ANDAMENTO", conclusion=f.get("conclusion"), pending_items=f.get("pending_items"),
      client_contact=f.get("client_contact"),client_company=f.get("client_company"),client_role=f.get("client_role"),client_email=f.get("client_email"),client_phone=f.get("client_phone"),
      contacts_json=(f.get("contacts_json") or "[]"),
      client_observations=f.get("client_observations"),client_accepted=f.get("client_accepted")=="1",technician_id=session["user_id"],status=f.get("save_mode") or "RASCUNHO")
    db.session.add(visit); db.session.flush()
    sig=f.get("signature_data") or ""
    if sig.startswith("data:image/") and "," in sig:
        import base64
        try:
            raw=base64.b64decode(sig.split(",",1)[1]); name=f"visit_sig_{visit.id}_{uuid.uuid4().hex[:8]}.png"
            if _r2_available():
                key=f"hardware-visits/{datetime.utcnow().strftime('%Y/%m')}/{name}";_r2_put_bytes(key,raw,"image/png");name="r2__"+key
            else: (UPLOAD_DIR/name).write_bytes(raw)
            visit.signature_file=name; visit.signed_at=datetime.utcnow()
        except Exception:
            app.logger.exception("Falha ao persistir assinatura da visita %s",visit.id)
    if visit.status=="FINALIZADO" and not (visit.client_accepted and visit.signature_file):
        db.session.rollback(); return jsonify({"ok":False,"error":"Para finalizar, registre o aceite e a assinatura do cliente."}),400
    for ph in request.files.getlist("photos"):
        if not ph or not ph.filename: continue
        ext=Path(secure_filename(ph.filename)).suffix.lower() or ".jpg"; name=f"visit_{visit.id}_{uuid.uuid4().hex}{ext}"
        name=_store_uploaded_file(ph,"hardware-visits",name,ph.mimetype or mimetypes.guess_type(name)[0] or "application/octet-stream")
        db.session.add(HardwareFieldVisitPhoto(visit_id=visit.id,stored_name=name,original_name=ph.filename,category="EVIDÊNCIA"))
    db.session.add(AuditEvent(user_id=session.get("user_id"),event_type="HARDWARE_FIELD_VISIT_SAVE",entity_type="hardware_field_visit",entity_id=str(visit.id),detail=f"{visit.report_code} · {visit.status}"))
    db.session.commit()
    return jsonify({"ok":True,"id":visit.id,"report_code":f"RV-{visit.id:06d}","status":visit.status})

@app.get("/api/implantacao-hardware/visitas/<int:visit_id>")
@manager_required
def hardware_field_visit_detail_api(visit_id):
    v=db.session.get(HardwareFieldVisit,visit_id)
    if not v: return jsonify({"ok":False,"error":"Relatório não encontrado."}),404
    fields=("client","project","report_group","requester","has_topdesk","topdesk_ticket","location_type","location_name","city","state","address","start_time","end_time","reason","activities","activity_notes","technical_details","conclusion_status","conclusion","pending_items","client_contact","client_company","client_role","client_email","client_phone","client_observations","client_accepted","status")
    data={k:getattr(v,k,None) for k in fields}
    try: contacts=json.loads(getattr(v,"contacts_json",None) or "[]")
    except Exception: contacts=[]
    data.update({"id":v.id,"report_code":f"RV-{v.id:06d}","visit_date":v.visit_date.isoformat() if v.visit_date else "","has_signature":bool(v.signature_file),"contacts":contacts})
    return jsonify({"ok":True,"visit":data})

@app.post("/api/implantacao-hardware/visitas/<int:visit_id>/editar")
@manager_required
def hardware_field_visit_edit_api(visit_id):
    v=db.session.get(HardwareFieldVisit,visit_id)
    if not v: return jsonify({"ok":False,"error":"Relatório não encontrado."}),404
    f=request.form; before=f"{v.client} | {v.project} | {v.status} | {v.conclusion_status}"
    for key in ("client","project","requester","topdesk_ticket","location_name","city","state","address","reason","activities","activity_notes","technical_details","conclusion_status","conclusion","pending_items","client_contact","client_company","client_role","client_email","client_phone","client_observations"):
        if key in f: setattr(v,key,(f.get(key) or "").strip() or None)
    if "contacts_json" in f: v.contacts_json=f.get("contacts_json") or "[]"
    if f.get("visit_date"):
        try: v.visit_date=datetime.strptime(f.get("visit_date"),"%Y-%m-%d").date()
        except ValueError: return jsonify({"ok":False,"error":"Data inválida."}),400
    db.session.add(AuditEvent(user_id=session.get("user_id"),event_type="HARDWARE_FIELD_VISIT_EDIT",entity_type="hardware_field_visit",entity_id=str(v.id),detail=f"RV-{v.id:06d} | antes: {before} | revisão administrativa")); db.session.commit()
    return jsonify({"ok":True,"id":v.id,"report_code":f"RV-{v.id:06d}"})

@app.post("/api/implantacao-hardware/visitas/<int:visit_id>/excluir")
@manager_required
def hardware_field_visit_delete_api(visit_id):
    v=db.session.get(HardwareFieldVisit,visit_id)
    if not v: return jsonify({"ok":False,"error":"Relatório não encontrado."}),404
    reason=(request.form.get("reason") or "").strip()
    if len(reason)<3: return jsonify({"ok":False,"error":"Informe o motivo da exclusão."}),400
    report=f"RV-{v.id:06d}"; detail=f"{report} | cliente={v.client} | projeto={v.project} | motivo={reason}"
    try:
        for ph in HardwareFieldVisitPhoto.query.filter_by(visit_id=v.id).all(): db.session.delete(ph)
        db.session.flush()
        db.session.delete(v)
        db.session.flush()
        db.session.add(AuditEvent(user_id=session.get("user_id"),event_type="HARDWARE_FIELD_VISIT_DELETE",entity_type="hardware_field_visit",entity_id=str(visit_id),detail=detail))
        db.session.commit()
        return jsonify({"ok":True,"deleted":report})
    except Exception as exc:
        db.session.rollback()
        app.logger.exception("Falha ao excluir %s",report)
        return jsonify({"ok":False,"error":"Não foi possível excluir o relatório. Consulte Diagnóstico ADM.","detail":str(exc)[:300]}),500

def _visit_media_data_uri(stored_name):
    if not stored_name: return None
    try:
        if stored_name.startswith("r2__"):
            raw=_r2_get_bytes(stored_name[4:])
            ext=Path(stored_name[4:]).suffix.lower()
        else:
            p=UPLOAD_DIR/stored_name
            if not p.exists(): return None
            raw=p.read_bytes();ext=p.suffix.lower()
        mime={".jpg":"image/jpeg",".jpeg":"image/jpeg",".png":"image/png",".webp":"image/webp"}.get(ext,"image/jpeg")
        import base64
        return f"data:{mime};base64,"+base64.b64encode(raw).decode("ascii")
    except Exception:
        app.logger.exception("Falha ao carregar mídia do relatório: %s",stored_name)
        return None

@app.get("/implantacao-hardware/visitas/<int:visit_id>/relatorio")
@hardware_implantation_required
def hardware_field_visit_report(visit_id):
    v=db.session.get(HardwareFieldVisit,visit_id)
    if not v: return "Relatório não encontrado",404
    photos=HardwareFieldVisitPhoto.query.filter_by(visit_id=v.id).all()
    photo_media=[{"row":p,"data_uri":_visit_media_data_uri(p.stored_name)} for p in photos]
    signature_data_uri=_visit_media_data_uri(v.signature_file) if v.signature_file else None
    tech=db.session.get(User,v.technician_id)
    return render_template("hardware_field_visit_report.html",v=v,photos=photos,photo_media=photo_media,signature_data_uri=signature_data_uri,tech=tech)

def _bytes_human(value):
    try:
        n=float(value or 0)
    except Exception:
        n=0.0
    units=("B","KB","MB","GB","TB")
    i=0
    while n>=1024 and i<len(units)-1:
        n/=1024.0; i+=1
    return f"{n:.1f} {units[i]}" if i else f"{int(n)} B"


def _process_memory_snapshot():
    out={"rss_bytes":None,"peak_rss_bytes":None,"vm_size_bytes":None,"limit_bytes":None,"usage_pct":None,"workers":int(os.getenv("WEB_CONCURRENCY","1") or 1),"cpu_count":os.cpu_count() or 1,"load_pct":None,"pid":os.getpid(),"threads":None}
    try:
        status=Path('/proc/self/status').read_text(errors='ignore')
        for key,target in (("VmRSS","rss_bytes"),("VmHWM","peak_rss_bytes"),("VmSize","vm_size_bytes")):
            m=re.search(rf'^{key}:\s+(\d+)\s+kB',status,re.M)
            if m: out[target]=int(m.group(1))*1024
        m=re.search(r'^Threads:\s+(\d+)',status,re.M)
        if m: out["threads"]=int(m.group(1))
    except Exception: pass
    for fp in ('/sys/fs/cgroup/memory.max','/sys/fs/cgroup/memory/memory.limit_in_bytes'):
        try:
            raw=Path(fp).read_text().strip()
            if raw and raw!='max':
                val=int(raw)
                if 0 < val < 10**15:
                    out["limit_bytes"]=val; break
        except Exception: pass
    # REV1: quando disponível, o cgroup representa melhor a memória real do container.
    try:
        cur=Path('/sys/fs/cgroup/memory.current')
        if cur.exists(): out["container_current_bytes"]=int(cur.read_text().strip())
    except Exception: out["container_current_bytes"]=None
    basis=out.get("container_current_bytes") or out.get("rss_bytes")
    if basis is not None and out["limit_bytes"]:
        out["usage_pct"]=round(basis*100/out["limit_bytes"],1)
    try:
        out["load_pct"]=round(min(999.0,os.getloadavg()[0]*100/max(1,out["cpu_count"])),1)
    except Exception: pass
    return out


def _database_storage_snapshot():
    data={"engine":"postgresql" if database_url.startswith('postgresql') else "sqlite","total_bytes":None,"limit_bytes":None,"usage_pct":None,"tables":[],"connections":None,"max_connections":None,"note":None}
    try:
        if data["engine"]=='postgresql':
            data["total_bytes"]=int(db.session.execute(text("SELECT pg_database_size(current_database())")).scalar() or 0)
            try:
                data["connections"]=int(db.session.execute(text("SELECT count(*) FROM pg_stat_activity WHERE datname=current_database()")).scalar() or 0)
                data["max_connections"]=int(db.session.execute(text("SHOW max_connections")).scalar() or 0)
            except Exception: pass
            rows=db.session.execute(text("""
                SELECT io.relname AS relname,
                       pg_total_relation_size(io.relid) AS total_bytes,
                       pg_relation_size(io.relid) AS table_bytes,
                       pg_indexes_size(io.relid) AS index_bytes,
                       COALESCE(st.n_live_tup,0) AS estimated_rows
                FROM pg_catalog.pg_statio_user_tables AS io
                LEFT JOIN pg_catalog.pg_stat_user_tables AS st ON st.relid = io.relid
                ORDER BY pg_total_relation_size(io.relid) DESC
                LIMIT 15
            """)).mappings().all()
            data["tables"]=[{"name":r["relname"],"total_bytes":int(r["total_bytes"] or 0),"table_bytes":int(r["table_bytes"] or 0),"index_bytes":int(r["index_bytes"] or 0),"estimated_rows":int(r["estimated_rows"] or 0)} for r in rows]
        else:
            fp=BASE_DIR/'inventario_local.db'
            data["total_bytes"]=fp.stat().st_size if fp.exists() else 0
    except Exception as exc:
        db.session.rollback()
        data["note"]=f"Métrica parcial: {str(exc)[:140]}"
    # Limite só é exibido quando explicitamente configurado pelo ambiente/provedor.
    try:
        lim_mb=float(os.getenv('DATABASE_STORAGE_LIMIT_MB','0') or 0)
        lim_gb=float(os.getenv('DATABASE_STORAGE_LIMIT_GB','0') or 0)
        if lim_mb>0: data["limit_bytes"]=int(lim_mb*1024*1024)
        elif lim_gb>0: data["limit_bytes"]=int(lim_gb*1024*1024*1024)
    except Exception: pass
    if data.get("total_bytes") is not None and data.get("limit_bytes"):
        data["usage_pct"]=round(data["total_bytes"]*100/data["limit_bytes"],1)
    if not data.get("limit_bytes") and not data.get("note"):
        data["note"]="Limite do plano não informado pela aplicação. Configure DATABASE_STORAGE_LIMIT_GB para exibir percentual disponível."
    return data


def _r2_storage_snapshot(force=False):
    now=time.time()
    with _STORAGE_CACHE_LOCK:
        cached=_STORAGE_CACHE.get("data")
        if cached and not force and now-float(_STORAGE_CACHE.get("at") or 0)<_STORAGE_CACHE_TTL:
            return dict(cached)
    data={"enabled":bool(_r2_available()),"objects":0,"total_bytes":0,"largest_bytes":0,"scanned_complete":True,"note":None}
    if data["enabled"]:
        try:
            client=r2_client(); token=None; pages=0
            while True:
                kw={"Bucket":os.environ["R2_BUCKET_NAME"],"MaxKeys":1000}
                if token: kw["ContinuationToken"]=token
                resp=client.list_objects_v2(**kw); pages+=1
                for obj in resp.get("Contents",[]):
                    sz=int(obj.get("Size") or 0); data["objects"]+=1; data["total_bytes"]+=sz; data["largest_bytes"]=max(data["largest_bytes"],sz)
                if not resp.get("IsTruncated"): break
                token=resp.get("NextContinuationToken")
                if not token or pages>=100:
                    data["scanned_complete"]=False; data["note"]="Leitura limitada às primeiras 100.000 mídias para proteger a telemetria."; break
        except Exception as exc:
            data["note"]=f"R2 indisponível para medição: {str(exc)[:140]}"
    else:
        data["note"]="R2 não configurado nesta instância."
    with _STORAGE_CACHE_LOCK:
        _STORAGE_CACHE["at"]=now; _STORAGE_CACHE["data"]=dict(data)
    return data


def _local_storage_snapshot():
    files=0; total=0
    try:
        for x in UPLOAD_DIR.rglob('*'):
            if x.is_file(): files+=1; total+=x.stat().st_size
    except Exception: pass
    return {"files":files,"total_bytes":total}


@app.get("/inteligencia-operacional")
@login_required
def operational_intelligence_page():
    if session.get("role") not in ("manager","manager_field"):
        return redirect(url_for("dashboard_landing"))
    return render_template("operational_intelligence.html", app_release=APP_RELEASE)

@app.get("/api/inteligencia-operacional")
@login_required
def operational_intelligence_api():
    if session.get("role") not in ("manager","manager_field"):
        return jsonify({"ok":False,"error":"Sem permissão."}),403
    now=datetime.utcnow(); active_roles=("technician","technician_implantation","manager_field")
    techs=User.query.filter(User.active.is_(True),User.role.in_(active_roles)).all(); tech_ids=[u.id for u in techs]
    since=now-timedelta(minutes=15)
    recent=dict(db.session.query(TechnicianPosition.user_id,func.max(TechnicianPosition.captured_at)).filter(TechnicianPosition.user_id.in_(tech_ids),TechnicianPosition.captured_at>=since).group_by(TechnicianPosition.user_id).all()) if tech_ids else {}
    open_status=("CONCLUIDO","CONCLUÍDO","FECHADO","RESOLVIDO")
    open_tickets=TopDeskTicket.query.filter(~func.upper(TopDeskTicket.work_status).in_(open_status)).all()
    assigned=sum(1 for t in open_tickets if t.assigned_technician_id); unassigned=len(open_tickets)-assigned
    available=len(recent); load=round(len(open_tickets)/available,1) if available else None
    # Heurística V66: prioridade explicável, sem caixa-preta.
    ranked=[]
    weights={"CRITICA":100,"CRÍTICA":100,"ALTA":70,"NORMAL":40,"BAIXA":20}
    for t in open_tickets:
        age_h=max(0,(now-(t.created_at or t.imported_at or now)).total_seconds()/3600)
        score=weights.get(normalize(t.priority),40)+min(60,age_h/2)+(25 if not t.assigned_technician_id else 0)
        ranked.append({"ticket":t.ticket_number,"priority":t.priority,"station":t.station_code or "","work_status":t.work_status,"age_h":round(age_h,1),"score":round(score,1),"reason":"sem técnico" if not t.assigned_technician_id else "idade + prioridade"})
    ranked.sort(key=lambda x:x["score"],reverse=True)
    anomalies=[]
    for u in techs:
        last=recent.get(u.id)
        if not last: anomalies.append({"type":"GPS","user":u.name,"detail":"Sem posição válida nos últimos 15 min"})
    capacity_status="CRÍTICA" if available==0 and open_tickets else ("ALTA CARGA" if load is not None and load>8 else "NORMAL")
    return jsonify({"ok":True,"release":APP_RELEASE,"generated_at":now.isoformat()+"Z","capacity":{"field_users":len(techs),"online_gps_15m":available,"open_work":len(open_tickets),"assigned":assigned,"unassigned":unassigned,"work_per_available":load,"status":capacity_status},"sla":{"at_risk":sum(1 for x in ranked if x["age_h"]>=24),"critical":sum(1 for x in ranked if x["age_h"]>=48)},"anomalies":anomalies[:20],"priorities":ranked[:20]})


def _v70_migration_snapshot():
    try:
        rows=SchemaMigration.query.order_by(SchemaMigration.applied_at.desc()).limit(20).all()
        return {"ok":True,"count":SchemaMigration.query.count(),"latest":[{"version":x.version,"description":x.description or "","applied_at":x.applied_at.isoformat() if x.applied_at else None} for x in rows]}
    except Exception as exc:
        return {"ok":False,"count":0,"latest":[],"error":str(exc)}

def _v70_index_snapshot():
    targets={
      "performance_metrics":{"ix_perf_created_route","ix_perf_route_created_status"},
      "topdesk_tickets":{"ix_topdesk_status_created","ix_topdesk_line_station_status"},
      "material_requests":{"ix_material_req_user_status"},
      "collaborator_documents":{"ix_collab_doc_user_status"},
      "technician_positions":{"ix_techpos_user_captured"},
    }
    out=[]
    try:
        insp=db.inspect(db.engine)
        for table,expected in targets.items():
            if not insp.has_table(table):
                out.append({"table":table,"expected":len(expected),"present":0,"status":"SEM TABELA"}); continue
            names={x.get("name") for x in insp.get_indexes(table)}
            present=len(expected & names)
            out.append({"table":table,"expected":len(expected),"present":present,"missing":sorted(expected-names),"status":"OK" if present==len(expected) else "ATENÇÃO"})
    except Exception as exc:
        return {"ok":False,"rows":out,"error":str(exc)}
    return {"ok":True,"rows":out}

@app.get("/telemetria")
@login_required
def telemetry_page():
    # V70: rota legada preservada. A navegação oficial é Saúde da Plataforma.
    if session.get("role") != "manager" and not _has_access("management.health"): abort(403)
    return redirect(url_for("platform_health_page"))

@app.get("/saude-plataforma")
@login_required
def platform_health_page():
    if session.get("role") != "manager" and not _has_access("management.health"): abort(403)
    return render_template("telemetry.html", app_release=APP_RELEASE)

@app.get("/api/telemetria/resumo")
@login_required
def telemetry_summary_api():
    if session.get("role") != "manager" and not _has_access("management.health"): return jsonify({"ok":False,"error":"Sem permissão para Saúde da Plataforma."}),403
    try:
        minutes=max(5,min(int(request.args.get("minutes") or 60),1440)); since=datetime.utcnow()-timedelta(minutes=minutes)
        q=PerformanceMetric.query.filter(PerformanceMetric.created_at>=since)
        rows=q.order_by(PerformanceMetric.created_at.desc()).limit(10000).all()
        vals=[float(x.duration_ms or 0) for x in rows]
        def pct(arr,p):
            if not arr:return 0
            a=sorted(arr); return round(a[min(len(a)-1,max(0,int((len(a)-1)*p)))],1)
        by={}
        for x in rows:
            d=by.setdefault(x.route,{"route":x.route,"count":0,"sum":0.0,"max":0.0,"errors":0,"vals":[],"sql_sum":0.0,"queries":0})
            d["count"]+=1; d["sum"]+=float(x.duration_ms or 0); d["max"]=max(d["max"],float(x.duration_ms or 0)); d["errors"]+=1 if x.status_code>=500 else 0; d["vals"].append(float(x.duration_ms or 0)); d["sql_sum"]+=float(getattr(x,"sql_ms",0) or 0); d["queries"]+=int(getattr(x,"query_count",0) or 0)
        route_rows=[]
        for d in by.values():
            route_rows.append({"route":d["route"],"count":d["count"],"avg_ms":round(d["sum"]/d["count"],1),"p95_ms":pct(d["vals"],.95),"max_ms":round(d["max"],1),"avg_sql_ms":round(d["sql_sum"]/d["count"],1),"avg_queries":round(d["queries"]/d["count"],1),"errors":d["errors"]})
        route_rows.sort(key=lambda x:(x["p95_ms"],x["avg_ms"]),reverse=True)
        # série em blocos de 5 minutos
        buckets={}
        for x in rows:
            dt=x.created_at; key=dt.replace(minute=(dt.minute//5)*5,second=0,microsecond=0)
            b=buckets.setdefault(key,[]); b.append(float(x.duration_ms or 0))
        timeline=[{"time":k.strftime("%H:%M"),"avg_ms":round(sum(v)/len(v),1),"p95_ms":pct(v,.95),"requests":len(v)} for k,v in sorted(buckets.items())]
        active_since=datetime.utcnow()-timedelta(minutes=15)
        active_users=db.session.query(func.count(func.distinct(SessionEvent.user_id))).filter(SessionEvent.created_at>=active_since).scalar() or 0
        table_counts={
          "Transações ATM":FinancialATMTransaction.query.count(),"Coletas":FinancialCashCollection.query.count(),"Inventário":Inventory.query.count(),"Posições GPS":TechnicianPosition.query.count(),"Chamados":TopDeskTicket.query.count(),"Auditoria":AuditEvent.query.count()
        }
        errors=sum(1 for x in rows if x.status_code>=500); avg=round(sum(vals)/len(vals),1) if vals else 0; p95=pct(vals,.95)
        health="NORMAL" if p95<1500 and errors==0 else ("ATENÇÃO" if p95<3000 and errors<3 else "CRÍTICO")
        top5=[dict(x,app_ms=max(0,round(float(x.get("avg_ms") or 0)-float(x.get("avg_sql_ms") or 0),1))) for x in route_rows[:5]]
        with PANORAMA_EXPORT_LOCK:
            _jobs=list(PANORAMA_EXPORT_JOBS.values())
        storage={"database":_database_storage_snapshot(),"r2":_r2_storage_snapshot(),"local":_local_storage_snapshot(),"runtime":_process_memory_snapshot(),"jobs":{"active":sum(1 for j in _jobs if j.get("status") in ("FILA","PROCESSANDO")),"ready":sum(1 for j in _jobs if j.get("status")=="PRONTO"),"errors":sum(1 for j in _jobs if j.get("status")=="ERRO")}}
        return jsonify({"ok":True,"release":APP_RELEASE,"generated_at":datetime.now(ZoneInfo("America/Sao_Paulo")).strftime("%d/%m/%Y %H:%M:%S"),"window_minutes":minutes,"health":health,"avg_ms":avg,"p95_ms":p95,"max_ms":round(max(vals),1) if vals else 0,"requests":len(rows),"errors_5xx":errors,"active_users_15m":int(active_users),"routes":route_rows[:20],"top5":top5,"timeline":timeline[-24:],"table_counts":table_counts,"storage":storage,"migrations":_v70_migration_snapshot(),"indexes":_v70_index_snapshot()})
    except Exception as exc:
        return jsonify({"ok":False,"error":str(exc)}),500

@app.get("/api/telemetria/export.xlsx")
@login_required
def telemetry_export_xlsx():
    if session.get("role") != "manager":
        abort(403)
    # Reaproveita exatamente o mesmo snapshot exibido na tela, respeitando a janela selecionada.
    result = telemetry_summary_api()
    response = result[0] if isinstance(result, tuple) else result
    status = result[1] if isinstance(result, tuple) and len(result) > 1 else getattr(response, "status_code", 200)
    if status >= 400:
        return result
    data = response.get_json(silent=True) or {}
    if not data.get("ok"):
        return jsonify(data), 500

    wb = Workbook()
    ws = wb.active
    ws.title = "Resumo"
    title_fill = PatternFill("solid", fgColor="17345D")
    head_fill = PatternFill("solid", fgColor="DCEAF7")
    title_font = Font(color="FFFFFF", bold=True, size=14)
    head_font = Font(bold=True, color="17345D")

    def style_header(sheet, row=1):
        for cell in sheet[row]:
            cell.fill = head_fill
            cell.font = head_font
            cell.alignment = Alignment(vertical="center")

    ws.append(["TELEMETRIA DO SISTEMA", data.get("release", APP_RELEASE)])
    ws["A1"].fill = title_fill; ws["A1"].font = title_font
    ws["B1"].fill = title_fill; ws["B1"].font = title_font
    summary_rows = [
        ("Gerado em", data.get("generated_at")),
        ("Janela (min)", data.get("window_minutes")),
        ("Saúde", data.get("health")),
        ("Tempo médio (ms)", data.get("avg_ms")),
        ("P95 (ms)", data.get("p95_ms")),
        ("Pico (ms)", data.get("max_ms")),
        ("Requisições medidas", data.get("requests")),
        ("Erros 5xx", data.get("errors_5xx")),
        ("Usuários ativos / 15 min", data.get("active_users_15m")),
    ]
    for k,v in summary_rows: ws.append([k,v])
    for row in ws.iter_rows(min_row=2, max_row=ws.max_row): row[0].font = Font(bold=True)

    wst = wb.create_sheet("Evolução 5min")
    wst.append(["Horário","Média (ms)","P95 (ms)","Requisições"]); style_header(wst)
    for x in data.get("timeline",[]): wst.append([x.get("time"),x.get("avg_ms"),x.get("p95_ms"),x.get("requests")])

    wsb = wb.create_sheet("Volumes das bases")
    wsb.append(["Base","Registros"]); style_header(wsb)
    for k,v in (data.get("table_counts") or {}).items(): wsb.append([k,v])

    st=data.get("storage") or {}; dbs=st.get("database") or {}; r2=st.get("r2") or {}; rt=st.get("runtime") or {}; loc=st.get("local") or {}; jobs=st.get("jobs") or {}
    wss=wb.create_sheet("Capacidade")
    wss.append(["Grupo","Métrica","Valor"]); style_header(wss)
    capacity=[
        ("PostgreSQL","Total bytes",dbs.get("total_bytes")),("PostgreSQL","Limite bytes",dbs.get("limit_bytes")),("PostgreSQL","Uso %",dbs.get("usage_pct")),("PostgreSQL","Conexões",dbs.get("connections")),("PostgreSQL","Máx conexões",dbs.get("max_connections")),
        ("R2","Total bytes",r2.get("total_bytes")),("R2","Objetos",r2.get("objects")),("R2","Maior arquivo bytes",r2.get("largest_bytes")),
        ("Aplicação","RAM atual bytes",rt.get("container_current_bytes") or rt.get("rss_bytes")),("Aplicação","Limite RAM bytes",rt.get("limit_bytes")),("Aplicação","Uso RAM %",rt.get("usage_pct")),("Aplicação","Workers",rt.get("workers")),("Aplicação","CPUs",rt.get("cpu_count")),
        ("Arquivos locais","Total bytes",loc.get("total_bytes")),("Arquivos locais","Arquivos",loc.get("files")),
        ("Processamentos","Ativos",jobs.get("active")),("Processamentos","Prontos",jobs.get("ready")),("Processamentos","Erros",jobs.get("errors")),
    ]
    for row in capacity: wss.append(list(row))

    wsdb=wb.create_sheet("Maiores tabelas DB")
    wsdb.append(["Tabela","Total bytes","Dados bytes","Índices bytes","Linhas estimadas"]); style_header(wsdb)
    for t in dbs.get("tables") or []: wsdb.append([t.get("name"),t.get("total_bytes"),t.get("table_bytes"),t.get("index_bytes"),t.get("estimated_rows")])

    ws5=wb.create_sheet("Top 5 gargalos")
    ws5.append(["Rota","P95 ms","Média ms","SQL médio ms","App ms","Queries/req","Chamadas","Erros 5xx"]); style_header(ws5)
    for x in data.get("top5") or []: ws5.append([x.get("route"),x.get("p95_ms"),x.get("avg_ms"),x.get("avg_sql_ms"),x.get("app_ms"),x.get("avg_queries"),x.get("count"),x.get("errors")])

    wsr=wb.create_sheet("Raio-X rotas")
    wsr.append(["Rota","Média ms","P95 ms","Máx ms","SQL médio ms","Queries/req","Chamadas","Erros 5xx","Diagnóstico"]); style_header(wsr)
    for x in data.get("routes") or []:
        p95=float(x.get("p95_ms") or 0)
        diag="CRÍTICO" if p95>5000 else ("LENTO" if p95>3000 else ("ATENÇÃO" if p95>1500 else "OK"))
        wsr.append([x.get("route"),x.get("avg_ms"),x.get("p95_ms"),x.get("max_ms"),x.get("avg_sql_ms"),x.get("avg_queries"),x.get("count"),x.get("errors"),diag])

    # Ajuste de largura para facilitar leitura e upload/análise posterior.
    for sh in wb.worksheets:
        sh.freeze_panes = "A2" if sh.max_row > 1 else None
        for col in range(1, sh.max_column+1):
            letter=get_column_letter(col)
            maxlen=0
            for cell in sh[letter]:
                try: maxlen=max(maxlen,len(str(cell.value or "")))
                except Exception: pass
            sh.column_dimensions[letter].width=min(max(maxlen+2,12),55)

    out=io.BytesIO(); wb.save(out); out.seek(0)
    stamp=datetime.now(ZoneInfo("America/Sao_Paulo")).strftime("%Y%m%d_%H%M%S")
    return send_file(out, as_attachment=True, download_name=f"telemetria_{APP_RELEASE.replace(' ','_')}_{stamp}.xlsx", mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")


@app.get("/diagnostico")
@login_required
def diagnostics_page():
    if session.get("role") not in ("manager","manager_field"): return redirect(url_for("dashboard_landing"))
    return render_template("diagnostics.html", app_release=APP_RELEASE)

@app.get("/api/diagnostico/resumo")
@login_required
def diagnostics_api():
    if session.get("role") not in ("manager","manager_field"): return jsonify({"ok":False,"error":"Sem permissão."}),403
    events=AuditEvent.query.order_by(AuditEvent.created_at.desc()).limit(150).all()
    user_ids={e.user_id for e in events if e.user_id}
    users={u.id:u.name for u in User.query.filter(User.id.in_(user_ids)).all()} if user_ids else {}
    media={"local_files":0,"local_bytes":0,"chip_photos":ChipSwapPhoto.query.count(),"emv_photos":EmvChipSwapPhoto.query.count(),"panorama_photos":PanoramaPointPhoto.query.count() if 'PanoramaPointPhoto' in globals() else 0,"r2_enabled":bool(_r2_available())}
    try:
        local=[x for x in UPLOAD_DIR.rglob('*') if x.is_file()]; media["local_files"]=len(local); media["local_bytes"]=sum(x.stat().st_size for x in local)
    except Exception: pass
    return jsonify({"ok":True,"release":APP_RELEASE,
        "database":{"users":User.query.count(),"inventory":Inventory.query.count(),"field_visits":HardwareFieldVisit.query.count(),"audit_events":AuditEvent.query.count(),"topdesk_tickets":TopDeskTicket.query.count()},
        "media":media,
        "events":[{"id":e.id,"created_at":e.created_at.isoformat() if e.created_at else None,"user":users.get(e.user_id,"Sistema"),"event_type":e.event_type,"entity_type":e.entity_type,"entity_id":e.entity_id or "","detail":e.detail or ""} for e in events]})

@app.get("/sobre")
@login_required
def about_page():
    # V52.1: Técnico Implantação acessa Sobre; o template oculta apenas o histórico de versões.
    return render_template(
        "about.html",
        app_release=APP_RELEASE,
        dashboard_release=DASHBOARD_RELEASE,
        base_version=BASE_DATA_VERSION,
        manager_version=DASHBOARD_RELEASE,
        teams_version=TEAMS_RELEASE,
    )



@app.get("/sw.js")
def service_worker():
    response = send_from_directory(STATIC_DIR, "sw.js", mimetype="application/javascript")
    response.headers["Service-Worker-Allowed"] = "/"
    response.headers["Cache-Control"] = "no-store, no-cache, must-revalidate, max-age=0"
    response.headers["Pragma"] = "no-cache"
    response.headers["Expires"] = "0"
    return response


@app.get("/manifest.webmanifest")
def web_manifest():
    response = send_from_directory(STATIC_DIR, "manifest.webmanifest", mimetype="application/manifest+json")
    response.headers["Cache-Control"] = "no-store, no-cache, must-revalidate, max-age=0"
    response.headers["Pragma"] = "no-cache"
    response.headers["Expires"] = "0"
    return response


@app.get("/offline")
def offline_page():
    return render_template("offline.html")


@app.get("/health")
def health():
    try:
        db.session.execute(db.text("SELECT 1"))
        return jsonify({"ok": True, "database": "connected", "release": "v6.0-central-operacional"})
    except Exception as exc:
        return jsonify({"ok": False, "database": "error", "detail": str(exc)}), 500


def _canonical_equipment_type(value):
    value = normalize(value)
    aliases = {
        "ATM": "ATM", "VALIDADOR": "VALIDADOR", "VALIDADOR DE RECARGA": "VALIDADOR",
        "POS": "POS", "POS DE BILHETERIA": "POS", "TDI": "TDI", "BLOQUEIO": "BLOQUEIO"
    }
    return aliases.get(value, value or "OUTRO")


def _base_asset_matches_location(asset, loc):
    if "INATIVO" in normalize(asset.base_status):
        return False
    asset_line = re.sub(r"^L(?=\\d{2}\\s*-)", "", normalize(asset.line))
    loc_line = re.sub(r"^L(?=\\d{2}\\s*-)", "", normalize(loc.line))
    if asset_line != loc_line:
        return False
    ac, lc = normalize(asset.company), normalize(loc.company)
    # V39.7.4: a base histórica de Validadores pode manter o operador anterior
    # (ex.: CPTM) mesmo quando a localidade hoje pertence à Via Mobilidade.
    # Para VALIDADORES, linha + estação/sigla são a referência operacional.
    if _canonical_equipment_type(asset.equipment_type) != "VALIDADOR":
        if ac and lc and lc not in ac and ac not in lc:
            return False
    station_text = normalize(loc.location)
    station_name = normalize(asset.locality)
    code = normalize(asset.location_code or asset.station_code)
    return bool((station_name and (station_name in station_text or station_text.endswith(station_name))) or
                (code and station_text.startswith(code + " ")))


def _normalize_line_key(value):
    return re.sub(r"^L(?=\d{2}\s*-)", "", normalize(value))


def _station_match_key(value):
    """Normaliza nome de estação para casar nomes abreviados e comerciais.

    Exemplos reais da base: SÉ <-> PRAÇA DA SÉ, ITAQUERA <-> CORINTHIANS-ITAQUERA
    e BARRA FUNDA <-> PALMEIRAS-BARRA FUNDA. O casamento é por tokens inteiros
    para evitar falsos positivos como SÉ dentro de BRESSER.
    """
    text_value = re.sub(r"[^A-Z0-9]+", " ", normalize(value))
    stopwords = {"ESTACAO", "EST", "PRACA", "DA", "DE", "DO", "DAS", "DOS", "PALMEIRAS", "CORINTHIANS"}
    return tuple(token for token in text_value.split() if token and token not in stopwords)


def _station_matches(asset_locality, asset_code, location_name):
    station_text = normalize(location_name)
    station_name = normalize(asset_locality)
    code = normalize(asset_code)

    # Regras exatas/legadas primeiro. Código também pode ser o próprio nome curto.
    if station_name and (station_name == station_text or station_name in station_text or station_text.endswith(station_name)):
        return True
    if code and (station_text == code or station_text.startswith(code + " ") or (" " + code + " ") in (" " + station_text + " ")):
        return True

    asset_tokens = set(_station_match_key(asset_locality))
    loc_tokens = set(_station_match_key(location_name))
    if not asset_tokens or not loc_tokens:
        return False
    smaller, larger = (asset_tokens, loc_tokens) if len(asset_tokens) <= len(loc_tokens) else (loc_tokens, asset_tokens)
    return smaller.issubset(larger)


def _invalidate_expected_cache():
    _expected_cache["at"] = 0.0
    _expected_cache["data"] = None


def _expected_assets_by_location(force=False):
    """Associa a base detalhada às localidades com índice por linha e cache em memória.

    A versão anterior comparava cada ativo com todas as localidades em toda chamada.
    Esta versão limita candidatos pela linha e reutiliza o resultado por 10 minutos.
    """
    now = time.monotonic()
    if (
        not force
        and _expected_cache.get("data") is not None
        and now - float(_expected_cache.get("at") or 0) < EXPECTED_CACHE_TTL_SECONDS
    ):
        return _expected_cache["data"]

    locations = Location.query.all()
    result = {
        loc.id: {"ATM": 0, "VALIDADOR": 0, "POS": 0, "TDI": 0, "BLOQUEIO": 0}
        for loc in locations
    }

    by_line = {}
    for loc in locations:
        by_line.setdefault(_normalize_line_key(loc.line), []).append(loc)

    assets = BaseAsset.query.all()
    valid_types = {"ATM", "VALIDADOR", "POS", "TDI", "BLOQUEIO"}
    for asset in assets:
        if "INATIVO" in normalize(asset.base_status):
            continue
        typ = _canonical_equipment_type(asset.equipment_type)
        if typ not in valid_types:
            continue

        candidates = by_line.get(_normalize_line_key(asset.line), ())
        if not candidates:
            continue

        qty = max(1, int(asset.quantity or 1))
        ac = normalize(asset.company)
        station_name = normalize(asset.locality)
        code = normalize(asset.location_code or asset.station_code)

        for loc in candidates:
            lc = normalize(loc.company)
            # V39.7.4: para Validadores, a referência da base é o TERMINAL
            # associado à linha/estação; não excluir por divergência de operador histórico.
            matched = _station_matches(asset.locality, asset.location_code or asset.station_code, loc.location)
            if matched:
                result[loc.id][typ] += qty
                break

    _expected_cache["at"] = now
    _expected_cache["data"] = result
    return result


def _haversine_m(lat1, lon1, lat2, lon2):
    from math import radians, sin, cos, sqrt, atan2
    r = 6371000.0
    p1, p2 = radians(lat1), radians(lat2)
    dphi = radians(lat2-lat1)
    dlambda = radians(lon2-lon1)
    a = sin(dphi/2)**2 + cos(p1)*cos(p2)*sin(dlambda/2)**2
    return 2*r*atan2(sqrt(a), sqrt(max(0.0, 1-a)))


def _observed_reference_stats(location_id):
    """V9: referência observada a partir de coletas GPS confiáveis.
    Não substitui automaticamente a referência oficial.
    """
    rows = (
        db.session.query(Inventory.latitude, Inventory.longitude, Inventory.gps_accuracy)
        .filter(
            Inventory.location_id == location_id,
            Inventory.latitude.isnot(None),
            Inventory.longitude.isnot(None),
            db.or_(Inventory.gps_accuracy.is_(None), Inventory.gps_accuracy <= 80),
        )
        .order_by(Inventory.created_at.desc())
        .limit(200)
        .all()
    )
    pts=[(float(a),float(b)) for a,b,_ in rows if a is not None and b is not None]
    if not pts:
        return {"count":0,"latitude":None,"longitude":None,"spread_m":None}
    # Mediana robusta primeiro; remove outliers acima de 400 m do centro observado.
    import statistics
    med_lat=statistics.median([x[0] for x in pts]); med_lon=statistics.median([x[1] for x in pts])
    clean=[x for x in pts if _haversine_m(x[0],x[1],med_lat,med_lon) <= 400]
    if not clean: clean=pts
    lat=sum(x[0] for x in clean)/len(clean); lon=sum(x[1] for x in clean)/len(clean)
    spread=max((_haversine_m(x[0],x[1],lat,lon) for x in clean), default=0)
    return {"count":len(clean),"latitude":lat,"longitude":lon,"spread_m":round(spread,1)}


@app.get("/api/locations")
@login_required
def api_locations():
    # V58: resposta padrão é leve. Referência observada por inventários GPS só é
    # calculada quando explicitamente solicitada (?observed=1), usada pelo Técnico.
    include_observed=request.args.get("observed") in ("1","true","yes")
    cache_slot="observed" if include_observed else "light"
    now=time.time(); cached=(_LOCATIONS_API_CACHE.get(cache_slot) or {}).get("payload")
    cached_at=float((_LOCATIONS_API_CACHE.get(cache_slot) or {}).get("at") or 0)
    if cached is not None and now-cached_at < _LOCATIONS_API_CACHE_TTL:
        resp=jsonify(cached); resp.headers["X-Autopass-Cache"]="HIT"; return resp
    expected_map = _expected_assets_by_location()

    # Inventário agregado em consultas pequenas, sem outer join pesado em todas as colunas.
    inv_rows = (
        db.session.query(
            Inventory.location_id,
            Inventory.equipment_type,
            func.count(Inventory.id),
            func.coalesce(func.sum(case((Inventory.operational_status == "Inoperante", 1), else_=0)), 0),
            func.coalesce(func.sum(case((
                Inventory.divergence.isnot(None)
                & (~Inventory.divergence.in_(["", "Não", "Nao"]))
            , 1), else_=0)), 0),
        )
        .group_by(Inventory.location_id, Inventory.equipment_type)
        .all()
    )

    inv_by_loc = {}
    for location_id, equipment_type, count, inop, divergence_count in inv_rows:
        bucket = inv_by_loc.setdefault(location_id, {
            "total": 0,
            "inoperative": 0,
            "divergences": 0,
            "by_type": {"ATM":0,"VALIDADOR":0,"POS":0,"TDI":0,"BLOQUEIO":0,"OUTRO":0},
        })
        canonical = _canonical_equipment_type(equipment_type)
        if canonical not in bucket["by_type"]:
            canonical = "OUTRO"
        bucket["by_type"][canonical] += int(count or 0)
        bucket["total"] += int(count or 0)
        bucket["inoperative"] += int(inop or 0)
        bucket["divergences"] += int(divergence_count or 0)

    rows = Location.query.order_by(Location.company, Location.line, Location.location).all()

    # REV4: calcula referências observadas em lote. A versão anterior executava
    # uma consulta por estação, elevando /api/locations para ~210 queries/request.
    observed_by_loc = {}
    if include_observed:
        obs_rows = (
            db.session.query(Inventory.location_id, Inventory.latitude, Inventory.longitude, Inventory.gps_accuracy)
            .filter(Inventory.latitude.isnot(None), Inventory.longitude.isnot(None), db.or_(Inventory.gps_accuracy.is_(None), Inventory.gps_accuracy <= 80))
            .order_by(Inventory.location_id, Inventory.created_at.desc())
            .all()
        )
        obs_points = {}
        for lid, lat, lon, acc in obs_rows:
            bucket=obs_points.setdefault(lid,[])
            if len(bucket)<200: bucket.append((float(lat),float(lon)))
        import statistics
        for lid,pts in obs_points.items():
            if not pts: continue
            med_lat=statistics.median([x[0] for x in pts]); med_lon=statistics.median([x[1] for x in pts])
            clean=[x for x in pts if _haversine_m(x[0],x[1],med_lat,med_lon)<=400] or pts
            lat=sum(x[0] for x in clean)/len(clean); lon=sum(x[1] for x in clean)/len(clean)
            spread=max((_haversine_m(x[0],x[1],lat,lon) for x in clean),default=0)
            observed_by_loc[lid]={"count":len(clean),"latitude":lat,"longitude":lon,"spread_m":round(spread,1)}

    out = []
    for loc in rows:
        inv = inv_by_loc.get(loc.id, {
            "total":0, "inoperative":0, "divergences":0,
            "by_type":{"ATM":0,"VALIDADOR":0,"POS":0,"TDI":0,"BLOQUEIO":0,"OUTRO":0}
        })
        exp_by_type = expected_map.get(loc.id, {})
        out.append({
            "id": loc.id,
            "company": loc.company,
            "line": loc.line,
            "location": loc.location,
            "base_status": loc.base_status,
            "expected_atm": loc.expected_atm,
            "expected_validator": loc.expected_validator,
            "expected_pos": loc.expected_pos,
            "expected_by_type": exp_by_type,
            "expected_total": sum(exp_by_type.values()),
            "inventoried_by_type": inv["by_type"],
            "survey_status": loc.survey_status,
            "started_at": loc.started_at.isoformat(timespec="seconds") if loc.started_at else None,
            "completed_at": loc.completed_at.isoformat(timespec="seconds") if loc.completed_at else None,
            "completed_by": loc.completed_by,
            "reference_latitude": loc.reference_latitude,
            "reference_longitude": loc.reference_longitude,
            "reference_source": loc.reference_source,
            "reference_updated_at": loc.reference_updated_at.isoformat(timespec="seconds") if loc.reference_updated_at else None,
            "observed_reference": observed_by_loc.get(loc.id,{"count":0,"latitude":None,"longitude":None,"spread_m":None}),
            "inventoried": int(inv["total"]),
            "inoperative": int(inv["inoperative"]),
            "divergences": int(inv["divergences"]),
        })
    _LOCATIONS_API_CACHE[cache_slot]={"payload":out,"at":time.time()}
    resp=jsonify(out); resp.headers["X-Autopass-Cache"]="MISS"; resp.headers["X-Autopass-Mode"]=("observed" if include_observed else "light"); return resp


@app.get("/api/location/<int:location_id>/inventory")
@login_required
def api_location_inventory(location_id):
    rows = (
        db.session.query(Inventory, User.name.label("technician"))
        .join(User, User.id == Inventory.technician_id)
        .filter(Inventory.location_id == location_id)
        .order_by(Inventory.created_at.desc())
        .all()
    )
    inv_ids=[inv.id for inv,_ in rows]
    attachment_counts={}
    if inv_ids:
        for iid,cnt in db.session.query(Attachment.inventory_id,func.count(Attachment.id)).filter(Attachment.inventory_id.in_(inv_ids)).group_by(Attachment.inventory_id).all(): attachment_counts[iid]=int(cnt or 0)
    loc=db.session.get(Location,location_id)
    survey_status=loc.survey_status if loc else ""
    out = []
    for inv, technician_name in rows:
        attachment_count = attachment_counts.get(inv.id,0)
        out.append({
            "id": inv.id,
            "location_id": inv.location_id,
            "equipment_type": inv.equipment_type,
            "base_asset_id": inv.base_asset_id,
            "asset_identifier": inv.asset_identifier,
            "serial": inv.serial,
            "supplier": inv.supplier,
            "model": inv.model,
            "application": inv.application,
            "bom_id": inv.bom_id,
            "bu_id": inv.bu_id,
            "validator_top_id": inv.validator_top_id,
            "software_version": inv.software_version,
            "exact_position": inv.exact_position,
            "mount": inv.mount,
            "operational_status": inv.operational_status,
            "connectivity": inv.connectivity,
            "network_id": inv.network_id,
            "label_status": inv.label_status,
            "in_base": inv.in_base,
            "divergence": inv.divergence,
            "notes": inv.notes,
            "created_at": inv.created_at.isoformat(timespec="seconds"),
            "updated_at": inv.updated_at.isoformat(timespec="seconds") if inv.updated_at else None,
            "latitude": inv.latitude,
            "longitude": inv.longitude,
            "gps_accuracy": inv.gps_accuracy,
            "gps_captured_at": inv.gps_captured_at.isoformat(timespec="seconds") if inv.gps_captured_at else None,
            "technician": technician_name,
            "attachments_count": attachment_count,
            "location_survey_status": survey_status,
            "can_manage": session.get("role") == "manager",
        })
    return jsonify(out)


_station_network_cache = None

def _load_station_network_rows():
    """V41.1: carrega uma vez a Tabela Estações usada nos dados técnicos de rede."""
    global _station_network_cache
    if _station_network_cache is not None:
        return _station_network_cache
    rows=[]
    path = BASE_DIR / "stations_network.xlsx"
    if not path.exists():
        _station_network_cache=[]
        return []
    try:
        wb=load_workbook(path, read_only=True, data_only=True); ws=wb.active
        it=ws.iter_rows(values_only=True); headers=[normalize(str(x or "")) for x in next(it)]
        for vals in it:
            d=dict(zip(headers, vals)); station=str(d.get("ESTACAO") or "").strip()
            if not station or not d.get("IP"):
                continue
            prefix=re.sub(r"\D", "", str(d.get("PREFIXO") or ""))
            line_logic=re.sub(r"\D", "", str(d.get("LINHA") or ""))
            rows.append({
                "station":station,"station_norm":normalize(station),"ip":str(d.get("IP") or ""),
                "mask":str(d.get("MASK") or ""),"gateway":str(d.get("GETAWAY") or d.get("GATEWAY") or ""),
                "dns1":str(d.get("DNS 1") or ""),"dns2":str(d.get("DNS 2") or ""),
                "group":str(d.get("GRUPO") or ""),"line_logic":line_logic,"prefix":prefix,
                "blocking_number":prefix[-2:] if len(prefix)>=2 else ""
            })
    except Exception:
        app.logger.exception("Falha ao ler stations_network.xlsx")
        rows=[]
    _station_network_cache=rows
    return rows

def _station_technical_config(loc):
    if not loc: return {}
    txt=normalize(loc.location); code=normalize(getattr(loc,'location_code',None) or getattr(loc,'station_code',None) or '')
    for r in _load_station_network_rows():
        st=r['station_norm']; st_code=st.split(' - ',1)[0].strip() if ' - ' in st else ''
        st_name=st.split(' - ',1)[1].strip() if ' - ' in st else st
        if (code and st_code and code==st_code) or (st_name and (st_name in txt or txt.endswith(st_name))) or (st_code and (txt.startswith(st_code+' ') or txt.startswith(st_code+' -'))):
            return dict(r)
    return {}

_block_config_cache = None
def _block_technical_config(prefix):
    """Retorna configuração pelo prefixo; V56-B REV2 mantém JSON em memória."""
    global _block_config_cache
    key = re.sub(r"\D", "", str(prefix or ""))
    if not key:
        return {}
    if _block_config_cache is None:
        try:
            source = DATA_DIR / "block_config_v18.json"
            payload = json.loads(source.read_text(encoding="utf-8")) if source.exists() else {}
            _block_config_cache = payload.get("by_prefix") or {}
        except Exception:
            _block_config_cache = {}
    hit=_block_config_cache.get(key, {}) or {}
    if hit: return hit
    for r in _load_station_network_rows():
        if r.get("prefix")==key:
            return dict(r)
    return {}


def _v20_block_asset_key(asset):
    """Corrige a identificação visual: 500502 => BLOQ02, preservando linha/localidade/prefixo."""
    prefix = re.sub(r"\D", "", str(asset.terminal_number or asset.top_id or asset.qrcode_id or ""))
    block_no = prefix[-2:] if len(prefix) >= 2 else ""
    original = str(asset.asset_key or "")
    if block_no and normalize(asset.equipment_type) == "BLOQUEIO":
        if re.match(r"(?i)^BLOQ\d+", original):
            return re.sub(r"(?i)^BLOQ\d+", f"BLOQ{block_no}", original)
        return f"BLOQ{block_no}|{asset.line or ''}|{asset.locality or ''}|{prefix}"
    return original


@app.get("/api/location/<int:location_id>/assets")
@login_required
def api_assets(location_id):
    loc = db.session.get(Location, location_id)
    if not loc:
        return jsonify([])

    requested_type = normalize(request.args.get("equipment_type", ""))
    type_aliases = {
        "ATM": "ATM",
        "VALIDADOR": "VALIDADOR",
        "VALIDADOR DE RECARGA": "VALIDADOR",
        "POS": "POS",
        "POS DE BILHETERIA": "POS",
        "TDI": "TDI",
        "BLOQUEIO": "BLOQUEIO",
    }
    requested_type = type_aliases.get(requested_type, requested_type)

    line = normalize(loc.line)
    company = normalize(loc.company)
    station_text = normalize(loc.location)

    already = {
        x[0]
        for x in db.session.query(Inventory.base_asset_id)
        .filter(
            Inventory.location_id == location_id,
            Inventory.base_asset_id.isnot(None)
        )
        .all()
    }

    out = []
    for a in BaseAsset.query.all():
        asset_type = normalize(a.equipment_type or "ATM")
        if requested_type and asset_type != requested_type:
            continue

        # Inventário de campo: não oferece estoque/inativos como ativo previsto.
        if "INATIVO" in normalize(a.base_status):
            continue

        asset_company = normalize(a.company)
        asset_line = normalize(a.line)
        # Algumas planilhas antigas usam L01 - AZUL; a localidade usa 01 - AZUL.
        asset_line_cmp = re.sub(r"^L(?=\d{2}\s*-)", "", asset_line)
        line_cmp = re.sub(r"^L(?=\d{2}\s*-)", "", line)
        if asset_line_cmp != line_cmp:
            continue
        # V39.7.4: Validadores são vinculados pela linha + estação/sigla.
        # Não bloquear por empresa porque a planilha de terminais pode conservar
        # o operador histórico enquanto a estação já está sob nova concessão.
        if asset_type != "VALIDADOR":
            if company not in asset_company and asset_company not in company:
                continue

        station_name = normalize(a.locality)
        code = normalize(a.location_code or a.station_code)
        station_match = (
            (station_name and (station_name in station_text or station_text.endswith(station_name)))
            or (code and station_text.startswith(code + " "))
        )
        if not station_match:
            continue

        out.append({
            "id": a.id,
            "equipment_type": a.equipment_type or "ATM",
            "asset_key": _v20_block_asset_key(a),
            "description": a.description,
            "company": a.company,
            "station_code": a.station_code,
            "location_code": a.location_code,
            "line": a.line,
            "locality": a.locality,
            "terminal_number": a.terminal_number,
            "serial": a.serial,
            "qrcode_id": a.qrcode_id,
            "top_id": a.top_id,
            "products": a.products,
            "model": a.model,
            "supplier": a.supplier,
            "transactions": a.transactions,
            "pix": a.pix,
            "mount": a.mount,
            "base_status": a.base_status,
            "application": a.application,
            "bom_id": a.bom_id,
            "bu_id": a.bu_id,
            "software_version": a.software_version,
            "quantity": a.quantity,
            "base_notes": a.base_notes,
            "leasing_status": a.leasing_status,
            "contract_end": a.contract_end,
            "installation_type": a.installation_type,
            "installation_date": a.installation_date,
            "technical_config": ({**_station_technical_config(loc), **_block_technical_config(a.terminal_number or a.top_id or a.qrcode_id or a.asset_key)} if asset_type in ("BLOQUEIO","ATM") else {}),
            "already_inventoried": a.id in already,
        })

    return jsonify(out)

def _optional_float(value):
    value = (value or "").strip()
    if not value:
        return None
    try:
        return float(value.replace(",", "."))
    except (TypeError, ValueError):
        return None


def _optional_iso_datetime(value):
    value = (value or "").strip()
    if not value:
        return None
    try:
        # O navegador envia ISO 8601. Normaliza "Z" para compatibilidade.
        return datetime.fromisoformat(value.replace("Z", "+00:00")).replace(tzinfo=None)
    except (TypeError, ValueError):
        return None


@app.post("/api/inventory")
@field_required
def create_inventory():
    if _activity_request_too_large(): return jsonify({"ok":False,"error":f"Envio excede {_ACTIVITY_REQUEST_MAX_MB} MB. Envie menos evidências por vez."}),413
    location_id = request.form.get("location_id", type=int)
    equipment_type = request.form.get("equipment_type", "").strip()
    base_asset_id = request.form.get("base_asset_id", type=int)
    serial = request.form.get("serial", "").strip()
    asset_identifier = request.form.get("asset_identifier", "").strip() or serial
    latitude = _optional_float(request.form.get("latitude"))
    longitude = _optional_float(request.form.get("longitude"))
    gps_accuracy = _optional_float(request.form.get("gps_accuracy"))
    gps_captured_at = _optional_iso_datetime(request.form.get("gps_captured_at"))
    client_uuid = request.form.get("client_uuid", "").strip() or None

    if client_uuid:
        existing_sync = Inventory.query.filter_by(sync_uuid=client_uuid).first()
        if existing_sync:
            return jsonify({"ok": True, "id": existing_sync.id, "idempotent": True})

    if not location_id or not equipment_type or not asset_identifier:
        return jsonify({"ok": False, "error": "Local, tipo e identificação/série são obrigatórios."}), 400

    loc = db.session.get(Location, location_id)
    if not loc:
        return jsonify({"ok": False, "error": "Local inválido."}), 400

    # V10: qualidade do GPS antes do confronto geográfico.
    gps_override_reason = (request.form.get("gps_override_reason") or "").strip()
    if gps_accuracy is not None and gps_accuracy > FIELD_GPS_MAX_ACCURACY_M and len(gps_override_reason) < 3:
        return jsonify({
            "ok": False,
            "error": f"Precisão GPS insuficiente ({round(gps_accuracy)} m). Atualize a localização ou informe uma justificativa para prosseguir.",
            "code": "GPS_LOW_ACCURACY",
            "accuracy_m": round(gps_accuracy),
            "max_accuracy_m": FIELD_GPS_MAX_ACCURACY_M
        }), 409

    # V10: confronto geográfico entre o local declarado e o GPS capturado.
    gps_distance_m = None
    gps_warn_m = float(os.getenv("FIELD_GPS_WARN_DISTANCE_M", "250"))
    gps_max_m = float(os.getenv("FIELD_GPS_MAX_DISTANCE_M", "600"))
    if latitude is not None and longitude is not None and loc.reference_latitude is not None and loc.reference_longitude is not None:
        gps_distance_m = _haversine_m(latitude, longitude, loc.reference_latitude, loc.reference_longitude)
        if gps_distance_m > gps_max_m and len(gps_override_reason) < 3:
            return jsonify({
                "ok": False,
                "error": f"GPS incompatível com {loc.location}: distância aproximada de {round(gps_distance_m)} m da referência. Confirme a localidade ou informe uma justificativa para exceção.",
                "code": "GPS_OUTSIDE_LOCATION",
                "distance_m": round(gps_distance_m),
                "warn_m": gps_warn_m,
                "max_m": gps_max_m
            }), 409

    duplicate = (
        db.session.query(Inventory, User.name.label("technician"))
        .join(User, User.id == Inventory.technician_id)
        .filter(
            Inventory.location_id == location_id,
            Inventory.equipment_type == equipment_type,
            func.upper(Inventory.asset_identifier) == asset_identifier.upper()
        )
        .first()
    )
    if duplicate:
        inv, technician_name = duplicate
        when = inv.created_at.strftime("%d/%m/%Y %H:%M")
        return jsonify({
            "ok": False,
            "duplicate": True,
            "error": f"Este equipamento já foi inventariado por {technician_name} em {when}."
        }), 409

    now = datetime.utcnow()
    inv = Inventory(
        location_id=location_id,
        equipment_type=equipment_type,
        base_asset_id=base_asset_id,
        asset_identifier=asset_identifier,
        serial=serial,
        supplier=request.form.get("supplier", ""),
        model=request.form.get("model", ""),
        application=request.form.get("application", ""),
        bom_id=request.form.get("bom_id", ""),
        bu_id=request.form.get("bu_id", ""),
        validator_top_id=request.form.get("validator_top_id", ""),
        software_version=request.form.get("software_version", ""),
        exact_position=request.form.get("exact_position", ""),
        mount=request.form.get("mount", ""),
        operational_status=request.form.get("operational_status", ""),
        connectivity=request.form.get("connectivity", ""),
        network_id=request.form.get("network_id", ""),
        teamviewer_id=request.form.get("teamviewer_id", ""),
        label_status=request.form.get("label_status", ""),
        in_base=request.form.get("in_base", ""),
        divergence=request.form.get("divergence", ""),
        notes=((request.form.get("notes", "") or "") + (("\n[Exceção GPS: " + gps_override_reason + "]") if gps_override_reason else "")).strip(),
        latitude=latitude,
        longitude=longitude,
        gps_accuracy=gps_accuracy,
        gps_captured_at=gps_captured_at,
        sync_uuid=client_uuid,
        technician_id=session["user_id"],
        created_at=now
    )

    try:
        db.session.add(inv)
        db.session.flush()

        # Armazenamento local temporário.
        # Em produção, vamos trocar esta parte por storage central (R2/S3).
        for f in request.files.getlist("attachments"):
            if not f or not f.filename:
                continue
            safe = secure_filename(f.filename)
            stored = f"{inv.id}_{secrets.token_hex(6)}_{safe}"
            stored = _store_uploaded_file(f, "inventory", stored, f.mimetype or "application/octet-stream")
            db.session.add(Attachment(
                inventory_id=inv.id,
                original_name=f.filename,
                stored_name=stored,
                mime_type=f.mimetype
            ))

        if loc.survey_status == "PENDENTE":
            loc.survey_status = "EM ANDAMENTO"
            loc.started_at = now

        db.session.commit()
    except IntegrityError:
        db.session.rollback()
        return jsonify({"ok": False, "duplicate": True, "error": "Registro duplicado para este local."}), 409

    return jsonify({"ok": True, "id": inv.id})


@app.get("/api/inventory/<int:inventory_id>")
@dashboard_required
def get_inventory_admin(inventory_id):
    if not _current_user_is_superadmin():
        return jsonify({"ok": False, "error": "Edição disponível somente para o administrador principal."}), 403
    inv = db.session.get(Inventory, inventory_id)
    if not inv:
        return jsonify({"ok": False, "error": "Registro não encontrado."}), 404
    loc = db.session.get(Location, inv.location_id)
    user = db.session.get(User, inv.technician_id)
    return jsonify({"ok": True, "inventory": {
        "id": inv.id, "location_id": inv.location_id, "location": loc.location if loc else "",
        "company": loc.company if loc else "", "line": loc.line if loc else "",
        "equipment_type": inv.equipment_type, "asset_identifier": inv.asset_identifier,
        "serial": inv.serial or "", "supplier": inv.supplier or "", "model": inv.model or "",
        "teamviewer_id": inv.teamviewer_id or "",
        "operational_status": inv.operational_status or "", "in_base": inv.in_base or "",
        "divergence": inv.divergence or "", "notes": inv.notes or "",
        "creator": user.name if user else "—", "creator_username": user.username if user else ""
    }})


@app.route("/api/inventory/<int:inventory_id>", methods=["PUT", "PATCH"])
@field_required
def update_inventory(inventory_id):
    if _activity_request_too_large(): return jsonify({"ok":False,"error":f"Envio excede {_ACTIVITY_REQUEST_MAX_MB} MB. Envie menos evidências por vez."}),413
    inv = db.session.get(Inventory, inventory_id)
    if not inv:
        return jsonify({"ok": False, "error": "Registro não encontrado."}), 404

    # V39.6: edição colaborativa. Técnico ou Gestor pode corrigir qualquer registro.
    # A autoria original permanece em technician_id e a alteração é auditada abaixo.

    location_id = request.form.get("location_id", type=int) or inv.location_id
    equipment_type = request.form.get("equipment_type", inv.equipment_type or "").strip()
    base_asset_id = request.form.get("base_asset_id", type=int)
    serial = request.form.get("serial", inv.serial or "").strip()
    asset_identifier = request.form.get("asset_identifier", inv.asset_identifier or "").strip() or serial

    if not location_id or not equipment_type or not asset_identifier:
        return jsonify({"ok": False, "error": "Local, tipo e identificação/série são obrigatórios."}), 400

    loc = db.session.get(Location, location_id)
    if not loc:
        return jsonify({"ok": False, "error": "Local inválido."}), 400

    duplicate = Inventory.query.filter(
        Inventory.id != inventory_id,
        Inventory.location_id == location_id,
        Inventory.equipment_type == equipment_type,
        func.upper(Inventory.asset_identifier) == asset_identifier.upper()
    ).first()
    if duplicate:
        return jsonify({
            "ok": False,
            "duplicate": True,
            "error": "Já existe outro equipamento com esta identificação neste local."
        }), 409

    inv.location_id = location_id
    inv.equipment_type = equipment_type
    inv.base_asset_id = base_asset_id
    inv.asset_identifier = asset_identifier
    inv.serial = serial
    inv.supplier = request.form.get("supplier", inv.supplier or "")
    inv.model = request.form.get("model", inv.model or "")
    inv.application = request.form.get("application", inv.application or "")
    inv.bom_id = request.form.get("bom_id", inv.bom_id or "")
    inv.bu_id = request.form.get("bu_id", inv.bu_id or "")
    inv.validator_top_id = request.form.get("validator_top_id", inv.validator_top_id or "")
    inv.software_version = request.form.get("software_version", inv.software_version or "")
    inv.exact_position = request.form.get("exact_position", inv.exact_position or "")
    inv.mount = request.form.get("mount", inv.mount or "")
    inv.operational_status = request.form.get("operational_status", inv.operational_status or "")

    # V42.4: ADM pode corrigir também a situação do levantamento da localidade
    # a partir da edição de qualquer registro. A alteração é preservada em auditoria.
    requested_survey_status = (request.form.get("location_survey_status") or "").strip().upper()
    survey_status_changed = False
    old_survey_status = loc.survey_status if loc else None
    if requested_survey_status and session.get("role") == "manager":
        allowed_survey_status = {"PENDENTE", "EM ANDAMENTO", "CONCLUIDA"}
        requested_survey_status = requested_survey_status.replace("CONCLUÍDA", "CONCLUIDA")
        if requested_survey_status not in allowed_survey_status:
            return jsonify({"ok": False, "error": "Status do levantamento inválido."}), 400
        if loc.survey_status != requested_survey_status:
            loc.survey_status = requested_survey_status
            survey_status_changed = True
            if requested_survey_status == "PENDENTE":
                loc.started_at = None
                loc.completed_at = None
                loc.completed_by = None
            elif requested_survey_status == "EM ANDAMENTO":
                loc.started_at = loc.started_at or datetime.utcnow()
                loc.completed_at = None
                loc.completed_by = None
            elif requested_survey_status == "CONCLUIDA":
                loc.completed_at = loc.completed_at or datetime.utcnow()
                loc.completed_by = session.get("user_id")

    inv.connectivity = request.form.get("connectivity", inv.connectivity or "")
    inv.network_id = request.form.get("network_id", inv.network_id or "")
    inv.teamviewer_id = request.form.get("teamviewer_id", inv.teamviewer_id or "")
    inv.label_status = request.form.get("label_status", inv.label_status or "")
    inv.in_base = request.form.get("in_base", inv.in_base or "")
    inv.divergence = request.form.get("divergence", inv.divergence or "")
    inv.notes = request.form.get("notes", inv.notes or "")

    if request.form.get("latitude") not in (None, ""):
        inv.latitude = _optional_float(request.form.get("latitude"))
    if request.form.get("longitude") not in (None, ""):
        inv.longitude = _optional_float(request.form.get("longitude"))
    if request.form.get("gps_accuracy") not in (None, ""):
        inv.gps_accuracy = _optional_float(request.form.get("gps_accuracy"))
    if request.form.get("gps_captured_at") not in (None, ""):
        inv.gps_captured_at = _optional_iso_datetime(request.form.get("gps_captured_at"))

    inv.updated_at = datetime.utcnow()
    added_photos = 0
    for f in request.files.getlist("attachments"):
        if not f or not f.filename:
            continue
        safe = secure_filename(f.filename) or f"evidencia_{secrets.token_hex(4)}.jpg"
        stored = f"{inv.id}_{secrets.token_hex(6)}_{safe}"
        stored = _store_uploaded_file(f, "inventory", stored, f.mimetype or "application/octet-stream")
        db.session.add(Attachment(inventory_id=inv.id, original_name=f.filename, stored_name=stored, mime_type=f.mimetype))
        added_photos += 1
    detail = f"Registro editado por {session.get('name','usuário')}; {added_photos} nova(s) evidência(s)"
    if survey_status_changed:
        detail += f"; status da localidade: {old_survey_status or '—'} → {loc.survey_status}"
    db.session.add(AuditEvent(user_id=session.get("user_id"), event_type="INVENTORY_EDIT", entity_type="inventory", entity_id=str(inv.id), detail=detail))
    if survey_status_changed:
        db.session.add(AuditEvent(user_id=session.get("user_id"), event_type="LOCATION_STATUS_OVERRIDE", entity_type="location", entity_id=str(loc.id), detail=f"ADM alterou status de {old_survey_status or '—'} para {loc.survey_status} a partir do inventário #{inv.id}."))

    try:
        db.session.commit()
        return jsonify({"ok": True, "id": inv.id, "message": "Cadastro atualizado com sucesso."})
    except IntegrityError:
        db.session.rollback()
        return jsonify({"ok": False, "duplicate": True, "error": "A alteração geraria um registro duplicado."}), 409
    except Exception as exc:
        db.session.rollback()
        return jsonify({"ok": False, "error": str(exc)}), 500


@app.delete("/api/inventory/<int:inventory_id>")
@manager_required
def delete_inventory(inventory_id):
    inv = db.session.get(Inventory, inventory_id)
    if not inv:
        return jsonify({"ok": False, "error": "Registro não encontrado."}), 404

    location_id = inv.location_id
    try:
        attachments = Attachment.query.filter_by(inventory_id=inventory_id).all()
        for attachment in attachments:
            try:
                if attachment.stored_name:
                    file_path = UPLOAD_DIR / attachment.stored_name
                    if file_path.exists():
                        file_path.unlink()
            except Exception:
                pass
            db.session.delete(attachment)

        db.session.add(AuditEvent(user_id=session.get("user_id"), event_type="INVENTORY_DELETE", entity_type="inventory", entity_id=str(inv.id), detail=f"Exclusão por {session.get('name','Gestor')}: {inv.equipment_type} · {inv.asset_identifier}"))
        db.session.delete(inv)
        db.session.flush()

        remaining = Inventory.query.filter_by(location_id=location_id).count()
        if remaining == 0:
            loc = db.session.get(Location, location_id)
            if loc:
                loc.survey_status = "PENDENTE"
                loc.started_at = None
                loc.completed_at = None
                loc.completed_by = None

        db.session.commit()
        return jsonify({"ok": True, "id": inventory_id, "message": "Cadastro excluído com sucesso."})
    except Exception as exc:
        db.session.rollback()
        return jsonify({"ok": False, "error": str(exc)}), 500


def _location_completion_check(loc):
    """V16: fechamento só ocorre após conciliação e evidência mínima de campo."""
    inventories = Inventory.query.filter_by(location_id=loc.id).all()
    base_assets = [a for a in BaseAsset.query.all() if _base_asset_matches_location(a, loc)]

    expected_by_type = {}
    for a in base_assets:
        typ = _canonical_equipment_type(a.equipment_type)
        expected_by_type[typ] = expected_by_type.get(typ, 0) + max(int(a.quantity or 1), 1)
    # Compatibilidade com localidades cuja previsão veio da base resumida.
    expected_by_type["ATM"] = max(expected_by_type.get("ATM", 0), int(loc.expected_atm or 0))
    expected_by_type["VALIDADOR"] = max(expected_by_type.get("VALIDADOR", 0), int(loc.expected_validator or 0))
    expected_by_type["POS"] = max(expected_by_type.get("POS", 0), int(loc.expected_pos or 0))

    reconciled_by_type = {}
    missing_evidence = []
    missing_justification = []
    for inv in inventories:
        typ = _canonical_equipment_type(inv.equipment_type)
        reconciled_by_type[typ] = reconciled_by_type.get(typ, 0) + 1
        media_count = Attachment.query.filter_by(inventory_id=inv.id).count()
        if media_count < 1:
            missing_evidence.append(inv.asset_identifier or f"registro {inv.id}")
        if normalize(inv.operational_status) == "NAO ENCONTRADO" and not (inv.notes or "").strip():
            missing_justification.append(inv.asset_identifier or f"registro {inv.id}")

    pending = []
    for typ, expected in expected_by_type.items():
        if expected <= 0:
            continue
        done = reconciled_by_type.get(typ, 0)
        if done < expected:
            pending.append({"type": typ, "expected": expected, "reconciled": done, "remaining": expected-done})

    errors = []
    if sum(expected_by_type.values()) == 0 and not inventories:
        errors.append("Nenhum equipamento foi registrado nesta localidade.")
    if pending:
        errors.append("Existem equipamentos previstos ainda não conciliados.")
    if missing_evidence:
        errors.append("Todos os registros precisam de ao menos uma foto/vídeo como evidência.")
    if missing_justification:
        errors.append("Itens marcados como Não encontrado exigem justificativa em Observações.")
    return {
        "ok": not errors, "errors": errors, "pending": pending,
        "missing_evidence": missing_evidence, "missing_justification": missing_justification,
        "registered": len(inventories), "expected": sum(expected_by_type.values())
    }


@app.get("/api/location/<int:location_id>/completion-check")
@field_required
def location_completion_check(location_id):
    loc = db.session.get(Location, location_id)
    if not loc:
        return jsonify({"ok": False, "error": "Local inválido."}), 404
    result = _location_completion_check(loc)
    return jsonify(result), (200 if result["ok"] else 409)


@app.post("/api/location/<int:location_id>/complete")
@field_required
def complete_location(location_id):
    loc = db.session.get(Location, location_id)
    if not loc:
        return jsonify({"ok": False, "error": "Local inválido."}), 404
    result = _location_completion_check(loc)
    if not result["ok"]:
        return jsonify({"ok": False, "error": "Localidade ainda não pode ser concluída.", **result}), 409

    loc.survey_status = "CONCLUIDA"
    loc.completed_at = datetime.utcnow()
    loc.completed_by = session["user_id"]
    db.session.commit()
    return jsonify({"ok": True, "message": "Localidade concluída após validação V16."})


@app.post("/api/location/<int:location_id>/reopen")
@manager_required
def reopen_location(location_id):
    loc = db.session.get(Location, location_id)
    if not loc:
        return jsonify({"ok": False, "error": "Local inválido."}), 404

    loc.survey_status = "EM ANDAMENTO"
    loc.completed_at = None
    loc.completed_by = None
    db.session.commit()
    return jsonify({"ok": True})



@app.post("/api/location/<int:location_id>/reference-position")
@manager_required
def save_location_reference_position(location_id):
    loc = db.session.get(Location, location_id)
    if not loc:
        return jsonify({"ok": False, "error": "Localidade não encontrada."}), 404

    data = request.get_json(silent=True) or {}
    try:
        lat = float(data.get("latitude"))
        lon = float(data.get("longitude"))
    except (TypeError, ValueError):
        return jsonify({"ok": False, "error": "Latitude/longitude inválidas."}), 400

    if not (-90 <= lat <= 90 and -180 <= lon <= 180):
        return jsonify({"ok": False, "error": "Coordenadas fora do intervalo válido."}), 400

    loc.reference_latitude = lat
    loc.reference_longitude = lon
    loc.reference_source = (data.get("source") or "Gestor - mapa").strip()[:120]
    loc.reference_updated_at = datetime.utcnow()
    db.session.commit()

    return jsonify({
        "ok": True,
        "location_id": loc.id,
        "reference_latitude": loc.reference_latitude,
        "reference_longitude": loc.reference_longitude,
        "reference_source": loc.reference_source
    })


@app.get("/api/gps/recent")
@dashboard_required
def api_recent_gps():
    """V66 REV4 — última posição diária em lote, sem consultas por técnico."""
    local_day=datetime.now(ZoneInfo("America/Sao_Paulo")).date()
    start_local=datetime.combine(local_day,datetime.min.time(),tzinfo=ZoneInfo("America/Sao_Paulo")); end_local=start_local+timedelta(days=1)
    start_utc=start_local.astimezone(ZoneInfo("UTC")).replace(tzinfo=None); end_utc=end_local.astimezone(ZoneInfo("UTC")).replace(tzinfo=None)
    users_q=User.query.filter(User.active.is_(True),User.role.in_(("technician","technician_implantation")))
    if session.get("role")=="technician": users_q=users_q.filter(User.id==session.get("user_id"))
    users=users_q.order_by(User.name).all(); user_ids=[u.id for u in users]
    if not user_ids:
        total_inventory=Inventory.query.count(); gps_inventory=Inventory.query.filter(Inventory.latitude.isnot(None),Inventory.longitude.isnot(None)).count()
        return jsonify({"date":local_day.isoformat(),"summary":{"total_inventory":total_inventory,"with_gps":gps_inventory,"without_gps":max(0,total_inventory-gps_inventory),"coverage_pct":round((gps_inventory/total_inventory*100),1) if total_inventory else 0,"technicians_today":0},"items":[]})

    pos_max=(db.session.query(TechnicianPosition.user_id.label("uid"),func.max(TechnicianPosition.captured_at).label("max_at"))
             .filter(TechnicianPosition.user_id.in_(user_ids),TechnicianPosition.captured_at>=start_utc,TechnicianPosition.captured_at<end_utc)
             .group_by(TechnicianPosition.user_id).subquery())
    pos_rows=(db.session.query(TechnicianPosition)
              .join(pos_max,and_(TechnicianPosition.user_id==pos_max.c.uid,TechnicianPosition.captured_at==pos_max.c.max_at)).all())
    pos_map={p.user_id:p for p in pos_rows}

    chk_max=(db.session.query(TechnicianCheckin.user_id.label("uid"),func.max(TechnicianCheckin.created_at).label("max_at"))
             .filter(TechnicianCheckin.user_id.in_(user_ids),TechnicianCheckin.created_at>=start_utc,TechnicianCheckin.created_at<end_utc)
             .group_by(TechnicianCheckin.user_id).subquery())
    check_rows=(db.session.query(TechnicianCheckin)
                .join(chk_max,and_(TechnicianCheckin.user_id==chk_max.c.uid,TechnicianCheckin.created_at==chk_max.c.max_at)).all())
    check_map={c.user_id:c for c in check_rows}
    loc_ids={c.location_id for c in check_rows if c.location_id}
    loc_map={loc.id:loc for loc in Location.query.filter(Location.id.in_(loc_ids)).all()} if loc_ids else {}

    items=[]
    for u in users:
        pos=pos_map.get(u.id)
        if not pos: continue
        checkin=check_map.get(u.id); loc=loc_map.get(checkin.location_id) if checkin else None
        items.append({"inventory_id":None,"location_id":loc.id if loc else None,"location_name":loc.location if loc else "Posição atual","company":loc.company if loc else (u.company or "Equipe de campo"),"line":loc.line if loc else "","equipment_type":"Última posição do dia","asset_identifier":"","technician":u.name,"technician_code":u.user_code or "","technician_user_id":u.id,"technician_photo_url":f"/usuarios/{u.id}/foto" if u.photo_url else None,"latitude":pos.latitude,"longitude":pos.longitude,"gps_accuracy":pos.accuracy,"gps_captured_at":pos.captured_at.isoformat(timespec="seconds")+"Z","created_at":pos.captured_at.isoformat(timespec="seconds")+"Z","_team_current":True})
    items.sort(key=lambda x:x["gps_captured_at"],reverse=True)
    total_inventory=Inventory.query.count(); gps_inventory=Inventory.query.filter(Inventory.latitude.isnot(None),Inventory.longitude.isnot(None)).count()
    resp=jsonify({"date":local_day.isoformat(),"summary":{"total_inventory":total_inventory,"with_gps":gps_inventory,"without_gps":max(0,total_inventory-gps_inventory),"coverage_pct":round((gps_inventory/total_inventory*100),1) if total_inventory else 0,"technicians_today":len(items)},"items":items})
    resp.headers["X-Autopass-Query-Mode"]="batch-latest-per-user"
    return resp


@app.get("/api/dashboard")
@dashboard_required
def dashboard():
    total = Location.query.count()
    pending = Location.query.filter_by(survey_status="PENDENTE").count()
    progress = Location.query.filter_by(survey_status="EM ANDAMENTO").count()
    completed = Location.query.filter_by(survey_status="CONCLUIDA").count()

    inv_rows = (
        db.session.query(Inventory.equipment_type, func.count(Inventory.id))
        .group_by(Inventory.equipment_type)
        .all()
    )
    inventoried_by_type = {"ATM":0,"VALIDADOR":0,"POS":0,"TDI":0,"BLOQUEIO":0}
    inventoried_total = 0
    for typ, count in inv_rows:
        count = int(count or 0)
        inventoried_total += count
        canonical = _canonical_equipment_type(typ)
        if canonical in inventoried_by_type:
            inventoried_by_type[canonical] += count

    classified_inventoried = sum(inventoried_by_type.values())
    unclassified = max(0, inventoried_total - classified_inventoried)
    inoperative = Inventory.query.filter_by(operational_status="Inoperante").count()
    divergences = Inventory.query.filter(
        Inventory.divergence.isnot(None),
        Inventory.divergence.notin_(["", "Não", "Nao"])
    ).count()

    technical_tdi_expected = TECHNICAL_TDI_TOTAL

    by_type = []
    for typ in ["ATM", "VALIDADOR", "POS", "TDI", "BLOQUEIO"]:
        # TDI continua como controle técnico separado do parque executivo oficial.
        exp = technical_tdi_expected if typ == "TDI" else int(OFFICIAL_PARK.get(typ, 0))
        inv = int(inventoried_by_type.get(typ, 0))
        by_type.append({
            "type": typ,
            "expected": exp,
            "inventoried": inv,
            "missing": max(0, exp - inv) if exp else 0,
            "coverage_pct": round(min(100, inv / exp * 100), 1) if exp else 0,
            "official": typ in OFFICIAL_PARK,
        })

    companies = (
        db.session.query(
            Location.company.label("company"),
            func.count(Location.id).label("total"),
            func.sum(case((Location.survey_status == "PENDENTE", 1), else_=0)).label("pending"),
            func.sum(case((Location.survey_status == "EM ANDAMENTO", 1), else_=0)).label("progress"),
            func.sum(case((Location.survey_status == "CONCLUIDA", 1), else_=0)).label("completed"),
        )
        .group_by(Location.company)
        .order_by(Location.company)
        .all()
    )

    official_inventoried = sum(
        min(int(inventoried_by_type.get(typ, 0)), int(exp))
        for typ, exp in OFFICIAL_PARK.items()
    )

    # V22.1 — os módulos executivos nunca podem derrubar os KPIs principais.
    today_utc = datetime.utcnow().date()
    start_14 = today_utc - timedelta(days=13)
    daily_counts = {(start_14 + timedelta(days=i)).isoformat(): 0 for i in range(14)}
    top_technicians = []
    evidence_kpis = {"visits": 0, "items": 0, "media": 0, "matched": 0, "review": 0, "unresolved_visits": 0}
    try:
        recent_rows = Inventory.query.filter(
            Inventory.created_at >= datetime.combine(start_14, datetime.min.time())
        ).all()
        tech_counts = {}
        for inv in recent_rows:
            if inv.created_at:
                key = inv.created_at.date().isoformat()
                if key in daily_counts:
                    daily_counts[key] += 1
            if inv.technician_id is not None:
                tech_counts[inv.technician_id] = tech_counts.get(inv.technician_id, 0) + 1
        ids = [uid for uid in tech_counts.keys() if uid is not None]
        user_names = {u.id: u.name for u in User.query.filter(User.id.in_(ids or [-1])).all()}
        top_technicians = [
            {"user_id": uid, "name": user_names.get(uid, f"Usuário {uid}"), "count": count}
            for uid, count in sorted(tech_counts.items(), key=lambda x: (-x[1], user_names.get(x[0], "")))[:8]
        ]
    except Exception as exc:
        app.logger.exception("V22.1: falha não crítica ao montar tendência/produtividade: %s", exc)
    try:
        evidence_kpis = _evidence_summary()
    except Exception as exc:
        app.logger.exception("V22.1: falha não crítica ao montar KPIs de evidências: %s", exc)

    # V37 — produtividade operacional por técnico, com volume, localidades e qualidade.
    productivity_37 = []
    try:
        start_dt = datetime.combine(start_14, datetime.min.time())
        rows37 = (db.session.query(
            Inventory.technician_id.label("uid"),
            func.count(Inventory.id).label("items"),
            func.count(func.distinct(Inventory.location_id)).label("locations"),
            func.count(func.distinct(func.date(Inventory.created_at))).label("days"),
            func.sum(case((Inventory.operational_status == "Inoperante", 1), else_=0)).label("inoperative"),
            func.sum(case((and_(Inventory.divergence.isnot(None), ~Inventory.divergence.in_(("", "Não", "Nao"))), 1), else_=0)).label("divergences"),
        ).filter(Inventory.created_at >= start_dt).group_by(Inventory.technician_id).all())
        uids=[r.uid for r in rows37 if r.uid]
        names={u.id:u.name for u in User.query.filter(User.id.in_(uids or [-1])).all()}
        for r in rows37:
            days=max(1,int(r.days or 0)); items=int(r.items or 0)
            productivity_37.append({"user_id":r.uid,"name":names.get(r.uid,f"Usuário {r.uid}"),"items":items,"locations":int(r.locations or 0),"active_days":int(r.days or 0),"avg_day":round(items/days,1),"divergences":int(r.divergences or 0),"inoperative":int(r.inoperative or 0)})
        productivity_37.sort(key=lambda x:(-x["items"],x["name"]))
    except Exception as exc:
        app.logger.exception("V37: falha não crítica na produtividade detalhada: %s", exc)

    return jsonify({
        "release": DASHBOARD_RELEASE,
        "official_park": {
            "total": OFFICIAL_PARK_TOTAL,
            "by_type": OFFICIAL_PARK,
            "technical_tdi": TECHNICAL_TDI_TOTAL,
            "note": "TDI é acompanhado separadamente e não compõe o total oficial de 3.801."
        },
        "technical_tdi": {
            "expected": TECHNICAL_TDI_TOTAL,
            "inventoried": int(inventoried_by_type.get("TDI", 0)),
            "missing": max(0, TECHNICAL_TDI_TOTAL - int(inventoried_by_type.get("TDI", 0))),
            "coverage_pct": round(min(100, int(inventoried_by_type.get("TDI", 0)) / TECHNICAL_TDI_TOTAL * 100), 1) if TECHNICAL_TDI_TOTAL else 0,
        },
        "totals": {
            "total": total,
            "pending": pending,
            "progress": progress,
            "completed": completed,
            "expected": OFFICIAL_PARK_TOTAL,
            "missing": max(0, OFFICIAL_PARK_TOTAL - official_inventoried),
        },
        "inventory": {
            "inventoried": inventoried_total,
            "official_inventoried": official_inventoried,
            "classified": classified_inventoried,
            "unclassified": unclassified,
            "inoperative": inoperative,
            "divergences": divergences,
        },
        "by_type": by_type,
        "trend_14d": [{"date": d, "count": c} for d, c in daily_counts.items()],
        "top_technicians_14d": top_technicians,
        "productivity_v37": productivity_37,
        "evidence": evidence_kpis,
        "by_company": [{
            "company": x.company,
            "total": int(x.total or 0),
            "pending": int(x.pending or 0),
            "progress": int(x.progress or 0),
            "completed": int(x.completed or 0),
        } for x in companies],
    })



def _v37_action_payload():
    actions=[]
    divs=Inventory.query.filter(Inventory.divergence.isnot(None), ~Inventory.divergence.in_(("", "Não", "Nao"))).order_by(Inventory.created_at.desc()).limit(25).all()
    for x in divs:
        loc=db.session.get(Location,x.location_id)
        actions.append({"key":f"DIV:{x.id}","category":"DIVERGENCIA","severity":"ALTA","title":f"Divergência · {x.asset_identifier}","detail":f"{loc.location if loc else 'Localidade'} · {x.divergence}","source_url":"/patrimonio"})
    stalled=Location.query.filter(Location.survey_status=="EM ANDAMENTO", Location.started_at.isnot(None), Location.started_at < datetime.utcnow()-timedelta(days=2)).order_by(Location.started_at).limit(20).all()
    for x in stalled:
        days=max(0,(datetime.utcnow()-x.started_at).days)
        actions.append({"key":f"PARADA:{x.id}","category":"LOCALIDADE_PARADA","severity":"MEDIA","title":f"Localidade em andamento há {days}d", "detail":f"{x.location} · {x.company} · {x.line}","source_url":"/gerencial"})
    return actions

@app.get("/api/v37/actions")
@dashboard_required
def v37_actions():
    generated=_v37_action_payload()
    keys=[x["key"] for x in generated]
    saved={x.action_key:x for x in OperationalAction.query.filter(OperationalAction.action_key.in_(keys or ["-"])).all()}
    out=[]
    for g in generated:
        row=saved.get(g["key"])
        out.append({**g,"status":row.status if row else "NOVO","owner_user_id":row.owner_user_id if row else None,"due_date":row.due_date.isoformat() if row and row.due_date else None})
    order={"NOVO":0,"EM ANALISE":1,"CORRIGIR EM CAMPO":2,"APROVADO":3,"DESCARTADO":4}
    out.sort(key=lambda x:(order.get(x["status"],9), 0 if x["severity"]=="ALTA" else 1))
    return jsonify({"ok":True,"items":out,"counts":{k:sum(1 for x in out if x["status"]==k) for k in order}})

@app.post("/api/v37/actions/update")
@dashboard_required
def v37_action_update():
    data=request.get_json(silent=True) or {}
    key=(data.get("key") or "")[:220]
    status=(data.get("status") or "NOVO").upper()
    allowed={"NOVO","EM ANALISE","CORRIGIR EM CAMPO","APROVADO","DESCARTADO"}
    if not key or status not in allowed: return jsonify({"ok":False,"error":"Dados inválidos"}),400
    gen=next((x for x in _v37_action_payload() if x["key"]==key),None)
    row=OperationalAction.query.filter_by(action_key=key).first()
    if not row:
        if not gen: return jsonify({"ok":False,"error":"Ação não encontrada"}),404
        row=OperationalAction(action_key=key,category=gen["category"],title=gen["title"],detail=gen["detail"],severity=gen["severity"],source_url=gen.get("source_url"))
        db.session.add(row)
    row.status=status; row.updated_at=datetime.utcnow()
    db.session.add(AuditEvent(user_id=session.get("user_id"),event_type="ACTION_STATUS",entity_type="OperationalAction",entity_id=key,detail=f"Status alterado para {status}"))
    db.session.commit()
    return jsonify({"ok":True})


@app.get("/api/export/excel")
@dashboard_required
def export_dashboard_excel():
    company = (request.args.get("company") or "").strip()
    line = (request.args.get("line") or "").strip()
    equipment_type = _canonical_equipment_type(request.args.get("type") or "")
    if equipment_type == "OUTRO":
        equipment_type = ""

    expected_map = _expected_assets_by_location()
    locations = Location.query.order_by(Location.company, Location.line, Location.location).all()

    inv_rows = (
        db.session.query(
            Inventory.location_id,
            Inventory.equipment_type,
            func.count(Inventory.id),
            func.coalesce(func.sum(case((Inventory.operational_status == "Inoperante", 1), else_=0)), 0),
            func.coalesce(func.sum(case((
                Inventory.divergence.isnot(None)
                & (~Inventory.divergence.in_(["", "Não", "Nao"]))
            , 1), else_=0)), 0),
        )
        .group_by(Inventory.location_id, Inventory.equipment_type)
        .all()
    )

    inv_map = {}
    for loc_id, typ, count, inop, divs in inv_rows:
        b = inv_map.setdefault(loc_id, {
            "total":0, "inoperative":0, "divergences":0,
            "by_type":{"ATM":0,"VALIDADOR":0,"POS":0,"TDI":0,"BLOQUEIO":0,"OUTRO":0}
        })
        canonical = _canonical_equipment_type(typ)
        if canonical not in b["by_type"]:
            canonical = "OUTRO"
        b["by_type"][canonical] += int(count or 0)
        b["total"] += int(count or 0)
        b["inoperative"] += int(inop or 0)
        b["divergences"] += int(divs or 0)

    selected = []
    for loc in locations:
        if company and loc.company != company:
            continue
        if line and loc.line != line:
            continue
        exp = expected_map.get(loc.id, {})
        inv = inv_map.get(loc.id, {
            "total":0,"inoperative":0,"divergences":0,
            "by_type":{"ATM":0,"VALIDADOR":0,"POS":0,"TDI":0,"BLOQUEIO":0,"OUTRO":0}
        })
        if equipment_type and int(exp.get(equipment_type, 0) or 0) <= 0 and int(inv["by_type"].get(equipment_type,0) or 0) <= 0:
            continue
        selected.append((loc, exp, inv))

    wb = Workbook()
    ws = wb.active
    ws.title = "Resumo Executivo"

    title_fill = PatternFill("solid", fgColor="17365D")
    header_fill = PatternFill("solid", fgColor="D9EAF7")
    white_font = Font(color="FFFFFF", bold=True)
    bold = Font(bold=True)

    ws["A1"] = "Inventário Autopass — Resumo Executivo"
    ws["A1"].font = Font(size=16, bold=True, color="FFFFFF")
    ws["A1"].fill = title_fill
    ws.merge_cells("A1:F1")
    ws["A2"] = f"Gerado em {datetime.now().strftime('%d/%m/%Y %H:%M')}"
    ws["A3"] = f"Filtros: Empresa={company or 'Todas'} | Linha={line or 'Todas'} | Tipo={equipment_type or 'Todos'}"

    official_rows = [
        ("ATM", OFFICIAL_PARK["ATM"]),
        ("POS", OFFICIAL_PARK["POS"]),
        ("Recarga", OFFICIAL_PARK["VALIDADOR"]),
        ("Bloqueio", OFFICIAL_PARK["BLOQUEIO"]),
        ("TOTAL OFICIAL", OFFICIAL_PARK_TOTAL),
    ]
    ws.append([])
    ws.append(["Produto", "Parque oficial"])
    for cell in ws[5]:
        cell.fill = header_fill
        cell.font = bold
    for row in official_rows:
        ws.append(list(row))

    wsl = wb.create_sheet("Localidades")
    headers = [
        "Empresa","Linha","Localidade","Status","Tipo filtrado",
        "Previsto","Inventariado","Faltante","Cobertura %",
        "Divergências","Inoperantes"
    ]
    wsl.append(headers)
    for cell in wsl[1]:
        cell.fill = title_fill
        cell.font = white_font
        cell.alignment = Alignment(horizontal="center")

    for loc, exp, inv in selected:
        if equipment_type:
            expected = int(exp.get(equipment_type, 0) or 0)
            inventoried = int(inv["by_type"].get(equipment_type, 0) or 0)
        else:
            expected = int(sum(exp.values()))
            inventoried = int(inv["total"])
        missing = max(0, expected - inventoried)
        coverage = round(min(100, inventoried / expected * 100), 1) if expected else 0
        wsl.append([
            loc.company, loc.line, loc.location, loc.survey_status,
            equipment_type or "Todos", expected, inventoried, missing, coverage,
            int(inv["divergences"]), int(inv["inoperative"])
        ])

    wsi = wb.create_sheet("Inventário Realizado")
    inv_headers = [
        "ID","Empresa","Linha","Localidade","Tipo","Identificação","Série",
        "Modelo","Status","Divergência","Data","Latitude","Longitude"
    ]
    wsi.append(inv_headers)
    for cell in wsi[1]:
        cell.fill = title_fill
        cell.font = white_font

    query = (
        db.session.query(Inventory, Location)
        .join(Location, Location.id == Inventory.location_id)
        .order_by(Inventory.created_at.desc())
    )
    if company:
        query = query.filter(Location.company == company)
    if line:
        query = query.filter(Location.line == line)
    for inv, loc in query.all():
        canonical = _canonical_equipment_type(inv.equipment_type)
        if equipment_type and canonical != equipment_type:
            continue
        wsi.append([
            inv.id, loc.company, loc.line, loc.location, canonical,
            inv.asset_identifier, inv.serial, inv.model, inv.operational_status,
            inv.divergence, inv.created_at.strftime("%d/%m/%Y %H:%M") if inv.created_at else "",
            inv.latitude, inv.longitude
        ])


    # Evidências de campo (WhatsApp) — fonte separada do inventário canônico.
    if FieldEvidenceVisit.query.count():
        ws_ev = wb.create_sheet("Evidências de Campo")
        ev_headers = [
            "Data", "Hora", "Responsável", "Estação informada", "Linha informada",
            "Localidade associada", "Confiança", "Tipo", "Identificador", "Modelo",
            "Série", "Patrimônio", "Status operacional", "Auditoria", "Mídias"
        ]
        ws_ev.append(ev_headers)
        for c in ws_ev[1]:
            c.font = Font(bold=True)
        evidence_rows = (
            db.session.query(FieldEvidenceVisit, FieldEvidenceItem)
            .join(FieldEvidenceItem, FieldEvidenceItem.visit_id == FieldEvidenceVisit.id)
            .order_by(FieldEvidenceVisit.source_date, FieldEvidenceVisit.source_time)
            .all()
        )
        media_counts = dict(
            db.session.query(FieldEvidenceMedia.visit_id, func.count(FieldEvidenceMedia.id))
            .group_by(FieldEvidenceMedia.visit_id).all()
        )
        for visit, item in evidence_rows:
            loc = db.session.get(Location, visit.location_id) if visit.location_id else None
            ws_ev.append([
                visit.source_date, visit.source_time, visit.author, visit.station_raw, visit.line_raw,
                loc.location if loc else "", visit.match_confidence,
                item.equipment_type, item.identifier, item.model or "", item.serial or "",
                item.patrimony or "", item.operational_status or "", item.audit_status or "",
                int(media_counts.get(visit.id, 0)),
            ])

    for sheet in wb.worksheets:
        for col in range(1, sheet.max_column + 1):
            letter = get_column_letter(col)
            width = 10
            for row in range(1, min(sheet.max_row, 300) + 1):
                value = sheet.cell(row=row, column=col).value
                if value is not None:
                    width = max(width, min(42, len(str(value)) + 2))
            sheet.column_dimensions[letter].width = width
        sheet.freeze_panes = "A2" if sheet.title != "Resumo Executivo" else "A5"

    output = io.BytesIO()
    wb.save(output)
    output.seek(0)
    filename = f"autopass_dashboard_{datetime.now().strftime('%Y%m%d_%H%M')}.xlsx"
    return send_file(
        output,
        as_attachment=True,
        download_name=filename,
        mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )


def _media_thumbnail(raw, max_px=520, quality=68):
    from PIL import Image, ImageOps
    """V52.8: miniatura leve para galerias; original continua disponível sob demanda."""
    try:
        with Image.open(io.BytesIO(raw)) as im:
            im=ImageOps.exif_transpose(im)
            im.thumbnail((max_px,max_px), Image.Resampling.LANCZOS)
            if im.mode != "RGB":
                bg=Image.new("RGB",im.size,"white")
                if "A" in im.getbands(): bg.paste(im,mask=im.getchannel("A"))
                else: bg.paste(im)
                im=bg
            out=io.BytesIO(); im.save(out,format="JPEG",quality=quality,optimize=True,progressive=True); return out.getvalue()
    except Exception:
        return None

def _cached_media_response(raw, mimetype, filename=None, thumb=False):
    if thumb and (mimetype or "").startswith("image/"):
        compact=_media_thumbnail(raw)
        if compact:
            resp=send_file(io.BytesIO(compact),mimetype="image/jpeg",download_name=(Path(filename or "thumb.jpg").stem+"-thumb.jpg"),max_age=604800)
        else:
            resp=send_file(io.BytesIO(raw),mimetype=mimetype,download_name=filename,max_age=604800)
    else:
        resp=send_file(io.BytesIO(raw),mimetype=mimetype,download_name=filename,max_age=604800)
    resp.headers["Cache-Control"]="private, max-age=604800, immutable"
    return resp

@app.route("/uploads/<path:name>")
@login_required
def uploaded(name):
    # V62 REV2: nunca baixar objetos R2 grandes para a RAM do worker.
    # Após validar a sessão, o backend entrega uma URL temporária assinada e o
    # navegador baixa a mídia diretamente do R2. Isso evita OOM e libera o único
    # worker Gunicorn para atender a operação de campo.
    if name.startswith("r2__"):
        key = name[4:]
        try:
            url = r2_client().generate_presigned_url(
                "get_object",
                Params={"Bucket": os.environ["R2_BUCKET_NAME"], "Key": key},
                ExpiresIn=900,
            )
            resp = redirect(url, code=302)
            resp.headers["Cache-Control"] = "private, max-age=300"
            return resp
        except Exception:
            app.logger.exception("Falha ao gerar URL temporária R2 para %s", key)
            abort(404)
    path=UPLOAD_DIR/name
    # Para arquivos locais, miniaturas antigas continuam disponíveis, mas sem
    # materializar o original em bytes quando não houver necessidade.
    if request.args.get("thumb") == "1" and path.exists():
        mime = mimetypes.guess_type(name)[0] or "application/octet-stream"
        if mime.startswith("image/"):
            try:
                # Thumbnail local é usada apenas para arquivos locais; estes são
                # poucos e não são a origem do incidente de memória observado.
                return _cached_media_response(path.read_bytes(), mime, Path(name).name, thumb=True)
            except Exception:
                pass
    return send_from_directory(UPLOAD_DIR,name,max_age=604800)


@app.get("/api/inventory/<int:inventory_id>/attachments")
@login_required
def attachments(inventory_id):
    rows = Attachment.query.filter_by(inventory_id=inventory_id).all()
    return jsonify([{
        "id": a.id,
        "inventory_id": a.inventory_id,
        "original_name": a.original_name,
        "stored_name": a.stored_name,
        "mime_type": a.mime_type,
        "storage": "R2" if (a.stored_name or "").startswith("r2__") else "LOCAL",
        "url": url_for("uploaded", name=a.stored_name)
    } for a in rows])


@app.route('/perfis', methods=['GET','POST'])
@login_required
def system_profiles_page():
    if session.get('role')!='manager' and not _has_access('management.profiles'): abort(403)
    if request.method=='POST':
        name=(request.form.get('name') or '').strip(); base=(request.form.get('base_role') or 'none').strip()
        if not name: flash('Informe o nome do perfil.'); return redirect('/perfis')
        if base not in ('none','technician','technician_implantation','manager_field','consultation','dispatcher','hr','atm_financial_admin'): base='none'
        if SystemProfile.query.filter(func.lower(SystemProfile.name)==name.lower()).first(): flash('Já existe um perfil com esse nome.'); return redirect('/perfis')
        row=SystemProfile(name=name,base_role=base,created_by=session['user_id'])
        row.access_json=json.dumps([x for x in request.form.getlist('access') if x in ACCESS_SUBMODULES],ensure_ascii=False)
        db.session.add(row); db.session.commit(); flash('Perfil criado.'); return redirect('/perfis')
    profiles=SystemProfile.query.order_by(SystemProfile.active.desc(),SystemProfile.name).all()
    counts=dict(db.session.query(User.system_profile_id,func.count(User.id)).filter(User.system_profile_id.isnot(None)).group_by(User.system_profile_id).all())
    return render_template('system_profiles.html',profiles=profiles,profile_counts=counts,access_groups=ACCESS_GROUPS,access_labels=ACCESS_LABELS)

@app.post('/perfis/<int:pid>/salvar')
@login_required
def system_profile_save(pid):
    if session.get('role')!='manager' and not _has_access('management.profiles'): abort(403)
    p=db.session.get(SystemProfile,pid) or abort(404)
    name=(request.form.get('name') or '').strip(); base=(request.form.get('base_role') or p.base_role or 'none').strip()
    if base not in ('none','technician','technician_implantation','manager_field','consultation','dispatcher','hr','atm_financial_admin'): base='none'
    if name: p.name=name
    p.base_role=base; p.active=request.form.get('active')=='1'
    p.access_json=json.dumps([x for x in request.form.getlist('access') if x in ACCESS_SUBMODULES],ensure_ascii=False)
    db.session.commit(); flash('Perfil atualizado.'); return redirect('/perfis')

@app.post('/perfis/<int:pid>/excluir')
@login_required
def system_profile_delete(pid):
    if session.get('role')!='manager' and not _has_access('management.profiles'): abort(403)
    p=db.session.get(SystemProfile,pid)
    if p:
        linked=User.query.filter_by(system_profile_id=p.id).count()
        if linked: flash(f'Perfil vinculado a {linked} usuário(s). Inative-o antes de remover vínculos.'); return redirect('/perfis')
        db.session.delete(p); db.session.commit(); flash('Perfil excluído.')
    return redirect('/perfis')

def _resolve_role_selection(raw):
    raw=(raw or 'technician').strip(); profile=None
    if raw.startswith('custom:'):
        try: profile=db.session.get(SystemProfile,int(raw.split(':',1)[1]))
        except Exception: profile=None
        if not profile or not profile.active: return None,None
        role=profile.base_role if profile.base_role and profile.base_role!='none' else 'technician'
        return role,profile
    return raw,None

@app.route("/usuarios")
@user_admin_required
def users_page():
    if not _has_access("users.view"): abort(403)
    # V71.1 — acessos externos do Portal são administrados no Cadastro de Clientes,
    # não em RH / Usuários.
    active_q = User.query.filter(User.archived_at.is_(None), User.role != "customer")
    archived_q = User.query.filter(User.archived_at.isnot(None), User.role != "customer")
    if session.get("role") == "hr":
        active_q = active_q.filter(User.role.in_(("technician", "technician_implantation")))
        archived_q = archived_q.filter(User.role.in_(("technician", "technician_implantation")))
    active_users = active_q.order_by(User.active.desc(), User.name).all()
    archived_users = archived_q.order_by(User.archived_at.desc(), User.name).all()
    return render_template(
        "users.html",
        users=active_users,
        archived_users=archived_users,
        can_assign_sensitive_roles=_current_user_is_superadmin(),
        is_hr_admin=session.get("role") == "hr",
        customer_companies=CustomerCompany.query.filter_by(active=True).order_by(CustomerCompany.legal_name).all(),
        customer_company_map={u.id:_customer_company_ids(u) for u in active_users},
        system_profiles=SystemProfile.query.filter_by(active=True).order_by(SystemProfile.name).all(),
    )


def _normalize_optional_email(value):
    value = (value or "").strip().lower()
    return value or None


def _normalize_optional_phone(value):
    digits = re.sub(r"\D", "", value or "")
    return digits or None


def _next_user_code(role):
    prefixes = {
        "technician": "T",
        "technician_implantation": "TI",
        "manager_field": "GF",
        "manager": "G",
        "consultation": "C",
        "hr": "RH",
        "dispatcher": "D",
        "atm_financial_admin": "AF",
        "customer": "CL",
    }
    prefix = prefixes.get(role, "U")
    existing = (
        db.session.query(User.user_code)
        .filter(User.user_code.isnot(None), User.user_code.like(f"{prefix}%"))
        .all()
    )
    highest = 0
    for (code,) in existing:
        match = re.fullmatch(rf"{prefix}(\d+)", (code or "").upper())
        if match:
            highest = max(highest, int(match.group(1)))
    return f"{prefix}{highest + 1:03d}"


@app.post("/usuarios/novo")
@user_admin_required
def create_user():
    if not _has_access("users.create"): abort(403)
    name = request.form.get("name", "").strip()
    username = request.form.get("username", "").strip().lower()
    password = request.form.get("password", "")
    role_selection = request.form.get("role", "technician").strip()
    role, system_profile = _resolve_role_selection(role_selection)
    if not role:
        flash("Perfil de acesso inválido."); return redirect(url_for("users_page"))
    email = _normalize_optional_email(request.form.get("email"))
    phone = _normalize_optional_phone(request.form.get("phone"))
    company = request.form.get("company", "").strip() or None
    customer_company_ids=request.form.getlist("customer_company_ids") if role=="customer" else []
    customer_company_ids=[x for x in customer_company_ids if str(x).isdigit() and db.session.get(CustomerCompany,int(x))]
    if role=="customer" and len(customer_company_ids)>1 and not _current_user_is_superadmin():
        customer_company_ids=customer_company_ids[:1]
    if role=="customer" and customer_company_ids:
        c0=db.session.get(CustomerCompany,int(customer_company_ids[0])); company=c0.legal_name if c0 else company
    job_title = request.form.get("job_title", "").strip() or None
    personnel_status = request.form.get("personnel_status", "ATIVO").strip().upper() or "ATIVO"
    personnel_status_note = request.form.get("personnel_status_note", "").strip() or None
    gps_required = request.form.get("gps_required") == "1"
    work_schedule_type = request.form.get("work_schedule_type", "12x36").strip() or "12x36"
    work_shift = request.form.get("work_shift", "05:00-17:00").strip() or "05:00-17:00"
    work_anchor_status = request.form.get("work_anchor_status", "TRABALHA").strip() or "TRABALHA"
    work_anchor_date_raw = request.form.get("work_anchor_date", "").strip()
    try:
        work_anchor_date = datetime.strptime(work_anchor_date_raw, "%Y-%m-%d").date() if work_anchor_date_raw else None
    except ValueError:
        work_anchor_date = None

    if not name or not username or not password:
        flash("Nome, usuário e senha são obrigatórios.")
        return redirect(url_for("users_page"))

    if role not in ("manager", "manager_field", "technician", "technician_implantation", "consultation", "hr", "dispatcher", "atm_financial_admin", "customer"):
        flash("Perfil de acesso inválido.")
        return redirect(url_for("users_page"))
    if not _role_assignment_allowed(role):
        flash("Somente o Administrador principal pode criar ou atribuir os perfis Gestor, Consulta ou Dispatcher.")
        return redirect(url_for("users_page"))

    allowed_personnel_status = {"ATIVO", "FERIAS", "AFASTADO", "LICENCA", "FOLGA_PROGRAMADA", "OUTRO"}
    if personnel_status not in allowed_personnel_status:
        personnel_status = "ATIVO"

    if len(password) < 8:
        flash("A senha deve ter pelo menos 8 caracteres.")
        return redirect(url_for("users_page"))

    normalized_name = re.sub(r"\s+", " ", name).strip().lower()
    if User.query.filter(func.lower(func.trim(User.name)) == normalized_name).first():
        flash("Já existe um usuário cadastrado com este nome.")
        return redirect(url_for("users_page"))
    if User.query.filter(func.lower(User.username) == username).first():
        flash("Já existe um usuário com esse login.")
        return redirect(url_for("users_page"))

    if email and User.query.filter(func.lower(User.email) == email).first():
        flash("Já existe um usuário com esse e-mail.")
        return redirect(url_for("users_page"))

    if phone and User.query.filter(User.phone == phone).first():
        flash("Já existe um usuário com esse celular.")
        return redirect(url_for("users_page"))

    user_code = _next_user_code(role)
    user = User(
        name=name,
        username=username,
        password_hash=generate_password_hash(password),
        role=role,
        active=True,
        user_code=user_code,
        email=email,
        phone=phone,
        company=company,
        job_title=job_title,
        personnel_status=personnel_status,
        personnel_status_note=personnel_status_note,
        work_schedule_type=work_schedule_type if role in ("technician", "technician_implantation", "manager_field") else None,
        work_shift=work_shift if role in ("technician", "technician_implantation", "manager_field") else None,
        work_anchor_date=work_anchor_date if role in ("technician", "technician_implantation", "manager_field") and work_schedule_type == "12x36" else None,
        work_anchor_status=work_anchor_status if role in ("technician", "technician_implantation", "manager_field") and work_schedule_type == "12x36" else None,
        access_json=(system_profile.access_json if system_profile else json.dumps(_parse_access_form(role), ensure_ascii=False)),
        system_profile_id=(system_profile.id if system_profile else None),
        gps_required=(gps_required if role in ("technician","technician_implantation","manager_field","dispatcher") else False),
        customer_company_ids=json.dumps([int(x) for x in customer_company_ids]) if role=="customer" else None,
    )

    photo = request.files.get("photo")
    if photo and photo.filename and role in ("manager", "manager_field", "technician", "technician_implantation", "hr"):
        if not (photo.mimetype or "").startswith("image/"):
            flash("A foto do usuário deve ser uma imagem.")
            return redirect(url_for("users_page"))
        data = photo.read()
        if len(data) > 3 * 1024 * 1024:
            flash("A foto do usuário deve ter no máximo 3 MB.")
            return redirect(url_for("users_page"))
        safe_name = secure_filename(photo.filename) or "foto.jpg"
        object_key = f"usuarios/{user_code}/{uuid.uuid4().hex}-{safe_name}"
        _r2_put_bytes(object_key, data, photo.mimetype or "image/jpeg")
        user.photo_url = object_key

    try:
        db.session.add(user)
        db.session.flush()
        _sync_user_schedule_profile(user)
        db.session.commit()
    except IntegrityError:
        db.session.rollback()
        flash("Não foi possível criar o usuário porque existe um dado duplicado.")
        return redirect(url_for("users_page"))

    flash(f"Usuário {user.name} criado com sucesso. Código: {user.user_code}.")
    return redirect(url_for("users_page"))


@app.get("/usuarios/<int:user_id>/foto")
@login_required
def user_photo(user_id):
    user = db.session.get(User, user_id)
    if not user or not user.photo_url:
        return "", 404
    if session.get("role") == "hr" and user.role not in ("technician", "technician_implantation"):
        return "", 404

    try:
        obj = r2_client().get_object(
            Bucket=os.environ["R2_BUCKET_NAME"],
            Key=user.photo_url
        )
        raw=obj["Body"].read()
        mime=obj.get("ContentType") or "image/jpeg"
        if request.args.get("thumb")=="1":
            return _cached_media_response(raw,mime,"usuario.jpg",thumb=True)
        return Response(raw,mimetype=mime,headers={"Cache-Control":"private, max-age=86400, stale-while-revalidate=604800"})
    except Exception:
        return "", 404


def _active_manager_count(exclude_user_id=None):
    query = User.query.filter_by(role="manager", active=True)
    if exclude_user_id is not None:
        query = query.filter(User.id != exclude_user_id)
    return query.count()


@app.post("/usuarios/<int:user_id>/toggle")
@user_admin_required
def toggle_user(user_id):
    if not _has_access("users.activate"): abort(403)
    user = db.session.get(User, user_id)
    if not user:
        flash("Usuário não encontrado.")
        return redirect(url_for("users_page"))
    if not _hr_target_allowed(user):
        flash("RH pode administrar somente perfis operacionais autorizados.")
        return redirect(url_for("users_page"))

    if user.id == session.get("user_id"):
        flash("Você não pode desativar o próprio usuário enquanto está conectado.")
        return redirect(url_for("users_page"))

    if user.active and user.role == "manager" and _active_manager_count(exclude_user_id=user.id) == 0:
        flash("Não é possível desativar o último Gestor ativo.")
        return redirect(url_for("users_page"))

    user.active = not user.active
    db.session.commit()
    flash("Status do usuário atualizado.")
    return redirect(url_for("users_page"))


@app.post("/usuarios/<int:user_id>/editar")
@user_admin_required
def edit_user(user_id):
    if not _has_access("users.edit"): abort(403)
    user = db.session.get(User, user_id)
    if not user:
        flash("Usuário não encontrado.")
        return redirect(url_for("users_page"))
    if not _hr_target_allowed(user):
        flash("RH pode administrar somente perfis operacionais autorizados.")
        return redirect(url_for("users_page"))

    name = request.form.get("name", "").strip()
    username = request.form.get("username", "").strip().lower()
    role_selection = request.form.get("role", (f"custom:{user.system_profile_id}" if getattr(user,"system_profile_id",None) else user.role)).strip()
    role, system_profile = _resolve_role_selection(role_selection)
    if not role:
        flash("Perfil de acesso inválido."); return redirect(url_for("users_page"))
    user_code = request.form.get("user_code", "").strip().upper()
    email = _normalize_optional_email(request.form.get("email"))
    phone = _normalize_optional_phone(request.form.get("phone"))
    company = request.form.get("company", "").strip() or None
    customer_company_ids=request.form.getlist("customer_company_ids") if role=="customer" else []
    customer_company_ids=[x for x in customer_company_ids if str(x).isdigit() and db.session.get(CustomerCompany,int(x))]
    if role=="customer" and len(customer_company_ids)>1 and not _current_user_is_superadmin():
        customer_company_ids=customer_company_ids[:1]
    if role=="customer" and customer_company_ids:
        c0=db.session.get(CustomerCompany,int(customer_company_ids[0])); company=c0.legal_name if c0 else company
    job_title = request.form.get("job_title", user.job_title or "").strip() or None
    personnel_status = request.form.get("personnel_status", user.personnel_status or "ATIVO").strip().upper() or "ATIVO"
    personnel_status_note = request.form.get("personnel_status_note", user.personnel_status_note or "").strip() or None
    gps_required = request.form.get("gps_required") == "1"
    work_schedule_type = request.form.get("work_schedule_type", user.work_schedule_type or "12x36").strip() or "12x36"
    work_shift = request.form.get("work_shift", user.work_shift or "05:00-17:00").strip() or "05:00-17:00"
    work_anchor_status = request.form.get("work_anchor_status", user.work_anchor_status or "TRABALHA").strip() or "TRABALHA"
    work_anchor_date_raw = request.form.get("work_anchor_date", "").strip()
    try:
        work_anchor_date = datetime.strptime(work_anchor_date_raw, "%Y-%m-%d").date() if work_anchor_date_raw else user.work_anchor_date
    except ValueError:
        work_anchor_date = user.work_anchor_date

    if not name or not username:
        flash("Nome e usuário são obrigatórios.")
        return redirect(url_for("users_page"))

    if role not in ("manager", "manager_field", "technician", "technician_implantation", "consultation", "hr", "dispatcher", "atm_financial_admin", "customer"):
        flash("Perfil de acesso inválido.")
        return redirect(url_for("users_page"))
    if not _role_assignment_allowed(role):
        flash("Somente o Administrador principal pode atribuir os perfis Gestor ou Consulta.")
        return redirect(url_for("users_page"))

    allowed_personnel_status = {"ATIVO", "FERIAS", "AFASTADO", "LICENCA", "FOLGA_PROGRAMADA", "OUTRO"}
    if personnel_status not in allowed_personnel_status:
        personnel_status = "ATIVO"

    if not user_code:
        user_code = user.user_code or _next_user_code(role)

    normalized_name = re.sub(r"\s+", " ", name).strip().lower()
    duplicate_name = User.query.filter(User.id != user.id, func.lower(func.trim(User.name)) == normalized_name).first()
    if duplicate_name:
        flash("Já existe outro usuário cadastrado com este nome.")
        return redirect(url_for("users_page"))
    duplicate_username = User.query.filter(
        User.id != user.id,
        func.lower(User.username) == username
    ).first()
    if duplicate_username:
        flash("Já existe outro usuário com esse login.")
        return redirect(url_for("users_page"))

    duplicate_code = User.query.filter(
        User.id != user.id,
        func.upper(User.user_code) == user_code
    ).first()
    if duplicate_code:
        flash("Já existe outro usuário com esse código.")
        return redirect(url_for("users_page"))

    if email:
        duplicate_email = User.query.filter(
            User.id != user.id,
            func.lower(User.email) == email
        ).first()
        if duplicate_email:
            flash("Já existe outro usuário com esse e-mail.")
            return redirect(url_for("users_page"))

    if phone:
        duplicate_phone = User.query.filter(
            User.id != user.id,
            User.phone == phone
        ).first()
        if duplicate_phone:
            flash("Já existe outro usuário com esse celular.")
            return redirect(url_for("users_page"))

    # Evita retirar o último acesso administrativo do sistema.
    if user.role == "manager" and role != "manager":
        if user.id == session.get("user_id"):
            flash("Você não pode alterar o próprio perfil de Gestor enquanto está conectado.")
            return redirect(url_for("users_page"))
        if _active_manager_count(exclude_user_id=user.id) == 0:
            flash("Não é possível alterar o perfil do último Gestor ativo.")
            return redirect(url_for("users_page"))

    old_photo_key = user.photo_url
    new_photo_key = None

    photo = request.files.get("photo")
    if photo and photo.filename:
        if role == "consultation":
            flash("Foto não é necessária para o perfil Consulta.")
            return redirect(url_for("users_page"))
        if not (photo.mimetype or "").startswith("image/"):
            flash("A foto do usuário deve ser uma imagem.")
            return redirect(url_for("users_page"))

        data = photo.read()
        if len(data) > 3 * 1024 * 1024:
            flash("A foto do usuário deve ter no máximo 3 MB.")
            return redirect(url_for("users_page"))

        safe_name = secure_filename(photo.filename) or "foto.jpg"
        new_photo_key = f"usuarios/{user_code}/{uuid.uuid4().hex}-{safe_name}"
        _r2_put_bytes(new_photo_key, data, photo.mimetype or "image/jpeg")

    remove_photo = request.form.get("remove_photo") == "1"

    old_role = user.role
    user.name = name
    user.username = username
    user.role = role
    user.system_profile_id = system_profile.id if system_profile else None
    user.user_code = user_code
    user.email = email
    user.phone = phone
    user.company = company
    user.customer_company_ids=json.dumps([int(x) for x in customer_company_ids]) if role=="customer" else None
    user.job_title = job_title
    user.personnel_status = personnel_status
    user.personnel_status_note = personnel_status_note
    old_gps_required = bool(getattr(user, "gps_required", False))
    user.gps_required = gps_required if role in ("technician","technician_implantation","manager_field","dispatcher") else False
    if system_profile:
        user.access_json = system_profile.access_json
    elif session.get("role") == "manager":
        user.access_json = json.dumps(_parse_access_form(role), ensure_ascii=False)
    if role in ("technician", "technician_implantation", "manager_field"):
        user.work_schedule_type = work_schedule_type
        user.work_shift = work_shift
        user.work_anchor_date = work_anchor_date if work_schedule_type == "12x36" else None
        user.work_anchor_status = work_anchor_status if work_schedule_type == "12x36" else None
    else:
        user.work_schedule_type = None
        user.work_shift = None
        user.work_anchor_date = None
        user.work_anchor_status = None

    if new_photo_key:
        user.photo_url = new_photo_key
    elif remove_photo:
        user.photo_url = None

    try:
        _sync_user_schedule_profile(user)
        if old_gps_required != bool(user.gps_required):
            db.session.add(AuditEvent(user_id=session.get("user_id"),event_type="USER_GPS_REQUIREMENT_CHANGE",entity_type="user",entity_id=str(user.id),detail=f"{user.name} · GPS obrigatório: {old_gps_required} -> {bool(user.gps_required)}"))
        if old_role != role:
            db.session.add(AuditEvent(
                user_id=session.get("user_id"),
                event_type="USER_ROLE_CHANGE",
                entity_type="user",
                entity_id=str(user.id),
                detail=f"{user.name} · perfil {old_role} -> {role}",
            ))
        db.session.commit()
    except IntegrityError:
        db.session.rollback()
        if new_photo_key:
            try:
                r2_client().delete_object(
                    Bucket=os.environ["R2_BUCKET_NAME"],
                    Key=new_photo_key
                )
            except Exception:
                pass
        flash("Não foi possível salvar porque existe um dado duplicado.")
        return redirect(url_for("users_page"))

    # Remove do R2 a foto antiga somente depois de a atualização do banco ter sido confirmada.
    if old_photo_key and old_photo_key != user.photo_url:
        try:
            r2_client().delete_object(
                Bucket=os.environ["R2_BUCKET_NAME"],
                Key=old_photo_key
            )
        except Exception:
            pass

    if user.id == session.get("user_id"):
        session["name"] = user.name
        session["role"] = user.role

    flash(f"Usuário {user.name} atualizado com sucesso.")
    return redirect(url_for("users_page"))



def _user_operational_history_counts(user_id):
    return {
        "inventory": Inventory.query.filter_by(technician_id=user_id).count(),
        "gps": TechnicianPosition.query.filter_by(user_id=user_id).count(),
        "completed_locations": Location.query.filter_by(completed_by=user_id).count(),
    }


@app.post("/usuarios/<int:user_id>/excluir")
@manager_required
def delete_or_archive_user(user_id):
    if not _has_access("users.delete"): abort(403)
    user = db.session.get(User, user_id)
    if not user:
        flash("Usuário não encontrado.")
        return redirect(url_for("users_page"))

    if user.id == session.get("user_id"):
        flash("Você não pode excluir ou arquivar o próprio usuário enquanto está conectado.")
        return redirect(url_for("users_page"))

    if user.role == "manager" and user.active and _active_manager_count(exclude_user_id=user.id) == 0:
        flash("Não é possível excluir ou arquivar o último Gestor ativo.")
        return redirect(url_for("users_page"))

    history = _user_operational_history_counts(user.id)
    has_history = any(history.values())
    old_photo_key = user.photo_url
    schedule_profiles = TeamScheduleProfile.query.filter_by(user_id=user.id).all()

    if has_history:
        for profile in schedule_profiles:
            profile.active = False
            profile.user_id = None
            profile.updated_at = datetime.utcnow()

        user.active = False
        user.archived_at = datetime.utcnow()
        user.username = f"arquivado-{user.id}-{uuid.uuid4().hex[:10]}"
        user.password_hash = generate_password_hash(uuid.uuid4().hex + uuid.uuid4().hex)
        user.email = None
        user.phone = None
        user.photo_url = None
        db.session.commit()

        if old_photo_key:
            try:
                r2_client().delete_object(Bucket=os.environ["R2_BUCKET_NAME"], Key=old_photo_key)
            except Exception:
                pass

        flash(
            f"Usuário {user.name} arquivado. O acesso, foto e contatos foram removidos; "
            f"o nome permanece para rastreabilidade histórica."
        )
        return redirect(url_for("users_page"))

    for profile in schedule_profiles:
        db.session.delete(profile)
    db.session.delete(user)
    try:
        db.session.commit()
    except IntegrityError:
        # Há referência histórica não contemplada pelo contador resumido: preservar rastreabilidade e arquivar.
        db.session.rollback()
        user = db.session.get(User, user_id)
        if user:
            for profile in TeamScheduleProfile.query.filter_by(user_id=user.id).all():
                profile.active=False
                profile.user_id=None
                profile.updated_at=datetime.utcnow()
            user.active=False
            user.archived_at=datetime.utcnow()
            user.username=f"arquivado-{user.id}-{uuid.uuid4().hex[:10]}"
            user.password_hash=generate_password_hash(uuid.uuid4().hex + uuid.uuid4().hex)
            user.email=None; user.phone=None; user.photo_url=None
            db.session.commit()
            flash(f"Usuário {user.name} possui histórico vinculado e foi arquivado com segurança em vez de excluído.")
            return redirect(url_for("users_page"))
        raise

    if old_photo_key:
        try:
            r2_client().delete_object(Bucket=os.environ["R2_BUCKET_NAME"], Key=old_photo_key)
        except Exception:
            pass

    flash("Usuário sem histórico excluído definitivamente.")
    return redirect(url_for("users_page"))



@app.post("/usuarios/<int:user_id>/reativar")
@user_admin_required
def reactivate_user(user_id):
    if not _has_access("users.activate"): abort(403)
    user = db.session.get(User, user_id)
    if not user or not user.archived_at:
        flash("Usuário arquivado não encontrado.")
        return redirect(url_for("users_page"))
    if not _hr_target_allowed(user):
        flash("RH pode administrar somente perfis operacionais autorizados.")
        return redirect(url_for("users_page"))
    # RH não pode restaurar diretamente um perfil sensível.
    if user.role in ("manager", "consultation", "dispatcher", "hr") and not _current_user_is_superadmin():
        flash("Somente o Administrador principal pode reativar usuários Gestor ou Consulta.")
        return redirect(url_for("users_page"))
    base_login = normalize(user.name).lower().replace(" ", ".")[:60] or f"usuario{user.id}"
    base_login = re.sub(r"[^a-z0-9._-]+", "", base_login) or f"usuario{user.id}"
    candidate = base_login
    n=1
    while User.query.filter(User.id != user.id, func.lower(User.username)==candidate.lower()).first():
        n += 1
        candidate = f"{base_login}.{n}"
    user.username = candidate
    user.active = True
    user.archived_at = None
    user.personnel_status = "ATIVO"
    user.password_hash = generate_password_hash(secrets.token_urlsafe(12))
    db.session.commit()
    flash(f"{user.name} foi reativado. Atualize login, contato, escala e redefina a senha antes de liberar o acesso.")
    return redirect(url_for("users_page"))


@app.get("/usuarios/exportar.xlsx")
@user_admin_required
def export_users_excel():
    if not _has_access("users.export"): abort(403)
    wb = Workbook()
    ws = wb.active
    ws.title = "Usuários e acessos"
    base_headers = ["Código","Nome","Login","Perfil","Cargo","Empresa","E-mail","Celular","Status de acesso","Situação","Escala","Horário","Data referência","Trabalha/Folga","Arquivado em"]
    permission_keys=[]
    for group,(group_label,children) in ACCESS_GROUPS.items():
        for perm in children:
            if perm not in permission_keys: permission_keys.append(perm)
    permission_headers=[f"{ACCESS_GROUPS.get(p.split('.',1)[0],(p.split('.',1)[0],[]))[0]} > {ACCESS_LABELS.get(p,p)}" for p in permission_keys]
    ws.append(base_headers + permission_headers)
    for c in ws[1]:
        c.font = Font(bold=True, color="FFFFFF")
        c.fill = PatternFill("solid", fgColor="17345D")
        c.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
    role_label={"manager":"Gestor","technician":"Técnico de Campo","technician_implantation":"Técnico Implantação","manager_field":"Gestor Field","consultation":"Consulta","hr":"RH","dispatcher":"Dispatcher","atm_financial_admin":"ADM Financeiro","customer":"Cliente"}
    q=(User.query.filter(User.role.in_(("technician","technician_implantation"))) if session.get("role") == "hr" else User.query)
    term=(request.args.get("q") or "").strip(); role_f=(request.args.get("role") or "").strip(); company_f=(request.args.get("company") or "").strip(); status_f=(request.args.get("status") or "").strip().upper()
    if term:
        like=f"%{term}%"; q=q.filter(db.or_(User.name.ilike(like),User.username.ilike(like),User.user_code.ilike(like),User.job_title.ilike(like)))
    if role_f: q=q.filter(User.role==role_f)
    if company_f: q=q.filter(User.company==company_f)
    if status_f=="ATIVO": q=q.filter(User.active.is_(True),User.archived_at.is_(None))
    elif status_f=="INATIVO": q=q.filter(User.active.is_(False),User.archived_at.is_(None))
    elif status_f=="ARQUIVADO": q=q.filter(User.archived_at.isnot(None))
    exported=[]
    for u in q.order_by(User.name).all():
        exported.append(u)
        status = "ARQUIVADO" if u.archived_at else ("ATIVO" if u.active else "INATIVO")
        profile=db.session.get(SystemProfile,u.system_profile_id) if getattr(u,'system_profile_id',None) else None
        profile_name=profile.name if profile else role_label.get(u.role,u.role)
        effective=_user_access_set(u)
        perm_values=[]
        for perm in permission_keys:
            group=perm.split('.',1)[0]
            allowed=(u.role=='manager') or (perm in effective) or (group in effective)
            perm_values.append("SIM" if allowed else "NÃO")
        ws.append([
            u.user_code or "",u.name,u.username,profile_name,u.job_title or "",u.company or "",u.email or "",u.phone or "",
            status,u.personnel_status or "",u.work_schedule_type or "",u.work_shift or "",
            u.work_anchor_date.isoformat() if u.work_anchor_date else "",u.work_anchor_status or "",
            u.archived_at.strftime("%d/%m/%Y %H:%M") if u.archived_at else "",
        ] + perm_values)
    ws.auto_filter.ref=ws.dimensions
    ws.freeze_panes="A2"
    ws2=wb.create_sheet("Legenda de Permissões")
    ws2.append(["Área","Subcategoria / Permissão","Código técnico","Descrição"])
    for c in ws2[1]:
        c.font=Font(bold=True,color="FFFFFF"); c.fill=PatternFill("solid",fgColor="17345D")
    for perm in permission_keys:
        group=perm.split('.',1)[0]; group_label=ACCESS_GROUPS.get(group,(group,[]))[0]; label=ACCESS_LABELS.get(perm,perm)
        ws2.append([group_label,label,perm,f"SIM = usuário possui acesso efetivo a {label}; NÃO = acesso não concedido."])
    ws2.freeze_panes="A2"; ws2.auto_filter.ref=ws2.dimensions
    ws3=wb.create_sheet("Resumo")
    ws3.append(["Indicador","Quantidade"]); ws3.append(["Usuários exportados",len(exported)])
    ws3.append(["Usuários ativos",sum(1 for u in exported if u.active and not u.archived_at)])
    ws3.append(["Usuários inativos",sum(1 for u in exported if not u.active and not u.archived_at)])
    ws3.append(["Usuários arquivados",sum(1 for u in exported if u.archived_at)])
    for c in ws3[1]: c.font=Font(bold=True,color="FFFFFF"); c.fill=PatternFill("solid",fgColor="17345D")
    for sheet in wb.worksheets:
        for col in range(1,sheet.max_column+1):
            letter=get_column_letter(col); width=12
            for row in range(1,min(sheet.max_row,500)+1):
                v=sheet.cell(row=row,column=col).value
                if v is not None: width=max(width,min(38,len(str(v))+2))
            sheet.column_dimensions[letter].width=width
        sheet.freeze_panes="A2"
    out=io.BytesIO(); wb.save(out); out.seek(0)
    return send_file(out,as_attachment=True,download_name=f"usuarios_acessos_{datetime.now().strftime('%Y%m%d_%H%M')}.xlsx",mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")


@app.post("/usuarios/<int:user_id>/senha")
@manager_required
def reset_user_password(user_id):
    if not _has_access("users.password"): abort(403)
    user = db.session.get(User, user_id)
    if not user:
        flash("Usuário não encontrado.")
        return redirect(url_for("users_page"))

    password = request.form.get("password", "")
    if len(password) < 8:
        flash("A nova senha deve ter pelo menos 8 caracteres.")
        return redirect(url_for("users_page"))

    user.password_hash = generate_password_hash(password)
    db.session.commit()
    flash(f"Senha de {user.name} alterada.")
    return redirect(url_for("users_page"))



def _wa_parse_messages(chat_text):
    pattern = re.compile(
        r'^\[(\d{2}/\d{2}/\d{4}),\s*(\d{2}:\d{2}:\d{2})\]\s*([^:]+):\s?(.*)$'
    )
    messages = []
    current = None
    for raw in chat_text.splitlines():
        line = (
            raw.replace("\u200e", "")
               .replace("\ufeff", "")
               .replace("\u202a", "")
               .replace("\u202c", "")
        )
        m = pattern.match(line)
        if m:
            if current:
                messages.append(current)
            current = {
                "date": m.group(1),
                "time": m.group(2),
                "author": m.group(3).strip(),
                "text": m.group(4).strip()
            }
        elif current is not None:
            current["text"] += "\n" + line
    if current:
        messages.append(current)
    return messages


def _wa_extract_attachments(text):
    return re.findall(r'<anexado:\s*([^>]+)>', text or "", flags=re.I)


_WA_LINE_COLORS = {
    1: "AZUL", 2: "VERDE", 3: "VERMELHA", 4: "AMARELA", 5: "LILAS",
    7: "RUBI", 8: "DIAMANTE", 9: "ESMERALDA", 10: "TURQUESA",
    11: "CORAL", 12: "SAFIRA", 13: "JADE", 15: "PRATA", 17: "OURO",
}


def _wa_extract_station(text):
    clean = re.sub(r'<anexado:[^>]+>', '', text or '', flags=re.I)
    m = re.search(
        r'Est[aã]ção\s+(.+?)\s+Linha\s+(\d{1,2})(?:\s*[-–]\s*([A-Za-zÀ-ÿ]+))?',
        clean, flags=re.I | re.S
    )
    if not m:
        return None
    station = re.sub(r'\s+', ' ', m.group(1)).strip(" -")
    number = int(m.group(2))
    color = (m.group(3) or _WA_LINE_COLORS.get(number, "")).upper()
    tail = clean[m.end():].splitlines()[0].strip() if m.end() < len(clean) else ""
    if normalize(station) in {"MOTIVA", "METRO", "CPTM"} and tail:
        station = tail.split("*", 1)[0].strip(" -")
    return {
        "station": station,
        "line_number": number,
        "line": f"{number:02d} - {color}" if color else f"{number:02d}",
    }


def _wa_split_competition(text):
    m = re.search(r'\*?\s*Concorr[eê]ncia\s*\*?', text or "", flags=re.I)
    if not m:
        m = re.search(r'\*?\s*concorrencia\s*\*?', text or "", flags=re.I)
    if not m:
        return text or "", ""
    return (text or "")[:m.start()], (text or "")[m.end():]


def _wa_equipment_rows(text, visit_token="WA"):
    official, competition = _wa_split_competition(text)
    clean = re.sub(r'<anexado:[^>]+>', '', official or '', flags=re.I)
    rows = []
    provisional_index = 0

    for raw_line in clean.splitlines():
        line = re.sub(r'\s+', ' ', raw_line).strip()
        if not line:
            continue

        # POS: TOPs de 6-8 dígitos. SN e patrimônio são atributos, não identificadores.
        if re.search(r'\bPOS\b', line, flags=re.I):
            tmp = re.sub(r'\bSN\s*[A-Za-z0-9]+\b', ' ', line, flags=re.I)
            tmp = re.sub(r'Patrim[oô]nio\s*(?:off|[A-Za-z0-9]+)', ' ', tmp, flags=re.I)
            for ident in re.findall(r'\b\d{6,8}\b', tmp):
                rows.append({
                    "type": "POS", "identifier": ident, "model": "",
                    "serial": "", "patrimony": "", "status": "Não informado",
                    "source_line": line, "provisional": False
                })

        # ATM inclui equipamentos descritos em campo como TCI / TCI NEO / MK / MK NEO.
        if re.search(r'\b(?:ATM|TCI)\b', line, flags=re.I):
            tmp = re.sub(r'\(?\s*Patrim[oô]nio\s*(?:off|\d+)\s*\)?', ' ', line, flags=re.I)
            tmp = re.sub(r'\bSN\s*[A-Za-z0-9]+\b', ' ', tmp, flags=re.I)
            ids = re.findall(r'\b\d{4,8}\b', tmp)
            model_match = re.search(r'\b(?:ATM|TCI)\s*([^:]*?)(?::|$)', line, flags=re.I)
            model = model_match.group(1).strip() if model_match else ""
            if ids:
                for ident in ids:
                    rows.append({
                        "type": "ATM", "identifier": ident, "model": model,
                        "serial": "", "patrimony": "",
                        "status": "Inoperante" if re.search(r'inoperante', line, re.I) else "Não informado",
                        "source_line": line, "provisional": False
                    })
            elif re.search(r'inoperante', line, re.I):
                provisional_index += 1
                rows.append({
                    "type": "ATM",
                    "identifier": f"ATM-SID-{visit_token}-{provisional_index}",
                    "model": model, "serial": "", "patrimony": "",
                    "status": "Inoperante", "source_line": line, "provisional": True
                })

        # Validadores de recarga.
        if re.search(r'\bVALIDADOR', line, flags=re.I):
            nums = re.findall(r'\b\d{2,5}\b', line)
            for ident in [x for x in nums if int(x) > 20]:
                rows.append({
                    "type": "VALIDADOR", "identifier": ident, "model": "",
                    "serial": "", "patrimony": "", "status": "Não informado",
                    "source_line": line, "provisional": False
                })

    # POS em linhas detalhadas: TOP SN patrimônio.
    for m in re.finditer(
        r'\b(\d{6,8})\s+SN\s+([A-Za-z0-9]+)(?:\s+Patrim[oô]nio\s+([A-Za-z0-9]+))?',
        clean, flags=re.I
    ):
        ident, serial, patrimony = m.group(1), m.group(2), (m.group(3) or "")
        existing = next((x for x in rows if x["type"] == "POS" and x["identifier"] == ident), None)
        if existing:
            existing["serial"] = serial
            existing["patrimony"] = patrimony
        else:
            rows.append({
                "type": "POS", "identifier": ident, "model": "",
                "serial": serial, "patrimony": patrimony, "status": "Não informado",
                "source_line": m.group(0), "provisional": False
            })

    # Rack/Hack é guardado como evidência, sem contaminar os Big Numbers oficiais.
    rack_mentions = re.findall(r'\b(?:RACK|HACK)\b', clean, flags=re.I)
    for idx in range(len(rack_mentions)):
        rows.append({
            "type": "RACK", "identifier": f"RACK-{visit_token}-{idx + 1}",
            "model": "", "serial": "", "patrimony": "",
            "status": "Não informado", "source_line": "Rack citado no relatório",
            "provisional": True
        })

    seen = set()
    out = []
    for row in rows:
        key = (row["type"], row["identifier"])
        if key not in seen:
            seen.add(key)
            out.append(row)

    competition_summary = {}
    for type_name, regex in (
        ("ATM", r'(\d+)\s+ATMs?'),
        ("VALIDADOR", r'(\d+)\s+validadores?'),
        ("TERMINAL", r'(\d+)\s+terminais?'),
    ):
        vals = [int(x) for x in re.findall(regex, competition or "", flags=re.I)]
        if vals:
            competition_summary[type_name] = sum(vals)

    return out, (competition or "").strip(), competition_summary


def _wa_location_name(loc):
    if not loc:
        return ""
    name = normalize(loc.location)
    return name.split(" - ", 1)[1] if " - " in name else name


def _wa_match_location(station_name, line_number):
    aliases = {
        "BRAZ CUBAS": "BRAS CUBAS",
        "ENG MANOEL FEIO": "MANOEL FEIO",
        "ENG GOULART": "ENGENHEIRO GOULART",
        "GUARULHOS": "GUARULHOS CECAP",
        "AEROPORTO": "AEROPORTO GUARULHOS",
    }
    target = aliases.get(normalize(station_name), normalize(station_name))

    all_locations = Location.query.all()
    same_line = []
    other_line = []
    for loc in all_locations:
        m = re.match(r'(\d+)', loc.line or "")
        ln = int(m.group(1)) if m else None
        name = _wa_location_name(loc)
        score = SequenceMatcher(None, target, name).ratio()
        if target in name or name in target:
            score = max(score, 0.93)
        bucket = same_line if ln == line_number else other_line
        bucket.append((score, loc))

    same_line.sort(key=lambda x: x[0], reverse=True)
    other_line.sort(key=lambda x: x[0], reverse=True)

    if same_line and same_line[0][0] >= 0.72:
        score, loc = same_line[0]
        return loc, ("SEGURA" if score >= 0.82 else "REVISAR"), round(score, 3)

    # Estações compartilhadas podem aparecer no WhatsApp com uma linha que ainda não existe
    # como Location separada (ex.: Luz / Brás). Mantém como revisão, sem atribuição automática.
    if other_line and other_line[0][0] >= 0.90:
        score, _loc = other_line[0]
        return None, "REVISAR", round(score, 3)

    return None, "NAO IDENTIFICADA", 0.0


def _wa_visit_source_key(msg, station):
    seed = "|".join([
        msg["date"], msg["time"], normalize(msg["author"]),
        normalize(station["station"]), str(station["line_number"]),
        normalize(re.sub(r'<anexado:[^>]+>', '', msg["text"], flags=re.I))
    ])
    return hashlib.sha256(seed.encode("utf-8")).hexdigest()[:40]


def _wa_analyze_archive(raw):
    source = raw if isinstance(raw, (str, Path)) else io.BytesIO(raw)
    with zipfile.ZipFile(source) as z:
        names = z.namelist()
        chat_names = [n for n in names if n.lower().endswith("_chat.txt") or n.lower().endswith(".txt")]
        if not chat_names:
            raise ValueError("O ZIP não contém o arquivo _chat.txt.")

        chat_text = z.read(chat_names[0]).decode("utf-8", errors="replace")
        messages = _wa_parse_messages(chat_text)
        available_media = set(names)

        visits = []
        active_by_author = {}
        pending_media = {}
        seen_visit_keys = set()

        def msg_dt(msg):
            try:
                return datetime.strptime(f'{msg["date"]} {msg["time"]}', "%d/%m/%Y %H:%M:%S")
            except Exception:
                return datetime.utcnow()

        for msg in messages:
            author = msg["author"]
            dt = msg_dt(msg)
            station = _wa_extract_station(msg["text"])
            attachments = [a for a in _wa_extract_attachments(msg["text"]) if a in available_media]

            if station:
                source_key = _wa_visit_source_key(msg, station)
                active = active_by_author.get(author)

                # Concorrência/continuação da mesma estação é incorporada à mesma visita.
                if (
                    active
                    and normalize(active["station_raw"]) == normalize(station["station"])
                    and active["line_number"] == station["line_number"]
                    and (dt - active["_last_dt"]).total_seconds() <= 5400
                ):
                    signature = normalize(re.sub(r'<anexado:[^>]+>', '', msg["text"], flags=re.I))
                    if signature and signature not in active["_signatures"]:
                        active["_signatures"].add(signature)
                        active["messages"].append(msg)
                        eqs, comp_text, comp_summary = _wa_equipment_rows(msg["text"], active["source_key"][:10])
                        existing_keys = {(x["type"], x["identifier"]) for x in active["equipment"]}
                        for eq in eqs:
                            if (eq["type"], eq["identifier"]) not in existing_keys:
                                active["equipment"].append(eq)
                                existing_keys.add((eq["type"], eq["identifier"]))
                        if comp_text:
                            active["competition_text"] = (active["competition_text"] + "\n" + comp_text).strip()
                        for k, v in comp_summary.items():
                            active["competition_summary"][k] = active["competition_summary"].get(k, 0) + v
                    for name in attachments:
                        if name not in active["attachments"]:
                            active["attachments"].append(name)
                    active["_last_dt"] = dt
                    continue

                if source_key in seen_visit_keys:
                    continue
                seen_visit_keys.add(source_key)

                loc, confidence, score = _wa_match_location(station["station"], station["line_number"])
                eqs, comp_text, comp_summary = _wa_equipment_rows(msg["text"], source_key[:10])
                visit = {
                    "source_key": source_key,
                    "date": msg["date"],
                    "time": msg["time"],
                    "author": author,
                    "station_raw": station["station"],
                    "line_number": station["line_number"],
                    "line_raw": station["line"],
                    "location_id": loc.id if loc else None,
                    "location_name": loc.location if loc else "",
                    "company": loc.company if loc else "",
                    "confidence": confidence,
                    "match_score": score,
                    "messages": [msg],
                    "equipment": eqs,
                    "competition_text": comp_text,
                    "competition_summary": comp_summary,
                    "attachments": list(dict.fromkeys(attachments)),
                    "_last_dt": dt,
                    "_signatures": {
                        normalize(re.sub(r'<anexado:[^>]+>', '', msg["text"], flags=re.I))
                    },
                }

                # Fotos enviadas pouco antes do texto da estação.
                pending = pending_media.get(author, [])
                keep = []
                for pdt, name in pending:
                    if 0 <= (dt - pdt).total_seconds() <= 900:
                        if name not in visit["attachments"]:
                            visit["attachments"].append(name)
                    else:
                        keep.append((pdt, name))
                pending_media[author] = keep

                visits.append(visit)
                active_by_author[author] = visit

            else:
                active = active_by_author.get(author)

                if attachments:
                    if active and 0 <= (dt - active["_last_dt"]).total_seconds() <= 5400:
                        for name in attachments:
                            if name not in active["attachments"]:
                                active["attachments"].append(name)
                        active["_last_dt"] = dt
                    else:
                        pending_media.setdefault(author, [])
                        for name in attachments:
                            pending_media[author].append((dt, name))

                # POS e complementos enviados em mensagem seguinte pertencem à última estação
                # do mesmo técnico, desde que dentro de 90 minutos.
                if active and 0 <= (dt - active["_last_dt"]).total_seconds() <= 5400:
                    clean = re.sub(r'<anexado:[^>]+>', '', msg["text"]).strip()
                    if clean and clean != "Mensagem apagada":
                        eqs, comp_text, comp_summary = _wa_equipment_rows(msg["text"], active["source_key"][:10])
                        if eqs or re.search(r'\bPOS\b|\bSN\b|Patrim|concorr|autoriz', clean, flags=re.I):
                            signature = normalize(clean)
                            if signature and signature not in active["_signatures"]:
                                active["_signatures"].add(signature)
                                active["messages"].append(msg)
                            existing_keys = {(x["type"], x["identifier"]) for x in active["equipment"]}
                            for eq in eqs:
                                if (eq["type"], eq["identifier"]) not in existing_keys:
                                    active["equipment"].append(eq)
                                    existing_keys.add((eq["type"], eq["identifier"]))
                            if comp_text:
                                active["competition_text"] = (active["competition_text"] + "\n" + comp_text).strip()
                            for k, v in comp_summary.items():
                                active["competition_summary"][k] = active["competition_summary"].get(k, 0) + v
                            active["_last_dt"] = dt

        for visit in visits:
            visit.pop("_last_dt", None)
            visit.pop("_signatures", None)

            # Comparação preliminar com a base e inventário atual.
            for eq in visit["equipment"]:
                eq["base_status"] = "NÃO CONFRONTADO"
                eq["inventory_status"] = "NÃO CONFRONTADO"
                eq["audit_status"] = "PENDENTE"
                eq["base_asset_id"] = None
                eq["inventory_id"] = None

                if eq["type"] == "RACK":
                    eq["audit_status"] = "EVIDÊNCIA FORA DO PARQUE OFICIAL"
                    continue

                type_map = {
                    "ATM": "ATM",
                    "VALIDADOR": "VALIDADOR",
                    "POS": "POS",
                }
                base_type = type_map.get(eq["type"], eq["type"])
                candidates = BaseAsset.query.filter(
                    func.upper(BaseAsset.equipment_type) == base_type,
                    (
                        (func.upper(func.coalesce(BaseAsset.terminal_number, "")) == eq["identifier"].upper())
                        | (func.upper(func.coalesce(BaseAsset.top_id, "")) == eq["identifier"].upper())
                        | (func.upper(func.coalesce(BaseAsset.serial, "")) == eq["identifier"].upper())
                    )
                ).all()

                if candidates:
                    chosen = candidates[0]
                    eq["base_asset_id"] = chosen.id
                    if visit["location_id"]:
                        loc = db.session.get(Location, visit["location_id"])
                        same_line = normalize(chosen.line) == normalize(loc.line)
                        same_name = (
                            normalize(chosen.locality) in normalize(loc.location)
                            or normalize(_wa_location_name(loc)) in normalize(chosen.locality)
                        )
                        if same_line and same_name:
                            eq["base_status"] = "BASE CONFERE"
                        else:
                            eq["base_status"] = "BASE EM OUTRA LOCALIDADE"
                    else:
                        eq["base_status"] = "ENCONTRADO NA BASE"
                else:
                    eq["base_status"] = "NÃO PREVISTO NA BASE"

                if visit["location_id"]:
                    inv_type = {
                        "ATM": "ATM",
                        "VALIDADOR": "Validador de Recarga",
                        "POS": "POS de Bilheteria",
                    }.get(eq["type"], eq["type"])
                    inv = Inventory.query.filter(
                        Inventory.location_id == visit["location_id"],
                        Inventory.equipment_type == inv_type,
                        func.upper(Inventory.asset_identifier) == eq["identifier"].upper()
                    ).first()
                    if inv:
                        eq["inventory_id"] = inv.id
                        eq["inventory_status"] = "JÁ INVENTARIADO"
                    else:
                        eq["inventory_status"] = "AINDA NÃO INVENTARIADO"

                if eq["inventory_id"]:
                    eq["audit_status"] = "CONFORME / JÁ INVENTARIADO"
                elif eq["base_status"] == "BASE CONFERE":
                    eq["audit_status"] = "CONFIRMADO EM CAMPO / FALTA PROMOVER"
                elif eq["base_status"] == "BASE EM OUTRA LOCALIDADE":
                    eq["audit_status"] = "DIVERGÊNCIA DE LOCALIDADE"
                elif eq["base_status"] == "NÃO PREVISTO NA BASE":
                    eq["audit_status"] = "NOVO / NÃO PREVISTO"
                else:
                    eq["audit_status"] = "PENDENTE DE REVISÃO"

        summary = {
            "messages": len(messages),
            "media_total": len([n for n in names if not n.lower().endswith(".txt")]),
            "visits": len(visits),
            "equipment": sum(len(v["equipment"]) for v in visits),
            "safe": sum(1 for v in visits if v["confidence"] == "SEGURA"),
            "review": sum(1 for v in visits if v["confidence"] == "REVISAR"),
            "unmatched": sum(1 for v in visits if v["confidence"] == "NAO IDENTIFICADA"),
            "duplicates": sum(
                1 for v in visits for e in v["equipment"]
                if e.get("inventory_status") == "JÁ INVENTARIADO"
            ),
            "by_type": {
                t: sum(1 for v in visits for e in v["equipment"] if e["type"] == t)
                for t in ("ATM", "VALIDADOR", "POS", "RACK")
            },
            "audit": {
                status: sum(1 for v in visits for e in v["equipment"] if e["audit_status"] == status)
                for status in sorted({
                    e["audit_status"] for v in visits for e in v["equipment"]
                })
            },
            "media_linked": len({a for v in visits for a in v["attachments"]}),
        }
        return visits, summary


def _r2_put_bytes(key, data, content_type=None):
    kwargs = {
        "Bucket": os.environ["R2_BUCKET_NAME"],
        "Key": key,
        "Body": data
    }
    if content_type:
        kwargs["ContentType"] = content_type
    r2_client().put_object(**kwargs)


# V62 REV2 — uploads de evidências com memória limitada.
# Nunca materializa a foto inteira em bytes antes de enviá-la ao R2.
_ACTIVITY_UPLOAD_MAX_MB = max(1, int(os.environ.get("ACTIVITY_UPLOAD_MAX_MB", "15") or 15))
_ACTIVITY_REQUEST_MAX_MB = max(_ACTIVITY_UPLOAD_MAX_MB, int(os.environ.get("ACTIVITY_REQUEST_MAX_MB", "40") or 40))

def _uploaded_file_size(upload):
    try:
        stream = upload.stream
        pos = stream.tell()
        stream.seek(0, 2)
        size = int(stream.tell() or 0)
        stream.seek(pos)
        return size
    except Exception:
        return 0

def _activity_request_too_large():
    size = int(request.content_length or 0)
    return bool(size and size > _ACTIVITY_REQUEST_MAX_MB * 1024 * 1024)

def _store_uploaded_file(upload, folder, stored_name, content_type=None, max_mb=None):
    """Persiste FileStorage sem f.read(), evitando duplicar imagens grandes na RAM.
    Retorna o stored_name local ou r2__<key>.
    """
    if not upload or not getattr(upload, "filename", None):
        raise ValueError("Arquivo inválido.")
    limit_mb = int(max_mb or _ACTIVITY_UPLOAD_MAX_MB)
    size = _uploaded_file_size(upload)
    if size and size > limit_mb * 1024 * 1024:
        raise ValueError(f"O arquivo {upload.filename} excede {limit_mb} MB. Reduza a resolução e tente novamente.")
    mime = content_type or upload.mimetype or "application/octet-stream"
    try:
        upload.stream.seek(0)
    except Exception:
        pass
    if _r2_available():
        key = f"{folder}/{datetime.utcnow().strftime('%Y/%m')}/{stored_name}"
        kwargs = {
            "Bucket": os.environ["R2_BUCKET_NAME"],
            "Key": key,
            "Body": upload.stream,
            "ContentType": mime,
        }
        if size:
            kwargs["ContentLength"] = size
        r2_client().put_object(**kwargs)
        return "r2__" + key
    upload.save(UPLOAD_DIR / stored_name)
    return stored_name


def _r2_get_bytes(key):
    obj = r2_client().get_object(Bucket=os.environ["R2_BUCKET_NAME"], Key=key)
    return obj["Body"].read()


def _r2_available():
    required = ("R2_ENDPOINT_URL", "R2_ACCESS_KEY_ID", "R2_SECRET_ACCESS_KEY", "R2_BUCKET_NAME")
    return all(os.environ.get(x, "").strip() for x in required)


def _wa_stage_archive(raw, filename):
    # V34.4: o ZIP-fonte fica temporariamente no disco do worker durante a importação.
    # As mídias definitivas continuam no Cloudflare R2. Isso evita centenas de downloads
    # repetidos do ZIP e permite processar a carga em pequenos lotes HTTP.
    batch_id = uuid.uuid4().hex
    safe_name = secure_filename(filename) or "whatsapp.zip"
    staging_dir = UPLOAD_DIR / "whatsapp_staging"
    staging_dir.mkdir(parents=True, exist_ok=True)
    path = staging_dir / f"{batch_id}_{safe_name}"
    path.write_bytes(raw)
    return f"local:{path.name}", batch_id, "R2 (mídias) · staging local"


def _wa_load_staged(staging_key):
    if staging_key.startswith("r2:"):
        return _r2_get_bytes(staging_key[3:])
    if staging_key.startswith("local:"):
        path = UPLOAD_DIR / "whatsapp_staging" / staging_key[6:]
        return path.read_bytes()
    raise ValueError("Origem de importação inválida.")

def _wa_materialize_staged(staging_key):
    """Materializa o ZIP em disco temporário sem mantê-lo inteiro na RAM."""
    tmp = tempfile.NamedTemporaryFile(prefix="wa-import-", suffix=".zip", delete=False)
    tmp_path = Path(tmp.name)
    try:
        if staging_key.startswith("r2:"):
            obj = r2_client().get_object(Bucket=os.environ["R2_BUCKET_NAME"], Key=staging_key[3:])
            body = obj["Body"]
            while True:
                chunk = body.read(1024 * 1024)
                if not chunk:
                    break
                tmp.write(chunk)
        elif staging_key.startswith("local:"):
            src = UPLOAD_DIR / "whatsapp_staging" / staging_key[6:]
            with src.open("rb") as fh:
                shutil.copyfileobj(fh, tmp, length=1024 * 1024)
        else:
            raise ValueError("Origem de importação inválida.")
        tmp.flush(); tmp.close()
        return tmp_path
    except Exception:
        try: tmp.close()
        except Exception: pass
        tmp_path.unlink(missing_ok=True)
        raise

def _sha256_file(path):
    h = hashlib.sha256()
    with Path(path).open("rb") as fh:
        for chunk in iter(lambda: fh.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def _evidence_store_media(data, original_name, batch_id):
    """Persiste mídia de evidência. V22 também repara registros antigos locais quando o R2 está disponível."""
    sha = hashlib.sha256(data).hexdigest()
    safe_name = secure_filename(Path(original_name).name) or f"midia-{sha[:12]}"
    mime = mimetypes.guess_type(safe_name)[0] or "application/octet-stream"
    existing = FieldEvidenceMedia.query.filter_by(sha256=sha).first()

    if existing:
        if _r2_available():
            needs_repair = existing.storage_kind != "r2"
            if existing.storage_kind == "r2":
                try:
                    # V34.1: HEAD evita carregar a mídia existente inteira em memória apenas para validar sua existência.
                    r2_client().head_object(Bucket=os.environ["R2_BUCKET_NAME"], Key=existing.storage_key)
                except Exception:
                    needs_repair = True
            if needs_repair:
                key = existing.storage_key if existing.storage_kind == "r2" and existing.storage_key else f"whatsapp/evidencias/reparados/{sha[:12]}_{safe_name}"
                _r2_put_bytes(key, data, mime)
                existing.storage_kind = "r2"
                existing.storage_key = key
                existing.mime_type = mime
                existing.original_name = safe_name
                return existing, False, True
        return existing, False, False

    if _r2_available():
        key = f"whatsapp/evidencias/{batch_id}/{sha[:12]}_{safe_name}"
        _r2_put_bytes(key, data, mime)
        return FieldEvidenceMedia(
            sha256=sha, original_name=safe_name, mime_type=mime,
            storage_kind="r2", storage_key=key
        ), True, False

    evidence_dir = UPLOAD_DIR / "field_evidence"
    evidence_dir.mkdir(parents=True, exist_ok=True)
    stored = f"{sha[:16]}_{safe_name}"
    (evidence_dir / stored).write_bytes(data)
    return FieldEvidenceMedia(
        sha256=sha, original_name=safe_name, mime_type=mime,
        storage_kind="local", storage_key=stored
    ), True, False


def _evidence_summary():
    visits = FieldEvidenceVisit.query.count()
    items = FieldEvidenceItem.query.count()
    media = FieldEvidenceMedia.query.count()
    matched = FieldEvidenceItem.query.filter(
        FieldEvidenceItem.audit_status.in_((
            "CONFORME / JÁ INVENTARIADO",
            "CONFIRMADO EM CAMPO / FALTA PROMOVER",
        ))
    ).count()
    review = FieldEvidenceItem.query.filter(
        ~FieldEvidenceItem.audit_status.in_((
            "CONFORME / JÁ INVENTARIADO",
            "CONFIRMADO EM CAMPO / FALTA PROMOVER",
            "EVIDÊNCIA FORA DO PARQUE OFICIAL",
        ))
    ).count()
    competition = {}
    for row in FieldEvidenceVisit.query.with_entities(FieldEvidenceVisit.competition_text).all():
        text = (row[0] or "").strip()
        if not text:
            continue
        for label, regex in (("ATM", r"ATM\s*[:=-]?\s*(\d+)"), ("POS", r"POS\s*[:=-]?\s*(\d+)"), ("Recarga", r"(?:RECARGA|VALIDADOR)\s*[:=-]?\s*(\d+)")):
            vals = [int(x) for x in re.findall(regex, text, flags=re.I)]
            if vals: competition[label] = competition.get(label, 0) + sum(vals)
    return {
        "visits": visits,
        "items": items,
        "media": media,
        "matched": matched,
        "review": review,
        "competition": competition,
        "competition_total": sum(competition.values()),
        "unresolved_visits": FieldEvidenceVisit.query.filter(FieldEvidenceVisit.location_id.is_(None)).count(),
    }


@app.get("/r2-status")
@manager_required
def r2_status_v22():
    configured = _r2_available()
    payload = {
        "ok": configured,
        "release": APP_RELEASE,
        "storage": "Cloudflare R2",
        "bucket": os.environ.get("R2_BUCKET_NAME", ""),
        "endpoint_configured": bool(os.environ.get("R2_ENDPOINT_URL", "").strip()),
        "access_key_configured": bool(os.environ.get("R2_ACCESS_KEY_ID", "").strip()),
        "secret_key_configured": bool(os.environ.get("R2_SECRET_ACCESS_KEY", "").strip()),
    }
    if request.args.get("test") == "1" and configured:
        ok, message = r2_test_connection()
        payload["ok"] = ok
        payload["connection_test"] = message
    elif not configured:
        payload["message"] = "Configure R2_ENDPOINT_URL, R2_ACCESS_KEY_ID, R2_SECRET_ACCESS_KEY e R2_BUCKET_NAME."
    return jsonify(payload), (200 if payload["ok"] else 503)


@app.route("/importar-excel", methods=["GET", "POST"])
@manager_required
def import_excel_v34():
    from openpyxl import load_workbook
    preview=[]; error=None; filename=None; staging_key=None; import_result=None
    aliases={
      "location":["LOCALIDADE","ESTACAO","ESTAÇÃO","LOCAL","STATION"],
      "line":["LINHA","LINE"], "company":["EMPRESA","OPERADORA","COMPANY"],
      "type":["TIPO","EQUIPAMENTO","TIPO EQUIPAMENTO","TIPO DE EQUIPAMENTO"],
      "identifier":["IDENTIFICACAO","IDENTIFICAÇÃO","PATRIMONIO","PATRIMÔNIO","ATIVO","QR CODE","QRCODE"],
      "serial":["SERIE","SÉRIE","NUMERO DE SERIE","NÚMERO DE SÉRIE"],
      "supplier":["FABRICANTE","FORNECEDOR","FABRICANTE / FORNECEDOR"], "model":["MODELO"],
      "status":["STATUS","STATUS OPERACIONAL","SITUACAO","SITUAÇÃO"], "notes":["OBSERVACOES","OBSERVAÇÕES","OBS","NOTAS"]}
    def norm(v): return normalize(str(v or '')).replace('_',' ')
    def mapping(headers):
      out={}; nh=[norm(x) for x in headers]
      for key,vals in aliases.items():
        for i,h in enumerate(nh):
          if h in [norm(v) for v in vals] or any(norm(v) in h for v in vals if len(norm(v))>4): out[key]=i; break
      return out
    def analyze(raw):
      wb=load_workbook(io.BytesIO(raw),read_only=True,data_only=True)
      result=[]
      for ws in wb.worksheets:
        rows=ws.iter_rows(values_only=True); headers=list(next(rows,[]) or [])
        mp=mapping(headers); sample=[]; count=0; ready=0
        for row in rows:
          if not any(v not in (None,'') for v in row): continue
          count+=1
          if len(sample)<5: sample.append([str(v or '')[:90] for v in row])
          if mp.get('location') is not None and mp.get('identifier') is not None: ready+=1
        result.append({'sheet':ws.title,'rows':count,'headers':[str(x or '') for x in headers], 'sample':sample,'mapping':mp,'ready':ready})
      return result
    if request.method=='POST':
      action=request.form.get('action','analyze')
      try:
        if action=='analyze':
          up=request.files.get('excel_file')
          if not up or not up.filename.lower().endswith(('.xlsx','.xlsm')): raise ValueError('Selecione um arquivo .xlsx ou .xlsm.')
          raw=up.read(); filename=secure_filename(up.filename) or up.filename
          preview=analyze(raw); staging_key,_bid,_store=_wa_stage_archive(raw,filename)
        elif action=='import':
          staging_key=request.form.get('staging_key',''); raw=_wa_load_staged(staging_key); preview=analyze(raw)
          wb=load_workbook(io.BytesIO(raw),read_only=True,data_only=True)
          created=0; skipped=0; unresolved=0
          # usuário executor: mantém auditoria no gestor logado
          tech_id=session['user_id']
          locs=Location.query.all()
          def find_loc(name,line='',company=''):
            n=norm(name); ln=norm(line); co=norm(company); best=None; score=0
            for loc in locs:
              sc=0; a=norm(loc.location)
              if n and n==a: sc+=100
              elif n and (n in a or a in n): sc+=70
              if ln and ln==norm(loc.line): sc+=20
              if co and co==norm(loc.company): sc+=10
              if sc>score: score=sc; best=loc
            return best if score>=70 else None
          for ws in wb.worksheets:
            rows=ws.iter_rows(values_only=True); headers=list(next(rows,[]) or []); mp=mapping(headers)
            if 'location' not in mp or 'identifier' not in mp: continue
            for row in rows:
              if not any(v not in (None,'') for v in row): continue
              get=lambda k: str(row[mp[k]] or '').strip() if k in mp and mp[k]<len(row) else ''
              loc=find_loc(get('location'),get('line'),get('company')); ident=get('identifier')
              if not loc or not ident: unresolved+=1; continue
              typ=_canonical_equipment_type(get('type') or 'OUTRO') or 'OUTRO'
              exists=Inventory.query.filter_by(location_id=loc.id,equipment_type=typ,asset_identifier=ident).first()
              if exists: skipped+=1; continue
              inv=Inventory(location_id=loc.id,equipment_type=typ,asset_identifier=ident,serial=get('serial'),supplier=get('supplier'),model=get('model'),operational_status=get('status') or 'Não informado',notes=(get('notes')+'\n[Importado via Excel V34.1]').strip(),technician_id=tech_id,created_at=datetime.utcnow())
              db.session.add(inv); created+=1
          db.session.commit(); import_result={'created':created,'skipped':skipped,'unresolved':unresolved}
      except Exception as exc:
        db.session.rollback(); error=f'{type(exc).__name__}: {exc}'
    return render_template('import_excel.html',preview=preview,error=error,filename=filename,staging_key=staging_key,import_result=import_result,app_release=APP_RELEASE)



@app.post("/api/importar-whatsapp/lote")
@manager_required
def import_whatsapp_batch_v344():
    """V34.4 — processa poucas visitas por chamada para não estourar RAM/timeout do Render."""
    payload = request.get_json(silent=True) or {}
    staging_key = str(payload.get("staging_key") or "").strip()
    offset = max(0, int(payload.get("offset") or 0))
    limit = min(5, max(1, int(payload.get("limit") or 3)))
    if not staging_key:
        return jsonify({"ok": False, "error": "Arquivo de origem não encontrado."}), 400
    if not _r2_available():
        return jsonify({"ok": False, "error": "Cloudflare R2 não configurado."}), 503
    temp_created = False
    try:
        if staging_key.startswith("local:"):
            zip_path = UPLOAD_DIR / "whatsapp_staging" / staging_key[6:]
            if not zip_path.exists():
                return jsonify({"ok": False, "error": "Arquivo temporário expirou. Analise o ZIP novamente."}), 410
        else:
            zip_path = _wa_materialize_staged(staging_key); temp_created = True
        preview, summary = _wa_analyze_archive(zip_path)
        total = len(preview)
        batch_id = re.sub(r'[^a-f0-9]', '', Path(staging_key.split(':',1)[-1]).name.lower())[:16] or uuid.uuid4().hex[:16]
        counters = {"inserted_visits":0,"updated_visits":0,"inserted_items":0,"skipped_items":0,"media_uploaded":0,"media_repaired":0,"media_failed":0}
        with zipfile.ZipFile(zip_path) as z:
            names=set(z.namelist())
            for visit in preview[offset:offset+limit]:
                row=FieldEvidenceVisit.query.filter_by(source_key=visit["source_key"]).first()
                report_text="\n\n".join(f'[{m["date"]} {m["time"]}] {m["author"]}: {m["text"]}' for m in visit["messages"])
                if not row:
                    row=FieldEvidenceVisit(source_key=visit["source_key"],source_batch=batch_id,source_date=visit["date"],source_time=visit["time"],author=visit["author"],station_raw=visit["station_raw"],line_raw=visit["line_raw"],location_id=visit["location_id"],match_confidence=visit["confidence"],match_score=visit["match_score"],report_text=report_text,competition_text=visit["competition_text"],storage_source="R2")
                    db.session.add(row); db.session.flush(); counters["inserted_visits"]+=1
                else:
                    row.location_id=visit["location_id"]; row.match_confidence=visit["confidence"]; row.match_score=visit["match_score"]; row.report_text=report_text; row.competition_text=visit["competition_text"]; counters["updated_visits"]+=1
                for eq in visit["equipment"]:
                    exists=FieldEvidenceItem.query.filter_by(visit_id=row.id,equipment_type=eq["type"],identifier=eq["identifier"]).first()
                    if exists: counters["skipped_items"]+=1; continue
                    db.session.add(FieldEvidenceItem(visit_id=row.id,equipment_type=eq["type"],identifier=eq["identifier"],model=eq.get("model",""),serial=eq.get("serial",""),patrimony=eq.get("patrimony",""),operational_status=eq.get("status",""),source_line=eq.get("source_line",""),base_asset_id=eq.get("base_asset_id"),inventory_id=eq.get("inventory_id"),audit_status=eq.get("audit_status","PENDENTE"),audit_detail=f'Base: {eq.get("base_status","")}. Inventário: {eq.get("inventory_status","")}.' )); counters["inserted_items"]+=1
                for media_name in visit["attachments"]:
                    if media_name not in names: continue
                    try:
                        with z.open(media_name) as fh: data=fh.read()
                        media_obj,is_new,was_repaired=_evidence_store_media(data,media_name,batch_id)
                        if not media_obj.visit_id: media_obj.visit_id=row.id
                        if is_new: db.session.add(media_obj); counters["media_uploaded"]+=1
                        elif was_repaired: counters["media_repaired"]+=1
                        del data
                    except Exception as media_exc:
                        counters["media_failed"]+=1
                        app.logger.warning("Falha mídia WhatsApp %s: %s", media_name, media_exc)
                db.session.commit(); db.session.expire_all()
        next_offset=min(total,offset+limit)
        done=next_offset>=total
        if done and staging_key.startswith("local:"):
            try: zip_path.unlink(missing_ok=True)
            except Exception: pass
        return jsonify({"ok":True,"offset":offset,"next_offset":next_offset,"total":total,"done":done,"progress":round(next_offset/max(1,total)*100,1),"counters":counters,"summary":{"visits":summary.get("visits",total),"media_total":summary.get("media_total",0)}})
    except Exception as exc:
        db.session.rollback(); app.logger.exception("Falha lote WhatsApp V34.4")
        return jsonify({"ok":False,"error":f"{type(exc).__name__}: {exc}","offset":offset}),500
    finally:
        if temp_created:
            try: zip_path.unlink(missing_ok=True)
            except Exception: pass

@app.route("/importar-whatsapp", methods=["GET", "POST"])
@manager_required
def import_whatsapp():
    preview = None
    summary = None
    error = None
    staging_key = None
    import_result = None
    analyzed_filename = None
    storage_note = "Cloudflare R2" if _r2_available() else "armazenamento local temporário"

    if request.method == "POST":
        action = request.form.get("action", "analyze")

        if action == "analyze":
            upload = request.files.get("zip_file")
            if not upload or not upload.filename.lower().endswith(".zip"):
                error = "Selecione um arquivo ZIP exportado pelo WhatsApp."
            else:
                try:
                    raw = upload.read()
                    analyzed_filename = secure_filename(upload.filename) or upload.filename
                    preview, summary = _wa_analyze_archive(raw)
                    staging_key, _batch_id, storage_note = _wa_stage_archive(raw, upload.filename)
                    summary["filename"] = analyzed_filename
                    summary["archive_sha"] = hashlib.sha256(raw).hexdigest()[:16]
                except Exception as exc:
                    error = f"Não foi possível analisar o ZIP: {type(exc).__name__}: {exc}"

        elif action == "import":
            staging_key = request.form.get("staging_key", "").strip()
            if not staging_key:
                error = "Arquivo de origem não encontrado. Analise o ZIP novamente."
            else:
                try:
                    if not _r2_available():
                        raise RuntimeError("Cloudflare R2 não está configurado corretamente. A V22 não permite importação definitiva de mídias no armazenamento temporário do Render.")
                    temp_zip = _wa_materialize_staged(staging_key)
                    try:
                        preview, summary = _wa_analyze_archive(temp_zip)
                        full_sha = _sha256_file(temp_zip)
                        batch_id = full_sha[:16]
                        analyzed_filename = Path(staging_key.split(":",1)[-1]).name if staging_key else "arquivo analisado"
                        summary["filename"] = analyzed_filename
                        summary["archive_sha"] = batch_id
                        inserted_visits = 0
                        updated_visits = 0
                        inserted_items = 0
                        skipped_items = 0
                        media_uploaded = 0
                        media_repaired = 0
                        media_failed = 0

                        with zipfile.ZipFile(temp_zip) as z:
                            names = set(z.namelist())
                            total_visits = max(1, len(preview))
                            for visit_idx, visit in enumerate(preview, start=1):
                                row = FieldEvidenceVisit.query.filter_by(source_key=visit["source_key"]).first()
                                report_text = "\n\n".join(
                                    f'[{m["date"]} {m["time"]}] {m["author"]}: {m["text"]}'
                                    for m in visit["messages"]
                                )
                                if not row:
                                    row = FieldEvidenceVisit(
                                        source_key=visit["source_key"], source_batch=batch_id,
                                        source_date=visit["date"], source_time=visit["time"], author=visit["author"],
                                        station_raw=visit["station_raw"], line_raw=visit["line_raw"],
                                        location_id=visit["location_id"], match_confidence=visit["confidence"],
                                        match_score=visit["match_score"], report_text=report_text,
                                        competition_text=visit["competition_text"], storage_source="R2"
                                    )
                                    db.session.add(row); db.session.flush(); inserted_visits += 1
                                else:
                                    row.location_id=visit["location_id"]; row.match_confidence=visit["confidence"]
                                    row.match_score=visit["match_score"]; row.report_text=report_text
                                    row.competition_text=visit["competition_text"]; updated_visits += 1

                                for eq in visit["equipment"]:
                                    existing_item = FieldEvidenceItem.query.filter_by(
                                        visit_id=row.id, equipment_type=eq["type"], identifier=eq["identifier"]
                                    ).first()
                                    if existing_item:
                                        skipped_items += 1; continue
                                    db.session.add(FieldEvidenceItem(
                                        visit_id=row.id, equipment_type=eq["type"], identifier=eq["identifier"],
                                        model=eq.get("model", ""), serial=eq.get("serial", ""), patrimony=eq.get("patrimony", ""),
                                        operational_status=eq.get("status", ""), source_line=eq.get("source_line", ""),
                                        base_asset_id=eq.get("base_asset_id"), inventory_id=eq.get("inventory_id"),
                                        audit_status=eq.get("audit_status", "PENDENTE"),
                                        audit_detail=f'Base: {eq.get("base_status", "")}. Inventário: {eq.get("inventory_status", "")}.'
                                    )); inserted_items += 1

                                for media_name in visit["attachments"]:
                                    if media_name not in names:
                                        continue
                                    try:
                                        # Uma mídia por vez: o ZIP completo permanece no disco temporário.
                                        data = z.read(media_name)
                                        media_obj, is_new, was_repaired = _evidence_store_media(data, media_name, batch_id)
                                        if not media_obj.visit_id: media_obj.visit_id = row.id
                                        if is_new: db.session.add(media_obj); media_uploaded += 1
                                        elif was_repaired: media_repaired += 1
                                        del data
                                    except Exception:
                                        media_failed += 1

                                # V34.1: commit incremental. Reduz memória e torna a reexecução segura/idempotente.
                                db.session.commit()
                                db.session.expire_all()

                        import_result = {
                            "inserted_visits": inserted_visits, "updated_visits": updated_visits,
                            "inserted_items": inserted_items, "skipped_items": skipped_items,
                            "media_uploaded": media_uploaded, "media_repaired": media_repaired,
                            "media_failed": media_failed, "storage": "R2",
                        }
                    finally:
                        temp_zip.unlink(missing_ok=True)

                except Exception as exc:
                    db.session.rollback()
                    error = f"Erro durante a importação: {type(exc).__name__}: {exc}"

    return render_template(
        "import_whatsapp.html",
        preview=preview,
        summary=summary,
        error=error,
        staging_key=staging_key,
        import_result=import_result,
        storage_note=storage_note,
        app_release=APP_RELEASE,
        analyzed_filename=analyzed_filename or (summary.get("filename") if summary else None),
    )




def _topdesk_equipment_type(row):
    textv = normalize(" ".join(str(row.get(k) or "") for k in ("ID do objeto","Categoria","Subcategoria","Pedido")))
    if "ATM" in textv: return "ATM"
    if "BLOQ" in textv or "BLOQUEIO" in textv: return "BLOQUEIO"
    if "VALID" in textv or "RECARG" in textv: return "VALIDADOR"
    if "POS" in textv: return "POS"
    if "TDI" in textv: return "TDI"
    if "RACK" in textv: return "RACK"
    return "OUTRO"


def _topdesk_demand_type(ticket):
    t = normalize((ticket.incident_type or "") + " " + (ticket.category or ""))
    return "INCIDENTE" if "INCIDENT" in t or "FALHA" in t else "SOLICITACAO"


def _topdesk_match_location(row, loc_cache=None):
    raw = " ".join(str(row.get(k) or "") for k in ("ID do objeto","Pedido","Ação","Subcategoria"))
    nraw = normalize(raw)
    if not nraw: return None
    # V55: a lista de localidades é pré-carregada uma vez por importação.
    # Antes, Location.query.all() era executado para CADA chamado.
    if loc_cache is None:
        loc_cache=[(loc,normalize(loc.location),[x for x in normalize(loc.location).split() if len(x)>=4]) for loc in Location.query.all()]
    best=None; bestscore=0
    for loc,name,tokens in loc_cache:
        if name and len(name)>=4 and name in nraw:
            score=100+len(name)
        else:
            score=sum(8 for x in tokens if x in nraw)
        if score>bestscore:
            bestscore=score; best=loc
    return best if bestscore>=8 else None


def _cell_text_xml(cell):
    m=re.search(r"<t[^>]*>(.*?)</t>",cell,re.S)
    if m: return html_lib.unescape(re.sub(r"<.*?>","",m.group(1)))
    m=re.search(r"<v>(.*?)</v>",cell,re.S)
    return html_lib.unescape(m.group(1)) if m else ""


def _read_topdesk_xlsx(stream):
    data=stream.read()
    rows=[]
    try:
        wb=load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        ws=wb.active
        allrows=list(ws.iter_rows(values_only=True))
        if not allrows: return []
        headers=[str(x or "").strip() for x in allrows[0]]
        for values in allrows[1:]:
            row={headers[i]: values[i] if i<len(values) else None for i in range(len(headers))}
            if any(v not in (None,"") for v in row.values()): rows.append(row)
        return rows
    except Exception:
        pass
    # Alguns exports TopDesk trazem XML Excel com caracteres inválidos para parsers estritos.
    with zipfile.ZipFile(io.BytesIO(data)) as z:
        xml=z.read("xl/worksheets/sheet1.xml").decode("utf-8","replace")
    rawrows=re.findall(r"<row\b[^>]*>(.*?)</row>",xml,re.S)
    matrix=[]
    for rr in rawrows:
        vals=[]
        cells=re.findall(r"<c\b[^>]*>.*?</c>",rr,re.S)
        for c in cells: vals.append(_cell_text_xml(c))
        matrix.append(vals)
    if not matrix: return []
    headers=[str(x or "").strip() for x in matrix[0]]
    for vals in matrix[1:]:
        row={headers[i]: vals[i] if i<len(vals) else None for i in range(len(headers))}
        if any(v not in (None,"") for v in row.values()): rows.append(row)
    return rows


def _ticket_open(ticket):
    s=normalize(ticket.status or "")
    return not any(x in s for x in ("RESOLVID","FECHAD","CANCEL"))


def _topdesk_suggestions(ticket, limit=5):
    techs=User.query.filter(User.active.is_(True), User.role=="technician").order_by(User.name).all()
    loc=db.session.get(Location,ticket.location_id) if ticket.location_id else None
    latest={}
    if techs:
        ids=[u.id for u in techs]
        for p in TechnicianPosition.query.filter(TechnicianPosition.user_id.in_(ids)).order_by(TechnicianPosition.captured_at.desc()).all():
            latest.setdefault(p.user_id,p)
    out=[]
    for u in techs:
        active=TopDeskTicket.query.filter(TopDeskTicket.assigned_technician_id==u.id).all()
        load=sum(1 for t in active if _ticket_open(t) and t.work_status!="CONCLUIDO")
        p=latest.get(u.id); dist=None
        if p and loc and loc.reference_latitude is not None and loc.reference_longitude is not None:
            dist=_haversine_m(p.latitude,p.longitude,loc.reference_latitude,loc.reference_longitude)
        score=(dist if dist is not None else 999999)+(load*5000)
        out.append({"id":u.id,"name":u.name,"distance_m":round(dist) if dist is not None else None,"active":load,"score":score})
    out.sort(key=lambda x:(x["score"],x["name"]))
    return out[:limit]


@app.route("/topdesk", methods=["GET"])
@topdesk_required
def topdesk_page():
    technicians=User.query.filter(User.active.is_(True),User.role=="technician").order_by(User.name).all()
    return render_template("topdesk.html", technicians=technicians, app_release=APP_RELEASE)


@app.post("/topdesk/import")
@topdesk_required
def topdesk_import():
    """V52.8: inicia importação em background e devolve job_id para polling."""
    f=request.files.get("file")
    if not f or not f.filename:
        return jsonify({"ok":False,"error":"Selecione o arquivo exportado do TopDesk."}),400
    if not f.filename.lower().endswith((".xlsx",".xlsm")):
        return jsonify({"ok":False,"error":"Use arquivo XLSX exportado do TopDesk."}),400
    raw=f.read()
    if not raw:
        return jsonify({"ok":False,"error":"Arquivo vazio."}),400
    # evita duas cargas simultâneas no mesmo processo
    with TOPDESK_IMPORT_LOCK:
        active=[x for x in TOPDESK_IMPORT_JOBS.values() if x.get("status") in ("QUEUED","READING","PROCESSING")]
        if active:
            return jsonify({"ok":False,"error":"Já existe uma importação TopDesk em andamento."}),409
    job_id=uuid.uuid4().hex
    _td_job_update(job_id,status="QUEUED",phase="Preparando arquivo",filename=secure_filename(f.filename),progress=1,total=0,processed=0,inserted=0,updated=0,ignored=0,errors=0,started_at=datetime.utcnow().isoformat()+"Z",finished_at=None,user_id=session.get("user_id"))
    thread=threading.Thread(target=_topdesk_import_worker,args=(job_id,raw,secure_filename(f.filename),session.get("user_id")),daemon=True,name=f"topdesk-import-{job_id[:8]}")
    thread.start()
    return jsonify({"ok":True,"job_id":job_id,"status":"QUEUED"}),202


def _topdesk_import_worker(job_id, raw, filename, user_id):
    started=time.time()
    with app.app_context():
        try:
            _td_job_update(job_id,status="READING",phase="Lendo planilha",progress=5)
            rows=_read_topdesk_xlsx(io.BytesIO(raw))
            if not rows:
                raise ValueError("Arquivo sem registros válidos.")
            headers=set(rows[0].keys())
            novo_padrao={"Dia/hora da criação","Status","ID do objeto","Categoria","Subcategoria","Operador","Nível","Pedido","Ação","Anexos"}
            padrao_antigo="Número do incidente" in headers
            if not padrao_antigo and len(novo_padrao.intersection(headers)) < 6:
                raise ValueError("Formato TopDesk não reconhecido. Use o padrão Campo - Dogma.")
            total=len(rows); inserted=updated=ignored=errors=0
            _td_job_update(job_id,status="PROCESSING",phase="Preparando índices de importação",progress=8,total=total,processed=0)
            existing={t.ticket_number:t for t in TopDeskTicket.query.all()}
            loc_cache=[(loc,normalize(loc.location),[x for x in normalize(loc.location).split() if len(x)>=4]) for loc in Location.query.all()]
            _td_job_update(job_id,status="PROCESSING",phase="Importando chamados",progress=10,total=total,processed=0)
            for pos,row in enumerate(rows, start=1):
                try:
                    obj=str(row.get("ID do objeto") or "").strip()
                    created=str(row.get("Dia/hora da criação") or "").strip()
                    num=str(row.get("Número do incidente") or "").strip()
                    if not num:
                        stable="|".join([obj,created,str(row.get("Categoria") or ""),str(row.get("Subcategoria") or "")])
                        num="TD-"+hashlib.sha1(stable.encode("utf-8","ignore")).hexdigest()[:16].upper()
                    if not obj and not str(row.get("Pedido") or "").strip():
                        ignored+=1
                        continue
                    t=existing.get(num)
                    is_new=t is None
                    if is_new:
                        t=TopDeskTicket(ticket_number=num); db.session.add(t); existing[num]=t
                    t.object_id=obj
                    t.category=str(row.get("Categoria") or "").strip()
                    t.subcategory=str(row.get("Subcategoria") or "").strip()
                    t.incident_type=str(row.get("Tipo de incidente") or row.get("Nível") or "").strip()
                    t.status=str(row.get("Status") or "").strip()
                    t.operator=str(row.get("Operador") or "").strip()
                    t.created_at_text=created
                    t.created_at=_td_dt(created)
                    _line,_station,_model=_td_object_parts(obj)
                    t.line_code=_line or None
                    t.station_code=_station or None
                    t.model_code=_model or None
                    t.sla_target_text=str(row.get("Data alvo do SLA") or "").strip()
                    t.requester=str(row.get("Nome do solicitante") or "").strip()
                    t.request_text=str(row.get("Pedido") or "").strip()
                    t.action_text=str(row.get("Ação") or "").strip()
                    t.attachments_text=str(row.get("Anexos") or "").strip()
                    t.source_file=filename[:300]
                    t.equipment_type=_topdesk_equipment_type(row)
                    loc=_topdesk_match_location(row,loc_cache)
                    if loc: t.location_id=loc.id
                    t.priority="ALTA" if _topdesk_demand_type(t)=="INCIDENTE" else "NORMAL"
                    t.last_import_at=datetime.utcnow(); t.updated_at=datetime.utcnow()
                    if is_new: inserted+=1
                    else: updated+=1
                except Exception:
                    errors+=1
                if pos % 1000 == 0:
                    db.session.commit()
                    progress=min(99,10+int(pos/max(total,1)*89))
                    _td_job_update(job_id,processed=pos,inserted=inserted,updated=updated,ignored=ignored,errors=errors,progress=progress,elapsed_seconds=int(time.time()-started))
            db.session.commit()
            batch=TopDeskImportBatch(filename=filename,imported_by=user_id,row_count=total,inserted_count=inserted,updated_count=updated,error_count=errors)
            db.session.add(batch); db.session.commit()
            _td_cache_clear()
            _td_job_update(job_id,status="DONE",phase="Importação concluída",progress=100,total=total,processed=total,inserted=inserted,updated=updated,ignored=ignored,errors=errors,elapsed_seconds=int(time.time()-started),finished_at=datetime.utcnow().isoformat()+"Z",format="TopDesk legado" if padrao_antigo else "Campo - Dogma")
        except Exception as exc:
            db.session.rollback()
            _td_job_update(job_id,status="ERROR",phase="Importação interrompida",error=f"{type(exc).__name__}: {exc}",elapsed_seconds=int(time.time()-started),finished_at=datetime.utcnow().isoformat()+"Z")
        finally:
            db.session.remove()


@app.get("/api/topdesk/import/<job_id>/status")
@topdesk_required
def topdesk_import_status(job_id):
    job=_td_job_snapshot(job_id)
    if not job:
        return jsonify({"ok":False,"error":"Importação não encontrada ou processo reiniciado."}),404
    return jsonify({"ok":True,**job})


@app.get("/api/topdesk/import/active")
@login_required
def topdesk_import_active():
    """V55: permite acompanhar a carga mesmo navegando para outra tela."""
    uid=session.get('user_id')
    with TOPDESK_IMPORT_LOCK:
        active=[dict(v,job_id=k) for k,v in TOPDESK_IMPORT_JOBS.items() if v.get('status') in ('QUEUED','READING','PROCESSING') and (not v.get('user_id') or v.get('user_id')==uid)]
    if not active: return jsonify({'ok':True,'active':False})
    active.sort(key=lambda x:x.get('started_at') or '',reverse=True)
    return jsonify({'ok':True,'active':True,'job':active[0]})


@app.get("/api/topdesk/tickets")
@topdesk_required
def topdesk_tickets_api():
    q=TopDeskTicket.query
    status=request.args.get("status","").strip(); eq=request.args.get("equipment_type","").strip(); assigned=request.args.get("assigned","").strip(); search=request.args.get("q","").strip()
    if status: q=q.filter(TopDeskTicket.status==status)
    if eq: q=q.filter(TopDeskTicket.equipment_type==eq)
    if assigned=="unassigned": q=q.filter(TopDeskTicket.assigned_technician_id.is_(None))
    if search:
        like=f"%{search}%"; q=q.filter(db.or_(TopDeskTicket.ticket_number.ilike(like),TopDeskTicket.object_id.ilike(like),TopDeskTicket.request_text.ilike(like),TopDeskTicket.category.ilike(like)))
    rows=q.order_by(TopDeskTicket.last_import_at.desc()).limit(1000).all()
    users={u.id:u for u in User.query.filter(User.id.in_([x.assigned_technician_id for x in rows if x.assigned_technician_id] or [-1])).all()}
    locs={l.id:l for l in Location.query.filter(Location.id.in_([x.location_id for x in rows if x.location_id] or [-1])).all()}
    return jsonify({"ok":True,"tickets":[{
      "id":t.id,"number":t.ticket_number,"object_id":t.object_id,"category":t.category,"subcategory":t.subcategory,"incident_type":t.incident_type,"demand_type":_topdesk_demand_type(t),"status":t.status,"work_status":t.work_status,"priority":t.priority,"equipment_type":t.equipment_type,"operator":t.operator,"created_at":t.created_at_text,"sla":t.sla_target_text,"requester":t.requester,"request":t.request_text,"action":t.action_text,"assigned_id":t.assigned_technician_id,"assigned":users.get(t.assigned_technician_id).name if users.get(t.assigned_technician_id) else None,"location_id":t.location_id,"location":locs.get(t.location_id).location if locs.get(t.location_id) else None,"line":locs.get(t.location_id).line if locs.get(t.location_id) else None,"suggestions":_topdesk_suggestions(t,3)
    } for t in rows]})


@app.post("/api/topdesk/tickets/<int:ticket_id>/assign")
@topdesk_required
def topdesk_assign(ticket_id):
    t=db.session.get(TopDeskTicket,ticket_id)
    if not t: return jsonify({"ok":False,"error":"Chamado não encontrado."}),404
    data=request.get_json(silent=True) or {}; tech_id=data.get("technician_id")
    u=db.session.get(User,int(tech_id)) if tech_id not in (None,"") else None
    if u and (not u.active or u.role!="technician"): return jsonify({"ok":False,"error":"Técnico inválido."}),400
    old=t.assigned_technician_id; t.assigned_technician_id=u.id if u else None; t.work_status="RECEBIDO" if u else "NAO_ATRIBUIDO"; t.updated_at=datetime.utcnow()
    db.session.add(TopDeskActivity(ticket_id=t.id,technician_id=t.assigned_technician_id,actor_user_id=session.get("user_id"),event_type="REATRIBUIDO" if old else "ATRIBUIDO",status=t.work_status,note=f"{db.session.get(User,old).name if old and db.session.get(User,old) else 'Sem técnico'} → {u.name if u else 'Sem técnico'}"))
    db.session.commit(); return jsonify({"ok":True,"assigned":u.name if u else None})


def _td_dt(value):
    if isinstance(value, datetime): return value
    textv=str(value or '').strip()
    if not textv: return None
    for fmt in ("%Y-%m-%d %H:%M:%S","%d/%m/%Y %H:%M:%S","%d/%m/%Y %H:%M","%Y-%m-%dT%H:%M:%S"):
        try: return datetime.strptime(textv[:19],fmt)
        except Exception: pass
    try: return datetime.fromisoformat(textv.replace('Z',''))
    except Exception: return None


def _td_object_parts(object_id):
    parts=[x.strip() for x in str(object_id or '').split('-') if x.strip()]
    line=parts[0] if parts and re.match(r'^L\\d+',parts[0],re.I) else ''
    station=parts[1] if len(parts)>1 and line else ''
    model=parts[-1] if len(parts)>=4 else ''
    return line.upper(),station.upper(),model.upper()


def _td_filter_rows(args):
    """V56-A: filtros de analytics executados prioritariamente no banco.

    Datas, linha, estação e modelo passam a ter colunas normalizadas/indexadas.
    O fallback legado permanece apenas para registros antigos ainda não normalizados.
    """
    start=_td_dt(args.get('start')); end=_td_dt(args.get('end'))
    if start is None: start=datetime(2026,1,1)
    eq=normalize(args.get('equipment_type','')); linef=normalize(args.get('line','')); locf=normalize(args.get('location',''))
    modelf=normalize(args.get('model','')); catf=normalize(args.get('category','')); subf=normalize(args.get('subcategory',''))
    opf=normalize(args.get('operator','')); statusf=normalize(args.get('status','')); slaf=normalize(args.get('sla',''))
    search=normalize(args.get('q',''))

    q=(db.session.query(TopDeskTicket, Location.line, Location.location)
       .outerjoin(Location, TopDeskTicket.location_id==Location.id))
    if eq: q=q.filter(func.upper(func.coalesce(TopDeskTicket.equipment_type,''))==eq)
    if catf: q=q.filter(func.upper(func.coalesce(TopDeskTicket.category,''))==catf)
    if subf: q=q.filter(func.upper(func.coalesce(TopDeskTicket.subcategory,''))==subf)
    if opf: q=q.filter(func.upper(func.coalesce(TopDeskTicket.operator,''))==opf)
    if statusf: q=q.filter(func.upper(func.coalesce(TopDeskTicket.status,''))==statusf)
    if slaf=='COM SLA': q=q.filter(func.length(func.trim(func.coalesce(TopDeskTicket.sla_target_text,'')))>0)
    elif slaf=='SEM SLA': q=q.filter(func.length(func.trim(func.coalesce(TopDeskTicket.sla_target_text,'')))==0)
    if search:
        like=f"%{search}%"
        q=q.filter(db.or_(
            func.upper(func.coalesce(TopDeskTicket.ticket_number,'')).like(like),
            func.upper(func.coalesce(TopDeskTicket.object_id,'')).like(like),
            func.upper(func.coalesce(TopDeskTicket.request_text,'')).like(like),
            func.upper(func.coalesce(TopDeskTicket.category,'')).like(like),
            func.upper(func.coalesce(TopDeskTicket.subcategory,'')).like(like),
        ))

    # V56-A.1: compatibilidade durante o backfill. Registros ainda não normalizados
    # continuam elegíveis e usam os campos legados, evitando a queda 50k -> poucos tickets.
    if end:
        q=q.filter(db.or_(
            db.and_(TopDeskTicket.created_at.isnot(None), TopDeskTicket.created_at >= start, TopDeskTicket.created_at <= end.replace(hour=23,minute=59,second=59)),
            TopDeskTicket.created_at.is_(None)
        ))
    else:
        q=q.filter(db.or_(TopDeskTicket.created_at >= start, TopDeskTicket.created_at.is_(None)))

    out=[]
    for t,loc_line,loc_name in q.all():
        dt=t.created_at or _td_dt(t.created_at_text)
        if not dt or dt < start or (end and dt > end.replace(hour=23,minute=59,second=59)):
            continue
        legacy_line,legacy_station,legacy_model=_td_object_parts(t.object_id)
        line=t.line_code or legacy_line or (loc_line or '')
        location=(loc_name or '') or t.station_code or legacy_station
        model=t.model_code or legacy_model or ''
        if linef and normalize(line)!=linef: continue
        if locf and normalize(location)!=locf and normalize(t.station_code or legacy_station)!=locf: continue
        if modelf and normalize(model)!=modelf: continue
        out.append((t,dt,line,location,model))
    return out


def _td_recurrence(rows, days):
    groups={}
    for t,dt,*_ in rows:
        if not dt or not t.object_id: continue
        groups.setdefault((normalize(t.object_id),normalize(t.subcategory)),[]).append(dt)
    repeated=0; eligible=0
    for dates in groups.values():
        dates.sort()
        for i in range(1,len(dates)):
            eligible+=1
            if (dates[i]-dates[i-1]).total_seconds() <= days*86400: repeated+=1
    return round(repeated*100/eligible,1) if eligible else 0


def _td_rank_map(rows, dimension):
    out={}
    for t,dt,line,loc,model in rows:
        if dimension=='failure': key=t.subcategory or 'Sem subcategoria'
        elif dimension=='line': key=line or ''
        elif dimension=='location': key=loc or ''
        elif dimension=='model': key=model or ''
        elif dimension=='object': key=t.object_id or ''
        else: key=''
        if key: out[key]=out.get(key,0)+1
    return out


def _td_previous_args(args):
    start=_td_dt(args.get('start')); end=_td_dt(args.get('end'))
    if not start or not end or end < start: return None
    days=(end.date()-start.date()).days+1
    prev_end=start-timedelta(days=1); prev_start=prev_end-timedelta(days=days-1)
    d={k:v for k,v in args.items()}
    d['start']=prev_start.strftime('%Y-%m-%d'); d['end']=prev_end.strftime('%Y-%m-%d')
    return d


def _td_ranked(current_map, n=15, total=0, previous_map=None):
    vals=list(current_map.values()); avg=(sum(vals)/len(vals)) if vals else 0
    rows=[]
    for k,v in sorted(current_map.items(),key=lambda x:x[1],reverse=True)[:n]:
        prev=(previous_map or {}).get(k,0) if previous_map is not None else None
        delta=None
        if previous_map is not None:
            if prev: delta=round((v-prev)*100/prev,1)
            elif v: delta=100.0
            else: delta=0.0
        rows.append({'name':k,'count':v,'share':round(v*100/max(1,total),1),
                     'vs_average_pct':round((v-avg)*100/avg,1) if avg else 0,
                     'previous_count':prev,'vs_previous_pct':delta})
    return rows



@app.get("/topdesk/tv")
@topdesk_required
def topdesk_tv_page():
    return render_template("topdesk_tv.html", app_release=APP_RELEASE)

@app.get("/api/topdesk/analytics")
@dashboard_required
def topdesk_analytics_api():
    # Cache por querystring; filtros iguais deixam de reler/recalcular toda a base.
    cache_key='analytics|'+request.query_string.decode('utf-8','ignore')
    cached=_td_cache_get(cache_key)
    if cached is not None:
        payload=dict(cached); payload['cache']='HIT'; return jsonify(payload)

    rows=_td_filter_rows(request.args); total=len(rows)
    objects={normalize(t.object_id) for t,*_ in rows if t.object_id}; locations={normalize(loc) for *_,loc,_ in rows if loc}; operators={normalize(t.operator) for t,*_ in rows if t.operator}
    monthly={}; failures={}; lines={}; models={}; locs={}; hours={}; weekdays={}; heat={}; tech={}; objcount={}
    weekday_names=['Seg','Ter','Qua','Qui','Sex','Sáb','Dom']
    for t,dt,line,loc,model in rows:
        if dt:
            if dt.year >= 2026:
                mk=dt.strftime('%Y-%m'); monthly[mk]=monthly.get(mk,0)+1; hours[str(dt.hour).zfill(2)]=hours.get(str(dt.hour).zfill(2),0)+1; weekdays[weekday_names[dt.weekday()]]=weekdays.get(weekday_names[dt.weekday()],0)+1
            heat[(weekday_names[dt.weekday()],str(dt.hour).zfill(2))]=heat.get((weekday_names[dt.weekday()],str(dt.hour).zfill(2)),0)+1
        failures[t.subcategory or 'Sem subcategoria']=failures.get(t.subcategory or 'Sem subcategoria',0)+1
        if line: lines[line]=lines.get(line,0)+1
        if model: models[model]=models.get(model,0)+1
        if loc: locs[loc]=locs.get(loc,0)+1
        if t.object_id: objcount[t.object_id]=objcount.get(t.object_id,0)+1
        op=t.operator or 'Sem operador'; d=tech.setdefault(op,{'tickets':0,'days':set(),'objects':set(),'locations':set(),'lines':set()}); d['tickets']+=1
        if dt:d['days'].add(dt.date().isoformat())
        if t.object_id:d['objects'].add(t.object_id)
        if loc:d['locations'].add(loc)
        if line:d['lines'].add(line)

    byop={}
    for row in rows: byop.setdefault(row[0].operator or 'Sem operador',[]).append(row)
    prod=[]
    for op,d in tech.items():
        days=max(1,len(d['days'])); prod.append({'operator':op,'tickets':d['tickets'],'active_days':len(d['days']),'per_day':round(d['tickets']/days,1),'objects':len(d['objects']),'locations':len(d['locations']),'lines':len(d['lines']),'recurrence7':_td_recurrence(byop.get(op,[]),7)})
    prod.sort(key=lambda x:x['tickets'],reverse=True)

    # Referência: período imediatamente anterior quando o usuário definiu início/fim.
    prev_args=_td_previous_args(request.args)
    prev_rows=_td_filter_rows(prev_args) if prev_args else None
    prev_failure=_td_rank_map(prev_rows,'failure') if prev_rows is not None else None
    prev_line=_td_rank_map(prev_rows,'line') if prev_rows is not None else None
    prev_loc=_td_rank_map(prev_rows,'location') if prev_rows is not None else None
    prev_model=_td_rank_map(prev_rows,'model') if prev_rows is not None else None
    prev_obj=_td_rank_map(prev_rows,'object') if prev_rows is not None else None
    reference_label='Período anterior' if prev_rows is not None else 'Média das categorias no recorte'

    filters={'lines':sorted(set(x[2] for x in rows if x[2])),'locations':sorted(set(x[3] for x in rows if x[3])),'models':sorted(set(x[4] for x in rows if x[4])),'categories':sorted(set(t.category for t,*_ in rows if t.category)),'subcategories':sorted(set(t.subcategory for t,*_ in rows if t.subcategory)),'operators':sorted(set(t.operator for t,*_ in rows if t.operator)),'statuses':sorted(set(t.status for t,*_ in rows if t.status)),'equipment_types':sorted(set(t.equipment_type for t,*_ in rows if t.equipment_type))}
    chronic=[]
    obj_avg=(sum(objcount.values())/len(objcount)) if objcount else 0
    for k,v in sorted(objcount.items(),key=lambda x:x[1],reverse=True)[:20]:
        level='CRÔNICO' if v>=50 else 'CRÍTICO' if v>=25 else 'ATENÇÃO'
        prev=(prev_obj or {}).get(k,0) if prev_obj is not None else None
        chronic.append({'object_id':k,'count':v,'level':level,'share':round(v*100/max(1,total),1),
                        'times_average':round(v/max(obj_avg,0.01),1),'previous_count':prev,
                        'vs_previous_pct':(round((v-prev)*100/prev,1) if prev else (100.0 if prev==0 else None))})
    alerts=[]
    object_rows={}
    for rr in rows:
        if rr[0].object_id: object_rows.setdefault(rr[0].object_id,[]).append(rr)
    for x in chronic[:5]:
        line, station, model = _td_object_parts(x['object_id']); rr=object_rows.get(x['object_id'],[])
        parts=[p for p in str(x['object_id'] or '').split('-') if p]; atm_num=next((p for p in parts if p.isdigit()), str(x['object_id'] or ''))
        friendly=' · '.join(y for y in [f'ATM {atm_num}' if atm_num else '', model, line, station] if y)
        failmap={}; dates=[]
        for t,dt,*_rest in rr:
            fail=t.subcategory or t.category or 'Sem classificação'; failmap[fail]=failmap.get(fail,0)+1
            if dt: dates.append(dt)
        top_fail=max(failmap.items(),key=lambda z:z[1]) if failmap else ('Sem classificação',0); last=max(dates).strftime('%d/%m/%Y') if dates else '—'
        recent30=sum(1 for t,dt,*_ in rr if dt and dt >= datetime.now()-timedelta(days=30))
        alerts.append({'level':'critical' if x['level']=='CRÔNICO' else 'warning','title':f"{friendly or x['object_id']} · {x['level']}",'detail':f"{x['count']} chamados · {x['times_average']}× média · {x['share']}% do recorte",'cause':f"Principal: {top_fail[0]} ({top_fail[1]})",'recent':f"30 dias: {recent30} · último: {last}",'object_id':x['object_id']})

    payload={'ok':True,'release':APP_RELEASE,'cache':'MISS','reference_label':reference_label,
      'kpis':{'tickets':total,'objects':len(objects),'locations':len(locations),'operators':len(operators),'per_object':round(total/max(1,len(objects)),1),'recurrence24':_td_recurrence(rows,1),'recurrence7':_td_recurrence(rows,7),'recurrence30':_td_recurrence(rows,30)},
      'monthly':[{'month':k,'count':monthly[k]} for k in sorted(monthly)],
      'failures':_td_ranked(failures,30,total,prev_failure),'lines':_td_ranked(lines,30,total,prev_line),'locations':_td_ranked(locs,30,total,prev_loc),'models':_td_ranked(models,30,total,prev_model),
      'hours':[{'hour':str(h).zfill(2),'count':hours.get(str(h).zfill(2),0)} for h in range(24)],'weekdays':[{'day':d,'count':weekdays.get(d,0)} for d in weekday_names],
      'heatmap':[{'day':d,'hour':str(h).zfill(2),'count':heat.get((d,str(h).zfill(2)),0)} for d in weekday_names for h in range(24)],'productivity':prod[:30],'chronic':chronic,'alerts':alerts,'filters':filters}
    _td_cache_put(cache_key,payload)
    return jsonify(payload)


_TOPDESK_DASH_CACHE = {}
_TOPDESK_DASH_CACHE_TTL = int(os.getenv("TOPDESK_DASH_CACHE_TTL", "900"))

@app.get("/api/topdesk/dashboard")
@dashboard_required
def topdesk_dashboard_api():
    # V56-B REV2: cache curto por conjunto de filtros. Evita recalcular 50k chamados
    # várias vezes durante a mesma navegação/refresh da dashboard.
    # V58: ignora cache-busters/timestamps do navegador; somente filtros funcionais compõem a chave.
    ignored={"_","t","ts","timestamp","cache","cacheBust","cb"}
    cache_key="&".join(f"{k}={v}" for k,v in sorted(request.args.items()) if k not in ignored)
    now=time.time(); cached=_TOPDESK_DASH_CACHE.get(cache_key)
    if cached and now-cached[0] < _TOPDESK_DASH_CACHE_TTL:
        payload=dict(cached[1]); payload["cache"]="HIT"; return jsonify(payload)
    # V56-D: dashboard sem filtros usa agregações SQL em vez de materializar ~50 mil objetos Python.
    # Analytics detalhado continua usando o pipeline completo quando há filtros.
    if not request.args:
        # REV4.2: usar o status operacional indexado no resumo principal.
        # Evita UPPER/LIKE sobre toda a coluna textual do TopDesk em cada MISS.
        ws_expr=func.upper(func.coalesce(TopDeskTicket.work_status,''))
        total,resolved,assigned=db.session.query(
            func.count(TopDeskTicket.id),
            func.coalesce(func.sum(case((ws_expr.in_(("CONCLUIDO","CONCLUÍDO","FECHADO","RESOLVIDO","CANCELADO")),1),else_=0)),0),
            func.coalesce(func.sum(case((TopDeskTicket.assigned_technician_id.isnot(None),1),else_=0)),0)
        ).one()
        total=int(total or 0); resolved=int(resolved or 0); assigned=int(assigned or 0)
        by_status={k or 'Sem status':int(v) for k,v in db.session.query(TopDeskTicket.status,func.count(TopDeskTicket.id)).group_by(TopDeskTicket.status).all()}
        by_type={k or 'OUTRO':int(v) for k,v in db.session.query(TopDeskTicket.equipment_type,func.count(TopDeskTicket.id)).group_by(TopDeskTicket.equipment_type).all()}
        top_locations=[{'name':k,'count':int(v)} for k,v in db.session.query(Location.location,func.count(TopDeskTicket.id)).join(TopDeskTicket,TopDeskTicket.location_id==Location.id).group_by(Location.location).order_by(func.count(TopDeskTicket.id).desc()).limit(12).all()]
        payload={"ok":True,"cache":"MISS","mode":"SQL_AGG","total":total,"open":max(0,total-resolved),"resolved":resolved,"assigned":assigned,"unassigned":max(0,total-assigned),"by_status":by_status,"by_type":by_type,"top_locations":top_locations}
        _TOPDESK_DASH_CACHE[cache_key]=(now,payload)
        return jsonify(payload)
    rows=_td_filter_rows(request.args); tickets=[x[0] for x in rows]; total=len(tickets); openrows=[t for t in tickets if _ticket_open(t)]; resolved=total-len(openrows)
    by_status={}; by_type={}; by_location={}; assigned=0
    for t,dt,line,location,model in rows:
        by_status[t.status or "Sem status"]=by_status.get(t.status or "Sem status",0)+1
        by_type[t.equipment_type or "OUTRO"]=by_type.get(t.equipment_type or "OUTRO",0)+1
        if t.assigned_technician_id: assigned+=1
        if location: by_location[location]=by_location.get(location,0)+1
    top_locations=sorted(by_location.items(),key=lambda x:x[1],reverse=True)[:12]
    payload={"ok":True,"cache":"MISS","total":total,"open":len(openrows),"resolved":resolved,"assigned":assigned,"unassigned":total-assigned,"by_status":by_status,"by_type":by_type,"top_locations":[{"name":k,"count":v} for k,v in top_locations]}
    _TOPDESK_DASH_CACHE[cache_key]=(now,payload)
    if len(_TOPDESK_DASH_CACHE)>40:
        for k,v in list(_TOPDESK_DASH_CACHE.items()):
            if now-v[0] > _TOPDESK_DASH_CACHE_TTL: _TOPDESK_DASH_CACHE.pop(k,None)
    return jsonify(payload)


@app.get("/topdesk/export.xlsx")
@topdesk_required
def topdesk_export():
    rows=_td_filter_rows(request.args); wb=Workbook(); ws=wb.active; ws.title="Chamados"
    ws.append(["Número","Tipo","Status TopDesk","Status Campo","Prioridade","Equipamento","Objeto","Modelo","Categoria","Subcategoria","Operador","Criação","SLA","Linha","Localidade","Técnico","Pedido","Ação"])
    for t,dt,line,location,model in rows:
        tech=db.session.get(User,t.assigned_technician_id) if t.assigned_technician_id else None
        ws.append([t.ticket_number,_topdesk_demand_type(t),t.status,t.work_status,t.priority,t.equipment_type,t.object_id,model,t.category,t.subcategory,t.operator,t.created_at_text,t.sla_target_text,line,location,tech.name if tech else "",t.request_text,t.action_text])
    # Abas analíticas respeitando o mesmo recorte.
    wp=wb.create_sheet("Produtividade"); wp.append(["Operador","Chamados","Dias ativos","Chamados/dia","Equipamentos","Localidades","Linhas","Reincidência 7d %"])
    byop={}
    for row in rows: byop.setdefault(row[0].operator or 'Sem operador',[]).append(row)
    for op,rr in sorted(byop.items(),key=lambda x:len(x[1]),reverse=True):
        days={x[1].date() for x in rr if x[1]}; objs={x[0].object_id for x in rr if x[0].object_id}; locs={x[3] for x in rr if x[3]}; lines={x[2] for x in rr if x[2]}
        wp.append([op,len(rr),len(days),round(len(rr)/max(1,len(days)),1),len(objs),len(locs),len(lines),_td_recurrence(rr,7)])
    wc=wb.create_sheet("Crônicos"); wc.append(["Objeto","Chamados"])
    counts={}
    for t,*_ in rows:
        if t.object_id: counts[t.object_id]=counts.get(t.object_id,0)+1
    for k,v in sorted(counts.items(),key=lambda x:x[1],reverse=True): wc.append([k,v])
    for sh in wb.worksheets:
        sh.freeze_panes='A2'; sh.auto_filter.ref=sh.dimensions; sh.row_dimensions[1].height=24
        for cell in sh[1]: cell.font=Font(bold=True,color='FFFFFF'); cell.fill=PatternFill('solid',fgColor='17365D')
        for col in range(1,sh.max_column+1): sh.column_dimensions[get_column_letter(col)].width=min(45,max(12,max(len(str(sh.cell(r,col).value or '')) for r in range(1,min(sh.max_row,200)+1))+2))
    bio=io.BytesIO(); wb.save(bio); bio.seek(0); return send_file(bio,as_attachment=True,download_name=f"topdesk_dashboard_{datetime.now().strftime('%Y%m%d_%H%M')}.xlsx",mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")


@app.get("/api/minhas-atividades")
@field_required
def my_work_api():
    uid=session.get("user_id"); q=TopDeskTicket.query.filter(TopDeskTicket.assigned_technician_id==uid)
    rows=q.order_by(case((TopDeskTicket.priority=="ALTA",0),else_=1),TopDeskTicket.id.desc()).all()
    locs={l.id:l for l in Location.query.filter(Location.id.in_([x.location_id for x in rows if x.location_id] or [-1])).all()}
    return jsonify({"ok":True,"tickets":[{"id":t.id,"number":t.ticket_number,"demand_type":_topdesk_demand_type(t),"priority":t.priority,"status":t.status,"work_status":t.work_status,"equipment_type":t.equipment_type,"object_id":t.object_id,"location":locs.get(t.location_id).location if locs.get(t.location_id) else "Não vinculada","line":locs.get(t.location_id).line if locs.get(t.location_id) else "","request":t.request_text,"sla":t.sla_target_text} for t in rows]})


@app.post("/api/minhas-atividades/<int:ticket_id>/status")
@field_required
def my_work_status(ticket_id):
    t=db.session.get(TopDeskTicket,ticket_id); uid=session.get("user_id")
    if not t or (session.get("role")!="manager" and t.assigned_technician_id!=uid): return jsonify({"ok":False,"error":"Atividade não encontrada para este técnico."}),404
    data=request.get_json(silent=True) or {}; status=normalize(data.get("status"))
    allowed={"ACEITO","EM DESLOCAMENTO","NO LOCAL","EM ATENDIMENTO","PAUSADO","CONCLUIDO"}
    if status not in allowed: return jsonify({"ok":False,"error":"Status inválido."}),400
    lat=data.get("latitude"); lon=data.get("longitude"); t.work_status=status; t.updated_at=datetime.utcnow()
    db.session.add(TopDeskActivity(ticket_id=t.id,technician_id=t.assigned_technician_id,actor_user_id=uid,event_type="STATUS",status=status,note=str(data.get("note") or ""),latitude=float(lat) if lat not in (None,"") else None,longitude=float(lon) if lon not in (None,"") else None))
    db.session.commit(); return jsonify({"ok":True,"status":status})


@app.get("/api/v38/gps-config")
@login_required
def v38_gps_config():
    u=db.session.get(User, session.get("user_id"))
    required=bool(getattr(u,"gps_required",False)) if u else False
    interval=max(60, int(_v50_settings().get("gps_interval_seconds", os.getenv("TEAM_GPS_INTERVAL_SECONDS", "300"))))
    return jsonify({"ok": True, "enabled": required, "required": required, "interval_seconds": interval, "retention_days": max(1, int(os.getenv("TEAM_GPS_RETENTION_DAYS", "7"))), "session_token": session.get("gps_session_token") or str(session.get("user_id") or "")})

@app.get("/api/v38/diario-bordo")
@teams_view_required
def v38_diario_bordo():
    days = max(1, min(request.args.get("days", 7, type=int) or 7, 30))
    user_id = request.args.get("user_id", type=int)
    since = datetime.utcnow() - timedelta(days=days)
    uq = User.query.filter(User.active.is_(True))
    users = uq.filter(~User.role.in_(("manager", "hr"))).order_by(User.name).all()
    if not user_id:
        return jsonify({"ok": True, "days": days, "users": [{"id":u.id,"name":u.name,"role":u.role} for u in users], "events": [], "summary": {}})
    u = db.session.get(User, user_id)
    if not u or (session.get("role")=="dispatcher" and u.role in ("manager","hr")):
        return jsonify({"ok":False,"error":"Colaborador não disponível para este perfil"}),404
    positions = TechnicianPosition.query.filter(TechnicianPosition.user_id==user_id, TechnicianPosition.captured_at>=since).order_by(TechnicianPosition.captured_at).all()
    inventories = Inventory.query.filter(Inventory.technician_id==user_id, Inventory.created_at>=since).order_by(Inventory.created_at).all()
    locids={x.location_id for x in inventories}; locmap={x.id:x for x in Location.query.filter(Location.id.in_(locids)).all()} if locids else {}
    events=[]
    for x in positions: events.append({"kind":"GPS","at":x.captured_at.isoformat()+"Z","latitude":x.latitude,"longitude":x.longitude,"accuracy":x.accuracy,"detail":"Posição automática"})
    for x in inventories:
        loc=locmap.get(x.location_id); events.append({"kind":"ATIVIDADE","at":x.created_at.isoformat()+"Z","latitude":x.latitude,"longitude":x.longitude,"detail":f"{x.equipment_type} · {x.asset_identifier}","location":loc.location if loc else ""})
    for x in SessionEvent.query.filter(SessionEvent.user_id==user_id, SessionEvent.created_at>=since).order_by(SessionEvent.created_at).all():
        events.append({"kind":"SESSAO","at":x.created_at.isoformat()+"Z","detail":x.event_type})
    acts=TopDeskActivity.query.filter(TopDeskActivity.technician_id==user_id, TopDeskActivity.created_at>=since).order_by(TopDeskActivity.created_at).all()
    for a in acts:
        t=db.session.get(TopDeskTicket,a.ticket_id); loc=db.session.get(Location,t.location_id) if t and t.location_id else None
        events.append({"kind":"CHAMADO","at":a.created_at.isoformat()+"Z","latitude":a.latitude,"longitude":a.longitude,"detail":f"{t.ticket_number if t else a.ticket_id} · {a.status or a.event_type}","location":loc.location if loc else ""})
    events.sort(key=lambda x:x['at'])
    return jsonify({"ok":True,"days":days,"user":{"id":u.id,"name":u.name},"events":events,"summary":{"gps":len(positions),"equipment":len(inventories),"locations":len(locids),"first":events[0]['at'] if events else None,"last":events[-1]['at'] if events else None}})

@app.get("/api/v38/evidence-analytics")
@dashboard_required
def v38_evidence_analytics():
    summary=_evidence_summary()
    statuses=db.session.query(FieldEvidenceItem.audit_status,func.count(FieldEvidenceItem.id)).group_by(FieldEvidenceItem.audit_status).all()
    return jsonify({"ok":True,**summary,"statuses":[{"status":a or "SEM STATUS","count":int(n)} for a,n in statuses]})


@app.get("/api/evidencias-campo/resumo")
@dashboard_required
def field_evidence_summary_api():
    return jsonify({"ok": True, **_evidence_summary()})


@app.get("/evidencias-campo")
@dashboard_required
def field_evidence_page():
    q = request.args.get("q", "").strip()
    audit = request.args.get("audit", "").strip()
    location_id = request.args.get("location_id", type=int)
    visit_id = request.args.get("visit", type=int)

    query = FieldEvidenceVisit.query
    if location_id:
        query = query.filter(FieldEvidenceVisit.location_id == location_id)
    if q:
        like = f"%{q}%"
        query = query.filter(
            (FieldEvidenceVisit.station_raw.ilike(like))
            | (FieldEvidenceVisit.author.ilike(like))
            | (FieldEvidenceVisit.line_raw.ilike(like))
        )
    if audit:
        query = query.join(
            FieldEvidenceItem,
            FieldEvidenceItem.visit_id == FieldEvidenceVisit.id
        ).filter(FieldEvidenceItem.audit_status == audit).distinct()

    visits = query.order_by(
        FieldEvidenceVisit.source_date.desc(),
        FieldEvidenceVisit.source_time.desc()
    ).limit(250).all()

    selected = db.session.get(FieldEvidenceVisit, visit_id) if visit_id else None
    selected_items = []
    selected_media = []
    selected_location = None
    if selected:
        selected_items = FieldEvidenceItem.query.filter_by(visit_id=selected.id).order_by(
            FieldEvidenceItem.equipment_type, FieldEvidenceItem.identifier
        ).all()
        selected_media = FieldEvidenceMedia.query.filter_by(visit_id=selected.id).order_by(
            FieldEvidenceMedia.id
        ).all()
        selected_location = db.session.get(Location, selected.location_id) if selected.location_id else None

    audit_statuses = [
        x[0] for x in db.session.query(FieldEvidenceItem.audit_status)
        .filter(FieldEvidenceItem.audit_status.isnot(None))
        .distinct().order_by(FieldEvidenceItem.audit_status).all()
    ]

    visit_cards = []
    for visit in visits:
        item_count = FieldEvidenceItem.query.filter_by(visit_id=visit.id).count()
        media_count = FieldEvidenceMedia.query.filter_by(visit_id=visit.id).count()
        review_count = FieldEvidenceItem.query.filter(
            FieldEvidenceItem.visit_id == visit.id,
            ~FieldEvidenceItem.audit_status.in_((
                "CONFORME / JÁ INVENTARIADO",
                "CONFIRMADO EM CAMPO / FALTA PROMOVER",
                "EVIDÊNCIA FORA DO PARQUE OFICIAL",
            ))
        ).count()
        loc = db.session.get(Location, visit.location_id) if visit.location_id else None
        visit_cards.append({
            "visit": visit, "items": item_count, "media": media_count,
            "review": review_count, "location": loc
        })

    return render_template(
        "field_evidence.html",
        summary=_evidence_summary(),
        visit_cards=visit_cards,
        selected=selected,
        selected_items=selected_items,
        selected_media=selected_media,
        selected_location=selected_location,
        audit_statuses=audit_statuses,
        q=q,
        audit=audit,
    )


@app.get("/evidencias-campo/midia/<int:media_id>")
@dashboard_required
def field_evidence_media(media_id):
    media = db.session.get(FieldEvidenceMedia, media_id)
    if not media:
        return "Mídia não encontrada.", 404

    if media.storage_kind == "r2":
        try:
            raw = _r2_get_bytes(media.storage_key)
            return send_file(
                io.BytesIO(raw),
                mimetype=media.mime_type or "application/octet-stream",
                download_name=media.original_name,
                max_age=604800,
            )
        except Exception:
            return "Não foi possível recuperar a mídia do R2.", 502

    local_path = UPLOAD_DIR / "field_evidence" / media.storage_key
    if not local_path.exists():
        return "Mídia histórica registrada, mas o arquivo local temporário não existe mais. Reimporte a fonte após validar o R2.", 410
    return send_from_directory(
        UPLOAD_DIR / "field_evidence",
        media.storage_key,
        mimetype=media.mime_type,
        max_age=604800,
    )


@app.post("/evidencias-campo/item/<int:item_id>/promover")
@manager_required
def promote_evidence_item(item_id):
    item = db.session.get(FieldEvidenceItem, item_id)
    if not item:
        flash("Evidência não encontrada.")
        return redirect(url_for("field_evidence_page"))

    visit = db.session.get(FieldEvidenceVisit, item.visit_id)
    if not visit or not visit.location_id:
        flash("A evidência precisa estar associada a uma localidade antes de ser promovida.")
        return redirect(url_for("field_evidence_page", visit=visit.id if visit else None))

    type_map = {
        "ATM": "ATM",
        "VALIDADOR": "Validador de Recarga",
        "POS": "POS de Bilheteria",
    }
    inv_type = type_map.get(item.equipment_type)
    if not inv_type:
        flash("Esse tipo é mantido somente como evidência e não pode ser promovido automaticamente.")
        return redirect(url_for("field_evidence_page", visit=visit.id))

    existing = Inventory.query.filter(
        Inventory.location_id == visit.location_id,
        Inventory.equipment_type == inv_type,
        func.upper(Inventory.asset_identifier) == item.identifier.upper()
    ).first()
    if existing:
        item.inventory_id = existing.id
        item.audit_status = "CONFORME / JÁ INVENTARIADO"
        db.session.commit()
        flash("O equipamento já existia no inventário. A evidência foi vinculada.")
        return redirect(url_for("field_evidence_page", visit=visit.id))

    now = datetime.utcnow()
    inv = Inventory(
        location_id=visit.location_id,
        equipment_type=inv_type,
        base_asset_id=item.base_asset_id,
        asset_identifier=item.identifier,
        serial=item.serial or item.identifier,
        supplier="",
        model=item.model or "",
        exact_position="Importado após revisão da evidência de campo do WhatsApp.",
        mount="",
        operational_status=item.operational_status or "Não informado",
        connectivity="",
        network_id="",
        label_status="",
        in_base="Sim" if item.base_asset_id else "Não",
        divergence="Não" if item.base_asset_id else "Sim - identificação",
        notes=f"Evidência WhatsApp V4. Visita #{visit.id}. Responsável: {visit.author}.",
        technician_id=session["user_id"],
        created_at=now,
    )
    db.session.add(inv)
    db.session.flush()
    item.inventory_id = inv.id
    item.audit_status = "CONFORME / JÁ INVENTARIADO"
    db.session.commit()
    flash(f"{item.equipment_type} {item.identifier} promovido para o inventário.")
    return redirect(url_for("field_evidence_page", visit=visit.id))


@app.post("/evidencias-campo/visita/<int:visit_id>/associar")
@manager_required
def associate_evidence_visit(visit_id):
    visit = db.session.get(FieldEvidenceVisit, visit_id)
    if not visit:
        flash("Visita não encontrada.")
        return redirect(url_for("field_evidence_page"))

    location_id = request.form.get("location_id", type=int)
    loc = db.session.get(Location, location_id) if location_id else None
    if not loc:
        flash("Localidade inválida.")
        return redirect(url_for("field_evidence_page", visit=visit.id))

    visit.location_id = loc.id
    visit.match_confidence = "REVISADA PELO GESTOR"
    db.session.commit()
    flash(f"Evidência associada a {loc.location}.")
    return redirect(url_for("field_evidence_page", visit=visit.id))


@app.post("/evidencias-campo/visita/<int:visit_id>/excluir")
@manager_required
def delete_evidence_visit(visit_id):
    """V35.1: exclusão administrativa de uma visita importada e suas mídias, sem afetar a base oficial."""
    visit = db.session.get(FieldEvidenceVisit, visit_id)
    if not visit:
        flash("Visita de evidência não encontrada.")
        return redirect(url_for("field_evidence_page"))

    try:
        media_rows = FieldEvidenceMedia.query.filter_by(visit_id=visit.id).all()
        for media in media_rows:
            try:
                if media.storage_kind == "r2" and media.storage_key:
                    r2_client().delete_object(Bucket=os.environ["R2_BUCKET_NAME"], Key=media.storage_key)
                elif media.storage_key:
                    local_path = UPLOAD_DIR / "field_evidence" / media.storage_key
                    if local_path.exists():
                        local_path.unlink()
            except Exception:
                # A exclusão do registro não deve ser impedida por uma mídia histórica indisponível.
                pass
            db.session.delete(media)

        FieldEvidenceItem.query.filter_by(visit_id=visit.id).delete(synchronize_session=False)
        db.session.delete(visit)
        db.session.commit()
        flash("Visita de evidência excluída. A base oficial e a referência da estação foram preservadas.")
    except Exception as exc:
        db.session.rollback()
        flash(f"Não foi possível excluir a visita: {exc}")
    return redirect(url_for("field_evidence_page"))


@app.post("/api/admin/cleanup-test-gps")
@manager_required
def cleanup_test_gps():
    """V35.1: remove somente coordenadas de registros de teste, preservando o cadastro do equipamento.

    Por segurança, a limpeza automática considera apenas usuários/logins claramente de teste.
    O administrador principal não é incluído automaticamente.
    """
    test_users = User.query.filter(
        db.or_(
            func.lower(User.username).in_(("adil_tst", "adil_teste", "teste", "test")),
            func.lower(User.user_code).in_(("tst", "test")),
            func.lower(User.name).like("%teste%")
        )
    ).all()
    ids = [u.id for u in test_users]
    if not ids:
        return jsonify({"ok": True, "updated": 0, "message": "Nenhum usuário de teste encontrado."})
    rows = Inventory.query.filter(Inventory.technician_id.in_(ids)).all()
    updated = 0
    for inv in rows:
        if inv.latitude is not None or inv.longitude is not None:
            inv.latitude = None
            inv.longitude = None
            inv.gps_accuracy = None
            inv.gps_captured_at = None
            updated += 1
    db.session.commit()
    return jsonify({"ok": True, "updated": updated, "users": [u.username for u in test_users]})


def migrate_panorama_status_column():
    """V50.1: permite ao ADM sobrescrever o status da Visão Panorâmica sem interferir no status do Inventário."""
    with db.engine.begin() as conn:
        try:
            conn.execute(db.text("ALTER TABLE locations ADD COLUMN IF NOT EXISTS panorama_status_override VARCHAR(30)"))
        except Exception:
            # Compatibilidade com SQLite antigo, onde ADD COLUMN IF NOT EXISTS pode não existir.
            cols={row[1] for row in conn.execute(db.text("PRAGMA table_info(locations)"))} if db.engine.dialect.name=="sqlite" else set()
            if db.engine.dialect.name=="sqlite" and "panorama_status_override" not in cols:
                conn.execute(db.text("ALTER TABLE locations ADD COLUMN panorama_status_override VARCHAR(30)"))


def migrate_base_asset_columns():
    with db.engine.begin() as conn:
        statements = [
            "ALTER TABLE base_assets ADD COLUMN IF NOT EXISTS equipment_type VARCHAR(50)",
            "ALTER TABLE base_assets ADD COLUMN IF NOT EXISTS location_code VARCHAR(80)",
            "ALTER TABLE base_assets ADD COLUMN IF NOT EXISTS terminal_number VARCHAR(120)",
            "ALTER TABLE base_assets ADD COLUMN IF NOT EXISTS application VARCHAR(180)",
            "ALTER TABLE base_assets ADD COLUMN IF NOT EXISTS bom_id VARCHAR(120)",
            "ALTER TABLE base_assets ADD COLUMN IF NOT EXISTS bu_id VARCHAR(120)",
            "ALTER TABLE base_assets ADD COLUMN IF NOT EXISTS software_version VARCHAR(120)",
            "ALTER TABLE base_assets ADD COLUMN IF NOT EXISTS quantity INTEGER",
            "ALTER TABLE base_assets ADD COLUMN IF NOT EXISTS base_notes TEXT",
            "ALTER TABLE base_assets ADD COLUMN IF NOT EXISTS leasing_status VARCHAR(120)",
            "ALTER TABLE base_assets ADD COLUMN IF NOT EXISTS contract_end VARCHAR(80)",
            "ALTER TABLE base_assets ADD COLUMN IF NOT EXISTS installation_type VARCHAR(180)",
            "ALTER TABLE base_assets ADD COLUMN IF NOT EXISTS installation_date VARCHAR(80)",
            "ALTER TABLE base_assets ADD COLUMN IF NOT EXISTS teamviewer_enabled VARCHAR(20)",
            "ALTER TABLE base_assets ADD COLUMN IF NOT EXISTS teamviewer_id VARCHAR(120)",
            "ALTER TABLE base_assets ADD COLUMN IF NOT EXISTS address VARCHAR(500)",
            "ALTER TABLE base_assets ADD COLUMN IF NOT EXISTS ip_address VARCHAR(120)",
            "ALTER TABLE base_assets ADD COLUMN IF NOT EXISTS city_id VARCHAR(120)",
            "ALTER TABLE base_assets ADD COLUMN IF NOT EXISTS praja_id VARCHAR(120)",
            "ALTER TABLE base_assets ADD COLUMN IF NOT EXISTS cielo_code VARCHAR(120)",
            "ALTER TABLE base_assets ADD COLUMN IF NOT EXISTS printer_model VARCHAR(180)",
            "ALTER TABLE base_assets ADD COLUMN IF NOT EXISTS acceptor_model VARCHAR(180)",
            "ALTER TABLE base_assets ADD COLUMN IF NOT EXISTS motherboard VARCHAR(220)",
            "ALTER TABLE base_assets ADD COLUMN IF NOT EXISTS ownership_type VARCHAR(80)",
            "ALTER TABLE base_assets ADD COLUMN IF NOT EXISTS contract_name VARCHAR(180)",
            "CREATE INDEX IF NOT EXISTS ix_base_assets_equipment_type ON base_assets (equipment_type)",
            "CREATE INDEX IF NOT EXISTS ix_base_assets_teamviewer_id ON base_assets (teamviewer_id)",
            "CREATE INDEX IF NOT EXISTS ix_base_assets_contract_name ON base_assets (contract_name)",
        ]
        for statement in statements:
            conn.execute(db.text(statement))
        conn.execute(db.text("""
            UPDATE base_assets
            SET equipment_type = 'ATM'
            WHERE equipment_type IS NULL OR TRIM(equipment_type) = ''
        """))


def sync_atm_complement_v424():
    """Atualiza/insere a base ATM complementar sem importar versão de aplicação como dado mestre."""
    source=DATA_DIR/"atm_complement_20260820.json"
    if not source.exists(): return {"ok":False,"reason":"arquivo ausente"}
    rows=json.loads(source.read_text(encoding="utf-8")); inserted=updated=0
    def txt(v):
        if v is None: return ""
        return str(v).strip()
    for x in rows:
        key=txt(x.get("TOPDESK"))
        if not key: continue
        a=BaseAsset.query.filter_by(asset_key=key).first()
        if not a:
            a=BaseAsset(asset_key=key,equipment_type="ATM");db.session.add(a);inserted+=1
        else: updated+=1
        a.description=a.description or key
        a.company=txt(x.get("EMPRESA")) or a.company
        a.station_code=txt(x.get("SIGLAS")) or a.station_code
        a.line=txt(x.get("LINHAS")) or a.line
        a.locality=txt(x.get("LOCALIDADES")) or a.locality
        a.top_id=txt(x.get("ID TOP")) or a.top_id
        a.qrcode_id=txt(x.get("ID QRCODE")) or a.qrcode_id
        a.city_id=txt(x.get("ID CITY"))
        a.praja_id=txt(x.get("ID PRAJÁ"))
        a.products=txt(x.get("PRODUTOS")) or a.products
        a.serial=txt(x.get("Nº de Séries")) or a.serial
        a.model=txt(x.get("TIPO ATM")) or a.model
        a.supplier=txt(x.get("FORNECEDOR")) or a.supplier
        a.address=txt(x.get("ENDEREÇO"))
        a.transactions=txt(x.get("TRANSACIONA")) or a.transactions
        a.application=txt(x.get("APLICAÇÃO")) or a.application
        # V42.4: deliberadamente NÃO absorve VERSÃO DA APLICAÇÃO; é dado volátil.
        a.cielo_code=txt(x.get("CÓD.  CIELO"))
        a.mount=txt(x.get("FIXAÇÃO")) or a.mount
        a.printer_model=txt(x.get("IMPRESSORAS"))
        a.acceptor_model=txt(x.get("ACEITADOR"))
        a.motherboard=txt(x.get("PLACA-MÃE"))
        a.ip_address=txt(x.get("IP"))
        a.base_status=txt(x.get("STATUS")) or a.base_status
        a.quantity=int(float(txt(x.get("QTTD")) or "1")) if (txt(x.get("QTTD")) or "1").replace('.','',1).isdigit() else 1
        a.teamviewer_enabled=txt(x.get("TEAMVIEWER"))
        a.teamviewer_id=txt(x.get("ID TV"))
        leasing=txt(x.get("LEASING")); contract=txt(x.get("CONTRATO GPN"))
        a.ownership_type=("LEASING" if normalize(leasing) in ("SIM","LEASING") else ("PRÓPRIO" if normalize(leasing) in ("NAO","NÃO") else leasing))
        a.contract_name=contract if normalize(contract) not in ("NAO","NÃO","N/A","#N/A") else "SEM CONTRATO"
        a.leasing_status=a.contract_name
    db.session.commit();return {"ok":True,"inserted":inserted,"updated":updated}


def migrate_v421_columns():
    """V42.1: TeamViewer no inventário e grupo visual do relatório de visita."""
    inspector=db.inspect(db.engine)
    with db.engine.begin() as conn:
        inv_cols={c["name"] for c in inspector.get_columns("inventory")}
        if "teamviewer_id" not in inv_cols:
            conn.execute(db.text("ALTER TABLE inventory ADD COLUMN teamviewer_id VARCHAR(120)"))
        visit_cols={c["name"] for c in inspector.get_columns("hardware_field_visits")} if inspector.has_table("hardware_field_visits") else set()
        if visit_cols and "report_group" not in visit_cols:
            conn.execute(db.text("ALTER TABLE hardware_field_visits ADD COLUMN report_group VARCHAR(40) DEFAULT 'AUTOPASS'"))


def migrate_inventory_validator_columns():
    with db.engine.begin() as conn:
        statements = [
            "ALTER TABLE inventory ADD COLUMN IF NOT EXISTS application VARCHAR(180)",
            "ALTER TABLE inventory ADD COLUMN IF NOT EXISTS bom_id VARCHAR(120)",
            "ALTER TABLE inventory ADD COLUMN IF NOT EXISTS bu_id VARCHAR(120)",
            "ALTER TABLE inventory ADD COLUMN IF NOT EXISTS validator_top_id VARCHAR(120)",
            "ALTER TABLE inventory ADD COLUMN IF NOT EXISTS software_version VARCHAR(120)",
        ]
        for statement in statements:
            conn.execute(db.text(statement))


@app.get("/admin/migrar-base-assets")
@manager_required
def migrar_base_assets():
    try:
        migrate_base_asset_columns()
        return jsonify({
            "ok": True,
            "mensagem": "Estrutura base_assets atualizada com sucesso.",
            "dados_preservados": True
        })
    except Exception as exc:
        return jsonify({"ok": False, "erro": str(exc)}), 500


@app.get("/admin/migrar-inventory-validator")
@manager_required
def migrar_inventory_validator():
    try:
        migrate_inventory_validator_columns()
        return jsonify({
            "ok": True,
            "mensagem": "Estrutura Inventory atualizada para Validador.",
            "dados_preservados": True
        })
    except Exception as exc:
        return jsonify({"ok": False, "erro": str(exc)}), 500



def sync_base_assets_1408(force=False):
    """Sincroniza a base detalhada das abas ATM, VALIDADOR, POS, TDI e BLOQUEIO."""
    source = DATA_DIR / "base_assets_1408.json"
    if not source.exists():
        return {"ok": False, "reason": "arquivo ausente", "inserted": 0, "updated": 0}

    rows = json.loads(source.read_text(encoding="utf-8"))

    # Atualiza automaticamente quando a estrutura/mapeamento da base muda.
    # As sentinelas evitam regravar milhares de linhas em todo cold start.
    if not force:
        sentinels = {
            "L01-PPQ-POS-13301": ("POS", "01 - AZUL", "METRO"),
            "L03-PSE-VAL-325-TMB": ("VALIDADOR", "03 - VERMELHA", "METRO"),
            "TDI|SAC|220|1500070008": ("TDI", "02 - VERDE", "METRO"),
            "BLOQ|09 - ESMERALDA|OSASCO|356406": ("BLOQUEIO", "09 - ESMERALDA", "VIA MOBILIDADE"),
        }
        current_ok = True
        for key, expected in sentinels.items():
            obj = BaseAsset.query.filter_by(asset_key=key).first()
            if not obj:
                current_ok = False
                break
            if (normalize(obj.equipment_type), normalize(obj.line), normalize(obj.company)) != tuple(normalize(x) for x in expected):
                current_ok = False
                break
        if current_ok:
            return {"ok": True, "reason": "base 1408-5 já sincronizada", "inserted": 0, "updated": 0}

    existing = {x.asset_key: x for x in BaseAsset.query.all() if x.asset_key}
    allowed = {
        "equipment_type", "description", "company", "station_code", "location_code",
        "line", "locality", "terminal_number", "serial", "qrcode_id", "top_id",
        "products", "model", "supplier", "transactions", "pix", "mount",
        "base_status", "application", "bom_id", "bu_id", "software_version",
        "quantity", "base_notes", "leasing_status", "contract_end",
        "installation_type", "installation_date"
    }

    inserted = updated = 0
    for row in rows:
        key = (row.get("asset_key") or "").strip()
        if not key:
            continue
        obj = existing.get(key)
        if obj is None:
            obj = BaseAsset(asset_key=key)
            db.session.add(obj)
            existing[key] = obj
            inserted += 1
        else:
            updated += 1
        for field in allowed:
            if field in row:
                setattr(obj, field, row.get(field))

    db.session.commit()
    _invalidate_expected_cache()
    return {"ok": True, "reason": "sincronizada", "inserted": inserted, "updated": updated, "total": len(rows)}


@app.get("/api/base/summary")
@manager_required
def api_base_summary():
    rows = (
        db.session.query(BaseAsset.equipment_type, func.count(BaseAsset.id))
        .filter(~func.upper(func.coalesce(BaseAsset.base_status, "")).like("%INATIVO%"))
        .group_by(BaseAsset.equipment_type)
        .all()
    )
    return jsonify({
        "ok": True,
        "base_version": BASE_DATA_VERSION,
        "active_by_type": {str(k or "ATM"): int(v or 0) for k, v in rows}
    })


@app.get("/admin/sincronizar-base-1408")
@manager_required
def admin_sync_base_1408():
    try:
        return jsonify(sync_base_assets_1408(force=True))
    except Exception as exc:
        db.session.rollback()
        return jsonify({"ok": False, "erro": str(exc)}), 500

def migrate_location_reference_columns():
    with db.engine.begin() as conn:
        conn.execute(db.text("""
            ALTER TABLE locations
            ADD COLUMN IF NOT EXISTS reference_latitude DOUBLE PRECISION
        """))

        conn.execute(db.text("""
            ALTER TABLE locations
            ADD COLUMN IF NOT EXISTS reference_longitude DOUBLE PRECISION
        """))

        conn.execute(db.text("""
            ALTER TABLE locations
            ADD COLUMN IF NOT EXISTS reference_source VARCHAR(120)
        """))

        conn.execute(db.text("""
            ALTER TABLE locations
            ADD COLUMN IF NOT EXISTS reference_updated_at TIMESTAMP
        """))

def migrate_location_reference_columns():
    with db.engine.begin() as conn:
        conn.execute(db.text("""
            ALTER TABLE locations
            ADD COLUMN IF NOT EXISTS reference_latitude DOUBLE PRECISION
        """))

        conn.execute(db.text("""
            ALTER TABLE locations
            ADD COLUMN IF NOT EXISTS reference_longitude DOUBLE PRECISION
        """))

        conn.execute(db.text("""
            ALTER TABLE locations
            ADD COLUMN IF NOT EXISTS reference_source VARCHAR(120)
        """))

        conn.execute(db.text("""
            ALTER TABLE locations
            ADD COLUMN IF NOT EXISTS reference_updated_at TIMESTAMP
        """))

@app.route("/debug-conciliar-estacoes")
@manager_required
def debug_conciliar_estacoes():
    import json
    import re
    import unicodedata
    import urllib.parse
    import urllib.request
    from difflib import SequenceMatcher

    base_url = "https://wfs.geosampa.prefeitura.sp.gov.br/geoserver/ows"

    def normalizar(texto):
        texto = texto or ""
        texto = unicodedata.normalize("NFKD", texto)
        texto = "".join(c for c in texto if not unicodedata.combining(c))
        texto = texto.upper().strip()

        # Remove siglas internas, com ou sem espaços ao redor do hífen.
        # Exemplos: JOD- JOAO DIAS, LUZ - LUZ, BTO - SAO BENTO.
        m = re.match(r"^([A-Z0-9]{1,4})\s*-\s*(.+)$", texto)
        if m:
            texto = m.group(2).strip()

        texto = texto.replace("–", "-").replace("—", "-").replace("-", " ")
        texto = re.sub(r"[^A-Z0-9 ]", " ", texto)
        texto = re.sub(r"\s+", " ", texto).strip()

        aliases = {
            "BARRA FUNDA": "PALMEIRAS BARRA FUNDA",
            "PALMEIRAS BARRA FUNDA": "PALMEIRAS BARRA FUNDA",
            "LAPA A": "LAPA LINHA 7",
            "LAPA B": "LAPA LINHA 8",
            "MANOEL FEIO": "ENGENHEIRO MANOEL FEIO",
            "JARDIM HELENA": "JARDIM HELENA VILA MARA",
            "JARDIM SAO PAULO": "AYRTON SENNA JARDIM SAO PAULO",
            "LIBERDADE": "JAPAO LIBERDADE",
            "PORTUGUESA TIETE": "PORTUGUESA TIETE",
            "BRESSER MOOCA": "BRESSER MOOCA",
            "CORINTHIANS ITAQUERA": "CORINTHIANS ITAQUERA",
            "GUILHERMINA ESPERANCA": "GUILHERMINA ESPERANCA",
            "SANTOS IMIGRANTES": "SANTOS IMIGRANTES",
            "USP LESTE": "USP LESTE",
            "GUARULHOS CECAP": "GUARULHOS CECAP",
            "AEROPORTO GUARULHOS": "AEROPORTO GUARULHOS",
            "JOAO DIAS": "JOAO DIAS",
            "INTERLAGOS": "PRIMAVERA INTERLAGOS",
            "MENDES": "MENDES BRUNO COVAS",
            "SANTO AMARO": "SANTO AMARO LINHA 9",
        }
        return aliases.get(texto, texto)

    def normalizar_linha(texto):
        texto = normalizar(texto)
        cores = {
            "01 AZUL": "AZUL", "1 AZUL": "AZUL",
            "02 VERDE": "VERDE", "2 VERDE": "VERDE",
            "03 VERMELHA": "VERMELHA", "3 VERMELHA": "VERMELHA",
            "04 AMARELA": "AMARELA", "4 AMARELA": "AMARELA",
            "05 LILAS": "LILAS", "5 LILAS": "LILAS",
            "06 LARANJA": "LARANJA", "6 LARANJA": "LARANJA",
            "07 RUBI": "RUBI", "7 RUBI": "RUBI",
            "08 DIAMANTE": "DIAMANTE", "8 DIAMANTE": "DIAMANTE",
            "09 ESMERALDA": "ESMERALDA", "9 ESMERALDA": "ESMERALDA",
            "10 TURQUESA": "TURQUESA",
            "11 CORAL": "CORAL",
            "12 SAFIRA": "SAFIRA",
            "13 JADE": "JADE",
            "15 PRATA": "PRATA",
            "17 OURO": "OURO",
        }
        return cores.get(texto, texto)

    def carregar_camada(nome):
        params = {
            "service": "WFS",
            "version": "1.0.0",
            "request": "GetFeature",
            "typeName": nome,
            "outputFormat": "application/json",
            "srsName": "EPSG:4326",
        }
        url = base_url + "?" + urllib.parse.urlencode(params)
        req = urllib.request.Request(url, headers={"User-Agent": "InventarioAutopass/1.0"})
        with urllib.request.urlopen(req, timeout=30) as response:
            return json.loads(response.read().decode("utf-8")).get("features", [])

    def preparar_estacoes():
        features = carregar_camada("estacao_metro") + carregar_camada("estacao_trem")
        saida = []
        for f in features:
            p = f.get("properties", {})
            g = f.get("geometry", {})
            coords = g.get("coordinates", [])
            if len(coords) < 2:
                continue
            saida.append({
                "nome": p.get("nm_estacao_metro_trem"),
                "nome_norm": normalizar(p.get("nm_estacao_metro_trem")),
                "linha": p.get("nm_linha_metro_trem"),
                "linha_norm": normalizar_linha(p.get("nm_linha_metro_trem")),
                "empresa": p.get("nm_empresa_metro_trem"),
                "situacao": p.get("tx_situacao_metro_trem"),
                "longitude": coords[0],
                "latitude": coords[1],
            })
        return saida

    def empresa_compativel(company, estacao):
        company = normalizar(company)
        linha = estacao["linha_norm"]

        if company == "METRO":
            return linha in {
                "AZUL", "VERDE", "VERMELHA", "AMARELA", "LILAS",
                "LARANJA", "PRATA", "OURO",
            }
        if company == "CPTM":
            return linha in {"RUBI", "TURQUESA", "CORAL", "SAFIRA", "JADE"}
        if company == "VIA MOBILIDADE":
            return linha in {"LILAS", "DIAMANTE", "ESMERALDA", "OURO"}
        return False

    def melhor_match(loc, estacoes):
        nome_loc = normalizar(loc.location)
        linha_loc = normalizar_linha(loc.line)

        candidatos = [
            e for e in estacoes
            if e["situacao"] == "OPERANDO" and empresa_compativel(loc.company, e)
        ]

        mesma_linha = [e for e in candidatos if e["linha_norm"] == linha_loc]
        if mesma_linha:
            candidatos = mesma_linha

        # Match exato/alias tem prioridade absoluta.
        for e in candidatos:
            if nome_loc == e["nome_norm"]:
                return {
                    "location_id": loc.id,
                    "company": loc.company,
                    "line": loc.line,
                    "location": loc.location,
                    "estacao_encontrada": e["nome"],
                    "linha_geosampa": e["linha"],
                    "empresa_geosampa": e["empresa"],
                    "latitude": e["latitude"],
                    "longitude": e["longitude"],
                    "score": 1.0,
                    "confianca": "ALTA",
                }

        melhor = None
        melhor_score = 0.0
        for e in candidatos:
            score = SequenceMatcher(None, nome_loc, e["nome_norm"]).ratio()
            if e["linha_norm"] == linha_loc:
                score += 0.08
            if score > melhor_score:
                melhor_score = score
                melhor = e

        if not melhor:
            return None

        if melhor_score >= 0.93:
            confianca = "ALTA"
        elif melhor_score >= 0.78:
            confianca = "REVISAR"
        else:
            confianca = "BAIXA"

        return {
            "location_id": loc.id,
            "company": loc.company,
            "line": loc.line,
            "location": loc.location,
            "estacao_encontrada": melhor["nome"],
            "linha_geosampa": melhor["linha"],
            "empresa_geosampa": melhor["empresa"],
            "latitude": melhor["latitude"],
            "longitude": melhor["longitude"],
            "score": round(min(melhor_score, 1.0), 3),
            "confianca": confianca,
        }

    try:
        estacoes = preparar_estacoes()
        localidades = (
            Location.query
            .filter(Location.company.in_(["METRO", "CPTM", "VIA MOBILIDADE"]))
            .order_by(Location.company, Location.line, Location.location)
            .all()
        )

        resultados = []
        for loc in localidades:
            match = melhor_match(loc, estacoes)
            if match:
                resultados.append(match)
            else:
                resultados.append({
                    "location_id": loc.id,
                    "company": loc.company,
                    "line": loc.line,
                    "location": loc.location,
                    "estacao_encontrada": None,
                    "latitude": None,
                    "longitude": None,
                    "score": 0,
                    "confianca": "NAO ENCONTRADA",
                })

        resumo = {
            "total_localidades": len(resultados),
            "alta": sum(1 for x in resultados if x["confianca"] == "ALTA"),
            "revisar": sum(1 for x in resultados if x["confianca"] == "REVISAR"),
            "baixa": sum(1 for x in resultados if x["confianca"] == "BAIXA"),
            "nao_encontrada": sum(1 for x in resultados if x["confianca"] == "NAO ENCONTRADA"),
        }

        return jsonify({"ok": True, "resumo": resumo, "resultados": resultados})
    except Exception as e:
        return jsonify({"ok": False, "erro": str(e)}), 500


def cleanup_v352_test_reference():
    """Remove somente a referência manual de teste LUZ/01-AZUL criada no mapa.
    Não remove a localidade, inventários, evidências nem referências de outras linhas.
    """
    rows = Location.query.filter(
        func.upper(func.coalesce(Location.location, "")).like("%LUZ%"),
        func.upper(func.coalesce(Location.line, "")).like("%01%AZUL%"),
        func.upper(func.coalesce(Location.reference_source, "")).like("%GESTOR%MAPA%")
    ).all()
    changed = 0
    for loc in rows:
        if loc.reference_latitude is None and loc.reference_longitude is None:
            continue
        loc.reference_latitude = None
        loc.reference_longitude = None
        loc.reference_source = "V35.2 · referência de teste removida"
        loc.reference_updated_at = datetime.utcnow()
        changed += 1
    if changed:
        db.session.commit()
    return changed




@app.get("/gestao-360")
@login_required
def management_360_page():
    """V47: visão gerencial integrada, preservando FIELD e IMPLANTAÇÃO como domínios distintos."""
    if session.get("role") not in ("manager", "manager_field"):
        return redirect(url_for("dashboard_landing"))
    return render_template("management_360.html", app_release=APP_RELEASE)


@app.get("/api/gestao-360/resumo")
@login_required
def management_360_summary_api():
    if session.get("role") not in ("manager", "manager_field"):
        return jsonify({"ok": False, "error": "Sem permissão para a Central 360."}), 403

    # FIELD: inventário, Recarga e Visão Panorâmica. EMV não pertence ao Field.
    inv_total = Inventory.query.count()
    inv_today = Inventory.query.filter(func.date(Inventory.created_at) == datetime.utcnow().date()).count()
    chip_rows = ChipSwap.query.all()
    chip_done = sum(1 for x in chip_rows if (x.status or "").upper().startswith("CONCLU"))
    chip_pending = max(len(chip_rows) - chip_done, 0)
    pan_total = PanoramaPoint.query.count()

    field = {
        "inventory_records": inv_total,
        "inventory_today": inv_today,
        "chip_total": len(chip_rows),
        "chip_done": chip_done,
        "chip_pending": chip_pending,
        "panorama_points": pan_total,
    }

    # IMPLANTAÇÃO: Visitas/RV + Troca de Chips EMV Trilhos.
    visits = HardwareFieldVisit.query.all()
    visit_finalized = sum(1 for x in visits if (x.status or "").upper() == "FINALIZADO")
    visit_pending = sum(1 for x in visits if "PEND" in (x.conclusion_status or "").upper())
    emv_rows = EmvChipSwap.query.all()
    emv_done = sum(1 for x in emv_rows if (x.status or "").upper().startswith("CONCLU"))
    implantation = {
        "visits_total": len(visits),
        "visits_finalized": visit_finalized,
        "visits_with_pending": visit_pending,
        "emv_total": len(emv_rows),
        "emv_done": emv_done,
        "emv_pending": max(len(emv_rows) - emv_done, 0),
    }
    td_total=TopDeskTicket.query.count(); td_open=TopDeskTicket.query.filter(~func.upper(func.coalesce(TopDeskTicket.status,'')).in_(['RESOLVIDO','FECHADO','CONCLUÍDO','CONCLUIDO'])).count()
    td_objects=db.session.query(TopDeskTicket.object_id).filter(TopDeskTicket.object_id.isnot(None),TopDeskTicket.object_id!='').distinct().count()
    topdesk={"tickets":td_total,"open":td_open,"objects":td_objects}
    return jsonify({"ok": True, "release": APP_RELEASE, "field": field, "implantation": implantation, "topdesk":topdesk})


@app.get("/api/busca-global")
@login_required
def global_search_api():
    """V48: busca transversal, preservando a identificação do domínio de origem."""
    q=(request.args.get("q") or "").strip()
    if len(q) < 2:
        return jsonify({"ok": True, "results": []})
    like=f"%{q}%"
    out=[]
    # Ativos / FIELD
    assets=BaseAsset.query.filter(db.or_(BaseAsset.asset_key.ilike(like), BaseAsset.terminal_number.ilike(like), BaseAsset.serial.ilike(like), BaseAsset.teamviewer_id.ilike(like), BaseAsset.locality.ilike(like))).limit(12).all()
    for a in assets:
        out.append({"type":"ATIVO","domain":"FIELD","title":a.asset_key or a.terminal_number or f"Ativo #{a.id}","subtitle":" · ".join(x for x in [a.equipment_type,a.company,a.line,a.locality] if x),"url":f"/dashboard/atm?asset_id={a.id}" if (a.equipment_type or '').upper()=='ATM' else f"/patrimonio?asset={a.id}"})
    # Localidades são compartilhadas, mas a tela 360 mantém os domínios separados.
    locs=Location.query.filter(db.or_(Location.location.ilike(like),Location.line.ilike(like),Location.company.ilike(like))).limit(10).all()
    for x in locs:
        out.append({"type":"LOCALIDADE","domain":"360","title":x.location,"subtitle":" · ".join(y for y in [x.company,x.line] if y),"url":f"/localidade-360/{x.id}"})
    # Pessoas / equipes
    users=User.query.filter(db.or_(User.name.ilike(like),User.username.ilike(like))).limit(8).all()
    for u in users:
        out.append({"type":"USUÁRIO","domain":"EQUIPE","title":u.name or u.username,"subtitle":u.username or "","url":"/equipes"})
    # RV / IMPLANTAÇÃO
    visits=HardwareFieldVisit.query.filter(db.or_(HardwareFieldVisit.report_code.ilike(like),HardwareFieldVisit.location_name.ilike(like),HardwareFieldVisit.client.ilike(like),HardwareFieldVisit.project.ilike(like))).limit(10).all()
    for v in visits:
        out.append({"type":"RELATÓRIO RV","domain":"IMPLANTAÇÃO","title":v.report_code or f"RV-{v.id:06d}","subtitle":" · ".join(y for y in [v.client,v.project,v.location_name] if y),"url":f"/implantacao-hardware/visita/{v.id}"})
    # V55: chamados entram na pesquisa transversal / Central 360.
    tickets=TopDeskTicket.query.filter(db.or_(TopDeskTicket.ticket_number.ilike(like),TopDeskTicket.object_id.ilike(like),TopDeskTicket.subcategory.ilike(like))).limit(8).all()
    for t in tickets:
        out.append({"type":"CHAMADO","domain":"GESTÃO","title":t.ticket_number,"subtitle":" · ".join(y for y in [t.object_id,t.subcategory,t.status] if y),"url":f"/topdesk?q={t.ticket_number}"})
    return jsonify({"ok":True,"query":q,"results":out[:25]})


@app.get("/api/gestao-360/alertas")
@login_required
def management_360_alerts_api():
    if session.get("role") not in ("manager","manager_field"):
        return jsonify({"ok":False,"error":"Sem permissão."}),403
    # FIELD
    official_path=DATA_DIR / "atm_official_082026.json"
    try:
        official_atms=json.loads(official_path.read_text(encoding="utf-8"))
    except Exception:
        official_atms=[]
    atm_without_tv=sum(1 for x in official_atms if not str(x.get("teamviewer_id") or "").strip())
    inv_div=Inventory.query.filter(func.coalesce(Inventory.divergence,"")!="").count()
    loc_pending=Location.query.filter(func.upper(func.coalesce(Location.survey_status,"PENDENTE"))!="CONCLUIDO").count()
    # IMPLANTAÇÃO
    visits=HardwareFieldVisit.query.all()
    rv_unsigned=sum(1 for v in visits if (v.status or '').upper()=='FINALIZADO' and not v.signature_file)
    rv_pending=sum(1 for v in visits if 'PEND' in (v.conclusion_status or '').upper())
    emv_rows=EmvChipSwap.query.all()
    emv_pending=sum(1 for x in emv_rows if not (x.status or '').upper().startswith('CONCLU'))
    alerts=[
      {"domain":"FIELD","severity":"ATENCAO","label":"ATMs sem ID TeamViewer","count":atm_without_tv,"url":"/dashboard/atm?teamviewer_missing=1"},
      {"domain":"FIELD","severity":"ATENCAO","label":"Divergências de inventário","count":inv_div,"url":"/dashboard/field"},
      {"domain":"FIELD","severity":"INFO","label":"Localidades não concluídas","count":loc_pending,"url":"/dashboard/field"},
      {"domain":"IMPLANTAÇÃO","severity":"ATENCAO","label":"RVs finalizados sem assinatura","count":rv_unsigned,"url":"/implantacao-hardware/dashboard"},
      {"domain":"IMPLANTAÇÃO","severity":"ATENCAO","label":"Visitas com pendências","count":rv_pending,"url":"/implantacao-hardware/dashboard"},
      {"domain":"IMPLANTAÇÃO","severity":"INFO","label":"EMV Trilhos pendentes","count":emv_pending,"url":"/implantacao-hardware/dashboard"},
    ]
    return jsonify({"ok":True,"alerts":[a for a in alerts if a["count"]>0]})


@app.get("/api/auditoria/recente")
@login_required
def recent_audit_api():
    if session.get("role") not in ("manager","manager_field"):
        return jsonify({"ok":False,"error":"Sem permissão."}),403
    rows=AuditEvent.query.order_by(AuditEvent.created_at.desc()).limit(20).all()
    user_ids={x.user_id for x in rows if x.user_id}
    users={u.id:u for u in User.query.filter(User.id.in_(user_ids)).all()} if user_ids else {}
    return jsonify({"ok":True,"events":[{"id":x.id,"type":x.event_type,"entity":x.entity_type,"entity_id":x.entity_id,"detail":x.detail,"created_at":x.created_at.isoformat() if x.created_at else None,"user_id":x.user_id,"user":users.get(x.user_id).name if users.get(x.user_id) else "Responsável não identificado","user_role":users.get(x.user_id).role if users.get(x.user_id) else ""} for x in rows]})


@app.get("/atividades")
@login_required
def activities_page():
    if session.get("role") not in ("technician", "technician_implantation", "manager", "manager_field"):
        return redirect(url_for("manager" if session.get("role") in ("consultation", "dispatcher") else "teams_page"))
    return render_template("activities.html")


@app.get("/api/atividades/resumo")
@login_required
def activities_summary_api():
    uid=session.get("user_id")
    role=session.get("role")
    activities=[]
    if role in ("technician","manager","manager_field"):
        inv_today=Inventory.query.filter(Inventory.technician_id==uid, func.date(Inventory.created_at)==datetime.utcnow().date()).count()
        swaps=ChipSwap.query.filter_by(technician_id=uid).all()
        chip_done=sum(1 for x in swaps if x.status=="CONCLUÍDO")
        pan=PanoramaPoint.query.filter_by(created_by=uid).count()
        candidates = [
          ("field.inventory", {"key":"inventory","title":"Inventário / Lançamento","href":"/tecnico","done":inv_today,"label":"lançamentos hoje"}),
          ("field.chip_recarga", {"key":"chips","title":"Troca de Chip Recarga","href":"/troca-chips","done":chip_done,"label":"concluídos"}),
          ("field.panorama", {"key":"panorama","title":"Visão Panorâmica","href":"/visao-panoramica","done":pan,"label":"pontos registrados"})]
        activities += [item for perm,item in candidates if _has_access(perm)]
    if role in ("manager","technician_implantation"):
        emv_done=EmvChipSwap.query.filter_by(status="CONCLUÍDA").count()
        
        if _has_access("implantation.emv"):
            activities.append({"key":"emv","title":"Troca de Chips EMV - Trilhos","href":"/troca-chips-emv","done":emv_done,"label":"concluídos"})
        done=HardwareFieldVisit.query.filter_by(status="FINALIZADO").count()
        
        if _has_access("implantation.visits") or _has_access("implantation.reports"):
            activities.append({"key":"implantation","title":"Implantação de Hardware","href":"/implantacao-hardware","done":done,"label":"visitas finalizadas"})
    return jsonify({"ok":True,"activities":activities})


@app.get("/troca-chips")
@login_required
def chip_swap_page():
    if not _has_access("field.chip_recarga"): abort(403)
    if session.get("role") not in ("manager", "manager_field", "technician", "consultation", "dispatcher"):
        return redirect(url_for("teams_page"))
    return render_template("chip_swap.html")

def _chip_swap_asset_label(a):
    # V39.7.4: TERMINAL é a identificação operacional principal do Validador de Recarga.
    return (f"Terminal {a.terminal_number}" if a.terminal_number else None) or a.description or a.asset_key or a.serial or f"Validador #{a.id}"


def _chip_swap_asset_matches_location(asset, loc):
    """V39.7.5: casa Validador de Recarga pela linha + estação/sigla.
    O estabelecimento/empresa da planilha não bloqueia o vínculo, pois a base histórica
    pode trazer METRO/CPTM enquanto a operação atual usa outro nome de empresa.
    """
    if not asset or not loc:
        return False
    if _canonical_equipment_type(asset.equipment_type) != "VALIDADOR":
        return False
    if _normalize_line_key(asset.line) != _normalize_line_key(loc.line):
        return False
    return _station_matches(asset.locality, asset.location_code or asset.station_code, loc.location)

_chip_swap_tables_ready = False
_chip_swap_payload_cache = {"at": 0.0, "data": None}
CHIP_SWAP_CACHE_TTL_SECONDS = 120

def _ensure_chip_swap_tables():
    """Verifica/cria as tabelas uma única vez por processo do Gunicorn."""
    global _chip_swap_tables_ready
    if _chip_swap_tables_ready:
        return
    ChipSwap.__table__.create(bind=db.engine, checkfirst=True)
    ChipSwapPhoto.__table__.create(bind=db.engine, checkfirst=True)
    # V39.7.9: adiciona resultado do teste sem exigir migration manual.
    try:
        cols = {c["name"] for c in db.inspect(db.engine).get_columns("chip_swaps")}
        with db.engine.begin() as conn:
            if "test_result" not in cols:
                conn.execute(text("ALTER TABLE chip_swaps ADD COLUMN test_result VARCHAR(80)"))
            if "test_notes" not in cols:
                conn.execute(text("ALTER TABLE chip_swaps ADD COLUMN test_notes TEXT"))
            if "completed_by_id" not in cols:
                conn.execute(text("ALTER TABLE chip_swaps ADD COLUMN completed_by_id INTEGER"))
    except Exception as exc:
        app.logger.warning("V39.7.9 chip swap migration: %s", exc)
    _chip_swap_tables_ready = True

def _invalidate_chip_swap_cache():
    _chip_swap_payload_cache["at"] = 0.0
    _chip_swap_payload_cache["data"] = None

def _chip_swap_locations_payload(force=False):
    """Monta o progresso da troca de chips sem fazer produto cartesiano localidade x ativo.

    V39.7.2:
    - DDL é verificado uma única vez por worker;
    - ativos são associados às localidades em uma única passagem, indexados por linha;
    - fotos e usuários são carregados em lote;
    - payload é reutilizado por poucos segundos e invalidado após gravação.
    """
    _ensure_chip_swap_tables()
    now = time.monotonic()
    if (not force and _chip_swap_payload_cache.get("data") is not None
            and now - float(_chip_swap_payload_cache.get("at") or 0) < CHIP_SWAP_CACHE_TTL_SECONDS):
        return _chip_swap_payload_cache["data"]

    locations = Location.query.order_by(Location.company, Location.line, Location.location).all()
    loc_by_id = {loc.id: loc for loc in locations}
    by_line = {}
    for loc in locations:
        by_line.setdefault(_normalize_line_key(loc.line), []).append(loc)

    assets_by_loc = {loc.id: [] for loc in locations}

    # V56-A.4 HOTFIX3: garante que a base detalhada de validadores esteja realmente
    # sincronizada no banco persistente. Versões anteriores podiam considerar a base
    # 1408-5 "já sincronizada" olhando apenas sentinelas de POS/TDI/Bloqueio, deixando
    # os Validadores ausentes em bancos antigos. Isso gerava contador previsto > 0,
    # mas lista detalhada vazia na Troca de Chip Recarga.
    validator_sentinel = BaseAsset.query.filter_by(asset_key="L03-PSE-VAL-325-TMB").first()
    if validator_sentinel is None:
        # V60 REV2 PERFORMANCE: nunca executa sincronização pesada dentro de uma requisição.
        app.logger.warning("V60 REV2: base detalhada de validadores sem sentinela; sincronização deve ocorrer fora do request.")

    validator_assets = BaseAsset.query.filter(
        func.upper(func.coalesce(BaseAsset.equipment_type, '')).like('%VALID%')
    ).all()

    for asset in validator_assets:
        if _canonical_equipment_type(asset.equipment_type) != "VALIDADOR" or "INATIVO" in normalize(asset.base_status) or "FORA DO ESCOPO" in normalize(asset.base_status):
            continue
        candidates = by_line.get(_normalize_line_key(asset.line), ())
        if not candidates:
            continue
        for loc in candidates:
            if _chip_swap_asset_matches_location(asset, loc):
                assets_by_loc[loc.id].append(asset)
                break

    swaps_list = ChipSwap.query.all()
    swaps = {(x.location_id, x.base_asset_id): x for x in swaps_list}
    swap_ids = [x.id for x in swaps_list]
    photo_map = {}
    if swap_ids:
        photos = ChipSwapPhoto.query.filter(
            ChipSwapPhoto.chip_swap_id.in_(swap_ids)
        ).order_by(ChipSwapPhoto.created_at).all()
        for ph in photos:
            photo_map.setdefault(ph.chip_swap_id, []).append(ph)

    user_ids = {uid for x in swaps_list for uid in (x.technician_id, getattr(x,"completed_by_id",None)) if uid}
    users = {u.id: u for u in User.query.filter(User.id.in_(user_ids)).all()} if user_ids else {}

    # V66 REV1 PERFORMANCE: a base operacional da Recarga é carregada uma única
    # vez. Na V66 _op_active_map() era chamado dentro do loop de cada validador,
    # provocando ~1 SELECT por equipamento (mais de 560 queries/request).
    op_active = _op_active_map("recarga")

    rows = []
    for loc in locations:
        assets = assets_by_loc.get(loc.id, [])
        if not assets and int(loc.expected_validator or 0) <= 0:
            continue
        items = []
        for a in assets:
            sw = swaps.get((loc.id, a.id))
            photos = photo_map.get(sw.id, []) if sw else []
            op_item=op_active.get(str(a.terminal_number or ""))
            status = sw.status if sw else ((op_item.desired_status if op_item else None) or "PENDENTE")
            if photos and status != "CONCLUÍDA":
                status = "CONCLUÍDA"
            tech = users.get(sw.technician_id) if sw else None
            completed_by = users.get(getattr(sw,"completed_by_id",None)) if sw else None
            items.append({
                "base_asset_id": a.id,
                "label": _chip_swap_asset_label(a),
                "serial": a.serial or "",
                "model": a.model or "",
                "status": status,
                "swap_id": sw.id if sw else None,
                "technician": tech.name if tech else "",
                "completed_by": completed_by.name if completed_by else (tech.name if sw and sw.completed_at and tech else ""),
                "completed_by_role": completed_by.role if completed_by else (tech.role if sw and sw.completed_at and tech else ""),
                "completed_at": sw.completed_at.isoformat()+"Z" if sw and sw.completed_at else None,
                "photo_count": len(photos),
                "notes": sw.notes if sw else "",
                "test_result": sw.test_result if sw else "",
                "test_notes": sw.test_notes if sw else "",
                "photos": [{"id": ph.id, "url": "/uploads/"+ph.stored_name, "thumb_url": "/uploads/"+ph.stored_name+"?thumb=1", "name": ph.original_name} for ph in photos],
            })
        # V56-A.4 HOTFIX3: a lista detalhada é a fonte principal. O contador legado
        # expected_validator fica apenas como fallback quando ainda não há ativos detalhados.
        total = len(items) if items else int(loc.expected_validator or 0)
        concluded = sum(1 for i in items if i["status"] == "CONCLUÍDA")
        progress = sum(1 for i in items if i["status"] == "EM ANDAMENTO")
        pending = max(total - concluded - progress, 0)
        rows.append({
            "id": loc.id, "company": loc.company, "line": loc.line, "location": loc.location,
            "reference_latitude": loc.reference_latitude, "reference_longitude": loc.reference_longitude,
            "total": total, "concluded": concluded, "in_progress": progress, "pending": pending,
            "percent": round((concluded/total*100), 1) if total else 0, "validators": items,
        })

    _chip_swap_payload_cache["at"] = now
    _chip_swap_payload_cache["data"] = rows
    return rows

@app.get("/api/chip-swaps")
@login_required
def chip_swap_list_api():
    return jsonify({"ok": True, "locations": _chip_swap_locations_payload()})

@app.post("/api/chip-swaps/<int:location_id>/<int:base_asset_id>")
@field_required
def chip_swap_save_api(location_id, base_asset_id):
    if _activity_request_too_large(): return jsonify({"ok":False,"error":f"Envio excede {_ACTIVITY_REQUEST_MAX_MB} MB. Envie menos fotos por vez."}),413
    _ensure_chip_swap_tables()
    loc = db.session.get(Location, location_id)
    asset = db.session.get(BaseAsset, base_asset_id)
    if not loc or not asset or _canonical_equipment_type(asset.equipment_type) != "VALIDADOR" or not _chip_swap_asset_matches_location(asset, loc):
        return jsonify({"ok": False, "error": "Validador de recarga não encontrado nesta localidade."}), 404
    lat = _optional_float(request.form.get("latitude"))
    lon = _optional_float(request.form.get("longitude"))
    acc = _optional_float(request.form.get("gps_accuracy"))
    sw = ChipSwap.query.filter_by(location_id=location_id, base_asset_id=base_asset_id).first()
    if sw and (sw.status or "").upper().replace("CONCLUIDA","CONCLUÍDA") == "CONCLUÍDA" and session.get("role") in ("technician", "technician_implantation"):
        return jsonify({"ok": False, "error": "Registro concluído e bloqueado. Solicite ao Gestor/ADM a reabertura para EM ANDAMENTO."}), 409
    if not sw:
        sw = ChipSwap(location_id=location_id, base_asset_id=base_asset_id, technician_id=session["user_id"], status="EM ANDAMENTO", started_at=datetime.utcnow())
        db.session.add(sw)
        db.session.flush()
    sw.technician_id = session["user_id"]
    sw.notes = (request.form.get("notes") or sw.notes or "").strip()
    test_result = (request.form.get("test_result") or "").strip()
    allowed_results = {"TESTADO_OK", "TESTADO_COM_DEFEITO", "NAO_FOI_POSSIVEL_TESTAR", "EQUIPAMENTO_INOPERANTE", "OUTRO"}
    is_admin = session.get("role") == "manager"
    if not is_admin and test_result not in allowed_results:
        return jsonify({"ok": False, "error": "Informe o resultado do teste após a troca."}), 400
    test_notes = (request.form.get("test_notes") or "").strip()
    if not is_admin and test_result != "TESTADO_OK" and not test_notes and not sw.notes:
        return jsonify({"ok": False, "error": "Para resultado com pendência, informe uma observação."}), 400
    if test_result in allowed_results:
        sw.test_result = test_result
    sw.test_notes = test_notes or None
    sw.latitude = lat; sw.longitude = lon; sw.gps_accuracy = acc; sw.updated_at = datetime.utcnow()
    files = [f for f in request.files.getlist("photos") if f and f.filename]
    for f in files:
        safe = secure_filename(f.filename) or f"chip_{secrets.token_hex(4)}.jpg"
        stored = f"chip_{sw.id}_{secrets.token_hex(6)}_{safe}"
        stored = _store_uploaded_file(f, "chip-swaps", stored, f.mimetype or "application/octet-stream")
        db.session.add(ChipSwapPhoto(chip_swap_id=sw.id, original_name=f.filename, stored_name=stored, mime_type=f.mimetype, uploaded_by=session["user_id"]))
    db.session.flush()
    photo_count = ChipSwapPhoto.query.filter_by(chip_swap_id=sw.id).count()
    if photo_count > 0:
        sw.status = "CONCLUÍDA"; sw.completed_at = sw.completed_at or datetime.utcnow(); sw.completed_by_id = session.get("user_id")
    else:
        sw.status = "EM ANDAMENTO"
    db.session.add(AuditEvent(user_id=session.get("user_id"), event_type="CHIP_SWAP_UPDATE", entity_type="base_asset", entity_id=str(base_asset_id), detail=f"{loc.location} · {_chip_swap_asset_label(asset)} · {sw.status} · teste {sw.test_result or '—'} · {photo_count} foto(s)"))
    db.session.commit()
    _invalidate_chip_swap_cache()
    return jsonify({"ok": True, "status": sw.status, "photo_count": photo_count})



@app.post("/api/chip-swaps/<int:location_id>/<int:base_asset_id>/admin-status")
@dashboard_required
def chip_swap_admin_status_api(location_id, base_asset_id):
    if session.get("role") not in ("manager", "manager_field"):
        return jsonify({"ok":False,"error":"Alteração administrativa restrita ao ADM/Gestor."}),403
    loc=db.session.get(Location,location_id); asset=db.session.get(BaseAsset,base_asset_id)
    if not loc or not asset or _canonical_equipment_type(asset.equipment_type)!="VALIDADOR" or not _chip_swap_asset_matches_location(asset,loc):
        return jsonify({"ok":False,"error":"Validador não encontrado nesta localidade."}),404
    data=request.get_json(silent=True) or {}; new=(data.get("status") or "").strip().upper().replace("CONCLUIDA","CONCLUÍDA")
    reason=(data.get("reason") or "").strip()
    if new not in {"PENDENTE","EM ANDAMENTO","CONCLUÍDA"}:
        return jsonify({"ok":False,"error":"Status administrativo inválido."}),400
    if not reason:
        return jsonify({"ok":False,"error":"Informe o motivo da alteração administrativa."}),400
    sw=ChipSwap.query.filter_by(location_id=location_id,base_asset_id=base_asset_id).first()
    old_status=sw.status if sw else "PENDENTE"
    if not sw:
        sw=ChipSwap(location_id=location_id,base_asset_id=base_asset_id,technician_id=session["user_id"],status=new,started_at=datetime.utcnow())
        db.session.add(sw)
    sw.status=new; sw.updated_at=datetime.utcnow()
    if new=="CONCLUÍDA": sw.completed_at=sw.completed_at or datetime.utcnow()
    else: sw.completed_at=None
    db.session.add(AuditEvent(user_id=session.get("user_id"),event_type="CHIP_SWAP_ADMIN_STATUS",entity_type="base_asset",entity_id=str(base_asset_id),detail=f"{loc.location} · {_chip_swap_asset_label(asset)} · {old_status} -> {new} · motivo: {reason}"))
    db.session.commit(); _invalidate_chip_swap_cache()
    return jsonify({"ok":True,"status":new})


def _delete_stored_media(stored_name):
    if not stored_name:
        return
    try:
        if stored_name.startswith("r2__"):
            key=stored_name[4:]
            if key and _r2_available():
                r2_client().delete_object(Bucket=os.environ["R2_BUCKET_NAME"],Key=key)
        else:
            fp=UPLOAD_DIR/stored_name
            if fp.exists():
                fp.unlink()
    except Exception:
        app.logger.exception("Falha ao excluir mídia %s", stored_name)

@app.delete("/api/chip-swaps/photos/<int:photo_id>")
@field_required
def chip_swap_delete_photo_api(photo_id):
    ph=db.session.get(ChipSwapPhoto,photo_id)
    if not ph:
        return jsonify({"ok":False,"error":"Foto não encontrada."}),404
    sw=db.session.get(ChipSwap,ph.chip_swap_id)
    if not sw:
        return jsonify({"ok":False,"error":"Troca de chip não encontrada."}),404
    old_name=ph.original_name
    _delete_stored_media(ph.stored_name)
    db.session.delete(ph)
    db.session.flush()
    remaining=ChipSwapPhoto.query.filter_by(chip_swap_id=sw.id).count()
    sw.status="CONCLUÍDA" if remaining else "EM ANDAMENTO"
    if not remaining:
        sw.completed_at=None
    sw.updated_at=datetime.utcnow()
    db.session.add(AuditEvent(user_id=session.get("user_id"),event_type="CHIP_SWAP_PHOTO_DELETE",entity_type="chip_swap",entity_id=str(sw.id),detail=f"Foto {old_name} excluída · {remaining} restante(s)"))
    db.session.commit();_invalidate_chip_swap_cache()
    return jsonify({"ok":True,"remaining":remaining,"status":sw.status})

@app.post("/api/chip-swaps/photos/<int:photo_id>/replace")
@field_required
def chip_swap_replace_photo_api(photo_id):
    ph=db.session.get(ChipSwapPhoto,photo_id)
    if not ph:
        return jsonify({"ok":False,"error":"Foto não encontrada."}),404
    sw=db.session.get(ChipSwap,ph.chip_swap_id)
    f=request.files.get("photo")
    if not f or not f.filename:
        return jsonify({"ok":False,"error":"Selecione a nova foto."}),400
    safe=secure_filename(f.filename) or f"chip_{secrets.token_hex(4)}.jpg"
    stored=f"chip_{sw.id}_{secrets.token_hex(6)}_{safe}"
    try:
        if _uploaded_file_size(f) == 0:
            return jsonify({"ok":False,"error":"Arquivo vazio."}),400
        stored=_store_uploaded_file(f,"chip-swaps",stored,f.mimetype or "application/octet-stream")
        old_stored=ph.stored_name
        ph.original_name=f.filename;ph.stored_name=stored;ph.mime_type=f.mimetype;ph.uploaded_by=session["user_id"];ph.created_at=datetime.utcnow()
        sw.status="CONCLUÍDA";sw.completed_at=sw.completed_at or datetime.utcnow();sw.updated_at=datetime.utcnow()
        db.session.add(AuditEvent(user_id=session.get("user_id"),event_type="CHIP_SWAP_PHOTO_REPLACE",entity_type="chip_swap",entity_id=str(sw.id),detail=f"Foto substituída por {f.filename}"))
        db.session.commit();_delete_stored_media(old_stored);_invalidate_chip_swap_cache()
        return jsonify({"ok":True,"photo_id":ph.id,"status":sw.status})
    except Exception as exc:
        db.session.rollback();_delete_stored_media(stored)
        app.logger.exception("Falha ao substituir foto da troca %s", sw.id if sw else "—")
        return jsonify({"ok":False,"error":"Não foi possível substituir a foto.","detail":str(exc)[:160]}),500

@app.post("/api/chip-swaps/<int:location_id>/new-asset")
@field_required
def chip_swap_new_asset_api(location_id):
    loc = db.session.get(Location, location_id)
    if not loc:
        return jsonify({"ok": False, "error": "Localidade não encontrada."}), 404
    terminal = (request.form.get("terminal_number") or "").strip()
    if not terminal:
        return jsonify({"ok": False, "error": "Informe o terminal do Validador de Recarga."}), 400
    existing = BaseAsset.query.filter(
        func.upper(func.coalesce(BaseAsset.equipment_type, '')).like('%VALID%'),
        func.upper(func.coalesce(BaseAsset.line, '')) == (loc.line or '').upper(),
        func.upper(func.coalesce(BaseAsset.terminal_number, '')) == terminal.upper()
    ).first()
    if existing and _chip_swap_asset_matches_location(existing, loc):
        return jsonify({"ok": True, "base_asset_id": existing.id, "existing": True})
    asset = BaseAsset(
        company=loc.company, line=loc.line, locality=loc.location, equipment_type="VALIDADOR",
        terminal_number=terminal, asset_key=f"FIELD-{loc.id}-{terminal}-{secrets.token_hex(3)}", serial=(request.form.get("serial") or "").strip() or None,
        model=(request.form.get("model") or "").strip() or None, base_status="ATIVO",
        base_notes="Cadastrado em campo pela atividade Troca de Chips V39.7.10"
    )
    db.session.add(asset); db.session.flush()
    db.session.add(AuditEvent(user_id=session.get("user_id"), event_type="CHIP_SWAP_NEW_ASSET", entity_type="base_asset", entity_id=str(asset.id), detail=f"{loc.location} · Terminal {terminal} · cadastrado em campo"))
    db.session.commit(); _invalidate_chip_swap_cache()
    return jsonify({"ok": True, "base_asset_id": asset.id, "existing": False})


def _chip_operation_name(company):
    t = normalize(company)
    if "CPTM" in t: return "CPTM"
    if "METRO" in t: return "Metrô"
    if "VIA MOBILIDADE" in t: return "Via Mobilidade"
    if "VIAQUATRO" in t or "VIA QUATRO" in t: return "ViaQuatro"
    return company or "Outros"

@app.get("/api/chip-swaps/dashboard")
@login_required
def chip_swap_dashboard_api():
    rows = _chip_swap_locations_payload()
    total = sum(x["total"] for x in rows); done = sum(x["concluded"] for x in rows)
    progress = sum(x["in_progress"] for x in rows); pending = sum(x["pending"] for x in rows)

    swaps = ChipSwap.query.order_by(ChipSwap.updated_at.desc()).all()
    user_ids = {sw.technician_id for sw in swaps if sw.technician_id}
    users = {u.id: u for u in User.query.filter(User.id.in_(user_ids)).all()} if user_ids else {}
    tech = {}
    for sw in swaps:
        u = users.get(sw.technician_id); name = u.name if u else "—"
        t = tech.setdefault(name, {"name": name, "concluded": 0, "in_progress": 0, "total": 0})
        t["total"] += 1
        t["concluded"] += 1 if sw.status == "CONCLUÍDA" else 0
        t["in_progress"] += 1 if sw.status == "EM ANDAMENTO" else 0
    result_labels = {
        "TESTADO_OK": "Testado - OK",
        "TESTADO_COM_DEFEITO": "Testado - com defeito",
        "NAO_FOI_POSSIVEL_TESTAR": "Não foi possível testar",
        "EQUIPAMENTO_INOPERANTE": "Equipamento inoperante",
        "OUTRO": "Outro",
        "SEM_RESULTADO": "Sem resultado",
    }
    result_counts = {k: 0 for k in result_labels}
    technical_pending = []
    for x in rows:
        for v in x.get("validators", []):
            if not v.get("swap_id"):
                continue
            r = v.get("test_result") or "SEM_RESULTADO"
            result_counts[r] = result_counts.get(r, 0) + 1
            if r != "TESTADO_OK":
                technical_pending.append({
                    "operation": _chip_operation_name(x.get("company")),
                    "company": x.get("company"), "line": x.get("line"), "location": x.get("location"),
                    "terminal": v.get("label"), "base_asset_id": v.get("base_asset_id"),
                    "technician": v.get("technician"), "completed_at": v.get("completed_at"),
                    "test_result": r, "test_result_label": result_labels.get(r, "Sem resultado"),
                    "notes": v.get("test_notes") or v.get("notes") or "", "photo_count": v.get("photo_count", 0),
                })
    return jsonify({
        "ok": True,
        "summary": {"total": total, "concluded": done, "in_progress": progress, "pending": pending, "percent": round(done/total*100, 1) if total else 0},
        "test_results": {"counts": result_counts, "labels": result_labels, "technical_pending": len(technical_pending)},
        "technical_pending": technical_pending,
        "locations": rows,
        "technicians": sorted(tech.values(), key=lambda x: (-x["concluded"], x["name"])),
    })


@app.get("/api/chip-swaps/export.xlsx")
@login_required
def chip_swap_export_xlsx():
    rows = _chip_swap_locations_payload()
    operation = (request.args.get("operation") or "").strip()
    company = (request.args.get("company") or "").strip()
    line = (request.args.get("line") or "").strip()
    location = (request.args.get("location") or "").strip()
    test_result = (request.args.get("test_result") or "").strip()
    pending_only = (request.args.get("pending_only") or "").strip().lower() in ("1", "true", "yes")
    rows=[x for x in rows if (not operation or _chip_operation_name(x["company"])==operation) and (not company or x["company"]==company) and (not line or x["line"]==line) and (not location or x["location"]==location)]
    total=sum(x["total"] for x in rows); done=sum(x["concluded"] for x in rows); prog=sum(x["in_progress"] for x in rows); pend=sum(x["pending"] for x in rows)
    wb=Workbook(); ws=wb.active; ws.title="Resumo"
    ws.append(["Troca de Chips - Validadores de Recarga"]); ws["A1"].font=Font(bold=True,size=14)
    ws.append(["Operação",operation or "Todos"]); ws.append(["Empresa",company or "Todas"]); ws.append(["Linha",line or "Todas"]); ws.append(["Localidade",location or "Todas"]); ws.append([])
    ws.append(["Total previsto","Concluídos","Em andamento","Pendentes","Progresso %"]); ws.append([total,done,prog,pend,round(done/total*100,1) if total else 0])
    det=wb.create_sheet("Detalhamento"); det.append(["Operação","Empresa","Linha","Localidade","Terminal / ativo","Modelo","Série","Status","Resultado pós-troca","Observação do teste","Técnico","Data/hora conclusão","Fotos"])
    for x in rows:
        for v in x.get("validators",[]):
            vr = v.get("test_result") or ""
            if test_result and vr != test_result:
                continue
            if pending_only and (not v.get("swap_id") or vr == "TESTADO_OK"):
                continue
            det.append([_chip_operation_name(x["company"]),x["company"],x["line"],x["location"],v.get("label") or v.get("base_asset_id"),v.get("model","") ,v.get("serial","") ,v.get("status","PENDENTE"),vr,v.get("test_notes") or v.get("notes") or "",v.get("technician","") ,v.get("completed_at") or "",v.get("photo_count",0)])
    for sh in wb.worksheets:
        for cell in sh[1]: cell.font=Font(bold=True)
        for col in range(1,sh.max_column+1): sh.column_dimensions[get_column_letter(col)].width=min(42,max(12,max((len(str(sh.cell(r,col).value or "")) for r in range(1,sh.max_row+1)),default=12)+2))
    bio=io.BytesIO(); wb.save(bio); bio.seek(0)
    return send_file(bio,as_attachment=True,download_name="troca_chips.xlsx",mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")



# V63 REV2 — importador padrão para EMV, Recarga e Garagem.
_OP_MODULES = {"emv":"EMV", "recarga":"RECARGA", "garagem":"GARAGEM"}
_OP_STATUS_ALIASES = {
    "PENDENTE":"PENDENTE", "PENDING":"PENDENTE",
    "EM ANDAMENTO":"EM ANDAMENTO", "ANDAMENTO":"EM ANDAMENTO",
    "CONCLUIDA":"CONCLUÍDA", "CONCLUIDO":"CONCLUÍDA", "CONCLUÍDA":"CONCLUÍDA", "CONCLUÍDO":"CONCLUÍDA",
}

def _op_clean_terminal(value):
    if value is None: return ""
    if isinstance(value, float) and value.is_integer(): value=int(value)
    text_value=str(value).strip()
    if re.fullmatch(r"\d+\.0", text_value): text_value=text_value[:-2]
    return text_value[:120]

def _op_header(value):
    value=normalize(str(value or "")).lower().replace("-"," ").replace("_"," ")
    value=" ".join(value.split())
    aliases={"empresa":"company","terminal":"terminal","estacao":"station","localidade":"station","linha":"line","status":"status"}
    return aliases.get(value,value)

def _op_parse_upload(file_storage):
    if not file_storage or not getattr(file_storage,"filename",""):
        raise ValueError("Selecione uma planilha .xlsx.")
    if not str(file_storage.filename).lower().endswith(".xlsx"):
        raise ValueError("Use o modelo Excel .xlsx.")
    try:
        wb=load_workbook(file_storage,read_only=True,data_only=True)
        ws=wb.active
        raw_headers=[_op_header(x) for x in next(ws.iter_rows(values_only=True))]
    except Exception as exc:
        raise ValueError(f"Não foi possível ler a planilha: {exc}")
    required={"company","terminal","station","line","status"}
    missing=required-set(raw_headers)
    if missing:
        wb.close(); raise ValueError("Colunas obrigatórias: empresa, terminal, estação, linha, status.")
    seen={}; errors=[]
    for row_no,vals in enumerate(ws.iter_rows(min_row=2,values_only=True),start=2):
        d=dict(zip(raw_headers,vals)); terminal=_op_clean_terminal(d.get("terminal"))
        if not terminal and not any(v not in (None,"") for v in vals): continue
        if not terminal:
            errors.append(f"Linha {row_no}: terminal vazio."); continue
        status_key=normalize(str(d.get("status") or "PENDENTE")).upper().strip()
        status=_OP_STATUS_ALIASES.get(status_key)
        if not status:
            errors.append(f"Linha {row_no}: status inválido '{d.get('status')}'."); continue
        item={"company":" ".join(str(d.get("company") or "").strip().split()),"terminal":terminal,"station":" ".join(str(d.get("station") or "").strip().split()),"line":" ".join(str(d.get("line") or "").strip().split()),"status":status,"row":row_no}
        if not item["company"] or not item["station"] or not item["line"]:
            errors.append(f"Linha {row_no}: empresa, estação e linha são obrigatórias."); continue
        if terminal in seen:
            errors.append(f"Linha {row_no}: terminal {terminal} duplicado na planilha."); continue
        seen[terminal]=item
    wb.close()
    if errors: raise ValueError(" | ".join(errors[:20]) + (f" | +{len(errors)-20} erro(s)" if len(errors)>20 else ""))
    if not seen: raise ValueError("A planilha não contém registros válidos.")
    return list(seen.values())

def _op_rows(module):
    mod=_OP_MODULES.get(module,module.upper())
    return OperationalBaseItem.query.filter_by(module=mod).all()

def _op_active_map(module):
    mod=_OP_MODULES.get(module,module.upper())
    return {x.terminal:x for x in OperationalBaseItem.query.filter_by(module=mod,active=True).all()}

def _op_current_status(module, terminal):
    if module=="emv":
        sw=EmvChipSwap.query.filter_by(terminal=terminal).first(); return (sw.status if sw else "PENDENTE")
    if module=="garagem":
        base=GarageChipBase.query.filter_by(terminal=terminal).first()
        sw=GarageChipSwap.query.filter_by(base_id=base.id).first() if base else None
        return sw.status if sw else "PENDENTE"
    asset=BaseAsset.query.filter_by(terminal_number=terminal).first()
    sw=ChipSwap.query.filter_by(base_asset_id=asset.id).first() if asset else None
    return sw.status if sw else "PENDENTE"

def _op_existing_terminals(module):
    if module=="emv": return {str(x.get("terminal") or "") for x in _v41_emv_rows() if x.get("terminal")}
    if module=="garagem": return {x.terminal for x in GarageChipBase.query.all()}
    return {str(x.terminal_number or "") for x in BaseAsset.query.filter(func.upper(func.coalesce(BaseAsset.equipment_type,'')).like('%VALID%')).all() if x.terminal_number}

def _op_preview(module, rows):
    imported={x["terminal"]:x for x in rows}; existing=_op_existing_terminals(module)
    new=sum(1 for t in imported if t not in existing); reopened=0; already_pending=0; concluded=0
    for t,x in imported.items():
        st=normalize(_op_current_status(module,t)).upper()
        if x["status"]=="PENDENTE":
            if st=="CONCLUIDA": reopened+=1
            elif st=="PENDENTE": already_pending+=1
        elif x["status"]=="CONCLUÍDA": concluded+=1
    stale_pending=sum(1 for t in existing-imported.keys() if normalize(_op_current_status(module,t)).upper() in ("PENDENTE","EM ANDAMENTO"))
    return {"rows":len(rows),"new":new,"existing":len(rows)-new,"reopened":reopened,"already_pending":already_pending,"marked_concluded":concluded,"outside_scope_pending":stale_pending}

def _op_upsert_item(module,row,user_id):
    mod=_OP_MODULES[module]; item=OperationalBaseItem.query.filter_by(module=mod,terminal=row["terminal"]).first()
    if not item:
        item=OperationalBaseItem(module=mod,terminal=row["terminal"],created_at=datetime.utcnow()); db.session.add(item)
    item.company=row["company"]; item.station=row["station"]; item.line=row["line"]; item.desired_status=row["status"]; item.active=True; item.updated_by=user_id; item.updated_at=datetime.utcnow()
    return item

def _op_reopen_swap(sw, new_status, module, terminal, user_id):
    if not sw: return False
    old=(sw.status or "PENDENTE").upper().replace("CONCLUIDA","CONCLUÍDA")
    changed=old!=new_status
    if changed:
        before={"status":sw.status,"completed_at":sw.completed_at.isoformat() if getattr(sw,"completed_at",None) else None,"test_result":getattr(sw,"test_result",None)}
        sw.status=new_status; sw.updated_at=datetime.utcnow()
        if new_status!="CONCLUÍDA": sw.completed_at=None
        db.session.add(AuditEvent(user_id=user_id,event_type="OP_BASE_STATUS_SYNC",entity_type=module,entity_id=str(getattr(sw,"id","")),detail=json.dumps({"terminal":terminal,"before":before,"after":new_status,"source":"IMPORTADOR V63 REV2"},ensure_ascii=False)))
    return changed

def _op_apply(module, rows, user_id):
    imported={x["terminal"]:x for x in rows}; preview=_op_preview(module,rows); now=datetime.utcnow(); mod=_OP_MODULES[module]
    # O arquivo representa o recorte operacional vigente. Ausentes nunca são apagados.
    for item in OperationalBaseItem.query.filter_by(module=mod,active=True).all():
        if item.terminal not in imported: item.active=False; item.updated_at=now; item.updated_by=user_id
    for row in rows: _op_upsert_item(module,row,user_id)

    if module=="emv":
        for row in rows:
            sw=EmvChipSwap.query.filter_by(terminal=row["terminal"]).first()
            _op_reopen_swap(sw,row["status"],"emv_chip_swap",row["terminal"],user_id)
        # Itens legados pendentes ausentes ficam fora do escopo por filtragem da base operacional.
        global _emv_base_rows_cache,_emv_base_by_terminal_cache
        _emv_base_rows_cache=None; _emv_base_by_terminal_cache=None; _v63_invalidate_emv_cache()

    elif module=="garagem":
        bases={b.terminal:b for b in GarageChipBase.query.all()}
        for row in rows:
            b=bases.get(row["terminal"])
            if not b:
                b=GarageChipBase(company=row["company"],terminal=row["terminal"],active=True,sam_type="NÃO MIGRADO"); db.session.add(b); db.session.flush(); bases[b.terminal]=b
            b.company=row["company"]; b.active=True
            sw=GarageChipSwap.query.filter_by(base_id=b.id).first(); _op_reopen_swap(sw,row["status"],"garage_chip_swap",row["terminal"],user_id)
        # Pendentes antigos fora do arquivo saem da operação, sem exclusão.
        for terminal,b in bases.items():
            if terminal in imported: continue
            sw=GarageChipSwap.query.filter_by(base_id=b.id).first(); st=normalize(sw.status if sw else "PENDENTE").upper()
            if st in ("PENDENTE","EM ANDAMENTO"): b.active=False
        _v63_invalidate_garage_cache()

    else: # recarga
        assets={str(a.terminal_number or ""):a for a in BaseAsset.query.filter(func.upper(func.coalesce(BaseAsset.equipment_type,'')).like('%VALID%')).all() if a.terminal_number}
        for row in rows:
            loc=Location.query.filter_by(company=row["company"],line=row["line"],location=row["station"]).first()
            if not loc:
                loc=Location(company=row["company"],line=row["line"],location=row["station"],base_status="IMPORTADO V63 REV2"); db.session.add(loc); db.session.flush()
            a=assets.get(row["terminal"])
            if not a:
                base_key=f"IMPORT-RECARGA-{re.sub(r'[^A-Za-z0-9_-]','_',row['terminal'])}"[:255]
                a=BaseAsset(asset_key=base_key,company=row["company"],line=row["line"],locality=row["station"],terminal_number=row["terminal"],equipment_type="VALIDADOR",base_status="ATIVO"); db.session.add(a); db.session.flush(); assets[row["terminal"]]=a
            a.company=row["company"]; a.line=row["line"]; a.locality=row["station"]; a.base_status="ATIVO"
            sw=ChipSwap.query.filter_by(base_asset_id=a.id).first(); _op_reopen_swap(sw,row["status"],"chip_swap",row["terminal"],user_id)
        for terminal,a in assets.items():
            if terminal in imported: continue
            sw=ChipSwap.query.filter_by(base_asset_id=a.id).first(); st=normalize(sw.status if sw else "PENDENTE").upper()
            if st in ("PENDENTE","EM ANDAMENTO"): a.base_status="INATIVO - FORA DO ESCOPO IMPORTAÇÃO"
        try: _invalidate_chip_swap_cache()
        except Exception: pass

    db.session.add(AuditEvent(user_id=user_id,event_type="OPERATIONAL_BASE_IMPORT",entity_type=mod,entity_id=now.strftime("%Y%m%d%H%M%S"),detail=json.dumps({**preview,"rule":"ausentes pendentes = fora do escopo; nunca apagados"},ensure_ascii=False)))
    db.session.commit()
    if module=="emv": _v63_invalidate_emv_cache()
    if module=="garagem": _v63_invalidate_garage_cache()
    return preview

@app.get("/api/operational-base/template.xlsx")
@login_required
def operational_base_template():
    if session.get("role") not in ("manager","manager_field"): abort(403)
    module=(request.args.get("module") or "emv").lower()
    if module not in _OP_MODULES: return jsonify({"ok":False,"error":"Módulo inválido."}),400
    wb=Workbook(); ws=wb.active; ws.title="Base operacional"
    ws.append(["empresa","terminal","estação","linha","status"])
    ws.append(["METRO","531505","Camilo Haddad","Linha 15 Prata","PENDENTE"])
    for c in ws[1]: c.font=Font(bold=True,color="FFFFFF"); c.fill=PatternFill("solid",fgColor="17365D")
    ws.freeze_panes="A2"; widths=[34,18,30,28,18]
    for i,w in enumerate(widths,1): ws.column_dimensions[get_column_letter(i)].width=w
    bio=io.BytesIO(); wb.save(bio); bio.seek(0)
    return send_file(bio,as_attachment=True,download_name=f"modelo_base_{module}.xlsx",mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

@app.post("/api/operational-base/import/preview")
@login_required
def operational_base_import_preview():
    if session.get("role") not in ("manager","manager_field"): return jsonify({"ok":False,"error":"Importação restrita a Gestor/ADM."}),403
    module=(request.args.get("module") or "").lower()
    if module not in _OP_MODULES: return jsonify({"ok":False,"error":"Módulo inválido."}),400
    try:
        rows=_op_parse_upload(request.files.get("file")); return jsonify({"ok":True,"module":module,"summary":_op_preview(module,rows)})
    except ValueError as exc: return jsonify({"ok":False,"error":str(exc)}),400
    except Exception as exc: app.logger.exception("V63 REV2 preview importador"); return jsonify({"ok":False,"error":str(exc)[:220]}),500

@app.post("/api/operational-base/import")
@login_required
def operational_base_import_apply():
    if session.get("role") not in ("manager","manager_field"): return jsonify({"ok":False,"error":"Importação restrita a Gestor/ADM."}),403
    module=(request.args.get("module") or "").lower()
    if module not in _OP_MODULES: return jsonify({"ok":False,"error":"Módulo inválido."}),400
    try:
        rows=_op_parse_upload(request.files.get("file")); summary=_op_apply(module,rows,session.get("user_id")); return jsonify({"ok":True,"summary":summary,"message":"Base atualizada. Registros ausentes não foram apagados; pendentes ausentes saíram do escopo operacional."})
    except ValueError as exc: db.session.rollback(); return jsonify({"ok":False,"error":str(exc)}),400
    except Exception as exc: db.session.rollback(); app.logger.exception("V63 REV2 importador operacional"); return jsonify({"ok":False,"error":"Falha ao aplicar a importação.","detail":str(exc)[:220]}),500

_emv_base_rows_cache = None
def _v41_emv_rows():
    """V62: base EMV consolidada. Mescla a base histórica com a planilha completa de não migrados.
    Normaliza espaços/aliases de empresa e infere Linha 8/9 pelo prefixo das estações 5xx/6xx.
    """
    global _emv_base_rows_cache
    if _emv_base_rows_cache:
        return _emv_base_rows_cache
    candidates=[BASE_DIR / "data_emv.xlsx", DATA_DIR / "data_emv.xlsx", BASE_DIR / "CHIPS NÃO MIGRADOS.xlsx", DATA_DIR / "CHIPS NÃO MIGRADOS.xlsx"]
    network_by_prefix={str(r.get("prefix") or ""):r for r in _load_station_network_rows() if r.get("prefix")}
    block_cfg={}
    try:
        source=DATA_DIR / "block_config_v18.json"
        payload=json.loads(source.read_text(encoding="utf-8")) if source.exists() else {}
        block_cfg=payload.get("by_prefix") or {}
    except Exception:
        block_cfg={}
    merged={}
    def clean_company(v):
        x=" ".join(str(v or "").strip().split())
        aliases={
          "VIA MOBILIDADE LINHAS 8 E 9":"VIA MOBILIDADE LINHAS 8 E 9",
          "VIA MOBILIDADE LINHA 8 E 9":"VIA MOBILIDADE LINHAS 8 E 9",
          "VIAMOBILIDADE LINHAS 8 E 9":"VIA MOBILIDADE LINHAS 8 E 9",
        }
        return aliases.get(x.upper(),x)
    def infer_line(company,station,raw_line):
        if str(raw_line or "").strip(): return str(raw_line).strip()
        st=str(station or "").strip(); co=clean_company(company).upper()
        if co=="VIA MOBILIDADE LINHAS 8 E 9":
            m=re.match(r"(\d{3})",st)
            if m:
                n=int(m.group(1))
                if 500 <= n < 600: return "Linha 8 Diamante"
                if 600 <= n < 700: return "Linha 9 Esmeralda"
            return "Linhas 8 e 9"
        return ""
    for path in candidates:
        if not path.exists(): continue
        try:
            wb=load_workbook(path,read_only=True,data_only=True); ws=wb.active
            headers=[str(x or "").strip().lower() for x in next(ws.iter_rows(values_only=True))]
            for vals in ws.iter_rows(min_row=2,values_only=True):
                d=dict(zip(headers,vals)); terminal=str(d.get("terminal") or "").split('.')[0].strip()
                if not terminal: continue
                key=re.sub(r"\D","",terminal); cfg=block_cfg.get(key) or network_by_prefix.get(key) or {}
                company=clean_company(d.get("empresa") or "")
                station=str(d.get("estação") or d.get("estacao") or "").strip()
                line=infer_line(company,station,d.get("linha"))
                # V62 REV1 — Linha 7 Rubi operada pela TIC TRENS S.A.
                if "7" in normalize(line).upper() and "RUBI" in normalize(line).upper():
                    company="TIC TRENS S.A."
                row={"tp_id":d.get("tp_id"),"company":company,"terminal":terminal,"version":str(d.get("versão") or d.get("versao") or "").strip(),"ip":str(d.get("ip") or cfg.get("ip") or "").strip(),"station":station,"line":line,"mask":cfg.get("mask") or "","gateway":cfg.get("gateway") or "","dns1":cfg.get("dns1") or "","dns2":cfg.get("dns2") or "","group":cfg.get("group") or "","blocking_number":cfg.get("blocking_number") or terminal[-2:]}
                if terminal in merged:
                    old=merged[terminal]
                    for k,v in row.items():
                        if v not in (None,""): old[k]=v
                else: merged[terminal]=row
            wb.close()
        except Exception as exc:
            app.logger.warning("V62 falha lendo base EMV %s: %s",path.name,exc)
    # V63 REV1 — fonte operacional mais recente exportada em 27/08/2026.
    # Ela complementa/substitui por terminal sem remover itens históricos que não constem do recorte.
    try:
        ref_path=DATA_DIR / "trilho_270826.json"
        ref=json.loads(ref_path.read_text(encoding="utf-8")) if ref_path.exists() else {}
        for d in (ref.get("rows") or []):
            terminal=str(d.get("terminal") or "").strip()
            if not terminal: continue
            company=clean_company(d.get("company") or "")
            station=str(d.get("station") or "").strip()
            line=infer_line(company,station,d.get("line"))
            if "7" in normalize(line).upper() and "RUBI" in normalize(line).upper(): company="TIC TRENS S.A."
            cfg=block_cfg.get(re.sub(r"\D","",terminal)) or {}
            row={"tp_id":None,"company":company,"terminal":terminal,"version":str(d.get("version") or "").strip(),"ip":str(d.get("ip") or cfg.get("ip") or "").strip(),"station":station,"line":line,"mask":cfg.get("mask") or "","gateway":cfg.get("gateway") or "","dns1":cfg.get("dns1") or "","dns2":cfg.get("dns2") or "","group":cfg.get("group") or "","blocking_number":cfg.get("blocking_number") or terminal[-2:],"sam_type":d.get("sam_type") or "NÃO MIGRADO","source":"TRILHO_270826"}
            if terminal in merged:
                merged[terminal].update({k:v for k,v in row.items() if v not in (None,"")})
            else:
                merged[terminal]=row
    except Exception as exc:
        app.logger.warning("V63 REV1 falha lendo trilho_270826.json: %s",exc)
    # V63 REV2: quando existe importação administrativa, ela define os pendentes ativos.
    # Concluídos históricos fora do arquivo permanecem visíveis; pendentes ausentes ficam fora do escopo, nunca apagados.
    try:
        imported=_op_active_map("emv")
    except Exception:
        imported={}
    if imported:
        swaps_status={x.terminal:normalize(x.status or "").upper() for x in EmvChipSwap.query.all()}
        active_rows={}
        for terminal,item in imported.items():
            base=dict(merged.get(terminal) or {})
            base.update({"company":item.company or base.get("company","") ,"terminal":terminal,"station":item.station or base.get("station","") ,"line":item.line or base.get("line","") ,"_base_status":item.desired_status or "PENDENTE"})
            base.setdefault("version",""); base.setdefault("ip",""); base.setdefault("mask",""); base.setdefault("gateway",""); base.setdefault("dns1",""); base.setdefault("dns2",""); base.setdefault("group",""); base.setdefault("blocking_number",terminal[-2:])
            active_rows[terminal]=base
        for terminal,base in merged.items():
            if terminal not in active_rows and swaps_status.get(terminal)=="CONCLUIDA": active_rows[terminal]=base
        merged=active_rows
    rows=list(merged.values())
    if not rows:
        try:
            assets=BaseAsset.query.filter(func.upper(BaseAsset.equipment_type).like("%BLOQ%")).all()
            for a in assets:
                terminal=str(a.terminal_number or a.asset_identifier or "").strip()
                if terminal: rows.append({"tp_id":None,"company":clean_company(a.company),"terminal":terminal,"version":"","ip":getattr(a,"ip_address",None) or "","station":a.locality or "","line":a.line or "","mask":"","gateway":"","dns1":"","dns2":"","group":"","blocking_number":terminal[-2:]})
        except Exception as exc: app.logger.warning("Fallback base EMV indisponível: %s",exc)
    if rows: _emv_base_rows_cache=rows
    return rows

_emv_base_by_terminal_cache = None
def _v41_emv_by_terminal():
    global _emv_base_by_terminal_cache
    if _emv_base_by_terminal_cache is None:
        _emv_base_by_terminal_cache = {str(x.get("terminal") or ""): x for x in _v41_emv_rows() if x.get("terminal")}
    return _emv_base_by_terminal_cache

def _sync_emv_trilho_rev1_once():
    """Aplica uma única vez o recorte Trilhos 27/08: reabre concluídos presentes no recorte.
    Histórico, fotos, resultado e datas anteriores permanecem preservados; a retificação fica auditada.
    Novos terminais entram pela fonte base trilho_270826.json, portanto não exigem swap fictício.
    """
    marker="EMV_SYNC_TRILHO_270826_REV1"
    try:
        if AuditEvent.query.filter_by(event_type=marker).first(): return
        path=DATA_DIR / "trilho_270826.json"
        data=json.loads(path.read_text(encoding="utf-8")) if path.exists() else {}
        terms={str(x.get("terminal") or "").strip() for x in (data.get("rows") or []) if str(x.get("terminal") or "").strip()}
        if not terms: return
        swaps=EmvChipSwap.query.filter(EmvChipSwap.terminal.in_(terms)).all()
        reopened=0
        for sw in swaps:
            st=normalize(sw.status or "")
            if st in ("CONCLUIDA","CONCLUÍDA"):
                before={"status":sw.status,"test_result":sw.test_result,"completed_at":sw.completed_at.isoformat() if sw.completed_at else None,"completed_by_id":sw.completed_by_id}
                sw.status="PENDENTE"; sw.updated_at=datetime.utcnow(); reopened+=1
                db.session.add(AuditEvent(user_id=None,event_type="EMV_REOPEN_TRILHO_270826",entity_type="emv_chip_swap",entity_id=str(sw.id),detail=json.dumps({"terminal":sw.terminal,"before":before,"reason":"Sincronização Trilhos 27/08/2026"},ensure_ascii=False)))
        db.session.add(AuditEvent(user_id=None,event_type=marker,entity_type="emv_reference",entity_id="2026-08-27",detail=json.dumps({"reference_terminals":len(terms),"reopened":reopened,"mode":"non_destructive"},ensure_ascii=False)))
        db.session.commit(); _v63_invalidate_emv_cache()
        app.logger.info("V63 REV1 EMV Trilhos 27/08 sincronizado: %s terminais; %s reabertos",len(terms),reopened)
    except Exception as exc:
        db.session.rollback(); app.logger.exception("V63 REV1 falha sincronização EMV Trilhos: %s",exc)

_EMV_SCHEMA_READY=False
_EMV_SCHEMA_LOCK=threading.Lock()
def _ensure_emv_tables():
    """V66 REV4: valida/migra o schema EMV uma vez por processo.

    A versão anterior fazia checkfirst + inspeção de colunas em cada leitura fria
    da API. Em produção isso adicionava trabalho administrativo ao caminho quente.
    """
    global _EMV_SCHEMA_READY
    if _EMV_SCHEMA_READY: return
    with _EMV_SCHEMA_LOCK:
        if _EMV_SCHEMA_READY: return
        EmvChipSwap.__table__.create(bind=db.engine,checkfirst=True); EmvChipSwapPhoto.__table__.create(bind=db.engine,checkfirst=True)
        try:
            cols={c["name"] for c in db.inspect(db.engine).get_columns("emv_chip_swaps")}
            additions=[]
            if "completed_by_id" not in cols: additions.append(("completed_by_id","INTEGER"))
            if "manual_entry" not in cols: additions.append(("manual_entry","BOOLEAN NOT NULL DEFAULT FALSE"))
            if "company" not in cols: additions.append(("company","VARCHAR(120)"))
            if "line" not in cols: additions.append(("line","VARCHAR(120)"))
            if "station" not in cols: additions.append(("station","VARCHAR(180)"))
            if "block_number" not in cols: additions.append(("block_number","VARCHAR(80)"))
            if additions:
                with db.engine.begin() as conn:
                    for name,typ in additions: conn.execute(text(f"ALTER TABLE emv_chip_swaps ADD COLUMN {name} {typ}"))
            _EMV_SCHEMA_READY=True
        except Exception as exc:
            app.logger.warning("V66 REV4 EMV schema migration: %s", exc)
    # V63 REV2: atualizações de base passam a ocorrer somente pelo importador administrativo.

# V60 REV3 — Troca de Chips Garagem, base inicial MIGRACAO_SAM GARAGENS.xlsx.
GARAGE_CHIP_SEED = [{'company': 'ABC SISTEMA DE TRANSPORTE SPE S.A', 'terminal': '9906', 'ip': '10.109.22.167', 'model': 'B80', 'sam_type': 'NÃO MIGRADO'}, {'company': 'ABC SISTEMA DE TRANSPORTE SPE S.A', 'terminal': '9907', 'ip': '10.109.18.234', 'model': 'B80', 'sam_type': 'NÃO MIGRADO'}, {'company': 'ALTO TIETE TRANSPORTES FILIAL', 'terminal': '41680', 'ip': '10.180.136.59', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'ALTO TIETE TRANSPORTES FILIAL', 'terminal': '47153', 'ip': '10109120173', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'ALTO TIETE TRANSPORTES FILIAL', 'terminal': '47353', 'ip': '10.109.8.187', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'ALTO TIETE TRANSPORTES FILIAL', 'terminal': '47465', 'ip': '10.109.8.222', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'ALTO TIETE TRANSPORTES FILIAL', 'terminal': '47487', 'ip': '10.109.82.32', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'ALTO TIETE TRANSPORTES LTDA', 'terminal': '47022', 'ip': '10.109.97.201', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'ALTO TIETE TRANSPORTES LTDA', 'terminal': '47091', 'ip': '10190133202', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'ALTO TIETE TRANSPORTES LTDA', 'terminal': '47099', 'ip': '10.109.21.210', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'ALTO TIETE TRANSPORTES LTDA', 'terminal': '47207', 'ip': '10.109.96.241', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'ALTO TIETE TRANSPORTES LTDA', 'terminal': '47225', 'ip': '10.109.99.247', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'ALTO TIETE TRANSPORTES LTDA', 'terminal': '47267', 'ip': '10190132161', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'ALTO TIETE TRANSPORTES LTDA', 'terminal': '47595', 'ip': '10180136140', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'ALTO TIETE TRANSPORTES LTDA', 'terminal': '47654', 'ip': '10.109.8.204', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'ALTO TIETE TRANSPORTES LTDA', 'terminal': '47664', 'ip': '10.109.43.100', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'ARUJA TRANS COLT - MUNICIPAL', 'terminal': '2301', 'ip': '10.109.74.108', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'ARUJA TRANS COLT - MUNICIPAL', 'terminal': '2304', 'ip': '10.109.21.122', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'ARUJA TRANS COLT - MUNICIPAL', 'terminal': '2309', 'ip': '10.109.97.167', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'ARUJA TRANS COLT - MUNICIPAL', 'terminal': '2313', 'ip': '10.109.4.182', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'ARUJA TRANS COLT - MUNICIPAL', 'terminal': '2317', 'ip': '10.109.32.83', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'ARUJA TRANSPORTES COLETIVOS LTDA', 'terminal': '39564', 'ip': '10.190.11.181', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'ARUJA TRANSPORTES COLETIVOS LTDA', 'terminal': '39581', 'ip': '10.180.137.89', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'ARUJA TRANSPORTES COLETIVOS LTDA', 'terminal': '39583', 'ip': '10.109.61.142', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'ARUJA TRANSPORTES COLETIVOS LTDA', 'terminal': '39588', 'ip': '10190132119', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'ARUJA TRANSPORTES COLETIVOS LTDA', 'terminal': '39591', 'ip': '10.109.30.56', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'AUTO ÔNIBUS MORATENSE LTDA', 'terminal': '26058', 'ip': '10.109.22.221', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'AUTO ÔNIBUS MORATENSE LTDA', 'terminal': '26061', 'ip': '10.109.99.26', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'AUTO ÔNIBUS MORATENSE LTDA', 'terminal': '26064', 'ip': '10.109.81.184', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'AUTO VIAÇÃO URUBUPUNGÁ LTDA', 'terminal': '20003', 'ip': '172.100.78.35', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'AUTO VIAÇÃO URUBUPUNGÁ LTDA', 'terminal': '20037', 'ip': '10.109.96.239', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'AUTO VIAÇÃO URUBUPUNGÁ LTDA', 'terminal': '20117', 'ip': '10.109.99.209', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'AUTO VIAÇÃO URUBUPUNGÁ LTDA', 'terminal': '20127', 'ip': '10.109.98.191', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'AUTO VIAÇÃO URUBUPUNGÁ LTDA', 'terminal': '20153', 'ip': '10.109.36.134', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'AUTO VIAÇÃO URUBUPUNGÁ LTDA', 'terminal': '20467', 'ip': '10.109.42.160', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'AUTO VIAÇÃO URUBUPUNGÁ LTDA', 'terminal': '20488', 'ip': '10.109.39.203', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'AUTO VIAÇÃO URUBUPUNGÁ LTDA', 'terminal': '20490', 'ip': '10.109.18.252', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'AUTO VIAÇÃO URUBUPUNGÁ LTDA', 'terminal': '20509', 'ip': '10180137250', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'AUTO VIAÇÃO URUBUPUNGÁ LTDA', 'terminal': '20539', 'ip': '10.109.19.66', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'AUTO VIAÇÃO URUBUPUNGÁ LTDA', 'terminal': '20596', 'ip': '10.109.98.120', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'AUTO VIAÇÃO URUBUPUNGÁ LTDA', 'terminal': '20640', 'ip': '10.109.15.113', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'AUTO VIAÇÃO URUBUPUNGÁ LTDA', 'terminal': '20648', 'ip': '10.109.96.191', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'AUTO VIAÇÃO URUBUPUNGÁ LTDA', 'terminal': '20654', 'ip': '10.109.38.109', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'AUTO VIAÇÃO URUBUPUNGÁ LTDA', 'terminal': '20668', 'ip': '10.109.97.27', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'AUTO VIAÇÃO URUBUPUNGÁ LTDA', 'terminal': '20685', 'ip': '10.109.98.78', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'AUTO VIAÇÃO URUBUPUNGÁ LTDA', 'terminal': '20689', 'ip': '10180137140', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'AUTO VIAÇÃO URUBUPUNGÁ LTDA', 'terminal': '20720', 'ip': '10109120163', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'AUTO VIAÇÃO URUBUPUNGÁ LTDA', 'terminal': '20764', 'ip': '10.109.82.236', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'AUTO VIAÇÃO URUBUPUNGÁ LTDA', 'terminal': '20768', 'ip': '172.100.81.32', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'AUTO VIAÇÃO URUBUPUNGÁ LTDA', 'terminal': '20769', 'ip': '10.109.97.145', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'AUTO VIAÇÃO URUBUPUNGÁ LTDA', 'terminal': '20775', 'ip': '10.109.9.31', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'AUTO VIAÇÃO URUBUPUNGÁ LTDA', 'terminal': '20785', 'ip': '10.109.6.10', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'AUTO VIAÇÃO URUBUPUNGÁ LTDA', 'terminal': '20787', 'ip': '10.109.35.6', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'AUTO VIAÇÃO URUBUPUNGÁ LTDA', 'terminal': '20967', 'ip': '10180137244', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'BB TRANSPORTE E TURISMO LTDA', 'terminal': '27443', 'ip': '10.109.98.93', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'BB TRANSPORTE E TURISMO LTDA', 'terminal': '27447', 'ip': '10.109.96.72', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'BB TRANSPORTE E TURISMO LTDA', 'terminal': '27483', 'ip': '10.109.96.57', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'BB TRANSPORTE E TURISMO LTDA', 'terminal': '27487', 'ip': '10.109.60.179', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'BB TRANSPORTE E TURISMO LTDA', 'terminal': '27513', 'ip': '10.109.30.174', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'BB TRANSPORTE E TURISMO LTDA', 'terminal': '27519', 'ip': '10.109.99.132', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'BB TRANSPORTE E TURISMO LTDA', 'terminal': '27521', 'ip': '10.109.21.175', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'BB TRANSPORTE E TURISMO LTDA', 'terminal': '27533', 'ip': '10.109.12.142', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'BB TRANSPORTE E TURISMO LTDA', 'terminal': '27575', 'ip': '10.109.96.232', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'BB TRANSPORTE E TURISMO LTDA', 'terminal': '27589', 'ip': '10.109.8.168', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'BB TRANSPORTE E TURISMO LTDA', 'terminal': '27601', 'ip': '10.109.20.172', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'BB TRANSPORTE E TURISMO LTDA', 'terminal': '27645', 'ip': '10.109.96.195', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'BB TRANSPORTE E TURISMO LTDA', 'terminal': '27671', 'ip': '10.109.32.128', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'BB TRANSPORTE E TURISMO LTDA', 'terminal': '88879', 'ip': '0.0.0.0', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'DANUBIO AZUL - SUB', 'terminal': '19063', 'ip': '10.109.98.54', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'DANUBIO AZUL - SUB', 'terminal': '99033', 'ip': '10.109.17.65', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'DEL REY TRANSPORTES LTDA', 'terminal': '25193', 'ip': '0.0.0.0', 'model': 'B80', 'sam_type': 'NÃO MIGRADO'}, {'company': 'DEL REY TRANSPORTES LTDA', 'terminal': '25194', 'ip': '0.0.0.0', 'model': 'B80', 'sam_type': 'NÃO MIGRADO'}, {'company': 'DEL REY TRANSPORTES LTDA', 'terminal': '25211', 'ip': '10.109.29.205', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'DEL REY TRANSPORTES LTDA', 'terminal': '25213', 'ip': '10.190.132.10', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA DE ÔNIBUS VILA GALVÃO ', 'terminal': '30097', 'ip': '10.109.41.74', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA DE ÔNIBUS VILA GALVÃO ', 'terminal': '30587', 'ip': '172.70.119.123', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA DE ÔNIBUS VILA GALVÃO ', 'terminal': '30596', 'ip': '10.109.38.160', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA DE ÔNIBUS VILA GALVÃO ', 'terminal': '30609', 'ip': '172.70.119.145', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA DE ÔNIBUS VILA GALVÃO ', 'terminal': '30625', 'ip': '10.109.99.140', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA DE ÔNIBUS VILA GALVÃO ', 'terminal': '30631', 'ip': '10.109.36.196', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA DE ÔNIBUS VILA GALVÃO ', 'terminal': '30668', 'ip': '10.109.38.94', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA DE ÔNIBUS VILA GALVÃO ', 'terminal': '30670', 'ip': '10.109.12.225', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA DE ÔNIBUS VILA GALVÃO ', 'terminal': '30694', 'ip': '10109125101', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA DE ÔNIBUS VILA GALVÃO ', 'terminal': '30715', 'ip': '10.109.60.176', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA DE ÔNIBUS VILA GALVÃO ', 'terminal': '30716', 'ip': '10.109.14.181', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA DE ÔNIBUS VILA GALVÃO ', 'terminal': '30813', 'ip': '10.109.6.166', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA DE ÔNIBUS VILA GALVÃO ', 'terminal': '30823', 'ip': '10.109.97.76', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA DE ÔNIBUS VILA GALVÃO ', 'terminal': '30870', 'ip': '10.109.13.20', 'model': 'B80', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA DE ÔNIBUS VILA GALVÃO ', 'terminal': '30878', 'ip': '10.109.20.22', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA DE ÔNIBUS VILA GALVÃO ', 'terminal': '30883', 'ip': '10.109.20.60', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA DE ÔNIBUS VILA GALVÃO ', 'terminal': '30890', 'ip': '10.109.12.96', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA DE ÔNIBUS VILA GALVÃO ', 'terminal': '30903', 'ip': '10.109.98.213', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA DE TRANSPORTES MAIRIPORÃ', 'terminal': '39003', 'ip': '10.109.8.86', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA DE TRANSPORTES MAIRIPORÃ', 'terminal': '39180', 'ip': '10.109.14.27', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA DE TRANSPORTES MAIRIPORÃ', 'terminal': '39192', 'ip': '10.109.62.105', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA DE TRANSPORTES MAIRIPORÃ', 'terminal': '39202', 'ip': '10.109.10.232', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA DE TRANSPORTES MAIRIPORÃ', 'terminal': '39208', 'ip': '10.109.19.112', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA DE TRANSPORTES MAIRIPORÃ', 'terminal': '39218', 'ip': '10.109.38.5', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA TRANSPORTES E TURISMO', 'terminal': '24416', 'ip': '10.109.23.115', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA TRANSPORTES E TURISMO', 'terminal': '24500', 'ip': '10.109.81.134', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA TRANSPORTES E TURISMO', 'terminal': '24518', 'ip': '10.109.81.24', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA TRANSPORTES E TURISMO', 'terminal': '24548', 'ip': '10.109.39.184', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA TRANSPORTES E TURISMO', 'terminal': '24752', 'ip': '10.109.15.35', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA TRANSPORTES E TURISMO', 'terminal': '24758', 'ip': '10.109.14.56', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA TRANSPORTES E TURISMO', 'terminal': '24834', 'ip': '10.109.75.80', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'EMPRESA TRANSPORTES E TURISMO', 'terminal': '24842', 'ip': '10109126204', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'GUARULHOS TRANSPORTES S/A', 'terminal': '33002', 'ip': '10.109.43.149', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'GUARULHOS TRANSPORTES S/A', 'terminal': '33624', 'ip': '10.109.16.7', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'GUARULHOS TRANSPORTES S/A', 'terminal': '33632', 'ip': '10.109.15.133', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'GUARULHOS TRANSPORTES S/A', 'terminal': '33638', 'ip': '10.109.97.160', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'GUARULHOS TRANSPORTES S/A', 'terminal': '33645', 'ip': '10.109.75.225', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'GUARULHOS TRANSPORTES S/A', 'terminal': '33651', 'ip': '10.109.22.16', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'GUARULHOS TRANSPORTES S/A', 'terminal': '33663', 'ip': '10.109.30.248', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'GUARULHOS TRANSPORTES S/A', 'terminal': '33670', 'ip': '10.109.97.140', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'GUARULHOS TRANSPORTES S/A', 'terminal': '33723', 'ip': '10.109.34.118', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'GUARULHOS TRANSPORTES S/A', 'terminal': '33724', 'ip': '10.109.37.158', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'GUARULHOS TRANSPORTES S/A', 'terminal': '33747', 'ip': '10.109.98.14', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'GUARULHOS TRANSPORTES S/A', 'terminal': '33749', 'ip': '10.190.133.37', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'GUARULHOS TRANSPORTES S/A', 'terminal': '33772', 'ip': '10.109.99.124', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'GUARULHOS TRANSPORTES S/A', 'terminal': '33777', 'ip': '10.109.97.87', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'GUARULHOS TRANSPORTES S/A', 'terminal': '33780', 'ip': '10.109.43.54', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'GUARULHOS TRANSPORTES S/A', 'terminal': '33785', 'ip': '10.109.96.39', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'GUARULHOS TRANSPORTES S/A', 'terminal': '33786', 'ip': '10.109.16.194', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'GUARULHOS TRANSPORTES S/A', 'terminal': '33789', 'ip': '10.109.32.36', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'GUARULHOS TRANSPORTES S/A', 'terminal': '33795', 'ip': '10.109.36.81', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'GUARULHOS TRANSPORTES S/A', 'terminal': '33800', 'ip': '10.109.40.204', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'GUARULHOS TRANSPORTES S/A', 'terminal': '33830', 'ip': '10.109.21.198', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'GUARULHOS TRANSPORTES S/A', 'terminal': '34091', 'ip': '10.180.2.127', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '80027', 'ip': '10.109.30.106', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '80211', 'ip': '10.109.33.179', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '80217', 'ip': '10.109.125.60', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '80251', 'ip': '10.109.99.232', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '80427', 'ip': '10180136105', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '80433', 'ip': '10.109.33.198', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '80633', 'ip': '10.109.17.105', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '80639', 'ip': '10.109.121.60', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '80859', 'ip': '10.109.40.94', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '80861', 'ip': '10.109.19.150', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '80907', 'ip': '10190132210', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '80921', 'ip': '10180137183', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '80923', 'ip': '10.109.4.134', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '80943', 'ip': '10.109.5.62', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '81047', 'ip': '10.109.22.192', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '81157', 'ip': '10.180.136.61', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '81159', 'ip': '10.109.72.164', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '81307', 'ip': '10.109.4.80', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '81361', 'ip': '10.109.83.139', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '81371', 'ip': '10.109.33.191', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '81383', 'ip': '10.180.136.50', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '81393', 'ip': '10.109.6.66', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '81611', 'ip': '10.109.38.8', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '81647', 'ip': '10.109.30.41', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '81727', 'ip': '10.109.14.59', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '81743', 'ip': '10.109.7.21', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '81753', 'ip': '10.109.99.152', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '81867', 'ip': '10.109.29.238', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '81917', 'ip': '10.109.40.62', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '81943', 'ip': '10.109.7.185', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '81967', 'ip': '10180136193', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '82411', 'ip': '10.180.136.88', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST REMANESCE', 'terminal': '82503', 'ip': '10.109.42.150', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '3', 'ip': '10190132229', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '5430', 'ip': '10.109.82.128', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '5439', 'ip': '10.109.14.76', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '6', 'ip': '10.109.40.71', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '7', 'ip': '10.109.38.151', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '7400', 'ip': '10.109.12.29', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '7401', 'ip': '10.109.98.114', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '7715', 'ip': '10.109.22.5', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '8', 'ip': '10.109.18.135', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '8106', 'ip': '10.109.11.199', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '8116', 'ip': '10.109.34.167', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '8150', 'ip': '10.180.136.79', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '8164', 'ip': '10190133116', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '8170', 'ip': '10.109.13.97', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '8251', 'ip': '10190133230', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '8254', 'ip': '10109113222', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '8255', 'ip': '10.109.96.100', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '8256', 'ip': '10.109.19.37', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '8265', 'ip': '10.109.30.83', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '8316', 'ip': '10.109.82.160', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '8326', 'ip': '10.109.83.128', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '8328', 'ip': '10.109.96.238', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '8335', 'ip': '10.109.98.151', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '8354', 'ip': '10.109.83.17', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '8369', 'ip': '10.109.11.29', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '8370', 'ip': '10.109.15.88', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '8380', 'ip': '10.109.10.121', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '9122', 'ip': '10.109.21.222', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '9212', 'ip': '10.109.41.77', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '9213', 'ip': '10.109.98.133', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '9232', 'ip': '10.109.19.111', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '9234', 'ip': '10180137118', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '9432', 'ip': '10.109.13.62', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '9613', 'ip': '10.109.18.70', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '9673', 'ip': '10190132104', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '9905', 'ip': '10.109.83.253', 'model': 'B80', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '9910', 'ip': '10.109.22.210', 'model': 'B80', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '9911', 'ip': '10.109.22.27', 'model': 'B80', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '9912', 'ip': '10.109.36.110', 'model': 'B80', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '9915', 'ip': '10.109.20.181', 'model': 'B80', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '9916', 'ip': '10.109.81.141', 'model': 'B80', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '9917', 'ip': '10.109.62.204', 'model': 'B80', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '9920', 'ip': '10109112126', 'model': 'B80', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '9921', 'ip': '10.109.82.233', 'model': 'B80', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '9922', 'ip': '10.109.31.241', 'model': 'B80', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '9925', 'ip': '10.109.99.59', 'model': 'B80', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '9926', 'ip': '10.109.96.2', 'model': 'B80', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '9927', 'ip': '10.109.18.27', 'model': 'B80', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '9932', 'ip': '10.109.72.30', 'model': 'B80', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '9935', 'ip': '10.109.99.199', 'model': 'B80', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '9936', 'ip': '10109127110', 'model': 'B80', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '9937', 'ip': '10.109.74.198', 'model': 'B80', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '9940', 'ip': '10.109.80.217', 'model': 'B80', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '9942', 'ip': '10.109.122.44', 'model': 'B80', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '9943', 'ip': '10.109.80.204', 'model': 'B80', 'sam_type': 'NÃO MIGRADO'}, {'company': 'NEXT MOBILIDADE SIST.EXISTENTE', 'terminal': '9945', 'ip': '10.109.29.16', 'model': 'B80', 'sam_type': 'NÃO MIGRADO'}, {'company': 'PÁSSARO MARRON SANTA ISABEL', 'terminal': '37653', 'ip': '10.109.81.139', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'PÁSSARO MARRON SANTA ISABEL', 'terminal': '37654', 'ip': '10.109.29.222', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'PÁSSARO MARRON SANTA ISABEL', 'terminal': '37656', 'ip': '10.109.10.163', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'PÁSSARO MARRON SANTA ISABEL', 'terminal': '37657', 'ip': '10.109.34.247', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'PÁSSARO MARRON SANTA ISABEL', 'terminal': '37809', 'ip': '10.109.11.218', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'PÁSSARO MARRON SANTA ISABEL', 'terminal': '37813', 'ip': '10.109.97.249', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'PÁSSARO MARRON UNILESTE', 'terminal': '45016', 'ip': '10.109.96.116', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'PÁSSARO MARRON UNILESTE', 'terminal': '45019', 'ip': '10.109.83.48', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'PÁSSARO MARRON UNILESTE', 'terminal': '45511', 'ip': '10.190.133.78', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'RADIAL SUZANO INTERMUNICIPAL ', 'terminal': '41023', 'ip': '10.109.18.109', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'RADIAL SUZANO INTERMUNICIPAL ', 'terminal': '41361', 'ip': '10.109.37.2', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'RADIAL SUZANO INTERMUNICIPAL ', 'terminal': '41663', 'ip': '10.109.60.56', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'RADIAL SUZANO INTERMUNICIPAL ', 'terminal': '41728', 'ip': '10.109.4.191', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'RADIAL SUZANO INTERMUNICIPAL ', 'terminal': '41742', 'ip': '10190133106', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'RALIP TRANSPORTES RODOVIÁRIOS ', 'terminal': '23065', 'ip': '10.109.40.172', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'RALIP TRANSPORTES RODOVIÁRIOS ', 'terminal': '23073', 'ip': '10.109.99.186', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'RTO-EMTU CECON', 'terminal': '0', 'ip': '0.0.0.0', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'RTO-EMTU CECON', 'terminal': '2048', 'ip': '10180137124', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'RTO-EMTU CECON', 'terminal': '210', 'ip': '0.0.0.0', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'RTO-EMTU CECON', 'terminal': '211', 'ip': '0.0.0.0', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'RTO-EMTU CECON', 'terminal': '212', 'ip': '0.0.0.0', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'RTO-EMTU CECON', 'terminal': '213', 'ip': '0.0.0.0', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'RTO-EMTU CECON', 'terminal': '216', 'ip': '0.0.0.0', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'RTO-EMTU CECON', 'terminal': '220', 'ip': '0.0.0.0', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'RTO-EMTU CECON', 'terminal': '221', 'ip': '0.0.0.0', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'RTO-EMTU CECON', 'terminal': '222', 'ip': '0.0.0.0', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'RTO-EMTU CECON', 'terminal': '223', 'ip': '0.0.0.0', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'RTO-EMTU CECON', 'terminal': '224', 'ip': '0.0.0.0', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'RTO-EMTU CECON', 'terminal': '225', 'ip': '0.0.0.0', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'RTO-EMTU CECON', 'terminal': '226', 'ip': '0.0.0.0', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'RTO-EMTU CECON', 'terminal': '228', 'ip': '0.0.0.0', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'RTO-EMTU CECON', 'terminal': '229', 'ip': '0.0.0.0', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'RTO-EMTU CECON', 'terminal': '232', 'ip': '0.0.0.0', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'RTO-EMTU CECON', 'terminal': '233', 'ip': '0.0.0.0', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'RTO-EMTU CECON', 'terminal': '235', 'ip': '0.0.0.0', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'RTO-EMTU CECON', 'terminal': '999', 'ip': '10190132197', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'SERVENG TRANSPORTES LTDA', 'terminal': '37301', 'ip': '10180137203', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'SERVENG TRANSPORTES LTDA', 'terminal': '37302', 'ip': '10190133200', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'SERVENG TRANSPORTES LTDA', 'terminal': '37404', 'ip': '10180136153', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'TIPBUS TRANSPORTE INTERMUNICIPAL', 'terminal': '36160', 'ip': '10.109.20.237', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'TIPBUS TRANSPORTE INTERMUNICIPAL', 'terminal': '36166', 'ip': '10.109.17.177', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'TIPBUS TRANSPORTE INTERMUNICIPAL', 'terminal': '36173', 'ip': '10.109.36.99', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'TIPBUS TRANSPORTE INTERMUNICIPAL', 'terminal': '36174', 'ip': '10.180.137.60', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'TIPBUS TRANSPORTE INTERMUNICIPAL', 'terminal': '36206', 'ip': '10.109.22.101', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'TIPBUS TRANSPORTE INTERMUNICIPAL', 'terminal': '36207', 'ip': '10.109.39.234', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'TIPBUS TRANSPORTE INTERMUNICIPAL', 'terminal': '88857', 'ip': '10.109.18.63', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'TIPBUS TRANSPORTE INTERMUNICIPAL', 'terminal': '88865', 'ip': '10.109.30.250', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'TIPBUS TRANSPORTE INTERMUNICIPAL', 'terminal': '888895', 'ip': '10.109.23.57', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO CIDADE DE CAIEIRAS LTDA', 'terminal': '22137', 'ip': '172.168.86.121', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO CIDADE DE CAIEIRAS LTDA', 'terminal': '22148', 'ip': '10.109.99.184', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO CIDADE DE CAIEIRAS LTDA', 'terminal': '22153', 'ip': '10.109.97.189', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO CIDADE DE CAIEIRAS LTDA', 'terminal': '22154', 'ip': '10.109.39.191', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO CIDADE DE CAIEIRAS LTDA', 'terminal': '22998', 'ip': '10.109.36.79', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO FERVIMA LTDA MUNICIPAL', 'terminal': '600', 'ip': '10.109.40.146', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO FERVIMA LTDA MUNICIPAL', 'terminal': '612', 'ip': '10.109.39.181', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO FERVIMA LTDA MUNICIPAL', 'terminal': '677', 'ip': '10.109.98.35', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO FERVIMA LTDA MUNICIPAL', 'terminal': '706', 'ip': '10.109.97.4', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO FERVIMA LTDA MUNICIPAL', 'terminal': '710', 'ip': '10.109.96.164', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO FERVIMA LTDA MUNICIPAL', 'terminal': '724', 'ip': '10.109.30.201', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO FERVIMA LTDA MUNICIPAL', 'terminal': '729', 'ip': '10.109.98.67', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO JACAREÍ LTDA', 'terminal': '44009', 'ip': '10.109.23.12', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO JACAREÍ LTDA', 'terminal': '44787', 'ip': '10.109.4.84', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '1032', 'ip': '10.109.98.8', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '1033', 'ip': '10.109.98.7', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '1050', 'ip': '10.109.98.137', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '1056', 'ip': '10.109.5.210', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '1058', 'ip': '10.109.72.29', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '1065', 'ip': '10.109.18.117', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '1073', 'ip': '10.109.20.179', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '1087', 'ip': '10.109.38.122', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '1088', 'ip': '10.109.16.243', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '1093', 'ip': '10.109.43.140', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '1096', 'ip': '10.109.12.255', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '1125', 'ip': '10.109.15.105', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '1127', 'ip': '10.109.42.81', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '1150', 'ip': '10.109.97.198', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '1188', 'ip': '10109121205', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '1189', 'ip': '10.109.20.143', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '1195', 'ip': '10.109.28.137', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '1199', 'ip': '10.109.97.238', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '1201', 'ip': '10.109.9.21', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '1214', 'ip': '10.109.6.245', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '1241', 'ip': '10.109.10.154', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '1259', 'ip': '10.109.21.247', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '1265', 'ip': '10.109.32.193', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '1269', 'ip': '10.109.14.232', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '1273', 'ip': '10.109.13.155', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '1281', 'ip': '10.109.38.164', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '15107', 'ip': '10190132129', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '15108', 'ip': '10180136166', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '15109', 'ip': '10.180.136.92', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '15110', 'ip': '10.180.136.48', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '15111', 'ip': '10190132174', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '15112', 'ip': '10.190.133.90', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '15701', 'ip': '10190133252', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '15702', 'ip': '10190132183', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '15703', 'ip': '10.180.136.28', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '15904', 'ip': '10.109.42.70', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '15942', 'ip': '10.109.37.123', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '15987', 'ip': '10180136116', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO MIRACATIBA LTDA', 'terminal': '421', 'ip': '10190133141', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO LTDA - FILIAL', 'terminal': '21000', 'ip': '10190132140', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO LTDA - FILIAL', 'terminal': '21415', 'ip': '10.190.132.16', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO LTDA - FILIAL', 'terminal': '21420', 'ip': '10.109.99.87', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO LTDA - FILIAL', 'terminal': '21421', 'ip': '10.180.137.37', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO LTDA - FILIAL', 'terminal': '21428', 'ip': '10.190.132.51', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO LTDA - FILIAL', 'terminal': '21504', 'ip': '10180137191', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO LTDA - FILIAL', 'terminal': '21561', 'ip': '10.109.36.193', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO LTDA - FILIAL', 'terminal': '21574', 'ip': '10190133192', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO LTDA - FILIAL', 'terminal': '21585', 'ip': '10180137219', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO LTDA - FILIAL', 'terminal': '21586', 'ip': '10190133151', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO LTDA - FILIAL', 'terminal': '21587', 'ip': '10190132214', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO LTDA - FILIAL', 'terminal': '21592', 'ip': '10.180.137.98', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO LTDA - FILIAL', 'terminal': '21594', 'ip': '10.190.132.99', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO LTDA - FILIAL', 'terminal': '21602', 'ip': '10.190.133.58', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO LTDA - FILIAL', 'terminal': '21628', 'ip': '10180136248', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO LTDA - FILIAL', 'terminal': '21689', 'ip': '10109126131', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO LTDA - FILIAL', 'terminal': '21701', 'ip': '10.109.38.73', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO LTDA - FILIAL', 'terminal': '21703', 'ip': '10180136133', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO LTDA - FILIAL', 'terminal': '21705', 'ip': '10.109.39.246', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO LTDA - FILIAL', 'terminal': '21710', 'ip': '10.109.96.143', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO LTDA - FILIAL', 'terminal': '21718', 'ip': '172.102.84.214', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO LTDA - FILIAL', 'terminal': '21976', 'ip': '10109123157', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO LTDA - FILIAL', 'terminal': '21980', 'ip': '10180137192', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO LTDA - FILIAL', 'terminal': '21982', 'ip': '10.109.72.83', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO LTDA - FILIAL', 'terminal': '21987', 'ip': '10180137113', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO LTDA - FILIAL', 'terminal': '4', 'ip': '10.109.124.55', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO MATRIZ LTDA', 'terminal': '2', 'ip': '10.109.38.116', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO MATRIZ LTDA', 'terminal': '21278', 'ip': '10.109.83.97', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO MATRIZ LTDA', 'terminal': '21335', 'ip': '10.109.21.93', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO MATRIZ LTDA', 'terminal': '21348', 'ip': '10.109.99.176', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO MATRIZ LTDA', 'terminal': '21351', 'ip': '10.109.12.82', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO MATRIZ LTDA', 'terminal': '21352', 'ip': '10.109.81.234', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO MATRIZ LTDA', 'terminal': '21373', 'ip': '10.109.36.60', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO MATRIZ LTDA', 'terminal': '21382', 'ip': '10.109.98.89', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO MATRIZ LTDA', 'terminal': '21392', 'ip': '10190133132', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO MATRIZ LTDA', 'terminal': '21403', 'ip': '10.180.136.65', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO MATRIZ LTDA', 'terminal': '21406', 'ip': '10190132148', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO MATRIZ LTDA', 'terminal': '21409', 'ip': '10.109.31.114', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO MATRIZ LTDA', 'terminal': '21554', 'ip': '10180137147', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO OSASCO MATRIZ LTDA', 'terminal': '5', 'ip': '10.109.23.210', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11010', 'ip': '10.190.133.46', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11401', 'ip': '10.109.32.48', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11402', 'ip': '10.109.99.28', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11404', 'ip': '10.109.98.22', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11407', 'ip': '10.109.22.228', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11408', 'ip': '10.109.33.48', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11411', 'ip': '10.109.98.154', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11413', 'ip': '10.109.23.76', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11417', 'ip': '10.109.83.125', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11424', 'ip': '10.109.20.241', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11428', 'ip': '10.109.20.63', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11430', 'ip': '10.109.35.75', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11432', 'ip': '10.109.99.204', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11434', 'ip': '10.109.35.189', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11441', 'ip': '10.109.83.133', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11448', 'ip': '10.109.32.143', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11454', 'ip': '10.109.10.103', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11457', 'ip': '10.109.80.99', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11460', 'ip': '10.109.43.11', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11461', 'ip': '10.109.21.209', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11467', 'ip': '10.109.81.19', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11472', 'ip': '10.109.32.155', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11481', 'ip': '10.109.9.7', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11484', 'ip': '10.109.5.103', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11489', 'ip': '10109123122', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11508', 'ip': '10.109.10.141', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11515', 'ip': '10.109.96.49', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11521', 'ip': '10.109.32.54', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11526', 'ip': '10.109.19.54', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11529', 'ip': '10.109.98.51', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11564', 'ip': '10.109.96.228', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11566', 'ip': '10109121135', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11567', 'ip': '10.109.97.135', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11588', 'ip': '10.109.43.23', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11599', 'ip': '10.109.82.24', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11610', 'ip': '10109124134', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11741', 'ip': '10.109.8.19', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11751', 'ip': '10.109.30.194', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11755', 'ip': '10.109.28.140', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA', 'terminal': '11799', 'ip': '10.109.32.31', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA MUNICIPAL', 'terminal': '1800', 'ip': '10.109.11.240', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA MUNICIPAL', 'terminal': '1801', 'ip': '10.180.136.36', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA MUNICIPAL', 'terminal': '1871', 'ip': '10.109.10.28', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA MUNICIPAL', 'terminal': '1884', 'ip': '10.190.133.49', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO PIRAJUÇARA LTDA MUNICIPAL', 'terminal': '1908', 'ip': '10.109.11.100', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO RAPOSO TAVARES LTDA', 'terminal': '12116', 'ip': '10.109.31.196', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO RAPOSO TAVARES LTDA', 'terminal': '12154', 'ip': '10.109.37.91', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO RAPOSO TAVARES LTDA', 'terminal': '122307', 'ip': '10.109.9.237', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO RAPOSO TAVARES LTDA', 'terminal': '122312', 'ip': '10.190.133.72', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO RAPOSO TAVARES LTDA', 'terminal': '122315', 'ip': '10190133251', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO RAPOSO TAVARES LTDA', 'terminal': '12243', 'ip': '10.109.96.245', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO RAPOSO TAVARES LTDA', 'terminal': '122512', 'ip': '10.109.28.142', 'model': 'B80', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO RAPOSO TAVARES LTDA', 'terminal': '122529', 'ip': '10.109.43.72', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO RAPOSO TAVARES LTDA', 'terminal': '122531', 'ip': '10190132184', 'model': 'B80', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO RAPOSO TAVARES LTDA', 'terminal': '122532', 'ip': '10.109.6.26', 'model': 'B80', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO RAPOSO TAVARES LTDA', 'terminal': '122540', 'ip': '10.180.134.79', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO RAPOSO TAVARES LTDA', 'terminal': '122558', 'ip': '10.109.38.21', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO RAPOSO TAVARES LTDA', 'terminal': '122559', 'ip': '10.109.37.252', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO RAPOSO TAVARES LTDA', 'terminal': '122564', 'ip': '10.180.136.91', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO RAPOSO TAVARES LTDA', 'terminal': '122567', 'ip': '10.180.137.93', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO RAPOSO TAVARES LTDA', 'terminal': '12316', 'ip': '10.109.31.41', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO RAPOSO TAVARES LTDA', 'terminal': '12729', 'ip': '10.109.9.75', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO RAPOSO TAVARES LTDA', 'terminal': '9997', 'ip': '10.109.97.80', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIACAO RAPOSO TAVARES MUNICIPAL', 'terminal': '1314', 'ip': '10.109.97.159', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIACAO RAPOSO TAVARES MUNICIPAL', 'terminal': '2114', 'ip': '10.109.9.54', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIACAO RAPOSO TAVARES MUNICIPAL', 'terminal': '2352', 'ip': '10.109.28.16', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIACAO RAPOSO TAVARES MUNICIPAL', 'terminal': '2406', 'ip': '10.109.10.243', 'model': 'B80', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIACAO RAPOSO TAVARES MUNICIPAL', 'terminal': '2414', 'ip': '10.109.42.82', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIACAO RAPOSO TAVARES MUNICIPAL', 'terminal': '9996', 'ip': '10.109.21.173', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO TALISMÃ - MUNICIPAL', 'terminal': '1200', 'ip': '10.109.18.72', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO TALISMÃ - MUNICIPAL', 'terminal': '1400', 'ip': '10.109.6.56', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO TALISMÃ - MUNICIPAL', 'terminal': '140505', 'ip': '172.70.36.217', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO TALISMÃ - MUNICIPAL', 'terminal': '151041', 'ip': '10.109.43.248', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO TALISMÃ - MUNICIPAL', 'terminal': '191249', 'ip': '10.109.10.239', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO TALISMÃ - MUNICIPAL', 'terminal': '220895', 'ip': '10.109.81.158', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO TALISMÃ - MUNICIPAL', 'terminal': '311224', 'ip': '10.109.5.190', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO TALISMÃ - MUNICIPAL', 'terminal': '50712', 'ip': '10.109.18.78', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO TALISMÃ - MUNICIPAL', 'terminal': '61025', 'ip': '10.109.98.0', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO TALISMÃ - MUNICIPAL', 'terminal': '71087', 'ip': '10.109.14.200', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO TALISMÃ - MUNICIPAL', 'terminal': '801', 'ip': '10.109.33.72', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO TALISMÃ - MUNICIPAL', 'terminal': '821', 'ip': '10.109.96.185', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO TALISMÃ - MUNICIPAL', 'terminal': '88804', 'ip': '10.109.8.253', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO TALISMÃ - MUNICIPAL', 'terminal': '901', 'ip': '10.190.5.127', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO TRANSDUTRA LTDA', 'terminal': '32003', 'ip': '10.109.82.239', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO TRANSDUTRA LTDA', 'terminal': '32004', 'ip': '10.109.11.172', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO TRANSDUTRA LTDA', 'terminal': '32005', 'ip': '10.109.82.56', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO TRANSDUTRA LTDA', 'terminal': '32009', 'ip': '10.109.96.25', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO TRANSDUTRA LTDA', 'terminal': '32027', 'ip': '10.109.22.181', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO TRANSDUTRA LTDA', 'terminal': '32029', 'ip': '10.109.34.241', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO TRANSDUTRA LTDA', 'terminal': '32037', 'ip': '10109114128', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO TRANSDUTRA LTDA', 'terminal': '32042', 'ip': '10.109.99.2', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO TRANSDUTRA LTDA', 'terminal': '32052', 'ip': '10.109.61.120', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO TRANSDUTRA LTDA', 'terminal': '32641', 'ip': '10.109.16.36', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO TRANSDUTRA LTDA', 'terminal': '32642', 'ip': '10.109.98.240', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO TRANSDUTRA LTDA', 'terminal': '32652', 'ip': '10.109.35.85', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO TRANSDUTRA LTDA', 'terminal': '32659', 'ip': '10.109.82.171', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO TRANSDUTRA LTDA', 'terminal': '32662', 'ip': '10.180.8.162', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO TRANSDUTRA LTDA', 'terminal': '32667', 'ip': '10.109.80.221', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}, {'company': 'VIAÇÃO TRANSDUTRA LTDA', 'terminal': '32670', 'ip': '10.109.5.183', 'model': 'V3695', 'sam_type': 'NÃO MIGRADO'}]

def _seed_garage_chip_base():
    if GarageChipBase.query.count(): return
    for r in GARAGE_CHIP_SEED:
        db.session.add(GarageChipBase(**r))
    db.session.commit()

def _garage_payload(force=False):
    now=time.monotonic()
    with _V63_GARAGE_CACHE_LOCK:
        cached=_V63_GARAGE_CACHE.get('payload')
        if cached is not None and not force and now-float(_V63_GARAGE_CACHE.get('at') or 0)<V63_GARAGE_CACHE_TTL: return cached
    _seed_garage_chip_base()
    bases=GarageChipBase.query.filter_by(active=True).order_by(GarageChipBase.company,GarageChipBase.terminal).all()
    swaps=GarageChipSwap.query.all(); sm={x.base_id:x for x in swaps}
    uids={x.technician_id for x in swaps if x.technician_id}; users={u.id:u for u in User.query.filter(User.id.in_(uids)).all()} if uids else {}
    sids=[x.id for x in swaps]; photos={}
    if sids:
        for p in GarageChipPhoto.query.filter(GarageChipPhoto.swap_id.in_(sids)).all(): photos.setdefault(p.swap_id,[]).append(p)
    # V66 REV4 Performance 2.0: mapa operacional carregado uma única vez.
    # Antes, _op_active_map('garagem') era executado para cada terminal (N+1).
    op_map=_op_active_map('garagem')
    out=[]
    for b in bases:
        sw=sm.get(b.id); ph=photos.get(sw.id,[]) if sw else []; op_item=op_map.get(str(b.terminal or '')); st=sw.status if sw else ((op_item.desired_status if op_item else None) or 'PENDENTE'); u=users.get(sw.technician_id) if sw else None
        out.append({'id':b.id,'company':b.company,'terminal':b.terminal,'model':b.model or '', 'status':st,'technician':u.name if u else '', 'test_result':sw.test_result if sw else '', 'notes':sw.notes if sw else '', 'photo_count':len(ph),'photos':[{'id':p.id,'url':'/uploads/'+p.stored_name,'thumb_url':'/uploads/'+p.stored_name+'?thumb=1'} for p in ph]})
    with _V63_GARAGE_CACHE_LOCK:
        _V63_GARAGE_CACHE['at']=now; _V63_GARAGE_CACHE['payload']=out
    return out

@app.get('/troca-chips-garagem')
@login_required
def garage_chip_page():
    return render_template('garage_chip_swap.html',app_release=APP_RELEASE)

@app.get('/api/garage-chip-swaps')
@login_required
def garage_chip_list_api():
    rows=_garage_payload(); total=len(rows); done=sum(x['status']=='CONCLUÍDA' for x in rows); prog=sum(x['status']=='EM ANDAMENTO' for x in rows)
    return jsonify({'ok':True,'rows':rows,'summary':{'total':total,'concluded':done,'in_progress':prog,'pending':total-done-prog,'percent':round(done*100/total,1) if total else 0}})

@app.post('/api/garage-chip-swaps/<int:base_id>')
@field_required
def garage_chip_save_api(base_id):
    if _activity_request_too_large(): return jsonify({"ok":False,"error":f"Envio excede {_ACTIVITY_REQUEST_MAX_MB} MB. Envie menos fotos por vez."}),413
    b=db.session.get(GarageChipBase,base_id)
    if not b:return jsonify({'ok':False,'error':'Terminal não encontrado na base.'}),404
    sw=GarageChipSwap.query.filter_by(base_id=base_id).first()
    if not sw:
        sw=GarageChipSwap(base_id=base_id,technician_id=session['user_id'],status='EM ANDAMENTO',started_at=datetime.utcnow());db.session.add(sw);db.session.flush()
    sw.technician_id=session['user_id']; sw.test_result=(request.form.get('test_result') or '').strip(); sw.notes=(request.form.get('notes') or '').strip(); sw.updated_at=datetime.utcnow()
    files=[f for f in request.files.getlist('photos') if f and f.filename]
    for f in files:
        safe=secure_filename(f.filename) or f'garage_{secrets.token_hex(4)}.jpg'; stored=f'garage_{sw.id}_{secrets.token_hex(6)}_{safe}'
        stored=_store_uploaded_file(f,'garage-chip-swaps',stored,f.mimetype or 'application/octet-stream')
        db.session.add(GarageChipPhoto(swap_id=sw.id,original_name=f.filename,stored_name=stored,mime_type=f.mimetype,uploaded_by=session['user_id']))
    if files or GarageChipPhoto.query.filter_by(swap_id=sw.id).count(): sw.status='CONCLUÍDA'; sw.completed_at=sw.completed_at or datetime.utcnow()
    db.session.commit(); _v63_invalidate_garage_cache(); return jsonify({'ok':True,'status':sw.status})

@app.get('/api/garage-chip-swaps/dashboard')
@login_required
def garage_chip_dashboard_api():
    rows=_garage_payload(); total=len(rows); done=sum(x['status']=='CONCLUÍDA' for x in rows); prog=sum(x['status']=='EM ANDAMENTO' for x in rows)
    by_company={}; by_tech={}; by_model={}; by_status={}; by_result={}
    for x in rows:
        c=by_company.setdefault(x['company'],{'company':x['company'],'name':x['company'],'total':0,'concluded':0,'count':0});c['total']+=1;c['count']+=1;c['concluded']+=x['status']=='CONCLUÍDA'
        if x['technician']:
            t=by_tech.setdefault(x['technician'],{'name':x['technician'],'total':0,'concluded':0,'count':0});t['total']+=1;t['count']+=1;t['concluded']+=x['status']=='CONCLUÍDA'
        mk=x.get('model') or 'Não informado'; by_model[mk]=by_model.get(mk,0)+1
        st=x.get('status') or 'PENDENTE'; by_status[st]=by_status.get(st,0)+1
        rs=x.get('test_result') or 'Não informado'; by_result[rs]=by_result.get(rs,0)+1
    pack=lambda d:[{'name':k,'count':v} for k,v in sorted(d.items(),key=lambda z:z[1],reverse=True)]
    return jsonify({'ok':True,'summary':{'total':total,'concluded':done,'in_progress':prog,'pending':total-done-prog,'percent':round(done*100/total,1) if total else 0},'companies':list(by_company.values()),'technicians':list(by_tech.values()),'models':pack(by_model),'statuses':pack(by_status),'results':pack(by_result)})

@app.get('/api/garage-chip-swaps/export.xlsx')
@login_required
def garage_chip_export_xlsx():
    rows=_garage_payload(); company=(request.args.get('company') or '').strip(); model=(request.args.get('model') or '').strip(); status=(request.args.get('status') or '').strip(); technician=(request.args.get('technician') or '').strip()
    rows=[x for x in rows if (not company or x.get('company')==company) and (not model or x.get('model')==model) and (not status or x.get('status')==status) and (not technician or x.get('technician')==technician)]
    wb=Workbook(); ws=wb.active; ws.title='Troca Chips Garagem'; ws.append(['Empresa','Terminal','Modelo','Status','Resultado','Técnico','Observações','Evidências'])
    for x in rows: ws.append([x.get('company'),x.get('terminal'),x.get('model'),x.get('status'),x.get('test_result'),x.get('technician'),x.get('notes'),len(x.get('photos') or [])])
    for cell in ws[1]: cell.font=Font(bold=True,color='FFFFFF'); cell.fill=PatternFill('solid',fgColor='17365D')
    ws.freeze_panes='A2'; ws.auto_filter.ref=ws.dimensions
    for col in range(1,ws.max_column+1): ws.column_dimensions[get_column_letter(col)].width=24
    bio=io.BytesIO(); wb.save(bio); bio.seek(0); return send_file(bio,as_attachment=True,download_name=f"troca_chips_garagem_{datetime.now().strftime('%Y%m%d_%H%M')}.xlsx",mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')

@app.get('/troca-chips-garagem/dashboard')
@login_required
def garage_chip_dashboard_page():
    # V69.3.2: Dashboard Garagem faz parte da Central /gerencial, como Recarga e EMV.
    return redirect(url_for('manager', view='garage'))

@app.get("/troca-chips-emv")
@emv_field_required
def emv_chip_page():
    if not _has_access("implantation.emv"): abort(403)
    return render_template("emv_chip_swap.html",app_release=APP_RELEASE)

def _v63_build_emv_payload(force=False, include_photos=True):
    """V66 REV4: monta EMV em dois modos de cache.

    O dashboard usa o modo slim e não consulta nem materializa evidências. A tela
    operacional usa o modo full. Isso evita pagar o custo de fotos/URLs em toda
    abertura de dashboard e reduz serialização/alocação de memória.
    """
    mode="full" if include_photos else "slim"
    now=time.monotonic()
    with _V63_EMV_CACHE_LOCK:
        slot=_V63_EMV_CACHE.get(mode) or {}
        cached=slot.get("payload")
        if cached is not None and not force and now-float(slot.get("at") or 0)<V63_EMV_CACHE_TTL:
            return cached
    _ensure_emv_tables()
    swap_rows=EmvChipSwap.query.all(); swaps={x.terminal:x for x in swap_rows}; rows=[]
    user_ids={uid for x in swap_rows for uid in (x.technician_id,getattr(x,"completed_by_id",None)) if uid}
    users={u.id:u for u in User.query.filter(User.id.in_(user_ids)).all()} if user_ids else {}
    station_names={}
    for nr in _load_station_network_rows():
        sn=str(nr.get("station") or "")
        if " - " in sn: station_names.setdefault(sn.split(" - ",1)[0].strip().upper(),sn.split(" - ",1)[1])
    photo_map={}
    if include_photos:
        swap_ids=[x.id for x in swap_rows]
        if swap_ids:
            for ph in EmvChipSwapPhoto.query.filter(EmvChipSwapPhoto.swap_id.in_(swap_ids)).order_by(EmvChipSwapPhoto.created_at).all():
                photo_map.setdefault(ph.swap_id,[]).append(ph)
    base_rows=_v41_emv_rows()
    for r in base_rows:
        sw=swaps.get(r["terminal"]); d=dict(r)
        station_code=str(r.get("station") or "").split('-')[-1].strip().upper(); station_name=station_names.get(station_code,"")
        photos=[]
        if include_photos and sw:
            photos=[{"id":ph.id,"name":ph.original_name,"url":url_for("uploaded",name=ph.stored_name),"thumb_url":url_for("uploaded",name=ph.stored_name,thumb=1),"created_at":ph.created_at.isoformat() if ph.created_at else None} for ph in photo_map.get(sw.id,[])]
        tech=users.get(sw.technician_id) if sw else None; completer=users.get(getattr(sw,"completed_by_id",None)) if sw else None
        d.update({"status":sw.status if sw else (d.get("_base_status") or "PENDENTE"),"test_result":sw.test_result if sw else None,"notes":sw.notes if sw else "","swap_id":sw.id if sw else None,"station_name":station_name,"photos":photos if include_photos else [],"photo_count":len(photos) if include_photos else None,"technician":tech.name if tech else "","completed_by":completer.name if completer else (tech.name if sw and sw.completed_at and tech else ""),"completed_by_role":completer.role if completer else (tech.role if sw and sw.completed_at and tech else ""),"completed_at":sw.completed_at.isoformat() if sw and sw.completed_at else None,"updated_at":sw.updated_at.isoformat() if sw and sw.updated_at else None})
        rows.append(d)
    base_terms={str(x.get("terminal") or "") for x in base_rows}
    has_op_base=OperationalBaseItem.query.filter_by(module="EMV").with_entities(OperationalBaseItem.id).first() is not None
    op_active_terms=set(_op_active_map("emv").keys()) if has_op_base else set()
    for sw in swap_rows:
        if not getattr(sw,"manual_entry",False) or sw.terminal in base_terms: continue
        if op_active_terms and sw.terminal not in op_active_terms and normalize(sw.status or "").upper()!="CONCLUIDA": continue
        photos=[]
        if include_photos:
            photos=[{"id":ph.id,"name":ph.original_name,"url":url_for("uploaded",name=ph.stored_name),"thumb_url":url_for("uploaded",name=ph.stored_name,thumb=1),"created_at":ph.created_at.isoformat() if ph.created_at else None} for ph in photo_map.get(sw.id,[])]
        tech=users.get(sw.technician_id); completer=users.get(getattr(sw,"completed_by_id",None))
        rows.append({"company":sw.company or "","line":sw.line or "","station":sw.station or "","station_name":sw.station or "","terminal":sw.terminal,"block_number":sw.block_number or "","version":"","ip":"","mask":"","gateway":"","dns1":"","dns2":"","group":"","manual_entry":True,"status":sw.status or "PENDENTE","test_result":sw.test_result,"notes":sw.notes or "","swap_id":sw.id,"photos":photos if include_photos else [],"photo_count":len(photos) if include_photos else None,"technician":tech.name if tech else "","completed_by":completer.name if completer else "","completed_at":sw.completed_at.isoformat() if sw.completed_at else None,"updated_at":sw.updated_at.isoformat() if sw.updated_at else None})
    with _V63_EMV_CACHE_LOCK:
        _V63_EMV_CACHE[mode]={"at":time.monotonic(),"payload":rows}
    return rows

@app.get("/api/emv-chip-swaps/")
@login_required
def emv_chip_list_legacy_slash():
    # V66 REV4.2: compatibilidade para clientes antigos sem redirecionamento.
    # A REV4.1 mostrou que o 308 mantinha uma segunda trilha de chamadas na
    # telemetria. A URL legada agora devolve diretamente o payload SLIM, sem
    # fotos/evidências e sem executar o pipeline full. Clientes atuais usam a
    # rota canônica sem barra.
    rows=list(_v63_build_emv_payload(include_photos=False))
    slim=[]
    for x in rows:
        slim.append({k:x.get(k) for k in (
            "company","line","station","station_name","terminal",
            "block_number","model","version","status","test_result",
            "technician","completed_by","completed_at","manual_entry"
        )})
    resp=jsonify({"ok":True,"rows":slim,"release":APP_RELEASE,"legacy":True})
    resp.headers["X-Autopass-Canonical"]="/api/emv-chip-swaps"
    resp.headers["X-Autopass-Deprecated"]="1"
    resp.headers["X-Autopass-Payload-Mode"]="slim-legacy"
    return resp

@app.get("/api/emv-chip-swaps", strict_slashes=True)
@login_required
def emv_chip_list():
    company=(request.args.get("company") or "").strip(); line=(request.args.get("line") or "").strip(); station=(request.args.get("station") or "").strip(); status=(request.args.get("status") or "").strip(); terminal=(request.args.get("terminal") or "").strip(); compact=request.args.get("compact") in ("1","true","yes"); include_photos=request.args.get("include_photos","1") not in ("0","false","no")
    # Dashboard/compact nunca precisa montar evidências completas.
    build_photos=include_photos and not compact
    rows=list(_v63_build_emv_payload(include_photos=build_photos))
    if company: rows=[x for x in rows if x.get("company")==company]
    if line: rows=[x for x in rows if x.get("line")==line]
    if station: rows=[x for x in rows if station in (x.get("station"),x.get("station_name"))]
    if status: rows=[x for x in rows if normalize(x.get("status"))==normalize(status)]
    if terminal: rows=[x for x in rows if terminal.lower() in str(x.get("terminal") or "").lower()]
    if compact or not include_photos:
        slim=[]
        for x in rows:
            if compact:
                d={k:x.get(k) for k in ("company","line","station","station_name","terminal","block_number","model","version","status","test_result","technician","completed_by","completed_at","manual_entry")}
            else: d=dict(x)
            d.pop("photos",None)
            # No modo slim a contagem de fotos não é consultada propositalmente.
            if d.get("photo_count") is None: d.pop("photo_count",None)
            slim.append(d)
        rows=slim
    resp=jsonify({"ok":True,"rows":rows,"release":APP_RELEASE})
    resp.headers["X-Autopass-Cache-TTL"]=str(V63_EMV_CACHE_TTL)
    resp.headers["X-Autopass-Payload-Mode"]="full" if build_photos else "slim"
    return resp

@app.post("/api/emv-chip-swaps/manual")
@emv_field_required
def emv_chip_manual_create():
    _ensure_emv_tables(); d=request.get_json(silent=True) or {}
    company=(d.get("company") or "").strip(); line=(d.get("line") or "").strip(); station=(d.get("station") or "").strip(); block=(d.get("block_number") or "").strip()
    if not all((company,line,station,block)): return jsonify({"ok":False,"error":"Informe operadora, linha, estação/localidade e número do bloqueio."}),400
    terminal=(d.get("terminal") or f"MANUAL-{normalize(company)[:12]}-{normalize(line)[:12]}-{normalize(station)[:16]}-{normalize(block)}").strip()[:120]
    if EmvChipSwap.query.filter_by(terminal=terminal).first(): return jsonify({"ok":False,"error":"Este bloqueio manual já foi incluído."}),409
    sw=EmvChipSwap(terminal=terminal,technician_id=session["user_id"],status="PENDENTE",manual_entry=True,company=company,line=line,station=station,block_number=block,notes=(d.get("notes") or "").strip(),updated_at=datetime.utcnow())
    db.session.add(sw);db.session.flush();db.session.add(AuditEvent(user_id=session.get("user_id"),event_type="EMV_MANUAL_CREATE",entity_type="emv_chip_swap",entity_id=str(sw.id),detail=f"{company} · {line} · {station} · bloqueio {block}"));db.session.commit()
    _v63_invalidate_emv_cache()
    return jsonify({"ok":True,"terminal":terminal,"id":sw.id})

@app.post("/api/emv-chip-swaps/<terminal>")
@emv_field_required
def emv_chip_save(terminal):
    # V62 REV2: proteção contra OOM em upload de fotos.
    # O limite é validado antes de request.form/request.files serem materializados.
    if _activity_request_too_large():
        return jsonify({"ok":False,"error":f"Envio excede {_ACTIVITY_REQUEST_MAX_MB} MB. Envie menos fotos por vez."}),413
    _ensure_emv_tables()
    base=_v41_emv_by_terminal().get(str(terminal))
    sw=EmvChipSwap.query.filter_by(terminal=terminal).first()
    if not base and not (sw and getattr(sw,"manual_entry",False)):
        return jsonify({"ok":False,"error":"Bloqueio EMV não encontrado na base."}),404
    if sw and (sw.status or "").upper().replace("CONCLUIDA","CONCLUÍDA")=="CONCLUÍDA" and session.get("role") in ("technician","technician_implantation"):
        return jsonify({"ok":False,"error":"Registro concluído e bloqueado. Solicite ao Gestor/ADM a reabertura para EM ANDAMENTO."}),409
    files=[x for x in request.files.getlist("photos") if x and x.filename]
    try:
        # Valida cada arquivo antes de alterar o registro.
        for f in files:
            size=_uploaded_file_size(f)
            if size and size > _ACTIVITY_UPLOAD_MAX_MB*1024*1024:
                return jsonify({"ok":False,"error":f"A foto {f.filename} excede {_ACTIVITY_UPLOAD_MAX_MB} MB. Reduza a resolução e tente novamente."}),413
        if not sw:
            sw=EmvChipSwap(terminal=terminal,technician_id=session["user_id"])
            db.session.add(sw);db.session.flush()
        sw.technician_id=session["user_id"]
        sw.test_result=(request.form.get("test_result") or "").strip()
        sw.notes=(request.form.get("notes") or "").strip()
        sw.latitude=_optional_float(request.form.get("latitude"))
        sw.longitude=_optional_float(request.form.get("longitude"))
        sw.gps_accuracy=_optional_float(request.form.get("gps_accuracy"))
        sw.updated_at=datetime.utcnow()
        if sw.test_result and sw.test_result!="TESTADO_OK" and not sw.notes:
            db.session.rollback()
            return jsonify({"ok":False,"error":"Para resultado diferente de OK, a observação é obrigatória."}),400
        for f in files:
            safe=secure_filename(f.filename) or f"emv_{secrets.token_hex(4)}.jpg"
            stored=f"emv_{sw.id}_{secrets.token_hex(6)}_{safe}"
            stored=_store_uploaded_file(f,"emv-chip-swaps",stored,f.mimetype or "application/octet-stream")
            db.session.add(EmvChipSwapPhoto(swap_id=sw.id,original_name=f.filename,stored_name=stored,mime_type=f.mimetype,uploaded_by=session["user_id"]))
            # Libera referências do upload antes de processar a próxima evidência.
            try: f.close()
            except Exception: pass
        db.session.flush()
        photos=EmvChipSwapPhoto.query.filter_by(swap_id=sw.id).count()
        sw.status="CONCLUÍDA" if photos else "PENDENTE"
        sw.completed_at=datetime.utcnow() if photos else None
        sw.completed_by_id=session.get("user_id") if photos else None
        db.session.add(AuditEvent(user_id=session.get("user_id"),event_type="EMV_CHIP_SWAP_UPDATE",entity_type="emv_chip_swap",entity_id=str(sw.id),detail=f"{terminal} · {sw.status} · teste {sw.test_result or '—'} · {photos} foto(s)"))
        db.session.commit()
        _v63_invalidate_emv_cache()
        return jsonify({"ok":True,"status":sw.status,"photo_count":photos})
    except ValueError as exc:
        db.session.rollback()
        return jsonify({"ok":False,"error":str(exc)}),413
    except Exception as exc:
        db.session.rollback()
        app.logger.exception("V62 REV2 falha controlada ao salvar EMV %s",terminal)
        return jsonify({"ok":False,"error":"Não foi possível salvar a troca EMV. Tente novamente com menos fotos.","detail":str(exc)[:160]}),500

@app.post("/api/emv-chip-swaps/<terminal>/admin-status")
@login_required
def emv_chip_admin_status(terminal):
    role=session.get("role")
    if role not in ("manager","manager_field"): return jsonify({"ok":False,"error":"Alteração administrativa restrita ao Gestor/ADM."}),403
    _ensure_emv_tables(); sw=EmvChipSwap.query.filter_by(terminal=terminal).first(); d=request.get_json(silent=True) or {}; new=(d.get("status") or "").strip().upper().replace("CONCLUIDA","CONCLUÍDA"); reason=(d.get("reason") or "").strip()
    if not sw: return jsonify({"ok":False,"error":"Registro não encontrado."}),404
    if new not in {"PENDENTE","EM ANDAMENTO","CONCLUÍDA"}: return jsonify({"ok":False,"error":"Status inválido."}),400
    # ADM (role=manager) pode retificar o status sem justificativa; Gestor Field
    # mantém justificativa obrigatória para preservar a auditoria operacional.
    if role != "manager" and not reason: return jsonify({"ok":False,"error":"Informe o motivo da alteração."}),400
    old=sw.status; sw.status=new; sw.updated_at=datetime.utcnow()
    if new=="CONCLUÍDA": sw.completed_at=sw.completed_at or datetime.utcnow(); sw.completed_by_id=sw.completed_by_id or session.get("user_id")
    else: sw.completed_at=None
    db.session.add(AuditEvent(user_id=session.get("user_id"),event_type="EMV_CHIP_SWAP_ADMIN_STATUS",entity_type="emv_chip_swap",entity_id=str(sw.id),detail=f"{terminal} · {old} -> {new} · motivo: {reason or 'retificação ADM'}"));db.session.commit();return jsonify({"ok":True,"status":new})

@app.delete("/api/emv-chip-swaps/photos/<int:photo_id>")
@emv_field_required
def emv_chip_delete_photo(photo_id):
    _ensure_emv_tables(); ph=db.session.get(EmvChipSwapPhoto,photo_id)
    if not ph: return jsonify({"ok":False,"error":"Foto não encontrada."}),404
    sw=db.session.get(EmvChipSwap,ph.swap_id)
    _delete_stored_media(ph.stored_name); db.session.delete(ph); db.session.flush()
    remaining=EmvChipSwapPhoto.query.filter_by(swap_id=sw.id).count() if sw else 0
    if sw and remaining==0:
        sw.status="PENDENTE"; sw.completed_at=None; sw.updated_at=datetime.utcnow()
    db.session.commit()
    _v63_invalidate_emv_cache()
    return jsonify({"ok":True,"remaining":remaining,"status":sw.status if sw else "PENDENTE"})

@app.delete("/api/emv-chip-swaps/<terminal>")
@emv_field_required
def emv_chip_delete(terminal):
    _ensure_emv_tables(); sw=EmvChipSwap.query.filter_by(terminal=terminal).first()
    if not sw: return jsonify({"ok":False,"error":"Registro não encontrado."}),404
    for ph in EmvChipSwapPhoto.query.filter_by(swap_id=sw.id).all():
        _delete_stored_media(ph.stored_name); db.session.delete(ph)
    db.session.delete(sw); db.session.commit()
    _v63_invalidate_emv_cache()
    return jsonify({"ok":True})

def _cleanup_activity_photos(kind):
    if session.get('role') not in ('manager','manager_field'):
        return jsonify({'ok':False,'error':'Limpeza restrita ao Gestor/ADM.'}),403
    data=request.get_json(silent=True) or {}; confirm=str(data.get('confirm') or '').strip().upper()
    if confirm!='CONFIRMAR': return jsonify({'ok':False,'error':'Digite CONFIRMAR para remover as evidências temporárias.'}),400
    if kind=='recarga':
        swaps=ChipSwap.query.all(); total=len(swaps); done=sum(1 for x in swaps if (x.status or '').upper().startswith('CONCLU'))
        if total and done<total: return jsonify({'ok':False,'error':f'Campanha ainda não concluída: {done}/{total} registros concluídos.'}),409
        photos=ChipSwapPhoto.query.all()
    else:
        _ensure_emv_tables(); swaps=EmvChipSwap.query.all(); total=len(swaps); done=sum(1 for x in swaps if (x.status or '').upper().startswith('CONCLU'))
        if total and done<total: return jsonify({'ok':False,'error':f'Campanha ainda não concluída: {done}/{total} registros concluídos.'}),409
        photos=EmvChipSwapPhoto.query.all()
    removed=0
    for ph in photos:
        _delete_stored_media(ph.stored_name); db.session.delete(ph); removed+=1
    db.session.add(AuditEvent(user_id=session.get('user_id'),event_type='CAMPAIGN_EVIDENCE_PURGE',entity_type=kind,entity_id='',detail=f'{removed} foto(s) temporárias removidas após encerramento; registros operacionais preservados.'))
    db.session.commit(); _v63_invalidate_emv_cache(); return jsonify({'ok':True,'removed':removed,'records_preserved':total})

@app.post('/api/chip-swaps/purge-evidence')
@login_required
def chip_swap_purge_evidence(): return _cleanup_activity_photos('recarga')

@app.post('/api/emv-chip-swaps/purge-evidence')
@login_required
def emv_chip_purge_evidence(): return _cleanup_activity_photos('emv')


@app.get("/api/emv-chip-swaps/export.xlsx")
@login_required
def emv_chip_export():
    _ensure_emv_tables(); swaps={x.terminal:x for x in EmvChipSwap.query.all()}; wb=Workbook();ws=wb.active;ws.title="Troca EMV";headers=["Operadora","Linha","Estação","Terminal","Versão","IP","Máscara","Gateway","DNS 1","DNS 2","Grupo","Status","Resultado","Observação"];ws.append(headers)
    company=(request.args.get("company") or "").strip(); line=(request.args.get("line") or "").strip(); station=(request.args.get("station") or "").strip(); status_filter=(request.args.get("status") or "").strip()
    for r in _v41_emv_rows():
        sw=swaps.get(r["terminal"]); status=sw.status if sw else "PENDENTE"
        if company and r["company"] != company: continue
        if line and r["line"] != line: continue
        if station and r["station"] != station: continue
        if status_filter and status != status_filter: continue
        ws.append([r["company"],r["line"],r["station"],r["terminal"],r["version"],r["ip"],r["mask"],r["gateway"],r["dns1"],r["dns2"],r["group"],status,sw.test_result if sw else "",sw.notes if sw else ""])
    bio=io.BytesIO();wb.save(bio);bio.seek(0);return send_file(bio,as_attachment=True,download_name="troca_chips_emv.xlsx",mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")



def migrate_financial_v524_columns():
    """V52.5: amplia lançamentos sem perder histórico existente."""
    try:
        inspector=db.inspect(db.engine)
        if "financial_monthly_costs" not in inspector.get_table_names():
            return
        cols={c["name"] for c in inspector.get_columns("financial_monthly_costs")}
        commands=[]
        if "forecast_amount" not in cols: commands.append("ALTER TABLE financial_monthly_costs ADD COLUMN forecast_amount FLOAT")
        if "cost_center" not in cols: commands.append("ALTER TABLE financial_monthly_costs ADD COLUMN cost_center VARCHAR(60) DEFAULT 'SUPORTE_CAMPO'")
        if "project" not in cols: commands.append("ALTER TABLE financial_monthly_costs ADD COLUMN project VARCHAR(220)")
        if "service_text" not in cols: commands.append("ALTER TABLE financial_monthly_costs ADD COLUMN service_text VARCHAR(300)")
        if "invoice_number" not in cols: commands.append("ALTER TABLE financial_monthly_costs ADD COLUMN invoice_number VARCHAR(120)")
        if "cost_center_id" not in cols: commands.append("ALTER TABLE financial_monthly_costs ADD COLUMN cost_center_id VARCHAR(20)")
        user_cols={c["name"] for c in inspector.get_columns("users")} if "users" in inspector.get_table_names() else set()
        if "access_json" not in user_cols: commands.append("ALTER TABLE users ADD COLUMN access_json TEXT")
        if "system_profile_id" not in user_cols: commands.append("ALTER TABLE users ADD COLUMN system_profile_id INTEGER")
        # V55.2: cadastro financeiro enriquecido para importação de empresas.
        sup_cols={c["name"] for c in inspector.get_columns("financial_suppliers")} if "financial_suppliers" in inspector.get_table_names() else set()
        sup_commands=[]
        for col,sql in (("trade_name","VARCHAR(180)"),("cnpj","VARCHAR(30)"),("primary_cost_center","VARCHAR(60)"),("cost_center_id","VARCHAR(20)"),("contact_name","VARCHAR(180)"),("phone","VARCHAR(40)"),("email","VARCHAR(180)"),("pending_profile","BOOLEAN NOT NULL DEFAULT FALSE")):
            if col not in sup_cols: sup_commands.append(f"ALTER TABLE financial_suppliers ADD COLUMN {col} {sql}")
        commands.extend(sup_commands)
        # V56-B: dados físicos do processamento TBForte para análise de numerário/cédulas.
        cash_cols={c["name"] for c in inspector.get_columns("financial_cash_collections")} if "financial_cash_collections" in inspector.get_table_names() else set()
        for col,sql in (("processed_note_count","INTEGER"),("processed_media_type","VARCHAR(40)"),("processing_charge","FLOAT")):
            if col not in cash_cols: commands.append(f"ALTER TABLE financial_cash_collections ADD COLUMN {col} {sql}")
        # Índices de leitura pesada observados na Telemetria V60.
        tables=set(inspector.get_table_names())
        if "financial_atm_transactions" in tables:
            commands.extend([
                "CREATE INDEX IF NOT EXISTS ix_fin_tx_imported_at ON financial_atm_transactions (imported_at)",
                "CREATE INDEX IF NOT EXISTS ix_fin_tx_source_file ON financial_atm_transactions (source_file)",
            ])
        commands.extend([
            "CREATE INDEX IF NOT EXISTS ix_fin_monthly_supplier_comp ON financial_monthly_costs (supplier_id, competence)",
            "CREATE INDEX IF NOT EXISTS ix_fin_monthly_center_comp ON financial_monthly_costs (cost_center, competence)",
        ])
        for command in commands: db.session.execute(db.text(command))
        if commands:
            db.session.execute(db.text("UPDATE financial_monthly_costs SET cost_center='SUPORTE_CAMPO' WHERE cost_center IS NULL OR cost_center=''"))
            # Backfill dos IDs oficiais de centro de custo informados para V60 REV2.
            for cc in FIN_COST_CENTERS:
                db.session.execute(db.text("UPDATE financial_monthly_costs SET cost_center_id=:ccid WHERE UPPER(cost_center)=:ckey AND (cost_center_id IS NULL OR cost_center_id='')"), {"ccid":cc["id"],"ckey":cc["key"]})
                db.session.execute(db.text("UPDATE financial_suppliers SET cost_center_id=:ccid WHERE UPPER(primary_cost_center)=:ckey AND (cost_center_id IS NULL OR cost_center_id='')"), {"ccid":cc["id"],"ckey":cc["key"]})
            db.session.commit()
    except Exception:
        db.session.rollback(); raise



def _fin_norm_supplier_name(value):
    return re.sub(r"[^A-Z0-9]+","", normalize(value or ""))

def _fin_cost_center_id_for_key(key):
    return (FIN_COST_CENTER_BY_KEY.get((key or "").strip().upper()) or {}).get("id","")

def _fin_supplier_pending(row):
    return not bool((row.name or "").strip() and getattr(row,"cnpj",None) and getattr(row,"primary_cost_center",None) and getattr(row,"cost_center_id",None) and (getattr(row,"contact_name",None) or getattr(row,"email",None) or getattr(row,"phone",None)))

def _fin_supplier_duplicates(suppliers):
    groups={}
    for s in suppliers:
        cnpj=re.sub(r"\D","",getattr(s,"cnpj",None) or "")
        key=("CNPJ:"+cnpj) if len(cnpj)>=8 else ("NOME:"+_fin_norm_supplier_name(s.name))
        if key and key not in ("NOME:","CNPJ:"): groups.setdefault(key,[]).append(s)
    return {k:v for k,v in groups.items() if len(v)>1}

def _financial_admin_allowed():
    return session.get("role") in ("manager", "atm_financial_admin")


@app.get("/financeiro/dashboard")
@login_required
def financial_dashboard_page():
    if not _has_access("finance.dashboard"):
        abort(403)
    # V55.2: para gestores, a Dashboard Financeiro é painel do shell gerencial.
    if session.get("role") in ("manager","manager_field"):
        return redirect("/gerencial?view=financial-dashboard")
    return render_template("financial_dashboard.html", app_release=APP_RELEASE)

@app.get("/financeiro/dashboard/embed")
@login_required
def financial_dashboard_embed():
    if session.get("role") not in ("manager","manager_field","atm_financial_admin"):
        abort(403)
    return render_template("financial_dashboard.html", app_release=APP_RELEASE, embedded=True)

@app.get("/financeiro")
@login_required
def financial_home():
    return redirect(url_for("financial_cost_management_page"))

@app.get("/financeiro/implantacao")
@login_required
def financial_implantation_page():
    if session.get("role") not in ("manager","manager_field","atm_financial_admin"): return redirect(url_for("dashboard_landing"))
    return render_template("financial_area_placeholder.html", app_release=APP_RELEASE, area="Implantação")

@app.get("/financeiro/assistencia-tecnica")
@login_required
def financial_assistance_page():
    if session.get("role") not in ("manager","manager_field","atm_financial_admin"): return redirect(url_for("dashboard_landing"))
    return render_template("financial_area_placeholder.html", app_release=APP_RELEASE, area="Assistência Técnica")

@app.get("/financeiro/coleta-valores")
@app.get("/financeiro/suporte-campo/coleta-valores")
@login_required
def financial_cash_collection_page():
    if not _has_access("finance.collection"):
        return redirect(url_for("dashboard_landing"))
    if session.get("role") not in ("manager", "manager_field", "atm_financial_admin"):
        return redirect(url_for("dashboard_landing"))
    return render_template("financial_cash_collection.html", app_release=APP_RELEASE)

@app.get("/financeiro/lancamentos")
@app.get("/financeiro-atm/gestao")
@app.get("/financeiro/suporte-campo")
@login_required
def financial_cost_management_page():
    if not (_has_access("finance.entries") or _has_access("finance.support") or _has_access("finance.suppliers")) or not _financial_admin_allowed():
        return redirect(url_for("dashboard_landing"))
    return render_template("financial_cost_management.html", app_release=APP_RELEASE)

@app.get("/financeiro/lancamentos/embed")
@login_required
def financial_cost_management_embed():
    if session.get("role") not in ("manager","manager_field","atm_financial_admin"):
        abort(403)
    return render_template("financial_cost_management.html", app_release=APP_RELEASE, embedded=True)

@app.get("/api/v56a/performance")
@login_required
def v56a_performance_status():
    """Diagnóstico leve da normalização V56-A sem varrer objetos ORM completos."""
    if session.get("role") not in ("manager","manager_field","atm_financial_admin"):
        abort(403)
    total=db.session.query(func.count(TopDeskTicket.id)).scalar() or 0
    normalized=db.session.query(func.count(TopDeskTicket.id)).filter(TopDeskTicket.created_at.isnot(None)).scalar() or 0
    return jsonify({
      "ok":True,"release":APP_RELEASE,"topdesk":{
        "tickets":int(total),"normalized":int(normalized),"pending":int(max(0,total-normalized)),
        "normalized_pct":round(normalized*100/max(1,total),1),
        "dashboard_eligible":int(total),
        "backfill":{"running":bool(_V56A_BACKFILL.get("running")),"processed_this_boot":int(_V56A_BACKFILL.get("processed") or 0),"error":_V56A_BACKFILL.get("error")},
        "analytics_cache_ttl_seconds":TOPDESK_ANALYTICS_TTL,
        "dimensions":["created_at","line_code","station_code","model_code"]
      }
    })


@app.get("/api/release/routes-v553")
@login_required
def release_routes_v553():
    """Diagnóstico simples para confirmar que o backend V55.3.1 carregou as novas rotas."""
    checks = [
        "/financeiro/dashboard",
        "/financeiro/dashboard/embed",
        "/financeiro/lancamentos",
        "/financeiro/lancamentos/embed",
        "/financeiro/coleta-valores",
        "/dashboard/implantacao",
        "/dashboard/implantacao/embed",
    ]
    registered = {rule.rule for rule in app.url_map.iter_rules()}
    return jsonify({"ok": True, "release": APP_RELEASE, "routes": {x: x in registered for x in checks}})

@app.get("/api/financeiro/cadastros")
@login_required
def financial_catalog_api():
    if session.get("role") not in ("manager","manager_field","atm_financial_admin"):
        return jsonify({"ok":False,"error":"Sem permissão."}),403
    suppliers=FinancialSupplier.query.order_by(FinancialSupplier.name).all()
    services=FinancialService.query.order_by(FinancialService.name).all()
    raw_products=[x[0] for x in db.session.query(BaseAsset.equipment_type).distinct().all() if x and x[0]]
    aliases={"VALIDADOR":"RECARGA","VALIDADOR DE RECARGA":"RECARGA","POS DE BILHETERIA":"POS"}
    products=[]
    for value in raw_products:
        key=aliases.get(str(value).strip().upper(),str(value).strip().upper())
        if key and key not in products: products.append(key)
    for fallback in ("ATM","POS","RECARGA","RACK","BLOQUEIO","TDI","OUTROS"):
        if fallback not in products: products.append(fallback)
    dup_groups=_fin_supplier_duplicates(suppliers)
    dup_ids={s.id for grp in dup_groups.values() for s in grp}
    sup_rows=[]
    dirty=False
    for x in suppliers:
        pending=_fin_supplier_pending(x)
        if bool(getattr(x,"pending_profile",False))!=pending:
            x.pending_profile=pending; dirty=True
        sup_rows.append({"id":x.id,"name":x.name,"trade_name":getattr(x,"trade_name",None),"cnpj":getattr(x,"cnpj",None),
            "primary_cost_center":getattr(x,"primary_cost_center",None),"cost_center_id":getattr(x,"cost_center_id",None),
            "contact_name":getattr(x,"contact_name",None),"phone":getattr(x,"phone",None),"email":getattr(x,"email",None),
            "pending_profile":pending,"duplicate":x.id in dup_ids,"active":x.active})
    if dirty: db.session.commit()
    return jsonify({"ok":True,"suppliers":sup_rows,
        "services":[{"id":x.id,"supplier_id":x.supplier_id,"name":x.name,"description":x.description or "","category":x.category or "OUTROS","active":x.active} for x in services],
        "products":products,"cost_centers":FIN_COST_CENTERS,"duplicate_groups":len(dup_groups)})

@app.post("/api/financeiro/fornecedores")
@login_required
def financial_supplier_create_api():
    if not _financial_admin_allowed(): return jsonify({"ok":False,"error":"Sem permissão."}),403
    d=request.get_json(silent=True) or {}; name=(d.get("name") or "").strip()
    if not name: return jsonify({"ok":False,"error":"Informe a empresa/fornecedor."}),400
    cnpj=re.sub(r"\D","",(d.get("cnpj") or ""))
    candidates=FinancialSupplier.query.all()
    row=next((x for x in candidates if _fin_norm_supplier_name(x.name)==_fin_norm_supplier_name(name)),None)
    if not row and cnpj:
        row=next((x for x in candidates if re.sub(r"\D","",getattr(x,"cnpj",None) or "")==cnpj),None)
    if row:
        if not row.active: row.active=True; db.session.commit()
        return jsonify({"ok":True,"id":row.id,"existing":True,"message":"Cadastro já existente; o registro original foi preservado."})
    center=(d.get("primary_cost_center") or "").strip().upper() or None
    center_id=(d.get("cost_center_id") or _fin_cost_center_id_for_key(center)).strip().upper() or None
    row=FinancialSupplier(name=name,trade_name=(d.get("trade_name") or "").strip() or None,cnpj=(d.get("cnpj") or "").strip() or None,
        primary_cost_center=center,cost_center_id=center_id,contact_name=(d.get("contact_name") or "").strip() or None,
        phone=(d.get("phone") or "").strip() or None,email=(d.get("email") or "").strip() or None,created_by=session.get("user_id"))
    row.pending_profile=_fin_supplier_pending(row)
    db.session.add(row);db.session.flush()
    db.session.add(AuditEvent(user_id=session.get("user_id"),event_type="FIN_SUPPLIER_CREATE",entity_type="financial_supplier",entity_id=str(row.id),detail=name));db.session.commit()
    return jsonify({"ok":True,"id":row.id})

@app.put("/api/financeiro/fornecedores/<int:row_id>")
@login_required
def financial_supplier_update_api(row_id):
    if not _financial_admin_allowed(): return jsonify({"ok":False,"error":"Sem permissão."}),403
    row=db.session.get(FinancialSupplier,row_id)
    if not row: return jsonify({"ok":False,"error":"Fornecedor não encontrado."}),404
    d=request.get_json(silent=True) or {}
    old=row.name
    row.name=(d.get("name") or row.name or "").strip()
    row.trade_name=(d.get("trade_name") or "").strip() or None
    row.cnpj=(d.get("cnpj") or "").strip() or None
    row.primary_cost_center=(d.get("primary_cost_center") or "").strip().upper() or None
    row.cost_center_id=(d.get("cost_center_id") or _fin_cost_center_id_for_key(row.primary_cost_center)).strip().upper() or None
    row.contact_name=(d.get("contact_name") or "").strip() or None
    row.phone=(d.get("phone") or "").strip() or None
    row.email=(d.get("email") or "").strip() or None
    if "active" in d: row.active=bool(d.get("active"))
    row.pending_profile=_fin_supplier_pending(row)
    db.session.add(AuditEvent(user_id=session.get("user_id"),event_type="FIN_SUPPLIER_UPDATE",entity_type="financial_supplier",entity_id=str(row.id),detail=f"{old} -> {row.name}"));db.session.commit()
    return jsonify({"ok":True})

@app.get("/api/financeiro/pendencias-cadastro")
@login_required
def financial_pending_profiles_api():
    if session.get("role") not in ("manager","manager_field","atm_financial_admin"): return jsonify({"ok":False,"error":"Sem permissão."}),403
    suppliers=FinancialSupplier.query.filter(FinancialSupplier.active.isnot(False)).order_by(FinancialSupplier.name).all()
    duplicates=_fin_supplier_duplicates(suppliers); dup_ids={s.id for grp in duplicates.values() for s in grp}
    rows=[]
    for s in suppliers:
        missing=[]
        if not getattr(s,"cnpj",None): missing.append("CNPJ")
        if not getattr(s,"primary_cost_center",None): missing.append("Centro de custo")
        if not getattr(s,"cost_center_id",None): missing.append("ID centro de custo")
        if not (getattr(s,"contact_name",None) or getattr(s,"email",None) or getattr(s,"phone",None)): missing.append("Contato")
        if s.id in dup_ids: missing.append("Possível duplicidade")
        if missing:
            rows.append({"id":s.id,"name":s.name,"missing":missing,"active":s.active is not False,"duplicate":s.id in dup_ids})
    return jsonify({"ok":True,"rows":rows,"count":len(rows)})

@app.post("/api/financeiro/fornecedores/consolidar-duplicados")
@login_required
def financial_supplier_merge_duplicates_api():
    if session.get("role")!="manager": return jsonify({"ok":False,"error":"Consolidação restrita ao ADM."}),403
    suppliers=FinancialSupplier.query.order_by(FinancialSupplier.id).all()
    groups=_fin_supplier_duplicates(suppliers); merged=0
    try:
        for _,grp in groups.items():
            primary=sorted(grp,key=lambda s:(not bool(getattr(s,"cnpj",None)), not bool(getattr(s,"contact_name",None) or getattr(s,"email",None)), s.id))[0]
            for dup in grp:
                if dup.id==primary.id: continue
                services=FinancialService.query.filter_by(supplier_id=dup.id).all()
                for svc in services:
                    existing=FinancialService.query.filter(FinancialService.supplier_id==primary.id,func.lower(FinancialService.name)==(svc.name or "").lower()).first()
                    if existing:
                        FinancialMonthlyCost.query.filter_by(service_id=svc.id).update({"service_id":existing.id,"supplier_id":primary.id},synchronize_session=False)
                        db.session.delete(svc)
                    else:
                        svc.supplier_id=primary.id
                        FinancialMonthlyCost.query.filter_by(service_id=svc.id).update({"supplier_id":primary.id},synchronize_session=False)
                FinancialMonthlyCost.query.filter_by(supplier_id=dup.id).update({"supplier_id":primary.id},synchronize_session=False)
                for attr in ("trade_name","cnpj","primary_cost_center","cost_center_id","contact_name","phone","email"):
                    if not getattr(primary,attr,None) and getattr(dup,attr,None): setattr(primary,attr,getattr(dup,attr))
                db.session.delete(dup); merged+=1
            primary.pending_profile=_fin_supplier_pending(primary)
        db.session.add(AuditEvent(user_id=session.get("user_id"),event_type="FIN_SUPPLIER_DEDUP",entity_type="financial_supplier",entity_id="batch",detail=f"{merged} cadastro(s) duplicado(s) consolidado(s)"))
        db.session.commit()
        return jsonify({"ok":True,"merged":merged,"groups":len(groups)})
    except Exception as exc:
        db.session.rollback(); return jsonify({"ok":False,"error":str(exc)}),500

@app.get("/api/financeiro/fornecedores/<int:row_id>/padroes")
@login_required
def financial_supplier_patterns_api(row_id):
    if session.get("role") not in ("manager","manager_field","atm_financial_admin"): return jsonify({"ok":False,"error":"Sem permissão."}),403
    supplier=db.session.get(FinancialSupplier,row_id)
    if not supplier:return jsonify({"ok":False,"error":"Fornecedor não encontrado."}),404
    rows=FinancialMonthlyCost.query.filter_by(supplier_id=row_id).order_by(FinancialMonthlyCost.competence.desc(),FinancialMonthlyCost.updated_at.desc()).limit(8).all()
    sups={row_id:supplier.name}; svcs={x.id:x.name for x in FinancialService.query.filter_by(supplier_id=row_id).all()}
    return jsonify({"ok":True,"services":[{"id":x.id,"name":x.name} for x in FinancialService.query.filter_by(supplier_id=row_id,active=True).order_by(FinancialService.name).all()],
        "patterns":[_fin_payload(x,{},sups,svcs) for x in rows]})

@app.delete("/api/financeiro/fornecedores/<int:row_id>")
@login_required
def financial_supplier_delete_api(row_id):
    if session.get("role") not in ("manager","atm_financial_admin"): return jsonify({"ok":False,"error":"Sem permissão para excluir fornecedor."}),403
    row=db.session.get(FinancialSupplier,row_id)
    if not row: return jsonify({"ok":False,"error":"Fornecedor não encontrado."}),404
    name=row.name
    used=FinancialMonthlyCost.query.filter_by(supplier_id=row.id).count()
    services=FinancialService.query.filter_by(supplier_id=row.id).count()
    if used:
        row.active=False
        db.session.add(AuditEvent(user_id=session.get("user_id"),event_type="FIN_SUPPLIER_INACTIVATE",entity_type="financial_supplier",entity_id=str(row_id),detail=f"{name} · preservados {used} lançamento(s)"));db.session.commit()
        return jsonify({"ok":True,"inactivated":True,"linked_launches":used})
    FinancialService.query.filter_by(supplier_id=row.id).delete(synchronize_session=False)
    db.session.delete(row)
    db.session.add(AuditEvent(user_id=session.get("user_id"),event_type="FIN_SUPPLIER_DELETE",entity_type="financial_supplier",entity_id=str(row_id),detail=f"{name} · sem histórico financeiro"));db.session.commit()
    return jsonify({"ok":True,"deleted":True})

@app.delete("/api/financeiro/servicos/<int:row_id>")
@login_required
def financial_service_delete_api(row_id):
    if session.get("role") not in ("manager","atm_financial_admin"): return jsonify({"ok":False,"error":"Sem permissão para excluir serviço."}),403
    row=db.session.get(FinancialService,row_id)
    if not row: return jsonify({"ok":False,"error":"Serviço não encontrado."}),404
    used=FinancialMonthlyCost.query.filter_by(service_id=row.id).count()
    if used: return jsonify({"ok":False,"error":"Serviço possui lançamentos. Edite/exclua os lançamentos primeiro."}),409
    name=row.name;db.session.delete(row);db.session.add(AuditEvent(user_id=session.get("user_id"),event_type="FIN_SERVICE_DELETE",entity_type="financial_service",entity_id=str(row_id),detail=name));db.session.commit();return jsonify({"ok":True})

def _fin_service_for_text(supplier_id, service_text):
    text=(service_text or "").strip()
    if not text: return None
    row=FinancialService.query.filter(FinancialService.supplier_id==supplier_id,func.lower(FinancialService.name)==text.lower()).first()
    if not row:
        row=FinancialService(supplier_id=supplier_id,name=text,description=text,category="OUTROS",created_by=session.get("user_id"));db.session.add(row);db.session.flush()
    return row

def _fin_payload(row, users=None, sups=None, svcs=None):
    users={u.id:u.name for u in User.query.all()} if users is None else users; sups={x.id:x.name for x in FinancialSupplier.query.all()} if sups is None else sups; svcs={x.id:x.name for x in FinancialService.query.all()} if svcs is None else svcs
    try: alloc=json.loads(row.allocation_json or "{}")
    except: alloc={}
    return {"id":row.id,"competence":row.competence, "cost_center":getattr(row,"cost_center",None) or "SUPORTE_CAMPO","cost_center_id":getattr(row,"cost_center_id",None) or _fin_cost_center_id_for_key(getattr(row,"cost_center",None) or "SUPORTE_CAMPO"),"project":getattr(row,"project",None) or "","supplier_id":row.supplier_id,"supplier":sups.get(row.supplier_id,""),"service":getattr(row,"service_text",None) or svcs.get(row.service_id,""),"amount":round(float(row.amount or 0),2),"forecast_amount":None if getattr(row,"forecast_amount",None) is None else round(float(row.forecast_amount),2),"allocation":alloc,"invoice_number":getattr(row,"invoice_number",None) or "","notes":row.notes or "","updated_by":users.get(row.updated_by or row.created_by,""),"updated_at":row.updated_at.isoformat()+"Z"}

@app.route("/api/financeiro/lancamentos",methods=["GET","POST"])
@login_required
def financial_monthly_costs_api():
    if request.method=="GET":
        if session.get("role") not in ("manager","manager_field","atm_financial_admin"): return jsonify({"ok":False,"error":"Sem permissão."}),403
        comp=(request.args.get("competence") or "").strip();q=FinancialMonthlyCost.query
        if comp:q=q.filter_by(competence=comp)
        users={u.id:u.name for u in User.query.all()}; sups={x.id:x.name for x in FinancialSupplier.query.all()}; svcs={x.id:x.name for x in FinancialService.query.all()}
        return jsonify({"ok":True,"rows":[_fin_payload(x,users,sups,svcs) for x in q.order_by(FinancialMonthlyCost.competence.desc(),FinancialMonthlyCost.id.desc()).all()]})
    if not _financial_admin_allowed(): return jsonify({"ok":False,"error":"Sem permissão."}),403
    d=request.get_json(silent=True) or {}; comp=(d.get("competence") or "").strip(); sid=int(d.get("supplier_id") or 0); service_text=(d.get("service") or "").strip(); center=(d.get("cost_center") or "SUPORTE_CAMPO").strip().upper(); center_id=(d.get("cost_center_id") or _fin_cost_center_id_for_key(center)).strip().upper(); project=(d.get("project") or "").strip()
    try:
        amount=round(float(str(d.get("amount") or 0).replace(",",".")),2); forecast=d.get("forecast_amount"); forecast=None if forecast in (None,"") else round(float(str(forecast).replace(",",".")),2); alloc={str(k).upper():round(float(v or 0),2) for k,v in (d.get("allocation") or {}).items() if float(v or 0)>0}
    except: return jsonify({"ok":False,"error":"Valor/rateio inválido."}),400
    if len(comp)!=7 or comp[4]!="-" or amount<0: return jsonify({"ok":False,"error":"Competência e valor são obrigatórios."}),400
    supplier=db.session.get(FinancialSupplier,sid)
    if not supplier or not service_text: return jsonify({"ok":False,"error":"Fornecedor e serviço são obrigatórios."}),400
    total=round(sum(alloc.values()),2)
    if not alloc or abs(total-100)>0.001: return jsonify({"ok":False,"error":f"O rateio deve totalizar 100%. Atual: {total:.2f}%."}),400
    service=_fin_service_for_text(sid,service_text)
    row=FinancialMonthlyCost(competence=comp,supplier_id=sid,service_id=service.id,service_text=service_text,amount=amount,forecast_amount=forecast,cost_center=center,cost_center_id=center_id,project=project,allocation_json=json.dumps(alloc,ensure_ascii=False),invoice_number=(d.get("invoice_number") or "").strip(),notes=(d.get("notes") or "").strip(),created_by=session.get("user_id"),updated_by=session.get("user_id"));db.session.add(row);db.session.flush()
    db.session.add(AuditEvent(user_id=session.get("user_id"),event_type="FIN_MONTHLY_COST_CREATE",entity_type="financial_monthly_cost",entity_id=str(row.id),detail=f"{comp} · {center} · {supplier.name} · {service_text} · R$ {amount:.2f}"));db.session.commit();return jsonify({"ok":True,"id":row.id})

@app.route("/api/financeiro/lancamentos/<int:row_id>",methods=["PUT","DELETE"])
@login_required
def financial_monthly_cost_update_api(row_id):
    row=db.session.get(FinancialMonthlyCost,row_id)
    if not row: return jsonify({"ok":False,"error":"Lançamento não encontrado."}),404
    if request.method=="DELETE":
        if session.get("role")!="manager": return jsonify({"ok":False,"error":"Somente ADM pode excluir lançamento."}),403
        detail=f"{row.competence} · R$ {row.amount:.2f}";db.session.delete(row);db.session.add(AuditEvent(user_id=session.get("user_id"),event_type="FIN_MONTHLY_COST_DELETE",entity_type="financial_monthly_cost",entity_id=str(row_id),detail=detail));db.session.commit();return jsonify({"ok":True})
    if not _financial_admin_allowed(): return jsonify({"ok":False,"error":"Sem permissão."}),403
    d=request.get_json(silent=True) or {}; sid=int(d.get("supplier_id") or row.supplier_id); service_text=(d.get("service") or getattr(row,"service_text",None) or "").strip(); center=(d.get("cost_center") or getattr(row,"cost_center",None) or "SUPORTE_CAMPO").strip().upper(); center_id=(d.get("cost_center_id") or getattr(row,"cost_center_id",None) or _fin_cost_center_id_for_key(center)).strip().upper(); project=(d.get("project",getattr(row,"project",None)) or "").strip(); comp=(d.get("competence") or row.competence).strip()
    try:
        amount=round(float(str(d.get("amount",row.amount)).replace(",",".")),2); forecast=d.get("forecast_amount",getattr(row,"forecast_amount",None)); forecast=None if forecast in (None,"") else round(float(str(forecast).replace(",",".")),2); alloc={str(k).upper():round(float(v or 0),2) for k,v in (d.get("allocation") or json.loads(row.allocation_json or "{}")).items() if float(v or 0)>0}
    except: return jsonify({"ok":False,"error":"Valor/rateio inválido."}),400
    total=round(sum(alloc.values()),2)
    if abs(total-100)>0.001:return jsonify({"ok":False,"error":f"O rateio deve totalizar 100%. Atual: {total:.2f}%."}),400
    supplier=db.session.get(FinancialSupplier,sid); service=_fin_service_for_text(sid,service_text) if supplier else None
    if not supplier or not service:return jsonify({"ok":False,"error":"Fornecedor/serviço inválido."}),400
    before=f"{row.competence} · R$ {row.amount:.2f}"
    row.competence=comp;row.supplier_id=sid;row.service_id=service.id;row.service_text=service_text;row.amount=amount;row.forecast_amount=forecast;row.cost_center=center;row.cost_center_id=center_id;row.project=project;row.allocation_json=json.dumps(alloc,ensure_ascii=False);row.invoice_number=(d.get("invoice_number",getattr(row,"invoice_number",None)) or "").strip();row.notes=(d.get("notes",row.notes) or "").strip();row.updated_by=session.get("user_id");row.updated_at=datetime.utcnow()
    db.session.add(AuditEvent(user_id=session.get("user_id"),event_type="FIN_MONTHLY_COST_EDIT",entity_type="financial_monthly_cost",entity_id=str(row.id),detail=f"antes: {before} | depois: {comp} · {center} · R$ {amount:.2f}"));db.session.commit();return jsonify({"ok":True})


@app.route("/api/financeiro/lancamentos-lote", methods=["GET","POST"])
@login_required
def financial_monthly_batch_api():
    if not _financial_admin_allowed() or not _has_access("finance"):
        return jsonify({"ok":False,"error":"Sem permissão."}),403
    if request.method=="GET":
        comp=(request.args.get("competence") or "").strip(); center=(request.args.get("cost_center") or "SUPORTE_CAMPO").strip().upper()
        if len(comp)!=7 or comp[4]!="-": return jsonify({"ok":False,"error":"Informe uma competência válida."}),400
        center_aliases={center,center.replace("_"," ")}
        if center=="SUPORTE_CAMPO": center_aliases.update({"SUPORTE A CAMPO","SUPORTE_CAMPO","SUPORTE DE CAMPO"})
        q=(FinancialMonthlyCost.query.join(FinancialSupplier,FinancialMonthlyCost.supplier_id==FinancialSupplier.id).filter(FinancialMonthlyCost.competence==comp,func.upper(FinancialMonthlyCost.cost_center).in_({x.upper() for x in center_aliases})))
        rows=q.order_by(func.upper(FinancialSupplier.name),func.upper(FinancialMonthlyCost.service_text)).all(); source=comp
        template=False
        if not rows:
            prev=(db.session.query(FinancialMonthlyCost.competence).filter(FinancialMonthlyCost.competence<comp,func.upper(FinancialMonthlyCost.cost_center).in_({x.upper() for x in center_aliases})).order_by(FinancialMonthlyCost.competence.desc()).first())
            if prev:
                source=prev[0]; rows=(FinancialMonthlyCost.query.join(FinancialSupplier,FinancialMonthlyCost.supplier_id==FinancialSupplier.id).filter(FinancialMonthlyCost.competence==source,func.upper(FinancialMonthlyCost.cost_center).in_({x.upper() for x in center_aliases})).order_by(func.upper(FinancialSupplier.name),func.upper(FinancialMonthlyCost.service_text)).all()); template=True
        users={u.id:u.name for u in User.query.all()}; sups={x.id:x.name for x in FinancialSupplier.query.all()}; svcs={x.id:x.name for x in FinancialService.query.all()}
        payload=[]
        for x in rows:
            d=_fin_payload(x,users,sups,svcs); d["source_id"]=x.id; d["id"]=None if template else x.id; d["competence"]=comp
            if template: d["amount"]=0; d["invoice_number"]=""
            payload.append(d)
        return jsonify({"ok":True,"rows":payload,"template":template,"source_competence":source,"competence":comp})
    d=request.get_json(silent=True) or {}; comp=(d.get("competence") or "").strip(); center=(d.get("cost_center") or "SUPORTE_CAMPO").strip().upper(); items=d.get("rows") or []
    if len(comp)!=7 or comp[4]!="-": return jsonify({"ok":False,"error":"Competência inválida."}),400
    saved=deleted=0
    try:
        for item in items:
            if item.get("deleted") and item.get("id"):
                row=db.session.get(FinancialMonthlyCost,int(item["id"]));
                if row and row.competence==comp and row.cost_center==center:
                    db.session.delete(row); deleted+=1
                continue
            if item.get("deleted"): continue
            sid=int(item.get("supplier_id") or 0); service_text=(item.get("service") or "").strip(); supplier=db.session.get(FinancialSupplier,sid)
            if not supplier or not service_text: raise ValueError("Fornecedor e serviço são obrigatórios em todas as linhas.")
            amount=round(float(item.get("amount") or 0),2); forecast=item.get("forecast_amount"); forecast=None if forecast in (None,"") else round(float(forecast),2)
            alloc={str(k).upper():round(float(v or 0),2) for k,v in (item.get("allocation") or {}).items() if float(v or 0)>0}
            if not alloc or abs(sum(alloc.values())-100)>0.01: raise ValueError(f"Rateio de {supplier.name} deve totalizar 100%.")
            service=_fin_service_for_text(sid,service_text)
            row=db.session.get(FinancialMonthlyCost,int(item["id"])) if item.get("id") else None
            if row is None:
                row=FinancialMonthlyCost(competence=comp,supplier_id=sid,service_id=service.id,cost_center=center,cost_center_id=(item.get("cost_center_id") or _fin_cost_center_id_for_key(center)),created_by=session.get("user_id")); db.session.add(row)
            row.competence=comp; row.supplier_id=sid; row.service_id=service.id; row.service_text=service_text; row.amount=amount; row.forecast_amount=forecast; row.cost_center=center; row.cost_center_id=(item.get("cost_center_id") or _fin_cost_center_id_for_key(center)); row.project=(item.get("project") or "").strip(); row.invoice_number=(item.get("invoice_number") or "").strip(); row.allocation_json=json.dumps(alloc,ensure_ascii=False); row.notes=(item.get("notes") or "").strip(); row.updated_by=session.get("user_id"); row.updated_at=datetime.utcnow(); saved+=1
        db.session.add(AuditEvent(user_id=session.get("user_id"),event_type="FIN_MONTHLY_BATCH_SAVE",entity_type="financial_monthly_cost",entity_id=comp,detail=f"{center} · {saved} salvo(s) · {deleted} excluído(s)")); db.session.commit()
        return jsonify({"ok":True,"saved":saved,"deleted":deleted})
    except Exception as exc:
        db.session.rollback(); return jsonify({"ok":False,"error":str(exc)}),400


def _fin_terminal(value, point_name=""):
    raw=str(value or "").strip()
    if raw and raw.lower() not in ("none","nan") and re.search(r"\d",raw):
        raw=re.sub(r"\.0$","",raw)
        return re.sub(r"\s+","",raw)
    m=re.search(r"\bATM\s+(\d+)\b",str(point_name or ""),re.I)
    return m.group(1) if m else ""

def _fin_gtv(value):
    return re.sub(r"\D","",str(value or ""))

def _fin_parse_date(value):
    if isinstance(value,datetime): return value.date()
    if hasattr(value,"year") and hasattr(value,"month") and hasattr(value,"day"):
        try:return value
        except:return None
    textv=str(value or "").strip()
    for fmt in ("%d/%m/%Y","%Y-%m-%d","%d/%m/%y"):
        try:return datetime.strptime(textv,fmt).date()
        except:pass
    return None

def _fin_parse_time(value):
    if isinstance(value,datetime): return value.time().replace(microsecond=0)
    if hasattr(value,"hour") and hasattr(value,"minute"):
        try:return value.replace(microsecond=0)
        except:return value
    textv=str(value or "").strip()
    for fmt in ("%H:%M:%S","%H:%M"):
        try:return datetime.strptime(textv,fmt).time()
        except:pass
    return None

def _fin_dt(date_value,time_value):
    d=_fin_parse_date(date_value); t=_fin_parse_time(time_value)
    return datetime.combine(d,t) if d and t else (datetime.combine(d,datetime.min.time()) if d else None)

def _fin_hash(*parts):
    return hashlib.sha256("|".join(str(x or "").strip() for x in parts).encode("utf-8","ignore")).hexdigest()

def _fin_bulk_ignore(model, mappings, chunk=2000):
    if not mappings:return 0
    inserted=0
    dialect=db.engine.dialect.name
    for start in range(0,len(mappings),chunk):
        part=mappings[start:start+chunk]
        if dialect=="postgresql":
            from sqlalchemy.dialects.postgresql import insert as pg_insert
            stmt=pg_insert(model.__table__).values(part).on_conflict_do_nothing(index_elements=["source_hash"])
            res=db.session.execute(stmt); inserted += max(0,res.rowcount or 0)
        elif dialect=="sqlite":
            from sqlalchemy.dialects.sqlite import insert as sqlite_insert
            stmt=sqlite_insert(model.__table__).values(part).on_conflict_do_nothing(index_elements=["source_hash"])
            res=db.session.execute(stmt); inserted += max(0,res.rowcount or 0)
        else:
            hashes=[x["source_hash"] for x in part]
            existing={x[0] for x in db.session.query(model.source_hash).filter(model.source_hash.in_(hashes)).all()}
            fresh=[x for x in part if x["source_hash"] not in existing]
            if fresh: db.session.bulk_insert_mappings(model,fresh); inserted+=len(fresh)
    return inserted

def _fin_import_tbf_wb(wb, filename, user_id):
    result={"kind":"TBFORTE","collections":0,"processed_updated":0,"errors":0}
    transport=[ws for ws in wb.worksheets if ws.title.upper().startswith("TRANSPORTE")]
    processing=[ws for ws in wb.worksheets if ws.title.upper().startswith("PROCESSAMENTO")]
    rows=[]
    for ws in transport:
        headers={str(v or "").strip().upper():i for i,v in enumerate(next(ws.iter_rows(min_row=1,max_row=1,values_only=True)))}
        def gv(row,key):
            i=headers.get(key.upper()); return row[i] if i is not None and i<len(row) else None
        for row in ws.iter_rows(min_row=2,values_only=True):
            try:
                point=str(gv(row,"PONTO ATENDIMENTO") or "").strip(); terminal=_fin_terminal(gv(row,"TERMINAL"),point)
                d=_fin_parse_date(gv(row,"DATA")); start_at=_fin_dt(d,gv(row,"Hora inicial")); end_at=_fin_dt(d,gv(row,"Hora final"))
                if not terminal or not d or not end_at: continue
                gtv=_fin_gtv(gv(row,"GTV")); amount=float(gv(row,"Valor RECOLHIDO") or 0)
                sh=_fin_hash("COL",terminal,end_at.isoformat(),gtv,amount)
                rows.append({"terminal":terminal,"point_name":point,"collection_date":d,"start_at":start_at,"end_at":end_at,"collected_amount":amount,"gtv":gtv or None,"route":str(gv(row,"Rota") or "").strip() or None,"municipality":str(gv(row,"MUNICÍPIO") or "").strip() or None,"source_file":filename,"source_hash":sh,"imported_by":user_id,"imported_at":datetime.utcnow()})
            except Exception: result["errors"]+=1
    result["collections"]=_fin_bulk_ignore(FinancialCashCollection,rows); db.session.flush()
    # Atualiza valor declarado/apurado pela GTV. O processamento não precisa criar outra coleta.
    for ws in processing:
        headers={str(v or "").strip().upper():i for i,v in enumerate(next(ws.iter_rows(min_row=1,max_row=1,values_only=True)))}
        def gv(row,key):
            i=headers.get(key.upper()); return row[i] if i is not None and i<len(row) else None
        for row in ws.iter_rows(min_row=2,values_only=True):
            try:
                gtv=_fin_gtv(gv(row,"GTV")); point=str(gv(row,"PONTO ATENDIMENTO") or "").strip(); terminal=_fin_terminal(None,point)
                if not gtv: continue
                candidates=FinancialCashCollection.query.filter_by(gtv=gtv).all()
                if not candidates and terminal:
                    pd=_fin_parse_date(gv(row,"Data do Processamento")); candidates=FinancialCashCollection.query.filter_by(terminal=terminal,collection_date=pd).all() if pd else []
                if not candidates: continue
                declared=float(gv(row,"Valor Declarado (R$)") or 0); processed=float(gv(row,"Valor Apurado (R$)") or 0); pd=_fin_parse_date(gv(row,"Data do Processamento"))
                media=str(gv(row,"MOEDA OU CEDULA") or "").strip()
                try: qty_processed=int(float(gv(row,"Qtde Process") or 0))
                except: qty_processed=0
                try: processing_charge=float(gv(row,"Valor Total") or 0)
                except: processing_charge=0
                for c in candidates:
                    c.declared_amount=declared; c.processed_amount=processed; c.processed_at=datetime.combine(pd,datetime.min.time()) if pd else c.processed_at
                    c.processed_media_type=media or None
                    c.processed_note_count=qty_processed if "CEDULA" in normalize(media).upper() else None
                    c.processing_charge=processing_charge
                    result["processed_updated"]+=1
            except Exception: result["errors"]+=1
    return result

def _fin_import_transactions_wb(wb, filename, user_id, job_id=None):
    ws=wb[wb.sheetnames[0]]; header=next(ws.iter_rows(min_row=1,max_row=1,values_only=True)); headers={str(v or "").strip().upper():i for i,v in enumerate(header)}
    required={"ATM","CTM_DATETIME_TZ","CTD_VALUE"}
    if not required.issubset(headers): raise ValueError("Planilha de transações sem colunas ATM, ctm_datetime_tz e ctd_value.")
    out=[]; errors=0; total=0; inserted_total=0; started=time.monotonic(); total_rows=max(0,(ws.max_row or 1)-1)
    def gv(row,key):
        i=headers.get(key.upper()); return row[i] if i is not None and i<len(row) else None
    def publish(force=False):
        if not job_id: return
        elapsed=max(0.001,time.monotonic()-started); rate=total/elapsed
        pct=(total/max(total_rows,1))*88.0 if total_rows else 0
        eta=max(0,(total_rows-total)/rate) if total_rows and rate>0 else None
        _fin_job_update(job_id,stage="TRANSACOES",rows_total=total_rows,rows_processed=total,rows_inserted=inserted_total,rows_errors=errors,rows_per_second=round(rate,1),eta_seconds=None if eta is None else int(eta),heartbeat_at=datetime.utcnow().isoformat()+"Z",progress=min(94,5+int(pct)),message=f"Transações: {total:,} de {total_rows:,} linhas processadas".replace(",","."))
    for row in ws.iter_rows(min_row=2,values_only=True):
        total+=1
        try:
            terminal=_fin_terminal(gv(row,"ATM")); dt=gv(row,"ctm_datetime_tz")
            if not isinstance(dt,datetime):
                d=_fin_parse_date(gv(row,"data")); dt=_fin_dt(d,gv(row,"hora"))
            if not terminal or not dt:
                if total%2000==0: publish()
                continue
            status=str(gv(row,"ctm_status") or "").strip().upper(); value=float(gv(row,"ctd_value") or 0); cpm=str(gv(row,"cpm_id") or "").strip()
            sh=_fin_hash("TX",terminal,dt.isoformat(),status,value,cpm)
            out.append({"terminal":terminal,"transaction_at":dt,"status":status or None,"value":value,"cpm_id":cpm or None,"source_file":filename,"source_hash":sh,"imported_by":user_id,"imported_at":datetime.utcnow()})
            if len(out)>=2000:
                inserted_total += _fin_bulk_ignore(FinancialATMTransaction,out,chunk=1000); db.session.commit(); out=[]; publish()
        except Exception:
            errors+=1
            if total%2000==0: publish()
    if out:
        inserted_total += _fin_bulk_ignore(FinancialATMTransaction,out,chunk=1000); db.session.commit()
    publish(True)
    persisted=db.session.query(func.count(FinancialATMTransaction.id)).filter(FinancialATMTransaction.source_file==filename).scalar() or 0
    return {"kind":"TRANSACOES","rows_read":total,"transactions":int(persisted),"inserted":int(inserted_total),"duplicates":max(0,total-errors-inserted_total),"errors":errors}

@app.get("/financeiro/apuracao")
@login_required
def financial_cash_reconciliation_page():
    if not _financial_admin_allowed() or not _has_access("finance.apuracao"):
        return redirect(url_for("dashboard_landing"))
    return render_template("financial_cash_reconciliation.html",app_release=APP_RELEASE)

def _financial_import_worker(job_id, paths, filenames, user_id):
    with app.app_context():
        results=[]
        try:
            total_files=len(paths)
            _fin_job_update(job_id,status="PROCESSANDO",progress=2,message="Processamento iniciado",stage="PREPARANDO",current_file=1,total_files=total_files,rows_total=0,rows_processed=0,rows_inserted=0,rows_errors=0,rows_per_second=0,eta_seconds=None,heartbeat_at=datetime.utcnow().isoformat()+"Z")
            for idx,(path,filename) in enumerate(zip(paths,filenames),start=1):
                base_pct=int((idx-1)/max(total_files,1)*90)+5
                _fin_job_update(job_id,current_file=idx,current_filename=filename,progress=base_pct,stage="LENDO_ARQUIVO",message=f"Lendo {filename}",heartbeat_at=datetime.utcnow().isoformat()+"Z")
                wb=load_workbook(path,read_only=True,data_only=True)
                upper=[x.upper() for x in wb.sheetnames]
                if any(x.startswith("TRANSPORTE") for x in upper):
                    result=_fin_import_tbf_wb(wb,filename,user_id)
                else:
                    result=_fin_import_transactions_wb(wb,filename,user_id,job_id=job_id)
                wb.close(); db.session.commit(); results.append(result)
                _fin_job_update(job_id,results=results,progress=min(95,int(idx/max(total_files,1)*90)+5),message=f"{filename} concluído")
            db.session.add(AuditEvent(user_id=user_id,event_type="FIN_APURACAO_IMPORT",entity_type="financial_cash_reconciliation",entity_id=str(len(paths)),detail=json.dumps(results,ensure_ascii=False)[:4000])); db.session.commit()
            tx_count=db.session.query(func.count(FinancialATMTransaction.id)).scalar() or 0
            col_count=db.session.query(func.count(FinancialCashCollection.id)).scalar() or 0
            _fin_job_update(job_id,status="CONCLUIDO",progress=100,message="Importação concluída",results=results,transactions=int(tx_count),collections=int(col_count),finished_at=datetime.utcnow().isoformat()+"Z")
        except Exception as exc:
            db.session.rollback()
            app.logger.exception("Falha importação financeira background")
            _fin_job_update(job_id,status="FALHOU",progress=0,message="Importação interrompida",error=str(exc),results=results,finished_at=datetime.utcnow().isoformat()+"Z")
        finally:
            for path in paths:
                try: Path(path).unlink(missing_ok=True)
                except Exception: pass

@app.post("/api/financeiro/apuracao/importar")
@login_required
def financial_cash_reconciliation_import():
    if not _financial_admin_allowed() or not _has_access("finance.apuracao"):
        return jsonify({"ok":False,"error":"Sem permissão."}),403
    uploaded=request.files.getlist("files") or ([request.files.get("file")] if request.files.get("file") else [])
    uploaded=[f for f in uploaded if f and f.filename]
    if not uploaded:return jsonify({"ok":False,"error":"Selecione uma ou mais planilhas Excel."}),400
    job_id=uuid.uuid4().hex
    paths=[]; names=[]
    try:
        for i,f in enumerate(uploaded,1):
            name=secure_filename(f.filename) or f"arquivo_{i}.xlsx"
            path=FIN_IMPORT_DIR / f"{job_id}_{i}_{name}"
            f.save(path); paths.append(str(path)); names.append(name)
        with FIN_IMPORT_LOCK:
            FIN_IMPORT_JOBS[job_id]={"job_id":job_id,"user_id":session.get("user_id"),"status":"NA_FILA","progress":1,"stage":"UPLOAD_CONCLUIDO","message":"Arquivos recebidos. Preparando importação...","files":names,"current_file":0,"total_files":len(names),"results":[],"rows_total":0,"rows_processed":0,"rows_inserted":0,"rows_errors":0,"rows_per_second":0,"eta_seconds":None,"heartbeat_at":datetime.utcnow().isoformat()+"Z","created_at":datetime.utcnow().isoformat()+"Z","updated_at":datetime.utcnow().isoformat()+"Z"}
        threading.Thread(target=_financial_import_worker,args=(job_id,paths,names,session.get("user_id")),daemon=True,name=f"fin-import-{job_id[:8]}").start()
        return jsonify({"ok":True,"background":True,"job_id":job_id,"message":"Importação iniciada em segundo plano. Você pode sair desta página."}),202
    except Exception as exc:
        for path in paths:
            try: Path(path).unlink(missing_ok=True)
            except Exception: pass
        return jsonify({"ok":False,"error":str(exc)}),400

@app.get("/api/financeiro/apuracao/importar/<job_id>/status")
@login_required
def financial_cash_reconciliation_import_status(job_id):
    if not _financial_admin_allowed() or not _has_access("finance.apuracao"):
        return jsonify({"ok":False,"error":"Sem permissão."}),403
    job=_fin_job_snapshot(job_id)
    if not job:return jsonify({"ok":False,"error":"Importação não encontrada ou servidor reiniciado."}),404
    return jsonify({"ok":True,**job})

@app.get("/api/financeiro/apuracao/importar/active")
@login_required
def financial_cash_reconciliation_import_active():
    if not _financial_admin_allowed() or not _has_access("finance.apuracao"):
        return jsonify({"ok":False,"error":"Sem permissão."}),403
    uid=session.get("user_id")
    with FIN_IMPORT_LOCK:
        jobs=[dict(v) for v in FIN_IMPORT_JOBS.values() if v.get("user_id")==uid and v.get("status") in ("NA_FILA","PROCESSANDO")]
    jobs.sort(key=lambda x:x.get("created_at") or "",reverse=True)
    return jsonify({"ok":True,"job":jobs[0] if jobs else None})

@app.get("/api/financeiro/apuracao/terminais")
@login_required
def financial_cash_reconciliation_terminals():
    if not _financial_admin_allowed() or not _has_access("finance.apuracao"): return jsonify({"ok":False,"error":"Sem permissão."}),403
    now=time.time(); cached=_FIN_TERMINALS_CACHE.get("payload")
    if cached is not None and now-float(_FIN_TERMINALS_CACHE.get("at") or 0) < _FIN_TERMINALS_CACHE_TTL:
        resp=jsonify(cached); resp.headers["X-Autopass-Cache"]="HIT"; return resp
    # V58: terminais e metadados via JOIN/subquery; evita IN enorme e objetos ORM desnecessários.
    term_sub=db.session.query(FinancialCashCollection.terminal.label("terminal")).distinct().subquery()
    terms=[x[0] for x in db.session.query(term_sub.c.terminal).order_by(term_sub.c.terminal).all() if x and x[0]]
    asset_rows=(db.session.query(BaseAsset.terminal_number,func.max(BaseAsset.locality),func.max(BaseAsset.company),func.max(BaseAsset.line))
        .join(term_sub,BaseAsset.terminal_number==term_sub.c.terminal).group_by(BaseAsset.terminal_number).all())
    locmap={_fin_terminal(t):{"locality":loc or "","company":comp or "","line":line or ""} for t,loc,comp,line in asset_rows if _fin_terminal(t)}
    # V60 REV2 PERFORMANCE: evita COUNT(*) em >500 mil transações.
    # O importador é append-only; MAX(id) fornece o volume operacional sem varredura integral.
    tx_count=int(db.session.query(func.max(FinancialATMTransaction.id)).scalar() or 0)
    latest_at=db.session.query(func.max(FinancialATMTransaction.transaction_at)).scalar()
    latest_import_at=db.session.query(func.max(FinancialATMTransaction.imported_at)).scalar()
    latest_tx=None; latest_import=None
    if latest_at:
        row=db.session.query(FinancialATMTransaction.terminal,FinancialATMTransaction.source_file).filter(FinancialATMTransaction.transaction_at==latest_at).first()
        latest_tx={"terminal":row[0] if row else "","at":latest_at.isoformat(),"source_file":row[1] if row else ""}
    if latest_import_at:
        src=db.session.query(FinancialATMTransaction.source_file).filter(FinancialATMTransaction.imported_at==latest_import_at).first()
        latest_import={"at":latest_import_at.isoformat(),"source_file":src[0] if src else ""}
    payload={"ok":True,"terminals":[{"terminal":t,**locmap.get(t,{"locality":"","company":"","line":""})} for t in terms],
        "collections":int(db.session.query(func.count(FinancialCashCollection.id)).scalar() or 0),"transactions":int(tx_count or 0),
        "latest_transaction":latest_tx,"latest_import":latest_import,"integrity":{"stored_transactions":int(tx_count or 0),"exact_duplicates_in_table":0,"unique_key":"terminal + data/hora + status + valor + cpm_id","note":"Duplicidades exatas são ignoradas no importador pela chave única source_hash."}}
    _FIN_TERMINALS_CACHE["payload"]=payload; _FIN_TERMINALS_CACHE["at"]=time.time()
    resp=jsonify(payload); resp.headers["X-Autopass-Cache"]="MISS"; return resp


@app.get("/api/financeiro/apuracao/coletas")
@login_required
def financial_cash_reconciliation_collections():
    if not _financial_admin_allowed() or not _has_access("finance.apuracao"): return jsonify({"ok":False,"error":"Sem permissão."}),403
    terminal=_fin_terminal(request.args.get("terminal"));
    if not terminal:return jsonify({"ok":True,"rows":[]})
    rows=FinancialCashCollection.query.filter_by(terminal=terminal).order_by(FinancialCashCollection.end_at).all()
    return jsonify({"ok":True,"rows":[{"id":x.id,"terminal":x.terminal,"date":x.end_at.isoformat(),"date_label":x.end_at.strftime("%d/%m/%y %H:%M"),"recollected_amount":round(float(x.collected_amount or 0),2),"collected_amount":round(float(x.collected_amount or 0),2),"declared_amount":None if x.declared_amount is None else round(float(x.declared_amount),2),"processed_amount":None if x.processed_amount is None else round(float(x.processed_amount),2),"processed_note_count":getattr(x,"processed_note_count",None),"processed_media_type":getattr(x,"processed_media_type",None) or "","point_name":x.point_name or "","gtv":x.gtv or ""} for x in rows]})

@app.get("/api/financeiro/apuracao/calcular")
@login_required
def financial_cash_reconciliation_calculate():
    if not _financial_admin_allowed() or not _has_access("finance.apuracao"): return jsonify({"ok":False,"error":"Sem permissão."}),403
    terminal=_fin_terminal(request.args.get("terminal")); a=db.session.get(FinancialCashCollection,int(request.args.get("initial_id") or 0)); b=db.session.get(FinancialCashCollection,int(request.args.get("final_id") or 0))
    if not terminal or not a or not b or a.terminal!=terminal or b.terminal!=terminal:return jsonify({"ok":False,"error":"Selecione terminal, coleta inicial e coleta final válidos."}),400
    if b.end_at<=a.end_at:return jsonify({"ok":False,"error":"A coleta final deve ser posterior à coleta inicial."}),400
    base=FinancialATMTransaction.query.filter(FinancialATMTransaction.terminal==terminal,FinancialATMTransaction.transaction_at>a.end_at,FinancialATMTransaction.transaction_at<=b.end_at,FinancialATMTransaction.status.in_(["V","A"]))
    tx_count=base.count()
    tx_sum=float(db.session.query(func.coalesce(func.sum(FinancialATMTransaction.value),0)).filter(FinancialATMTransaction.terminal==terminal,FinancialATMTransaction.transaction_at>a.end_at,FinancialATMTransaction.transaction_at<=b.end_at,FinancialATMTransaction.status.in_(["V","A"])).scalar() or 0)
    status_rows=db.session.query(FinancialATMTransaction.status,func.count(FinancialATMTransaction.id),func.coalesce(func.sum(FinancialATMTransaction.value),0)).filter(FinancialATMTransaction.terminal==terminal,FinancialATMTransaction.transaction_at>a.end_at,FinancialATMTransaction.transaction_at<=b.end_at).group_by(FinancialATMTransaction.status).all()
    recollected=float(b.collected_amount or 0); declared=None if b.declared_amount is None else float(b.declared_amount); processed=None if b.processed_amount is None else float(b.processed_amount)
    diff_tx_ap=round(tx_sum-processed,2) if processed is not None else None; pct_tx_ap=round((diff_tx_ap/processed*100),4) if diff_tx_ap is not None and processed else None
    diff_tx_dec=round(tx_sum-declared,2) if declared is not None else None; diff_ap_dec=round(processed-declared,2) if processed is not None and declared is not None else None
    duration_hours=max(0,(b.end_at-a.end_at).total_seconds()/3600); duration_days=duration_hours/24 if duration_hours else 0
    avg_ticket=round(tx_sum/tx_count,2) if tx_count else 0; avg_day=round(tx_sum/duration_days,2) if duration_days else 0
    note_count=int(getattr(b,"processed_note_count",0) or 0); avg_note=round(processed/note_count,2) if processed is not None and note_count else None
    abs_pct=abs(pct_tx_ap) if pct_tx_ap is not None else None
    if tx_count==0: diagnosis="SEM_TRANSACOES"; diagnosis_text="Não há transações no intervalo selecionado."
    elif diff_tx_ap is None: diagnosis="SEM_APURACAO"; diagnosis_text="O período possui transações, mas não há Valor Apurado vinculado à coleta final."
    elif abs(diff_tx_ap)<0.01: diagnosis="CONCILIADO"; diagnosis_text="Valor das transações coincide com o Valor Apurado do período."
    elif abs_pct is not None and abs_pct<=0.5: diagnosis="ATENCAO"; diagnosis_text="Há pequena diferença entre Valor das Transações e Valor Apurado."
    else: diagnosis="DIVERGENCIA_RELEVANTE"; diagnosis_text="Há diferença relevante entre Valor das Transações e Valor Apurado."
    observations=[f"Intervalo entre coletas: {duration_days:.1f} dia(s) ({duration_hours:.1f} h)."]
    if tx_count: observations.append(f"{tx_count:,} transação(ões), considerando somente status V + A; ticket médio R$ {avg_ticket:,.2f}.".replace(",","X").replace(".",",").replace("X","."))
    if note_count: observations.append(f"TBForte processou {note_count:,} cédula(s) vinculadas à coleta final.".replace(",","."))
    details=base.order_by(FinancialATMTransaction.transaction_at).limit(500).all()
    return jsonify({"ok":True,"terminal":terminal,"initial":{"id":a.id,"at":a.end_at.isoformat(),"label":a.end_at.strftime("%d/%m/%Y %H:%M"),"amount":None if a.processed_amount is None else round(float(a.processed_amount),2),"processed_amount":None if a.processed_amount is None else round(float(a.processed_amount),2),"declared_amount":None if a.declared_amount is None else round(float(a.declared_amount),2),"recollected_amount":round(float(a.collected_amount or 0),2)},"final":{"id":b.id,"at":b.end_at.isoformat(),"label":b.end_at.strftime("%d/%m/%Y %H:%M"),"amount":None if processed is None else round(processed,2),"processed_amount":None if processed is None else round(processed,2),"declared_amount":None if declared is None else round(declared,2),"recollected_amount":round(recollected,2),"processed_note_count":note_count,"processed_media_type":getattr(b,"processed_media_type",None) or "","processing_charge":round(float(getattr(b,"processing_charge",0) or 0),2)},"transaction_count":tx_count,"transaction_sum":round(tx_sum,2),"all_status_count":tx_count,"recollected_amount":round(recollected,2),"collected_amount":round(recollected,2),"declared_amount":None if declared is None else round(declared,2),"processed_amount":None if processed is None else round(processed,2),"processed_note_count":note_count,"difference":diff_tx_ap,"difference_pct":pct_tx_ap,"difference_tx_apurado":diff_tx_ap,"difference_tx_declarado":diff_tx_dec,"difference_apurado_declarado":diff_ap_dec,"duration_hours":round(duration_hours,2),"duration_days":round(duration_days,2),"average_ticket":avg_ticket,"average_per_day":avg_day,"average_value_per_processed_note":avg_note,"status":"CONCILIADO" if diff_tx_ap is not None and abs(diff_tx_ap)<0.01 else "DIVERGENCIA","diagnosis":diagnosis,"diagnosis_text":diagnosis_text,"observations":observations,"status_breakdown":[{"status":st or "—","count":int(n),"amount":round(float(v or 0),2)} for st,n,v in status_rows],"transactions":[{"at":x.transaction_at.isoformat(),"label":x.transaction_at.strftime("%d/%m/%Y %H:%M:%S"),"status":x.status or "","value":round(float(x.value or 0),2)} for x in details],"transactions_truncated":tx_count>500,"rule":"Valor das Transações considera somente status V + A; após o horário final da coleta inicial e até o horário final da coleta final"})

@app.get("/api/financeiro/apuracao/calcular-multiplos")
@login_required
def financial_cash_reconciliation_multi():
    if not _financial_admin_allowed() or not _has_access("finance.apuracao"): return jsonify({"ok":False,"error":"Sem permissão."}),403
    terminals=[_fin_terminal(x) for x in request.args.getlist("terminal") if _fin_terminal(x)]
    primary=_fin_terminal(request.args.get("primary")); a=db.session.get(FinancialCashCollection,int(request.args.get("initial_id") or 0)); b=db.session.get(FinancialCashCollection,int(request.args.get("final_id") or 0))
    if not terminals or not primary or not a or not b or a.terminal!=primary or b.terminal!=primary or b.end_at<=a.end_at: return jsonify({"ok":False,"error":"Selecione terminais e o intervalo de coletas do terminal de referência."}),400
    out=[]
    for terminal in terminals[:500]:
        ia=FinancialCashCollection.query.filter(FinancialCashCollection.terminal==terminal,FinancialCashCollection.end_at<=a.end_at).order_by(FinancialCashCollection.end_at.desc()).first()
        fb=FinancialCashCollection.query.filter(FinancialCashCollection.terminal==terminal,FinancialCashCollection.end_at>=b.end_at).order_by(FinancialCashCollection.end_at.asc()).first()
        if not ia or not fb or fb.end_at<=ia.end_at:
            out.append({"terminal":terminal,"status":"SEM JANELA COMPATÍVEL"}); continue
        q=db.session.query(func.count(FinancialATMTransaction.id),func.coalesce(func.sum(FinancialATMTransaction.value),0)).filter(FinancialATMTransaction.terminal==terminal,FinancialATMTransaction.transaction_at>ia.end_at,FinancialATMTransaction.transaction_at<=fb.end_at,FinancialATMTransaction.status.in_(["V","A"])).first()
        count=int(q[0] or 0); total=float(q[1] or 0); recollected=float(fb.collected_amount or 0); declared=None if fb.declared_amount is None else float(fb.declared_amount); processed=None if fb.processed_amount is None else float(fb.processed_amount)
        diff=round(total-processed,2) if processed is not None else None; pct=round(diff/processed*100,4) if diff is not None and processed else None
        hours=(fb.end_at-ia.end_at).total_seconds()/3600; days=hours/24; avg_day=round(total/days,2) if days else 0; avg_ticket=round(total/count,2) if count else 0
        asset=BaseAsset.query.filter(BaseAsset.terminal_number==terminal).first(); note_count=int(getattr(fb,"processed_note_count",0) or 0)
        diag="SEM APURAÇÃO" if diff is None else ("IGUAL" if abs(diff)<.01 else ("MAIOR" if diff>0 else "MENOR"))
        out.append({"terminal":terminal,"locality":(asset.locality if asset else (fb.point_name or "")),"initial":ia.end_at.strftime("%d/%m/%Y %H:%M"),"final":fb.end_at.strftime("%d/%m/%Y %H:%M"),"duration_days":round(days,2),"transaction_count":count,"transaction_sum":round(total,2),"recollected_amount":round(recollected,2),"collected_amount":round(recollected,2),"declared_amount":None if declared is None else round(declared,2),"processed_amount":None if processed is None else round(processed,2),"processed_note_count":note_count,"average_ticket":avg_ticket,"average_per_day":avg_day,"difference":diff,"difference_pct":pct,"difference_tx_declarado":None if declared is None else round(total-declared,2),"difference_apurado_declarado":None if processed is None or declared is None else round(processed-declared,2),"status":diag})
    out.sort(key=lambda x: abs(float(x.get("difference") or 0)),reverse=True)
    summary={"total":len(out),"equal":0,"greater":0,"less":0,"no_window":0,"sum_greater":0.0,"sum_less":0.0,"sum_net":0.0}
    for x in out:
        st=x.get("status"); dv=x.get("difference")
        if st=="IGUAL": summary["equal"]+=1
        elif st=="MAIOR": summary["greater"]+=1; summary["sum_greater"]+=float(dv or 0)
        elif st=="MENOR": summary["less"]+=1; summary["sum_less"]+=float(dv or 0)
        elif st=="SEM JANELA COMPATÍVEL": summary["no_window"]+=1
        if dv is not None: summary["sum_net"]+=float(dv)
    summary={k:(round(v,2) if isinstance(v,float) else v) for k,v in summary.items()}
    return jsonify({"ok":True,"rows":out,"summary":summary,"rule":"Valor das Transações considera somente status V + A. Para cada ATM, compara a janela equivalente do terminal de referência."})

@app.get("/api/financeiro/apuracao/calcular-periodo")
@login_required
def financial_cash_reconciliation_period():
    if not _financial_admin_allowed() or not _has_access("finance.apuracao"):
        return jsonify({"ok":False,"error":"Sem permissão."}),403
    terminals=list(dict.fromkeys([_fin_terminal(x) for x in request.args.getlist("terminal") if _fin_terminal(x)]))[:500]
    try:
        start=datetime.fromisoformat((request.args.get("start") or "").strip())
        end=datetime.fromisoformat((request.args.get("end") or "").strip())
    except Exception:
        return jsonify({"ok":False,"error":"Informe data e hora inicial e final válidas."}),400
    if not terminals:return jsonify({"ok":False,"error":"Selecione uma ou mais ATMs."}),400
    if end<=start:return jsonify({"ok":False,"error":"O período final deve ser posterior ao inicial."}),400
    out=[]
    for terminal in terminals:
        ia=FinancialCashCollection.query.filter(FinancialCashCollection.terminal==terminal,FinancialCashCollection.end_at<=start).order_by(FinancialCashCollection.end_at.desc()).first()
        fb=FinancialCashCollection.query.filter(FinancialCashCollection.terminal==terminal,FinancialCashCollection.end_at<=end,FinancialCashCollection.end_at>start).order_by(FinancialCashCollection.end_at.desc()).first()
        q=db.session.query(func.count(FinancialATMTransaction.id),func.coalesce(func.sum(FinancialATMTransaction.value),0)).filter(FinancialATMTransaction.terminal==terminal,FinancialATMTransaction.transaction_at>=start,FinancialATMTransaction.transaction_at<=end,FinancialATMTransaction.status.in_(["V","A"])).first()
        count=int(q[0] or 0); total=float(q[1] or 0); processed=None if not fb or fb.processed_amount is None else float(fb.processed_amount); recollected=None if not fb else float(fb.collected_amount or 0); declared=None if not fb or fb.declared_amount is None else float(fb.declared_amount)
        diff=round(total-processed,2) if processed is not None else None; pct=round(diff/processed*100,4) if diff is not None and processed else None
        asset=BaseAsset.query.filter(BaseAsset.terminal_number==terminal).first(); locality=asset.locality if asset else (fb.point_name if fb else "")
        status="SEM APURAÇÃO" if diff is None else ("IGUAL" if abs(diff)<.01 else ("MAIOR" if diff>0 else "MENOR"))
        out.append({"terminal":terminal,"locality":locality or "","initial":start.strftime("%d/%m/%Y %H:%M"),"final":end.strftime("%d/%m/%Y %H:%M"),"duration_days":round((end-start).total_seconds()/86400,2),"transaction_count":count,"transaction_sum":round(total,2),"recollected_amount":None if recollected is None else round(recollected,2),"declared_amount":None if declared is None else round(declared,2),"processed_amount":None if processed is None else round(processed,2),"processed_note_count":int(getattr(fb,"processed_note_count",0) or 0) if fb else 0,"difference":diff,"difference_pct":pct,"difference_tx_declarado":None if declared is None else round(total-declared,2),"difference_apurado_declarado":None if processed is None or declared is None else round(processed-declared,2),"status":status})
    summary={"total":len(out),"equal":0,"greater":0,"less":0,"no_window":0,"sum_greater":0.0,"sum_less":0.0,"sum_net":0.0}
    for x in out:
        st=x["status"]; dv=x["difference"]
        if st=="IGUAL":summary["equal"]+=1
        elif st=="MAIOR":summary["greater"]+=1;summary["sum_greater"]+=float(dv or 0)
        elif st=="MENOR":summary["less"]+=1;summary["sum_less"]+=float(dv or 0)
        else:summary["no_window"]+=1
        if dv is not None:summary["sum_net"]+=float(dv)
    for k in ("sum_greater","sum_less","sum_net"):summary[k]=round(summary[k],2)
    return jsonify({"ok":True,"rows":out,"summary":summary,"rule":"Período independente · 1, várias ou todas as ATMs · Valor das Transações somente status V + A."})

@app.get("/api/financeiro/apuracao/exportar.xlsx")
@login_required
def financial_cash_reconciliation_export():
    if not _financial_admin_allowed() or not _has_access("finance.apuracao"): abort(403)
    terminal=_fin_terminal(request.args.get("terminal")); a=db.session.get(FinancialCashCollection,int(request.args.get("initial_id") or 0)); b=db.session.get(FinancialCashCollection,int(request.args.get("final_id") or 0))
    if not terminal or not a or not b or a.terminal!=terminal or b.terminal!=terminal or b.end_at<=a.end_at: return "Filtros inválidos",400
    txs=FinancialATMTransaction.query.filter(FinancialATMTransaction.terminal==terminal,FinancialATMTransaction.transaction_at>a.end_at,FinancialATMTransaction.transaction_at<=b.end_at,FinancialATMTransaction.status.in_(["V","A"])).order_by(FinancialATMTransaction.transaction_at).all()
    total=sum(float(x.value or 0) for x in txs); recollected=float(b.collected_amount or 0); declared=None if b.declared_amount is None else float(b.declared_amount); processed=None if b.processed_amount is None else float(b.processed_amount); diff=(total-processed) if processed is not None else None; pct=(diff/processed) if diff is not None and processed else None
    hours=(b.end_at-a.end_at).total_seconds()/3600; days=hours/24; asset=BaseAsset.query.filter(BaseAsset.terminal_number==terminal).first(); locality=asset.locality if asset else (b.point_name or ""); note_count=int(getattr(b,"processed_note_count",0) or 0)
    wb=Workbook(); ws=wb.active; ws.title="Apuração"; ws.append(["Terminal","Localidade","Coleta inicial","Valor apurado inicial","Coleta final","Valor apurado final","Dias","Qtd transações - status V+A","Valor transações","Valor recolhido","Valor declarado","Valor apurado","Qtde cédulas processadas","Dif. transações x apurado","Dif. transações x declarado","Dif. apurado x declarado","% transações x apurado","Status"])
    status="SEM APURAÇÃO" if diff is None else ("IGUAL" if abs(diff)<.01 else ("MAIOR" if diff>0 else "MENOR"))
    ws.append([terminal,locality,a.end_at,a.processed_amount,b.end_at,processed,days,len(txs),total,recollected,declared,processed,note_count,diff,(total-declared if declared is not None else None),(processed-declared if processed is not None and declared is not None else None),pct,status])
    wd=wb.create_sheet("Transações"); wd.append(["Terminal","Localidade","Data/hora","Status","Valor"]); [wd.append([terminal,locality,x.transaction_at,x.status,x.value]) for x in txs]
    for sh in wb.worksheets:
        sh.freeze_panes="A2"; sh.auto_filter.ref=sh.dimensions; [setattr(c,'font',Font(bold=True)) for c in sh[1]]
    bio=io.BytesIO(); wb.save(bio); bio.seek(0); return send_file(bio,as_attachment=True,download_name=f"apuracao_{terminal}_{datetime.now().strftime('%Y%m%d_%H%M')}.xlsx",mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

@app.get("/api/financeiro/apuracao/exportar-multiplos.xlsx")
@login_required
def financial_cash_reconciliation_export_multi():
    if not _financial_admin_allowed() or not _has_access("finance.apuracao"): abort(403)
    terminals=[_fin_terminal(x) for x in request.args.getlist("terminal") if _fin_terminal(x)]; primary=_fin_terminal(request.args.get("primary")); a=db.session.get(FinancialCashCollection,int(request.args.get("initial_id") or 0)); b=db.session.get(FinancialCashCollection,int(request.args.get("final_id") or 0))
    if not terminals or not primary or not a or not b or a.terminal!=primary or b.terminal!=primary or b.end_at<=a.end_at:return "Filtros inválidos",400
    wb=Workbook(); ws=wb.active; ws.title="Comparativo"; ws.append(["Terminal","Localidade","Coleta inicial","Coleta final","Dias","Qtd transações - status V+A","Valor transações","Valor recolhido","Valor declarado","Valor apurado","Qtde cédulas","Dif. transações x apurado","Dif. transações x declarado","Dif. apurado x declarado","% transações x apurado","Status"])
    for terminal in terminals[:500]:
        ia=FinancialCashCollection.query.filter(FinancialCashCollection.terminal==terminal,FinancialCashCollection.end_at<=a.end_at).order_by(FinancialCashCollection.end_at.desc()).first(); fb=FinancialCashCollection.query.filter(FinancialCashCollection.terminal==terminal,FinancialCashCollection.end_at>=b.end_at).order_by(FinancialCashCollection.end_at.asc()).first()
        if not ia or not fb or fb.end_at<=ia.end_at: continue
        cnt,total=db.session.query(func.count(FinancialATMTransaction.id),func.coalesce(func.sum(FinancialATMTransaction.value),0)).filter(FinancialATMTransaction.terminal==terminal,FinancialATMTransaction.transaction_at>ia.end_at,FinancialATMTransaction.transaction_at<=fb.end_at,FinancialATMTransaction.status.in_(["V","A"])).first(); cnt=int(cnt or 0); total=float(total or 0)
        recollected=float(fb.collected_amount or 0); declared=None if fb.declared_amount is None else float(fb.declared_amount); processed=None if fb.processed_amount is None else float(fb.processed_amount); diff=(total-processed) if processed is not None else None; days=(fb.end_at-ia.end_at).total_seconds()/86400; asset=BaseAsset.query.filter(BaseAsset.terminal_number==terminal).first(); locality=asset.locality if asset else (fb.point_name or ""); pct=(diff/processed) if diff is not None and processed else None; status="SEM APURAÇÃO" if diff is None else ("IGUAL" if abs(diff)<.01 else ("MAIOR" if diff>0 else "MENOR"))
        ws.append([terminal,locality,ia.end_at,fb.end_at,days,cnt,total,recollected,declared,processed,int(getattr(fb,"processed_note_count",0) or 0),diff,(total-declared if declared is not None else None),(processed-declared if processed is not None and declared is not None else None),pct,status])
    ws.freeze_panes="A2"; ws.auto_filter.ref=ws.dimensions; [setattr(c,'font',Font(bold=True)) for c in ws[1]]
    bio=io.BytesIO(); wb.save(bio); bio.seek(0); return send_file(bio,as_attachment=True,download_name=f"apuracao_multiplos_{datetime.now().strftime('%Y%m%d_%H%M')}.xlsx",mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

@app.post("/api/financeiro/importar.xlsx")
@login_required
def financial_import_xlsx():
    if not _financial_admin_allowed(): return jsonify({"ok":False,"error":"Sem permissão."}),403
    f=request.files.get("file")
    if not f or not f.filename: return jsonify({"ok":False,"error":"Selecione o arquivo Excel."}),400
    try: wb=load_workbook(f.stream,data_only=True)
    except Exception as exc: return jsonify({"ok":False,"error":f"Excel inválido: {exc}"}),400
    if "Empresas" not in wb.sheetnames or "Lançamentos" not in wb.sheetnames:
        return jsonify({"ok":False,"error":"O arquivo deve conter as abas Empresas e Lançamentos."}),400
    def headers(ws): return {str(c.value or '').strip():i+1 for i,c in enumerate(ws[4])}
    def val(ws,row,h,name):
        i=h.get(name); return ws.cell(row,i).value if i else None
    created_sup=updated_sup=pending_sup=created_cost=duplicates=errors=0
    ws=wb["Empresas"]; h=headers(ws)
    for r in range(5,ws.max_row+1):
        name=str(val(ws,r,h,"Razão Social / Nome") or '').strip()
        if not name: continue
        sup=FinancialSupplier.query.filter(func.lower(FinancialSupplier.name)==name.lower()).first()
        if not sup: sup=FinancialSupplier(name=name,created_by=session.get('user_id'));db.session.add(sup);created_sup+=1
        else: updated_sup+=1
        sup.trade_name=str(val(ws,r,h,"Nome Fantasia") or '').strip() or None; sup.cnpj=str(val(ws,r,h,"CNPJ") or '').strip() or None
        sup.primary_cost_center=str(val(ws,r,h,"Centro de Custo Principal") or '').strip().upper() or None; sup.cost_center_id=str(val(ws,r,h,"ID Centro de Custo") or _fin_cost_center_id_for_key(sup.primary_cost_center)).strip().upper() or None; sup.contact_name=str(val(ws,r,h,"Contato") or '').strip() or None
        sup.phone=str(val(ws,r,h,"Telefone") or '').strip() or None; sup.email=str(val(ws,r,h,"E-mail") or '').strip() or None
        sup.pending_profile=_fin_supplier_pending(sup); pending_sup+=1 if sup.pending_profile else 0
    db.session.flush()
    ws=wb["Lançamentos"]; h=headers(ws)
    prodcols=[("Produto 1","Rateio 1 %"),("Produto 2","Rateio 2 %"),("Produto 3","Rateio 3 %"),("Produto 4","Rateio 4 %"),("Produto 5","Rateio 5 %"),("Produto 6","Rateio 6 %")]
    for r in range(5,ws.max_row+1):
        try:
            supplier_name=str(val(ws,r,h,"Fornecedor / Empresa") or '').strip(); service=str(val(ws,r,h,"Serviço / Atividade") or '').strip()
            if not supplier_name and not service: continue
            if not supplier_name or not service: errors+=1; continue
            cv=val(ws,r,h,"Competência")
            if isinstance(cv,datetime): comp=cv.strftime('%Y-%m')
            else:
                raw=str(cv or '').strip(); m=re.search(r'(20\d{2})[-/](\d{1,2})',raw); comp=f"{m.group(1)}-{int(m.group(2)):02d}" if m else raw
            amount=float(val(ws,r,h,"Valor Realizado (R$)") or 0); forecast=val(ws,r,h,"Valor Forecast (R$)"); forecast=None if forecast in (None,'') else float(forecast)
            center=str(val(ws,r,h,"Centro de Custo") or 'SUPORTE_CAMPO').strip().upper().replace(' ','_').replace('Ê','E').replace('É','E').replace('Ç','C').replace('Ã','A')
            center={'SUPORTE_A_CAMPO':'SUPORTE_CAMPO','ASSISTENCIA_TECNICA':'ASSISTENCIA_TECNICA','IMPLANTACAO_DE_HARDWARE':'IMPLANTACAO_HARDWARE'}.get(center,center)
            sup=FinancialSupplier.query.filter(func.lower(FinancialSupplier.name)==supplier_name.lower()).first()
            if not sup:
                sup=FinancialSupplier(name=supplier_name,pending_profile=True,created_by=session.get('user_id'));db.session.add(sup);db.session.flush();created_sup+=1;pending_sup+=1
            alloc={}
            for pc,rc in prodcols:
                pr=str(val(ws,r,h,pc) or '').strip().upper(); rv=val(ws,r,h,rc)
                if pr and rv not in (None,''): alloc[pr]=round(float(rv),2)
            if alloc and abs(sum(alloc.values())-100)>0.01: errors+=1; continue
            if not alloc: alloc={'OUTROS':100.0}
            exists=FinancialMonthlyCost.query.filter_by(competence=comp,supplier_id=sup.id).filter(func.abs(FinancialMonthlyCost.amount-amount)<0.01,func.lower(FinancialMonthlyCost.service_text)==service.lower()).first()
            if exists: duplicates+=1; continue
            svc=_fin_service_for_text(sup.id,service)
            row=FinancialMonthlyCost(competence=comp,supplier_id=sup.id,service_id=svc.id,service_text=service,amount=round(amount,2),forecast_amount=None if forecast is None else round(forecast,2),cost_center=center,cost_center_id=str(val(ws,r,h,"ID Centro de Custo") or _fin_cost_center_id_for_key(center)).strip().upper(),project=str(val(ws,r,h,"Projeto") or '').strip(),invoice_number=str(val(ws,r,h,"NF / Documento") or '').strip(),allocation_json=json.dumps(alloc,ensure_ascii=False),notes=str(val(ws,r,h,"Observação / Justificativa") or '').strip(),created_by=session.get('user_id'),updated_by=session.get('user_id'))
            db.session.add(row); created_cost+=1
        except Exception: errors+=1
    db.session.commit(); _td_cache_clear()
    return jsonify({"ok":True,"companies_created":created_sup,"companies_updated":updated_sup,"companies_pending":pending_sup,"launches_created":created_cost,"duplicates":duplicates,"errors":errors})

@app.get("/api/financeiro/export.xlsx")
@login_required
def financial_export_xlsx():
    if session.get("role") not in ("manager","manager_field","atm_financial_admin"): abort(403)
    comps=set(x for x in (request.args.get('competences') or '').split(',') if x); center=(request.args.get('center') or 'ALL'); product=(request.args.get('product') or 'ALL').upper()
    q=FinancialMonthlyCost.query
    if comps:q=q.filter(FinancialMonthlyCost.competence.in_(comps))
    if center!='ALL':q=q.filter_by(cost_center=center)
    rows=q.order_by(FinancialMonthlyCost.competence,FinancialMonthlyCost.id).all(); sups={x.id:x.name for x in FinancialSupplier.query.all()}
    wb=Workbook(); ws=wb.active;ws.title='Lançamentos';ws.append(['Competência','Centro de Custo','ID Centro de Custo','Fornecedor','Serviço','Projeto','Realizado','Forecast','Produto','Rateio %','Valor Rateado'])
    for x in rows:
        alloc=json.loads(x.allocation_json or '{}')
        pairs=alloc.items() if product=='ALL' else [(product,alloc.get(product,0))]
        for pr,pct in pairs:
            if not pct: continue
            ws.append([x.competence,x.cost_center,getattr(x,'cost_center_id',None) or _fin_cost_center_id_for_key(x.cost_center),sups.get(x.supplier_id,''),x.service_text or '',x.project or '',x.amount,x.forecast_amount,pr,pct,round(x.amount*float(pct)/100,2)])
    bio=io.BytesIO();wb.save(bio);bio.seek(0);return send_file(bio,as_attachment=True,download_name='dashboard_financeiro_v55_2.xlsx',mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')

@app.get("/api/panoramas/export.xlsx")
@login_required
def panorama_export():
    wb=Workbook();ws=wb.active;ws.title="Visões Panorâmicas";ws.append(["Empresa","Linha","Localidade","Status","Pontos","Fotos"])
    for x in _panorama_payload(): ws.append([x.get("company"),x.get("line"),x.get("location"),x.get("status"),x.get("point_count",0),x.get("photo_count",0)])
    bio=io.BytesIO();wb.save(bio);bio.seek(0);return send_file(bio,as_attachment=True,download_name="visoes_panoramicas.xlsx",mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

def _panorama_media_bytes(stored_name):
    if not stored_name:
        return None
    try:
        if stored_name.startswith("r2__"):
            return _r2_get_bytes(stored_name[4:])
        path=UPLOAD_DIR/stored_name
        return path.read_bytes() if path.exists() else None
    except Exception:
        return None


def _pptx_compact_image(raw, max_width=1280, max_height=900, quality=68):
    from PIL import Image, ImageOps
    """Reduz memória/tamanho do PPTX antes de inserir a imagem."""
    if not raw:
        return None
    try:
        source = io.BytesIO(raw)
        with Image.open(source) as im:
            im = ImageOps.exif_transpose(im)
            im.thumbnail((max_width, max_height), Image.Resampling.LANCZOS)
            if im.mode not in ("RGB", "L"):
                bg = Image.new("RGB", im.size, "white")
                if "A" in im.getbands():
                    bg.paste(im, mask=im.getchannel("A"))
                else:
                    bg.paste(im)
                im = bg
            elif im.mode == "L":
                im = im.convert("RGB")
            out = io.BytesIO()
            im.save(out, format="JPEG", quality=quality, optimize=True, progressive=True)
            out.seek(0)
            return out
    except Exception:
        return None

def _panorama_export_job_update(job_id, **changes):
    with PANORAMA_EXPORT_LOCK:
        job=PANORAMA_EXPORT_JOBS.get(job_id)
        if job:
            job.update(changes); job["updated_at"]=datetime.utcnow().isoformat()+"Z"


def _panorama_export_cleanup():
    cutoff=time.time()-PANORAMA_EXPORT_MAX_AGE_SECONDS
    with PANORAMA_EXPORT_LOCK:
        old=[]
        for jid,job in list(PANORAMA_EXPORT_JOBS.items()):
            ts=float(job.get("created_ts") or time.time())
            if ts<cutoff and job.get("status") not in ("PROCESSANDO","FILA"):
                old.append((jid,job.get("path")))
        for jid,path in old:
            PANORAMA_EXPORT_JOBS.pop(jid,None)
            try:
                if path: Path(path).unlink(missing_ok=True)
            except Exception: pass


def _generate_panorama_pptx(company,line,status,search,output_path,progress_cb=None):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    from PIL import Image
    rows=[x for x in _panorama_payload() if (not company or x.get("company")==company) and (not line or x.get("line")==line) and (not status or x.get("status")==status) and (not search or search in (x.get("location") or "").lower())]
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    navy=(18,52,93); teal=(51,190,190); dark=(25,38,58); light=(241,246,250)
    def textbox(slide,x,y,w,h,text,size=18,bold=False,rgb=dark,align=PP_ALIGN.LEFT):
        box=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=box.text_frame; tf.clear(); p=tf.paragraphs[0]; p.text=str(text); p.alignment=align; run=p.runs[0]; run.font.size=Pt(size); run.font.bold=bold; run.font.name="Arial"; run.font.color.rgb=__import__('pptx').dml.color.RGBColor(*rgb); return box
    def brand(slide):
        logo=STATIC_DIR/'autopass-logo.png'
        if logo.exists(): slide.shapes.add_picture(str(logo),Inches(.45),Inches(.28),width=Inches(1.65))
        textbox(slide,10.7,.34,2.1,.3,"VISÃO PANORÂMICA",10,True,navy,PP_ALIGN.RIGHT)
    sl=prs.slides.add_slide(prs.slide_layouts[6]); brand(sl); textbox(sl,.7,2.15,11.9,.7,"Visão Panorâmica das Estações",30,True,navy,PP_ALIGN.CENTER); textbox(sl,.8,3.0,11.7,.45,"Apresentação executiva do acervo fotográfico",17,False,dark,PP_ALIGN.CENTER)
    recorte=" · ".join([x for x in [company or "Todas as empresas",line or "Todas as linhas",status or "Todos os status"] if x]); textbox(sl,.8,4.0,11.7,.35,recorte,12,False,dark,PP_ALIGN.CENTER); textbox(sl,.8,6.65,11.7,.3,datetime.now().strftime("Gerado em %d/%m/%Y %H:%M"),9,False,dark,PP_ALIGN.CENTER)
    sl=prs.slides.add_slide(prs.slide_layouts[6]); brand(sl); textbox(sl,.55,1.0,12,.45,"Resumo executivo",24,True,navy); total=len(rows); done=sum(x['status']=='CONCLUÍDA' for x in rows); prog=sum(x['status']=='EM ANDAMENTO' for x in rows); pend=sum(x['status']=='PENDENTE' for x in rows); photos=sum(x.get('photo_count',0) for x in rows)
    vals=[("Localidades",total),("Concluídas",done),("Em andamento",prog),("Pendentes",pend),("Fotos",photos)]
    for i,(lab,val) in enumerate(vals):
        x=.55+i*2.5; sh=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.0),Inches(2.2),Inches(1.35)); sh.fill.solid(); sh.fill.fore_color.rgb=__import__('pptx').dml.color.RGBColor(*light); sh.line.color.rgb=__import__('pptx').dml.color.RGBColor(205,218,230); textbox(sl,x+.12,2.18,1.95,.28,lab,10,True,dark); textbox(sl,x+.12,2.55,1.95,.48,val,24,True,navy)
    pct=round(done/total*100) if total else 0; textbox(sl,.65,4.25,12,.45,f"Progresso geral: {pct}%",20,True,teal)
    processed=0; image_errors=0
    for loc in rows:
        photo_items=[]
        for pt in loc.get('points',[]):
            for ph in pt.get('photos',[]): photo_items.append((pt,ph))
        if not photo_items: photo_items=[(None,None)]
        for start in range(0,len(photo_items),4):
            batch=photo_items[start:start+4]; sl=prs.slides.add_slide(prs.slide_layouts[6]); brand(sl); textbox(sl,.55,.88,12,.42,loc.get('location') or 'Localidade',23,True,navy); meta=f"{loc.get('company','—')} · {loc.get('line','—')} · {loc.get('status','—')}"; textbox(sl,.55,1.34,12,.3,meta,11,False,dark); tech=', '.join(loc.get('technicians') or []) or '—'; textbox(sl,.55,1.68,12,.25,f"Técnico(s): {tech}",10,False,dark)
            positions=[(.55,2.15,6.0,2.05),(6.78,2.15,6.0,2.05),(.55,4.65,6.0,2.05),(6.78,4.65,6.0,2.05)]
            for j,(pt,ph) in enumerate(batch):
                x,y,w,h=positions[j]
                if ph:
                    raw=_panorama_media_bytes(ph.get('stored_name'))
                    if raw:
                        try:
                            # REV6: compactação mais agressiva para proteger RAM/CPU no plano de 1 worker.
                            bio=_pptx_compact_image(raw,max_width=960,max_height=720,quality=55); raw=None
                            if bio is None: raise ValueError("Imagem inválida")
                            with Image.open(bio) as im: iw,ih=im.size
                            bio.seek(0); frame_w=w; frame_h=h-0.32; ratio=min(frame_w/max(iw,1),frame_h/max(ih,1)); pic_w=max(0.1,iw*ratio); pic_h=max(0.1,ih*ratio); px=x+(frame_w-pic_w)/2; py=y+(frame_h-pic_h)/2
                            sl.shapes.add_picture(bio,Inches(px),Inches(py),width=Inches(pic_w),height=Inches(pic_h)); bio.close()
                        except Exception:
                            image_errors+=1; textbox(sl,x,y,w,h-0.3,"Imagem indisponível",12,False,dark,PP_ALIGN.CENTER)
                    else:
                        image_errors+=1; textbox(sl,x,y,w,h-0.3,"Imagem indisponível",12,False,dark,PP_ALIGN.CENTER)
                    textbox(sl,x,y+h-.28,w,.25,(pt.get('name') if pt else 'Foto')+" · "+(ph.get('uploaded_by') or '—'),9,False,dark); processed+=1
                    if progress_cb and photos: progress_cb(processed,photos)
                else: textbox(sl,x,y,w,h,"Nenhuma foto anexada nesta localidade.",13,False,dark,PP_ALIGN.CENTER)
    prs.save(str(output_path))
    return {"locations":total,"photos":photos,"processed":processed,"image_errors":image_errors,"slides":len(prs.slides),"bytes":Path(output_path).stat().st_size if Path(output_path).exists() else 0}


class _PanoramaExportCancelled(Exception):
    pass

def _panorama_job_cancelled(job_id):
    with PANORAMA_EXPORT_LOCK:
        job=PANORAMA_EXPORT_JOBS.get(job_id) or {}
        return bool(job.get("cancel_requested") or job.get("status")=="CANCELADO")

def _panorama_export_worker(job_id,company,line,status,search):
    started=time.time(); path=PANORAMA_EXPORT_DIR/f"{job_id}.pptx"
    _panorama_export_job_update(job_id,status="PROCESSANDO",started_at=datetime.utcnow().isoformat()+"Z",message="Preparando dados...",progress=1)
    try:
        with app.app_context():
            def progress(done,total):
                if _panorama_job_cancelled(job_id):
                    raise _PanoramaExportCancelled()
                pct=max(2,min(95,round(done*95/max(1,total))))
                _panorama_export_job_update(job_id,progress=pct,message=f"Processando fotos: {done}/{total}")
            if _panorama_job_cancelled(job_id): raise _PanoramaExportCancelled()
            stats=_generate_panorama_pptx(company,line,status,search,path,progress)
        if _panorama_job_cancelled(job_id): raise _PanoramaExportCancelled()
        _panorama_export_job_update(job_id,status="PRONTO",progress=100,message="PowerPoint pronto para baixar.",path=str(path),filename=f"visao_panoramica_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx",stats=stats,duration_s=round(time.time()-started,1),finished_at=datetime.utcnow().isoformat()+"Z")
    except _PanoramaExportCancelled:
        try:
            if path.exists(): path.unlink()
        except Exception: pass
        _panorama_export_job_update(job_id,status="CANCELADO",progress=0,message="Geração cancelada pelo usuário.",finished_at=datetime.utcnow().isoformat()+"Z")
    except Exception as exc:
        app.logger.exception("Falha na exportação assíncrona do PowerPoint panorâmico")
        _panorama_export_job_update(job_id,status="ERRO",progress=0,message=str(exc)[:300],duration_s=round(time.time()-started,1),finished_at=datetime.utcnow().isoformat()+"Z")


@app.post("/api/panoramas/export.pptx/jobs")
@login_required
def panorama_export_pptx_job_create():
    if session.get("role") not in ("manager","manager_field","technician","consultation"):
        return jsonify({"ok":False,"error":"Sem permissão."}),403
    _panorama_export_cleanup()
    company=(request.form.get("company") or request.args.get("company") or "").strip(); line=(request.form.get("line") or request.args.get("line") or "").strip(); status=(request.form.get("status") or request.args.get("status") or "").strip(); search=(request.form.get("search") or request.args.get("search") or "").strip().lower()
    filter_key=json.dumps({"company":company,"line":line,"status":status,"search":search},sort_keys=True,ensure_ascii=False,separators=(",",":"))
    with PANORAMA_EXPORT_LOCK:
        active=next(((jid,j) for jid,j in PANORAMA_EXPORT_JOBS.items() if j.get("status") in ("FILA","PROCESSANDO")),None)
        if active:
            jid,j=active
            if j.get("user_id") not in (None,session.get("user_id")) and session.get("role") not in ("manager","manager_field"):
                return jsonify({"ok":False,"error":"Há outra exportação em processamento. Tente novamente em alguns instantes."}),409
            if j.get("filter_key")==filter_key:
                return jsonify({"ok":True,"job_id":jid,"status":j.get("status"),"reused":True,"filter_key":filter_key,"message":"Já existe um PowerPoint em processamento para estes filtros."}),202
            return jsonify({"ok":False,"error":"Há um PowerPoint em geração com outros filtros. Cancele a geração atual antes de iniciar outra."}),409
        job_id=uuid.uuid4().hex
        PANORAMA_EXPORT_JOBS[job_id]={"id":job_id,"status":"FILA","progress":0,"message":"Exportação na fila.","created_ts":time.time(),"created_at":datetime.utcnow().isoformat()+"Z","user_id":session.get("user_id"),"filter_key":filter_key,"filters":{"company":company,"line":line,"status":status,"search":search}}
    threading.Thread(target=_panorama_export_worker,args=(job_id,company,line,status,search),daemon=True,name=f"panorama-pptx-{job_id[:8]}").start()
    return jsonify({"ok":True,"job_id":job_id,"status":"FILA","filter_key":filter_key,"message":"PowerPoint sendo preparado em segundo plano."}),202


@app.post("/api/panoramas/export.pptx/jobs/<job_id>/cancel")
@login_required
def panorama_export_pptx_job_cancel(job_id):
    with PANORAMA_EXPORT_LOCK:
        job=PANORAMA_EXPORT_JOBS.get(job_id)
        if not job: return jsonify({"ok":False,"error":"Exportação não encontrada ou expirada."}),404
        if session.get("role") not in ("manager","manager_field") and job.get("user_id") not in (None,session.get("user_id")):
            return jsonify({"ok":False,"error":"Sem permissão para este processamento."}),403
        if job.get("status") in ("FILA","PROCESSANDO"):
            job["cancel_requested"]=True
            job["message"]="Cancelamento solicitado..."
            return jsonify({"ok":True,"status":"CANCELANDO","job_id":job_id})
        return jsonify({"ok":True,"status":job.get("status"),"job_id":job_id})

@app.get("/api/panoramas/export.pptx/jobs/<job_id>")
@login_required
def panorama_export_pptx_job_status(job_id):
    _panorama_export_cleanup()
    with PANORAMA_EXPORT_LOCK:
        job=dict(PANORAMA_EXPORT_JOBS.get(job_id) or {})
    if not job: return jsonify({"ok":False,"error":"Exportação não encontrada ou expirada."}),404
    if session.get("role") not in ("manager","manager_field") and job.get("user_id") not in (None,session.get("user_id")):
        return jsonify({"ok":False,"error":"Sem permissão para este processamento."}),403
    # Não expõe caminho físico interno.
    job.pop("path",None); job.pop("user_id",None); job.pop("created_ts",None)
    job["ok"]=True
    if job.get("status")=="PRONTO": job["download_url"]=f"/api/panoramas/export.pptx/jobs/{job_id}/download"
    return jsonify(job)


@app.get("/api/processamentos")
@login_required
def v63_processamentos_api():
    _panorama_export_cleanup()
    with PANORAMA_EXPORT_LOCK:
        jobs=[]; uid=session.get("user_id"); is_admin=session.get("role") in ("manager","manager_field")
        for jid,j in PANORAMA_EXPORT_JOBS.items():
            if not is_admin and j.get("user_id") not in (None,uid): continue
            item={k:v for k,v in dict(j).items() if k not in ("path","created_ts","user_id")}; item["id"]=jid; item["type"]="POWERPOINT_PANORAMA"
            if item.get("status")=="PRONTO": item["download_url"]=f"/api/panoramas/export.pptx/jobs/{jid}/download"
            jobs.append(item)
    jobs.sort(key=lambda x:x.get("created_at") or "",reverse=True)
    return jsonify({"ok":True,"release":APP_RELEASE,"jobs":jobs[:20],"active":sum(1 for j in jobs if j.get("status") in ("FILA","PROCESSANDO"))})

@app.get("/api/panoramas/export.pptx/jobs/<job_id>/download")
@login_required
def panorama_export_pptx_job_download(job_id):
    with PANORAMA_EXPORT_LOCK:
        job=dict(PANORAMA_EXPORT_JOBS.get(job_id) or {})
    if not job or job.get("status")!="PRONTO": return jsonify({"ok":False,"error":"Arquivo ainda não está pronto."}),409
    if session.get("role") not in ("manager","manager_field") and job.get("user_id") not in (None,session.get("user_id")):
        return jsonify({"ok":False,"error":"Sem permissão para este processamento."}),403
    path=Path(job.get("path") or "")
    if not path.exists(): return jsonify({"ok":False,"error":"Arquivo temporário expirou."}),404
    return send_file(path,as_attachment=True,download_name=job.get("filename") or "visao_panoramica.pptx",mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")


@app.get("/api/panoramas/export.pptx")
@login_required
def panorama_export_pptx_legacy():
    # REV6: evita bloquear o único worker por vários minutos. O frontend usa o job assíncrono.
    return jsonify({"ok":False,"error":"Exportação direta desativada por segurança. Use a geração em segundo plano."}),409


@app.get("/visao-panoramica")
@login_required
def panorama_page():
    if not _has_access("field.panorama"): abort(403)
    if session.get("role") not in ("manager", "manager_field", "technician", "consultation"):
        return redirect(url_for("activities_page" if session.get("role") == "technician_implantation" else "teams_page"))
    return render_template("panorama.html")


def _panorama_payload():
    # V42.2.3.1 PERFORMANCE — elimina consultas N+1.
    locations = Location.query.order_by(Location.company, Location.line, Location.location).all()
    if not locations: return []
    loc_ids=[x.id for x in locations]
    points=(PanoramaPoint.query.filter(PanoramaPoint.location_id.in_(loc_ids)).order_by(PanoramaPoint.location_id,PanoramaPoint.point_name).all())
    point_ids=[x.id for x in points]
    photos=((PanoramaPhoto.query.filter(PanoramaPhoto.point_id.in_(point_ids)).order_by(PanoramaPhoto.point_id,PanoramaPhoto.created_at).all()) if point_ids else [])
    user_ids={x.created_by for x in points if x.created_by}; user_ids.update(x.uploaded_by for x in photos if x.uploaded_by)
    users=({u.id:u.name for u in User.query.filter(User.id.in_(user_ids)).all()} if user_ids else {})
    photos_by_point={}
    for ph in photos: photos_by_point.setdefault(ph.point_id,[]).append(ph)
    points_by_location={}
    for pt in points: points_by_location.setdefault(pt.location_id,[]).append(pt)
    rows=[]
    for loc in locations:
        p_out=[]; total=0
        for pt in points_by_location.get(loc.id,[]):
            pp=photos_by_point.get(pt.id,[]); total+=len(pp); creator=users.get(pt.created_by,"—")
            p_out.append({"id":pt.id,"name":pt.point_name,"notes":pt.notes or "","technician":creator,"status":"CONCLUÍDA" if pp else "EM ANDAMENTO","photos":[{"id":ph.id,"url":"/uploads/"+ph.stored_name,"thumb_url":"/uploads/"+ph.stored_name+"?thumb=1","name":ph.original_name,"stored_name":ph.stored_name,"uploaded_by":users.get(ph.uploaded_by,"—"),"created_at":ph.created_at.isoformat()+"Z" if ph.created_at else None,"latitude":ph.latitude,"longitude":ph.longitude} for ph in pp]})
        auto_status="PENDENTE" if not p_out else ("CONCLUÍDA" if all(x["photos"] for x in p_out) else "EM ANDAMENTO")
        override=(loc.panorama_status_override or "").strip().upper().replace("CONCLUIDA","CONCLUÍDA")
        status=override if override in ("PENDENTE","EM ANDAMENTO","CONCLUÍDA") else auto_status
        techs=sorted({x["technician"] for x in p_out if x.get("technician") and x["technician"]!="—"})
        rows.append({"id":loc.id,"company":loc.company,"line":loc.line,"location":loc.location,"reference_latitude":loc.reference_latitude,"reference_longitude":loc.reference_longitude,"status":status,"auto_status":auto_status,"status_override":bool(override),"photo_count":total,"technicians":techs,"points":p_out})
    return rows

@app.get("/api/panoramas")
@login_required
def panorama_list_api():
    return jsonify({"ok":True,"locations":_panorama_payload()})

@app.post("/api/panoramas/<int:location_id>/points")
@field_required
def panorama_upload_api(location_id):
    if _activity_request_too_large(): return jsonify({"ok":False,"error":f"Envio excede {_ACTIVITY_REQUEST_MAX_MB} MB. Envie menos fotos por vez."}),413
    # V40.1: upload panorâmico com resposta JSON garantida para o mobile.
    # Antes, qualquer falha de storage/R2 escapava como HTML 500 e o JavaScript
    # ficava sem feedback ao tocar em "Salvar fotos".
    try:
        loc=db.session.get(Location,location_id)
        if not loc:
            return jsonify({"ok":False,"error":"Localidade não encontrada."}),404

        point_name=(request.form.get("point_name") or "Visão geral").strip() or "Visão geral"
        files=[f for f in request.files.getlist("photos") if f and f.filename]
        if not files:
            return jsonify({"ok":False,"error":"Selecione ou tire pelo menos uma foto antes de salvar."}),400

        # Limite individual evita estouro de memória em aparelhos móveis e no servidor.
        max_file_bytes=25*1024*1024
        for f in files:
            try:
                pos=f.stream.tell()
                f.stream.seek(0,2); size=f.stream.tell(); f.stream.seek(pos)
            except Exception:
                size=0
            if size and size>max_file_bytes:
                return jsonify({"ok":False,"error":f"A foto {f.filename} excede 25 MB. Reduza a resolução e tente novamente."}),413

        pt=PanoramaPoint.query.filter(
            func.lower(PanoramaPoint.point_name)==point_name.lower(),
            PanoramaPoint.location_id==location_id
        ).first()
        if not pt:
            pt=PanoramaPoint(
                location_id=location_id,point_name=point_name,
                notes=(request.form.get("notes") or "").strip(),
                created_by=session["user_id"]
            )
            db.session.add(pt); db.session.flush()
        elif request.form.get("notes") is not None:
            pt.notes=(request.form.get("notes") or "").strip()

        lat=_optional_float(request.form.get("latitude")); lon=_optional_float(request.form.get("longitude"))
        added=0
        for f in files:
            safe=secure_filename(f.filename) or f"panorama_{secrets.token_hex(4)}.jpg"
            stored=f"pan_{pt.id}_{secrets.token_hex(6)}_{safe}"
            mime=f.mimetype or "application/octet-stream"
            stored=_store_uploaded_file(f,"panorama",stored,mime)
            db.session.add(PanoramaPhoto(
                point_id=pt.id,original_name=f.filename,stored_name=stored,mime_type=mime,
                uploaded_by=session["user_id"],latitude=lat,longitude=lon
            ))
            added+=1

        if not added:
            db.session.rollback()
            return jsonify({"ok":False,"error":"Nenhuma foto válida foi recebida pelo servidor."}),400

        pt.updated_at=datetime.utcnow()
        db.session.add(AuditEvent(
            user_id=session.get("user_id"),event_type="PANORAMA_UPLOAD",
            entity_type="location",entity_id=str(location_id),
            detail=f"{point_name}: {added} foto(s)"
        ))
        db.session.commit()
        return jsonify({"ok":True,"point_id":pt.id,"photos_added":added,"message":f"{added} foto(s) salva(s) com sucesso."})
    except Exception as exc:
        db.session.rollback()
        app.logger.exception("Falha no upload panorâmico da localidade %s", location_id)
        return jsonify({"ok":False,"error":"Não foi possível salvar as fotos. Tente novamente. Se persistir, informe o horário do erro.","detail":str(exc)[:180]}),500

@app.post("/api/panoramas/import-whatsapp")
@manager_required
def panorama_import_whatsapp_api():
    zf=request.files.get("zip")
    if not zf or not zf.filename.lower().endswith(".zip"):
        return jsonify({"ok":False,"error":"Selecione um arquivo ZIP exportado do WhatsApp."}),400
    raw=zf.read()
    if len(raw)>150*1024*1024: return jsonify({"ok":False,"error":"ZIP acima do limite de 150 MB."}),413
    locs=Location.query.all(); aliases=[]
    def nrm(v): return re.sub(r"[^A-Z0-9]+"," ",unicodedata.normalize("NFD",str(v or "")).encode("ascii","ignore").decode().upper()).strip()
    for loc in locs: aliases.append((loc,nrm(loc.location)))
    imported=[]; unresolved=[]; duplicates=0
    try:
        with zipfile.ZipFile(io.BytesIO(raw)) as z:
            names=z.namelist(); text=""
            for n in names:
                if n.lower().endswith(".txt"):
                    try: text += "\n"+z.read(n).decode("utf-8",errors="ignore")
                    except Exception: pass
            for name in names:
                ext=Path(name).suffix.lower()
                if ext not in (".jpg",".jpeg",".png",".webp",".heic"): continue
                base=nrm(Path(name).stem); context=nrm(text[max(0,text.upper().find(Path(name).name.upper())-350):text.upper().find(Path(name).name.upper())+350]) if Path(name).name.upper() in text.upper() else ""
                hay=base+" "+context; matches=[(len(a),loc) for loc,a in aliases if a and a in hay]
                if not matches: unresolved.append(name); continue
                loc=max(matches,key=lambda x:x[0])[1]
                data=z.read(name); digest=hashlib.sha256(data).hexdigest()[:20]
                original=Path(name).name
                if PanoramaPhoto.query.filter_by(original_name=original).first(): duplicates+=1; continue
                pt=PanoramaPoint.query.filter_by(location_id=loc.id,point_name="Importado do WhatsApp").first()
                if not pt:
                    pt=PanoramaPoint(location_id=loc.id,point_name="Importado do WhatsApp",notes="Importação ZIP WhatsApp",created_by=session["user_id"]);db.session.add(pt);db.session.flush()
                stored=f"pan_{pt.id}_{digest}_{secure_filename(original)}"
                if _r2_available():
                    key=f"panorama/{datetime.utcnow().strftime('%Y/%m')}/{stored}";_r2_put_bytes(key,data,mimetypes.guess_type(original)[0] or "application/octet-stream");stored="r2__"+key
                else: (UPLOAD_DIR/stored).write_bytes(data)
                db.session.add(PanoramaPhoto(point_id=pt.id,original_name=original,stored_name=stored,mime_type=mimetypes.guess_type(original)[0],uploaded_by=session["user_id"]))
                imported.append({"file":original,"location":loc.location})
        db.session.add(AuditEvent(user_id=session.get("user_id"),event_type="PANORAMA_WHATSAPP_IMPORT",entity_type="panorama",entity_id="zip",detail=f"{len(imported)} importadas; {len(unresolved)} não identificadas; {duplicates} duplicadas"));db.session.commit()
    except zipfile.BadZipFile: return jsonify({"ok":False,"error":"Arquivo ZIP inválido."}),400
    return jsonify({"ok":True,"imported":imported,"unresolved":unresolved,"duplicates":duplicates})


@app.post("/api/panoramas/<int:location_id>/status")
@manager_required
def panorama_status_override_api(location_id):
    loc=db.session.get(Location,location_id)
    if not loc:
        return jsonify({"ok":False,"error":"Localidade não encontrada."}),404
    requested=(request.get_json(silent=True) or {}).get("status")
    requested=(requested or "").strip().upper().replace("CONCLUIDA","CONCLUÍDA")
    if requested in ("AUTOMATICO","AUTOMÁTICO",""):
        new_value=None
    elif requested in ("PENDENTE","EM ANDAMENTO","CONCLUÍDA"):
        new_value=requested
    else:
        return jsonify({"ok":False,"error":"Status inválido."}),400
    old=loc.panorama_status_override
    loc.panorama_status_override=new_value
    db.session.add(AuditEvent(user_id=session.get("user_id"),event_type="PANORAMA_STATUS_OVERRIDE",entity_type="location",entity_id=str(location_id),detail=f"Visão Panorâmica: status manual {old or 'AUTOMÁTICO'} → {new_value or 'AUTOMÁTICO'}."))
    db.session.commit()
    row=next((x for x in _panorama_payload() if x["id"]==location_id),None)
    return jsonify({"ok":True,"status":row["status"] if row else (new_value or "PENDENTE"),"auto_status":row.get("auto_status") if row else None,"status_override":bool(new_value)})


@app.delete("/api/panoramas/photos/<int:photo_id>")
@manager_required
def panorama_delete_photo_api(photo_id):
    ph=db.session.get(PanoramaPhoto,photo_id)
    if not ph:
        return jsonify({"ok":False,"error":"Foto não encontrada."}),404
    pt=db.session.get(PanoramaPoint,ph.point_id)
    location_id=pt.location_id if pt else None
    original=ph.original_name
    _delete_stored_media(ph.stored_name)
    db.session.delete(ph)
    db.session.flush()
    point_removed=False
    if pt and PanoramaPhoto.query.filter_by(point_id=pt.id).count()==0:
        db.session.delete(pt);point_removed=True
    db.session.add(AuditEvent(user_id=session.get("user_id"),event_type="PANORAMA_DELETE",entity_type="panorama_photo",entity_id=str(photo_id),detail=f"{original} · ponto vazio removido: {'sim' if point_removed else 'não'}"))
    db.session.commit()
    status="PENDENTE"
    if location_id:
        locrow=next((x for x in _panorama_payload() if x["id"]==location_id),None)
        if locrow:
            status=locrow["status"]
    return jsonify({"ok":True,"location_id":location_id,"status":status,"point_removed":point_removed})


# ===== V50.0 — Command Center 360 3.0 / Configuração ADM / Contratos 360 =====
V50_SETTINGS_PATH = DATA_DIR / "v50_admin_settings.json"

def _v50_settings():
    defaults = {
        "alert_activity_days": 7,
        "alert_pending_days": 7,
        "gps_radius_m": FIELD_NEARBY_RADIUS_M,
        "dashboard_refresh_seconds": 60,
        "gps_interval_seconds": 300,
        "forecast_days_emv": 3,
        "forecast_days_garagem": 3,
        "forecast_days_recarga": 3,
    }
    try:
        if V50_SETTINGS_PATH.exists():
            saved = json.loads(V50_SETTINGS_PATH.read_text(encoding="utf-8"))
            defaults.update({k:v for k,v in saved.items() if k in defaults})
    except Exception:
        app.logger.exception("Falha ao carregar configurações V50")
    return defaults

@app.get("/configuracoes")
@manager_required
def v50_settings_page():
    return render_template("admin_settings_v50.html", app_release=APP_RELEASE)

@app.route("/api/configuracoes", methods=["GET","POST"])
@manager_required
def v50_settings_api():
    if request.method == "GET":
        return jsonify({"ok":True,"settings":_v50_settings(),"release":APP_RELEASE})
    payload=request.get_json(silent=True) or {}
    current=_v50_settings()
    ranges={"alert_activity_days":(1,90),"alert_pending_days":(1,90),"gps_radius_m":(100,20000),"dashboard_refresh_seconds":(15,900),"gps_interval_seconds":(60,1800),"forecast_days_emv":(1,30),"forecast_days_garagem":(1,30),"forecast_days_recarga":(1,30)}
    for key,(lo,hi) in ranges.items():
        if key in payload:
            try: val=int(payload[key])
            except Exception: return jsonify({"ok":False,"error":f"Valor inválido para {key}."}),400
            if not lo <= val <= hi: return jsonify({"ok":False,"error":f"{key} deve ficar entre {lo} e {hi}."}),400
            current[key]=val
    V50_SETTINGS_PATH.write_text(json.dumps(current,ensure_ascii=False,indent=2),encoding="utf-8")
    db.session.add(AuditEvent(user_id=session.get("user_id"),event_type="CONFIG_UPDATE",entity_type="settings",entity_id="v50",detail=json.dumps(current,ensure_ascii=False)))
    db.session.commit()
    return jsonify({"ok":True,"settings":current})

@app.get("/api/command-center/hoje")
@login_required
def v50_command_center_today():
    if session.get("role") not in ("manager","manager_field"):
        return jsonify({"ok":False,"error":"Sem permissão."}),403
    today=datetime.utcnow().date()
    inv_today=Inventory.query.filter(func.date(Inventory.created_at)==today).count()
    chip_today=ChipSwap.query.filter(func.date(ChipSwap.started_at)==today).count()
    chip_done_today=ChipSwap.query.filter(func.date(ChipSwap.completed_at)==today).count()
    visits_today=HardwareFieldVisit.query.filter(HardwareFieldVisit.visit_date==today).count()
    visits_final_today=HardwareFieldVisit.query.filter(HardwareFieldVisit.visit_date==today,HardwareFieldVisit.status=="FINALIZADO").count()
    emv_today=EmvChipSwap.query.filter(func.date(EmvChipSwap.started_at)==today).count()
    emv_done_today=EmvChipSwap.query.filter(func.date(EmvChipSwap.completed_at)==today).count()
    tech_today=db.session.query(TechnicianPosition.user_id).filter(func.date(TechnicianPosition.captured_at)==today).distinct().count()
    return jsonify({"ok":True,"date":today.isoformat(),"field":{"inventory":inv_today,"chip_started":chip_today,"chip_done":chip_done_today,"technicians":tech_today},"implantation":{"visits":visits_today,"visits_finalized":visits_final_today,"emv_started":emv_today,"emv_done":emv_done_today}})


@app.get("/api/operational-forecast")
@login_required
def operational_forecast_api():
    module=(request.args.get("module") or "").strip().lower()
    cfg=_v50_settings()
    mapping={
        "emv": (EmvChipSwap, "forecast_days_emv"),
        "garagem": (GarageChipSwap, "forecast_days_garagem"),
        "recarga": (ChipSwap, "forecast_days_recarga"),
    }
    if module not in mapping:
        return jsonify({"ok":False,"error":"Módulo inválido."}),400
    model,key=mapping[module]; days=max(1,int(cfg.get(key,3) or 3))
    now=datetime.utcnow(); cutoff=now-timedelta(days=days)
    done_status=("CONCLUÍDA","CONCLUIDA","CONCLUÍDO","CONCLUIDO")
    active_terminals={str(x[0]).strip() for x in db.session.query(OperationalBaseItem.terminal).filter(OperationalBaseItem.module==module,OperationalBaseItem.active.is_(True)).all() if x[0]}
    if module=="emv":
        completed_terminals={str(x[0]).strip() for x in db.session.query(EmvChipSwap.terminal).filter(EmvChipSwap.status.in_(done_status)).all() if x[0]}
    elif module=="garagem":
        completed_terminals={str(x[0]).strip() for x in db.session.query(GarageChipBase.terminal).join(GarageChipSwap,GarageChipSwap.base_id==GarageChipBase.id).filter(GarageChipSwap.status.in_(done_status)).all() if x[0]}
    else:
        completed_terminals={str(x[0]).strip() for x in db.session.query(BaseAsset.terminal_number).join(ChipSwap,ChipSwap.base_asset_id==BaseAsset.id).filter(ChipSwap.status.in_(done_status)).all() if x[0]}
    total=len(active_terminals|completed_terminals)
    done=len(completed_terminals)
    in_progress=model.query.filter(model.status=="EM ANDAMENTO").count()
    if total < done+in_progress: total=done+in_progress
    pending=max(0,total-done-in_progress)
    recent=model.query.filter(model.completed_at.isnot(None),model.completed_at>=cutoff).count()
    daily=recent/days if days else 0
    eta_days=(pending/daily) if daily>0 and pending>0 else (0 if pending==0 else None)
    eta_date=(now+timedelta(days=eta_days)).date().isoformat() if eta_days is not None else None
    half=max(1.0,days/2); split=now-timedelta(days=half)
    recent2=model.query.filter(model.completed_at.isnot(None),model.completed_at>=split).count()
    previous=max(0,recent-recent2)
    r2=recent2/half; r1=previous/max(0.5,days-half)
    if recent < 2: trend="SEM DADOS"
    elif r2 > r1*1.12: trend="ACELERANDO"
    elif r2 < r1*0.88: trend="DESACELERANDO"
    else: trend="ESTÁVEL"
    confidence="ALTA" if recent>=max(10,days*3) else ("MÉDIA" if recent>=max(3,days) else "BAIXA")
    return jsonify({"ok":True,"release":APP_RELEASE,"module":module,"window_days":days,"total":total,"done":done,"in_progress":in_progress,"pending":pending,"completed_in_window":recent,"daily_rate":round(daily,2),"eta_days":round(eta_days,1) if eta_days is not None else None,"eta_date":eta_date,"trend":trend,"confidence":confidence})

@app.get("/notificacoes")
@login_required
def notifications_page():
    if session.get('role') not in ('manager','manager_field'):
        return redirect(url_for('dashboard_landing'))
    return render_template('notifications.html', app_release=APP_RELEASE)


@app.get("/api/notificacoes")
@login_required
def v50_notifications_api():
    if session.get("role") not in ("manager","manager_field"):
        return jsonify({"ok":True,"count":0,"items":[]})
    # Notificações derivadas das exceções operacionais, sem duplicar uma nova tabela nesta versão.
    official_path=DATA_DIR / "atm_official_082026.json"
    try:
        official_atms=json.loads(official_path.read_text(encoding="utf-8"))
    except Exception:
        official_atms=[]
    atm_without_tv=sum(1 for x in official_atms if not str(x.get("teamviewer_id") or "").strip())
    inv_div=Inventory.query.filter(func.coalesce(Inventory.divergence,"")!="").count()
    visits=HardwareFieldVisit.query.all()
    rv_pending=sum(1 for v in visits if 'PEND' in (v.conclusion_status or '').upper())
    emv_pending=EmvChipSwap.query.filter(~func.upper(func.coalesce(EmvChipSwap.status,"" )).startswith("CONCLU")).count()
    items=[
      {"domain":"FIELD","severity":"ATENÇÃO","title":"ATMs sem TeamViewer","count":atm_without_tv,"url":"/dashboard/atm?teamviewer_missing=1"},
      {"domain":"FIELD","severity":"ATENÇÃO","title":"Divergências de inventário","count":inv_div,"url":"/dashboard/field"},
      {"domain":"IMPLANTAÇÃO","severity":"ATENÇÃO","title":"Visitas com pendências","count":rv_pending,"url":"/implantacao-hardware/dashboard"},
      {"domain":"IMPLANTAÇÃO","severity":"INFO","title":"EMV Trilhos pendentes","count":emv_pending,"url":"/implantacao-hardware/dashboard"},
    ]
    items=[x for x in items if x["count"]>0]
    return jsonify({"ok":True,"count":sum(x["count"] for x in items),"items":items})

@app.get("/dashboard/contratos-atm")
@manager_required
def v50_contracts_page():
    return render_template("contracts_atm_v50.html", app_release=APP_RELEASE)

@app.get("/api/dashboard/contratos-atm")
@manager_required
def v50_contracts_api():
    official_path=DATA_DIR / "atm_official_082026.json"
    financial_path=DATA_DIR / "atm_financial_082026.json"
    try: assets=json.loads(official_path.read_text(encoding="utf-8"))
    except Exception: assets=[]
    try: financial=json.loads(financial_path.read_text(encoding="utf-8"))
    except Exception: financial={"models":[],"contracts":[]}
    unit={str(x.get("model") or "").upper():float(x.get("unit_value") or 0) for x in financial.get("models",[])}
    by_contract={}
    for a in assets:
        contract=(a.get("contract") or "Sem contrato").strip() or "Sem contrato"
        own=(a.get("ownership") or "Não informado").strip() or "Não informado"
        model=(a.get("model") or "Não informado").strip() or "Não informado"
        leasing="LEAS" in own.upper()
        monthly=0.0 if leasing else unit.get(model.upper(),0.0)
        row=by_contract.setdefault(contract,{"contract":contract,"qty":0,"leasing":0,"rental":0,"monthly":0.0,"locations":set(),"models":{}})
        row["qty"]+=1; row["leasing"]+=1 if leasing else 0; row["rental"]+=0 if leasing else 1; row["monthly"]+=monthly
        if a.get("locality"): row["locations"].add(a.get("locality"))
        row["models"][model]=row["models"].get(model,0)+1
    rows=[]
    for row in by_contract.values():
        row["location_count"]=len(row.pop("locations")); row["annual"]=row["monthly"]*12; rows.append(row)
    rows.sort(key=lambda x:x["monthly"],reverse=True)
    total_month=sum(x["monthly"] for x in rows)
    return jsonify({"ok":True,"release":APP_RELEASE,"park_total":len(assets),"rental":sum(x["rental"] for x in rows),"leasing":sum(x["leasing"] for x in rows),"monthly":total_month,"annual":total_month*12,"contracts":rows})


def _v56a_backfill_worker():
    global _V56A_BACKFILL
    with app.app_context():
        try:
            _V56A_BACKFILL.update({"running":True,"processed":0,"error":None})
            while True:
                batch=(TopDeskTicket.query.filter(TopDeskTicket.created_at.is_(None)).order_by(TopDeskTicket.id).limit(250).all())
                if not batch: break
                for t in batch:
                    parsed=_td_dt(t.created_at_text)
                    line,station,model=_td_object_parts(t.object_id)
                    t.created_at=parsed or t.imported_at or t.last_import_at or datetime.utcnow()
                    t.line_code=line or None; t.station_code=station or None; t.model_code=model or None
                db.session.commit(); _V56A_BACKFILL["processed"]+=len(batch); db.session.expire_all()
                time.sleep(0.20)
        except Exception as exc:
            db.session.rollback(); _V56A_BACKFILL["error"]=str(exc); app.logger.warning('V56-A.2 TopDesk background backfill: %s',exc)
        finally:
            _V56A_BACKFILL["running"]=False

def _start_v56a_backfill():
    if _V56A_BACKFILL.get("running"): return
    with _V56A_BACKFILL_LOCK:
        if _V56A_BACKFILL.get("running"): return
        _V56A_BACKFILL["running"]=True
        threading.Thread(target=_v56a_backfill_worker,name='v56a-topdesk-backfill',daemon=True).start()

def migrate_v56a_topdesk_dimensions():
    """V56-A: adiciona e normaliza dimensões TopDesk usadas pelos dashboards.

    Migração aditiva/idempotente. O backfill ocorre apenas em registros sem created_at
    e é commitado em lotes para não manter uma transação gigante.
    """
    inspector=db.inspect(db.engine)
    if not inspector.has_table('topdesk_tickets'): return
    cols={c['name'] for c in inspector.get_columns('topdesk_tickets')}
    dialect=db.engine.dialect.name
    dt_type='TIMESTAMP' if dialect=='postgresql' else 'DATETIME'
    additions=[
      ('created_at',dt_type),('line_code','VARCHAR(40)'),('station_code','VARCHAR(180)'),('model_code','VARCHAR(80)')
    ]
    try:
        with db.engine.begin() as conn:
            for name,sqltype in additions:
                if name not in cols:
                    conn.execute(text(f'ALTER TABLE topdesk_tickets ADD COLUMN {name} {sqltype}'))
    except Exception as exc:
        app.logger.warning('V56-A TopDesk columns: %s',exc); return

    commands=[
      'CREATE INDEX IF NOT EXISTS ix_topdesk_operator ON topdesk_tickets (operator)',
      'CREATE INDEX IF NOT EXISTS ix_topdesk_location_status ON topdesk_tickets (location_id, status)',
      'CREATE INDEX IF NOT EXISTS ix_topdesk_equipment_status ON topdesk_tickets (equipment_type, status)',
      'CREATE INDEX IF NOT EXISTS ix_topdesk_category_subcategory ON topdesk_tickets (category, subcategory)',
      'CREATE INDEX IF NOT EXISTS ix_topdesk_created_at ON topdesk_tickets (created_at)',
      'CREATE INDEX IF NOT EXISTS ix_topdesk_line_created ON topdesk_tickets (line_code, created_at)',
      'CREATE INDEX IF NOT EXISTS ix_topdesk_station_created ON topdesk_tickets (station_code, created_at)',
      'CREATE INDEX IF NOT EXISTS ix_topdesk_model_created ON topdesk_tickets (model_code, created_at)',
      'CREATE INDEX IF NOT EXISTS ix_topdesk_operator_created ON topdesk_tickets (operator, created_at)',
    ]
    try:
        with db.engine.begin() as conn:
            for cmd in commands: conn.execute(text(cmd))
    except Exception as exc: app.logger.warning('V56-A TopDesk indexes: %s',exc)

    # Após garantir colunas/índices, normaliza em background sem bloquear o boot.
    _start_v56a_backfill()



def migrate_v56a3_visit_contacts():
    """V56-A.3: coluna aditiva para múltiplos contatos do Relatório de Visita."""
    try:
        inspector=db.inspect(db.engine)
        if not inspector.has_table("hardware_field_visits"): return
        cols={c["name"] for c in inspector.get_columns("hardware_field_visits")}
        if "contacts_json" not in cols:
            db.session.execute(db.text("ALTER TABLE hardware_field_visits ADD COLUMN contacts_json TEXT")); db.session.commit()
    except Exception:
        db.session.rollback(); raise

def migrate_v55_performance_indexes():
    # Compatibilidade com chamadas/patches anteriores.
    return migrate_v56a_topdesk_dimensions()



# V62 — Dashboard Builder. Alteração/criação exclusivamente ADM (role=manager).
DASHBOARD_CATALOG = {
    "TOPDESK": {
        "label":"Chamados TOPdesk",
        "filters":["period","line","location","equipment_type","category","subcategory","operator","status"],
        "dimensions":["failure","line","location","model","object","operator"],
        "widgets":["kpi_tickets","kpi_objects","kpi_locations","monthly","rank_failure","rank_location","rank_model","rank_object","productivity","table"]
    },
    "GARAGE": {
        "label":"Troca de Chips Garagem",
        "filters":["company","model","status","technician"],
        "dimensions":["company","model","status","technician"],
        "widgets":["kpi_total","kpi_concluded","kpi_pending","progress","rank_company","rank_technician"]
    },
    "EMV": {
        "label":"Troca de Chips EMV – Trilhos",
        "filters":["company","line","station","status","test_result","technician"],
        "dimensions":["company","line","station","status","test_result","technician"],
        "widgets":["kpi_total","kpi_concluded","kpi_pending","progress","rank_station","rank_technician"]
    }
}

def _dashboard_admin_required():
    if session.get("role") != "manager": abort(403)

def _dash_cfg(row):
    try: cfg=json.loads(row.config_json or "{}")
    except Exception: cfg={}
    try: roles=json.loads(row.allowed_roles_json or "[]")
    except Exception: roles=[]
    return {"id":row.id,"name":row.name,"slug":row.slug,"data_source":row.data_source,"config":cfg,"published":bool(row.published),"tv_enabled":bool(row.tv_enabled),"tv_order":row.tv_order or 0,"tv_seconds":row.tv_seconds or 30,"allowed_roles":roles}

def _dashboard_visible(row):
    if session.get("role")=="manager": return True
    if not row.published: return False
    try: roles=json.loads(row.allowed_roles_json or "[]")
    except Exception: roles=[]
    return not roles or session.get("role") in roles

@app.context_processor
def _v62_dashboard_context():
    try:
        rows=DashboardDefinition.query.filter_by(published=True).order_by(DashboardDefinition.name).all() if session.get("user_id") else []
        return {"custom_dashboards":[_dash_cfg(x) for x in rows if _dashboard_visible(x)]}
    except Exception:
        return {"custom_dashboards":[]}

BUILTIN_DASHBOARD_CATALOG = [
    {"key":"overview","label":"Visão Geral","group":"VISÃO GERAL","icon":"▦","roles":[]},
    {"key":"execution","label":"Inventário","group":"ATIVIDADES","icon":"▥","roles":[]},
    {"key":"atm-inventory","label":"Dashboard ATM","group":"ATIVIDADES","icon":"▦","roles":[]},
    {"key":"financial-dashboard","label":"Dashboard Financeiro","group":"ATIVIDADES","icon":"$","roles":["manager","manager_field","atm_financial_admin"]},
    {"key":"journal","label":"Diário de bordo","group":"ATIVIDADES","icon":"◷","roles":[]},
    {"key":"topdesk","label":"Dashboard Chamados","group":"SUPORTE A CAMPO","icon":"⚡","roles":[]},
    {"key":"map","label":"Mapa","group":"ATIVIDADES","icon":"⌖","roles":[]},
    {"key":"panorama","label":"Visões panorâmicas","group":"ATIVIDADES","icon":"▤","roles":[]},
    {"key":"chips","label":"Troca de Chip Recarga","group":"ATIVIDADES","icon":"▣","roles":[]},
    {"key":"emv","label":"Troca Chips EMV · Trilhos","group":"IMPLANTAÇÃO","icon":"▣","roles":[]},
    {"key":"garage","label":"Dashboard Garagem","group":"IMPLANTAÇÃO","icon":"▦","roles":[]},
    {"key":"implantation-dashboard","label":"Dashboard Implantação","group":"IMPLANTAÇÃO","icon":"⌁","roles":[]},
    {"key":"ranking","label":"Ranking","group":"OPERAÇÃO","icon":"↗","roles":[]},
    {"key":"competition","label":"Concorrência","group":"OPERAÇÃO","icon":"◉","roles":[]},
]

def _builtin_dashboard_menu_items():
    try: saved={x.dashboard_key:x for x in BuiltinDashboardSetting.query.all()}
    except Exception: saved={}
    out=[]
    for idx,item in enumerate(BUILTIN_DASHBOARD_CATALOG,1):
        row=saved.get(item['key']); visible=True if row is None else bool(row.visible); order=(idx*10 if row is None else row.order_index)
        try: allowed=json.loads(row.allowed_roles_json or '[]') if row else list(item.get('roles') or [])
        except Exception: allowed=list(item.get('roles') or [])
        if not visible: continue
        role=session.get('role')
        if allowed and role not in allowed: continue
        out.append({**item,'order':order,'allowed_roles':allowed})
    return sorted(out,key=lambda x:(x['order'],x['label']))

@app.context_processor
def _builtin_dashboards_context():
    return {'builtin_dashboard_menu_items':_builtin_dashboard_menu_items() if session.get('user_id') else []}

@app.route('/gestao/configuracao-dashboards',methods=['GET','POST'])
@login_required
def builtin_dashboard_settings_page():
    if session.get('role')!='manager' and not _has_access('management.dashboard_config'): abort(403)
    if request.method=='POST':
        for idx,item in enumerate(BUILTIN_DASHBOARD_CATALOG,1):
            key=item['key']; row=BuiltinDashboardSetting.query.filter_by(dashboard_key=key).first()
            if not row: row=BuiltinDashboardSetting(dashboard_key=key); db.session.add(row)
            row.visible=request.form.get(f'visible__{key}')=='1'
            try: row.order_index=int(request.form.get(f'order__{key}') or idx*10)
            except Exception: row.order_index=idx*10
            roles=[x for x in request.form.getlist(f'roles__{key}') if x in ('manager','manager_field','technician','technician_implantation','consultation','hr','dispatcher','atm_financial_admin')]
            row.allowed_roles_json=json.dumps(roles,ensure_ascii=False); row.updated_by=session.get('user_id')
        db.session.commit()
        if not any(x.visible for x in BuiltinDashboardSetting.query.all()):
            first=BuiltinDashboardSetting.query.filter_by(dashboard_key='overview').first(); first.visible=True; db.session.commit(); flash('Visão Geral foi mantida ativa para evitar uma Central vazia.')
        else: flash('Configuração de dashboards atualizada.')
        return redirect('/gestao/configuracao-dashboards')
    saved={x.dashboard_key:x for x in BuiltinDashboardSetting.query.all()}; rows=[]
    for idx,item in enumerate(BUILTIN_DASHBOARD_CATALOG,1):
        row=saved.get(item['key'])
        try: roles=json.loads(row.allowed_roles_json or '[]') if row else list(item.get('roles') or [])
        except Exception: roles=[]
        rows.append({**item,'visible':True if row is None else row.visible,'order':idx*10 if row is None else row.order_index,'allowed_roles':roles})
    return render_template('builtin_dashboard_settings.html',rows=rows,app_release=APP_RELEASE)

@app.get('/configuracoes/dashboards')
@login_required
def dashboard_builder_page():
    _dashboard_admin_required()
    return render_template('dashboard_builder.html',app_release=APP_RELEASE,catalog=DASHBOARD_CATALOG)

@app.get('/api/dashboard-configs')
@login_required
def dashboard_configs_api():
    _dashboard_admin_required()
    return jsonify({"ok":True,"catalog":DASHBOARD_CATALOG,"dashboards":[_dash_cfg(x) for x in DashboardDefinition.query.order_by(DashboardDefinition.name).all()]})

@app.post('/api/dashboard-configs')
@login_required
def dashboard_configs_save():
    _dashboard_admin_required(); p=request.get_json(silent=True) or {}; name=str(p.get('name') or '').strip()
    if not name: return jsonify({'error':'Informe o nome da dashboard.'}),400
    slug=re.sub(r'[^a-z0-9]+','-',unicodedata.normalize('NFKD',name).encode('ascii','ignore').decode().lower()).strip('-') or f'dashboard-{int(time.time())}'
    rid=p.get('id'); row=db.session.get(DashboardDefinition,int(rid)) if rid else None
    if not row:
        base=slug; n=2
        while DashboardDefinition.query.filter_by(slug=slug).first(): slug=f'{base}-{n}'; n+=1
        row=DashboardDefinition(name=name,slug=slug,created_by=session['user_id']); db.session.add(row)
    row.name=name; row.data_source=str(p.get('data_source') or 'TOPDESK').upper(); row.config_json=json.dumps(p.get('config') or {},ensure_ascii=False)
    row.published=bool(p.get('published')); row.tv_enabled=bool(p.get('tv_enabled')); row.tv_order=int(p.get('tv_order') or 0); row.tv_seconds=max(10,min(600,int(p.get('tv_seconds') or 30)))
    row.allowed_roles_json=json.dumps(p.get('allowed_roles') or [],ensure_ascii=False); row.updated_at=datetime.utcnow(); db.session.commit()
    return jsonify({'ok':True,'dashboard':_dash_cfg(row)})

@app.delete('/api/dashboard-configs/<int:dashboard_id>')
@login_required
def dashboard_configs_delete(dashboard_id):
    _dashboard_admin_required(); row=db.session.get(DashboardDefinition,dashboard_id)
    if not row: return jsonify({'error':'Dashboard não encontrada.'}),404
    db.session.delete(row); db.session.commit(); return jsonify({'ok':True})

@app.get('/dashboards/<slug>')
@login_required
def custom_dashboard_page(slug):
    row=DashboardDefinition.query.filter_by(slug=slug).first_or_404()
    if not _dashboard_visible(row): abort(403)
    return render_template('custom_dashboard.html',app_release=APP_RELEASE,dashboard=_dash_cfg(row),catalog=DASHBOARD_CATALOG)

@app.get('/modo-tv/dashboards')
@login_required
def custom_dashboard_tv_page():
    rows=DashboardDefinition.query.filter_by(published=True,tv_enabled=True).order_by(DashboardDefinition.tv_order,DashboardDefinition.name).all()
    rows=[_dash_cfg(x) for x in rows if _dashboard_visible(x)]
    return render_template('dashboard_tv.html',app_release=APP_RELEASE,dashboards=rows)

@app.get('/api/dashboard-builder/source/<source>')
@login_required
def dashboard_builder_source(source):
    source=source.upper(); cfg=DASHBOARD_CATALOG.get(source)
    if not cfg: return jsonify({'error':'Fonte não suportada'}),404
    # TOPdesk usa o analytics já otimizado e aceita os mesmos filtros.
    if source=='TOPDESK':
        return topdesk_analytics_api()
    if source=='GARAGE':
        return garage_chip_dashboard_api()
    if source=='EMV':
        rows=emv_chip_list().get_json().get('rows',[])
        total=len(rows); done=sum(x.get('status')=='CONCLUÍDA' for x in rows); prog=sum(x.get('status')=='EM ANDAMENTO' for x in rows)
        def rank(key):
            d={}
            for x in rows:
                k=str(x.get(key) or 'Não informado'); d[k]=d.get(k,0)+1
            return [{'name':k,'count':v} for k,v in sorted(d.items(),key=lambda z:z[1],reverse=True)[:30]]
        return jsonify({'ok':True,'kpis':{'total':total,'concluded':done,'pending':total-done-prog,'in_progress':prog},'stations':rank('station'),'companies':rank('company'),'lines':rank('line'),'technicians':rank('technician')})
    return jsonify({'ok':True})





def seed_v67_materials():
    """Carga inicial configurável do Kit Técnico Field N2; depois tudo é administrado pela UI."""
    if MaterialCatalogItem.query.count()>0: return
    rows=[
      ('N2-001','Mochila de ferramentas STANLEY','STANLEY',''),('N2-002','Chave de Fenda - 3/16 x 4”','Gedore',''),('N2-003','Chave de Fenda - 1/8 x 4”','Gedore',''),('N2-004','Chave Philips - 3/16 x 4”','Gedore',''),('N2-005','Chave Philips - 1/8 x 4”','Gedore',''),('N2-006','Chave canhão - 7/32 mm','Gedore',''),('N2-007','Chave canhão - 8 mm','Gedore',''),('N2-008','Jogo de chaves allen c/ estojo 1,5mm à 10mm','MTX',''),('N2-009','Alicate Crimpador - RJ11 - RJ12 - RJ45','Universal',''),('N2-010','Alicate de Bico redondo longo Hikari HK - 507','Hikari','HK-507'),('N2-011','Mini Alicate universal Hikari HK - 502','Hikari','HK-502'),('N2-012','Decapador de fio 501 universal','Universal',''),('N2-013','Caixa organizadora mini','Universal',''),('N2-014','Pincel Antiestático ESD HK - 217','Hikari','HK-217'),('N2-015','Trincha 2 ¹/2 (pincel)','Vonder',''),('N2-016','Chave multiteste digital com display LCD','Sparta',''),('N2-017','Estilete emborrachado','Universal',''),('N2-018','Testador de cabos de rede','Universal',''),('N2-019','Multímetro digital HM 1001 com bateria 9V','Hikari','HM 1001'),('N2-020','Alicate de corte rente HK - 170','Hikari','HK-170'),('N2-021','Miniteclado USB com fio','Knupp',''),('N2-022','Extractor de pic San PLCC','Pinça',''),('N2-023','Alicate de corte diagonal 6”','Hikari',''),('N2-024','Pinça','Hikari',''),('N2-025','Alicate Puntch Down universal','Universal','')]
    uid=None
    for code,desc,brand,model in rows:
        db.session.add(MaterialCatalogItem(code=code,category='FERRAMENTA',description=desc,brand=brand,model=model,unit='UN',control_type='DEVOLVIVEL',quantity_mode='INTEIRO',active=True,created_by=uid))
    db.session.flush(); kit=MaterialKit(name='Kit Técnico Field N2',description='Kit inicial com os 25 itens do Termo de Responsabilidade de Equipamentos.',active=True,created_by=uid);db.session.add(kit);db.session.flush()
    for m in MaterialCatalogItem.query.filter(MaterialCatalogItem.code.like('N2-%')).all():db.session.add(MaterialKitItem(kit_id=kit.id,material_id=m.id,quantity=1))
    db.session.commit()

# ==================== V67 · DOSSIÊ & MATERIAIS ====================
def _materials_require(permission):
    if not _has_access(permission): abort(403)

def _material_quantity(m, raw, field='Quantidade'):
    try:
        qty=float(raw)
    except Exception:
        raise ValueError(f'{field} inválida.')
    if qty <= 0:
        raise ValueError(f'{field} deve ser maior que zero.')
    mode=(getattr(m,'quantity_mode',None) or ('DECIMAL' if (getattr(m,'category','') or '').upper()=='CONSUMIVEL' else 'INTEIRO')).upper()
    if mode=='INTEIRO' and abs(qty-round(qty)) > 1e-9:
        raise ValueError(f'{m.description}: a quantidade deve ser um número inteiro.')
    return float(round(qty)) if mode=='INTEIRO' else qty

def _material_qty_text(qty):
    try:
        v=float(qty or 0)
        return str(int(round(v))) if abs(v-round(v)) < 1e-9 else ('%g' % v)
    except Exception:
        return str(qty or '')

def _doc_payload(d):
    u=db.session.get(User,d.user_id)
    items=CollaboratorDocumentItem.query.filter_by(document_id=d.id).order_by(CollaboratorDocumentItem.id).all()
    return {"id":d.id,"code":d.document_code or f"DOC-{d.id:06d}","type":d.document_type,"title":d.title,"user_id":d.user_id,"user_name":u.name if u else "—","status":d.status,"delivery_date":d.delivery_date.isoformat() if d.delivery_date else None,"sent_at":d.sent_at.isoformat() if d.sent_at else None,"signed_at":d.signed_at.isoformat() if d.signed_at else None,"notes":d.notes or "","correction_note":d.correction_note or "","has_pdf":bool(d.pdf_file),"items":[{"id":x.id,"material_id":x.material_id,"description":x.description,"brand":x.brand or "","model":x.model or "","quantity":x.quantity,"unit":x.unit or "UN","condition":x.condition or "BOM","notes":x.notes or ""} for x in items]}

def _material_pdf_bytes(doc):
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.enums import TA_CENTER, TA_LEFT
        from reportlab.lib import colors
        from reportlab.lib.utils import ImageReader
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image, PageBreak, Image, KeepTogether
    except Exception as exc:
        raise RuntimeError("Dependência reportlab não instalada") from exc
    out=io.BytesIO(); styles=getSampleStyleSheet()
    body=ParagraphStyle('trm_body',parent=styles['BodyText'],fontSize=9,leading=12,spaceAfter=5)
    small=ParagraphStyle('trm_small',parent=body,fontSize=8,leading=10)
    center=ParagraphStyle('trm_center',parent=styles['Heading2'],alignment=TA_CENTER,fontSize=13,leading=16,spaceAfter=8)
    title=ParagraphStyle('trm_title',parent=styles['Heading1'],alignment=TA_CENTER,fontSize=14,leading=17,spaceAfter=6)
    pdf=SimpleDocTemplate(out,pagesize=A4,rightMargin=32,leftMargin=32,topMargin=28,bottomMargin=28,title=doc.document_code or 'Termo de Responsabilidade')
    u=db.session.get(User,doc.user_id); rows=CollaboratorDocumentItem.query.filter_by(document_id=doc.id).order_by(CollaboratorDocumentItem.id).all()
    story=[]
    logo=Path(app.root_path)/'static'/'logo-tecsomobi.png'
    if logo.exists():
        try:
            im=Image(str(logo),width=115,height=38); im.hAlign='LEFT'; story += [im,Spacer(1,4)]
        except Exception: pass
    is_return=(doc.document_type=='TERMO_DEVOLUCAO')
    heading="TERMO DE DEVOLUÇÃO DE MATERIAIS / FERRAMENTAS" if is_return else "TERMO DE RESPONSABILIDADE DE EQUIPAMENTOS"
    intro="Declaro o registro da devolução à Tecsomobi dos materiais, ferramentas e equipamentos abaixo especificados, nas quantidades e condições conferidas no ato do recebimento." if is_return else "Declaro que recebi da Tecsomobi Fábrica de Software e Inteligência Digital Ltda., a título de empréstimo e/ou para uso profissional, os materiais, ferramentas e equipamentos abaixo especificados."
    story += [Paragraph(heading,title),Paragraph(intro,body),Spacer(1,5)]
    meta=[[Paragraph(f"<b>Documento:</b> {doc.document_code or '—'}",small),Paragraph(f"<b>Colaborador:</b> {u.name if u else '—'}",small)],
          [Paragraph(f"<b>Empresa:</b> {getattr(u,'company',None) or '—'}",small),Paragraph(f"<b>Cargo:</b> {getattr(u,'job_title',None) or '—'}",small)],
          [Paragraph(f"<b>{'Data da devolução' if is_return else 'Data da entrega'}:</b> {doc.delivery_date.strftime('%d/%m/%Y') if doc.delivery_date else '—'}",small),Paragraph(f"<b>Itens:</b> {len(rows)}",small)]]
    mt=Table(meta,colWidths=[260,260]); mt.setStyle(TableStyle([('BOX',(0,0),(-1,-1),.5,colors.HexColor('#b8c5d3')),('INNERGRID',(0,0),(-1,-1),.25,colors.HexColor('#d9e2ec')),('BACKGROUND',(0,0),(-1,-1),colors.HexColor('#f7f9fc')),('VALIGN',(0,0),(-1,-1),'TOP'),('PADDING',(0,0),(-1,-1),6)])); story += [mt,Spacer(1,9)]
    data=[["Item","Descrição","Marca / Modelo","Qtd.","Un.","Estado"]]
    for i,x in enumerate(rows,1): data.append([str(i),Paragraph(x.description,small),Paragraph(" / ".join(y for y in (x.brand,x.model) if y) or "—",small),_material_qty_text(x.quantity),x.unit or "UN",x.condition or "BOM"])
    t=Table(data,colWidths=[28,205,135,42,38,62],repeatRows=1);t.setStyle(TableStyle([('BACKGROUND',(0,0),(-1,0),colors.HexColor('#e8eef5')),('GRID',(0,0),(-1,-1),.4,colors.HexColor('#9aa9b8')),('FONTNAME',(0,0),(-1,0),'Helvetica-Bold'),('VALIGN',(0,0),(-1,-1),'TOP'),('FONTSIZE',(0,0),(-1,-1),8),('PADDING',(0,0),(-1,-1),4)])); story += [t,Spacer(1,10)]
    terms=[
      "Responsabilizo-me por manter os materiais, ferramentas e equipamentos sob minha guarda em adequado estado de conservação e funcionamento, utilizando-os exclusivamente para as atividades profissionais e observando as orientações internas da empresa.",
      "Em caso de dano, inutilização, perda ou extravio, comprometo-me a comunicar imediatamente a empresa/Recursos Humanos para registro, avaliação e providências cabíveis.",
      "Ao término dos serviços, mudança de função que dispense o uso, solicitação da empresa ou rescisão do vínculo, comprometo-me a devolver os bens devolvíveis que estiverem sob minha responsabilidade, considerando o desgaste normal decorrente do uso adequado.",
      "Os materiais e equipamentos sob minha responsabilidade poderão ser conferidos ou inspecionados pela empresa durante o período de custódia.",
      "Declaro que conferi a relação acima antes do aceite. Eventuais divergências devem ser apontadas por meio da opção Solicitar correção antes da assinatura."
    ]
    if is_return:
        story.append(Paragraph("<b>REGISTRO DA DEVOLUÇÃO</b>",body))
        story.append(Paragraph("• Os itens relacionados acima foram registrados como devolvidos nas condições informadas. Itens vinculados à carga do colaborador tiveram a respectiva quantidade baixada; devoluções avulsas permanecem identificadas para auditoria.",body))
    else:
        story.append(Paragraph("<b>RESPONSABILIDADES E ORIENTAÇÕES</b>",body))
        for tx in terms: story.append(Paragraph("• "+tx,body))
    if doc.notes: story += [Spacer(1,3),Paragraph(f"<b>{'Observação da devolução' if is_return else 'Observação da entrega'}:</b> {doc.notes}",small)]
    story += [Spacer(1,9),Paragraph('<b>CARIMBO DE RECEBIMENTO / DEVOLUÇÃO</b>' if is_return else '<b>ACEITE ELETRÔNICO</b>',body),Paragraph(f"{'Devolução registrada em' if is_return else 'Aceite realizado em'}: <b>{doc.signed_at.strftime('%d/%m/%Y %H:%M:%S') if doc.signed_at else '—'}</b>",small),Paragraph(f"Colaborador: <b>{u.name if u else '—'}</b>",small)]
    sigraw=None
    try:
        if doc.signature_file:
            sigraw=_r2_get_bytes(doc.signature_file[4:]) if doc.signature_file.startswith('r2__') else (UPLOAD_DIR/doc.signature_file).read_bytes()
    except Exception: sigraw=None
    if sigraw:
        try:
            sig=Image(io.BytesIO(sigraw),width=180,height=56); sig.hAlign='LEFT'; story += [Spacer(1,4),sig,Paragraph("Assinatura manuscrita eletrônica do colaborador",small)]
        except Exception: pass
    if is_return and getattr(doc,'return_receiver_signature_file',None):
        try:
            rf=doc.return_receiver_signature_file; rraw=_r2_get_bytes(rf[4:]) if rf.startswith('r2__') else (UPLOAD_DIR/rf).read_bytes()
            receiver=db.session.get(User,doc.return_receiver_id) if doc.return_receiver_id else None
            rsig=Image(io.BytesIO(rraw),width=180,height=56); rsig.hAlign='LEFT'
            story += [Spacer(1,8),rsig,Paragraph(f"Recebido por: <b>{receiver.name if receiver else 'Responsável'}</b>",small),Paragraph(f"Conferência em: <b>{doc.return_received_at.strftime('%d/%m/%Y %H:%M:%S') if doc.return_received_at else '—'}</b>",small)]
        except Exception: pass
    story += [Spacer(1,5),Paragraph(f"{'Carimbo de devolução' if is_return else 'Carimbo de aceite'}: {doc.document_code or '—'} · usuário #{doc.user_id} · {doc.signed_at.strftime('%d/%m/%Y %H:%M:%S UTC') if doc.signed_at else '—'}",small),Paragraph("Este PDF é a versão assinada vinculada ao dossiê digital do colaborador. O original eletrônico e sua trilha de auditoria permanecem armazenados no sistema.",small)]
    pdf.build(story); return out.getvalue()

@app.get('/documentos-materiais/raio-x')
@login_required
def materials_xray_page():
    _materials_require('materials.dossier.view'); return render_template('materials_xray.html',app_release=APP_RELEASE)

def _materials_xray_rows():
    users=(User.query.filter(User.active.is_(True), ~User.role.in_(['customer','consultation'])).order_by(User.name).all())
    user_ids=[u.id for u in users]; doc_stats={}
    if user_ids:
        for uid,st in db.session.query(CollaboratorDocument.user_id,CollaboratorDocument.status).filter(CollaboratorDocument.user_id.in_(user_ids)).all():
            x=doc_stats.setdefault(uid,{'total':0,'signed':0,'pending':0}); x['total']+=1
            if st=='ASSINADO': x['signed']+=1
            # V69.5: pendência do colaborador = próxima ação depende dele.
            # CORRECAO_SOLICITADA/RASCUNHO dependem da gestão, não do colaborador.
            if st == 'AGUARDANDO_ACEITE': x['pending']+=1
    out=[]
    for u in users:
        ds=doc_stats.get(u.id,{'total':0,'signed':0,'pending':0})
        reqs=MaterialRequest.query.filter_by(user_id=u.id).filter(MaterialRequest.status.in_(['SOLICITADO','EM_ANALISE','APROVADO','EM_COMPRA','DISPONIVEL'])).all()
        mov=MaterialMovement.query.filter_by(user_id=u.id).all(); bal={}
        for m in mov: bal[m.material_id]=bal.get(m.material_id,0)+(m.quantity if m.movement_type in ('ENTREGA','SUBSTITUICAO_ENTRADA','TRANSFERENCIA_ENTRADA') else -m.quantity)
        load=sum(1 for q in bal.values() if q>0); issues=ds['pending']; reasons=[]
        if ds['pending']: reasons.append(f"{ds['pending']} termo/documento(s) aguardando ação")
        status='REGULAR' if issues==0 else 'PENDENTE'
        out.append({'id':u.id,'name':u.name,'company':u.company or '', 'job_title':u.job_title or '', 'documents_total':ds['total'],'terms_signed':ds['signed'],'terms_pending':ds['pending'],'materials_open':load,'requests_open':len(reqs),'issues':issues,'status':status,'reason':' · '.join(reasons) if reasons else 'Nenhuma pendência operacional'})
    return out

@app.get('/api/materials/xray')
@login_required
def materials_xray_api():
    _materials_require('materials.dossier.view'); out=_materials_xray_rows()
    return jsonify({'ok':True,'rows':out,'summary':{'total':len(out),'regular':sum(1 for x in out if x['status']=='REGULAR'),'pending':sum(1 for x in out if x['status']=='PENDENTE'),'without_term':sum(1 for x in out if x['terms_signed']==0),'materials_open':sum(1 for x in out if x['materials_open']>0)}})

@app.get('/api/materials/xray/<int:user_id>/pending')
@login_required
def materials_xray_pending_api(user_id):
    _materials_require('materials.dossier.view'); u=db.session.get(User,user_id) or abort(404); items=[]
    for d in CollaboratorDocument.query.filter_by(user_id=user_id).filter(CollaboratorDocument.status=='AGUARDANDO_ACEITE').order_by(CollaboratorDocument.created_at.desc()).all():
        items.append({'type':'Termo / Documento','code':d.document_code or '—','description':d.title or d.document_type or 'Documento','status':d.status,'date':d.created_at.strftime('%d/%m/%Y') if d.created_at else '—'})
    return jsonify({'ok':True,'user':{'id':u.id,'name':u.name,'company':u.company or '', 'job_title':u.job_title or ''},'items':items})

@app.get('/documentos-materiais/raio-x/exportar.xlsx')
@login_required
def materials_xray_export():
    _materials_require('materials.dossier.view'); rows=_materials_xray_rows(); wb=Workbook(); ws=wb.active; ws.title='Raio-X Colaboradores'
    ws.append(['Colaborador','Empresa','Cargo','Documentos','Termos assinados','Termos pendentes','Materiais em carga','Solicitações','Situação','Motivo'])
    for c in ws[1]: c.font=Font(bold=True)
    for x in rows: ws.append([x['name'],x['company'],x['job_title'],x['documents_total'],x['terms_signed'],x['terms_pending'],x['materials_open'],x['requests_open'],x['status'],x['reason']])
    for col in range(1,ws.max_column+1): ws.column_dimensions[get_column_letter(col)].width=min(45,max(12,max(len(str(ws.cell(r,col).value or '')) for r in range(1,ws.max_row+1))+2))
    bio=io.BytesIO(); wb.save(bio); bio.seek(0); return send_file(bio,as_attachment=True,download_name='raio_x_colaboradores.xlsx',mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')

@app.get('/documentos-materiais')
@login_required
def materials_home_page():
    return render_template('materials.html',app_release=APP_RELEASE)

@app.get('/api/materials/catalog')
@login_required
def materials_catalog_api():
    _materials_require('materials.catalog.view')
    rows=MaterialCatalogItem.query.order_by(MaterialCatalogItem.active.desc(),MaterialCatalogItem.description).all()
    return jsonify({'ok':True,'rows':[{'id':x.id,'code':x.code,'category':x.category,'description':x.description,'brand':x.brand or '', 'model':x.model or '', 'unit':x.unit,'control_type':x.control_type,'quantity_mode':x.quantity_mode or 'INTEIRO','active':x.active} for x in rows]})

@app.post('/api/materials/catalog')
@login_required
def materials_catalog_save_api():
    _materials_require('materials.catalog.manage'); d=request.get_json(silent=True) or {}; mid=d.get('id'); row=db.session.get(MaterialCatalogItem,int(mid)) if mid else MaterialCatalogItem(created_by=session['user_id'])
    if not row: return jsonify({'ok':False,'error':'Item não encontrado.'}),404
    code=(d.get('code') or '').strip().upper(); desc=(d.get('description') or '').strip()
    if not code or not desc:return jsonify({'ok':False,'error':'Código e descrição são obrigatórios.'}),400
    dup=MaterialCatalogItem.query.filter(func.upper(MaterialCatalogItem.code)==code)
    if row.id:dup=dup.filter(MaterialCatalogItem.id!=row.id)
    if dup.first():return jsonify({'ok':False,'error':'Código já cadastrado.'}),409
    row.code=code;row.description=desc;row.category=(d.get('category') or 'FERRAMENTA').upper();row.brand=(d.get('brand') or '').strip();row.model=(d.get('model') or '').strip();row.unit=(d.get('unit') or 'UN').upper();row.control_type=(d.get('control_type') or 'DEVOLVIVEL').upper();row.quantity_mode=(d.get('quantity_mode') or ('DECIMAL' if row.category=='CONSUMIVEL' else 'INTEIRO')).upper();row.active=bool(d.get('active',True));db.session.add(row);db.session.add(AuditEvent(user_id=session['user_id'],event_type='MATERIAL_CATALOG_SAVE',entity_type='material',entity_id=str(row.id or ''),detail=f'{code} · {desc}'));db.session.commit();return jsonify({'ok':True,'id':row.id})

@app.delete('/api/materials/catalog/<int:mid>')
@login_required
def materials_catalog_delete_api(mid):
    _materials_require('materials.catalog.manage'); row=db.session.get(MaterialCatalogItem,mid)
    if not row:return jsonify({'ok':False,'error':'Item não encontrado.'}),404
    used=CollaboratorDocumentItem.query.filter_by(material_id=mid).first() or MaterialMovement.query.filter_by(material_id=mid).first()
    if used: row.active=False; action='INATIVADO'
    else: db.session.delete(row); action='EXCLUIDO'
    db.session.add(AuditEvent(user_id=session['user_id'],event_type='MATERIAL_CATALOG_'+action,entity_type='material',entity_id=str(mid),detail=row.description));db.session.commit();return jsonify({'ok':True,'action':action})

@app.get('/api/materials/kits')
@login_required
def materials_kits_api():
    _materials_require('materials.catalog.view'); kits=MaterialKit.query.order_by(MaterialKit.active.desc(),MaterialKit.name).all(); out=[]
    for k in kits:
        its=db.session.query(MaterialKitItem,MaterialCatalogItem).join(MaterialCatalogItem,MaterialCatalogItem.id==MaterialKitItem.material_id).filter(MaterialKitItem.kit_id==k.id).all();out.append({'id':k.id,'name':k.name,'description':k.description or '', 'active':k.active,'items':[{'material_id':mi.material_id,'quantity':mi.quantity,'description':m.description,'brand':m.brand or '', 'model':m.model or '', 'unit':m.unit,'quantity_mode':m.quantity_mode or 'INTEIRO'} for mi,m in its]})
    return jsonify({'ok':True,'rows':out})

@app.post('/api/materials/kits')
@login_required
def materials_kits_save_api():
    _materials_require('materials.kits.manage'); d=request.get_json(silent=True) or {}; kid=d.get('id'); k=db.session.get(MaterialKit,int(kid)) if kid else MaterialKit(created_by=session['user_id'])
    if not k:return jsonify({'ok':False,'error':'Kit não encontrado.'}),404
    k.name=(d.get('name') or '').strip();k.description=(d.get('description') or '').strip();k.active=bool(d.get('active',True))
    if not k.name:return jsonify({'ok':False,'error':'Informe o nome do kit.'}),400
    db.session.add(k);db.session.flush();MaterialKitItem.query.filter_by(kit_id=k.id).delete()
    for x in d.get('items') or []:
        try:
            mid=int(x.get('material_id')); m=db.session.get(MaterialCatalogItem,mid); qty=_material_quantity(m,x.get('quantity') or 1) if m else 0
        except (ValueError,TypeError): continue
        if qty>0:db.session.add(MaterialKitItem(kit_id=k.id,material_id=mid,quantity=qty))
    db.session.commit();return jsonify({'ok':True,'id':k.id})

@app.get('/api/materials/collaborators')
@login_required
def materials_collaborators_api():
    _materials_require('materials.delivery.create'); rows=User.query.filter(User.active==True,User.role.in_(['technician','technician_implantation','manager_field','dispatcher'])).order_by(User.name).all();return jsonify({'ok':True,'rows':[{'id':u.id,'name':u.name,'role':u.role,'company':u.company or '', 'job_title':u.job_title or ''} for u in rows]})

@app.post('/api/materials/deliveries')
@login_required
def materials_delivery_create_api():
    _materials_require('materials.delivery.create'); d=request.get_json(silent=True) or {}; users=[]
    for uid in d.get('user_ids') or []:
        try:
            u=db.session.get(User,int(uid));
            if u and u.active:users.append(u)
        except:pass
    items=d.get('items') or []
    if not users or not items:return jsonify({'ok':False,'error':'Selecione colaborador(es) e pelo menos um item.'}),400
    created=[]
    try: delivery_date=datetime.strptime(d.get('delivery_date') or '', '%Y-%m-%d').date()
    except: delivery_date=datetime.now().date()
    for u in users:
        doc=CollaboratorDocument(user_id=u.id,created_by=session['user_id'],delivery_date=delivery_date,notes=(d.get('notes') or '').strip(),status='RASCUNHO');db.session.add(doc);db.session.flush();doc.document_code=f'TRM-{doc.id:06d}'
        for x in items:
            try: mid=int(x.get('material_id'))
            except: continue
            m=db.session.get(MaterialCatalogItem,mid)
            if not m: continue
            try: qty=_material_quantity(m,x.get('quantity') or 0)
            except ValueError as exc: return jsonify({'ok':False,'error':str(exc)}),400
            db.session.add(CollaboratorDocumentItem(document_id=doc.id,material_id=m.id,description=m.description,brand=m.brand,model=m.model,quantity=qty,unit=m.unit,condition=(x.get('condition') or 'BOM'),notes=(x.get('notes') or '').strip()))
        if d.get('send_now',True):doc.status='AGUARDANDO_ACEITE';doc.sent_at=datetime.utcnow()
        created.append(doc.document_code)
    db.session.add(AuditEvent(user_id=session['user_id'],event_type='MATERIAL_DELIVERY_BATCH',entity_type='collaborator_document',entity_id=','.join(created),detail=f'{len(created)} termo(s) criado(s)'));db.session.commit();return jsonify({'ok':True,'created':created,'count':len(created)})

@app.get('/api/materials/pending-count')
@login_required
def materials_pending_count_api():
    # V69.3.2: COUNT somente nas colunas estáveis; evita SELECT do modelo inteiro.
    n=(db.session.query(func.count(CollaboratorDocument.id))
       .filter(CollaboratorDocument.user_id==session['user_id'],CollaboratorDocument.status=='AGUARDANDO_ACEITE')
       .scalar() or 0)
    return jsonify({'ok':True,'count':int(n)})

@app.put('/api/materials/documents/<int:did>/items')
@login_required
def materials_document_items_update_api(did):
    _materials_require('materials.delivery.manage'); d=db.session.get(CollaboratorDocument,did)
    if not d:return jsonify({'ok':False,'error':'Documento não encontrado.'}),404
    if d.status=='ASSINADO':return jsonify({'ok':False,'error':'Documento assinado é imutável.'}),409
    payload=request.get_json(silent=True) or {}; items=payload.get('items') or []
    if 'notes' in payload: d.notes=(payload.get('notes') or '').strip()
    if payload.get('delivery_date'):
        try: d.delivery_date=datetime.strptime(payload.get('delivery_date'), '%Y-%m-%d').date()
        except: pass
    CollaboratorDocumentItem.query.filter_by(document_id=d.id).delete()
    for x in items:
        try: mid=int(x.get('material_id'))
        except: continue
        m=db.session.get(MaterialCatalogItem,mid)
        if not m: continue
        try: qty=_material_quantity(m,x.get('quantity') or 0)
        except ValueError as exc: return jsonify({'ok':False,'error':str(exc)}),400
        db.session.add(CollaboratorDocumentItem(document_id=d.id,material_id=m.id,description=m.description,brand=m.brand,model=m.model,quantity=qty,unit=m.unit,condition=(x.get('condition') or 'BOM'),notes=(x.get('notes') or '').strip()))
    d.status='RASCUNHO'; d.correction_note=None; db.session.add(AuditEvent(user_id=session['user_id'],event_type='MATERIAL_DOCUMENT_CORRECTED',entity_type='collaborator_document',entity_id=str(d.id),detail=d.document_code));db.session.commit();return jsonify({'ok':True})

@app.delete('/api/materials/documents/<int:did>')
@login_required
def materials_document_delete_api(did):
    _materials_require('materials.delivery.manage')
    d=db.session.get(CollaboratorDocument,did)
    if not d:return jsonify({'ok':False,'error':'Documento não encontrado.'}),404
    payload=request.get_json(silent=True) or {}; reason=(payload.get('reason') or '').strip()
    if len(reason)<3:return jsonify({'ok':False,'error':'Informe o motivo da exclusão.'}),400
    code=d.document_code or str(d.id); title=d.title or d.document_type
    # V69.5: exclusão administrativa de documento de teste/indevido também desfaz
    # os movimentos vinculados ao próprio documento, mantendo AuditEvent permanente.
    MaterialMovement.query.filter_by(document_id=d.id).delete(synchronize_session=False)
    CollaboratorDocumentItem.query.filter_by(document_id=d.id).delete(synchronize_session=False)
    db.session.add(AuditEvent(user_id=session['user_id'],event_type='MATERIAL_DOCUMENT_DELETED',entity_type='collaborator_document',entity_id=str(d.id),detail=f'{code} · {title} · Motivo: {reason[:500]}'))
    db.session.delete(d); db.session.commit()
    return jsonify({'ok':True,'code':code})

@app.get('/api/materials/documents')
@login_required
def materials_documents_api():
    own=request.args.get('own')=='1'; q=CollaboratorDocument.query
    if own or not _has_access('materials.dossier.view'):q=q.filter_by(user_id=session['user_id'])
    rows=q.order_by(CollaboratorDocument.created_at.desc()).limit(500).all();return jsonify({'ok':True,'rows':[_doc_payload(x) for x in rows]})

@app.get('/api/materials/documents/<int:did>')
@login_required
def materials_document_detail_api(did):
    d=db.session.get(CollaboratorDocument,did)
    if not d:return jsonify({'ok':False,'error':'Documento não encontrado.'}),404
    if d.user_id!=session['user_id'] and not _has_access('materials.dossier.view'):abort(403)
    return jsonify({'ok':True,'document':_doc_payload(d)})

@app.post('/api/materials/documents/<int:did>/send')
@login_required
def materials_document_send_api(did):
    _materials_require('materials.delivery.manage');d=db.session.get(CollaboratorDocument,did)
    if not d:return jsonify({'ok':False,'error':'Documento não encontrado.'}),404
    if d.status=='ASSINADO':return jsonify({'ok':False,'error':'Documento assinado é imutável.'}),409
    d.status='AGUARDANDO_ACEITE';d.sent_at=datetime.utcnow();d.correction_note=None;db.session.commit();return jsonify({'ok':True})

@app.post('/api/materials/documents/<int:did>/correction')
@login_required
def materials_document_correction_api(did):
    d=db.session.get(CollaboratorDocument,did)
    if not d or d.user_id!=session['user_id']:abort(404)
    if d.status!='AGUARDANDO_ACEITE':return jsonify({'ok':False,'error':'Documento não está aguardando aceite.'}),409
    note=((request.get_json(silent=True) or {}).get('note') or '').strip()
    if not note:return jsonify({'ok':False,'error':'Informe o que precisa ser corrigido.'}),400
    d.status='CORRECAO_SOLICITADA';d.correction_note=note;db.session.add(AuditEvent(user_id=session['user_id'],event_type='MATERIAL_DOCUMENT_CORRECTION',entity_type='collaborator_document',entity_id=str(d.id),detail=note[:500]));db.session.commit();return jsonify({'ok':True})

@app.post('/api/materials/documents/<int:did>/accept')
@login_required
def materials_document_accept_api(did):
    d=db.session.get(CollaboratorDocument,did)
    if not d or d.user_id!=session['user_id']:abort(404)
    if d.status!='AGUARDANDO_ACEITE':return jsonify({'ok':False,'error':'Documento não está aguardando aceite.'}),409
    sig=((request.get_json(silent=True) or {}).get('signature_data') or '')
    if not sig.startswith('data:image/'):return jsonify({'ok':False,'error':'Assinatura obrigatória.'}),400
    try:
        head,b64=sig.split(',',1);raw=base64.b64decode(b64);key=f'dossie/{d.user_id}/{d.document_code}_assinatura.png';
        if _r2_available():_r2_put_bytes(key,raw,'image/png');d.signature_file='r2__'+key
        else:
            name=f'{d.document_code}_assinatura.png';(UPLOAD_DIR/name).write_bytes(raw);d.signature_file=name
        d.status='ASSINADO';d.signed_at=datetime.utcnow();db.session.flush()
        pdfraw=_material_pdf_bytes(d);pkey=f'dossie/{d.user_id}/{d.document_code}.pdf'
        if _r2_available():_r2_put_bytes(pkey,pdfraw,'application/pdf');d.pdf_file='r2__'+pkey
        else:
            name=f'{d.document_code}.pdf';(UPLOAD_DIR/name).write_bytes(pdfraw);d.pdf_file=name
        for x in CollaboratorDocumentItem.query.filter_by(document_id=d.id).all():
            if x.material_id:db.session.add(MaterialMovement(user_id=d.user_id,material_id=x.material_id,document_id=d.id,movement_type='ENTREGA',quantity=x.quantity,condition=x.condition,notes=x.notes,created_by=session['user_id']))
        db.session.add(AuditEvent(user_id=session['user_id'],event_type='MATERIAL_DOCUMENT_SIGNED',entity_type='collaborator_document',entity_id=str(d.id),detail=d.document_code));db.session.commit();return jsonify({'ok':True,'code':d.document_code})
    except Exception as exc:
        db.session.rollback();app.logger.exception('Falha aceite dossiê');return jsonify({'ok':False,'error':str(exc)}),500

@app.get('/api/materials/documents/<int:did>/pdf')
@login_required
def materials_document_pdf_api(did):
    d=db.session.get(CollaboratorDocument,did)
    if not d or not d.pdf_file:abort(404)
    if d.user_id!=session['user_id'] and not _has_access('materials.dossier.view'):abort(403)
    raw=_r2_get_bytes(d.pdf_file[4:]) if d.pdf_file.startswith('r2__') else (UPLOAD_DIR/d.pdf_file).read_bytes();return send_file(io.BytesIO(raw),mimetype='application/pdf',download_name=f'{d.document_code}.pdf',as_attachment=False)

@app.get('/api/materials/my-load')
@login_required
def materials_my_load_api():
    uid=session['user_id']; rows=db.session.query(MaterialMovement.material_id,func.sum(case((MaterialMovement.movement_type.in_(['ENTREGA','SUBSTITUICAO_ENTRADA','TRANSFERENCIA_ENTRADA']),MaterialMovement.quantity),else_=-MaterialMovement.quantity))).filter(MaterialMovement.user_id==uid).group_by(MaterialMovement.material_id).all(); mids=[x[0] for x in rows]; mm={m.id:m for m in MaterialCatalogItem.query.filter(MaterialCatalogItem.id.in_(mids)).all()} if mids else {};return jsonify({'ok':True,'rows':[{'material_id':mid,'description':mm[mid].description if mid in mm else 'Item','quantity':float(qty or 0),'unit':mm[mid].unit if mid in mm else 'UN','control_type':mm[mid].control_type if mid in mm else 'DEVOLVIVEL','quantity_mode':mm[mid].quantity_mode if mid in mm else 'INTEIRO'} for mid,qty in rows if float(qty or 0)>0]})

@app.get('/api/materials/summary')
@login_required
def materials_summary_api():
    _materials_require('materials.dossier.view'); total=CollaboratorDocument.query.count();pending=CollaboratorDocument.query.filter_by(status='AGUARDANDO_ACEITE').count();signed=CollaboratorDocument.query.filter_by(status='ASSINADO').count();corr=CollaboratorDocument.query.filter_by(status='CORRECAO_SOLICITADA').count();req=MaterialRequest.query.filter(MaterialRequest.status.in_(['SOLICITADO','EM_ANALISE','APROVADO','EM_COMPRA','DISPONIVEL'])).count();return jsonify({'ok':True,'total':total,'pending':pending,'signed':signed,'corrections':corr,'requests_open':req})

@app.get('/api/materials/requests')
@login_required
def materials_requests_api():
    q=MaterialRequest.query
    if not _has_access('materials.delivery.manage'):q=q.filter_by(user_id=session['user_id'])
    rows=q.order_by(MaterialRequest.created_at.desc()).limit(300).all(); mids={x.material_id for x in rows};mm={m.id:m for m in MaterialCatalogItem.query.filter(MaterialCatalogItem.id.in_(mids)).all()} if mids else {};uu={u.id:u for u in User.query.filter(User.id.in_({x.user_id for x in rows})).all()} if rows else {};return jsonify({'ok':True,'rows':[{'id':x.id,'code':x.request_code or f'SM-{x.id:06d}','user_id':x.user_id,'user_name':uu[x.user_id].name if x.user_id in uu else '—','job_title':uu[x.user_id].job_title if x.user_id in uu else '', 'company':uu[x.user_id].company if x.user_id in uu else '', 'material_id':x.material_id,'material':mm[x.material_id].description if x.material_id in mm else '—','category':mm[x.material_id].category if x.material_id in mm else '', 'quantity':x.quantity,'urgency':x.urgency,'reason':x.reason or '', 'notes':x.notes or '', 'status':x.status,'created_at':x.created_at.isoformat()} for x in rows]})

@app.post('/api/materials/requests')
@login_required
def materials_request_create_api():
    _materials_require('materials.request');d=request.get_json(silent=True) or {}
    try: mid=int(d.get('material_id'))
    except: return jsonify({'ok':False,'error':'Item/quantidade inválidos.'}),400
    m=db.session.get(MaterialCatalogItem,mid)
    if not m: return jsonify({'ok':False,'error':'Item/quantidade inválidos.'}),400
    try: qty=_material_quantity(m,d.get('quantity') or 1)
    except ValueError as exc: return jsonify({'ok':False,'error':str(exc)}),400
    r=MaterialRequest(user_id=session['user_id'],material_id=mid,quantity=qty,reason=(d.get('reason') or '').strip(),urgency=(d.get('urgency') or 'NORMAL').upper(),notes=(d.get('notes') or '').strip());db.session.add(r);db.session.flush();r.request_code=f'SM-{r.id:06d}';db.session.commit();return jsonify({'ok':True,'code':r.request_code})

@app.get('/api/materials/requests/export.xlsx')
@login_required
def materials_requests_export_api():
    if not (_has_access('materials.delivery.manage') or _has_access('materials.dossier.view')): abort(403)
    q=MaterialRequest.query.order_by(MaterialRequest.created_at.desc())
    rows=q.all(); mids={x.material_id for x in rows}; uids={x.user_id for x in rows}
    mm={m.id:m for m in MaterialCatalogItem.query.filter(MaterialCatalogItem.id.in_(mids)).all()} if mids else {}
    uu={u.id:u for u in User.query.filter(User.id.in_(uids)).all()} if uids else {}
    # Os filtros são repetidos no servidor para que o Excel corresponda à visão administrativa.
    name=(request.args.get('name') or '').strip().lower(); status=(request.args.get('status') or '').strip().upper(); urgency=(request.args.get('urgency') or '').strip().upper(); category=(request.args.get('category') or '').strip().upper(); qtext=(request.args.get('q') or '').strip().lower(); dt_from=(request.args.get('from') or '').strip(); dt_to=(request.args.get('to') or '').strip()
    def keep(x):
        u=uu.get(x.user_id); m=mm.get(x.material_id); created=x.created_at.date().isoformat() if x.created_at else ''
        if name and name not in ((u.name if u else '') or '').lower(): return False
        if status and x.status!=status: return False
        if urgency and x.urgency!=urgency: return False
        if category and (m.category if m else '')!=category: return False
        if dt_from and created<dt_from: return False
        if dt_to and created>dt_to: return False
        hay=' '.join([x.request_code or '',u.name if u else '',u.company if u else '',u.job_title if u else '',m.description if m else '',m.category if m else '',x.reason or '',x.notes or '']).lower()
        return not qtext or qtext in hay
    rows=[x for x in rows if keep(x)]
    wb=Workbook(); ws=wb.active; ws.title='Solicitações'
    headers=['Solicitação','Data','Colaborador','Cargo','Empresa','Categoria','Material','Quantidade','Urgência','Status','Motivo','Observação']
    ws.append(headers)
    for c in ws[1]: c.font=Font(bold=True); c.fill=PatternFill('solid',fgColor='D9EAF7'); c.alignment=Alignment(horizontal='center')
    for x in rows:
        u=uu.get(x.user_id); m=mm.get(x.material_id)
        ws.append([x.request_code or f'SM-{x.id:06d}',x.created_at.strftime('%d/%m/%Y %H:%M') if x.created_at else '',u.name if u else '',u.job_title if u else '',u.company if u else '',m.category if m else '',m.description if m else '',x.quantity,x.urgency,x.status,x.reason or '',x.notes or ''])
    for i,w in enumerate([16,18,28,22,22,18,42,12,14,18,38,38],1): ws.column_dimensions[get_column_letter(i)].width=w
    bio=io.BytesIO(); wb.save(bio); bio.seek(0)
    return send_file(bio,mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',as_attachment=True,download_name='solicitacoes_materiais.xlsx')

@app.get('/api/materials/load/<int:uid>')
@login_required
def materials_user_load_api(uid):
    _materials_require('materials.delivery.manage')
    rows=db.session.query(MaterialMovement.material_id,func.sum(case((MaterialMovement.movement_type.in_(['ENTREGA','SUBSTITUICAO_ENTRADA','TRANSFERENCIA_ENTRADA']),MaterialMovement.quantity),else_=-MaterialMovement.quantity))).filter(MaterialMovement.user_id==uid).group_by(MaterialMovement.material_id).all()
    mids=[x[0] for x in rows]; mm={m.id:m for m in MaterialCatalogItem.query.filter(MaterialCatalogItem.id.in_(mids)).all()} if mids else {}
    return jsonify({'ok':True,'rows':[{'material_id':mid,'description':mm[mid].description if mid in mm else 'Item','quantity':float(qty or 0),'unit':mm[mid].unit if mid in mm else 'UN','control_type':mm[mid].control_type if mid in mm else 'DEVOLVIVEL','quantity_mode':mm[mid].quantity_mode if mid in mm else 'INTEIRO'} for mid,qty in rows if float(qty or 0)>0]})

@app.post('/api/materials/returns')
@login_required
def materials_return_create_api():
    _materials_require('materials.delivery.manage'); data=request.get_json(silent=True) or {}
    try: uid=int(data.get('user_id'))
    except: return jsonify({'ok':False,'error':'Selecione o colaborador.'}),400
    u=db.session.get(User,uid)
    if not u:return jsonify({'ok':False,'error':'Colaborador não encontrado.'}),404
    items=data.get('items') or []
    if not items:return jsonify({'ok':False,'error':'Informe ao menos um item devolvido.'}),400
    # saldo atual para impedir carga negativa; item avulso é permitido e não baixa saldo inexistente.
    balances=dict(db.session.query(MaterialMovement.material_id,func.sum(case((MaterialMovement.movement_type.in_(['ENTREGA','SUBSTITUICAO_ENTRADA','TRANSFERENCIA_ENTRADA']),MaterialMovement.quantity),else_=-MaterialMovement.quantity))).filter(MaterialMovement.user_id==uid).group_by(MaterialMovement.material_id).all())
    d=CollaboratorDocument(document_type='TERMO_DEVOLUCAO',user_id=uid,status='DEVOLVIDO',delivery_date=date.today(),title='Termo de Devolução de Materiais / Ferramentas',notes=(data.get('notes') or '').strip(),created_by=session['user_id'],signed_at=datetime.utcnow())
    db.session.add(d);db.session.flush();d.document_code=f'TDV-{d.id:06d}'
    for it in items:
        try: mid=int(it.get('material_id'))
        except: continue
        m=db.session.get(MaterialCatalogItem,mid)
        if not m: continue
        try: qty=_material_quantity(m,it.get('quantity') or 0)
        except ValueError as exc: return jsonify({'ok':False,'error':str(exc)}),400
        current=float(balances.get(mid) or 0); avulsa=bool(it.get('avulsa')) or current<=0
        if not avulsa and qty>current+1e-9: return jsonify({'ok':False,'error':f'{m.description}: devolução maior que a carga atual ({current:g}).'}),400
        note=(it.get('notes') or '').strip(); note=('DEVOLUÇÃO AVULSA — item não localizado na carga atual. '+note).strip() if avulsa else note
        db.session.add(CollaboratorDocumentItem(document_id=d.id,material_id=mid,description=m.description,brand=m.brand,model=m.model,quantity=qty,unit=m.unit,condition=(it.get('condition') or 'BOM').upper(),notes=note))
        if not avulsa: db.session.add(MaterialMovement(user_id=uid,material_id=mid,document_id=d.id,movement_type='DEVOLUCAO',quantity=qty,condition=(it.get('condition') or 'BOM').upper(),notes=note,created_by=session['user_id']))
        else: db.session.add(MaterialMovement(user_id=uid,material_id=mid,document_id=d.id,movement_type='DEVOLUCAO_AVULSA',quantity=0,condition=(it.get('condition') or 'BOM').upper(),notes=note,created_by=session['user_id']))
    db.session.flush()
    pdfraw=_material_pdf_bytes(d); pkey=f'dossie/{d.user_id}/{d.document_code}.pdf'
    if _r2_available(): _r2_put_bytes(pkey,pdfraw,'application/pdf'); d.pdf_file='r2__'+pkey
    else:
        name=f'{d.document_code}.pdf'; (UPLOAD_DIR/name).write_bytes(pdfraw); d.pdf_file=name
    db.session.add(AuditEvent(user_id=session['user_id'],event_type='MATERIAL_RETURN',entity_type='collaborator_document',entity_id=str(d.id),detail=d.document_code));db.session.commit()
    return jsonify({'ok':True,'code':d.document_code,'document_id':d.id})

@app.post('/api/materials/returns/request')
@login_required
def materials_return_request_api():
    """Colaborador assina a intenção de devolução. A carga só baixa após conferência e assinatura do responsável."""
    _materials_require('materials.my_documents')
    data=request.get_json(silent=True) or {}; uid=session['user_id']; items=data.get('items') or []; sig=(data.get('signature_data') or '')
    if not items:return jsonify({'ok':False,'error':'Selecione ao menos um item para devolução.'}),400
    if not sig.startswith('data:image/'):return jsonify({'ok':False,'error':'Assinatura do colaborador é obrigatória para solicitar a devolução.'}),400
    try:
        balances=dict(db.session.query(MaterialMovement.material_id,func.sum(case((MaterialMovement.movement_type.in_(['ENTREGA','SUBSTITUICAO_ENTRADA','TRANSFERENCIA_ENTRADA']),MaterialMovement.quantity),else_=-MaterialMovement.quantity))).filter(MaterialMovement.user_id==uid).group_by(MaterialMovement.material_id).all())
        d=CollaboratorDocument(document_type='TERMO_DEVOLUCAO',user_id=uid,status='AGUARDANDO_RECEBIMENTO',delivery_date=date.today(),title='Solicitação de Devolução de Materiais / Ferramentas',notes=(data.get('notes') or '').strip(),created_by=uid,signed_at=datetime.utcnow())
        db.session.add(d);db.session.flush();d.document_code=f'TDV-{d.id:06d}'
        valid=0
        for it in items:
            try: mid=int(it.get('material_id'))
            except: continue
            m=db.session.get(MaterialCatalogItem,mid)
            if not m or m.control_type=='CONSUMIVEL': continue
            qty=_material_quantity(m,it.get('quantity') or 0); current=float(balances.get(mid) or 0)
            if qty>current+1e-9: raise ValueError(f'{m.description}: devolução maior que a carga atual ({_material_qty_text(current)}).')
            db.session.add(CollaboratorDocumentItem(document_id=d.id,material_id=mid,description=m.description,brand=m.brand,model=m.model,quantity=qty,unit=m.unit,condition='A_CONFERIR',notes=(it.get('notes') or '').strip()));valid+=1
        if not valid: raise ValueError('Nenhum item devolvível válido foi informado.')
        head,b64=sig.split(',',1); raw=base64.b64decode(b64); key=f'dossie/{uid}/{d.document_code}_assinatura_colaborador.png'
        if _r2_available(): _r2_put_bytes(key,raw,'image/png'); d.signature_file='r2__'+key
        else:
            name=f'{d.document_code}_assinatura_colaborador.png'; (UPLOAD_DIR/name).write_bytes(raw); d.signature_file=name
        db.session.add(AuditEvent(user_id=uid,event_type='MATERIAL_RETURN_REQUEST_SIGNED',entity_type='collaborator_document',entity_id=str(d.id),detail=d.document_code));db.session.commit()
        return jsonify({'ok':True,'code':d.document_code,'document_id':d.id,'status':'AGUARDANDO_RECEBIMENTO'})
    except ValueError as exc:
        db.session.rollback(); return jsonify({'ok':False,'error':str(exc)}),400
    except Exception as exc:
        db.session.rollback(); app.logger.exception('Falha ao criar solicitação de devolução'); return jsonify({'ok':False,'error':'Não foi possível registrar a devolução. Consulte o log do sistema.','detail':str(exc)}),500

@app.get('/api/materials/returns/pending')
@login_required
def materials_returns_pending_api():
    _materials_require('materials.delivery.manage')
    rows=CollaboratorDocument.query.filter_by(document_type='TERMO_DEVOLUCAO',status='AGUARDANDO_RECEBIMENTO').order_by(CollaboratorDocument.created_at.asc()).all()
    return jsonify({'ok':True,'rows':[_doc_payload(x) for x in rows]})

@app.post('/api/materials/returns/<int:did>/confirm')
@login_required
def materials_return_confirm_api(did):
    _materials_require('materials.delivery.manage'); d=db.session.get(CollaboratorDocument,did)
    if not d or d.document_type!='TERMO_DEVOLUCAO' or d.status!='AGUARDANDO_RECEBIMENTO': return jsonify({'ok':False,'error':'Devolução pendente não encontrada.'}),404
    payload=request.get_json(silent=True) or {}; condition=(payload.get('condition') or 'BOM').upper(); receiver_note=(payload.get('notes') or '').strip(); sig=(payload.get('signature_data') or '')
    if not sig.startswith('data:image/'): return jsonify({'ok':False,'error':'Assinatura do responsável pelo recebimento é obrigatória.'}),400
    try:
        balances=dict(db.session.query(MaterialMovement.material_id,func.sum(case((MaterialMovement.movement_type.in_(['ENTREGA','SUBSTITUICAO_ENTRADA','TRANSFERENCIA_ENTRADA']),MaterialMovement.quantity),else_=-MaterialMovement.quantity))).filter(MaterialMovement.user_id==d.user_id).group_by(MaterialMovement.material_id).all())
        its=CollaboratorDocumentItem.query.filter_by(document_id=d.id).all()
        for x in its:
            m=db.session.get(MaterialCatalogItem,x.material_id) if x.material_id else None
            if not m: continue
            qty=_material_quantity(m,x.quantity); current=float(balances.get(m.id) or 0)
            if qty>current+1e-9: raise ValueError(f'{m.description}: carga atual insuficiente para confirmar a devolução.')
            x.condition=condition
            db.session.add(MaterialMovement(user_id=d.user_id,material_id=m.id,document_id=d.id,movement_type='DEVOLUCAO',quantity=qty,condition=condition,notes=receiver_note or x.notes,created_by=session['user_id']))
        head,b64=sig.split(',',1); raw=base64.b64decode(b64); key=f'dossie/{d.user_id}/{d.document_code}_assinatura_recebedor.png'
        if _r2_available(): _r2_put_bytes(key,raw,'image/png'); d.return_receiver_signature_file='r2__'+key
        else:
            name=f'{d.document_code}_assinatura_recebedor.png'; (UPLOAD_DIR/name).write_bytes(raw); d.return_receiver_signature_file=name
        d.return_receiver_id=session['user_id']; d.return_received_at=datetime.utcnow(); d.status='DEVOLVIDO'; d.notes=(' | '.join(x for x in [d.notes,receiver_note] if x)).strip()
        db.session.flush();pdfraw=_material_pdf_bytes(d);pkey=f'dossie/{d.user_id}/{d.document_code}.pdf'
        if _r2_available():_r2_put_bytes(pkey,pdfraw,'application/pdf');d.pdf_file='r2__'+pkey
        else:
            name=f'{d.document_code}.pdf';(UPLOAD_DIR/name).write_bytes(pdfraw);d.pdf_file=name
        db.session.add(AuditEvent(user_id=session['user_id'],event_type='MATERIAL_RETURN_CONFIRMED_SIGNED',entity_type='collaborator_document',entity_id=str(d.id),detail=d.document_code));db.session.commit()
        return jsonify({'ok':True,'code':d.document_code})
    except ValueError as exc:
        db.session.rollback(); return jsonify({'ok':False,'error':str(exc)}),409
    except Exception as exc:
        db.session.rollback(); app.logger.exception('Falha ao confirmar devolução'); return jsonify({'ok':False,'error':'Não foi possível confirmar a devolução.','detail':str(exc)}),500

@app.post('/api/materials/requests/<int:rid>/status')
@login_required
def materials_request_status_api(rid):
    _materials_require('materials.delivery.manage');r=db.session.get(MaterialRequest,rid)
    if not r:return jsonify({'ok':False,'error':'Solicitação não encontrada.'}),404
    status=((request.get_json(silent=True) or {}).get('status') or '').upper();allowed={'SOLICITADO','EM_ANALISE','APROVADO','EM_COMPRA','DISPONIVEL','ENTREGUE','RECEBIDO','REJEITADO'}
    if status not in allowed:return jsonify({'ok':False,'error':'Status inválido.'}),400
    r.status=status;db.session.commit();return jsonify({'ok':True})


# V69.2.1 HOTFIX2 — colunas aditivas do documento fiscal do Portal do Cliente.
try:
    with db.engine.begin() as conn:
        insp=db.inspect(db.engine); cols={c['name'] for c in insp.get_columns('customer_appointments')}
        for col,typ in [('invoice_number','VARCHAR(120)'),('invoice_file','VARCHAR(600)'),('invoice_original_name','VARCHAR(255)')]:
            if col not in cols: conn.execute(text(f'ALTER TABLE customer_appointments ADD COLUMN {col} {typ}'))
except Exception as exc:
    app.logger.warning('HOTFIX2: não foi possível validar colunas fiscais do Portal: %s',exc)

# ---------------- V69 Portal do Cliente ----------------
def _customer_company_ids(u):
    try: return [int(x) for x in json.loads(getattr(u,'customer_company_ids',None) or '[]') if str(x).isdigit()]
    except Exception: return []

def _customer_companies_for_user(u):
    ids=_customer_company_ids(u)
    if ids: return CustomerCompany.query.filter(CustomerCompany.id.in_(ids),CustomerCompany.active.is_(True)).order_by(CustomerCompany.trade_name,CustomerCompany.legal_name).all()
    if (u.company or '').strip(): return CustomerCompany.query.filter(CustomerCompany.active.is_(True),func.lower(CustomerCompany.legal_name)==u.company.strip().lower()).all()
    return []

def _portal_internal():
    return session.get('role') != 'customer' and (_has_access('portal.receive') or _has_access('portal.manage'))

def _portal_can_see(a):
    if _portal_internal(): return True
    u=db.session.get(User,session.get('user_id'))
    return bool(u and u.role=='customer' and (a.customer_company or '').strip().casefold() in {(c.legal_name or '').strip().casefold() for c in _customer_companies_for_user(u)})

def _portal_code(aid): return f"AG-{datetime.utcnow().year}-{aid:06d}"

def _portal_equipment_code(a,item): return f"{a.code or _portal_code(a.id)}-{item.item_no:02d}"

def _portal_store_dataurl(dataurl, prefix):
    if not dataurl or ',' not in dataurl: return None
    import base64
    head,raw=dataurl.split(',',1); blob=base64.b64decode(raw); ext='png' if 'png' in head else 'jpg'
    key=f"portal-cliente/{datetime.utcnow().strftime('%Y/%m')}/{prefix}-{uuid.uuid4().hex}.{ext}"
    try: _r2_put_bytes(key,blob,'image/png' if ext=='png' else 'image/jpeg'); return 'r2__'+key
    except Exception:
        name=key.replace('/','_'); (UPLOAD_DIR/name).write_bytes(blob); return name

def _portal_stored_bytes(stored):
    if not stored: return b''
    if stored.startswith('r2__'): return _r2_get_bytes(stored[4:])
    return (UPLOAD_DIR/stored).read_bytes()

def _portal_store_upload(f,prefix,allowed_images_only=False):
    if not f or not f.filename: return None,None
    mime=(f.mimetype or '').lower()
    if allowed_images_only and not mime.startswith('image/'): raise ValueError('A evidência deve ser uma imagem.')
    if not allowed_images_only and not (mime=='application/pdf' or mime.startswith('image/')): raise ValueError('Documento deve ser PDF, JPG ou PNG.')
    raw=f.read(); maxb=5*1024*1024 if allowed_images_only else 12*1024*1024
    if len(raw)>maxb: raise ValueError('Arquivo excede o tamanho permitido.')
    safe=secure_filename(f.filename) or ('foto.jpg' if allowed_images_only else 'documento')
    key=f"portal-cliente/{datetime.utcnow().strftime('%Y/%m')}/{prefix}-{uuid.uuid4().hex}-{safe}"
    try: _r2_put_bytes(key,raw,mime or 'application/octet-stream'); return 'r2__'+key,f.filename
    except Exception:
        name=key.replace('/','_'); (UPLOAD_DIR/name).write_bytes(raw); return name,f.filename

def _portal_pdf(a, item=None):
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image, PageBreak
    import io
    out=io.BytesIO(); doc=SimpleDocTemplate(out,pagesize=A4,rightMargin=30,leftMargin=30,topMargin=30,bottomMargin=30)
    st=getSampleStyleSheet(); title=ParagraphStyle('pt',parent=st['Title'],fontSize=18,leading=22,textColor=colors.HexColor('#123b68')); body=ParagraphStyle('pb',parent=st['BodyText'],fontSize=9,leading=12)
    story=[Paragraph('AGENDAMENTO DE ENTRADA DE EQUIPAMENTO',title),Spacer(1,8)]
    story.append(Paragraph(f"<b>Agendamento:</b> {a.code} &nbsp;&nbsp; <b>Cliente:</b> {a.customer_company} &nbsp;&nbsp; <b>Responsável:</b> {a.responsible_name}",body))
    cc=CustomerCompany.query.filter(func.lower(CustomerCompany.legal_name)==(a.customer_company or '').lower()).first()
    if cc:
        story.append(Paragraph(f"<b>Razão social:</b> {cc.legal_name} &nbsp;&nbsp; <b>CNPJ:</b> {cc.cnpj or '—'} &nbsp;&nbsp; <b>Contato:</b> {cc.contact_name or '—'} &nbsp;&nbsp; <b>Telefone:</b> {cc.phone or cc.mobile or '—'}",body))
        story.append(Paragraph(f"<b>E-mail:</b> {cc.email or '—'} &nbsp;&nbsp; <b>Endereço:</b> {cc.address or '—'} {cc.city or ''}/{cc.state or ''}",body))
    story.append(Paragraph(f"<b>Data do agendamento:</b> {a.scheduled_date.strftime('%d/%m/%Y') if a.scheduled_date else '—'} &nbsp;&nbsp; <b>Aceite:</b> {a.accepted_at.strftime('%d/%m/%Y %H:%M') if a.accepted_at else '—'}",body))
    if getattr(a,'invoice_number',None) or getattr(a,'invoice_original_name',None): story.append(Paragraph(f"<b>Nota Fiscal / Documento:</b> {getattr(a,'invoice_number',None) or '—'} &nbsp;&nbsp; {getattr(a,'invoice_original_name',None) or ''}",body))
    story.append(Spacer(1,10))
    rows=[item] if item else CustomerAppointmentEquipment.query.filter_by(appointment_id=a.id).order_by(CustomerAppointmentEquipment.item_no).all()
    data=[["Protocolo","Série","Equipamento","Versão","EOD","Defeito","Observação"]]
    for x in rows: data.append([_portal_equipment_code(a,x),x.serial_number,x.equipment or '—',x.version or '—',x.eod or '—',Paragraph(x.defect,body),Paragraph(x.notes or '—',body)])
    t=Table(data,colWidths=[78,68,65,42,48,105,105],repeatRows=1); t.setStyle(TableStyle([('BACKGROUND',(0,0),(-1,0),colors.HexColor('#eaf2fb')),('TEXTCOLOR',(0,0),(-1,0),colors.HexColor('#123b68')),('FONTNAME',(0,0),(-1,0),'Helvetica-Bold'),('GRID',(0,0),(-1,-1),.4,colors.HexColor('#b9c7d6')),('VALIGN',(0,0),(-1,-1),'TOP'),('FONTSIZE',(0,0),(-1,-1),7),('PADDING',(0,0),(-1,-1),4)])); story += [t,Spacer(1,10)]
    if a.notes: story.append(Paragraph(f"<b>Observação geral:</b> {a.notes}",body))
    for x in rows:
        if x.photo_file:
            try:
                raw=_portal_stored_bytes(x.photo_file); im=Image(io.BytesIO(raw)); im._restrictSize(500,300); story += [Spacer(1,8),Paragraph(f"<b>Evidência do equipamento:</b> {_portal_equipment_code(a,x)} · Série {x.serial_number}",body),Spacer(1,4),im]
            except Exception: app.logger.exception('Falha ao incluir evidência no PDF do equipamento %s',x.id)
    story += [Spacer(1,8),Paragraph(f"Aceite registrado por <b>{a.accepted_name or a.responsible_name}</b>{' com assinatura eletrônica' if a.signature_file else ' sem assinatura eletrônica'}. Documento bloqueado para edição após o envio.",body)]
    doc.build(story); return out.getvalue()

def _portal_send_email(a):
    host=os.environ.get('SMTP_HOST','').strip(); to=os.environ.get('ASSISTENCIA_EMAIL','').strip(); sender=os.environ.get('SMTP_FROM','').strip() or os.environ.get('SMTP_USER','').strip()
    if not host or not to or not sender: return 'NAO_CONFIGURADO','Defina SMTP_HOST, SMTP_FROM/SMTP_USER e ASSISTENCIA_EMAIL.'
    msg=EmailMessage(); msg['Subject']=f"Novo agendamento {a.code} — {a.customer_company}"; msg['From']=sender; msg['To']=to
    count=CustomerAppointmentEquipment.query.filter_by(appointment_id=a.id).count(); msg.set_content(f"Novo agendamento de equipamentos.\
\
Agendamento: {a.code}\
Cliente: {a.customer_company}\
Responsável: {a.responsible_name}\
Equipamentos: {count}\
Data do agendamento: {a.scheduled_date.strftime('%d/%m/%Y') if a.scheduled_date else '—'}\
\
Acesse o Portal do Cliente para visualizar e baixar os PDFs individuais.")
    try:
        port=int(os.environ.get('SMTP_PORT','587')); user=os.environ.get('SMTP_USER',''); pwd=os.environ.get('SMTP_PASSWORD','');
        with smtplib.SMTP(host,port,timeout=15) as srv:
            if os.environ.get('SMTP_TLS','1')!='0': srv.starttls()
            if user: srv.login(user,pwd)
            srv.send_message(msg)
        return 'ENVIADO',f'Enviado para {to}'
    except Exception as exc:
        app.logger.exception('Falha ao enviar e-mail do agendamento %s',a.code); return 'FALHA',str(exc)[:450]


@app.get('/portal-cliente/cadastro-clientes/exportar.xlsx')
@login_required
def portal_customer_export_xlsx():
    if session.get('role')!='manager': abort(403)
    rows=CustomerCompany.query.order_by(CustomerCompany.active.desc(),CustomerCompany.legal_name).all()
    wb=Workbook(); ws=wb.active; ws.title='Clientes e Garagens'
    headers=['Razão Social','Nome Fantasia / Garagem','CNPJ','Inscrição Estadual','Contato','Cargo / Função','Telefone','Celular','E-mail','Endereço de Retirada','Cidade / Região','UF','CEP','Status','Cadastro','Pendências','Observações']
    ws.append(headers)
    for c in rows:
        missing=[]
        if not c.contact_name: missing.append('Contato')
        if not c.email: missing.append('E-mail')
        if not (c.phone or c.mobile): missing.append('Telefone/Celular')
        if not c.address: missing.append('Endereço')
        if not c.city: missing.append('Cidade/Região')
        ws.append([c.legal_name,c.trade_name or '',c.cnpj or '',c.state_registration or '',c.contact_name or '',c.contact_role or '',c.phone or '',c.mobile or '',c.email or '',c.address or '',c.city or '',c.state or '',c.zip_code or '','ATIVO' if c.active else 'INATIVO','PENDENTE' if missing else 'COMPLETO',', '.join(missing),c.notes or ''])
    for cell in ws[1]: cell.font=Font(bold=True,color='FFFFFF'); cell.fill=PatternFill('solid',fgColor='0B5FC7'); cell.alignment=Alignment(horizontal='center')
    widths=[28,28,20,18,22,18,22,18,32,42,22,8,14,12,13,35,40]
    for i,w in enumerate(widths,1): ws.column_dimensions[get_column_letter(i)].width=w
    ws.freeze_panes='A2'; ws.auto_filter.ref=ws.dimensions
    bio=io.BytesIO(); wb.save(bio); bio.seek(0)
    return send_file(bio,as_attachment=True,download_name=f'CLIENTES_GARAGENS_{datetime.now().strftime("%Y%m%d_%H%M")}.xlsx',mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')

@app.get('/portal-cliente/cadastro-clientes/modelo.xlsx')
@login_required
def portal_customer_model_xlsx():
    if session.get('role')!='manager': abort(403)
    wb=Workbook(); ws=wb.active; ws.title='Importar Clientes'
    headers=['Razão Social *','Nome Fantasia','CNPJ','Inscrição Estadual','Contato Principal','Cargo/Função','Telefone','Celular','E-mail','Endereço','Cidade','UF','CEP','Observações','Status']
    ws.append(headers); ws.append(['Empresa Exemplo Ltda','Empresa Exemplo','00.000.000/0001-00','','Contato','Gestor','','','','','São Paulo','SP','','','ATIVO'])
    for c in ws[1]: c.font=Font(bold=True); c.fill=PatternFill('solid',fgColor='D9EAF7')
    for i,w in enumerate([28,24,20,20,22,18,18,18,28,32,20,8,14,34,12],1): ws.column_dimensions[get_column_letter(i)].width=w
    ins=wb.create_sheet('Instruções'); ins.append(['IMPORTAÇÃO DE CLIENTES - V69.3']); ins.append(['Razão Social é obrigatória. CNPJ é usado para detectar duplicidades. Status: ATIVO ou INATIVO.']); ins.column_dimensions['A'].width=110
    bio=io.BytesIO(); wb.save(bio); bio.seek(0); return send_file(bio,as_attachment=True,download_name='MODELO_IMPORTACAO_CLIENTES_PORTAL.xlsx',mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')

@app.post('/portal-cliente/cadastro-clientes/importar')
@login_required
def portal_customer_import_xlsx():
    if session.get('role')!='manager': abort(403)
    f=request.files.get('file')
    if not f or not (f.filename or '').lower().endswith('.xlsx'): flash('Selecione uma planilha .xlsx válida.'); return redirect('/portal-cliente/cadastro-clientes')
    try: wb=load_workbook(f,read_only=True,data_only=True); ws=wb['Importar Clientes'] if 'Importar Clientes' in wb.sheetnames else wb.active
    except Exception: flash('Não foi possível ler a planilha.'); return redirect('/portal-cliente/cadastro-clientes')
    created=0; skipped=[]
    for n,row in enumerate(ws.iter_rows(min_row=2,values_only=True),2):
        vals=list(row)+[None]*15; legal=str(vals[0] or '').strip(); cnpj=str(vals[2] or '').strip()
        if not legal: continue
        norm=re.sub(r'\\D','',cnpj)
        dup=None
        if norm:
            for c in CustomerCompany.query.filter(CustomerCompany.cnpj.isnot(None)).all():
                if re.sub(r'\\D','',c.cnpj or '')==norm: dup=c; break
        if dup: skipped.append(f'linha {n}: CNPJ já cadastrado'); continue
        c=CustomerCompany(legal_name=legal,trade_name=str(vals[1] or '').strip() or None,cnpj=cnpj or None,state_registration=str(vals[3] or '').strip() or None,contact_name=str(vals[4] or '').strip() or None,contact_role=str(vals[5] or '').strip() or None,phone=str(vals[6] or '').strip() or None,mobile=str(vals[7] or '').strip() or None,email=str(vals[8] or '').strip() or None,address=str(vals[9] or '').strip() or None,city=str(vals[10] or '').strip() or None,state=str(vals[11] or '').strip().upper() or None,zip_code=str(vals[12] or '').strip() or None,notes=str(vals[13] or '').strip() or None,active=str(vals[14] or 'ATIVO').strip().upper()!='INATIVO')
        db.session.add(c); created+=1
    db.session.commit(); flash(f'Importação concluída: {created} cliente(s) incluído(s).'+(f' {len(skipped)} linha(s) ignorada(s) por duplicidade.' if skipped else '')); return redirect('/portal-cliente/cadastro-clientes')

@app.route('/portal-cliente/cadastro-clientes', methods=['GET','POST'])
@login_required
def portal_customer_companies():
    if session.get('role')!='manager': abort(403)
    if request.method=='POST':
        d=request.form; legal=(d.get('legal_name') or '').strip()
        if not legal: flash('Razão social é obrigatória.'); return redirect(request.path)
        c=CustomerCompany(legal_name=legal,trade_name=(d.get('trade_name') or '').strip() or None,cnpj=(d.get('cnpj') or '').strip() or None,state_registration=(d.get('state_registration') or '').strip() or None,contact_name=(d.get('contact_name') or '').strip() or None,contact_role=(d.get('contact_role') or '').strip() or None,phone=(d.get('phone') or '').strip() or None,mobile=(d.get('mobile') or '').strip() or None,email=(d.get('email') or '').strip() or None,address=(d.get('address') or '').strip() or None,city=(d.get('city') or '').strip() or None,state=(d.get('state') or '').strip().upper() or None,zip_code=(d.get('zip_code') or '').strip() or None,notes=(d.get('notes') or '').strip() or None)
        db.session.add(c); db.session.commit(); flash('Cliente cadastrado.'); return redirect(request.path)
    companies=CustomerCompany.query.order_by(CustomerCompany.active.desc(),CustomerCompany.legal_name).all(); customer_users=User.query.filter(User.role=='customer',User.archived_at.is_(None)).order_by(User.active.desc(),User.name).all(); pending_count=sum(1 for c in companies if (not c.contact_name) or (not c.email) or (not (c.phone or c.mobile)) or (not c.address) or (not c.city)); return render_template('customer_companies.html',companies=companies,customer_users=customer_users,pending_count=pending_count,customer_company_map={u.id:_customer_company_ids(u) for u in customer_users})

@app.post('/portal-cliente/cadastro-clientes/<int:cid>/editar')
@login_required
def portal_customer_company_edit(cid):
    if session.get('role')!='manager': abort(403)
    c=db.session.get(CustomerCompany,cid)
    if not c: abort(404)
    d=request.form; legal=(d.get('legal_name') or '').strip()
    if not legal: flash('Razão social é obrigatória.'); return redirect('/portal-cliente/cadastro-clientes')
    for attr in ('trade_name','cnpj','state_registration','contact_name','contact_role','phone','mobile','email','address','city','zip_code','notes'):
        setattr(c,attr,(d.get(attr) or '').strip() or None)
    c.legal_name=legal; c.state=(d.get('state') or '').strip().upper() or None; c.active=d.get('active')=='1'
    db.session.commit(); flash('Cadastro do cliente atualizado.'); return redirect('/portal-cliente/cadastro-clientes')

@app.post('/portal-cliente/cadastro-clientes/<int:cid>/excluir')
@login_required
def portal_customer_company_delete(cid):
    if session.get('role')!='manager': abort(403)
    c=db.session.get(CustomerCompany,cid)
    if not c: abort(404)
    linked_users=[u for u in User.query.all() if cid in _customer_company_ids(u)]
    linked_appts=CustomerAppointment.query.filter(func.lower(CustomerAppointment.customer_company)==(c.legal_name or '').lower()).count()
    if linked_users or linked_appts:
        c.active=False; db.session.commit(); flash('Cliente possui vínculos e foi inativado para preservar o histórico.')
    else:
        db.session.delete(c); db.session.commit(); flash('Cliente excluído.')
    return redirect('/portal-cliente/cadastro-clientes')


# V71.1 — Acessos externos do Portal ficam junto ao cadastro do cliente.
@app.post('/portal-cliente/cadastro-clientes/acessos/novo')
@login_required
def portal_customer_access_create():
    if session.get('role')!='manager': abort(403)
    d=request.form
    try: cid=int(d.get('company_id') or 0)
    except Exception: cid=0
    company=db.session.get(CustomerCompany,cid)
    if not company or not company.active:
        flash('Selecione uma empresa/garagem ativa.'); return redirect('/portal-cliente/cadastro-clientes#acessos')
    name=(d.get('name') or '').strip(); username=(d.get('username') or d.get('email') or '').strip().lower(); password=d.get('password') or ''
    email=_normalize_optional_email(d.get('email')); phone=_normalize_optional_phone(d.get('phone'))
    if not name or not username or len(password)<8:
        flash('Nome, login e senha com pelo menos 8 caracteres são obrigatórios.'); return redirect('/portal-cliente/cadastro-clientes#acessos')
    if User.query.filter(func.lower(User.username)==username).first():
        flash('Já existe um usuário com esse login.'); return redirect('/portal-cliente/cadastro-clientes#acessos')
    if email and User.query.filter(func.lower(User.email)==email).first():
        flash('Já existe um usuário com esse e-mail.'); return redirect('/portal-cliente/cadastro-clientes#acessos')
    if phone and User.query.filter(User.phone==phone).first():
        flash('Já existe um usuário com esse celular.'); return redirect('/portal-cliente/cadastro-clientes#acessos')
    u=User(name=name,username=username,password_hash=generate_password_hash(password),role='customer',active=True,user_code=_next_user_code('customer'),email=email,phone=phone,company=company.legal_name,personnel_status='ATIVO',access_json=json.dumps(sorted(_default_access_for_role('customer')),ensure_ascii=False),customer_company_ids=json.dumps([company.id]))
    db.session.add(u); db.session.add(AuditEvent(user_id=session.get('user_id'),event_type='PORTAL_CUSTOMER_ACCESS_CREATED',entity_type='user',entity_id=username,detail=f'{name} · {company.legal_name}')); db.session.commit()
    flash('Acesso do cliente criado com sucesso.'); return redirect('/portal-cliente/cadastro-clientes#acessos')

@app.post('/portal-cliente/cadastro-clientes/acessos/<int:uid>/editar')
@login_required
def portal_customer_access_edit(uid):
    if session.get('role')!='manager': abort(403)
    u=db.session.get(User,uid)
    if not u or u.role!='customer': abort(404)
    d=request.form
    try: cid=int(d.get('company_id') or 0)
    except Exception: cid=0
    company=db.session.get(CustomerCompany,cid)
    if not company: flash('Empresa inválida.'); return redirect('/portal-cliente/cadastro-clientes#acessos')
    name=(d.get('name') or '').strip(); username=(d.get('username') or '').strip().lower(); email=_normalize_optional_email(d.get('email')); phone=_normalize_optional_phone(d.get('phone')); password=d.get('password') or ''
    if not name or not username: flash('Nome e login são obrigatórios.'); return redirect('/portal-cliente/cadastro-clientes#acessos')
    if User.query.filter(User.id!=uid,func.lower(User.username)==username).first(): flash('Login já utilizado.'); return redirect('/portal-cliente/cadastro-clientes#acessos')
    if email and User.query.filter(User.id!=uid,func.lower(User.email)==email).first(): flash('E-mail já utilizado.'); return redirect('/portal-cliente/cadastro-clientes#acessos')
    if phone and User.query.filter(User.id!=uid,User.phone==phone).first(): flash('Celular já utilizado.'); return redirect('/portal-cliente/cadastro-clientes#acessos')
    if password and len(password)<8: flash('A nova senha deve ter pelo menos 8 caracteres.'); return redirect('/portal-cliente/cadastro-clientes#acessos')
    u.name=name;u.username=username;u.email=email;u.phone=phone;u.company=company.legal_name;u.customer_company_ids=json.dumps([company.id]);u.active=d.get('active')=='1';u.personnel_status='ATIVO' if u.active else 'INATIVO';u.access_json=json.dumps(sorted(_default_access_for_role('customer')),ensure_ascii=False)
    if password:u.password_hash=generate_password_hash(password)
    db.session.add(AuditEvent(user_id=session.get('user_id'),event_type='PORTAL_CUSTOMER_ACCESS_UPDATED',entity_type='user',entity_id=str(u.id),detail=f'{u.name} · {company.legal_name} · {"ATIVO" if u.active else "INATIVO"}'));db.session.commit()
    flash('Acesso do cliente atualizado.'); return redirect('/portal-cliente/cadastro-clientes#acessos')


# ---------------- V71 Logística / Leva e Traz ----------------
_V71_WEEKDAYS={"segunda":0,"terça":1,"terca":1,"quarta":2,"quinta":3,"sexta":4}
_V71_WEEKDAY_LABELS={0:"Segunda-feira",1:"Terça-feira",2:"Quarta-feira",3:"Quinta-feira",4:"Sexta-feira"}

def _v71_norm(value):
    txt=unicodedata.normalize('NFKD',str(value or '')).encode('ascii','ignore').decode('ascii').strip().casefold()
    return re.sub(r'\s+',' ',txt)

def _v71_route_for_company(company):
    key=_v71_norm(company)
    if not key: return None
    rows=LogisticsGarageRoute.query.filter_by(active=True).all()
    exact=next((r for r in rows if _v71_norm(r.garage_name)==key),None)
    if exact: return exact
    c=CustomerCompany.query.filter(CustomerCompany.active.is_(True)).all()
    match=next((x for x in c if _v71_norm(x.legal_name)==key or _v71_norm(x.trade_name)==key),None)
    if match:
        byid=next((r for r in rows if r.customer_company_id==match.id),None)
        if byid:return byid
        names={_v71_norm(match.legal_name),_v71_norm(match.trade_name)}-{''}
        return next((r for r in rows if _v71_norm(r.garage_name) in names),None)
    return None

def _v71_next_operational_date(company, start=None):
    route=_v71_route_for_company(company)
    if not route:return None,route
    cur=start or datetime.now(ZoneInfo('America/Sao_Paulo')).date()
    blocked={x.blocked_date for x in LogisticsBlockedDate.query.filter_by(active=True).all()}
    # Inclui hoje se hoje já é o dia regular: não há horário/cutoff na operação.
    for offset in range(0,45):
        d=cur+timedelta(days=offset)
        if d.weekday()==route.weekday and d not in blocked:return d,route
    return None,route

def _v71_link_customer(route):
    if route.customer_company_id:return
    key=_v71_norm(route.garage_name)
    for c in CustomerCompany.query.filter_by(active=True).all():
        if key in {_v71_norm(c.legal_name),_v71_norm(c.trade_name)}:
            route.customer_company_id=c.id;return

def _v711_sync_customers_from_workbook(wb):
    """Cria clientes/garagens ausentes e enriquece campos vazios usando Contato Garagem."""
    if 'Contato Garagem' not in wb.sheetnames: return {'created':0,'enriched':0,'matched':0}
    ws=wb['Contato Garagem']; created=enriched=matched=0
    companies=CustomerCompany.query.all()
    byname={_v71_norm(c.legal_name):c for c in companies}
    for c in companies:
        if c.trade_name: byname.setdefault(_v71_norm(c.trade_name),c)
    for row in ws.iter_rows(min_row=2,values_only=True):
        vals=list(row)+[None]*7
        garage=str(vals[0] or '').strip(); responsible=str(vals[1] or '').strip(); email=str(vals[2] or '').strip(); phone=str(vals[3] or '').strip(); address=str(vals[4] or '').strip(); region=str(vals[5] or '').strip()
        if not garage: continue
        key=_v71_norm(garage); c=byname.get(key)
        if c is None:
            c=CustomerCompany(legal_name=garage,trade_name=garage,contact_name=responsible or None,email=email or None,phone=phone or None,address=address or None,city=region or None,state='SP',notes='Importado da Matriz Leva e Traz V71.1',active=True)
            db.session.add(c);byname[key]=c;created+=1
        else:
            matched+=1; changed=False
            fills={'trade_name':garage,'contact_name':responsible,'email':email,'phone':phone,'address':address,'city':region,'state':'SP'}
            for attr,val in fills.items():
                if val and not (getattr(c,attr,None) or '').strip(): setattr(c,attr,val);changed=True
            if not c.active: c.active=True; changed=True
            if changed: enriched+=1
    db.session.flush()
    return {'created':created,'enriched':enriched,'matched':matched}

def _v71_import_resumo(file_storage):
    wb=load_workbook(file_storage,read_only=True,data_only=True)
    if 'Resumo' not in wb.sheetnames: raise ValueError('A planilha padrão precisa conter a aba Resumo.')
    customer_sync=_v711_sync_customers_from_workbook(wb)
    ws=wb['Resumo']; current_day=None; parsed=[]
    for row in ws.iter_rows(min_row=2,values_only=True):
        garage=str(row[0] or '').strip(); contact=str(row[1] or '').strip(); address=str(row[2] or '').strip(); region=str(row[3] or '').strip()
        if not garage: continue
        day_key=_v71_norm(garage)
        if day_key in _V71_WEEKDAYS:
            current_day=_V71_WEEKDAYS[day_key];continue
        if current_day is None: continue
        parsed.append({'garage':garage,'contact':contact,'address':address,'region':region,'weekday':current_day})
    if not parsed: raise ValueError('Nenhuma garagem válida foi encontrada na aba Resumo.')

    # Normaliza duplicidades na própria planilha antes de tocar no banco.
    # Repetições da mesma garagem no mesmo dia são consolidadas; dias conflitantes são bloqueados.
    dedup={}
    for item in parsed:
        key=_v71_norm(item['garage'])
        prev=dedup.get(key)
        if prev:
            if prev['weekday'] != item['weekday']:
                raise ValueError(f"A garagem '{item['garage']}' aparece em mais de um dia da semana na aba Resumo. Corrija a planilha antes de importar.")
            # Mantém o primeiro nome e aproveita dados preenchidos na repetição.
            for fld in ('contact','address','region'):
                if not prev.get(fld) and item.get(fld): prev[fld]=item[fld]
            continue
        dedup[key]=dict(item)
    parsed=list(dedup.values())

    # V71.1 — toda garagem da matriz também passa a existir no Cadastro de Clientes.
    # A aba Contato Garagem enriquece os dados; a aba Resumo garante que nenhuma rota fique sem cadastro.
    companies=CustomerCompany.query.all(); company_names={}
    for c in companies:
        company_names[_v71_norm(c.legal_name)]=c
        if c.trade_name: company_names.setdefault(_v71_norm(c.trade_name),c)
    for item in parsed:
        key=_v71_norm(item['garage']); c=company_names.get(key)
        if c is None:
            c=CustomerCompany(legal_name=item['garage'],trade_name=item['garage'],contact_name=item['contact'] or None,address=item['address'] or None,city=item['region'] or None,state='SP',notes='Importado da Matriz Leva e Traz V71.1',active=True)
            db.session.add(c);db.session.flush();company_names[key]=c;customer_sync['created']+=1
        else:
            changed_customer=False
            for attr,val in (('contact_name',item['contact']),('address',item['address']),('city',item['region'])):
                if val and not (getattr(c,attr,None) or '').strip(): setattr(c,attr,val);changed_customer=True
            if changed_customer: customer_sync['enriched']+=1
    db.session.flush()

    # A planilha é uma fotografia completa da matriz: ausentes são inativados, histórico é preservado.
    incoming=set(dedup); changed=created=unchanged=0
    existing=LogisticsGarageRoute.query.all()
    byname={_v71_norm(x.garage_name):x for x in existing}
    for item in parsed:
        key=_v71_norm(item['garage']); r=byname.get(key)
        is_new = r is None
        if is_new:
            r=LogisticsGarageRoute(garage_name=item['garage'],created_by=session.get('user_id'))
            db.session.add(r)
            # Atualiza o mapa imediatamente para impedir um segundo INSERT da mesma garagem
            # durante a mesma importação, antes do flush/commit.
            byname[key]=r
            created+=1
            before=None
        else:
            before=(r.weekday,r.contact_name or '',r.address or '',r.region or '',r.active)
        r.weekday=item['weekday'];r.contact_name=item['contact'] or None;r.address=item['address'] or None;r.region=item['region'] or None;r.active=True;r.source_import=secure_filename(file_storage.filename or 'Resumo.xlsx')
        # A consulta de vínculo não deve disparar autoflush enquanto a linha nova ainda está sendo montada.
        with db.session.no_autoflush:
            _v71_link_customer(r)
        if not is_new:
            after=(r.weekday,r.contact_name or '',r.address or '',r.region or '',r.active)
            if before==after: unchanged+=1
            else: changed+=1
    deactivated=0
    for r in existing:
        if _v71_norm(r.garage_name) not in incoming and r.active:
            r.active=False;deactivated+=1
    db.session.flush()
    # Recalcula somente solicitações futuras ainda em programação; recebidas/concluídas nunca mudam.
    impacted=0
    for a in CustomerAppointment.query.filter(CustomerAppointment.status.in_(['NA_PROGRAMACAO','ENVIADO','RASCUNHO'])).all():
        if a.status=='RASCUNHO':continue
        nxt,_=_v71_next_operational_date(a.customer_company,a.request_date or date.today())
        if nxt and a.expected_date!=nxt:
            a.expected_date=nxt
            if not a.programmed_date:a.scheduled_date=nxt
            impacted+=1
    db.session.add(AuditEvent(user_id=session.get('user_id'),event_type='V71_LOGISTICS_MATRIX_IMPORTED',entity_type='logistics_matrix',entity_id='Resumo',detail=f'{created} novas · {changed} alteradas · {deactivated} inativadas · {impacted} agendamentos recalculados'))
    db.session.commit()
    return {'created':created,'changed':changed,'unchanged':unchanged,'deactivated':deactivated,'impacted':impacted,'total':len(parsed),'customers_created':customer_sync['created'],'customers_enriched':customer_sync['enriched']}

@app.get('/portal-cliente/logistica/modelo.xlsx')
@login_required
def v71_logistics_model():
    if not _portal_internal():abort(403)
    path=DATA_DIR/'MODELO_LEVA_E_TRAZ_PADRAO.xlsx'
    if not path.exists():abort(404)
    return send_file(path,as_attachment=True,download_name='MODELO_LEVA_E_TRAZ_PADRAO.xlsx')

@app.post('/portal-cliente/logistica/importar')
@login_required
def v71_logistics_import():
    if not _portal_internal():abort(403)
    f=request.files.get('file')
    if not f or not (f.filename or '').lower().endswith('.xlsx'):
        flash('Selecione a planilha padrão .xlsx.');return redirect('/portal-cliente/gestao-agendamentos')
    try:
        result=_v71_import_resumo(f)
        flash(f"Matriz atualizada: {result['total']} garagens · {result['created']} novas · {result['changed']} alteradas · {result['deactivated']} inativadas · {result['impacted']} agendamentos futuros ajustados · {result.get('customers_created',0)} cadastros de clientes criados · {result.get('customers_enriched',0)} enriquecidos.")
    except Exception as exc:
        db.session.rollback();app.logger.exception('V71: falha na importação Leva e Traz');flash(f'Falha na importação: {exc}')
    return redirect('/portal-cliente/gestao-agendamentos')

@app.get('/api/portal/logistics/next-date')
@login_required
def v71_logistics_next_date():
    if not (_has_access('portal.appointments') or _portal_internal()):abort(403)
    company=(request.args.get('company') or '').strip();u=db.session.get(User,session['user_id'])
    if not _portal_internal():
        allowed={(c.legal_name or '').strip() for c in _customer_companies_for_user(u)}
        if company not in allowed: company=next(iter(allowed),'')
    nxt,route=_v71_next_operational_date(company)
    today=datetime.now(ZoneInfo('America/Sao_Paulo')).date()
    return jsonify({'ok':True,'request_date':today.isoformat(),'expected_date':nxt.isoformat() if nxt else '', 'weekday':_V71_WEEKDAY_LABELS.get(route.weekday,'') if route else '', 'garage':route.garage_name if route else '', 'has_route':bool(route)})

@app.post('/api/portal/logistics/blocked-dates')
@login_required
def v71_blocked_date_create():
    if not _portal_internal():abort(403)
    data=request.get_json(silent=True) or {}; raw=data.get('date');desc=(data.get('description') or '').strip()
    try:d=date.fromisoformat(raw)
    except Exception:return jsonify({'ok':False,'error':'Data inválida.'}),400
    row=LogisticsBlockedDate.query.filter_by(blocked_date=d).first() or LogisticsBlockedDate(blocked_date=d,created_by=session['user_id'])
    row.description=desc;row.active=True;db.session.add(row);db.session.commit();return jsonify({'ok':True})

@app.delete('/api/portal/logistics/blocked-dates/<int:bid>')
@login_required
def v71_blocked_date_delete(bid):
    if not _portal_internal():abort(403)
    row=db.session.get(LogisticsBlockedDate,bid)
    if not row:abort(404)
    row.active=False;db.session.commit();return jsonify({'ok':True})

@app.get('/portal-cliente/gestao-agendamentos')
@login_required
def v71_schedule_management_page():
    if not _portal_internal():abort(403)
    return render_template('logistics_schedule.html',app_release=APP_RELEASE)

@app.post('/api/portal/logistics/share-whatsapp')
@login_required
def portal_logistics_share_whatsapp():
    if not _portal_internal(): abort(403)
    data=request.get_json(silent=True) or {}; day=(data.get('date') or '').strip(); count=int(data.get('appointments') or 0)
    try:
        db.session.add(AuditEvent(user_id=session.get('user_id'),event_type='PORTAL_LOGISTICS_WHATSAPP_SHARE',entity_type='logistics_day',entity_id=day or 'sem-data',detail=('ATIVIDADE_AGENDADA' if count else 'NOVA_ATIVIDADE')+f' | {count} agendamento(s)'))
        db.session.commit()
    except Exception:
        db.session.rollback(); app.logger.exception('Falha ao auditar compartilhamento WhatsApp da logística')
    return jsonify({'ok':True,'activity_type':'ATIVIDADE_AGENDADA' if count else 'NOVA_ATIVIDADE'})

@app.get('/api/portal/logistics/management')
@login_required
def v71_schedule_management_api():
    if not _portal_internal():abort(403)
    today=datetime.now(ZoneInfo('America/Sao_Paulo')).date();month=(request.args.get('month') or today.strftime('%Y-%m'))
    try:y,m=[int(x) for x in month.split('-',1)];start=date(y,m,1);end=(date(y+1,1,1) if m==12 else date(y,m+1,1))
    except Exception:start=date(today.year,today.month,1);end=(date(today.year+1,1,1) if today.month==12 else date(today.year,today.month+1,1))
    appts=CustomerAppointment.query.filter(CustomerAppointment.status!='RASCUNHO').order_by(CustomerAppointment.created_at.desc()).all()
    ids=[a.id for a in appts];counts=dict(db.session.query(CustomerAppointmentEquipment.appointment_id,func.count(CustomerAppointmentEquipment.id)).filter(CustomerAppointmentEquipment.appointment_id.in_(ids or [-1])).group_by(CustomerAppointmentEquipment.appointment_id).all());received_counts=dict(db.session.query(CustomerAppointmentEquipment.appointment_id,func.count(CustomerAppointmentEquipment.id)).filter(CustomerAppointmentEquipment.appointment_id.in_(ids or [-1]),CustomerAppointmentEquipment.received.is_(True)).group_by(CustomerAppointmentEquipment.appointment_id).all())
    rows=[];daily={}
    for a in appts:
        d=a.programmed_date or a.expected_date or a.scheduled_date
        item={'id':a.id,'code':a.code,'company':a.customer_company,'request_date':(a.request_date or a.created_at.date()).isoformat(),'expected_date':a.expected_date.isoformat() if a.expected_date else '', 'programmed_date':a.programmed_date.isoformat() if a.programmed_date else '', 'date':d.isoformat() if d else '', 'status':a.status,'count':int(counts.get(a.id,0)),'received':int(received_counts.get(a.id,0)),'alternate':bool(a.alternate_date_requested),'alternate_reason':a.alternate_reason or '', 'cancellation_reason':a.cancellation_reason or ''}
        rows.append(item)
        if d and start<=d<end and str(a.status or '').upper()!='CANCELADO':
            k=d.isoformat();bucket=daily.setdefault(k,{'appointments':0,'equipment':0,'garages':set()});bucket['appointments']+=1;bucket['equipment']+=item['count'];bucket['garages'].add(a.customer_company)
    daily={k:{'appointments':v['appointments'],'equipment':v['equipment'],'garages':len(v['garages'])} for k,v in daily.items()}
    requested_today=[x for x in rows if x['request_date']==today.isoformat() and str(x.get('status') or '').upper()!='CANCELADO']
    programmed_today=[x for x in rows if x['date']==today.isoformat() and str(x.get('status') or '').upper()!='CANCELADO']
    routes=[{'id':r.id,'garage':r.garage_name,'weekday':r.weekday,'weekday_label':_V71_WEEKDAY_LABELS.get(r.weekday,''),'contact':r.contact_name or '', 'address':r.address or '', 'region':r.region or '', 'active':r.active} for r in LogisticsGarageRoute.query.order_by(LogisticsGarageRoute.weekday,LogisticsGarageRoute.garage_name).all()]
    blocked=[{'id':x.id,'date':x.blocked_date.isoformat(),'description':x.description or ''} for x in LogisticsBlockedDate.query.filter_by(active=True).order_by(LogisticsBlockedDate.blocked_date).all()]
    return jsonify({'ok':True,'today':today.isoformat(),'rows':rows,'daily':daily,'requested_today':requested_today,'programmed_today':programmed_today,'routes':routes,'blocked':blocked,'can_manage_appointments':session.get('role')=='manager'})

@app.patch('/api/portal/appointments/<int:aid>/programacao')
@login_required
def v71_program_appointment(aid):
    if not _portal_internal():abort(403)
    a=db.session.get(CustomerAppointment,aid)
    if not a:abort(404)
    data=request.get_json(silent=True) or {};raw=data.get('programmed_date')
    try:d=date.fromisoformat(raw) if raw else None
    except Exception:return jsonify({'ok':False,'error':'Data inválida.'}),400
    if d and LogisticsBlockedDate.query.filter_by(blocked_date=d,active=True).first():return jsonify({'ok':False,'error':'A data está bloqueada no calendário operacional.'}),409
    a.programmed_date=d;a.scheduled_date=d or a.expected_date
    if a.status not in ('RECEBIDO','RECEBIMENTO_PARCIAL','CONCLUIDO'):a.status='NA_PROGRAMACAO'
    db.session.add(AuditEvent(user_id=session['user_id'],event_type='V71_APPOINTMENT_PROGRAMMED',entity_type='customer_appointment',entity_id=str(a.id),detail=f'{a.code} · {d.isoformat() if d else "data padrão"}'))
    db.session.commit();return jsonify({'ok':True,'date':(a.programmed_date or a.expected_date).isoformat() if (a.programmed_date or a.expected_date) else ''})

@app.post('/api/portal/appointments/<int:aid>/cancel')
@login_required
def v712_cancel_appointment(aid):
    if not _portal_internal() or session.get('role')!='manager': abort(403)
    a=db.session.get(CustomerAppointment,aid)
    if not a: abort(404)
    if str(a.status or '').upper()=='CANCELADO': return jsonify({'ok':True,'status':'CANCELADO'})
    data=request.get_json(silent=True) or {}; reason=(data.get('reason') or '').strip()
    if not reason: return jsonify({'ok':False,'error':'Informe o motivo do cancelamento.'}),400
    a.status='CANCELADO'; a.cancelled_at=datetime.utcnow(); a.cancelled_by=session.get('user_id'); a.cancellation_reason=reason
    db.session.add(AuditEvent(user_id=session.get('user_id'),event_type='V71_2_APPOINTMENT_CANCELLED',entity_type='customer_appointment',entity_id=str(a.id),detail=f'{a.code} · {reason[:350]}'))
    db.session.commit(); return jsonify({'ok':True,'status':'CANCELADO'})

@app.delete('/api/portal/appointments/<int:aid>')
@login_required
def v712_delete_appointment(aid):
    if not _portal_internal() or session.get('role')!='manager': abort(403)
    a=db.session.get(CustomerAppointment,aid)
    if not a: abort(404)
    received=CustomerAppointmentEquipment.query.filter_by(appointment_id=a.id,received=True).count()
    if received: return jsonify({'ok':False,'error':'Agendamento com equipamento recebido não pode ser excluído. Use cancelamento/encerramento administrativo.'}),409
    code=a.code or str(a.id); company=a.customer_company or ''
    # Exclui os itens explicitamente para manter compatibilidade com bancos legados sem CASCADE físico.
    CustomerAppointmentEquipment.query.filter_by(appointment_id=a.id).delete(synchronize_session=False)
    db.session.delete(a)
    db.session.flush()
    db.session.add(AuditEvent(user_id=session.get('user_id'),event_type='V71_2_APPOINTMENT_DELETED',entity_type='customer_appointment',entity_id=str(aid),detail=f'{code} · {company}'))
    db.session.commit(); return jsonify({'ok':True})

@app.get('/portal-cliente')
@login_required
def portal_cliente_page():
    if not (_has_access('portal.appointments') or _portal_internal()): abort(403)
    u=db.session.get(User,session['user_id']); return render_template('customer_portal.html',app_release=APP_RELEASE,internal=_portal_internal(),customer_companies=_customer_companies_for_user(u) if u and u.role=='customer' else CustomerCompany.query.filter_by(active=True).order_by(CustomerCompany.legal_name).all())

@app.get('/api/portal/appointments')
@login_required
def portal_appointments_list():
    if not (_has_access('portal.appointments') or _portal_internal()): abort(403)
    q=CustomerAppointment.query
    if not _portal_internal():
        u=db.session.get(User,session['user_id']); allowed=[c.legal_name for c in _customer_companies_for_user(u)]; q=q.filter(CustomerAppointment.customer_company.in_(allowed or ['__SEM_EMPRESA__']))
    rows=q.order_by(CustomerAppointment.created_at.desc()).limit(500).all(); out=[]
    for a in rows:
        items=CustomerAppointmentEquipment.query.filter_by(appointment_id=a.id).all()
        out.append({'id':a.id,'code':a.code,'company':a.customer_company,'responsible':a.responsible_name,'scheduled_date':(a.programmed_date or a.expected_date or a.scheduled_date).isoformat() if (a.programmed_date or a.expected_date or a.scheduled_date) else '', 'request_date':(a.request_date or a.created_at.date()).isoformat(),'expected_date':a.expected_date.isoformat() if a.expected_date else '', 'programmed_date':a.programmed_date.isoformat() if a.programmed_date else '', 'alternate':bool(a.alternate_date_requested),'alternate_reason':a.alternate_reason or '', 'status':a.status,'created_at':a.created_at.isoformat(),'count':len(items),'received':sum(1 for x in items if x.received),'email_status':a.email_status or ''})
    return jsonify({'ok':True,'rows':out,'internal':_portal_internal()})

@app.post('/api/portal/appointments')
@login_required
def portal_appointments_create():
    if not _has_access('portal.appointments'): abort(403)
    u=db.session.get(User,session['user_id']); data=request.get_json(silent=True) or {}; items=data.get('items') or []
    if not items: return jsonify({'ok':False,'error':'Adicione pelo menos um equipamento.'}),400
    allowed=_customer_companies_for_user(u); requested=(data.get('company') or '').strip(); company=(requested if requested and requested in [c.legal_name for c in allowed] else (allowed[0].legal_name if allowed else (u.company or '').strip())); responsible=(data.get('responsible') or u.name or '').strip()
    if not company:return jsonify({'ok':False,'error':'Usuário sem cliente/empresa vinculada.'}),400
    today=datetime.now(ZoneInfo('America/Sao_Paulo')).date();expected,_route=_v71_next_operational_date(company,today)
    alt=bool(data.get('alternate_date_requested'));alt_reason=(data.get('alternate_reason') or '').strip()
    a=CustomerAppointment(customer_company=company,responsible_name=responsible,responsible_email=(data.get('email') or u.email or '').strip(),responsible_phone=(data.get('phone') or u.phone or '').strip(),scheduled_date=expected,request_date=today,expected_date=expected,programmed_date=None,alternate_date_requested=alt,alternate_reason=alt_reason or None,notes=(data.get('notes') or '').strip(),status='RASCUNHO',created_by=u.id); db.session.add(a);db.session.flush();a.code=_portal_code(a.id)
    for i,x in enumerate(items,1):
        serial=(x.get('serial') or '').strip(); defect=(x.get('defect') or '').strip()
        if not serial or not defect: db.session.rollback(); return jsonify({'ok':False,'error':f'Equipamento {i}: série e defeito são obrigatórios.'}),400
        db.session.add(CustomerAppointmentEquipment(appointment_id=a.id,item_no=i,serial_number=serial,equipment=(x.get('equipment') or '').strip(),version=(x.get('version') or '').strip(),eod=(x.get('eod') or '').strip(),defect=defect,notes=(x.get('notes') or '').strip()))
    db.session.commit(); return jsonify({'ok':True,'id':a.id,'code':a.code})

@app.post('/api/portal/appointments/<int:aid>/attachments')
@login_required
def portal_appointment_attachments(aid):
    a=db.session.get(CustomerAppointment,aid)
    if not a or not _portal_can_see(a) or a.created_by!=session['user_id'] or a.status!='RASCUNHO': abort(404)
    try:
        a.invoice_number=(request.form.get('invoice_number') or '').strip() or None
        inv=request.files.get('invoice')
        if inv and inv.filename:
            a.invoice_file,a.invoice_original_name=_portal_store_upload(inv,f'NF-{a.code}')
        items=CustomerAppointmentEquipment.query.filter_by(appointment_id=a.id).order_by(CustomerAppointmentEquipment.item_no).all()
        for x in items:
            f=request.files.get(f'photo_{x.item_no}')
            if f and f.filename: x.photo_file,_=_portal_store_upload(f,f'{_portal_equipment_code(a,x)}-foto',True)
        db.session.commit(); return jsonify({'ok':True})
    except ValueError as exc:
        db.session.rollback(); return jsonify({'ok':False,'error':str(exc)}),400

@app.get('/api/portal/equipments/<int:eid>/photo')
@login_required
def portal_equipment_photo(eid):
    x=db.session.get(CustomerAppointmentEquipment,eid); a=db.session.get(CustomerAppointment,x.appointment_id) if x else None
    if not a or not _portal_can_see(a) or not x.photo_file: abort(404)
    return Response(_portal_stored_bytes(x.photo_file),mimetype='image/jpeg')

@app.get('/api/portal/appointments/<int:aid>/invoice')
@login_required
def portal_appointment_invoice(aid):
    a=db.session.get(CustomerAppointment,aid)
    if not a or not _portal_can_see(a) or not getattr(a,'invoice_file',None): abort(404)
    raw=_portal_stored_bytes(a.invoice_file); name=a.invoice_original_name or f'{a.code}-documento'
    mime='application/pdf' if name.lower().endswith('.pdf') else ('image/png' if name.lower().endswith('.png') else 'image/jpeg')
    return Response(raw,mimetype=mime,headers={'Content-Disposition':f'inline; filename="{secure_filename(name)}"'})

@app.post('/api/portal/appointments/<int:aid>/submit')
@login_required
def portal_appointments_submit(aid):
    a=db.session.get(CustomerAppointment,aid)
    if not a or not _portal_can_see(a) or a.created_by!=session['user_id']: abort(404)
    if a.status!='RASCUNHO': return jsonify({'ok':False,'error':'Agendamento já finalizado.'}),409
    data=request.get_json(silent=True) or {}; sig=data.get('signature_data'); name=(data.get('accepted_name') or a.responsible_name).strip()
    # V69.2.1: assinatura eletrônica permanece disponível, porém não bloqueia a finalização.
    a.signature_file=_portal_store_dataurl(sig,f'assinatura-{a.code}') if sig else None
    a.accepted_name=name;a.accepted_at=datetime.utcnow();a.request_date=a.request_date or datetime.now(ZoneInfo('America/Sao_Paulo')).date();a.expected_date=a.expected_date or _v71_next_operational_date(a.customer_company,a.request_date)[0];a.scheduled_date=a.programmed_date or a.expected_date;a.status='NA_PROGRAMACAO'
    pdf=_portal_pdf(a); key=f"portal-cliente/{datetime.utcnow().strftime('%Y/%m')}/{a.code}.pdf"
    try:_r2_put_bytes(key,pdf,'application/pdf');a.pdf_file='r2__'+key
    except Exception:
        namef=a.code+'.pdf';(UPLOAD_DIR/namef).write_bytes(pdf);a.pdf_file=namef
    db.session.commit(); st,detail=_portal_send_email(a);a.email_status=st;a.email_detail=detail;db.session.commit(); return jsonify({'ok':True,'code':a.code,'email_status':st,'request_date':a.request_date.isoformat() if a.request_date else '', 'expected_date':a.expected_date.isoformat() if a.expected_date else '', 'status':a.status})

@app.get('/api/portal/appointments/<int:aid>')
@login_required
def portal_appointment_detail(aid):
    a=db.session.get(CustomerAppointment,aid)
    if not a or not _portal_can_see(a):abort(404)
    items=CustomerAppointmentEquipment.query.filter_by(appointment_id=a.id).order_by(CustomerAppointmentEquipment.item_no).all()
    download_events=AuditEvent.query.filter(AuditEvent.entity_type=='customer_appointment_equipment',AuditEvent.entity_id.in_([str(x.id) for x in items]),AuditEvent.event_type.in_(['PORTAL_EQUIPMENT_PDF_DOWNLOADED','PORTAL_EQUIPMENT_PDF_DOWNLOADED_ALL'])).order_by(AuditEvent.created_at.desc()).all() if items else []
    latest_download={}
    for ev in download_events:
        latest_download.setdefault(ev.entity_id,ev)
    return jsonify({'ok':True,'appointment':{'id':a.id,'code':a.code,'company':a.customer_company,'responsible':a.responsible_name,'date':(a.programmed_date or a.expected_date or a.scheduled_date).isoformat() if (a.programmed_date or a.expected_date or a.scheduled_date) else '', 'request_date':(a.request_date or a.created_at.date()).isoformat(),'expected_date':a.expected_date.isoformat() if a.expected_date else '', 'programmed_date':a.programmed_date.isoformat() if a.programmed_date else '', 'alternate':bool(a.alternate_date_requested),'alternate_reason':a.alternate_reason or '', 'notes':a.notes or '', 'status':a.status,'email_status':a.email_status or '', 'invoice_number':getattr(a,'invoice_number',None) or '', 'invoice_name':getattr(a,'invoice_original_name',None) or '', 'has_invoice':bool(getattr(a,'invoice_file',None))},'items':[{'id':x.id,'item_no':x.item_no,'protocol':_portal_equipment_code(a,x),'serial':x.serial_number,'equipment':x.equipment or '', 'version':x.version or '', 'eod':x.eod or '', 'defect':x.defect,'notes':x.notes or '', 'has_photo':bool(x.photo_file), 'photo_url':(f'/api/portal/equipments/{x.id}/photo' if x.photo_file else ''), 'received':x.received,'pdf_downloaded':str(x.id) in latest_download,'pdf_downloaded_at':latest_download[str(x.id)].created_at.replace(tzinfo=ZoneInfo('UTC')).astimezone(ZoneInfo('America/Sao_Paulo')).strftime('%d/%m/%Y %H:%M') if str(x.id) in latest_download else ''} for x in items]})

@app.get('/api/portal/appointments/<int:aid>/pdf')
@login_required
def portal_appointment_pdf(aid):
    a=db.session.get(CustomerAppointment,aid)
    if not a or not _portal_can_see(a):abort(404)
    return Response(_portal_pdf(a),mimetype='application/pdf',headers={'Content-Disposition':f'inline; filename="{a.code}.pdf"'})

def _portal_safe_filename_part(value):
    value=re.sub(r'[^A-Za-z0-9._-]+','-',str(value or '').strip()).strip('-_.')
    return value[:100] or 'SEM-SN'


def _portal_refresh_status(a):
    items=CustomerAppointmentEquipment.query.filter_by(appointment_id=a.id).all()
    if not items: return
    if all(x.received for x in items):
        ids=[str(x.id) for x in items]
        downloaded={r[0] for r in db.session.query(AuditEvent.entity_id).filter(AuditEvent.entity_type=='customer_appointment_equipment',AuditEvent.entity_id.in_(ids),AuditEvent.event_type.in_(['PORTAL_EQUIPMENT_PDF_DOWNLOADED','PORTAL_EQUIPMENT_PDF_DOWNLOADED_ALL'])).distinct().all()}
        a.status='CONCLUIDO' if len(downloaded)==len(ids) else 'RECEBIDO'

def _portal_equipment_pdf_filename(a,x):
    return f"{_portal_equipment_code(a,x)}_SN-{_portal_safe_filename_part(x.serial_number)}.pdf"

def _portal_record_pdf_download(a,x,event_type='PORTAL_EQUIPMENT_PDF_DOWNLOADED'):
    db.session.add(AuditEvent(user_id=session.get('user_id'),event_type=event_type,entity_type='customer_appointment_equipment',entity_id=str(x.id),detail=f'{a.code} · série {x.serial_number} · {_portal_equipment_pdf_filename(a,x)}')); db.session.flush(); _portal_refresh_status(a)

@app.get('/api/portal/equipments/<int:eid>/pdf')
@login_required
def portal_equipment_pdf(eid):
    x=db.session.get(CustomerAppointmentEquipment,eid);a=db.session.get(CustomerAppointment,x.appointment_id) if x else None
    if not a or not _portal_can_see(a):abort(404)
    fn=_portal_equipment_pdf_filename(a,x)
    _portal_record_pdf_download(a,x);db.session.commit()
    return Response(_portal_pdf(a,x),mimetype='application/pdf',headers={'Content-Disposition':f'attachment; filename="{fn}"'})

@app.get('/api/portal/appointments/<int:aid>/pdfs.zip')
@login_required
def portal_appointment_pdfs_zip(aid):
    a=db.session.get(CustomerAppointment,aid)
    if not a or not _portal_can_see(a):abort(404)
    items=CustomerAppointmentEquipment.query.filter_by(appointment_id=a.id).order_by(CustomerAppointmentEquipment.item_no).all()
    if not items:return jsonify({'ok':False,'error':'Agendamento sem equipamentos.'}),404
    out=io.BytesIO()
    with zipfile.ZipFile(out,'w',compression=zipfile.ZIP_DEFLATED) as zf:
        if getattr(a,'invoice_file',None):
            try: zf.writestr('Documento_'+secure_filename(a.invoice_original_name or 'nota-fiscal'),_portal_stored_bytes(a.invoice_file))
            except Exception: app.logger.exception('Falha ao incluir documento do agendamento no ZIP %s',a.code)
        for x in items:
            zf.writestr(_portal_equipment_pdf_filename(a,x),_portal_pdf(a,x))
            _portal_record_pdf_download(a,x,'PORTAL_EQUIPMENT_PDF_DOWNLOADED_ALL')
    db.session.commit();out.seek(0)
    return send_file(out,as_attachment=True,download_name=f'{a.code}_PDFs.zip',mimetype='application/zip')

@app.post('/api/portal/equipments/<int:eid>/receive')
@login_required
def portal_equipment_receive(eid):
    if not _portal_internal():abort(403)
    x=db.session.get(CustomerAppointmentEquipment,eid);a=db.session.get(CustomerAppointment,x.appointment_id) if x else None
    if not a:abort(404)
    x.received=True;x.received_at=datetime.utcnow();x.received_by=session['user_id'];db.session.flush()
    total=CustomerAppointmentEquipment.query.filter_by(appointment_id=a.id).count(); rec=CustomerAppointmentEquipment.query.filter_by(appointment_id=a.id,received=True).count()
    a.status='RECEBIDO' if rec>=total else 'RECEBIMENTO_PARCIAL';a.received_at=datetime.utcnow() if rec>=total else None;a.received_by=session['user_id'] if rec>=total else None
    if rec>=total: _portal_refresh_status(a)
    db.session.add(AuditEvent(user_id=session['user_id'],event_type='PORTAL_EQUIPMENT_RECEIVED',entity_type='customer_appointment',entity_id=str(a.id),detail=f'{a.code} · série {x.serial_number}'));db.session.commit();return jsonify({'ok':True,'status':a.status})

# V70 — Performance & Banco: migrações versionadas, aditivas e idempotentes.
def _apply_v70_migrations():
    migrations=[
      ("V70-001", "Índices de performance e operação", (
        "CREATE INDEX IF NOT EXISTS ix_perf_route_created_status ON performance_metrics (route, created_at, status_code)",
        "CREATE INDEX IF NOT EXISTS ix_topdesk_status_created ON topdesk_tickets (status, created_at)",
        "CREATE INDEX IF NOT EXISTS ix_topdesk_line_station_status ON topdesk_tickets (line_code, station_code, status)",
        "CREATE INDEX IF NOT EXISTS ix_material_req_user_status ON material_requests (user_id, status)",
        "CREATE INDEX IF NOT EXISTS ix_collab_doc_user_status ON collaborator_documents (user_id, status)",
      )),
      ("V70.1-001", "Dashboard Chamados e diagnóstico de memória", ()),
    ]
    try:
        db.metadata.create_all(bind=db.engine, tables=[SchemaMigration.__table__], checkfirst=True)
        for version,description,commands in migrations:
            if SchemaMigration.query.filter_by(version=version).first():
                continue
            with db.engine.begin() as conn:
                for sql in commands:
                    try: conn.execute(text(sql))
                    except Exception as exc: app.logger.warning("%s índice não aplicado: %s",version,exc)
            db.session.add(SchemaMigration(version=version,description=description)); db.session.commit()
            app.logger.info("Migração %s aplicada: %s",version,description)
    except Exception:
        try: db.session.rollback()
        except Exception: pass
        app.logger.exception("V70: falha no controle de migrações")


# V71 — schema aditivo da logística e programação de agendamentos.
def _apply_v71_migrations():
    try:
        db.metadata.create_all(bind=db.engine,tables=[LogisticsGarageRoute.__table__,LogisticsBlockedDate.__table__],checkfirst=True)
        insp=db.inspect(db.engine)
        if insp.has_table('customer_appointments'):
            cols={c['name'] for c in insp.get_columns('customer_appointments')}
            additions=(
                ('request_date','DATE'),('expected_date','DATE'),('programmed_date','DATE'),
                ('alternate_date_requested','BOOLEAN DEFAULT FALSE'),('alternate_reason','VARCHAR(300)'),
                # Compatibilidade com bancos legados do Portal: o ORM já possui estes campos,
                # mas instalações anteriores podem não ter recebido as colunas físicas.
                ('invoice_number','VARCHAR(120)'),('invoice_file','VARCHAR(600)'),('invoice_original_name','VARCHAR(255)'),
                ('cancelled_at','TIMESTAMP'),('cancelled_by','INTEGER'),('cancellation_reason','VARCHAR(500)'),
            )
            with db.engine.begin() as conn:
                for col,typ in additions:
                    if col not in cols:conn.execute(text(f'ALTER TABLE customer_appointments ADD COLUMN {col} {typ}'))
                conn.execute(text('CREATE INDEX IF NOT EXISTS ix_customer_appt_request_date ON customer_appointments (request_date)'))
                conn.execute(text('CREATE INDEX IF NOT EXISTS ix_customer_appt_expected_date ON customer_appointments (expected_date)'))
                conn.execute(text('CREATE INDEX IF NOT EXISTS ix_customer_appt_programmed_date ON customer_appointments (programmed_date)'))
        if not SchemaMigration.query.filter_by(version='V71-001').first():
            db.session.add(SchemaMigration(version='V71-001',description='Logística Leva e Traz + programação de agendamentos'));db.session.commit()
        if insp.has_table('customer_companies'):
            # V71.1 HF2 — contatos da planilha podem conter ramal e mais de um telefone.
            with db.engine.begin() as conn:
                conn.execute(text('ALTER TABLE customer_companies ALTER COLUMN phone TYPE VARCHAR(120)'))
                conn.execute(text('ALTER TABLE customer_companies ALTER COLUMN mobile TYPE VARCHAR(120)'))
        if not SchemaMigration.query.filter_by(version='V71-HF1').first():
            db.session.add(SchemaMigration(version='V71-HF1',description='Compatibilidade fiscal do Portal + importação idempotente Leva e Traz'));db.session.commit()
        if not SchemaMigration.query.filter_by(version='V71.1-HF2').first():
            db.session.add(SchemaMigration(version='V71.1-HF2',description='Dashboard 2.0 Portal + ampliação de telefone de clientes'));db.session.commit()
        # V71.1 HF4 — feriados oficiais de 2026 para a operação em São Paulo.
        # A carga é executada uma única vez; depois o ADM continua podendo remover
        # ou acrescentar bloqueios manualmente pelo calendário operacional.
        if not SchemaMigration.query.filter_by(version='V71.1-HF4').first():
            holidays_2026=(
                ('2026-01-01','Confraternização Universal'),
                ('2026-01-25','Aniversário da Cidade de São Paulo'),
                ('2026-04-03','Paixão de Cristo'),
                ('2026-04-21','Tiradentes'),
                ('2026-05-01','Dia do Trabalho'),
                ('2026-06-04','Corpus Christi'),
                ('2026-07-09','Revolução Constitucionalista / Data Magna de SP'),
                ('2026-09-07','Independência do Brasil'),
                ('2026-10-12','Nossa Senhora Aparecida'),
                ('2026-11-02','Finados'),
                ('2026-11-15','Proclamação da República'),
                ('2026-11-20','Dia da Consciência Negra'),
                ('2026-12-25','Natal'),
            )
            for raw_date,description in holidays_2026:
                holiday=date.fromisoformat(raw_date)
                row=LogisticsBlockedDate.query.filter_by(blocked_date=holiday).first()
                if row is None:
                    row=LogisticsBlockedDate(blocked_date=holiday,description=description,active=True,created_by=None)
                else:
                    row.description=description
                    row.active=True
                db.session.add(row)
            db.session.add(SchemaMigration(version='V71.1-HF4',description='Bloqueio padrão dos feriados oficiais de São Paulo em 2026'))
            db.session.commit()
        if not SchemaMigration.query.filter_by(version='V71.2-001').first():
            db.session.add(SchemaMigration(version='V71.2-001',description='Dashboard Inventário unificado + cancelamento administrativo de agendamentos'))
            db.session.commit()
    except Exception:
        try:db.session.rollback()
        except Exception:pass
        app.logger.exception('V71: falha na migração logística')

# V56-B — índices aditivos para leituras críticas.
try:
    with db.engine.begin() as conn:
        for sql in (
            "CREATE INDEX IF NOT EXISTS ix_fin_collection_terminal_end ON financial_cash_collections (terminal, end_at)",
            "CREATE INDEX IF NOT EXISTS ix_techpos_user_captured ON technician_positions (user_id, captured_at)",
            "CREATE INDEX IF NOT EXISTS ix_session_user_created ON session_events (user_id, created_at)",
            "CREATE INDEX IF NOT EXISTS ix_monthly_cost_center_comp ON financial_monthly_costs (cost_center, competence)",
            "CREATE INDEX IF NOT EXISTS ix_perf_created_route ON performance_metrics (created_at, route)",
            "CREATE INDEX IF NOT EXISTS ix_session_user_event_created ON session_events (user_id, event_type, created_at)",
            "CREATE INDEX IF NOT EXISTS ix_team_profile_active_user ON team_schedule_profiles (active, user_id)",
            "CREATE INDEX IF NOT EXISTS ix_fin_atm_tx_imported_at ON financial_atm_transactions (imported_at)",
            "CREATE INDEX IF NOT EXISTS ix_fin_atm_tx_terminal_status_at ON financial_atm_transactions (terminal, status, transaction_at)"
        ):
            try: conn.execute(text(sql))
            except Exception: pass
except Exception:
    pass

with app.app_context():
    # V56-B REV: migração aditiva da telemetria detalhada.
    try:
        insp=db.inspect(db.engine)
        if insp.has_table("performance_metrics"):
            cols={c["name"] for c in insp.get_columns("performance_metrics")}
            with db.engine.begin() as conn:
                if "sql_ms" not in cols: conn.execute(text("ALTER TABLE performance_metrics ADD COLUMN sql_ms FLOAT DEFAULT 0"))
                if "query_count" not in cols: conn.execute(text("ALTER TABLE performance_metrics ADD COLUMN query_count INTEGER DEFAULT 0"))
    except Exception:
        try: db.session.rollback()
        except Exception: pass

    # V69.2.1 HOTFIX1 — cadastro de clientes e vínculos multiempresa.
    try:
        insp=db.inspect(db.engine)
        if insp.has_table('users') and 'customer_company_ids' not in {c['name'] for c in insp.get_columns('users')}:
            with db.engine.begin() as conn: conn.execute(text("ALTER TABLE users ADD COLUMN customer_company_ids TEXT"))
    except Exception: app.logger.exception('Falha migração vínculo Cliente x Empresa')
    migrate_location_reference_columns()
    migrate_panorama_status_column()
    # V39.7.1: não deixa a criação das novas tabelas de Troca de Chips bloquear o startup.
    core_tables=[t for t in db.metadata.sorted_tables if t.name not in ("chip_swaps","chip_swap_photos")]
    db.metadata.create_all(bind=db.engine, tables=core_tables, checkfirst=True)
    _apply_v70_migrations()
    _apply_v71_migrations()
    # V70.1 — registra uma única vez o Dashboard Chamados como dashboard nativa visível.
    try:
        seed_done=SchemaMigration.query.filter_by(version='V70.1-002').first()
        if not seed_done:
            topdesk_cfg=BuiltinDashboardSetting.query.filter_by(dashboard_key='topdesk').first()
            if not topdesk_cfg:
                topdesk_cfg=BuiltinDashboardSetting(dashboard_key='topdesk',visible=True,order_index=55,allowed_roles_json='[]')
                db.session.add(topdesk_cfg)
            else:
                topdesk_cfg.visible=True
                if not topdesk_cfg.order_index: topdesk_cfg.order_index=55
            db.session.add(SchemaMigration(version='V70.1-002',description='Dashboard Chamados habilitado na Central'))
            db.session.commit()
    except Exception:
        try: db.session.rollback()
        except Exception: pass
        app.logger.exception('V70.1: falha ao registrar Dashboard Chamados')
    # V69.3.2 — compatibilidade PostgreSQL do Dossiê/Documentos.
    # HOTFIX2 adicionou estes campos ao ORM, mas o banco legado podia não recebê-los.
    try:
        insp=db.inspect(db.engine)
        if insp.has_table('collaborator_documents'):
            cols={c['name'] for c in insp.get_columns('collaborator_documents')}
            missing=[
                ('invoice_number','VARCHAR(120)'),
                ('invoice_file','VARCHAR(600)'),
                ('invoice_original_name','VARCHAR(255)'),
            ]
            with db.engine.begin() as conn:
                for col,typ in missing:
                    if col not in cols:
                        conn.execute(text(f'ALTER TABLE collaborator_documents ADD COLUMN {col} {typ}'))
            # Reinspeciona e registra eventual divergência para diagnóstico de deploy.
            cols_after={c['name'] for c in db.inspect(db.engine).get_columns('collaborator_documents')}
            unresolved=[col for col,_ in missing if col not in cols_after]
            if unresolved: app.logger.error('V69.3.2: colunas do Dossiê ainda ausentes: %s', unresolved)
            else: app.logger.info('V69.3.2: schema collaborator_documents validado.')
    except Exception:
        app.logger.exception('V69.3.2: falha na migração aditiva collaborator_documents')
    # V68 REV1 — modo de quantidade por item (inteiro/decimal), migração aditiva/idempotente.
    try:
        insp=db.inspect(db.engine)
        if insp.has_table('material_catalog_items'):
            cols={c['name'] for c in insp.get_columns('material_catalog_items')}
            if 'quantity_mode' not in cols:
                with db.engine.begin() as conn:
                    conn.execute(text("ALTER TABLE material_catalog_items ADD COLUMN quantity_mode VARCHAR(20) DEFAULT 'INTEIRO'"))
                    conn.execute(text("UPDATE material_catalog_items SET quantity_mode = CASE WHEN UPPER(COALESCE(category,''))='CONSUMIVEL' THEN 'DECIMAL' ELSE 'INTEIRO' END WHERE quantity_mode IS NULL"))
    except Exception:
        app.logger.exception('Falha na migração V68 REV1 quantity_mode')
    # V68 REV2 — dupla assinatura na devolução (colaborador + responsável).
    try:
        insp=db.inspect(db.engine)
        if insp.has_table('collaborator_documents'):
            cols={c['name'] for c in insp.get_columns('collaborator_documents')}
            with db.engine.begin() as conn:
                if 'return_receiver_signature_file' not in cols: conn.execute(text("ALTER TABLE collaborator_documents ADD COLUMN return_receiver_signature_file VARCHAR(600)"))
                if 'return_receiver_id' not in cols: conn.execute(text("ALTER TABLE collaborator_documents ADD COLUMN return_receiver_id INTEGER REFERENCES users(id)"))
                if 'return_received_at' not in cols: conn.execute(text("ALTER TABLE collaborator_documents ADD COLUMN return_received_at TIMESTAMP"))
    except Exception:
        app.logger.exception('Falha na migração V68 REV2 devolução assinada')
    migrate_v56a_topdesk_dimensions()
    migrate_team_schedule_columns()
    migrate_inventory_sync_uuid()
    migrate_user_archive_column()
    migrate_user_v23_columns()
    migrate_user_gps_required_column()
    migrate_financial_v524_columns()
    migrate_v56a3_visit_contacts()
    migrate_base_asset_columns()
    migrate_inventory_validator_columns()
    migrate_v421_columns()
    seed_data()
    try: seed_v67_materials()
    except Exception:
        db.session.rollback(); app.logger.exception("Falha ao semear catálogo V67")
    try: _seed_garage_chip_base()
    except Exception:
        db.session.rollback(); app.logger.exception("Falha ao semear base Troca de Chips Garagem")
    sync_base_assets_1408(force=False)
    sync_atm_complement_v424()
    cleanup_v352_test_reference()



# -----------------------------------------------------------------------------
# V71.3 - Pendências Operacionais consolidadas
# Base prevista - atividade concluída = pendência real.
# -----------------------------------------------------------------------------
def _v713_pending_rows(module="todos"):
    module=(module or "todos").strip().lower()
    out=[]
    if module in ("todos","recarga"):
        for loc in _chip_swap_locations_payload():
            for v in (loc.get("validators") or []):
                st=(v.get("status") or "PENDENTE").upper().replace("CONCLUIDA","CONCLUÍDA")
                if st != "CONCLUÍDA":
                    out.append({"module":"RECARGA","company":loc.get("company") or "","line":loc.get("line") or "","location":loc.get("location") or "","asset":v.get("label") or str(v.get("base_asset_id") or ""),"status":st,"reason":"Não realizado" if st=="PENDENTE" else "Atividade iniciada/incompleta","technician":v.get("technician") or ""})
    if module in ("todos","emv","mv"):
        swaps={str(x.terminal):x for x in EmvChipSwap.query.all()}
        users={u.id:u.name for u in User.query.filter(User.id.in_({x.technician_id for x in swaps.values() if x.technician_id})).all()} if swaps else {}
        for b in _v41_emv_rows():
            terminal=str(b.get("terminal") or ""); sw=swaps.get(terminal)
            st=((sw.status if sw else b.get("_base_status")) or "PENDENTE").upper().replace("CONCLUIDA","CONCLUÍDA")
            if st != "CONCLUÍDA":
                out.append({"module":"MV / EMV TRILHOS","company":b.get("company") or "","line":b.get("line") or "","location":b.get("station") or "","asset":terminal,"status":st,"reason":"Não realizado" if st=="PENDENTE" else "Atividade iniciada/incompleta","technician":users.get(sw.technician_id,"") if sw else ""})
    if module in ("todos","panorama","panoramica"):
        for x in _panorama_payload():
            st=(x.get("status") or "PENDENTE").upper().replace("CONCLUIDA","CONCLUÍDA")
            if st != "CONCLUÍDA":
                out.append({"module":"PANORÂMICA","company":x.get("company") or "","line":x.get("line") or "","location":x.get("location") or "","asset":"Visão panorâmica","status":st,"reason":"Sem registro panorâmico" if st=="PENDENTE" else "Evidências/pontos incompletos","technician":", ".join(x.get("technicians") or [])})
    if module in ("todos","garagem"):
        for x in _garage_payload():
            st=(x.get("status") or "PENDENTE").upper().replace("CONCLUIDA","CONCLUÍDA")
            if st != "CONCLUÍDA":
                out.append({"module":"GARAGEM","company":x.get("company") or "","line":"","location":x.get("company") or "","asset":x.get("terminal") or "","status":st,"reason":"Não realizado" if st=="PENDENTE" else "Atividade iniciada/incompleta","technician":x.get("technician") or ""})
    return out

@app.get('/pendencias-operacionais')
@login_required
def v713_pending_page():
    return render_template('pending_operations.html',app_release=APP_RELEASE)

@app.get('/api/pendencias-operacionais')
@login_required
def v713_pending_api():
    rows=_v713_pending_rows(request.args.get('module') or 'todos')
    return jsonify({'ok':True,'rows':rows,'summary':{'total':len(rows),'pending':sum(1 for x in rows if x['status']=='PENDENTE'),'in_progress':sum(1 for x in rows if x['status']=='EM ANDAMENTO'),'locations':len({(x['module'],x['company'],x['line'],x['location']) for x in rows})}})

@app.get('/api/pendencias-operacionais/export.xlsx')
@login_required
def v713_pending_export():
    rows=_v713_pending_rows(request.args.get('module') or 'todos')
    wb=Workbook(); ws=wb.active; ws.title='Pendências'
    ws.append(['Módulo','Empresa','Linha','Localidade','Equipamento / Ativo','Situação','Motivo','Técnico'])
    for x in rows: ws.append([x['module'],x['company'],x['line'],x['location'],x['asset'],x['status'],x['reason'],x['technician']])
    for cell in ws[1]: cell.font=Font(bold=True)
    ws.freeze_panes='A2'; ws.auto_filter.ref=ws.dimensions
    widths=[22,28,24,30,26,18,32,26]
    for i,w in enumerate(widths,1): ws.column_dimensions[get_column_letter(i)].width=w
    bio=BytesIO(); wb.save(bio); bio.seek(0)
    return send_file(bio,as_attachment=True,download_name=f"pendencias_operacionais_{datetime.now().strftime('%Y%m%d_%H%M')}.xlsx",mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=int(os.environ.get("PORT", "5000")), debug=False)
